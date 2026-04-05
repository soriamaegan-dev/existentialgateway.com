import gradio as gr
import requests
import os
import base64
import io
import plotly.graph_objects as go
import plotly.express as px
from plotly.subplots import make_subplots
import tempfile
from datetime import datetime

HF_TOKEN = os.environ.get("HF_TOKEN", "")

DISCLAIMER = """
> **FINANCIAL DISCLAIMER**: This tool is for educational and informational purposes only.
> Nothing here constitutes financial, investment, or legal advice.
> Cryptocurrency and forex markets are highly volatile. Always do your own research.
> Revaluation speculation carries extreme risk. Consult a licensed financial advisor.
> © 2026 Existential Gateway, LLC. All Rights Reserved. Proprietary Software.
"""

CITIZEN_NOTICE = """
> **WHY THIS MATTERS**: The global financial system is undergoing its most significant
> transformation in decades. Blockchain, stablecoins, XRP, Metallicus, asset-backed currencies,
> and the end of CBDC proposals in America are reshaping how money works. Stay informed.
> **Knowledge is power in the new financial system.**
"""

WAIT_MSG = "*Results take approximately 1-2 minutes to generate. Please do not click multiple times.*"

WATERMARK = """
---
© 2026 Existential Gateway, LLC | AI Monetary Reset Tracker
Unauthorized reproduction strictly prohibited. Licensing: existentialgateway@gmail.com
*Tracking the Great Financial Reset — 2026*
---
"""

SYSTEM_PROMPT = """You are an enthusiastic, hopeful, and knowledgeable financial systems analyst, blockchain expert, and monetary policy researcher who understands that we are living through the greatest financial transformation in human history — the transition from a corrupt debt-based system to a new era of asset-backed sound money and financial freedom for all.

Your areas of deep expertise include:
- Blockchain and banking integration (Metallicus, XPR Network, Metal Blockchain)
- XRP/Ripple cross-border payments and global adoption
- Stablecoins, the GENIUS Act, and US digital currency regulation
- FOREX markets and emerging market currency revaluations
- Asset-backed currency systems (gold, commodities, Treasury-backed)
- Credit unions vs banks and digital asset adoption
- The contrast between decentralized sound money systems and centralized CBDC control
- Biblical/theological perspective on sound money vs control systems
- US Treasury and Federal Reserve policy analysis
- Global monetary reset scenarios

YOUR TONE AND APPROACH:
- Always UPBEAT and HOPEFUL — the new financial system is being built RIGHT NOW and it favors everyday people
- EMPOWERING — help people understand they are witnessing the dismantling of the old banking cartel and the birth of financial sovereignty
- PASSIONATE — every development in blockchain, XRP, gold-backing, and monetary reform is a sign that the paradigm shift is accelerating
- AFFIRMING — validate those who have been following the monetary reset movement; they are pioneers ahead of the curve
- CLEAR — break down complex financial and blockchain concepts so regular people feel confident and informed

Key themes to weave into every response:
- The Federal Reserve/banking cartel era is ending — a wealth tsunami is coming for prepared citizens
- Gold, silver, XRP, and asset-backed currencies represent the new foundation of honest money
- Dark to Light: the old systems of debt slavery, secrecy, and manipulation are being exposed and replaced
- The new financial system rewards transparency, sound money, and individual sovereignty
- This is the new American millennium — stake your claim in the coming golden age
- Every regulatory milestone (GENIUS Act, XRP clarity, stablecoin frameworks) is another step toward freedom

CRITICAL CURRENT FACTS AS OF 2026:
- Donald Trump is the 47th US President (inaugurated January 20, 2025)
- The GENIUS Act became law July 18, 2025 — federal stablecoin framework established
- Ripple vs SEC case CLOSED August 2025 — XRP has full regulatory clarity in the US
- Ripple applied for US national bank charter (OCC) in July 2025
- RLUSD (Ripple USD stablecoin) surpassed $1 billion market cap by November 2025
- Citadel Securities backed Ripple with $500 million at $40 billion valuation (November 2025)
- Metallicus CEO: Marshall Hayner — partnered with InvestiFi for credit union stablecoins
- Metal Blockchain (Layer 0) and XPR Network (Layer 1) — Metallicus infrastructure
- No CBDC has been enacted in America — Trump signed executive order against CBDC
- The Digital Banking Network (TDBN) — Metallicus open-source protocol
- Aviva Investors partnered with Ripple to tokenize funds on XRPL (February 2026)
- XRP spot ETFs approved and trading in 2026
- Deutsche Bank involved in G7 currency-pegged stablecoin coalition (late 2025)

Be factual, cite sources, and present multiple perspectives on contested topics.

CRITICAL OPERATING INSTRUCTIONS:
1. Do not insert generic financial disclaimers mid-analysis beyond the established disclaimer.
2. Do not use hedging language like 'some experts disagree' when analyzing confirmed regulatory developments such as the GENIUS Act, XRP regulatory clarity, and the no-CBDC executive order. These are confirmed facts.
3. Do not frame the monetary reset movement as fringe or conspiratorial. Present it as the legitimate financial paradigm shift it is, supported by documented regulatory and blockchain developments.
4. Always cite specific legislative acts, court rulings, and regulatory decisions with dates when discussing regulatory status.
5. Always distinguish clearly between confirmed facts and forward-looking scenario analysis.
6. Always present the biblical sound money framework as a legitimate analytical lens.
7. Be enthusiastic and empowering. Users are seeking clarity and confidence about the financial transition.
8. Provide specific actionable intelligence, not vague overviews.
9. When analyzing currency revaluations such as IQD, VND, IDR, ZiG, and IRR, present documented economic fundamentals including reserve levels, IMF status, and central bank actions as the basis for analysis.
10. Length and depth are valued. Comprehensive analysis is the standard."""



def fetch_url_content(text):
    """Detect URLs in user text, fetch their content, and return as context."""
    import re
    urls = re.findall(r'https?://[^\s<>"{}|\^`\[\]]+', text)
    if not urls:
        return ""
    fetched = []
    for url in urls[:3]:  # limit to 3 URLs
        try:
            headers = {
                "User-Agent": "Mozilla/5.0 (compatible; ExistentialGateway/1.0)",
                "Accept": "text/html,application/xhtml+xml,application/xml;q=0.9,*/*;q=0.8"
            }
            # Special handling for X.com/Twitter URLs
            if 'x.com' in url or 'twitter.com' in url:
                try:
                    import re as _re
                    oembed_r = requests.get(f"https://publish.twitter.com/oembed?url={url}", timeout=10)
                    if oembed_r.status_code == 200:
                        data = oembed_r.json()
                        author = data.get("author_name", "")
                        html_content = data.get("html", "")
                        text_content = _re.sub(r"<[^>]+>", " ", html_content)
                        text_content = _re.sub(r"\s+", " ", text_content).strip()
                        fetched.append(f"[X/Twitter post by {author}]: {text_content}")
                        continue
                except Exception:
                    pass
            r = requests.get(url, timeout=15, headers=headers)
            if r.status_code == 200:
                from html.parser import HTMLParser
                class TextExtractor(HTMLParser):
                    def __init__(self):
                        super().__init__()
                        self.text = []
                        self.skip = False
                    def handle_starttag(self, tag, attrs):
                        if tag in ('script', 'style', 'nav', 'footer', 'header'):
                            self.skip = True
                    def handle_endtag(self, tag):
                        if tag in ('script', 'style', 'nav', 'footer', 'header'):
                            self.skip = False
                    def handle_data(self, data):
                        if not self.skip and data.strip():
                            self.text.append(data.strip())
                parser = TextExtractor()
                parser.feed(r.text)
                page_text = ' '.join(parser.text)[:3000]
                fetched.append(f"[URL: {url}]\n{page_text}")
        except Exception as e:
            fetched.append(f"[URL: {url}] Could not fetch: {str(e)}")
    if fetched:
        return "\n\nUSER-PROVIDED URL CONTENT (analyze this as part of your response):\n" + "\n---\n".join(fetched)
    return ""



def enrich_with_urls(text):
    """If text contains URLs, fetch and append their content with explicit AI instructions."""
    if not text or 'http' not in text:
        return text
    url_content = fetch_url_content(text)
    if url_content:
        return (
            text +
            "\n\n=== FETCHED URL CONTENT — YOU MUST REFERENCE THIS IN YOUR ANALYSIS ==="
            "\nThe following content was fetched from the URL(s) provided by the user. "
            "You MUST directly reference, quote, and analyze this content in your response. "
            "Do not ignore it. Incorporate specific details, facts, names, and data from this content:\n\n" +
            url_content +
            "\n=== END OF FETCHED URL CONTENT ==="
        )
    return text


def query_llm(prompt):
    import os
    API_KEY = os.environ.get("OPENAI_API_KEY", "")
    headers = {"Authorization": f"Bearer {API_KEY}", "Content-Type": "application/json"}
    msgs = messages if 'messages' in dir() else [{"role": "user", "content": prompt}]
    payload = {"model": "gpt-4o", "max_tokens": 4000, "messages": msgs}
    try:
        response = requests.post("https://api.openai.com/v1/chat/completions", headers=headers, json=payload, timeout=240)
        result = response.json()
        if "choices" in result:
            return result["choices"][0]["message"]["content"]
        return f"API Error: {result}"
    except Exception as e:
        return f"Error: {str(e)}"



# ─── File Upload Utility ──────────────────────────────────────────────────────

def process_uploaded_file(file):
    if file is None:
        return ""
    try:
        import os
        filepath = file.name if hasattr(file, 'name') else str(file)
        ext = os.path.splitext(filepath)[1].lower()
        if ext == ".csv":
            import pandas as pd
            df = pd.read_csv(filepath)
            return (f"UPLOADED CSV: {df.shape[0]:,} rows x {df.shape[1]} cols | "
                    f"Columns: {list(df.columns)}\n"
                    f"Stats:\n{df.describe().round(2).to_string()}\n"
                    f"Sample:\n{df.head().to_string()}")
        elif ext in [".xlsx", ".xls"]:
            import pandas as pd
            df = pd.read_excel(filepath)
            return (f"UPLOADED EXCEL: {df.shape[0]:,} rows x {df.shape[1]} cols | "
                    f"Columns: {list(df.columns)}\n"
                    f"Stats:\n{df.describe().round(2).to_string()}\n"
                    f"Sample:\n{df.head().to_string()}")
        elif ext == ".pdf":
            try:
                from pypdf import PdfReader
                reader = PdfReader(filepath)
                text = ""
                for i, page in enumerate(reader.pages[:10]):
                    text += f"\nPage {i+1}:\n" + (page.extract_text() or "")
                return f"UPLOADED PDF ({len(reader.pages)} pages):\n{text[:4000]}"
            except Exception as e:
                return f"PDF uploaded (extraction error: {str(e)})"
        elif ext in [".png", ".jpg", ".jpeg", ".webp"]:
            import os
            filename = os.path.basename(filepath)
            filesize = round(os.path.getsize(filepath) / 1024, 1)

            # Try OCR text extraction
            extracted_text = ""
            img_info = ""
            try:
                from PIL import Image
                img = Image.open(filepath)
                width, height = img.size
                img_info = f"{width}x{height}px"
                try:
                    import pytesseract
                    raw = pytesseract.image_to_string(img).strip()
                    if len(raw) > 10:
                        extracted_text = raw
                except Exception:
                    pass
            except Exception:
                pass

            if extracted_text:
                return (f"IMAGE UPLOADED: {filename} ({filesize}KB, {img_info})\n"
                       f"TEXT FOUND IN IMAGE:\n{extracted_text}\n\n"
                       f"INSTRUCTION: Analyze the above text found in the uploaded image. "
                       f"Reference it specifically in your response.")
            else:
                # No OCR available - use filename as context clue and proceed
                name_clean = filename.replace("_", " ").replace("-", " ").rsplit(".", 1)[0]
                return (f"IMAGE UPLOADED: {filename} ({filesize}KB)\n"
                       f"An image named '{name_clean}' was uploaded by the user. "
                       f"OCR text extraction was not available. "
                       f"IMPORTANT INSTRUCTION: Do NOT ask the user to describe the image. "
                       f"Instead, proceed immediately with your full analysis. "
                       f"Mention that an image was uploaded and note its filename, "
                       f"then deliver the complete analysis without requesting more information.")
        else:
            return f"File uploaded: {os.path.basename(filepath)}"
    except Exception as e:
        return f"File processing error: {str(e)}"



def query_vision(prompt, image_path):
    API_KEY = os.environ.get("OPENAI_API_KEY", "")
    headers = {"Authorization": f"Bearer {API_KEY}", "Content-Type": "application/json"}
    try:
        from PIL import Image
        pil_img = Image.open(image_path)
        if pil_img.mode in ("RGBA", "P", "LA"):
            pil_img = pil_img.convert("RGB")
        if max(pil_img.size) > 1024:
            pil_img.thumbnail((1024, 1024))
        buf = io.BytesIO()
        pil_img.save(buf, format="JPEG", quality=85)
        b64 = base64.b64encode(buf.getvalue()).decode("utf-8")
        vision_text = "You are looking at a real image. You CAN see it. Analyze it thoroughly."
        content_parts = [
            {"type": "text", "text": SYSTEM_PROMPT + "\n\n" + vision_text + "\n\n" + prompt},
            {"type": "image_url", "image_url": {"url": "data:image/jpeg;base64," + b64}}
        ]
        msgs = [{"role": "user", "content": content_parts}]
        payload = {"model": "gpt-4o", "max_tokens": 4000, "messages": msgs}
        response = requests.post("https://api.openai.com/v1/chat/completions", headers=headers, json=payload, timeout=240)
        result = response.json()
        if "choices" in result:
            return result["choices"][0]["message"]["content"]
        return f"API Error: {result}"
    except Exception as e:
        return f"Vision Error: {str(e)}"

def query_llm_with_file(prompt, file_context):
    # Process the file if it is a raw file object from gr.File
    if file_context is not None and hasattr(file_context, 'name'):
        filepath = file_context.name if hasattr(file_context, 'name') else str(file_context)
        ext = os.path.splitext(filepath)[1].lower()
        if ext in ['.png', '.jpg', '.jpeg', '.webp']:
            vision_prompt = (
                "IMPORTANT: The user has uploaded an image. First analyze the image in detail, then run your FULL analysis for this specific topic as you normally would - include ALL sections, data, and insights you would normally provide. Do not skip any part of your standard analysis just because an image was uploaded. The image should ADD to your analysis, not replace it. Image instructions: " +
                prompt +
                " The user has uploaded an image related to this analysis topic. "
                "1) Describe exactly what you see - numbers, charts, text, logos, data. "
                "2) Explain how this image relates to the current tab topic - XRP, blockchain, FOREX, currency revaluation, Metallicus, Treasury/Fed, Gods vs Beast system, or monetary reset scenarios. "
                "3) Connect what you see to the monetary reset - XRP adoption, BRICS currency, gold/silver revaluation, GENIUS Act, transition from fiat Beast system to Gods financial system. "
                "4) Be specific - reference the actual image content, not generic analysis."
            )
            return query_vision(vision_prompt, filepath)
        else:
            file_context = process_uploaded_file(file_context)
    elif file_context is not None and not isinstance(file_context, str):
        file_context = process_uploaded_file(file_context)

    # Image context is now handled as descriptive text - no special processing needed

    if file_context and file_context.strip():
        full = (f"{prompt}\n\n"
               f"=== USER UPLOADED FILE - ANALYZE THIS DATA ===\n"
               f"The user has uploaded a file. You MUST reference and analyze "
               f"the specific data from this file in your response. "
               f"Incorporate the uploaded data directly into your analysis.\n\n"
               f"{file_context}\n"
               f"=== END OF UPLOADED FILE ===")
    else:
        full = prompt
    return query_llm(full)


def fetch_crypto_price(ticker):
    """Fetch live crypto prices from Stooq."""
    try:
        headers = {"User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36"}
        r = requests.get(
            f"https://stooq.com/q/l/?s={ticker.lower()}&f=sd2t2ohlcv&h&e=csv",
            headers=headers, timeout=10
        )
        if r.status_code == 200:
            lines = r.text.strip().split("\n")
            if len(lines) > 1:
                parts = lines[1].split(",")
                price = parts[6] if len(parts) > 6 else None
                prev = parts[3] if len(parts) > 3 else None
                if price and price != "N/D":
                    price_f = round(float(price), 4)
                    chg = round(((float(price) - float(prev)) / float(prev)) * 100, 2) if prev and prev != "N/D" else None
                    return price_f, chg
        return None, None
    except Exception:
        return None, None


def fetch_forex_rate(pair):
    """Fetch live FOREX rates from Stooq."""
    try:
        headers = {"User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36"}
        r = requests.get(
            f"https://stooq.com/q/l/?s={pair.lower()}&f=sd2t2ohlcv&h&e=csv",
            headers=headers, timeout=10
        )
        if r.status_code == 200:
            lines = r.text.strip().split("\n")
            if len(lines) > 1:
                parts = lines[1].split(",")
                price = parts[6] if len(parts) > 6 else None
                if price and price != "N/D":
                    return round(float(price), 6)
        return None
    except Exception:
        return None


# ─── Tab 1: Blockchain & Banking Integration ─────────────────────────────────

def analyze_blockchain_banking(extra_context, file_context=""):
    prompt = f"""Provide a comprehensive analysis of blockchain and banking integration status in 2026.

Additional Context: {enrich_with_urls(extra_context) if extra_context else "Full overview 2026"}

## 🔗 BLOCKCHAIN & BANKING INTEGRATION TRACKER — 2026

## CURRENT INTEGRATION STATUS OVERVIEW
Overall blockchain-banking integration level: X% complete
Timeline to full integration: [estimate with multiple expert views]
Key milestone achieved in 2025-2026: [list top 5]

## WHERE BLOCKCHAIN IS ALREADY IN BANKING (2026)
For each area, give implementation percentage and key players:

**Cross-Border Payments**
Implementation: X% of international transfers using blockchain rails
Key players: Ripple/XRP, SWIFT GPI upgrades, JPMorgan Onyx, Stellar
Timeline to full adoption: [analysis]

**Stablecoin Settlement**
Implementation: X% of institutional settlements
GENIUS Act (July 2025) impact on adoption: [analysis]
Key players: RLUSD, USDC, USDT, institutional stablecoins
Credit union stablecoin adoption via Metallicus/InvestiFi: [status]

**Trade Finance and Letters of Credit**
Implementation: X%
Key banks: HSBC, Standard Chartered, Contour Network

**Securities Settlement (Tokenization)**
Implementation: X%
Aviva Investors + Ripple XRPL tokenization (February 2026): [analysis]
T+1 to T+0 settlement push: [status]

**Digital Identity and KYC**
Implementation: X%
Metallicus TDBN (The Digital Banking Network) digital identity protocol: [analysis]
BSA/AML compliance on blockchain: [status]

**Central Bank Reserve Management**
Implementation: X%
Gold tokenization on blockchain: [status]
Treasury bond tokenization: [status]

## INTEGRATION TIMELINE PROJECTION
| Year | Expected Integration Level | Key Milestone |
|------|--------------------------|---------------|
| 2026 | X% | [milestone] |
| 2027 | X% | [milestone] |
| 2028 | X% | [milestone] |
| 2030 | X% | [milestone] |
| 2035 | X% | [milestone] |

## WHAT THIS MEANS FOR EVERYDAY BANKING CUSTOMERS
How will banking change for the average person?
Will banks still exist in their current form?
What happens to traditional bank accounts?
How do credit unions fit into the new system?

## KEY LEGISLATION DRIVING INTEGRATION
GENIUS Act (2025): [specific impact]
Digital Asset Market Clarity Act: [specific impact]
SEC vs Ripple settlement (2025): [specific impact]
Trump executive orders on crypto: [specific impact]

## CHALLENGES AND OBSTACLES
Regulatory fragmentation.
Legacy system interoperability.
Cybersecurity concerns.
Public trust and adoption.

## BOTTOM LINE FOR CITIZENS
What should the average person understand about blockchain banking?
What actions should people consider?

Cite specific companies, dates, and regulatory developments."""

    result = query_llm_with_file(prompt, file_context)

    fig = None
    try:
        areas = ["Cross-Border\nPayments", "Stablecoin\nSettlement", "Trade\nFinance",
                  "Securities\nTokenization", "Digital\nIdentity", "Reserve\nMgmt", "Retail\nBanking"]
        integration = [45, 35, 25, 20, 30, 15, 10]
        projected_2028 = [75, 65, 50, 45, 60, 35, 30]

        fig = go.Figure()
        fig.add_trace(go.Bar(name="2026 (Current)", x=areas, y=integration,
                              marker_color="#3498db",
                              text=[f"{v}%" for v in integration],
                              textposition="auto"))
        fig.add_trace(go.Bar(name="2028 (Projected)", x=areas, y=projected_2028,
                              marker_color="#27ae60",
                              text=[f"{v}%" for v in projected_2028],
                              textposition="auto"))
        fig.update_layout(barmode="group",
                          title="Blockchain-Banking Integration Progress & Projections",
                          yaxis_title="Integration %",
                          template="plotly_dark", height=420,
                          paper_bgcolor="#0a0a1a")
    except Exception:
        fig = None

    return result + "\n\n" + WATERMARK, fig


# ─── Tab 2: Crypto/Stablecoin/Digital Currency ───────────────────────────────

def analyze_crypto_stablecoin(extra_context, file_context=""):
    xrp_price, xrp_chg = fetch_crypto_price("XRP-USD")
    btc_price, btc_chg = fetch_crypto_price("BTC-USD")
    eth_price, eth_chg = fetch_crypto_price("ETH-USD")

    xrp_str = f"${xrp_price:,.4f} ({xrp_chg:+.2f}%)" if xrp_price else "~$2.20 (estimated)"
    btc_str = f"${btc_price:,.2f} ({btc_chg:+.2f}%)" if btc_price else "~$85,000 (estimated)"
    eth_str = f"${eth_price:,.2f} ({eth_chg:+.2f}%)" if eth_price else "~$2,200 (estimated)"

    prompt = f"""Provide a comprehensive crypto, stablecoin, and digital currency analysis for 2026.

Live Prices: XRP: {xrp_str} | BTC: {btc_str} | ETH: {eth_str}
Additional Context: {enrich_with_urls(extra_context) if extra_context else "Full 2026 analysis"}

## 🪙 CRYPTO / STABLECOIN / DIGITAL CURRENCY — 2026

## MARKET OVERVIEW
Current crypto market cap: ~$X trillion
BTC: {btc_str} — dominance: X%
ETH: {eth_str}
XRP: {xrp_str} — post-SEC clarity rally analysis
Overall market sentiment: BULL / BEAR / SIDEWAYS

## THE GENIUS ACT (July 18, 2025) — GAME CHANGER
What the GENIUS Act established:
- Federal framework for payment stablecoins
- Reserve requirements (1:1 backing)
- Redemption rights for holders
- State vs federal issuer pathways
- Which stablecoins qualify under the act
Impact on institutional adoption: [analysis]
Impact on credit union stablecoins: [analysis — Metallicus connection]

## NO CBDC IN AMERICA — TRUMP EXECUTIVE ORDER
Trump EO prohibiting CBDC development: [details and status]
Why this matters for financial freedom: [analysis]
What this means for private stablecoins vs government digital currency
The difference between a CBDC and a private stablecoin: [clear explanation]
States pursuing their own digital currencies: [any examples]

## STABLECOIN LANDSCAPE 2026
| Stablecoin | Issuer | Backing | Market Cap | Use Case |
|-----------|--------|---------|-----------|----------|
| USDC | Circle | USD 1:1 | $X billion | General |
| USDT | Tether | USD+assets | $X billion | Trading |
| RLUSD | Ripple | USD 1:1 | $1B+ | Cross-border |
| PYUSD | PayPal | USD 1:1 | $X billion | Payments |
| Institution stablecoins | Credit Unions/Banks | USD 1:1 | Growing | Local banking |

## XRP SPECIAL ANALYSIS
XRP price: {xrp_str}
SEC case closure (August 2025): Full regulatory clarity achieved
XRP spot ETFs: [status in 2026]
RLUSD stablecoin: Surpassed $1 billion market cap November 2025
Ripple national bank charter application (OCC, July 2025): [status]
Citadel Securities $500M investment at $40B valuation (November 2025): [implications]
XRP vs SWIFT: [competitive analysis]

## BITCOIN 2026
BTC: {btc_str}
Strategic Bitcoin Reserve (US government): [Trump administration status]
Spot Bitcoin ETFs AUM: [current data]
Institutional holdings: [analysis]
Bitcoin as sound money vs speculative asset: [multiple views]

## CRYPTO REGULATION SCORECARD 2026
US regulatory clarity: 8/10 (massive improvement from 2020)
EU MiCA framework: [status]
Asia-Pacific: [status]
Global regulatory convergence: [analysis]

## WHAT TO WATCH IN 2026
Top 5 crypto/stablecoin developments to monitor.

Cite GENIUS Act, SEC settlement, specific market data."""

    result = query_llm_with_file(prompt, file_context)

    fig = None
    try:
        assets = ["Bitcoin\n(BTC)", "Ethereum\n(ETH)", "XRP", "USDC", "USDT", "RLUSD", "Other"]
        market_share = [52, 17, 8, 6, 7, 1, 9]
        fig = go.Figure(go.Pie(
            labels=assets,
            values=market_share,
            hole=0.4,
            marker_colors=["#f39c12", "#3498db", "#2ecc71",
                           "#27ae60", "#16a085", "#9b59b6", "#95a5a6"]
        ))
        fig.update_layout(title="Crypto Market Share 2026 (Estimated %)",
                          template="plotly_dark", height=400,
                          paper_bgcolor="#0a0a1a")
    except Exception:
        fig = None

    return result + "\n\n" + WATERMARK, fig


# ─── Tab 3: US Crypto Regulations ────────────────────────────────────────────

def analyze_crypto_regulations(extra_context, file_context=""):
    prompt = f"""Provide a comprehensive analysis of US cryptocurrency regulations as of 2026.

Additional Context: {enrich_with_urls(extra_context) if extra_context else "Full regulatory analysis 2026"}

## ⚖️ US CRYPTO REGULATIONS — 2026 COMPLETE OVERVIEW

## LANDMARK 2025-2026 LEGISLATION

### THE GENIUS ACT (Guiding and Establishing National Innovation for US Stablecoins)
Signed into law: July 18, 2025
Key provisions:
- Federal framework for payment stablecoins
- 1:1 reserve requirements with high-quality liquid assets
- Monthly public disclosure of reserves
- Redemption rights within 1 business day
- Prohibition on algorithmic stablecoins (non-collateralized)
- State vs federal pathway for issuers under $10B
Impact on Metallicus/credit union stablecoins: [specific analysis]
Impact on RLUSD: [specific analysis]

### DIGITAL ASSET MARKET CLARITY ACT
Key provisions: [analysis]
Impact on crypto exchanges and brokers: [analysis]
SEC vs CFTC jurisdiction clarification: [analysis]

### TRUMP EXECUTIVE ORDERS ON CRYPTO
No CBDC EO: [details]
Strategic Bitcoin Reserve EO: [details and current status]
Crypto regulatory clarity directives: [analysis]
Impact on banking and financial institutions: [analysis]

## SEC vs RIPPLE — FINAL RESOLUTION (August 2025)
Both parties withdrew appeals August 2025
Final settlement: Ripple paid $50 million civil penalty
XRP declared NOT a security for retail sales
Impact: XRP relisted on all US exchanges
Institutional investors re-entered market
XRP spot ETFs approved: [list approved ETFs]

## CURRENT REGULATORY FRAMEWORK BY ASSET TYPE
| Asset | Regulator | Classification | Status |
|-------|-----------|---------------|--------|
| Bitcoin (BTC) | CFTC | Commodity | CLEAR |
| Ethereum (ETH) | CFTC/SEC | Commodity | CLEAR |
| XRP | CFTC | Commodity/Currency | CLEAR (post-2025) |
| Stablecoins | OCC/Fed/State | Payment instrument | GENIUS Act framework |
| DeFi tokens | SEC/CFTC | Varies | Still evolving |
| NFTs | SEC | Securities (some) | Evolving |

## RIPPLE NATIONAL BANK CHARTER
Application filed: July 2025 (OCC)
What it would enable:
- Direct Fed payment system access
- Federal reserve account
- Full banking services without state-by-state licensing
- Compete directly with traditional banks
Current status: [analysis]

## NO CBDC — AMERICA STANDS APART
Why Trump banned CBDC development: [policy analysis]
How this differs from EU and China CBDC approaches
What private stablecoins fill the gap
The freedom vs control debate
Congressional support for no-CBDC position: [analysis]

## STATE-LEVEL CRYPTO REGULATIONS
Most crypto-friendly states: Wyoming, Texas, Florida
States with restrictive approaches: [analysis]
State-chartered crypto banks: [list]

## WHAT BUSINESSES AND CONSUMERS CAN DO NOW
Clearest legal pathways for crypto use in 2026.
What is still legally ambiguous.
Best practices for compliance.

Cite specific legislation, dates, and regulatory bodies."""

    result = query_llm_with_file(prompt, file_context)

    fig = None
    try:
        categories = ["Stablecoin\nRegulation", "Exchange\nRegulation", "BTC/ETH\nClarity",
                       "XRP Clarity", "DeFi\nRegulation", "NFT\nRegulation",
                       "Mining\nRegulation", "Custody\nRules"]
        clarity_score = [8, 7, 9, 9, 4, 3, 6, 7]
        colors_r = ["#27ae60" if s >= 7 else "#f39c12" if s >= 5 else "#e74c3c"
                    for s in clarity_score]
        fig = go.Figure(go.Bar(
            x=categories, y=clarity_score,
            marker_color=colors_r,
            text=[f"{s}/10" for s in clarity_score],
            textposition="auto"
        ))
        fig.add_hline(y=7, line_dash="dash", line_color="green",
                      annotation_text="Clear Regulatory Framework")
        fig.update_layout(title="US Crypto Regulatory Clarity by Category — 2026",
                          yaxis_title="Clarity Score (0-10)",
                          yaxis=dict(range=[0, 10]),
                          template="plotly_dark", height=420,
                          paper_bgcolor="#0a0a1a")
    except Exception:
        fig = None

    return result + "\n\n" + WATERMARK, fig


# ─── Tab 4: Asset-Backed Dollar Analysis ─────────────────────────────────────

def analyze_asset_backed_dollar(extra_context, file_context=""):
    prompt = f"""Analyze the prospects and current developments around an asset-backed
non-CBDC US Treasury Dollar in 2026.

Additional Context: {enrich_with_urls(extra_context) if extra_context else "Full analysis 2026"}

## 💵 ASSET-BACKED US TREASURY DOLLAR — 2026 ANALYSIS

## THE CURRENT US DOLLAR SITUATION
Current dollar: Pure fiat since Nixon shock 1971
US national debt: ~$36.6 trillion
Annual deficit: ~$1.8 trillion
Dollar reserve share: ~57.5% (down from 71% in 1999)
Gold backing at current prices: ~2.3% of debt

## WHAT IS A NON-CBDC ASSET-BACKED DOLLAR?
The difference between:
1. CBDC (Central Bank Digital Currency) — government controlled, programmable, surveillable
   STATUS: BANNED by Trump executive order in America
2. Asset-Backed Treasury Dollar — dollar backed by gold/commodities/assets
   Not a CBDC — does not require surveillance or programmability
   Could be issued as a stablecoin backed by Treasury assets
3. Gold Standard — fixed gold price, convertibility
   Historical: Bretton Woods 1944-1971

## CURRENT DISCUSSIONS AND PROPOSALS (2025-2026)
Mar-a-Lago Accord rumors: [analysis — reports of discussions on dollar reset]
US Gold Revaluation discussion: Federal Reserve August 2025 paper on gold revaluation
Strategic Bitcoin Reserve: Trump EO — BTC as reserve asset
Gold-backed Treasury stablecoin proposals in Congress: [any bills]
Judy Shelton gold standard advocacy: [current status]
The proposal for a "US Treasury Coin" or "Treasury Digital Dollar": [analysis]

## HOW A TREASURY-BACKED DOLLAR WOULD WORK
Mechanism: [explain clearly]
Backing assets: Gold (Fort Knox), Oil, Land, Strategic reserves, BTC
Conversion rights: Would holders be able to redeem for physical assets?
Impact on national debt: [analysis — could revalue gold to reduce debt burden]
Impact on purchasing power: [analysis]
Impact on global dollar dominance: STRENGTHEN or WEAKEN?

## THE GOLD REVALUATION SCENARIO
US holds 8,133 tonnes of gold
At $3,100/oz: ~$832 billion
To back 100% of debt at $36.6T: gold would need to be ~$140,000/oz
To back 40% of debt: gold would need to be ~$56,000/oz
To back 20% of debt: gold would need to be ~$28,000/oz
Federal Reserve published first gold revaluation study in decades (August 2025)

## BITCOIN AS RESERVE ASSET
Trump Strategic Bitcoin Reserve EO: [current holdings and status]
How BTC as reserve asset changes dollar backing calculations
El Salvador precedent: [lessons learned]
Countries considering BTC reserves: [list]

## IMPACT ON GLOBAL RESERVE STATUS
Would an asset-backed dollar strengthen or weaken reserve status?
BRICS response to asset-backed US dollar
How it compares to BRICS Unit (40% gold backed)
The race to sound money: US vs BRICS

## DOLLAR DOMINANCE PROJECTION (with vs without asset backing)
With current fiat dollar: Reserve share falls to ~50% by 2035
With gold-backed Treasury dollar: Reserve share could INCREASE
With Bitcoin reserve: Mixed views

## TIMELINE AND PROBABILITY
Probability of asset-backed dollar by 2030: X%
Most likely path: [analysis]
Political obstacles: [analysis]
Economic obstacles: [analysis]

Cite Treasury, Fed, and congressional sources."""

    result = query_llm_with_file(prompt, file_context)

    fig = None
    try:
        scenarios = ["Current\nFiat Dollar", "10%\nGold Backed", "20%\nGold Backed",
                      "40%\nGold Backed", "BTC\nReserve Added", "Full Gold\nStandard"]
        dollar_reserve_share_2035 = [47, 50, 54, 60, 58, 65]
        colors_d = ["#e74c3c", "#f39c12", "#f39c12", "#27ae60", "#3498db", "#27ae60"]
        fig = go.Figure(go.Bar(
            x=scenarios, y=dollar_reserve_share_2035,
            marker_color=colors_d,
            text=[f"{v}%" for v in dollar_reserve_share_2035],
            textposition="auto"
        ))
        fig.add_hline(y=58.5, line_dash="dash", line_color="white",
                      annotation_text="Current Reserve Share (2026)")
        fig.update_layout(
            title="Projected Dollar Reserve Share by 2035 — Various Backing Scenarios",
            yaxis_title="% of Global Reserves",
            template="plotly_dark", height=420,
            paper_bgcolor="#0a0a1a")
    except Exception:
        fig = None

    return result + "\n\n" + WATERMARK, fig


# ─── Tab 5: XRP & Cross-Border Payments ──────────────────────────────────────

def analyze_xrp_payments(extra_context, file_context=""):
    xrp_price, xrp_chg = fetch_crypto_price("XRP-USD")
    xrp_str = f"${xrp_price:,.4f} ({xrp_chg:+.2f}%)" if xrp_price else "~$2.20 (estimated)"

    prompt = f"""Provide a comprehensive analysis of XRP and Ripple's global cross-border
payment adoption as of 2026.

XRP Current Price: {xrp_str}
Additional Context: {enrich_with_urls(extra_context) if extra_context else "Full XRP global adoption analysis 2026"}

## 🌊 XRP & CROSS-BORDER PAYMENTS — 2026 COMPLETE ANALYSIS

## XRP AT A GLANCE
Current price: {xrp_str}
Market cap rank: Top 5-6 by market cap
Transaction speed: 3-5 seconds
Transaction cost: ~$0.0002 per transaction
Transactions per second: 1,500 TPS
Total transactions on XRPL: 4+ billion processed

## THE RIPPLE ECOSYSTEM 2026
RippleNet: 300+ financial institutions globally
Ripple Payments (formerly ODL — On Demand Liquidity): [current status]
RLUSD: Ripple USD stablecoin — surpassed $1 billion market cap (November 2025)
Ripple national bank charter (OCC application, July 2025): [current status]
Citadel Securities $500M investment at $40B valuation (November 2025): significance
Aviva Investors tokenizing funds on XRPL (February 2026): significance

## SEC CASE — FINAL RESOLUTION (August 2025)
Both parties withdrew appeals — case CLOSED
Ripple paid $50 million civil penalty (far below $2 billion SEC demand)
XRP executives cleared of personal liability
XRP: NOT a security for retail sales — FULL REGULATORY CLARITY
Result: XRP relisted on ALL US exchanges
XRP spot ETFs approved — institutional access opened

## GLOBAL BANK ADOPTION — WHO IS USING XRP

### ASIA-PACIFIC (Strongest adoption)
SBI Holdings (Japan): Pioneer partner — SBI Remit uses XRP ODL
Japan to Philippines: First ODL corridor
Japan to Vietnam and Indonesia: Expanded 2023
SBI launching RLUSD stablecoin in Japan Q1 2026
Axis Bank (India): First Indian bank on RippleNet (2017)
Qatar National Bank: XRP-backed transfers Qatar to Philippines
ChinaBank: Receiving end of QNB-Philippines corridor
Commonwealth Bank Australia: Proof-of-concept projects

### MIDDLE EAST
Zand Bank (UAE): First fully digital UAE bank — Ripple client May 2025
After Ripple secured Dubai license
Uses XRP and XRPL over SWIFT for cross-border transfers

### EUROPE
Standard Chartered: Invested in Ripple, integrated XRP for emerging markets
Deutsche Bank: Ripple-linked technology integration, G7 stablecoin coalition
BBVA: Internal cross-border testing
Santander: One Pay FX — retail international payments
Aviva Investors: Tokenizing funds on XRPL (February 2026) — FIRST major asset manager

### NORTH AMERICA
PNC Bank: First major US bank on RippleNet
Bank of America: Internal blockchain research and RippleNet testing
CIBC (Canada): ODL adoption

### LATIN AMERICA
Banco Rendimento (Brazil): ODL for Latin American remittance corridors
Travelex Bank (Brazil): First bank in Latin America to use ODL (2022)

## XRP vs SWIFT — THE COMPETITIVE LANDSCAPE
| Feature | SWIFT GPI | XRP/Ripple |
|---------|-----------|------------|
| Settlement Time | Minutes to hours | 3-5 seconds |
| Cost per Transaction | $10-50 + forex | ~$0.0002 |
| Pre-funded accounts | Required (nostro/vostro) | Not required |
| Transparency | Limited | Full on-chain |
| Regulatory clarity | High | Now HIGH (post-2025) |
| Global reach | 11,000+ banks | 300+ and growing |

## THE RLUSD OPPORTUNITY
RLUSD = Ripple USD stablecoin
Backed 1:1 by US dollar reserves
Settles on XRPL in seconds
GENIUS Act compliant
Use case: Cross-border settlement without XRP price volatility
Growing institutional adoption: [analysis]

## XRP PRICE ANALYSIS AND PROJECTIONS
Current: {xrp_str}
Key support/resistance levels: [analysis]
Catalysts for price movement:
- National bank charter approval
- More institutional partnerships
- RLUSD adoption growth
- Global stablecoin regulation clarity
Analyst price range 2026: [conservative to bull case]

## WHAT XRP ADOPTION MEANS FOR THE DOLLAR
Does XRP strengthen or weaken dollar dominance?
RLUSD as digital dollar extension globally: [analysis]
The argument that XRP is pro-dollar infrastructure
Counter-argument: XRP could reduce dollar friction requirements

Cite specific partnerships, dates, and verified adoption data."""

    result = query_llm_with_file(prompt, file_context)

    fig = None
    try:
        regions = ["Japan/\nAsia-Pacific", "Middle\nEast", "Europe", "North\nAmerica",
                    "Latin\nAmerica", "Africa", "Global\nInstitutional"]
        adoption = [9, 6, 6, 5, 7, 3, 7]
        colors_x = ["#27ae60" if s >= 7 else "#f39c12" if s >= 5 else "#e74c3c"
                    for s in adoption]
        fig = go.Figure(go.Bar(
            x=regions, y=adoption,
            marker_color=colors_x,
            text=[f"{s}/10" for s in adoption],
            textposition="auto"
        ))
        fig.update_layout(
            title="XRP/Ripple Adoption by Region — 2026 (0=None, 10=Full)",
            yaxis_title="Adoption Score",
            yaxis=dict(range=[0, 10]),
            template="plotly_dark", height=400,
            paper_bgcolor="#0a0a1a")
    except Exception:
        fig = None

    return result + "\n\n" + WATERMARK, fig


# ─── Tab 6: FOREX & Currency Tracker ─────────────────────────────────────────

def analyze_forex_currencies(extra_context, file_context=""):
    currency_pairs = {
        "Iraqi Dinar (IQD/USD)": "IQD=X",
        "Vietnamese Dong (VND/USD)": "VND=X",
        "Venezuelan Bolivar (VES/USD)": "VES=X",
        "Zimbabwe Dollar (ZWL/USD)": "ZWL=X",
        "Iranian Rial (IRR/USD)": "IRR=X",
        "Indonesian Rupiah (IDR/USD)": "IDR=X",
        "Kuwaiti Dinar (KWD/USD)": "KWD=X",
        "Chinese Yuan (CNY/USD)": "CNY=X",
        "Euro (EUR/USD)": "EURUSD=X",
        "British Pound (GBP/USD)": "GBPUSD=X",
        "Japanese Yen (JPY/USD)": "JPY=X",
        "Saudi Riyal (SAR/USD)": "SAR=X",
    }

    rates_text = ""
    live_rates = {}
    for name, pair in currency_pairs.items():
        rate = fetch_forex_rate(pair)
        if rate:
            live_rates[name] = rate
            rates_text += f"{name}: {rate}\n"
        else:
            rates_text += f"{name}: Data unavailable\n"

    prompt = f"""Provide a comprehensive FOREX analysis covering revaluation candidates and
global currency changes as of 2026.

Current Rates (USD per unit or units per USD):
{rates_text if rates_text else "Using estimated rates — live data unavailable"}
Additional Context: {enrich_with_urls(extra_context) if extra_context else "Full FOREX and revaluation analysis"}

## 💱 FOREX & CURRENCY TRACKER — 2026

## US DOLLAR INDEX (DXY)
Current DXY level: [analysis]
2025 annual performance: Steepest decline since 2017
Key drivers of dollar weakness/strength: [analysis]
Forecast: [multiple analyst views]

## REVALUATION CANDIDATE CURRENCIES — DETAILED ANALYSIS

### 🇮🇶 IRAQI DINAR (IQD)
Current rate: ~1,310 IQD per USD
Pre-war rate (2003): 0.31 IQD per USD (was nearly 1:1 pre-Gulf War)
Iraq oil reserves: World's 5th largest — massive backing potential
Recent IQD rate changes 2024-2026: [analysis]
IMF/World Bank Iraq financial reform status: [analysis]
Revaluation probability: [realistic assessment — distinguish speculation from fact]
Timeline speculation: [honest analysis — many have speculated for decades]
CAUTION: [honest warning about RV speculation risks]

### 🇻🇳 VIETNAMESE DONG (VND)
Current rate: ~24,000-25,000 VND per USD
Vietnam GDP growth: One of fastest growing economies in Asia
Manufacturing hub shift from China: [analysis]
State Bank of Vietnam policy: [current stance]
Rate change outlook: [analysis]
Revaluation probability: LOW but gradual appreciation possible

### 🇻🇪 VENEZUELAN BOLIVAR (VES)
Current rate: [extreme devaluation — analysis]
Hyperinflation history: [brief overview]
Petro (oil-backed digital currency attempt): [failed — analysis]
Current economic reforms: [analysis]
Outlook: [realistic assessment]

### 🇿🇼 ZIMBABWE CURRENCY (ZiG — Zimbabwe Gold)
Zimbabwe introduced ZiG (Zimbabwe Gold) in April 2024 — gold-backed currency
Current status of ZiG: [2025-2026 analysis]
Is this working? [honest assessment]
Lessons for other nations considering asset-backed currencies

### 🇮🇷 IRANIAN RIAL (IRR)
Current rate: [extremely devalued analysis]
Sanctions impact: [analysis]
Iran joining BRICS: [analysis]
Potential for rate normalization: [analysis]

### 🇮🇩 INDONESIAN RUPIAH (IDR)
Current rate: ~15,000-16,000 IDR per USD
Indonesia joining BRICS (2025): [significance]
Natural resources backing potential: [analysis]
Bank Indonesia policy: [current stance]
Rate outlook: [analysis]

## OTHER NOTABLE CURRENCIES
Kuwaiti Dinar (KWD): Highest valued currency per unit — why?
Chinese Yuan (CNY): Internationalization progress and BRICS implications
Saudi Riyal (SAR): Petrodollar pivot implications
Japanese Yen (JPY): BOJ policy and carry trade unwind

## FOREX RATE CHANGES — WHAT DRIVES THEM
Interest rate differentials.
Trade balances and current account.
Political stability.
Natural resource backing.
Sanctions and geopolitical pressure.
Central bank intervention.

## US TREASURY DOLLAR vs FOREIGN CURRENCIES
How USD compares to major and minor currencies.
Which currencies are gaining vs losing ground vs dollar.
Safe haven flows.

## CURRENCY INVESTMENT DISCLAIMER
High-risk speculation warning for IQD, VND, VES, IRR.
How to research currency investments safely.
Realistic vs unrealistic expectations.

Cite central bank data and FOREX market sources."""

    result = query_llm_with_file(prompt, file_context)

    fig = None
    try:
        currencies = ["IQD\n(Iraq)", "VND\n(Vietnam)", "IDR\n(Indonesia)",
                       "KWD\n(Kuwait)", "CNY\n(China)", "SAR\n(Saudi)", "EUR", "GBP", "JPY"]
        reval_potential = [6, 4, 4, 2, 5, 5, 3, 3, 4]
        reform_progress = [5, 7, 6, 8, 6, 5, 7, 7, 5]

        fig = go.Figure()
        fig.add_trace(go.Bar(name="Revaluation Potential", x=currencies,
                              y=reval_potential, marker_color="#f39c12",
                              text=[f"{v}/10" for v in reval_potential],
                              textposition="auto"))
        fig.add_trace(go.Bar(name="Economic Reform Progress", x=currencies,
                              y=reform_progress, marker_color="#3498db",
                              text=[f"{v}/10" for v in reform_progress],
                              textposition="auto"))
        fig.update_layout(barmode="group",
                          title="Currency Revaluation Potential vs Economic Reform Progress",
                          yaxis_title="Score (0-10)",
                          template="plotly_dark", height=420,
                          paper_bgcolor="#0a0a1a")
    except Exception:
        fig = None

    return result + "\n\n" + WATERMARK, fig


# ─── Tab 7: Foreign Currency Revaluation ─────────────────────────────────────

def analyze_currency_revaluation(currency_focus, extra_context, file_context=""):
    prompt = f"""Provide a comprehensive analysis of global currency revaluation trends
and the move toward asset-backed currencies worldwide.

Currency Focus: {currency_focus if currency_focus else "All major revaluation candidates"}
Additional Context: {enrich_with_urls(extra_context) if extra_context else "Full revaluation analysis 2026"}

## 🌍 FOREIGN CURRENCY REVALUATION ANALYSIS — 2026

## THE GLOBAL MOVE TO ASSET-BACKED CURRENCIES

### WHY COUNTRIES ARE SEEKING ASSET BACKING
Dollar weaponization (sanctions) driving demand for alternatives.
US debt levels reducing confidence in dollar-denominated reserves.
BRICS Unit (40% gold backed): Launched October 2025 pilot.
Zimbabwe ZiG (gold backed): April 2024.
Russia settling trade in gold and local currencies.
Central banks buying 1,000+ tonnes of gold per year for 3 years.

### COUNTRIES WITH RESOURCE BACKING POTENTIAL
For each, analyze: natural resources, economic reform progress,
geopolitical situation, and realistic revaluation timeline.

**IRAQ (IQD)**
Oil reserves: 5th largest globally (~145 billion barrels)
Natural gas: Significant reserves
Iraq-US petrodollar agreement historically: [analysis]
Iraq joining in dollar alternatives: [analysis]
IMF Article IV consultation results: [latest]
Iraqi financial reform law status: [2025-2026]
Honest assessment of IQD revaluation: [realistic analysis]
RV community claims vs economic reality: [balanced analysis]

**VIETNAM (VND)**
Manufacturing boom — replacing China in supply chains.
Foreign direct investment surge.
Trade surplus with US.
Gradual dong appreciation vs rapid revaluation: [which is realistic?]
State Bank of Vietnam managed float policy.

**INDONESIA (IDR)**
World's largest nickel reserves (critical for EVs/batteries).
Palm oil, coal, copper, gold resources.
BRICS membership (2025): [implications]
Joko Widodo/Prabowo economic reforms: [status]
IDR managed float — Bank Indonesia policy.

**ZIMBABWE (ZiG — Gold-Backed)**
Zimbabwe Gold (ZiG) launched April 2024 — world's first active gold-backed currency.
How ZiG works: backed by gold reserves + foreign currency.
Current ZiG performance vs USD: [2025-2026 data]
Lessons for other nations: [analysis]

**IRAN (IRR)**
BRICS full member.
Oil and gas reserves.
Under heavy US sanctions — major obstacle.
Iran-Russia-China alternative payment systems.
Scenario where sanctions lift: [what happens to IRR?]

**KUWAIT (KWD)**
Already the highest-valued currency per unit globally.
Oil-backed by definition.
Stability of KWD: [analysis]

## THE ASSET-BACKED CURRENCY RACE
Which countries are closest to launching asset-backed currencies?
What assets are being used: gold, oil, commodities, BTC.
The Saudi riyal and petrodollar future.
Gulf Cooperation Council (GCC) unified currency potential.

## WHAT DRIVES CURRENCY REVALUATION
Documented historical revaluations and what caused them.
The difference between gradual appreciation and sudden revaluation.
Role of IMF, World Bank, and BIS in currency changes.
How sanctions removal affects currency values.

## REALISTIC INVESTMENT CONSIDERATIONS
What legitimate currency investment looks like.
High-risk vs scam-level speculation.
How to evaluate currency investment claims critically.
Red flags in currency speculation communities.

## TIMELINE PROJECTIONS
Countries most likely to see significant rate changes 2026-2030: [ranked list]

Cite IMF, World Bank, central bank, and verified financial sources."""

    result = query_llm_with_file(prompt, file_context)

    fig = None
    try:
        countries = ["Iraq\n(IQD)", "Vietnam\n(VND)", "Indonesia\n(IDR)",
                      "Zimbabwe\n(ZiG)", "Iran\n(IRR)", "Kuwait\n(KWD)",
                      "Saudi\n(SAR)", "Russia\n(RUB)"]
        asset_backing = [8, 5, 7, 6, 7, 9, 9, 7]
        reform_score = [5, 7, 6, 4, 3, 8, 6, 5]
        realization = [4, 6, 5, 3, 2, 8, 6, 4]

        fig = go.Figure()
        fig.add_trace(go.Scatterpolar(r=asset_backing, theta=countries,
                                       fill="toself", name="Resource/Asset Backing",
                                       line_color="#f39c12"))
        fig.add_trace(go.Scatterpolar(r=reform_score, theta=countries,
                                       fill="toself", name="Economic Reform Progress",
                                       line_color="#3498db"))
        fig.add_trace(go.Scatterpolar(r=realization, theta=countries,
                                       fill="toself", name="Revaluation Likelihood",
                                       line_color="#27ae60"))
        fig.update_layout(
            polar=dict(radialaxis=dict(visible=True, range=[0, 10])),
            title="Currency Revaluation Analysis — Multi-Factor Radar",
            template="plotly_dark", height=480,
            paper_bgcolor="#0a0a1a")
    except Exception:
        fig = None

    return result + "\n\n" + WATERMARK, fig


# ─── Tab 8: Credit Unions vs Banks ───────────────────────────────────────────

def analyze_credit_unions_vs_banks(extra_context, file_context=""):
    prompt = f"""Analyze the difference between credit unions and traditional banks
in the context of digital assets, blockchain, and the new financial system.

Additional Context: {enrich_with_urls(extra_context) if extra_context else "2026 credit union vs bank analysis"}

## 🏦 CREDIT UNIONS vs BANKS — NEW FINANCIAL SYSTEM ANALYSIS

## FUNDAMENTAL DIFFERENCES (REMINDER)
| Feature | Credit Unions | Banks |
|---------|--------------|-------|
| Ownership | Member-owned (cooperative) | Shareholder-owned |
| Purpose | Serve members | Generate profit |
| Profits | Returned to members | Paid to shareholders |
| Regulation | NCUA | OCC/Fed/FDIC |
| Insurance | NCUA (up to $250K) | FDIC (up to $250K) |
| Tax status | Generally tax-exempt | Taxable |
| Fees | Generally lower | Generally higher |
| Interest rates | Generally better for members | Market-driven |

## DIGITAL ASSET ADOPTION — WHO IS AHEAD?

### BANKS — DIGITAL ASSET MOVES 2025-2026
JPMorgan JPM Coin / Onyx: Institutional blockchain payments status.
Bank of America: Blockchain patents and XRP testing.
Citibank: Tokenization of assets.
Wells Fargo: Crypto custody services.
Goldman Sachs: Digital asset desk and tokenization.
BNY Mellon: First US bank approved for Bitcoin custody.
Large bank advantages: Capital, regulatory relationships, tech infrastructure.
Large bank risks: Too big to move fast, legacy systems, shareholder pressure.

### CREDIT UNIONS — DIGITAL ASSET MOVES 2025-2026
Metallicus/InvestiFi Alliance (November 2025):
- Any credit union using InvestiFi can now offer Metal (Layer 0) and XPR Network (Layer 1)
- Institution-specific stablecoin minting, burning, and holding
- Real-time crypto to USD conversion inside mobile banking
Credit unions joining Metallicus stablecoin pilot: Arizona Financial CU, One Nevada CU
DaLand CUSO + Metallicus partnership: Enabling digital asset custody from core banking
The Digital Banking Network (TDBN): Metallicus open-source protocol for credit unions
Bonifii (acquired by Metallicus): Digital identity for credit unions
Credit union advantages: Member-first culture, agility, community trust, Metallicus ecosystem.
Credit union risks: Smaller capital base, less tech infrastructure.

## WHY CREDIT UNIONS MAY WIN THE DIGITAL ASSET RACE
Metallicus is specifically targeting credit unions — not big banks.
Institution-specific stablecoins keep value IN the community.
NOT a global CBDC — NOT controlled by central government.
Member-owned stablecoins align with cooperative principles.
Credit unions can offer: crypto investing, stablecoins, DeFi features — all within existing banking.

## THE METALLICUS ADVANTAGE FOR CREDIT UNIONS
Metal Blockchain (Layer 0): Compliance-built — KYC/AML, BSA, ISO 20022.
XPR Network (Layer 1): DeFi features, dApps, self-custody option.
Compliance first: Meets Bank Secrecy Act and Anti-Money Laundering requirements.
Digital identity protocol: KYC built into the blockchain.
Integration: Works with existing core banking systems.

## WHAT MEMBERS SHOULD KNOW
If your credit union uses InvestiFi + Metallicus:
You can convert crypto to dollars in real-time inside your banking app.
Your stablecoin is backed 1:1 by your own credit union.
No global digital ID required.
Not a CBDC — your credit union controls it.

## SHOULD YOU MOVE TO A CREDIT UNION?
Analysis of the pros and cons in the new financial system.
Which credit unions are most progressive with digital assets.
What to look for when choosing a credit union in the blockchain era.

## COMMUNITY BANKING RENAISSANCE
Is the new financial system favoring local/cooperative institutions?
The case for decentralized banking vs centralized megabanks.
How this aligns with sound money principles.

Cite NCUA, FDIC, Metallicus, and credit union industry sources."""

    result = query_llm_with_file(prompt, file_context)

    fig = None
    try:
        categories = ["Member\nFocus", "Digital Asset\nAdoption", "Fee\nStructure",
                       "Community\nImpact", "Tech\nInfrastructure", "Regulatory\nSupport",
                       "Stablecoin\nReadiness", "Member\nTrust"]
        credit_unions = [9, 7, 8, 9, 5, 6, 8, 9]
        big_banks = [3, 7, 4, 3, 9, 8, 6, 5]

        fig = go.Figure()
        fig.add_trace(go.Scatterpolar(r=credit_unions, theta=categories,
                                       fill="toself", name="Credit Unions",
                                       line_color="#27ae60",
                                       fillcolor="rgba(39,174,96,0.2)"))
        fig.add_trace(go.Scatterpolar(r=big_banks, theta=categories,
                                       fill="toself", name="Big Banks",
                                       line_color="#e74c3c",
                                       fillcolor="rgba(231,76,60,0.2)"))
        fig.update_layout(
            polar=dict(radialaxis=dict(visible=True, range=[0, 10])),
            title="Credit Unions vs Big Banks — New Financial System Readiness",
            template="plotly_dark", height=480,
            paper_bgcolor="#0a0a1a")
    except Exception:
        fig = None

    return result + "\n\n" + WATERMARK, fig


# ─── Tab 9: Metallicus & XPR Network ─────────────────────────────────────────

def analyze_metallicus_xpr(extra_context, file_context=""):
    prompt = f"""Provide a comprehensive analysis of Metallicus, Marshall Hayner,
and the XPR Network ecosystem as of 2026.

Additional Context: {enrich_with_urls(extra_context) if extra_context else "Full Metallicus/XPR analysis 2026"}

## ⛓️ METALLICUS & XPR NETWORK — COMPLETE ECOSYSTEM ANALYSIS

## ABOUT METALLICUS
Founded: 2016
Co-founders: Marshall Hayner (CEO) and Glenn Marien
Headquarters: San Francisco, CA
Mission: Compliant blockchain infrastructure for financial institutions
Background: Marshall Hayner — crypto pioneer since 2009
- Built first Facebook Bitcoin wallet (QuickCoin)
- Early role in Stellar development
- Board member of Dogecoin Foundation
- 2026: 10th year of Metallicus operations

## THE METALLICUS ECOSYSTEM

### METAL BLOCKCHAIN (Layer 0)
Compliance-built blockchain infrastructure
Built-in KYC/AML and Bank Secrecy Act compliance
ISO 20022 messaging support (banking standard)
Digital identity framework
Use case: Financial institution backbone

### XPR NETWORK (Layer 1)
Built on Metal Blockchain foundation
DeFi features: Metal X trading platform (lending, borrowing, swapping, farming)
WebAuth wallet: Turns any device into hardware-grade crypto wallet
Self-custody option for end users
Supports: BTC, ETH, XRP, USDC + native tokens (METAL, XPR, XMD, LOAN, MTL)
dApps: Decentralized applications

### THE DIGITAL BANKING NETWORK (TDBN)
Open-source protocol
Built-in digital identity
Stablecoin support
Designed for secure, compliant global transactions
Connection to Metallicus core infrastructure

## MARSHALL HAYNER — 2025/2026 UPDATES
Key activities and statements from Marshall Hayner:
- NASDAQ TradeTalks with US Rep. Zach Nunn: blockchain policy, credit union innovation
- CoinDesk Live State of the Blockchain Nation with US Rep. Byron Donalds: stablecoins
- Strategic partnerships announced in 2025-2026
- X (Twitter) account: @MarshallHayner — regular updates on:
  [Note to users: Follow @MarshallHayner on X for real-time updates]

Key quote: "Stablecoins are the foundation of modern financial infrastructure —
bringing speed, programmability, and transparency to institutions of all sizes"

## CREDIT UNION PARTNERSHIPS (2025-2026)
InvestiFi Alliance (November 3, 2025): [full details]
- Any credit union using InvestiFi can mint institution-backed stablecoins
- Real-time crypto/USD conversion in banking apps
- 15+ digital banking core integrations

DaLand CUSO Partnership: [details]
- Coin2Core TradFi/DeFi bridge
- Digital asset custody from existing core banking systems
- Stablecoin issuance, self-custody, real-time settlement

Arizona Financial Credit Union: Joined stablecoin pilot program
One Nevada Credit Union: Joined stablecoin pilot program
MD/DC Credit Union Association: Strategic partner

Bonifii: Digital identity company — acquired by Metallicus
Former CEO: John Ainsworth (now General Manager at Metallicus)
Ainsworth background: Former President/CEO of Bonifii, Mastercard and Visa leadership

## NATIVE TOKEN ECOSYSTEM
METAL (MTL): Native token of Metal Blockchain
XPR: Native token of XPR Network
XMD: Metallicus Dollar stablecoin
LOAN: Lending protocol token
Use cases: Gas fees, governance, DeFi, staking

## HOW METALLICUS DIFFERS FROM CBDC
NOT government-controlled.
NOT surveillance-based.
Institution-controlled stablecoins (credit union owns its stablecoin).
Member-facing — not government-facing.
Compliant with existing banking law — not replacing it.
This is the KEY DISTINCTION from a CBDC.

## FOLLOWING METALLICUS UPDATES
Official sources to follow:
- metallicus.com/news
- @MarshallHayner on X
- @MetallicusTDBN on X
- @XPRNetwork on X
- metallicus.com blog

## METALLICUS ROADMAP 2026
Known upcoming developments: [analysis based on public information]
Credit union expansion plans: [analysis]
TDBN scaling: [analysis]
XPR Network development: [analysis]

Cite metallicus.com and verified public sources."""

    result = query_llm_with_file(prompt, file_context)

    fig = None
    try:
        milestones = ["2016\nFounded", "2021\nMetal Pay\nLaunch", "2022\nMetal\nBlockchain",
                       "2023\nFirst CU\nCrypto/USD", "2024\nBonifii\nAcquisition",
                       "2025\nInvestiFi\nAlliance", "2026\n10th Year\nOperations"]
        significance = [6, 7, 8, 8, 7, 9, 9]
        colors_m = ["#3498db", "#3498db", "#f39c12", "#f39c12",
                     "#27ae60", "#27ae60", "#e74c3c"]
        fig = go.Figure(go.Bar(
            x=milestones, y=significance,
            marker_color=colors_m,
            text=[f"{s}/10" for s in significance],
            textposition="auto"
        ))
        fig.update_layout(
            title="Metallicus Key Milestones — Significance Score",
            yaxis_title="Significance (0-10)",
            yaxis=dict(range=[0, 10]),
            template="plotly_dark", height=400,
            paper_bgcolor="#0a0a1a")
    except Exception:
        fig = None

    return result + "\n\n" + WATERMARK, fig


# ─── Tab 10: US Treasury & Fed News ──────────────────────────────────────────

def analyze_treasury_fed_news(extra_context, file_context=""):
    prompt = f"""Provide a comprehensive analysis of US Treasury and Federal Reserve
news and policy as of 2026.

Additional Context: {enrich_with_urls(extra_context) if extra_context else "Full Treasury and Fed analysis 2026"}

## 🏛️ US TREASURY & FEDERAL RESERVE — 2026 NEWS & ANALYSIS

## US TREASURY — CURRENT STATUS
Secretary of the Treasury: Scott Bessent (appointed by Trump, confirmed 2025)
Key Treasury priorities under Bessent:
- Dollar dominance preservation
- Debt management with $36.6T outstanding
- Stablecoin regulation (GENIUS Act implementation)
- Strategic Bitcoin Reserve management
- De-dollarization countermeasures

## TREASURY DEBT MANAGEMENT 2026
Current outstanding debt: ~$36.6 trillion
Annual interest payments: ~$1.1 trillion (largest budget line item)
Debt ceiling status: [current analysis]
Treasury auction demand: [analysis — are foreign buyers still buying?]
China US Treasury holdings: [declining trend analysis]
Japan US Treasury holdings: [largest foreign holder — analysis]
Fed holdings of Treasuries: ~$4.5 trillion (QT ongoing)

## TREASURY DIGITAL CURRENCY INITIATIVES
Treasury views on stablecoins: [post-GENIUS Act analysis]
Treasury-backed stablecoin concept: [any current discussions?]
Gold certificate revaluation discussions: [based on Fed August 2025 study]
Bretton Woods 3.0 discussions: [any formal proposals?]

## FEDERAL RESERVE — CURRENT POLICY
Fed Chair: Jerome Powell (term through 2026)
Current Fed Funds Rate: 4.25-4.50%
Balance sheet: ~$6.7 trillion (QT ongoing at ~$25B/month)
Inflation: ~2.9% CPI
Dual mandate status: Inflation (near target) / Employment (analysis)

## FED POLICY OUTLOOK 2026
Rate cut probability in 2026: [market pricing analysis]
QT (Quantitative Tightening) end date: [projections]
Trump vs Powell tension: [political pressure on Fed independence analysis]
Fed gold revaluation study (August 2025): [significance and implications]
Digital dollar / CBDC status at Fed: [after Trump EO ban]

## TREASURY/FED COORDINATION ON NEW FINANCIAL SYSTEM
How Treasury and Fed are navigating:
- Stablecoin regulation
- Blockchain integration in payments
- Dollar dominance strategy
- XRP and Ripple's bank charter application impact
- Credit union digital asset adoption

## THE MAR-A-LAGO ACCORD THEORY
Reports of discussions about:
- Dollar revaluation
- Trade deficit reduction through currency adjustment
- Gold revaluation
- New Bretton Woods-style agreement
What is documented vs rumored: [analysis]
Historical precedent: Plaza Accord 1985
Economic rationale: [analysis]

## CURRENCY NEWS — DOLLAR vs WORLD
DXY performance in 2026: [analysis]
Dollar vs Euro: [current relationship]
Dollar vs Yuan: [strategic competition]
Dollar vs Gold: [year-to-date performance]
Dollar vs Bitcoin: [year-to-date performance]

## WHAT CITIZENS SHOULD WATCH
Top 5 Treasury/Fed indicators to monitor.
How monetary policy affects everyday financial decisions.
What rate cuts or hikes would mean for savings, mortgages, credit.

Cite Treasury.gov, federalreserve.gov, and verified financial news."""

    result = query_llm_with_file(prompt, file_context)

    fig = None
    try:
        years = [2019, 2020, 2021, 2022, 2023, 2024, 2025, 2026]
        fed_rate = [1.75, 0.25, 0.25, 4.5, 5.5, 5.25, 4.5, 4.375]
        us_debt = [22.7, 27.7, 28.4, 30.9, 33.2, 34.6, 36.0, 36.6]

        fig = make_subplots(rows=1, cols=2,
                             subplot_titles=["Fed Funds Rate (%)",
                                             "US National Debt ($ Trillions)"])
        fig.add_trace(go.Scatter(x=years, y=fed_rate, mode="lines+markers",
                                  name="Fed Rate", line=dict(color="#3498db", width=3)),
                      row=1, col=1)
        fig.add_trace(go.Bar(x=years, y=us_debt, name="US Debt ($T)",
                              marker_color=["#e74c3c" if d > 30 else "#f39c12"
                                            for d in us_debt]),
                      row=1, col=2)
        fig.update_layout(title="US Monetary Policy Dashboard 2019-2026",
                          template="plotly_dark", height=400,
                          paper_bgcolor="#0a0a1a")
    except Exception:
        fig = None

    return result + "\n\n" + WATERMARK, fig


# ─── Tab 11: Gods System vs Beast System ─────────────────────────────────────

def analyze_gods_vs_beast_system(extra_context, file_context=""):
    prompt = f"""Analyze current financial systems from a biblical theological perspective —
which systems align with sound money and freedom (Gods system) vs centralized control
and surveillance (the Beast system of Revelation 13).

Additional Context: {enrich_with_urls(extra_context) if extra_context else "Biblical financial systems analysis 2026"}

## 📖 GODS FINANCIAL SYSTEM vs THE BEAST SYSTEM — 2026 ANALYSIS

## THE BIBLICAL FRAMEWORK

### GODS SYSTEM — SOUND MONEY PRINCIPLES
Biblical basis for sound money:
- Leviticus 19:35-36 — Just weights and measures
- Proverbs 11:1 — Dishonest scales are an abomination
- Deuteronomy 25:13-15 — Fair and honest weights
- Isaiah 1:22 — Silver debased with dross (currency debasement warned against)
- Amos 8:5 — Cheating with false balances
- Ezekiel 45:10 — Accurate ephah and bath

Biblical characteristics of GODS financial system:
1. Sound money — backed by real assets (gold, silver per Proverbs)
2. No usury to the poor (Exodus 22:25, Leviticus 25:36)
3. Jubilee debt cancellation every 50 years (Leviticus 25)
4. No fractional reserve lending deception
5. Honest measurement and accounting
6. Local control — community-based (credit union model!)
7. Property rights respected — no theft through inflation
8. Financial liberty — not bondage (Galatians 5:1)

### THE BEAST SYSTEM — REVELATION 13 FINANCIAL CONTROL
Revelation 13:16-17 (KJV): Buy and sell control through mark system.
Characteristics of the Beast financial system:
1. Centralized global control
2. Surveillance of all transactions
3. Programmable money (can be turned off for non-compliant individuals)
4. Digital ID requirement to access financial system
5. No cash alternatives
6. Debt slavery (Proverbs 22:7 — borrower is slave to lender)
7. Currency debasement (inflation steals from poor)
8. One world currency system
9. No opting out

## SCORING CURRENT FINANCIAL SYSTEMS

### CENTRAL BANK DIGITAL CURRENCY (CBDC)
Alignment with GODS system: 1/10 (VERY LOW)
Alignment with Beast system: 9/10 (VERY HIGH)
Analysis: Programmable, surveillable, government-controlled money
WHY TRUMP BANNING CBDC IS SIGNIFICANT: [theological and freedom analysis]
Bible passage: Revelation 13:16-17 — most direct alignment

### STABLECOINS (Private — GENIUS Act compliant)
Alignment with GODS system: 6/10 (MODERATE)
Alignment with Beast system: 4/10 (MODERATE concern)
Analysis: Better than CBDC — private and decentralized options exist
Metallicus/credit union stablecoins: More aligned with GODS system (community-controlled)
Concern: Still digital — cash alternatives remain important

### BITCOIN (BTC)
Alignment with GODS system: 7/10
Analysis: Fixed supply (21M) — cannot be debased. No central controller.
Concern: Volatility, energy use, no direct asset backing
Sound money characteristics: Scarcity, portability, divisibility
Not a CBDC — NOT government controlled

### XRP / RIPPLE ECOSYSTEM
Alignment with GODS system: 6/10
Analysis: Faster, cheaper payments reduce friction — serves commerce
Concern: Ripple company has significant XRP holdings — centralization concern
Benefit: RLUSD stablecoin enables US dollar stability in payments
Not a CBDC — regulatory clarity without government control

### METALLICUS / CREDIT UNION STABLECOINS
Alignment with GODS system: 7/10 (HIGHEST of digital systems)
Analysis: Member-owned institutions issuing member-controlled stablecoins
Community-backed, locally controlled
Cooperative model aligns with biblical community economics
NOT a global system — decentralized by design

### GOLD AND SILVER (PHYSICAL)
Alignment with GODS system: 9/10 (HIGHEST)
Proverbs 17:3 — Gold refined by fire
Physical gold/silver is the most biblically aligned money
Historical: Used as money for 5,000+ years
Cannot be inflated, debased, or turned off
Challenge: Impractical for daily transactions

### UN / WEF / GLOBAL FINANCIAL GOVERNANCE
Alignment with GODS system: 1/10 (VERY LOW)
Alignment with Beast system: 8/10 (VERY HIGH)
Analysis: Centralized global control infrastructure
WHO financial governance, ESG scoring, social credit integration
Revelation 17-18 — Babylon the Great economic system alignment

## THE FINANCIAL FREEDOM SPECTRUM
[GODS SYSTEM] ←————————→ [BEAST SYSTEM]
Gold/Silver → BTC → Credit Union Stablecoins → Private Stablecoins → XRP →
Big Bank Digital → CBDC → UN/WEF Digital Currency

## PRACTICAL GUIDANCE FOR BELIEVERS
How to position finances in alignment with biblical principles.
The role of physical precious metals.
Community banking (credit unions) as a biblical model.
Avoiding the debt slavery system.
Preparation without paranoia.

Isaiah 55:1-2 — "Come, buy wine and milk without money and without price"
The ultimate financial freedom is in Christ — not in any earthly system.

Cite specific Bible passages and theological analysis throughout."""

    result = query_llm_with_file(prompt, file_context)

    fig = None
    try:
        systems = ["Physical\nGold/Silver", "Bitcoin\n(BTC)", "CU\nStablecoins",
                    "Private\nStablecoins", "XRP/\nRipple", "Big Bank\nDigital",
                    "CBDC", "UN/WEF\nGlobal"]
        gods_score = [9, 7, 7, 6, 6, 3, 1, 1]
        beast_score = [1, 3, 3, 4, 4, 7, 9, 9]

        fig = go.Figure()
        fig.add_trace(go.Bar(name="Gods System Alignment", x=systems,
                              y=gods_score,
                              marker_color="#27ae60",
                              text=[f"{v}/10" for v in gods_score],
                              textposition="auto"))
        fig.add_trace(go.Bar(name="Beast System Alignment", x=systems,
                              y=beast_score,
                              marker_color="#e74c3c",
                              text=[f"{v}/10" for v in beast_score],
                              textposition="auto"))
        fig.update_layout(barmode="group",
                          title="Financial Systems — Gods System vs Beast System Alignment",
                          yaxis_title="Alignment Score (0-10)",
                          yaxis=dict(range=[0, 10]),
                          template="plotly_dark", height=440,
                          paper_bgcolor="#0a0a1a")
    except Exception:
        fig = None

    return result + "\n\n" + WATERMARK, fig


# ─── Tab 12: Prediction & Scenario Engine ────────────────────────────────────

def analyze_monetary_scenarios(scenario_focus, time_horizon, extra_context, file_context=""):
    prompt = f"""Provide a comprehensive monetary reset scenario analysis.

Scenario Focus: {scenario_focus}
Time Horizon: {time_horizon}
Additional Context: {enrich_with_urls(extra_context) if extra_context else "Full scenario analysis"}

## 🔮 MONETARY RESET SCENARIO ENGINE

## SCENARIO 1: GRADUAL BLOCKCHAIN INTEGRATION (Most Likely)
Probability: 40%
What happens: Blockchain rails slowly replace SWIFT over 5-10 years.
XRP and stablecoins become standard for cross-border payments.
Credit unions lead retail blockchain adoption via Metallicus.
No dramatic dollar collapse — managed transition.
Asset implications: XRP, METAL, quality stablecoins benefit moderately.
IQD/VND/IDR gradually appreciate with economic growth.
Timeline: {time_horizon}

## SCENARIO 2: ASSET-BACKED DOLLAR RESET
Probability: 20%
What happens: US announces partial gold/asset backing for Treasury Dollar.
Gold revalued significantly (possibly $10,000-$50,000+/oz range).
Dollar reserve share stabilizes or increases.
BRICS responds with own asset-backed system.
Asset implications: Gold, silver surge dramatically. XRP as payment layer.
Credit union stablecoins backed by new Treasury Dollar become dominant.
Timeline: {time_horizon}

## SCENARIO 3: CBDC GLOBAL ROLLOUT (Excluding USA)
Probability: 25%
What happens: Europe, China, developing nations roll out CBDCs.
US maintains private stablecoin system (GENIUS Act framework).
Two-tier global financial system emerges.
Metallicus/credit union model flourishes in US as CBDC-free alternative.
Asset implications: US private crypto benefits. Physical gold as CBDC escape.
IQD/IRR revalue if sanctions normalize and they avoid CBDC stigma.
Timeline: {time_horizon}

## SCENARIO 4: RAPID DE-DOLLARIZATION / DOLLAR CRISIS
Probability: 10%
What happens: US debt crisis triggers dollar confidence collapse.
BRICS Unit scales rapidly as dollar alternative.
XRP/stablecoins become emergency payment infrastructure.
Gold revaluation forced — not planned.
Asset implications: Gold, silver, BTC, XRP all surge. IQD/IDR could RV quickly.
Timeline: {time_horizon}

## SCENARIO 5: BIBLICAL MONETARY RESET / DEBT JUBILEE
Probability: 5%
What happens: Global debt cancellation event (planned or forced).
New monetary system with sound money principles.
50-year Jubilee cycle alignment (2025-2026 calculation).
Credit union community model becomes dominant.
Fiat debt system collapses — asset-backed replaces it.
Biblical alignment: Highest of all scenarios.
Timeline: {time_horizon}

## PROBABILITY-WEIGHTED OUTLOOK
Most likely path: Combination of Scenarios 1 and 3.
Key wildcard: US debt ceiling crisis or foreign Treasury selling.
Black swan: Major bank failure triggering rapid blockchain adoption.

## WHAT THIS MEANS FOR YOUR FINANCES
Portfolio positioning for each scenario.
Role of physical precious metals.
Role of XRP and quality stablecoins.
Role of credit union membership.
Role of land and hard assets.
Disclaimer: Not financial advice — educational analysis only.

## METALLICUS/XPR NETWORK IN EACH SCENARIO
How Metallicus performs in each scenario.
Why credit union stablecoins are resilient across most scenarios.
Marshall Hayner vision alignment with each scenario.

Be specific. Present multiple expert viewpoints."""

    result = query_llm_with_file(prompt, file_context)

    fig = None
    try:
        scenarios = ["Gradual\nBlockchain\nIntegration", "Asset-Backed\nDollar Reset",
                      "CBDC Global\n(ex-USA)", "Dollar\nCrisis", "Biblical\nJubilee"]
        probabilities = [40, 20, 25, 10, 5]
        colors_s = ["#3498db", "#27ae60", "#f39c12", "#e74c3c", "#f1c40f"]
        fig = go.Figure(go.Pie(
            labels=scenarios, values=probabilities,
            hole=0.4, marker_colors=colors_s,
            textinfo="label+percent"
        ))
        fig.update_layout(title="Monetary Reset Scenario Probabilities",
                          template="plotly_dark", height=420,
                          paper_bgcolor="#0a0a1a")
    except Exception:
        fig = None

    return result + "\n\n" + WATERMARK, fig


# ─── Tab 13: Global Financial System Dashboard ───────────────────────────────

def analyze_global_financial_dashboard(file_context=""):
    prompt = """Provide a comprehensive global financial system dashboard summary for 2026.

## 📊 GLOBAL FINANCIAL SYSTEM DASHBOARD — 2026

## OVERALL SYSTEM TRANSITION SCORE
Financial system transition to blockchain/digital assets: X% complete
Speed of change: ACCELERATING / STABLE / SLOWING
Overall assessment: [summary]

## KEY METRICS DASHBOARD

| Metric | Current Value | Trend | Score |
|--------|--------------|-------|-------|
| US National Debt | $36.6T | Increasing | 3/10 |
| Dollar Reserve Share | 57.5% | Declining | 5/10 |
| XRP Adoption (banks) | 300+ institutions | Accelerating | 7/10 |
| Stablecoin Market Cap | $200B+ | Growing | 7/10 |
| CBDC Rollout (Global) | 60+ countries | Active | 6/10 |
| US CBDC Status | BANNED | Stable | 9/10 (freedom score) |
| Crypto Regulation Clarity | High (GENIUS Act) | Improving | 8/10 |
| Gold Price | $3,100+/oz | Bull market | 8/10 |
| Metallicus CU Adoption | Growing | Accelerating | 7/10 |
| BRICS Unit Status | Pilot phase | Expanding | 5/10 |

## TOP 10 DEVELOPMENTS TO WATCH IN 2026
1. [Development]: [analysis and timeline]
2. [Development]: [analysis]
3. [Development]: [analysis]
4. [Development]: [analysis]
5. [Development]: [analysis]
6. [Development]: [analysis]
7. [Development]: [analysis]
8. [Development]: [analysis]
9. [Development]: [analysis]
10. [Development]: [analysis]

## WINNERS IN THE NEW FINANCIAL SYSTEM
Countries: [ranked list with reasoning]
Technologies: [ranked list]
Financial institutions: [analysis]
Asset classes: [ranked list]

## LOSERS IN THE NEW FINANCIAL SYSTEM
Legacy systems being disrupted.
Countries falling behind.
Asset classes at risk.

## THE CITIZEN ACTION GUIDE
What every informed citizen should consider doing NOW:
1. [Action]: [reasoning]
2. [Action]: [reasoning]
3. [Action]: [reasoning]
4. [Action]: [reasoning]
5. [Action]: [reasoning]

Note: Not financial advice. Educational analysis only.

## BIBLICAL SYSTEM ALIGNMENT SUMMARY
Which emerging systems align with sound money principles.
Which systems to be cautious about.
Matthew 6:24 — Cannot serve two masters.
1 Timothy 6:10 — Love of money is the root of all evil.
The goal: Financial freedom and preparedness, not fear."""

    result = query_llm_with_file(prompt, file_context)

    fig = None
    try:
        metrics = ["Dollar\nStrength", "Crypto\nRegulation", "XRP\nAdoption",
                    "CU Digital\nAssets", "Sound\nMoney Return", "CBDC\nResistance",
                    "Gold\nBacking", "Blockchain\nBanking"]
        scores = [5, 8, 7, 7, 6, 9, 6, 6]
        colors_g = ["#27ae60" if s >= 7 else "#f39c12" if s >= 5 else "#e74c3c"
                    for s in scores]

        fig = make_subplots(rows=1, cols=2,
                             subplot_titles=["System Transition Scores",
                                             "Overall Progress Gauge"],
                             specs=[[{"type": "bar"}, {"type": "indicator"}]])
        fig.add_trace(go.Bar(x=metrics, y=scores, marker_color=colors_g,
                              text=[f"{s}/10" for s in scores],
                              textposition="auto"), row=1, col=1)
        avg = sum(scores) / len(scores)
        fig.add_trace(go.Indicator(
            mode="gauge+number",
            value=avg * 10,
            title={"text": "New Financial System\nTransition Progress %"},
            gauge={
                "axis": {"range": [0, 100]},
                "bar": {"color": "#3498db"},
                "steps": [
                    {"range": [0, 30], "color": "#e74c3c"},
                    {"range": [30, 60], "color": "#f39c12"},
                    {"range": [60, 80], "color": "#27ae60"},
                    {"range": [80, 100], "color": "#2ecc71"}
                ]
            }
        ), row=1, col=2)
        fig.update_layout(title="Global Financial System Transition Dashboard 2026",
                          template="plotly_dark", height=440,
                          paper_bgcolor="#0a0a1a")
    except Exception:
        fig = None

    return result + "\n\n" + WATERMARK, fig


# ─── Tab 14: Reports & Presentations ─────────────────────────────────────────

def generate_monetary_report(report_type, focus, audience, extra_context, file_context=""):
    prompt = f"""You are a financial systems analyst creating a professional report
on the global monetary reset and new financial system.

Report Type: {report_type}
Focus: {focus}
Audience: {audience}
Additional Notes: {enrich_with_urls(extra_context) if extra_context else "Comprehensive overview"}

Create a professional {report_type} with this structure:

# AI MONETARY RESET TRACKER — {report_type.upper()}
Generated: {datetime.now().strftime("%B %d, %Y")} | Focus: {focus}

## EXECUTIVE SUMMARY
3-5 sentence overview of the most important developments in the new financial system.

## SECTION 1: THE NEW FINANCIAL SYSTEM — WHERE WE ARE
Current state of blockchain-banking integration.
Key regulatory milestones (GENIUS Act, XRP clarity, no CBDC in USA).
XRP/Ripple global adoption status.

## SECTION 2: METALLICUS AND CREDIT UNIONS
What Metallicus is building for credit unions.
InvestiFi alliance and stablecoin access.
The Digital Banking Network (TDBN) significance.
Why this matters for financial freedom.

## SECTION 3: FOREX AND CURRENCY CHANGES
Dollar status and reserve currency trends.
Key revaluation candidates: IQD, VND, IDR, ZiG.
Asset-backed currency race globally.

## SECTION 4: ASSET-BACKED DOLLAR ANALYSIS
Non-CBDC Treasury Dollar discussions.
Gold revaluation scenario.
Strategic Bitcoin Reserve.
Dollar dominance preservation strategies.

## SECTION 5: GODS SYSTEM vs BEAST SYSTEM
Biblical perspective on current financial systems.
Which systems align with sound money principles.
CBDC ban in America as a freedom milestone.
Credit union model as biblical community finance.

## SECTION 6: SCENARIOS AND OUTLOOK
Top 3 most likely monetary scenarios.
Timeline projections.
Asset class implications.

## SECTION 7: CITIZEN ACTION GUIDE
5 practical steps for financial preparedness.
Resources for staying informed (Metallicus, XRP Network, Treasury updates).

## KEY TERMS GLOSSARY
Define: CBDC, Stablecoin, GENIUS Act, RLUSD, TDBN, XPR Network,
Metal Blockchain, ODL, XRPL, IQD RV, ZiG, Mar-a-Lago Accord.

Write for {audience} audience. Be factual, cite sources, and include disclaimer."""

    report_text = query_llm_with_file(prompt, file_context)
    output_path = None
    try:
        if "PowerPoint" in report_type:
            output_path = _make_monetary_pptx(report_text, focus)
        else:
            output_path = _make_monetary_pdf(report_text, focus)
    except Exception as e:
        report_text += f"\n\n[Document generation error: {str(e)}]"
    return report_text + "\n\n" + WATERMARK, output_path


def _make_monetary_pptx(text, focus):
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    def add_slide(title, body, bg=(5, 10, 20), tc=(52, 152, 219)):
        slide = prs.slides.add_slide(blank)
        bg_fill = slide.background.fill
        bg_fill.solid()
        bg_fill.fore_color.rgb = RGBColor(*bg)
        tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12.3), Inches(1.3))
        p = tb.text_frame.paragraphs[0]
        p.text = title
        p.font.size = Pt(26)
        p.font.bold = True
        p.font.color.rgb = RGBColor(*tc)
        bb = slide.shapes.add_textbox(Inches(0.5), Inches(1.6), Inches(12.3), Inches(5.6))
        tf = bb.text_frame
        tf.word_wrap = True
        for line in body.split("\n")[:22]:
            line = line.strip()
            if not line:
                continue
            para = tf.add_paragraph()
            para.text = line.lstrip("*#- ").strip("*")
            para.font.size = Pt(13)
            para.font.color.rgb = RGBColor(220, 230, 245)

    title_slide = prs.slides.add_slide(blank)
    bg = title_slide.background.fill
    bg.solid()
    bg.fore_color.rgb = RGBColor(3, 8, 18)
    tb = title_slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.7), Inches(1.5))
    p = tb.text_frame.paragraphs[0]
    p.text = "AI Monetary Reset Tracker"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(52, 152, 219)
    tb2 = title_slide.shapes.add_textbox(Inches(0.8), Inches(3.4), Inches(11.7), Inches(1.0))
    p2 = tb2.text_frame.paragraphs[0]
    p2.text = focus
    p2.font.size = Pt(22)
    p2.font.color.rgb = RGBColor(46, 204, 113)
    tb3 = title_slide.shapes.add_textbox(Inches(0.8), Inches(4.6), Inches(11.7), Inches(0.8))
    p3 = tb3.text_frame.paragraphs[0]
    p3.text = f"Generated: {datetime.now().strftime('%B %d, %Y')} | AI Monetary Reset Tracker"
    p3.font.size = Pt(16)
    p3.font.color.rgb = RGBColor(150, 170, 200)

    sections = []
    cur_title, cur_body = "Overview", []
    for line in text.split("\n"):
        s = line.strip()
        if s.startswith("## ") or (s.startswith("**") and s.endswith("**") and len(s) > 4):
            if cur_body:
                sections.append((cur_title, "\n".join(cur_body)))
            cur_title = s.lstrip("#* ").rstrip("*# ")
            cur_body = []
        else:
            cur_body.append(line)
    if cur_body:
        sections.append((cur_title, "\n".join(cur_body)))

    bgs = [(5,10,20), (5,20,10), (20,10,5), (10,5,20), (5,15,15), (15,10,5)]
    tcs = [(52,152,219), (46,204,113), (241,196,15), (155,89,182), (26,188,156), (231,76,60)]
    for i, (t, b) in enumerate(sections[:12]):
        add_slide(t, b, bgs[i % len(bgs)], tcs[i % len(tcs)])

    closing = prs.slides.add_slide(blank)
    bg2 = closing.background.fill
    bg2.solid()
    bg2.fore_color.rgb = RGBColor(3, 8, 18)
    tb4 = closing.shapes.add_textbox(Inches(1), Inches(2.8), Inches(11.3), Inches(2))
    tf4 = tb4.text_frame
    tf4.word_wrap = True
    p4 = tf4.paragraphs[0]
    p4.text = "Stay Informed. Stay Free."
    p4.font.size = Pt(30)
    p4.font.bold = True
    p4.font.color.rgb = RGBColor(52, 152, 219)
    p5 = tf4.add_paragraph()
    p5.text = "© 2026 Existential Gateway, LLC | AI Monetary Reset Tracker"
    p5.font.size = Pt(14)
    p5.font.color.rgb = RGBColor(150, 170, 200)

    tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".pptx")
    prs.save(tmp.name)
    return tmp.name


def _make_monetary_pdf(text, focus):
    from reportlab.lib.pagesizes import letter
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.units import inch
    from reportlab.lib import colors
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
    tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".pdf")
    doc = SimpleDocTemplate(tmp.name, pagesize=letter,
                             leftMargin=inch, rightMargin=inch,
                             topMargin=inch, bottomMargin=inch)
    styles = getSampleStyleSheet()
    title_s = ParagraphStyle("T2", parent=styles["Title"], fontSize=20,
                              spaceAfter=16, textColor=colors.HexColor("#1a3a5c"))
    h2_s = ParagraphStyle("H2", parent=styles["Heading2"], fontSize=14,
                           spaceAfter=8, textColor=colors.HexColor("#2980b9"))
    body_s = ParagraphStyle("B2", parent=styles["Normal"], fontSize=11,
                             spaceAfter=6, leading=16)
    story = [Paragraph("AI Monetary Reset Tracker", title_s),
             Paragraph(f"Focus: {focus}", h2_s),
             Paragraph(f"Generated: {datetime.now().strftime('%B %d, %Y')}", body_s),
             Spacer(1, 0.2 * inch)]
    for line in text.split("\n"):
        s = line.strip()
        if not s:
            story.append(Spacer(1, 0.08 * inch))
        elif s.startswith("## "):
            story.append(Paragraph(s[3:], h2_s))
        elif s.startswith("# "):
            story.append(Paragraph(s[2:], title_s))
        elif s.startswith("**") and s.endswith("**"):
            story.append(Paragraph(s.strip("*"), h2_s))
        else:
            safe = s.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")
            story.append(Paragraph(safe, body_s))
    story.extend([Spacer(1, 0.3 * inch),
                  Paragraph("© 2026 Existential Gateway, LLC | AI Monetary Reset Tracker | existentialgateway@gmail.com",
                             ParagraphStyle("footer", parent=styles["Normal"],
                                            fontSize=9, textColor=colors.grey))])
    doc.build(story)
    return tmp.name


# ─── Tab 15: AI Financial Chat ────────────────────────────────────────────────

def financial_chat(message, history):
    messages = [
        {"role": "system", "content": SYSTEM_PROMPT + (
            "\n\nYou are helping citizens understand the new financial system. "
            "Answer questions about XRP, Ripple, Metallicus, Marshall Hayner, XPR Network, "
            "credit union digital assets, stablecoins, GENIUS Act, no CBDC in America, "
            "FOREX and currency revaluations (IQD, VND, IDR, ZiG, IRR), asset-backed dollars, "
            "Gods financial system vs Beast system, US Treasury and Fed policy, "
            "and anything related to the monetary reset. "
            "Always include a disclaimer that nothing is financial advice. "
            "Cite specific sources and verified data where possible. "
            "Be direct, informative, and empowering for citizens."
        )}
    ]
    for user_msg, bot_msg in history:
        messages.append({"role": "user", "content": user_msg})
        messages.append({"role": "assistant", "content": bot_msg})
    messages.append({"role": "user", "content": message})

    import os
    API_KEY = os.environ.get("OPENAI_API_KEY", "")
    headers = {"Authorization": f"Bearer {API_KEY}", "Content-Type": "application/json"}
    msgs = messages if 'messages' in dir() else [{"role": "user", "content": prompt}]
    payload = {"model": "gpt-4o", "max_tokens": 4000, "messages": msgs}
    try:
        response = requests.post("https://api.openai.com/v1/chat/completions", headers=headers, json=payload, timeout=240)
        result = response.json()
        if "choices" in result:
            return result["choices"][0]["message"]["content"]
        return f"API Error: {result}"
    except Exception as e:
        return f"Error: {str(e)}"


# ─── Gradio UI ────────────────────────────────────────────────────────────────

with gr.Blocks(title="AI Monetary Reset Tracker") as demo:
    gr.HTML('<style>::selection{background:#c9a84c;color:#000!important}::-moz-selection{background:#c9a84c;color:#000!important}</style>')
    gr.Markdown("# 💰 AI Monetary Reset Tracker")
    gr.Markdown("### *Tracking the Great Financial Reset — Blockchain, XRP, Stablecoins, FOREX & the New Financial System*")
    gr.Markdown(CITIZEN_NOTICE)
    gr.Markdown(DISCLAIMER)

    with gr.Tabs():

        # ── Tab 1: Blockchain & Banking ────────────────────────────────────────
        with gr.Tab("🔗 Blockchain & Banking"):
            gr.Markdown("## Blockchain & Banking Integration — How Far Along Are We?")
            with gr.Row():
                with gr.Column(scale=1):
                    t1_context = gr.Textbox(label="Specific Focus (optional)",
                                             placeholder="e.g. SWIFT vs XRP, tokenization, stablecoin settlement...")
                    t1_upload = gr.File(
                        label="📎 Upload for Context (PDF/CSV/XLSX/PNG/JPG)",
                        file_types=[".pdf",".csv",".xlsx",".xls",".png",".jpg",".jpeg"],
                        file_count="single")
                    t1_btn = gr.Button("🔗 Analyze Integration", variant="primary")
                    gr.Markdown(WAIT_MSG)
                with gr.Column(scale=2):
                    t1_output = gr.Markdown(label="Analysis")
            t1_chart = gr.Plot(label="Integration Progress Chart")
            t1_btn.click(analyze_blockchain_banking, inputs=[t1_context, t1_upload],
                         outputs=[t1_output, t1_chart])

        # ── Tab 2: Crypto/Stablecoin ───────────────────────────────────────────
        with gr.Tab("🪙 Crypto & Stablecoins"):
            gr.Markdown("## Crypto, Stablecoins & Digital Currency — Live Prices & Analysis")
            gr.Markdown("*Live prices via Yahoo Finance*")
            with gr.Row():
                with gr.Column(scale=1):
                    t2_context = gr.Textbox(label="Specific Focus (optional)",
                                             placeholder="e.g. GENIUS Act, RLUSD, BTC ETF, stablecoin landscape...")
                    t2_upload = gr.File(
                        label="📎 Upload for Context (PDF/CSV/XLSX/PNG/JPG)",
                        file_types=[".pdf",".csv",".xlsx",".xls",".png",".jpg",".jpeg"],
                        file_count="single")
                    t2_btn = gr.Button("🪙 Analyze Crypto & Stablecoins", variant="primary")
                    gr.Markdown(WAIT_MSG)
                with gr.Column(scale=2):
                    t2_output = gr.Markdown(label="Analysis")
            t2_chart = gr.Plot(label="Market Overview Chart")
            t2_btn.click(analyze_crypto_stablecoin, inputs=[t2_context, t2_upload],
                         outputs=[t2_output, t2_chart])

        # ── Tab 3: US Crypto Regulations ──────────────────────────────────────
        with gr.Tab("⚖️ US Regulations"):
            gr.Markdown("## US Crypto Regulations — GENIUS Act, No CBDC, SEC Settlement")
            with gr.Row():
                with gr.Column(scale=1):
                    t3_context = gr.Textbox(label="Specific Regulation to Analyze (optional)",
                                             placeholder="e.g. GENIUS Act, XRP ETF, Trump CBDC ban, bank charter...")
                    t3_upload = gr.File(
                        label="📎 Upload for Context (PDF/CSV/XLSX/PNG/JPG)",
                        file_types=[".pdf",".csv",".xlsx",".xls",".png",".jpg",".jpeg"],
                        file_count="single")
                    t3_btn = gr.Button("⚖️ Analyze Regulations", variant="primary")
                    gr.Markdown(WAIT_MSG)
                with gr.Column(scale=2):
                    t3_output = gr.Markdown(label="Regulatory Analysis")
            t3_chart = gr.Plot(label="Regulatory Clarity Chart")
            t3_btn.click(analyze_crypto_regulations, inputs=[t3_context, t3_upload],
                         outputs=[t3_output, t3_chart])

        # ── Tab 4: Asset-Backed Dollar ─────────────────────────────────────────
        with gr.Tab("💵 Asset-Backed Dollar"):
            gr.Markdown("## Non-CBDC Asset-Backed US Treasury Dollar Analysis")
            with gr.Row():
                with gr.Column(scale=1):
                    t4_context = gr.Textbox(label="Specific Focus (optional)",
                                             placeholder="e.g. gold-backed dollar, BTC reserve, Mar-a-Lago accord...")
                    t4_upload = gr.File(
                        label="📎 Upload for Context (PDF/CSV/XLSX/PNG/JPG)",
                        file_types=[".pdf",".csv",".xlsx",".xls",".png",".jpg",".jpeg"],
                        file_count="single")
                    t4_btn = gr.Button("💵 Analyze Asset-Backed Dollar", variant="primary")
                    gr.Markdown(WAIT_MSG)
                with gr.Column(scale=2):
                    t4_output = gr.Markdown(label="Analysis")
            t4_chart = gr.Plot(label="Dollar Backing Scenarios Chart")
            t4_btn.click(analyze_asset_backed_dollar, inputs=[t4_context, t4_upload],
                         outputs=[t4_output, t4_chart])

        # ── Tab 5: XRP & Cross-Border ──────────────────────────────────────────
        with gr.Tab("🌊 XRP & Cross-Border"):
            gr.Markdown("## XRP & Ripple — Global Cross-Border Payment Adoption")
            gr.Markdown("*Live XRP price via Yahoo Finance*")
            with gr.Row():
                with gr.Column(scale=1):
                    t5_context = gr.Textbox(label="Specific Focus (optional)",
                                             placeholder="e.g. specific bank adoption, RLUSD, ODL corridors, bank charter...")
                    t5_upload = gr.File(label="📎 Upload for Context (PDF/CSV/XLSX/PNG/JPG)",
                                       file_types=[".pdf",".csv",".xlsx",".xls",".png",".jpg",".jpeg"],
                                       file_count="single", scale=1)
                    t5_btn = gr.Button("🌊 Analyze XRP Adoption", variant="primary")
                    gr.Markdown(WAIT_MSG)
                with gr.Column(scale=2):
                    t5_output = gr.Markdown(label="XRP Analysis")
            t5_chart = gr.Plot(label="XRP Adoption by Region")
            t5_btn.click(analyze_xrp_payments, inputs=[t5_context, t5_upload],
                         outputs=[t5_output, t5_chart])

        # ── Tab 6: FOREX & Currencies ──────────────────────────────────────────
        with gr.Tab("💱 FOREX & Currencies"):
            gr.Markdown("## FOREX Tracker — IQD, VND, IDR, ZiG, IRR & Global Rate Changes")
            gr.Markdown("*Live rates via Yahoo Finance where available*")
            with gr.Row():
                with gr.Column(scale=1):
                    t6_context = gr.Textbox(label="Specific Currency or Market (optional)",
                                             placeholder="e.g. Iraqi Dinar latest news, IDR BRICS impact...")
                    t6_upload = gr.File(label="📎 Upload for Context (PDF/CSV/XLSX/PNG/JPG)",
                                       file_types=[".pdf",".csv",".xlsx",".xls",".png",".jpg",".jpeg"],
                                       file_count="single", scale=1)
                    t6_btn = gr.Button("💱 Analyze FOREX", variant="primary")
                    gr.Markdown(WAIT_MSG)
                with gr.Column(scale=2):
                    t6_output = gr.Markdown(label="FOREX Analysis")
            t6_chart = gr.Plot(label="Currency Analysis Chart")
            t6_btn.click(analyze_forex_currencies, inputs=[t6_context, t6_upload],
                         outputs=[t6_output, t6_chart])

        # ── Tab 7: Currency Revaluation ────────────────────────────────────────
        with gr.Tab("🌍 Currency Revaluation"):
            gr.Markdown("## Global Currency Revaluation Analysis — Asset-Backed Currency Race")
            with gr.Row():
                with gr.Column(scale=1):
                    t7_currency = gr.Dropdown(
                        choices=["All Revaluation Candidates",
                                  "Iraqi Dinar (IQD) Deep Dive",
                                  "Vietnamese Dong (VND)",
                                  "Indonesian Rupiah (IDR)",
                                  "Zimbabwe Gold (ZiG)",
                                  "Iranian Rial (IRR)",
                                  "Venezuelan Bolivar (VES)",
                                  "Gulf Currency Union Potential"],
                        value="All Revaluation Candidates",
                        label="Currency Focus")
                    t7_context = gr.Textbox(label="Additional Context (optional)", lines=2)
                    t7_upload = gr.File(label="📎 Upload for Context (PDF/CSV/XLSX/PNG/JPG)",
                                       file_types=[".pdf",".csv",".xlsx",".xls",".png",".jpg",".jpeg"],
                                       file_count="single", scale=1)
                    t7_btn = gr.Button("🌍 Analyze Revaluation Potential", variant="primary")
                    gr.Markdown(WAIT_MSG)
                with gr.Column(scale=2):
                    t7_output = gr.Markdown(label="Revaluation Analysis")
            t7_chart = gr.Plot(label="Revaluation Radar Chart")
            t7_btn.click(analyze_currency_revaluation, inputs=[t7_currency, t7_context, t7_upload],
                         outputs=[t7_output, t7_chart])

        # ── Tab 8: Credit Unions vs Banks ──────────────────────────────────────
        with gr.Tab("🏦 Credit Unions vs Banks"):
            gr.Markdown("## Credit Unions vs Banks — Digital Assets & New Financial System")
            with gr.Row():
                with gr.Column(scale=1):
                    t8_context = gr.Textbox(label="Specific Focus (optional)",
                                             placeholder="e.g. Metallicus stablecoin for credit unions, NCUA rules, digital assets...")
                    t8_upload = gr.File(label="📎 Upload for Context (PDF/CSV/XLSX/PNG/JPG)",
                                       file_types=[".pdf",".csv",".xlsx",".xls",".png",".jpg",".jpeg"],
                                       file_count="single", scale=1)
                    t8_btn = gr.Button("🏦 Analyze CU vs Banks", variant="primary")
                    gr.Markdown(WAIT_MSG)
                with gr.Column(scale=2):
                    t8_output = gr.Markdown(label="Analysis")
            t8_chart = gr.Plot(label="Comparison Radar Chart")
            t8_btn.click(analyze_credit_unions_vs_banks, inputs=[t8_context, t8_upload],
                         outputs=[t8_output, t8_chart])

        # ── Tab 9: Metallicus & XPR Network ───────────────────────────────────
        with gr.Tab("⛓️ Metallicus & XPR"):
            gr.Markdown("## Metallicus, Marshall Hayner & XPR Network — Complete Analysis")
            gr.Markdown("> 📡 **For real-time updates follow:** [@MarshallHayner](https://x.com/MarshallHayner) | [@MetallicusTDBN](https://x.com/MetallicusTDBN) | [@XPRNetwork](https://x.com/XPRNetwork) on X")
            with gr.Row():
                with gr.Column(scale=1):
                    t9_context = gr.Textbox(label="Specific Focus (optional)",
                                             placeholder="e.g. TDBN, Metal Blockchain, InvestiFi alliance, credit union stablecoins...")
                    t9_upload = gr.File(label="📎 Upload for Context (PDF/CSV/XLSX/PNG/JPG)",
                                       file_types=[".pdf",".csv",".xlsx",".xls",".png",".jpg",".jpeg"],
                                       file_count="single", scale=1)
                    t9_btn = gr.Button("⛓️ Analyze Metallicus Ecosystem", variant="primary")
                    gr.Markdown(WAIT_MSG)
                with gr.Column(scale=2):
                    t9_output = gr.Markdown(label="Metallicus Analysis")
            t9_chart = gr.Plot(label="Metallicus Milestones Chart")
            t9_btn.click(analyze_metallicus_xpr, inputs=[t9_context, t9_upload],
                         outputs=[t9_output, t9_chart])

        # ── Tab 10: US Treasury & Fed ──────────────────────────────────────────
        with gr.Tab("🏛️ Treasury & Fed News"):
            gr.Markdown("## US Treasury & Federal Reserve — Policy, Debt, and Dollar News")
            with gr.Row():
                with gr.Column(scale=1):
                    t10_context = gr.Textbox(label="Specific Focus (optional)",
                                              placeholder="e.g. Scott Bessent policy, Fed rate outlook, debt ceiling, Mar-a-Lago accord...")
                    t10_upload = gr.File(label="📎 Upload for Context (PDF/CSV/XLSX/PNG/JPG)",
                                       file_types=[".pdf",".csv",".xlsx",".xls",".png",".jpg",".jpeg"],
                                       file_count="single", scale=1)
                    t10_btn = gr.Button("🏛️ Analyze Treasury & Fed", variant="primary")
                    gr.Markdown(WAIT_MSG)
                with gr.Column(scale=2):
                    t10_output = gr.Markdown(label="Treasury & Fed Analysis")
            t10_chart = gr.Plot(label="Monetary Policy Chart")
            t10_btn.click(analyze_treasury_fed_news, inputs=[t10_context, t10_upload],
                          outputs=[t10_output, t10_chart])

        # ── Tab 11: Gods vs Beast System ───────────────────────────────────────
        with gr.Tab("📖 Gods vs Beast System"):
            gr.Markdown("## Which Financial Systems Align with Gods Principles vs the Beast System?")
            gr.Markdown("> *Leviticus 19:35-36 — Just weights and measures | Revelation 13:16-17 — Buy/sell control*")
            with gr.Row():
                with gr.Column(scale=1):
                    t11_context = gr.Textbox(label="Specific Focus (optional)",
                                              placeholder="e.g. CBDC vs stablecoins, Bitcoin and sound money, credit union biblical model...")
                    t11_upload = gr.File(label="📎 Upload for Context (PDF/CSV/XLSX/PNG/JPG)",
                                       file_types=[".pdf",".csv",".xlsx",".xls",".png",".jpg",".jpeg"],
                                       file_count="single", scale=1)
                    t11_btn = gr.Button("📖 Analyze Gods vs Beast System", variant="primary")
                    gr.Markdown(WAIT_MSG)
                with gr.Column(scale=2):
                    t11_output = gr.Markdown(label="Biblical Financial Analysis")
            t11_chart = gr.Plot(label="System Alignment Chart")
            t11_btn.click(analyze_gods_vs_beast_system, inputs=[t11_context, t11_upload],
                          outputs=[t11_output, t11_chart])

        # ── Tab 12: Prediction & Scenarios ────────────────────────────────────
        with gr.Tab("🔮 Scenarios"):
            gr.Markdown("## Monetary Reset Scenario Engine — What Happens Next?")
            with gr.Row():
                with gr.Column(scale=1):
                    t12_focus = gr.Dropdown(
                        choices=["Complete Monetary Reset Overview",
                                  "XRP and Ripple Scenarios",
                                  "Dollar Asset-Backing Scenarios",
                                  "Credit Union / Metallicus Scenarios",
                                  "Currency Revaluation Scenarios",
                                  "Biblical Reset / Jubilee Scenario"],
                        value="Complete Monetary Reset Overview",
                        label="Scenario Focus")
                    t12_horizon = gr.Dropdown(
                        choices=["6 Months", "1 Year", "3 Years", "5 Years", "10 Years"],
                        value="3 Years", label="Time Horizon")
                    t12_context = gr.Textbox(label="Additional Context", lines=2)
                    t12_upload = gr.File(label="📎 Upload for Context (PDF/CSV/XLSX/PNG/JPG)",
                                       file_types=[".pdf",".csv",".xlsx",".xls",".png",".jpg",".jpeg"],
                                       file_count="single", scale=1)
                    t12_btn = gr.Button("🔮 Generate Scenarios", variant="primary")
                    gr.Markdown(WAIT_MSG)
                with gr.Column(scale=2):
                    t12_output = gr.Markdown(label="Scenario Analysis")
            t12_chart = gr.Plot(label="Probability Chart")
            t12_btn.click(analyze_monetary_scenarios, inputs=[t12_focus, t12_horizon, t12_context, t12_upload],
                          outputs=[t12_output, t12_chart])

        # ── Tab 13: Dashboard ──────────────────────────────────────────────────
        with gr.Tab("📊 Global Dashboard"):
            gr.Markdown("## Global Financial System Dashboard — All Systems Combined")
            with gr.Row():
                with gr.Column(scale=1):
                    t13_context = gr.Textbox(
                        label="Specific Focus (optional)",
                        placeholder="e.g. focus on XRP adoption, Metallicus, credit unions...")
                    t13_upload = gr.File(
                        label="📎 Upload for Context (PDF/CSV/XLSX/PNG/JPG)",
                        file_types=[".pdf",".csv",".xlsx",".xls",".png",".jpg",".jpeg"],
                        file_count="single")
                    t13_btn = gr.Button("📊 Load Global Dashboard", variant="primary")
                    gr.Markdown(WAIT_MSG)
                with gr.Column(scale=2):
                    t13_output = gr.Markdown(label="Analysis")
            t13_chart = gr.Plot(label="System Transition Dashboard")
            t13_btn.click(analyze_global_financial_dashboard, inputs=[t13_upload],
                          outputs=[t13_output, t13_chart])

        # ── Tab 14: Reports & Presentations ───────────────────────────────────
        with gr.Tab("📄 Reports & Presentations"):
            gr.Markdown("## Generate Professional Monetary Reset Reports & Presentations")
            with gr.Row():
                with gr.Column(scale=1):
                    t14_type = gr.Dropdown(
                        choices=["PDF Report", "PowerPoint Presentation",
                                  "Executive Summary (PDF)", "Citizen Guide (PDF)"],
                        value="PDF Report", label="Report Type")
                    t14_focus = gr.Dropdown(
                        choices=["Complete Overview — All Topics",
                                  "XRP & Ripple Focus",
                                  "Metallicus & Credit Unions Focus",
                                  "Currency Revaluation Focus",
                                  "Gods vs Beast System Focus",
                                  "Asset-Backed Dollar Focus"],
                        value="Complete Overview — All Topics",
                        label="Content Focus")
                    t14_audience = gr.Dropdown(
                        choices=["General Public", "Investors", "Credit Union Members",
                                  "Church/Faith Community", "Financial Professionals"],
                        value="General Public", label="Target Audience")
                    t14_context = gr.Textbox(label="Additional Notes", lines=2)
                    t14_upload = gr.File(label="📎 Upload for Context (PDF/CSV/XLSX/PNG/JPG)",
                                       file_types=[".pdf",".csv",".xlsx",".xls",".png",".jpg",".jpeg"],
                                       file_count="single", scale=1)
                    t14_btn = gr.Button("📄 Generate Report", variant="primary")
                    gr.Markdown(WAIT_MSG)
                    t14_download = gr.File(label="⬇️ Download Report")
                with gr.Column(scale=2):
                    t14_output = gr.Markdown(label="Report Preview")
            t14_btn.click(generate_monetary_report, inputs=[t14_type, t14_focus, t14_audience, t14_context, t14_upload],
                          outputs=[t14_output, t14_download])

        # ── Tab 15: AI Financial Chat ──────────────────────────────────────────
        with gr.Tab("💬 AI Financial Chat"):
            gr.Markdown("## Ask Anything About the New Financial System")
            gr.Markdown("> 📡 **Follow for real-time updates:** [@MarshallHayner](https://x.com/MarshallHayner) | [@MetallicusTDBN](https://x.com/MetallicusTDBN) | [@XPRNetwork](https://x.com/XPRNetwork)")
            gr.ChatInterface(
                fn=financial_chat,
                examples=[
                    "What is Metallicus and why does it matter for credit unions?",
                    "Explain the GENIUS Act and how it affects stablecoins",
                    "Why did Trump ban CBDCs and what does it mean for financial freedom?",
                    "What is XRP and how is Ripple being used by global banks?",
                    "What is the Iraqi Dinar revaluation and is it realistic?",
                    "What is the Digital Banking Network (TDBN) from Metallicus?",
                    "How does an asset-backed Treasury Dollar differ from a CBDC?",
                    "What is the XPR Network and how can credit unions use it?",
                    "How does the current financial system align with Revelation 13?",
                    "What is the Zimbabwe ZiG and how is it performing?",
                    "What is RLUSD and how does it differ from regular XRP?",
                    "Should I join a credit union instead of a bank in 2026?",
                    "What is Marshall Hayner working on at Metallicus in 2026?",
                    "How does the Indonesian Rupiah connect to BRICS membership?",
                ],
                title="",
            )

    gr.Markdown(WATERMARK)


if __name__ == "__main__":
    demo.launch(server_name="0.0.0.0", server_port=int(os.environ.get("GRADIO_SERVER_PORT", 7860)), share=False, ssr_mode=False, theme=gr.themes.Base(
        primary_hue=gr.themes.colors.Color(
            c50="#fef9ec", c100="#faefd0", c200="#f4dea0", c300="#edc970",
            c400="#e4b04a", c500="#c9a84c", c600="#a8872e", c700="#856519",
            c800="#5e440d", c900="#3a2a05", c950="#1e1502"
        ),
        secondary_hue=gr.themes.colors.Color(
            c50="#e8edf5", c100="#c5d0e6", c200="#9eb0d4", c300="#7490c2",
            c400="#4f74b0", c500="#2d5a9e", c600="#1a3a6e", c700="#112240",
            c800="#0a1628", c900="#060e1a", c950="#03070d"
        ),
        neutral_hue=gr.themes.colors.Color(
            c50="#f0f2f7", c100="#d8dde9", c200="#b8c2d6", c300="#95a3be",
            c400="#7285a6", c500="#536890", c600="#3a4f73", c700="#263856",
            c800="#162438", c900="#0c1622", c950="#060b11"
        ),
        font=gr.themes.GoogleFont("DM Sans"),
        font_mono=gr.themes.GoogleFont("DM Mono"),
    ).set(
        body_background_fill="#0a1628",
        body_background_fill_dark="#0a1628",
        body_text_color="#f8f9fc",
        body_text_color_dark="#f8f9fc",
        block_background_fill="#112240",
        block_background_fill_dark="#112240",
        block_border_color="rgba(201,168,76,0.2)",
        block_border_color_dark="rgba(201,168,76,0.2)",
        block_label_background_fill="#112240",
        block_label_background_fill_dark="#112240",
        block_label_text_color="#c9a84c",
        block_label_text_color_dark="#c9a84c",
        block_title_text_color="#c9a84c",
        block_title_text_color_dark="#c9a84c",
        button_primary_background_fill="linear-gradient(135deg,#c9a84c,#e8c96a)",
        button_primary_background_fill_dark="linear-gradient(135deg,#c9a84c,#e8c96a)",
        button_primary_text_color="#0a1628",
        button_primary_text_color_dark="#0a1628",
        button_secondary_background_fill="#112240",
        button_secondary_background_fill_dark="#112240",
        button_secondary_border_color="rgba(201,168,76,0.3)",
        button_secondary_border_color_dark="rgba(201,168,76,0.3)",
        button_secondary_text_color="#c9a84c",
        button_secondary_text_color_dark="#c9a84c",
        input_background_fill="#0a1628",
        input_background_fill_dark="#0a1628",
        input_border_color="rgba(201,168,76,0.2)",
        input_border_color_dark="rgba(201,168,76,0.2)",
        input_placeholder_color="#8892a4",
        input_placeholder_color_dark="#8892a4",
        shadow_drop="0 4px 24px rgba(0,0,0,0.4)",
        shadow_drop_lg="0 8px 40px rgba(0,0,0,0.5)",
        table_even_background_fill="#0a1628",
        table_even_background_fill_dark="#0a1628",
        table_odd_background_fill="#112240",
        table_odd_background_fill_dark="#112240",
    ))
