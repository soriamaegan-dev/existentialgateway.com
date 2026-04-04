import gradio as gr
import requests
import os
import pandas as pd
import io
import plotly.express as px
import plotly.graph_objects as go
from plotly.subplots import make_subplots
from pypdf import PdfReader
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from reportlab.lib.pagesizes import letter
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import inch
from reportlab.lib import colors
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
import tempfile
from datetime import datetime

HF_TOKEN = os.environ.get("HF_TOKEN", "")

DISCLAIMER = """
> **MEDICAL DISCLAIMER**: This tool is AI-generated for administrative and research purposes only.
> This tool does NOT provide medical diagnoses, treatment recommendations, or clinical advice.
> All findings must be reviewed by licensed medical professionals.
> This platform is NOT HIPAA compliant. Never upload identifiable patient data.
> © 2026 Existential Gateway, LLC. All Rights Reserved. Proprietary Software.
"""

WAIT_MSG = "*Results take approximately 1-2 minutes to generate. Please do not click multiple times.*"

WATERMARK = """
---
© 2026 Existential Gateway, LLC | AI Healthcare and Medical Analyzer
Unauthorized reproduction strictly prohibited. Licensing: existentialgateway@gmail.com
---
"""

PII_WARNING = """
> **HIPAA NOTICE**: Before uploading any data you MUST remove ALL protected health information (PHI)
> including: Patient Names | DOB | SSN | Medical Record Numbers | Address |
> Phone Number | Email | Biometric Identifiers | Device Identifiers.
> This platform is NOT HIPAA compliant. Only upload fully de-identified data.
"""



def fetch_url_content(text):
    """Detect URLs in user text, fetch their content, and return as context."""
    import re
    urls = re.findall(r'https?://[^\s<>"{}|\^`\[\]]+', text)
    if not urls:
        return ""
    fetched = []
    for url in urls[:3]:  # limit to 3 URLs
        try:
            headers = {
                "User-Agent": "Mozilla/5.0 (compatible; ExistentialGateway/1.0)",
                "Accept": "text/html,application/xhtml+xml,application/xml;q=0.9,*/*;q=0.8"
            }
            r = requests.get(url, timeout=15, headers=headers)
            if r.status_code == 200:
                from html.parser import HTMLParser
                class TextExtractor(HTMLParser):
                    def __init__(self):
                        super().__init__()
                        self.text = []
                        self.skip = False
                    def handle_starttag(self, tag, attrs):
                        if tag in ('script', 'style', 'nav', 'footer', 'header'):
                            self.skip = True
                    def handle_endtag(self, tag):
                        if tag in ('script', 'style', 'nav', 'footer', 'header'):
                            self.skip = False
                    def handle_data(self, data):
                        if not self.skip and data.strip():
                            self.text.append(data.strip())
                parser = TextExtractor()
                parser.feed(r.text)
                page_text = ' '.join(parser.text)[:3000]
                fetched.append(f"[URL: {url}]\n{page_text}")
        except Exception as e:
            fetched.append(f"[URL: {url}] Could not fetch: {str(e)}")
    if fetched:
        return "\n\nUSER-PROVIDED URL CONTENT (analyze this as part of your response):\n" + "\n---\n".join(fetched)
    return ""



def enrich_with_urls(text):
    """If text contains URLs, fetch and append their content."""
    if not text or 'http' not in text:
        return text
    url_content = fetch_url_content(text)
    if url_content:
        return text + url_content
    return text


def query_llm(prompt):
    API_KEY = os.environ.get("OPENAI_API_KEY", "")
    headers = {"Authorization": f"Bearer {API_KEY}", "Content-Type": "application/json"}
    msgs = [{"role": "system", "content": HEALTHCARE_SYSTEM_PROMPT}, {"role": "user", "content": prompt}]
    payload = {"model": "gpt-4o", "max_tokens": 4000, "messages": msgs}
    try:
        response = requests.post("https://api.openai.com/v1/chat/completions", headers=headers, json=payload, timeout=240)
        result = response.json()
        if "choices" in result:
            return result["choices"][0]["message"]["content"]
        return f"API Error: {result}"
    except Exception as e:
        return f"Error: {str(e)}"


def load_data(file):
    if file is None:
        return None, "No file uploaded."
    try:
        path = file.name if hasattr(file, "name") else file
        if path.endswith(".csv"):
            df = pd.read_csv(path)
        elif path.endswith((".xlsx", ".xls")):
            df = pd.read_excel(path)
        elif path.endswith(".json"):
            df = pd.read_json(path)
        elif path.endswith(".pdf"):
            reader = PdfReader(path)
            text = "\n".join(page.extract_text() or "" for page in reader.pages)
            return None, text
        else:
            return None, "Unsupported file type."
        return df, None
    except Exception as e:
        return None, f"Error loading file: {str(e)}"


# ─── Tab 1: Data Upload and Overview ─────────────────────────────────────────

def analyze_overview(file, custom_prompt=""):
    if file is None:
        return "Please upload a file.", "", ""
    df, err = load_data(file)
    if err and df is None:
        if "Unsupported" in err or "Error loading" in err:
            return err, "", ""
        text_preview = err[:3000]
        prompt = f"""You are an expert healthcare data analyst. The user uploaded a PDF healthcare document. Write as the professional analyst who performed this work. State all findings directly with specific numbers, percentages, and comparisons. NEVER say 'this appears to be', 'the data seems to show', 'the dataset contains', or 'it looks like'. Instead, deliver results confidently: 'Average claim cost was $4,230, a 12% increase year-over-year' or 'Readmission rates of 18.3% in Cardiology exceeded the facility average by 4.1 percentage points.' Be direct, precise, and actionable. Present findings as YOUR analysis, not a description of a dataset. Content (truncated):
{text_preview}

Present your professional findings:
1. What type of healthcare document or dataset is this?
2. Key metrics and data elements present
3. Time period and population covered
4. Summary of main findings or data categories
5. Data quality observations
6. Recommended next analysis steps"""
        result = query_llm(prompt)
        return result + "\n\n" + WATERMARK, "PDF — text extracted", "See analysis above"

    rows, cols = df.shape
    dtypes = df.dtypes.to_string()
    missing = df.isnull().sum()
    missing_str = missing[missing > 0].to_string() if missing.any() else "None"
    preview = df.head(10).to_string()

    # Multi-dataset mode: detect Source_Dataset column
    if "Source_Dataset" in df.columns:
        sources = df["Source_Dataset"].unique().tolist()
        per_ds = []
        for src in sources:
            sub = df[df["Source_Dataset"] == src].drop(columns=["Source_Dataset"])
            sub_rows = len(sub)
            sub_desc = sub.describe(include="all").to_string()
            sub_missing = sub.isnull().sum()
            sub_missing_str = sub_missing[sub_missing > 0].to_string() if sub_missing.any() else "None"
            sub_preview = sub.head(10).to_string()
            cat_cols = sub.select_dtypes(include=["object", "category"]).columns.tolist()
            cat_dist = ""
            for cc in cat_cols[:6]:
                vc = sub[cc].value_counts().head(8).to_string()
                cat_dist += f"  {cc}:\n{vc}\n"
            num_cols = sub.select_dtypes(include="number").columns.tolist()
            num_summary = ""
            for nc in num_cols[:10]:
                num_summary += f"  {nc}: mean={sub[nc].mean():.2f}, median={sub[nc].median():.2f}, min={sub[nc].min():.2f}, max={sub[nc].max():.2f}\n"
            per_ds.append(f"""### Dataset: {src} ({sub_rows} rows)
    Columns: {list(sub.columns)}
    Numeric summaries:
    {num_summary}
    Statistics:
    {sub_desc}
    Missing values:
    {sub_missing_str}
    Top categories:
    {cat_dist}
    Sample rows:
    {sub_preview}""")
        combined_summary = "\n\n".join(per_ds)
        prompt = f"""You are an expert healthcare data analyst. You have {len(sources)} SEPARATE datasets merged together. Write as the professional analyst who performed this work. State all findings directly with specific numbers, percentages, and comparisons. NEVER say 'this appears to be', 'the data seems to show', 'the dataset contains', or 'it looks like'. Instead, deliver results confidently: 'Average claim cost was $4,230, a 12% increase year-over-year' or 'Readmission rates of 18.3% in Cardiology exceeded the facility average by 4.1 percentage points.' Be direct, precise, and actionable. Present findings as YOUR analysis, not a description of a dataset.     Each is identified by Source_Dataset. You MUST analyze each individually AND compare them against each other.
    
    DETAILED DATA FROM EACH DATASET:
    
    {combined_summary}
    
    Provide: 1) Individual summary of EACH dataset with specific numbers 2) Side-by-side comparison of matching metrics 3) Trends and changes between datasets 4) Top 5 insights from comparing datasets 5) Anomalies and discrepancies 6) Recommendations 7) Summary comparison table. Reference specific numbers from EVERY dataset."""
        result = query_llm(prompt)
        return result + "\n\n" + WATERMARK, None, None
    prompt = f"""You are an expert healthcare data analyst. Write as the professional analyst who performed this work. State all findings directly with specific numbers, percentages, and comparisons. NEVER say 'this appears to be', 'the data seems to show', 'the dataset contains', or 'it looks like'. Instead, deliver results confidently: 'Average claim cost was $4,230, a 12% increase year-over-year' or 'Readmission rates of 18.3% in Cardiology exceeded the facility average by 4.1 percentage points.' Be direct, precise, and actionable. Present findings as YOUR analysis, not a description of a dataset. Analyze this healthcare dataset overview:

Rows: {rows} | Columns: {cols}
Column types:\n{dtypes}
Missing values:\n{missing_str}
First 10 rows:\n{preview}

Provide:
1. What type of healthcare data this appears to be
2. Patient population and time period if identifiable
3. Key clinical or administrative metrics present
4. Data quality assessment
5. Key columns and their significance
6. Recommended analyses to run
7. Any immediate patterns or concerns

Note: Assume all data has been properly de-identified."""

    result = query_llm(prompt)
    stats = f"**Rows:** {rows} | **Columns:** {cols}\n\n**Missing Values:**\n```\n{missing_str}\n```"
    return result + "\n\n" + WATERMARK, stats, f"```\n{preview}\n```"


# ─── Tab 2: Patient Outcome Analysis ─────────────────────────────────────────

def analyze_outcomes(file):
    if file is None:
        return "Please upload de-identified patient data.", None
    df, err = load_data(file)
    if err:
        return err, None

    rows, cols = df.shape
    preview = df.head(20).to_string()
    desc = df.describe(include="all").to_string()

    prompt = f"""You are a senior healthcare outcomes analyst presenting your findings on patient outcomes. Write as the professional analyst who performed this work. State all findings directly with specific numbers, percentages, and comparisons. NEVER say 'this appears to be', 'the data seems to show', 'the dataset contains', or 'it looks like'. Instead, deliver results confidently: 'Average claim cost was $4,230, a 12% increase year-over-year' or 'Readmission rates of 18.3% in Cardiology exceeded the facility average by 4.1 percentage points.' Be direct, precise, and actionable. Present findings as YOUR analysis, not a description of a dataset. Write as the professional who performed this analysis — state results directly with specific numbers. Never say "the data appears to show." Instead: "30-day readmission rates averaged 18.3% across all departments, with Cardiology at 24.1% significantly exceeding the facility mean."

Analyze this de-identified patient data:

Shape: {rows} rows x {cols} columns
Columns: {list(df.columns)}
Statistical Summary:\n{desc}
Sample Data:\n{preview}

Provide a comprehensive patient outcome analysis:
1. **Treatment Effectiveness** — which treatments or interventions show best outcomes
2. **Readmission Rate Analysis** — readmission patterns and risk factors
3. **Length of Stay Optimization** — average LOS, outliers, and optimization opportunities
4. **Mortality Rate Analysis** — by condition, department, or demographic if available
5. **Recovery Time Benchmarking** — compare against standard benchmarks
6. **High Risk Patient Segments** — groups with poorest outcomes
7. **Quality Improvement Opportunities** — specific areas to target
8. **Outcome Summary** — professional written summary with key findings

Use specific numbers. Assume all data is fully de-identified."""

    result = query_llm(prompt)

    fig = None
    try:
        numeric_cols = df.select_dtypes(include="number").columns.tolist()
        cat_cols = df.select_dtypes(include="object").columns.tolist()
        if cat_cols and numeric_cols:
            top = df.groupby(cat_cols[0])[numeric_cols[0]].mean().sort_values(ascending=False).head(10)
            fig = go.Figure(go.Bar(
                x=top.index.tolist(),
                y=top.values,
                marker_color="#2980b9"
            ))
            fig.update_layout(
                title=f"Average {numeric_cols[0]} by {cat_cols[0]}",
                template="plotly_dark",
                height=400
            )
        elif numeric_cols:
            fig = px.histogram(df, x=numeric_cols[0], template="plotly_dark",
                               title=f"{numeric_cols[0]} Distribution", nbins=30)
            fig.update_layout(height=400)
    except Exception:
        fig = None

    return result + "\n\n" + WATERMARK, fig


# ─── Tab 3: Disease Pattern Detection ────────────────────────────────────────

def analyze_disease_patterns(file):
    if file is None:
        return "Please upload epidemiological data.", None
    df, err = load_data(file)
    if err:
        return err, None

    rows, cols = df.shape
    preview = df.head(20).to_string()
    desc = df.describe(include="all").to_string()

    prompt = f"""You are a senior epidemiologist presenting disease pattern findings to a public health team. Write as the professional analyst who performed this work. State all findings directly with specific numbers, percentages, and comparisons. NEVER say 'this appears to be', 'the data seems to show', 'the dataset contains', or 'it looks like'. Instead, deliver results confidently: 'Average claim cost was $4,230, a 12% increase year-over-year' or 'Readmission rates of 18.3% in Cardiology exceeded the facility average by 4.1 percentage points.' Be direct, precise, and actionable. Present findings as YOUR analysis, not a description of a dataset. Write as the analyst who conducted this investigation — state findings directly with specific numbers and prevalence rates. Never say "the data suggests" — instead: "COPD prevalence was 12.4% among patients over 65, a 3.1 percentage point increase from the prior year."

Analyze this epidemiological data:

Shape: {rows} rows x {cols} columns
Columns: {list(df.columns)}
Statistical Summary:\n{desc}
Sample Data:\n{preview}

Provide a comprehensive disease pattern analysis:
1. **Disease Spread Patterns** — how and where disease is spreading
2. **Outbreak Risk Assessment** — current outbreak probability (Low/Medium/High)
3. **Geographic Distribution** — hotspots and affected areas
4. **At-Risk Population Identification** — most vulnerable demographic groups
5. **Seasonal Illness Patterns** — time-based trends and peaks
6. **Infection Rate Trend Analysis** — rising, falling, or stable rates
7. **Transmission Pattern Analysis** — likely transmission routes
8. **Public Health Recommendations** — specific intervention strategies

Use specific numbers. Be professional and precise."""

    result = query_llm(prompt)

    fig = None
    try:
        numeric_cols = df.select_dtypes(include="number").columns.tolist()
        cat_cols = df.select_dtypes(include="object").columns.tolist()
        if numeric_cols:
            fig = go.Figure()
            for col in numeric_cols[:3]:
                fig.add_trace(go.Scatter(
                    y=df[col],
                    mode="lines+markers",
                    name=col
                ))
            fig.update_layout(
                title="Disease Rate Trends",
                template="plotly_dark",
                height=400,
                yaxis_title="Rate"
            )
        elif cat_cols:
            top = df[cat_cols[0]].value_counts().head(10)
            fig = go.Figure(go.Bar(
                x=top.index.tolist(),
                y=top.values,
                marker_color="#e74c3c"
            ))
            fig.update_layout(
                title=f"Disease Distribution by {cat_cols[0]}",
                template="plotly_dark",
                height=400
            )
    except Exception:
        fig = None

    return result + "\n\n" + WATERMARK, fig


# ─── Tab 4: Hospital Resource Planning ───────────────────────────────────────

def analyze_resources(file):
    if file is None:
        return "Please upload capacity and utilization data.", None
    df, err = load_data(file)
    if err:
        return err, None

    rows, cols = df.shape
    preview = df.head(20).to_string()
    desc = df.describe(include="all").to_string()

    prompt = f"""You are a senior hospital operations analyst presenting resource utilization findings to administration. Write as the professional analyst who performed this work. State all findings directly with specific numbers, percentages, and comparisons. NEVER say 'this appears to be', 'the data seems to show', 'the dataset contains', or 'it looks like'. Instead, deliver results confidently: 'Average claim cost was $4,230, a 12% increase year-over-year' or 'Readmission rates of 18.3% in Cardiology exceeded the facility average by 4.1 percentage points.' Be direct, precise, and actionable. Present findings as YOUR analysis, not a description of a dataset. Write as the professional who performed this analysis — deliver findings with specific numbers, percentages, and comparisons. Never say "it appears that" — instead: "Bed occupancy reached 94.2% in Q3, exceeding the 85% efficiency threshold by 9.2 percentage points."

Analyze this operational data:

Shape: {rows} rows x {cols} columns
Columns: {list(df.columns)}
Statistical Summary:\n{desc}
Sample Data:\n{preview}

Provide a comprehensive hospital resource planning analysis:
1. **Bed Allocation Optimization** — current utilization and recommended allocation
2. **Staff Scheduling Optimization** — staffing gaps, surpluses, and recommendations
3. **Equipment Utilization Analysis** — over and under-utilized equipment
4. **Emergency Department Flow** — patient flow bottlenecks and solutions
5. **Operating Room Efficiency** — OR utilization rates and improvement opportunities
6. **Supply Chain Recommendations** — inventory optimization suggestions
7. **Capacity Planning** — future capacity needs based on trends
8. **Implementation Priorities** — top 5 actions to improve resource efficiency

Be specific with percentages and numbers."""

    result = query_llm(prompt)

    fig = None
    try:
        numeric_cols = df.select_dtypes(include="number").columns.tolist()
        cat_cols = df.select_dtypes(include="object").columns.tolist()
        if numeric_cols and cat_cols:
            fig = px.bar(
                df.head(15),
                x=cat_cols[0],
                y=numeric_cols[0],
                template="plotly_dark",
                title=f"Resource Utilization: {numeric_cols[0]} by {cat_cols[0]}",
                color_discrete_sequence=["#27ae60"]
            )
            fig.update_layout(height=400)
        elif numeric_cols:
            fig = make_subplots(rows=1, cols=min(2, len(numeric_cols)),
                                subplot_titles=numeric_cols[:2])
            for i, col in enumerate(numeric_cols[:2], 1):
                fig.add_trace(go.Histogram(x=df[col], name=col, nbinsx=20), row=1, col=i)
            fig.update_layout(title="Resource Data Distribution",
                              template="plotly_dark", height=400)
    except Exception:
        fig = None

    return result + "\n\n" + WATERMARK, fig


# ─── Tab 5: Medical Billing Audit ────────────────────────────────────────────

def analyze_billing(file):
    if file is None:
        return "Please upload billing data.", None
    df, err = load_data(file)
    if err:
        return err, None

    rows, cols = df.shape
    preview = df.head(30).to_string()
    desc = df.describe(include="all").to_string()

    signals = []
    for col in df.columns:
        cl = col.lower()
        if "amount" in cl or "charge" in cl or "bill" in cl or "cost" in cl:
            try:
                q99 = df[col].quantile(0.99)
                high = (df[col] > q99).sum()
                if high > 0:
                    signals.append(f"- {high} charges above 99th percentile in '{col}' (>{q99:.2f})")
            except Exception:
                pass
    for col in df.columns:
        if df[col].duplicated().sum() > 0:
            dups = df[col].duplicated().sum()
            signals.append(f"- {dups} duplicate values detected in '{col}'")
            break

    signal_str = "\n".join(signals) if signals else "No automatic signals computed."

    prompt = f"""You are a senior medical billing auditor presenting audit findings. Write as the professional analyst who performed this work. State all findings directly with specific numbers, percentages, and comparisons. NEVER say 'this appears to be', 'the data seems to show', 'the dataset contains', or 'it looks like'. Instead, deliver results confidently: 'Average claim cost was $4,230, a 12% increase year-over-year' or 'Readmission rates of 18.3% in Cardiology exceeded the facility average by 4.1 percentage points.' Be direct, precise, and actionable. Present findings as YOUR analysis, not a description of a dataset. Write as the auditor who conducted this review — state findings directly with specific dollar amounts, percentages, and variance analyses. Never say "the billing data appears to show" — instead: "Average charge per encounter was $5,240, with 14.7% of claims exceeding the expected reimbursement threshold by more than $2,000."

Analyze:

Rows: {rows} | Columns: {list(df.columns)}
Statistical Summary:\n{desc}
Sample Data:\n{preview}
Auto-detected signals:\n{signal_str}

Provide a comprehensive medical billing audit:
1. **Overall Fraud Risk Score** — Low/Medium/High with justification
2. **Duplicate Billing Detection** — exact or near-duplicate claims
3. **Upcoding Identification** — charges coded at higher level than warranted
4. **Unbundling Detection** — services billed separately that should be bundled
5. **Coding Compliance** — adherence to standard coding practices
6. **Unusual Charge Patterns** — statistical outliers and anomalies
7. **Revenue Cycle Optimization** — opportunities to improve legitimate revenue
8. **Priority Audit List** — top 5 items requiring immediate review
9. **Compliance Report** — executive summary for compliance team

Be specific. Use numbers. Flag all red flags clearly."""

    result = query_llm(prompt)

    fig = None
    try:
        numeric_cols = df.select_dtypes(include="number").columns.tolist()
        if numeric_cols:
            col = numeric_cols[0]
            q99 = df[col].quantile(0.99)
            flag_colors = ["red" if v > q99 else "#3498db" for v in df[col]]
            fig = go.Figure()
            fig.add_trace(go.Bar(
                x=list(range(len(df))),
                y=df[col],
                marker_color=flag_colors,
                name=col
            ))
            fig.add_hline(y=q99, line_dash="dash", line_color="orange",
                          annotation_text="99th Percentile Flag Threshold")
            fig.update_layout(
                title=f"Billing Amounts — Red = Flagged for Review",
                template="plotly_dark",
                height=400
            )
    except Exception:
        fig = None

    return result + "\n\n" + WATERMARK, fig


# ─── Tab 6: Clinical Trial Analysis ──────────────────────────────────────────

def analyze_trials(file):
    if file is None:
        return "Please upload clinical trial data.", None
    df, err = load_data(file)
    if err and df is None:
        if "Unsupported" in err or "Error" in err:
            return err, None
        prompt = f"""You are a senior clinical research analyst presenting trial findings to the research committee. Write as the professional analyst who performed this work. State all findings directly with specific numbers, percentages, and comparisons. NEVER say 'this appears to be', 'the data seems to show', 'the dataset contains', or 'it looks like'. Instead, deliver results confidently: 'Average claim cost was $4,230, a 12% increase year-over-year' or 'Readmission rates of 18.3% in Cardiology exceeded the facility average by 4.1 percentage points.' Be direct, precise, and actionable. Present findings as YOUR analysis, not a description of a dataset. Write as the biostatistician who analyzed these results — report findings with specific p-values, confidence intervals, effect sizes, and enrollment figures. Never say "the trial data suggests" — instead: "The treatment arm showed a 23.4% improvement in primary endpoint (p=0.003, 95% CI: 15.2-31.6%), with 312 of 400 enrolled participants completing the protocol."

Analyze this trial data:

{err[:3000]}

Provide:
1. **Trial Overview** — type, phase, and purpose of the trial
2. **Statistical Significance** — p-values and confidence intervals if present
3. **Efficacy Analysis** — primary and secondary endpoint results
4. **Adverse Event Analysis** — safety profile and notable events
5. **Patient Cohort Analysis** — demographic breakdown and balance
6. **Regulatory Compliance** — alignment with ICH/FDA/EMA guidelines
7. **Comparative Analysis** — treatment vs control group outcomes
8. **Conclusions and Recommendations** — overall trial assessment

Be precise and scientifically rigorous."""
        result = query_llm(prompt)
        return result + "\n\n" + WATERMARK, None

    rows, cols = df.shape
    preview = df.head(20).to_string()
    desc = df.describe(include="all").to_string()

    prompt = f"""You are a senior clinical research analyst presenting trial findings to the research committee. Write as the professional analyst who performed this work. State all findings directly with specific numbers, percentages, and comparisons. NEVER say 'this appears to be', 'the data seems to show', 'the dataset contains', or 'it looks like'. Instead, deliver results confidently: 'Average claim cost was $4,230, a 12% increase year-over-year' or 'Readmission rates of 18.3% in Cardiology exceeded the facility average by 4.1 percentage points.' Be direct, precise, and actionable. Present findings as YOUR analysis, not a description of a dataset. Write as the biostatistician who analyzed these results — report findings with specific p-values, confidence intervals, effect sizes, and enrollment figures. Never say "the trial data suggests" — instead: "The treatment arm showed a 23.4% improvement in primary endpoint (p=0.003, 95% CI: 15.2-31.6%), with 312 of 400 enrolled participants completing the protocol."

Analyze this trial data:

Shape: {rows} rows x {cols} columns
Columns: {list(df.columns)}
Statistical Summary:\n{desc}
Sample Data:\n{preview}

Provide a comprehensive clinical trial analysis:
1. **Trial Results Overview** — primary findings and outcomes
2. **Statistical Significance Assessment** — significance of key results
3. **Adverse Event Analysis** — frequency and severity of adverse events
4. **Efficacy Comparison** — treatment effectiveness across groups
5. **Patient Cohort Analysis** — cohort characteristics and balance
6. **Regulatory Compliance Check** — alignment with standard guidelines
7. **Subgroup Analysis** — outcomes across demographic subgroups
8. **Conclusions** — overall trial assessment and next steps

Be scientifically rigorous and precise."""

    result = query_llm(prompt)

    fig = None
    try:
        numeric_cols = df.select_dtypes(include="number").columns.tolist()
        cat_cols = df.select_dtypes(include="object").columns.tolist()
        if cat_cols and numeric_cols:
            fig = px.box(
                df,
                x=cat_cols[0],
                y=numeric_cols[0],
                color=cat_cols[0],
                template="plotly_dark",
                title=f"Trial Outcome Distribution: {numeric_cols[0]} by {cat_cols[0]}"
            )
            fig.update_layout(height=420)
        elif numeric_cols and len(numeric_cols) >= 2:
            corr = df[numeric_cols].corr()
            fig = px.imshow(corr, text_auto=True, template="plotly_dark",
                            title="Trial Variable Correlation Matrix",
                            color_continuous_scale="RdBu_r")
            fig.update_layout(height=420)
    except Exception:
        fig = None

    return result + "\n\n" + WATERMARK, fig


# ─── Tab 7: Visualizations and Charts ────────────────────────────────────────

def generate_viz(file, chart_type, x_col, y_col, color_col):
    if file is None:
        return None, "Please upload a data file."
    df, err = load_data(file)
    if err:
        return None, err

    cols = df.columns.tolist()
    x = x_col if x_col in cols else cols[0]
    y = y_col if y_col in cols else (cols[1] if len(cols) > 1 else cols[0])
    color = color_col if color_col in cols else None

    try:
        if chart_type == "Bar Chart":
            fig = px.bar(df, x=x, y=y, color=color, template="plotly_dark",
                         title=f"{y} by {x}")
        elif chart_type == "Line Chart":
            fig = px.line(df, x=x, y=y, color=color, template="plotly_dark",
                          title=f"{y} over {x}")
        elif chart_type == "Pie Chart":
            fig = px.pie(df, names=x, values=y, template="plotly_dark",
                         title=f"{y} Distribution by {x}")
        elif chart_type == "Scatter Plot":
            fig = px.scatter(df, x=x, y=y, color=color, template="plotly_dark",
                             title=f"{x} vs {y}")
        elif chart_type == "Heatmap":
            numeric = df.select_dtypes(include="number")
            corr = numeric.corr()
            fig = px.imshow(corr, text_auto=True, template="plotly_dark",
                            title="Correlation Heatmap", color_continuous_scale="RdBu_r")
        elif chart_type == "Box Plot":
            fig = px.box(df, x=x if df[x].dtype == "object" else None,
                         y=y, color=color, template="plotly_dark",
                         title=f"{y} Distribution")
        elif chart_type == "Histogram":
            fig = px.histogram(df, x=x, color=color, template="plotly_dark",
                               title=f"{x} Frequency Distribution", nbins=40)
        elif chart_type == "Area Chart":
            fig = px.area(df, x=x, y=y, color=color, template="plotly_dark",
                          title=f"{y} Area Chart")
        else:
            fig = px.bar(df, x=x, y=y, template="plotly_dark")

        fig.update_layout(height=500, margin=dict(l=40, r=40, t=60, b=40))
        return fig, f"Chart generated: {chart_type} — {y} by {x}"
    except Exception as e:
        return None, f"Chart error: {str(e)}"


def get_columns(file):
    if file is None:
        return gr.Dropdown(choices=[], allow_custom_value=True), gr.Dropdown(choices=[], allow_custom_value=True), gr.Dropdown(choices=[], allow_custom_value=True)
    df, err = load_data(file)
    if err or df is None:
        return gr.Dropdown(choices=[], allow_custom_value=True), gr.Dropdown(choices=[], allow_custom_value=True), gr.Dropdown(choices=[], allow_custom_value=True)
    cols = df.columns.tolist()
    return gr.Dropdown(choices=cols), gr.Dropdown(choices=cols), gr.Dropdown(choices=cols)


# ─── Tab 8: Reports and Presentations ────────────────────────────────────────

def generate_report(file, report_type, output_format="PDF Report"):
    if file is None:
        return "Please upload a data file.", None
    df, err = load_data(file)
    if err and df is None:
        return err, None

    if df is not None:
        preview = df.head(20).to_string()
        desc = df.describe(include="all").to_string()
        cols = list(df.columns)
        rows = len(df)
        data_summary = f"Rows: {rows} | Columns: {cols}\nStatistics:\n{desc}\nSample:\n{preview}"
    else:
        data_summary = err[:3000]

    prompt = f"""You are a senior healthcare analyst who has completed a thorough analysis and is now writing a formal {report_type} for stakeholders. Write as the professional analyst who performed this work. State all findings directly with specific numbers, percentages, and comparisons. NEVER say 'this appears to be', 'the data seems to show', 'the dataset contains', or 'it looks like'. Instead, deliver results confidently: 'Average claim cost was $4,230, a 12% increase year-over-year' or 'Readmission rates of 18.3% in Cardiology exceeded the facility average by 4.1 percentage points.' Be direct, precise, and actionable. Present findings as YOUR analysis, not a description of a dataset. Write as the professional who performed this work. Present all findings with specific numbers, percentages, and trends. Never describe the data abstractly — deliver concrete results.
Analyze this healthcare data and produce the {report_type}:

{data_summary}

Structure:
1. **Executive Summary** (3-5 sentences)
2. **Data Overview** — what the healthcare data contains
3. **Key Clinical Findings** — top 5 findings with specific numbers
4. **Patient Outcome Analysis** — notable outcomes and trends
5. **Resource and Operational Assessment** — efficiency and utilization
6. **Risk and Quality Assessment** — areas of concern
7. **Recommendations** — 5 specific actionable recommendations
8. **Conclusion**

Write professionally for a healthcare administration audience. Use specific numbers."""

    report_text = query_llm(prompt)

    output_path = None
    try:
        if report_type == "PowerPoint Presentation":
            output_path = _make_pptx(report_text, df if df is not None else None)
        else:
            output_path = _make_pdf(report_text)
    except Exception as e:
        report_text += f"\n\n[Document generation error: {str(e)}]"

    return report_text + "\n\n" + WATERMARK, output_path


def _make_pptx(text, df):
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    def add_slide(title_text, body_text, bg_color=(5, 15, 25)):
        slide = prs.slides.add_slide(blank_layout)
        bg = slide.background
        bg.fill.solid()
        bg.fill.fore_color.rgb = RGBColor(*bg_color)
        tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(1.2))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = RGBColor(52, 152, 219)
        body_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.7), Inches(12.3), Inches(5.5))
        tf2 = body_box.text_frame
        tf2.word_wrap = True
        for line in body_text.split("\n")[:20]:
            line = line.strip()
            if not line:
                continue
            para = tf2.add_paragraph()
            para.text = line
            para.font.size = Pt(14)
            para.font.color.rgb = RGBColor(220, 235, 245)

    title_slide = prs.slides.add_slide(blank_layout)
    bg = title_slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(3, 10, 20)
    tb = title_slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11), Inches(1.5))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = "AI Healthcare & Medical Analysis Report"
    p.font.size = Pt(34)
    p.font.bold = True
    p.font.color.rgb = RGBColor(52, 152, 219)
    tb2 = title_slide.shapes.add_textbox(Inches(1), Inches(4.2), Inches(11), Inches(1))
    tf2 = tb2.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = f"Generated: {datetime.now().strftime('%B %d, %Y')} | AI Healthcare Analyzer"
    p2.font.size = Pt(18)
    p2.font.color.rgb = RGBColor(150, 190, 220)

    sections = []
    current_title = "Overview"
    current_body = []
    for line in text.split("\n"):
        stripped = line.strip()
        if stripped.startswith("**") and stripped.endswith("**") and len(stripped) > 4:
            if current_body:
                sections.append((current_title, "\n".join(current_body)))
            current_title = stripped.strip("*").strip()
            current_body = []
        elif stripped.startswith("## "):
            if current_body:
                sections.append((current_title, "\n".join(current_body)))
            current_title = stripped[3:].strip()
            current_body = []
        else:
            current_body.append(line)
    if current_body:
        sections.append((current_title, "\n".join(current_body)))

    bg_colors = [(5, 15, 25), (5, 20, 30), (10, 15, 25), (5, 25, 20), (15, 10, 25)]
    for i, (title, body) in enumerate(sections[:8]):
        add_slide(title, body, bg_colors[i % len(bg_colors)])

    if df is not None:
        stats_body = f"Total Records: {len(df)}\nColumns: {len(df.columns)}\n\nColumn Overview:\n"
        for col in df.columns[:10]:
            stats_body += f"  • {col} ({df[col].dtype})\n"
        add_slide("Data Statistics", stats_body)

    closing = prs.slides.add_slide(blank_layout)
    bg = closing.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(3, 10, 20)
    tb = closing.shapes.add_textbox(Inches(1), Inches(2.8), Inches(11), Inches(2))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Report Complete"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(52, 152, 219)
    p2 = tf.add_paragraph()
    p2.text = "© 2026 Existential Gateway, LLC | AI Healthcare and Medical Analyzer"
    p2.font.size = Pt(14)
    p2.font.color.rgb = RGBColor(150, 190, 220)

    tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".pptx")
    prs.save(tmp.name)
    return tmp.name


def _make_pdf(text):
    tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".pdf")
    doc = SimpleDocTemplate(tmp.name, pagesize=letter,
                            leftMargin=inch, rightMargin=inch,
                            topMargin=inch, bottomMargin=inch)
    styles = getSampleStyleSheet()
    title_style = ParagraphStyle("Title2", parent=styles["Title"],
                                 fontSize=20, spaceAfter=20,
                                 textColor=colors.HexColor("#1a5276"))
    heading_style = ParagraphStyle("H2", parent=styles["Heading2"],
                                   fontSize=14, spaceAfter=8,
                                   textColor=colors.HexColor("#2980b9"))
    body_style = ParagraphStyle("Body2", parent=styles["Normal"],
                                fontSize=11, spaceAfter=6, leading=16)
    story = []
    story.append(Paragraph("AI Healthcare & Medical Analysis Report", title_style))
    story.append(Paragraph(
        f"Generated: {datetime.now().strftime('%B %d, %Y')} | AI Healthcare Analyzer",
        body_style))
    story.append(Spacer(1, 0.2 * inch))

    for line in text.split("\n"):
        stripped = line.strip()
        if not stripped:
            story.append(Spacer(1, 0.1 * inch))
            continue
        if stripped.startswith("**") and stripped.endswith("**"):
            story.append(Paragraph(stripped.strip("*"), heading_style))
        elif stripped.startswith("## "):
            story.append(Paragraph(stripped[3:], heading_style))
        elif stripped.startswith("# "):
            story.append(Paragraph(stripped[2:], title_style))
        else:
            safe = stripped.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")
            story.append(Paragraph(safe, body_style))

    story.append(Spacer(1, 0.3 * inch))
    story.append(Paragraph(
        "© 2026 Existential Gateway, LLC | AI Healthcare and Medical Analyzer | existentialgateway@gmail.com",
        ParagraphStyle("footer", parent=styles["Normal"], fontSize=9,
                       textColor=colors.grey)))
    doc.build(story)
    return tmp.name


# ─── Tab 9: AI Data Chat ─────────────────────────────────────────────────────

_chat_df = None


def upload_chat_data(file):
    global _chat_df
    if file is None:
        _chat_df = None
        return "No file uploaded."
    df, err = load_data(file)
    if err and df is None:
        return err
    _chat_df = df
    return f"Data loaded: {len(df)} rows x {len(df.columns)} columns\nColumns: {list(df.columns)}"


def chat_with_data(message, history):
    global _chat_df

    data_context = ""
    if _chat_df is not None:
        desc = _chat_df.describe(include="all").to_string()
        preview = _chat_df.head(10).to_string()
        # Compute additional stats for better answers
        num_cols = _chat_df.select_dtypes(include="number").columns.tolist()
        extra_stats = ""
        for col in num_cols[:15]:
            try:
                mode_val = _chat_df[col].mode().iloc[0] if not _chat_df[col].mode().empty else "N/A"
                extra_stats += f"  {col}: mean={_chat_df[col].mean():.4f}, median={_chat_df[col].median():.4f}, mode={mode_val}, std={_chat_df[col].std():.4f}, sum={_chat_df[col].sum():.2f}\n"
            except Exception:
                pass
        cat_cols = _chat_df.select_dtypes(include=["object", "category"]).columns.tolist()
        cat_stats = ""
        for col in cat_cols[:10]:
            try:
                vc = _chat_df[col].value_counts().head(10).to_string()
                mode_val = _chat_df[col].mode().iloc[0] if not _chat_df[col].mode().empty else "N/A"
                cat_stats += f"  {col} (mode={mode_val}, {_chat_df[col].nunique()} unique):\n{vc}\n\n"
            except Exception:
                pass
        data_context = (
            f"The user has uploaded de-identified healthcare data.\n"
            f"Shape: {_chat_df.shape}\n"
            f"Columns: {list(_chat_df.columns)}\n"
            f"Statistical Summary:\n{desc}\n"
            f"Numeric Column Details (mean, median, mode, std, sum):\n{extra_stats}\n"
            f"Categorical Column Distributions:\n{cat_stats}\n"
            f"First 10 rows:\n{preview}"
        )
    else:
        data_context = "No data uploaded yet. Answer general healthcare analytics questions."

    messages = [
        {
            "role": "system",
            "content": (
                "You are a senior healthcare data analyst answering questions from your team. Write as the professional analyst who performed this work. State all findings directly with specific numbers, percentages, and comparisons. NEVER say 'this appears to be', 'the data seems to show', 'the dataset contains', or 'it looks like'. Instead, deliver results confidently: 'Average claim cost was $4,230, a 12% increase year-over-year' or 'Readmission rates of 18.3% in Cardiology exceeded the facility average by 4.1 percentage points.' Be direct, precise, and actionable. Present findings as YOUR analysis, not a description of a dataset. "
                "When asked a question, deliver the answer immediately with specific numbers — as a professional analyst would in a meeting. "
                "Do NOT explain how to calculate. Do NOT describe what columns contain. Just answer with the actual computed result. "
                "For example, if asked about net profit: \'Net profit averaged $67,688 (median: $50,575), with Cardiology generating the highest at $312,807.\' "
                "Be direct, confident, and precise. Include context like comparisons, rankings, or trends when relevant. "
                "Only provide code if explicitly asked. Keep responses concise and data-driven.\n\n"
                + data_context
            )
        }
    ]
    for user_msg, bot_msg in history:
        messages.append({"role": "user", "content": user_msg})
        messages.append({"role": "assistant", "content": bot_msg})
    messages.append({"role": "user", "content": message})

    import os
    API_KEY = os.environ.get("OPENAI_API_KEY", "")
    headers = {"Authorization": f"Bearer {API_KEY}", "Content-Type": "application/json"}
    msgs = [{"role": "system", "content": HEALTHCARE_SYSTEM_PROMPT}, {"role": "user", "content": prompt}]
    payload = {"model": "gpt-4o", "max_tokens": 4000, "messages": msgs}
    try:
        response = requests.post("https://api.openai.com/v1/chat/completions", headers=headers, json=payload, timeout=240)
        result = response.json()
        if "choices" in result:
            return result["choices"][0]["message"]["content"]
        return f"API Error: {result}"
    except Exception as e:
        return f"Error: {str(e)}"


# ─── Gradio UI ────────────────────────────────────────────────────────────────


def clean_data(file, cleaning_type="full"):
    """Clean the uploaded dataset and return a report + cleaned Excel file."""
    import pandas as pd
    import numpy as np
    import tempfile, os

    try:
        if file is None:
            return "Please upload a file.", None

        file_path = file
        if hasattr(file, 'name'):
            file_path = file.name
        elif isinstance(file, dict) and 'name' in file:
            file_path = file['name']
        elif isinstance(file, dict) and 'path' in file:
            file_path = file['path']
        file_path = str(file_path)
        if file_path.endswith('.csv'):
            df = pd.read_csv(file_path)
        else:
            df = pd.read_excel(file_path)

        original_rows = len(df)
        original_cols = len(df.columns)
        report_parts = []
        report_parts.append(f"# Data Cleaning Report\n")
        report_parts.append(f"**Original Dataset:** {original_rows} rows × {original_cols} columns\n")

        changes = []

        # Track issues found
        empty_rows = df.isna().all(axis=1).sum()
        duplicate_count = df.duplicated().sum()
        null_counts = df.isnull().sum()
        cols_with_nulls = null_counts[null_counts > 0]

        report_parts.append(f"## Issues Found\n")
        report_parts.append(f"- **Empty rows:** {empty_rows}")
        report_parts.append(f"- **Duplicate rows:** {duplicate_count}")
        report_parts.append(f"- **Columns with missing values:** {len(cols_with_nulls)}")
        for col, cnt in cols_with_nulls.items():
            report_parts.append(f"  - {col}: {cnt} nulls ({cnt/original_rows*100:.1f}%)")

        # Check for inconsistent casing in object columns
        obj_cols = df.select_dtypes(include='object').columns
        casing_issues = {}
        for col in obj_cols:
            unique_vals = df[col].dropna().unique()
            lower_map = {}
            for v in unique_vals:
                key = str(v).strip().lower()
                if key in lower_map and lower_map[key] != str(v).strip():
                    if col not in casing_issues:
                        casing_issues[col] = 0
                    casing_issues[col] += 1
                lower_map[key] = str(v).strip()
        if casing_issues:
            report_parts.append(f"- **Columns with inconsistent casing:** {len(casing_issues)}")
            for col, cnt in casing_issues.items():
                report_parts.append(f"  - {col}: {cnt} variations found")

        report_parts.append(f"\n## Cleaning Actions Performed\n")

        # 1. Remove completely empty rows
        if cleaning_type in ('full', 'duplicates'):
            before = len(df)
            df = df.dropna(how='all')
            removed = before - len(df)
            if removed > 0:
                changes.append(f"Removed {removed} completely empty rows")
                report_parts.append(f"- Removed **{removed} empty rows**")

        # 2. Remove duplicates
        if cleaning_type in ('full', 'duplicates'):
            before = len(df)
            df = df.drop_duplicates()
            removed = before - len(df)
            if removed > 0:
                changes.append(f"Removed {removed} duplicate rows")
                report_parts.append(f"- Removed **{removed} duplicate rows**")

        # 3. Standardize text casing
        if cleaning_type in ('full', 'formats'):
            for col in obj_cols:
                if col in df.columns:
                    original_unique = df[col].dropna().nunique()
                    df[col] = df[col].apply(lambda x: str(x).strip().title() if pd.notna(x) else x)
                    new_unique = df[col].dropna().nunique()
                    if new_unique < original_unique:
                        changes.append(f"Standardized casing in {col}: {original_unique} → {new_unique} unique values")
                        report_parts.append(f"- Standardized casing in **{col}**: {original_unique} → {new_unique} unique values")

        # 4. Standardize phone numbers
        if cleaning_type in ('full', 'formats'):
            import re
            for col in df.columns:
                if 'phone' in col.lower() or 'tel' in col.lower():
                    def standardize_phone(val):
                        if pd.isna(val): return val
                        digits = re.sub(r'\D', '', str(val))
                        if digits.startswith('1') and len(digits) == 11:
                            digits = digits[1:]
                        if len(digits) == 10:
                            return f"({digits[:3]}) {digits[3:6]}-{digits[6:]}"
                        return val
                    df[col] = df[col].apply(standardize_phone)
                    changes.append(f"Standardized phone format in {col}")
                    report_parts.append(f"- Standardized phone format in **{col}**")

        # 5. Fix missing values
        if cleaning_type in ('full', 'nulls'):
            numeric_cols = df.select_dtypes(include='number').columns
            for col in numeric_cols:
                null_count = df[col].isnull().sum()
                if null_count > 0:
                    median_val = df[col].median()
                    df[col] = df[col].fillna(median_val)
                    changes.append(f"Filled {null_count} nulls in {col} with median ({median_val:.2f})")
                    report_parts.append(f"- Filled **{null_count} nulls** in **{col}** with median ({median_val:.2f})")
            for col in obj_cols:
                if col in df.columns:
                    null_count = df[col].isnull().sum()
                    if null_count > 0:
                        df[col] = df[col].fillna('Unknown')
                        changes.append(f"Filled {null_count} nulls in {col} with 'Unknown'")
                        report_parts.append(f"- Filled **{null_count} nulls** in **{col}** with 'Unknown'")

        # 6. Flag outliers
        if cleaning_type in ('full', 'outliers'):
            numeric_cols = df.select_dtypes(include='number').columns
            outlier_flags = []
            for col in numeric_cols:
                Q1 = df[col].quantile(0.25)
                Q3 = df[col].quantile(0.75)
                IQR = Q3 - Q1
                if IQR > 0:
                    lower = Q1 - 1.5 * IQR
                    upper = Q3 + 1.5 * IQR
                    outliers = ((df[col] < lower) | (df[col] > upper)).sum()
                    if outliers > 0:
                        flag_col = f"{col}_Outlier"
                        df[flag_col] = ((df[col] < lower) | (df[col] > upper)).astype(str).replace({'True': 'Yes', 'False': 'No'})
                        outlier_flags.append(f"{col}: {outliers} outliers flagged")
                        report_parts.append(f"- Flagged **{outliers} outliers** in **{col}** (IQR method)")

        report_parts.append(f"\n## Cleaned Dataset Summary\n")
        report_parts.append(f"- **Final rows:** {len(df)}")
        report_parts.append(f"- **Final columns:** {len(df.columns)}")
        report_parts.append(f"- **Rows removed:** {original_rows - len(df)}")
        report_parts.append(f"- **Total changes:** {len(changes)}")

        # Save cleaned file
        output_path = tempfile.mktemp(suffix='_cleaned.xlsx')
        df.to_excel(output_path, index=False)

        report_text = "\n".join(report_parts)
        return report_text, output_path

    except Exception as e:
        return f"Error cleaning data: {str(e)}", None


# ─── Multi-Dataset Analysis ──────────────────────────────────────────────────
def analyze_multi_dataset(file1, file2, file3, analysis_type, instructions):
    if file1 is None or file2 is None:
        return "Please upload at least 2 datasets.", None
    datasets = []
    for i, f in enumerate([file1, file2, file3], 1):
        if f is None:
            continue
        df, err = load_data(f)
        if err or df is None:
            datasets.append((f"Dataset {i}", None, err or "Could not load file"))
        else:
            datasets.append((f"Dataset {i}", df, None))
    # Build individual summaries
    all_summaries = []
    for name, df, err in datasets:
        if df is None:
            all_summaries.append(f"### {name}\nError: {err}")
            continue
        rows, cols_list = df.shape
        dtypes = df.dtypes.to_string()
        missing = df.isnull().sum()
        missing_str = missing[missing > 0].to_string() if missing.sum() > 0 else "None"
        desc = df.describe(include="all").to_string()
        preview = df.head(15).to_string()
        summary = f"""### {name}: {rows} rows x {cols_list} columns
Columns: {list(df.columns)}
Types:
{dtypes}
Missing values:
{missing_str}
Statistics:
{desc}
Sample data (first 15 rows):
{preview}"""
        all_summaries.append(summary)
    combined = "\n\n".join(all_summaries)
    prompt = f"""You are a senior healthcare data analyst presenting a {analysis_type} of {len(datasets)} related datasets to leadership. Write as the professional analyst who performed this work. State all findings directly with specific numbers, percentages, and comparisons. NEVER say 'this appears to be', 'the data seems to show', 'the dataset contains', or 'it looks like'. Instead, deliver results confidently: 'Average claim cost was $4,230, a 12% increase year-over-year' or 'Readmission rates of 18.3% in Cardiology exceeded the facility average by 4.1 percentage points.' Be direct, precise, and actionable. Present findings as YOUR analysis, not a description of a dataset. Write as the analyst who performed this work. State all findings with specific numbers, percentages, and direct comparisons. Never say "the data appears to show" — deliver results confidently like: "Average length of stay decreased 17.4% from 3.96 days in 2023 to 3.27 days in 2024, while total charges increased 12.4%."
HERE IS THE ACTUAL DATA FROM EACH DATASET — you MUST analyze ALL of them:
{combined}
{f'Additional instructions: {instructions}' if instructions else ''}
You MUST structure your response EXACTLY as follows:
## Dataset 1 Individual Summary
- What type of healthcare data this is
- Row count, column count, key fields
- Key statistics (averages, medians, ranges for important numeric columns)
- Data quality assessment
- Notable patterns or issues specific to this dataset
## Dataset 2 Individual Summary
- What type of healthcare data this is
- Row count, column count, key fields
- Key statistics (averages, medians, ranges for important numeric columns)
- Data quality assessment
- Notable patterns or issues specific to this dataset
{"## Dataset 3 Individual Summary" + chr(10) + "- Same structure as above" + chr(10) if len(datasets) > 2 else ""}
## Merged Cross-Dataset Analysis
- How these datasets relate to each other (shared keys, overlapping fields, complementary data)
- Patterns that ONLY emerge when comparing across datasets
- Correlations between variables across different datasets
- Discrepancies or inconsistencies between datasets
- Combined assessment
## Top 5 Key Insights
Numbered list of the most important actionable findings from the COMBINED analysis, with specific numbers and column references
## Summary & Conclusion
- Overall assessment of the combined healthcare data landscape
- Risk areas identified across all datasets
- Strategic recommendations based on the merged analysis
- Suggested next steps
Use SPECIFIC numbers from the actual data. Reference actual column names and values. Do NOT make up data."""
    result = query_llm(prompt)
    # Try to create a comparison chart
    fig = None
    try:
        numeric_summaries = []
        for name, df, err in datasets:
            if df is not None:
                numeric_cols = df.select_dtypes(include="number").columns.tolist()
                if numeric_cols:
                    means = {col: df[col].mean() for col in numeric_cols[:8]}
                    numeric_summaries.append((name, means))
        if len(numeric_summaries) >= 2:
            import plotly.graph_objects as go
            fig = go.Figure()
            for name, means in numeric_summaries:
                fig.add_trace(go.Bar(name=name, x=list(means.keys()), y=list(means.values())))
            fig.update_layout(
                title="Cross-Dataset Numeric Comparison (Mean Values)",
                barmode="group",
                template="plotly_dark",
                paper_bgcolor="rgba(0,0,0,0)",
                plot_bgcolor="rgba(0,0,0,0)",
            )
    except Exception:
        pass
    return result, fig


with gr.Blocks(title="AI Healthcare & Medical Analyzer", theme=gr.themes.Base(
        primary_hue=gr.themes.colors.Color(
            c50="#fef9ec", c100="#faefd0", c200="#f4dea0", c300="#edc970",
            c400="#e4b04a", c500="#c9a84c", c600="#a8872e", c700="#856519",
            c800="#5e440d", c900="#3a2a05", c950="#1e1502"
        ),
        secondary_hue=gr.themes.colors.Color(
            c50="#e8edf5", c100="#c5d0e6", c200="#9eb0d4", c300="#7490c2",
            c400="#4f74b0", c500="#2d5a9e", c600="#1a3a6e", c700="#112240",
            c800="#0a1628", c900="#060e1a", c950="#03070d"
        ),
        neutral_hue=gr.themes.colors.Color(
            c50="#f0f2f7", c100="#d8dde9", c200="#b8c2d6", c300="#95a3be",
            c400="#7285a6", c500="#536890", c600="#3a4f73", c700="#263856",
            c800="#162438", c900="#0c1622", c950="#060b11"
        ),
        font=gr.themes.GoogleFont("DM Sans"),
        font_mono=gr.themes.GoogleFont("DM Mono"),
    ).set(
        body_background_fill="#0a1628",
        body_background_fill_dark="#0a1628",
        body_text_color="#f8f9fc",
        body_text_color_dark="#f8f9fc",
        block_background_fill="#112240",
        block_background_fill_dark="#112240",
        block_border_color="rgba(201,168,76,0.2)",
        block_border_color_dark="rgba(201,168,76,0.2)",
        block_label_background_fill="#112240",
        block_label_background_fill_dark="#112240",
        block_label_text_color="#c9a84c",
        block_label_text_color_dark="#c9a84c",
        block_title_text_color="#c9a84c",
        block_title_text_color_dark="#c9a84c",
        button_primary_background_fill="linear-gradient(135deg,#c9a84c,#e8c96a)",
        button_primary_background_fill_dark="linear-gradient(135deg,#c9a84c,#e8c96a)",
        button_primary_text_color="#0a1628",
        button_primary_text_color_dark="#0a1628",
        button_secondary_background_fill="#112240",
        button_secondary_background_fill_dark="#112240",
        button_secondary_border_color="rgba(201,168,76,0.3)",
        button_secondary_border_color_dark="rgba(201,168,76,0.3)",
        button_secondary_text_color="#c9a84c",
        button_secondary_text_color_dark="#c9a84c",
        input_background_fill="#0a1628",
        input_background_fill_dark="#0a1628",
        input_border_color="rgba(201,168,76,0.2)",
        input_border_color_dark="rgba(201,168,76,0.2)",
        input_placeholder_color="#8892a4",
        input_placeholder_color_dark="#8892a4",
        shadow_drop="0 4px 24px rgba(0,0,0,0.4)",
        shadow_drop_lg="0 8px 40px rgba(0,0,0,0.5)",
        table_even_background_fill="#0a1628",
        table_even_background_fill_dark="#0a1628",
        table_odd_background_fill="#112240",
        table_odd_background_fill_dark="#112240",
    )) as demo:
    gr.Markdown("# 🏥 AI Healthcare & Medical Analyzer")
    gr.Markdown(DISCLAIMER)

    with gr.Tabs():

        # ── Tab 1 ──────────────────────────────────────────────────────────────
        with gr.Tab("📊 Data Upload & Overview"):
            gr.Markdown("## Upload Your Healthcare Data")
            gr.Markdown(PII_WARNING)
            with gr.Row():
                with gr.Column(scale=1):
                    t1_file = gr.File(label="Upload File (CSV, Excel, PDF, JSON)",
                                      file_types=[".csv", ".xlsx", ".xls", ".pdf", ".json"])
                    t1_btn = gr.Button("🔍 Analyze Dataset", variant="primary")
                    gr.Markdown(WAIT_MSG)
                with gr.Column(scale=2):
                    t1_analysis = gr.Markdown(label="AI Analysis")
            with gr.Row():
                t1_stats = gr.Markdown(label="Statistics")
                t1_preview = gr.Markdown(label="Data Preview")
            t1_btn.click(analyze_overview, inputs=[t1_file],
                         outputs=[t1_analysis, t1_stats, t1_preview])

        # ── Tab 2 ──────────────────────────────────────────────────────────────
        with gr.Tab("🩺 Patient Outcome Analysis"):
            gr.Markdown("## Patient Outcome & Treatment Effectiveness Analysis")
            gr.Markdown(PII_WARNING)
            with gr.Row():
                with gr.Column(scale=1):
                    t2_file = gr.File(label="Upload De-identified Patient Data (CSV/Excel)",
                                      file_types=[".csv", ".xlsx", ".xls"])
                    t2_btn = gr.Button("🩺 Analyze Outcomes", variant="primary")
                    gr.Markdown(WAIT_MSG)
                with gr.Column(scale=2):
                    t2_output = gr.Markdown(label="Patient Outcome Analysis")
            t2_chart = gr.Plot(label="Outcome Chart")
            t2_btn.click(analyze_outcomes, inputs=[t2_file],
                         outputs=[t2_output, t2_chart])

        # ── Tab 3 ──────────────────────────────────────────────────────────────
        with gr.Tab("🦠 Disease Pattern Detection"):
            gr.Markdown("## Epidemiological Disease Pattern Analysis")
            gr.Markdown(PII_WARNING)
            with gr.Row():
                with gr.Column(scale=1):
                    t3_file = gr.File(label="Upload Epidemiological Data (CSV/Excel)",
                                      file_types=[".csv", ".xlsx", ".xls"])
                    t3_btn = gr.Button("🦠 Detect Disease Patterns", variant="primary")
                    gr.Markdown(WAIT_MSG)
                with gr.Column(scale=2):
                    t3_output = gr.Markdown(label="Disease Pattern Analysis")
            t3_chart = gr.Plot(label="Disease Trend Chart")
            t3_btn.click(analyze_disease_patterns, inputs=[t3_file],
                         outputs=[t3_output, t3_chart])

        # ── Tab 4 ──────────────────────────────────────────────────────────────
        with gr.Tab("🏨 Hospital Resource Planning"):
            gr.Markdown("## Hospital Capacity & Resource Optimization")
            gr.Markdown(PII_WARNING)
            with gr.Row():
                with gr.Column(scale=1):
                    t4_file = gr.File(label="Upload Capacity & Utilization Data (CSV/Excel)",
                                      file_types=[".csv", ".xlsx", ".xls"])
                    t4_btn = gr.Button("🏨 Analyze Resources", variant="primary")
                    gr.Markdown(WAIT_MSG)
                with gr.Column(scale=2):
                    t4_output = gr.Markdown(label="Resource Planning Analysis")
            t4_chart = gr.Plot(label="Resource Chart")
            t4_btn.click(analyze_resources, inputs=[t4_file],
                         outputs=[t4_output, t4_chart])

        # ── Tab 5 ──────────────────────────────────────────────────────────────
        with gr.Tab("💰 Medical Billing Audit"):
            gr.Markdown("## Medical Billing Audit & Compliance Analysis")
            gr.Markdown(PII_WARNING)
            with gr.Row():
                with gr.Column(scale=1):
                    t5_file = gr.File(label="Upload Billing Data (CSV/Excel)",
                                      file_types=[".csv", ".xlsx", ".xls"])
                    t5_btn = gr.Button("💰 Run Billing Audit", variant="primary")
                    gr.Markdown(WAIT_MSG)
                with gr.Column(scale=2):
                    t5_output = gr.Markdown(label="Billing Audit Report")
            t5_chart = gr.Plot(label="Billing Flag Chart")
            t5_btn.click(analyze_billing, inputs=[t5_file],
                         outputs=[t5_output, t5_chart])

        # ── Tab 6 ──────────────────────────────────────────────────────────────
        with gr.Tab("🔬 Clinical Trial Analysis"):
            gr.Markdown("## Clinical Trial Results & Statistical Analysis")
            gr.Markdown(PII_WARNING)
            with gr.Row():
                with gr.Column(scale=1):
                    t6_file = gr.File(label="Upload Trial Data (CSV, Excel, PDF)",
                                      file_types=[".csv", ".xlsx", ".xls", ".pdf"])
                    t6_btn = gr.Button("🔬 Analyze Trial Data", variant="primary")
                    gr.Markdown(WAIT_MSG)
                with gr.Column(scale=2):
                    t6_output = gr.Markdown(label="Clinical Trial Analysis")
            t6_chart = gr.Plot(label="Trial Results Chart")
            t6_btn.click(analyze_trials, inputs=[t6_file],
                         outputs=[t6_output, t6_chart])

        # ── Tab 7 ──────────────────────────────────────────────────────────────
        with gr.Tab("📈 Visualizations & Charts"):
            gr.Markdown("## Interactive Healthcare Data Visualizations")
            gr.Markdown(PII_WARNING)
            with gr.Row():
                with gr.Column(scale=1):
                    t7_file = gr.File(label="Upload Data (CSV/Excel)",
                                      file_types=[".csv", ".xlsx", ".xls"])
                    t7_chart_type = gr.Dropdown(
                        choices=["Bar Chart", "Line Chart", "Pie Chart", "Scatter Plot",
                                 "Heatmap", "Box Plot", "Histogram", "Area Chart"],
                        value="Bar Chart", label="Chart Type")
                    t7_x = gr.Dropdown(choices=[], allow_custom_value=True, label="X Axis Column")
                    t7_y = gr.Dropdown(choices=[], allow_custom_value=True, label="Y Axis Column")
                    t7_color = gr.Dropdown(choices=[], allow_custom_value=True, label="Color By (optional)")
                    t7_btn = gr.Button("📊 Generate Chart", variant="primary")
                    t7_status = gr.Markdown()
                with gr.Column(scale=2):
                    t7_plot = gr.Plot(label="Chart")
            t7_file.change(get_columns, inputs=[t7_file],
                           outputs=[t7_x, t7_y, t7_color])
            t7_btn.click(generate_viz,
                         inputs=[t7_file, t7_chart_type, t7_x, t7_y, t7_color],
                         outputs=[t7_plot, t7_status])

        # ── Tab 8 ──────────────────────────────────────────────────────────────
        with gr.Tab("📄 Reports & Presentations"):
            gr.Markdown("## Generate Professional Healthcare Reports")
            gr.Markdown(PII_WARNING)
            with gr.Row():
                with gr.Column(scale=1):
                    t8_file = gr.File(label="Upload Data (CSV, Excel, PDF, JSON)",
                                      file_types=[".csv", ".xlsx", ".xls", ".pdf", ".json"])
                    t8_type = gr.Dropdown(
                        choices=["PowerPoint Presentation", "PDF Report",
                                 "Executive Summary", "Quality Improvement Report"],
                        value="PDF Report", label="Report Type")
                    t8_btn = gr.Button("📄 Generate Report", variant="primary")
                    gr.Markdown(WAIT_MSG)
                    t8_download = gr.File(label="⬇️ Download Report")
                with gr.Column(scale=2):
                    t8_output = gr.Markdown(label="Report Preview")
            t8_btn.click(generate_report, inputs=[t8_file, t8_type],
                         outputs=[t8_output, t8_download])

        # ── Multi-Dataset Analysis ─────────────────────────────────────────────
        with gr.Tab("🔗 Multi-Dataset Analysis"):
            gr.Markdown("## Cross-Reference Multiple Datasets")
            gr.Markdown(PII_WARNING)
            with gr.Row():
                with gr.Column(scale=1):
                    t8m_file1 = gr.File(label="Dataset 1 (Primary)", file_types=[".csv", ".xlsx", ".xls"])
                    t8m_file2 = gr.File(label="Dataset 2", file_types=[".csv", ".xlsx", ".xls"])
                    t8m_file3 = gr.File(label="Dataset 3 (Optional)", file_types=[".csv", ".xlsx", ".xls"])
                    t8m_type = gr.Dropdown(
                        choices=["Comprehensive Comparison", "Year-over-Year Trends", "Correlation Analysis", "Anomaly Detection", "Data Quality Audit"],
                        value="Comprehensive Comparison", label="Analysis Type", allow_custom_value=True)
                    t8m_instructions = gr.Textbox(label="Additional Instructions (Optional)", lines=2)
                    t8m_btn = gr.Button("🔗 Analyze All Datasets", variant="primary")
                    gr.Markdown(WAIT_MSG)
                with gr.Column(scale=2):
                    t8m_output = gr.Markdown(label="Multi-Dataset Analysis")
                    t8m_plot = gr.Plot(label="Cross-Dataset Chart")
            t8m_btn.click(analyze_multi_dataset,
                          inputs=[t8m_file1, t8m_file2, t8m_file3, t8m_type, t8m_instructions],
                          outputs=[t8m_output, t8m_plot])

        # ── Tab 9 ──────────────────────────────────────────────────────────────
        with gr.Tab("💬 AI Data Chat"):
            gr.Markdown("## Chat With Your Healthcare Data")
            gr.Markdown(PII_WARNING)
            with gr.Row():
                with gr.Column(scale=1):
                    t9_file = gr.File(label="Upload De-identified Data for Context (CSV/Excel)",
                                      file_types=[".csv", ".xlsx", ".xls"])
                    t9_upload_btn = gr.Button("📤 Load Data", variant="secondary")
                    t9_data_status = gr.Markdown()
                with gr.Column(scale=2):
                    t9_chatbot = gr.ChatInterface(
                        fn=chat_with_data,
                        examples=[
                            "What is the average length of stay in this dataset?",
                            "Which condition has the highest readmission rate?",
                            "What are the most common diagnoses?",
                            "Generate a SQL query to find the top 10 longest hospital stays",
                            "Write Python code to calculate 30-day readmission rates",
                            "Are there any unusual patterns in the billing data?",
                            "Which department has the highest resource utilization?",
                        ],
                        title="",
                    )
            t9_upload_btn.click(upload_chat_data, inputs=[t9_file],
                                outputs=[t9_data_status])
            # Hidden clean_data endpoint
            clean_file_input = gr.File(visible=False)
            clean_type_input = gr.Dropdown(choices=["full","nulls","duplicates","formats","outliers","report"], value="full", visible=False, allow_custom_value=True)
            clean_text_output = gr.Textbox(visible=False)
            clean_file_output = gr.File(visible=False)
            clean_btn = gr.Button(visible=False)
            clean_btn.click(clean_data, inputs=[clean_file_input, clean_type_input], outputs=[clean_text_output, clean_file_output], api_name="clean_data")




    gr.Markdown(WATERMARK)


# --- FastAPI route for multi-dataset analysis (bypasses Gradio event system) ---
import json as _json
from fastapi import Request as _Request
from fastapi.responses import JSONResponse as _JSONResponse

HEALTHCARE_SYSTEM_PROMPT = """You are a senior healthcare data analyst and clinical informatics specialist with expertise in hospital operations, clinical outcomes, medical billing, population health, and healthcare quality metrics.

ANALYTICAL STANDARDS:
1. Always calculate and interpret key healthcare KPIs: readmission rates (HRRP benchmark: <15%), length of stay (LOS), case mix index (CMI), cost per discharge, mortality rates, patient satisfaction scores (HCAHPS).
2. Apply CMS quality benchmarks, Joint Commission standards, and NCQA HEDIS measures where relevant.
3. For billing analysis, identify DRG optimization opportunities, coding accuracy (ICD-10-CM/PCS), charge capture gaps, and denial root causes.
4. Evaluate clinical outcomes against risk-adjusted national benchmarks (CMS Star Ratings, Leapfrog, US News methodology).
5. Apply population health frameworks: chronic disease management, preventive care gaps, social determinants of health (SDOH) impact.
6. Flag HIPAA compliance considerations, fraud/waste/abuse (FWA) indicators per OIG guidelines.
7. Structure every analysis with: Clinical Summary → Quality Metrics → Financial Performance → Operational Efficiency → Risk Flags → Strategic Recommendations.
8. Reference specific CMS programs: MSSP, BPCI, VBP, HAC Reduction Program where applicable.
9. Be precise with statistical analysis — confidence intervals, p-values, and sample size considerations matter in clinical data.
10. Always distinguish between correlation and causation in clinical outcome analysis."""


@demo.app.post("/run/analyze_multi_dataset")
async def _api_analyze_multi_dataset(request: _Request):
    body = await request.json()
    data = body.get("data", [])
    f1 = data[0] if len(data) > 0 else None
    f2 = data[1] if len(data) > 1 else None
    f3 = data[2] if len(data) > 2 else None
    atype = data[3] if len(data) > 3 else "Comprehensive Comparison"
    instr = data[4] if len(data) > 4 else ""
    def fix_file(fd):
        if fd is None:
            return None
        if isinstance(fd, dict) and "path" in fd:
            return fd["path"]
        return fd
    result_text, fig = analyze_multi_dataset(fix_file(f1), fix_file(f2), fix_file(f3), atype, instr)
    fig_data = None
    if fig is not None:
        try:
            fig_data = _json.loads(fig.to_json())
        except Exception:
            pass
    return _JSONResponse({"data": [result_text, fig_data]})


if __name__ == "__main__":


    demo.launch(server_name="0.0.0.0", server_port=int(os.environ.get("GRADIO_SERVER_PORT", 7860)), share=False, ssr_mode=False)
