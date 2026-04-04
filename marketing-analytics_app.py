"""
Marketing Analytics Suite - Gradio Backend App
Port: 7876
Uses Gradio 6.10.0 with OpenAI API (gpt-4o) and Plotly for visualizations
"""

import os
import json
import io
import pandas as pd
import plotly.graph_objects as go
import plotly.express as px
from datetime import datetime, timedelta
import gradio as gr
import requests as http_requests
from collections import Counter
import numpy as np
ANALYST_TONE = """CRITICAL ACCURACY RULES:
1. ONLY cite numbers, percentages, and statistics that are EXPLICITLY provided in the data below.
2. NEVER invent, estimate, fabricate, or hallucinate any statistic, percentage, or metric.
3. If a metric is not in the provided data, do NOT mention it. Do NOT assume demographics, conversion rates, ROI, or any other metric unless the actual computed value is given to you.
4. Every number you cite must trace directly back to the data provided. If you cannot point to the source number, do not include it.
5. State findings directly. Do NOT use phrases like: 'As the professional analyst', 'This dataset appears to portray', 'the data seems to show', 'I have identified', 'Under my analysis'.
6. Lead every section with hard numbers FROM THE DATA. No preamble, no self-referencing.
7. If the data is limited, say so — a short accurate analysis is better than a long fabricated one.
8. Recommendations should be grounded in the actual data patterns, not hypothetical scenarios.

REQUIRED ANALYSIS STRUCTURE — every response MUST include ALL of these sections:
## Key Metrics Summary
List the top 5-10 most important metrics from the data with exact values.

## Performance Analysis
Interpret what the numbers mean. Identify trends, patterns, anomalies, and correlations visible in the data. Compare highs vs lows. Identify top and bottom performers. Calculate derived metrics (rates, ratios, growth) where the raw numbers allow.

## Conclusions
State 3-5 definitive conclusions drawn directly from the data patterns. Be specific and decisive — not vague.

## Strategic Recommendations
Provide 4-6 specific, actionable recommendations grounded in the actual data. Each recommendation must reference a specific metric or finding that justifies it. Prioritize by potential impact.

## Watch List
Identify 2-3 metrics or trends that need monitoring or immediate attention, with specific thresholds from the data.

Example of BAD: 'Videos linked to promotional campaigns show a conversion rate of 12%' (if conversion rate was never computed)
Example of GOOD: 'Total impressions: 245,891 across 366 posts. Average engagement per post: 47.3. Shares peaked in September at 1,204. The 8.3x gap between best and worst performing posts (47,821 vs 5,731 impressions) indicates high content variability — standardizing the top post format could lift average performance significantly.'"""

WATERMARK = "\n\n---\n*Analysis powered by QuantusData.ai Marketing Analytics Suite*"


# ============================================================================
# UTILITY FUNCTIONS
# ============================================================================


def fetch_url_content(text):
    """Detect URLs in user text, fetch their content, and return as context."""
    import re
    urls = re.findall(r'https?://[^\s<>"{}|\^`\[\]]+', text)
    if not urls:
        return ""
    fetched = []
    for url in urls[:3]:  # limit to 3 URLs
        try:
            headers = {
                "User-Agent": "Mozilla/5.0 (compatible; ExistentialGateway/1.0)",
                "Accept": "text/html,application/xhtml+xml,application/xml;q=0.9,*/*;q=0.8"
            }
            r = requests.get(url, timeout=15, headers=headers)
            if r.status_code == 200:
                from html.parser import HTMLParser
                class TextExtractor(HTMLParser):
                    def __init__(self):
                        super().__init__()
                        self.text = []
                        self.skip = False
                    def handle_starttag(self, tag, attrs):
                        if tag in ('script', 'style', 'nav', 'footer', 'header'):
                            self.skip = True
                    def handle_endtag(self, tag):
                        if tag in ('script', 'style', 'nav', 'footer', 'header'):
                            self.skip = False
                    def handle_data(self, data):
                        if not self.skip and data.strip():
                            self.text.append(data.strip())
                parser = TextExtractor()
                parser.feed(r.text)
                page_text = ' '.join(parser.text)[:3000]
                fetched.append(f"[URL: {url}]\n{page_text}")
        except Exception as e:
            fetched.append(f"[URL: {url}] Could not fetch: {str(e)}")
    if fetched:
        return "\n\nUSER-PROVIDED URL CONTENT (analyze this as part of your response):\n" + "\n---\n".join(fetched)
    return ""



def enrich_with_urls(text):
    """If text contains URLs, fetch and append their content."""
    if not text or 'http' not in text:
        return text
    url_content = fetch_url_content(text)
    if url_content:
        return text + url_content
    return text


def fig_to_dict(fig):
    """Convert Plotly figure to JSON dict."""
    if fig is None:
        return None
    try:
        return json.loads(fig.to_json())
    except:
        return None


def query_llm(prompt, max_tokens=6000):
    """Query OpenAI-compatible endpoint for analysis"""
    api_key = os.getenv("OPENAI_API_KEY", "")
    try:
        resp = http_requests.post(
            "https://api.openai.com/v1/chat/completions",
            headers={"Authorization": f"Bearer {api_key}", "Content-Type": "application/json"},
            json={
                "model": "gpt-4o",
                "messages": [{"role": "user", "content": prompt}],
                "max_tokens": max_tokens,
                "temperature": 0.3
            },
            timeout=120
        )
        data = resp.json()
        return data["choices"][0]["message"]["content"]
    except Exception as e:
        return f"LLM Error: {str(e)}"


def load_data(file):
    """Load CSV or Excel file into DataFrame."""
    try:
        if file is None:
            return None, "No file uploaded"
        file_path = file.name if hasattr(file, 'name') else str(file)
        if file_path.endswith('.csv'):
            try:
                df = pd.read_csv(file_path)
            except Exception:
                df = pd.read_csv(file_path, on_bad_lines='skip')
            # Detect title rows: skip row 0 if it looks like a report title
            total_cols = len(df.columns)
            unnamed_count = sum(1 for c in df.columns if 'Unnamed' in str(c))
            if unnamed_count > total_cols * 0.5:
                try:
                    df = pd.read_csv(file_path, skiprows=1, on_bad_lines='skip')
                except Exception:
                    pass
            elif len(df) > 0:
                first_row_nulls = df.iloc[0].isnull().sum()
                if first_row_nulls > total_cols * 0.5:
                    try:
                        df = pd.read_csv(file_path, skiprows=1, on_bad_lines='skip')
                    except Exception:
                        pass
        elif file_path.endswith(('.xls', '.xlsx')):
            df = pd.read_excel(file_path)
            total_cols = len(df.columns)
            unnamed_count = sum(1 for c in df.columns if 'Unnamed' in str(c))
            if unnamed_count > total_cols * 0.5:
                try:
                    df = pd.read_excel(file_path, skiprows=1)
                except Exception:
                    pass
        else:
            return None, "Unsupported file format. Please use CSV or Excel."
        return df, None
    except Exception as e:
        return None, f"Error loading file: {str(e)}"
def normalize_columns(df):
    """Normalize column names to standard marketing metrics."""
    column_mapping = {
        # Campaign identifiers
        'campaign_name': ['campaign_name', 'campaign'],
        'campaign_id': ['campaign_id', 'campaign'],
        'channel': ['channel', 'media_channel', 'platform', 'ad_platform'],

        # Spend metrics
        'spend': ['spend', 'cost', 'budget', 'media_spend', 'ad_spend', 'total_spend'],

        # Impression metrics
        'impressions': ['impressions', 'imps', 'total_impressions'],

        # Click metrics
        'clicks': ['clicks', 'link_clicks', 'total_clicks'],

        # Conversion metrics
        'conversions': ['conversions', 'conv', 'purchases', 'leads', 'actions', 'total_conversions'],

        # Revenue metrics
        'revenue': ['revenue', 'sales', 'value', 'conversion_value', 'total_revenue'],

        # Rate metrics
        'ctr': ['ctr', 'click_through_rate', 'click_rate'],
        'cpm': ['cpm', 'cost_per_mille', 'cost_per_thousand'],
        'cpc': ['cpc', 'cost_per_click'],
        'cpa': ['cpa', 'cost_per_action', 'cost_per_acquisition', 'cost_per_lead'],
        'roas': ['roas', 'return_on_ad_spend', 'roi'],

        # Reach metrics
        'reach': ['reach', 'unique_reach'],
        'engagement': ['engagement', 'engagements', 'total_engagement'],

        # Social metrics
        'likes': ['likes', 'reactions'],
        'comments': ['comments', 'replies'],
        'shares': ['shares', 'retweets', 'reposts'],
        'saves': ['saves', 'bookmarks'],
        'followers': ['followers', 'followers_gained', 'new_followers'],
        'video_views': ['video_views', 'views', 'video_plays'],
        'completion_rate': ['completion_rate', 'video_completion', 'vcr'],

        # Creative metrics
        'creative_id': ['creative_id', 'ad_id'],
        'creative_name': ['creative_name', 'ad_name'],
        'creative_type': ['creative_type', 'ad_type', 'format', 'ad_format'],

        # Audience metrics
        'segment': ['segment', 'audience', 'audience_segment'],
        'age_group': ['age_group', 'age', 'age_range', 'demographic'],
        'gender': ['gender', 'sex'],
        'geo': ['geo', 'geography', 'location', 'region', 'state', 'dma', 'market'],
        'interest': ['interest', 'interest_category'],
        'behavior': ['behavior', 'user_behavior'],
        'ltv': ['ltv', 'lifetime_value'],

        # Attribution metrics
        'brand_awareness_pre': ['brand_awareness_pre', 'awareness_pre', 'pre_awareness'],
        'brand_awareness_post': ['brand_awareness_post', 'awareness_post', 'post_awareness'],
        'ad_recall': ['ad_recall', 'recall'],
        'purchase_intent': ['purchase_intent', 'intent'],
        'favorability': ['favorability', 'brand_favorability'],

        # Path metrics
        'touchpoint_order': ['touchpoint_order', 'touch_order', 'position'],
        'conversion_flag': ['conversion_flag', 'converted', 'is_conversion'],

        # Content metrics
        'post_id': ['post_id', 'content_id'],
        'post_type': ['post_type', 'content_type'],
        'date': ['date', 'report_date', 'day', 'period', 'timestamp']
    }

    df_normalized = df.copy()
    df_normalized.columns = df_normalized.columns.str.lower().str.strip()

    # Handle duplicate column names BEFORE renaming
    # If two columns have the same name, keep only the first one
    seen = {}
    new_cols = []
    drop_indices = []
    for i, col in enumerate(df_normalized.columns):
        if col in seen:
            drop_indices.append(i)
        else:
            seen[col] = i
            new_cols.append(col)
    if drop_indices:
        df_normalized = df_normalized.iloc[:, [i for i in range(len(df_normalized.columns)) if i not in drop_indices]]

    reverse_mapping = {}
    for standard_col, aliases in column_mapping.items():
        for alias in aliases:
            reverse_mapping[alias] = standard_col

    rename_dict = {}
    for col in df_normalized.columns:
        if col in reverse_mapping:
            rename_dict[col] = reverse_mapping[col]

    df_normalized.rename(columns=rename_dict, inplace=True)

    # Final safety check: if rename created duplicates, drop duplicate columns
    df_normalized = df_normalized.loc[:, ~df_normalized.columns.duplicated()]

    return df_normalized


# ============================================================================
# TAB 0: PAID MEDIA ANALYTICS
# ============================================================================

def analyze_paid_media(file):
    """Analyze paid channel performance (social, programmatic, linear TV, digital, display)."""
    df, error = load_data(file)
    if error:
        return f"Error: {error}", None

    df = normalize_columns(df)

    try:
        # Detect available columns
        has_channel = 'channel' in df.columns
        has_spend = 'spend' in df.columns
        has_impressions = 'impressions' in df.columns
        has_clicks = 'clicks' in df.columns
        has_conversions = 'conversions' in df.columns
        has_revenue = 'revenue' in df.columns

        # Ensure numeric columns
        for col in ['spend', 'impressions', 'clicks', 'conversions', 'revenue']:
            if col in df.columns:
                df[col] = pd.to_numeric(df[col], errors='coerce').fillna(0)

        # Find best groupby column
        group_col = None
        for candidate in ['channel', 'campaign_name', 'campaign_id', 'segment']:
            if candidate in df.columns and df[candidate].nunique() > 1:
                group_col = candidate
                break

        metrics_text = "PAID MEDIA ANALYTICS\n"
        metrics_text += f"Dataset: {len(df)} rows, {len(df.columns)} columns\n"
        metrics_text += f"Columns detected: {', '.join(df.columns.tolist())}\n\n"

        # Overall totals
        metrics_text += "TOTALS:\n"
        if has_spend:
            metrics_text += f"Total Spend: ${df['spend'].sum():,.2f}\n"
        if has_impressions:
            metrics_text += f"Total Impressions: {df['impressions'].sum():,.0f}\n"
        if has_clicks:
            metrics_text += f"Total Clicks: {df['clicks'].sum():,.0f}\n"
        if has_conversions:
            metrics_text += f"Total Conversions: {df['conversions'].sum():,.0f}\n"
        if has_revenue:
            metrics_text += f"Total Revenue: ${df['revenue'].sum():,.2f}\n"

        # Derived metrics
        if has_clicks and has_impressions and df['impressions'].sum() > 0:
            metrics_text += f"Overall CTR: {df['clicks'].sum() / df['impressions'].sum() * 100:.2f}%\n"
        if has_spend and has_impressions and df['impressions'].sum() > 0:
            metrics_text += f"Overall CPM: ${df['spend'].sum() / df['impressions'].sum() * 1000:.2f}\n"
        if has_spend and has_clicks and df['clicks'].sum() > 0:
            metrics_text += f"Overall CPC: ${df['spend'].sum() / df['clicks'].sum():.2f}\n"
        if has_spend and has_conversions and df['conversions'].sum() > 0:
            metrics_text += f"Overall CPA: ${df['spend'].sum() / df['conversions'].sum():.2f}\n"
        if has_revenue and has_spend and df['spend'].sum() > 0:
            metrics_text += f"Overall ROAS: {df['revenue'].sum() / df['spend'].sum():.2f}x\n"

        # Breakdowns by group (if available)
        if group_col:
            metrics_text += f"\nBY {group_col.upper()}:\n"
            agg_dict = {}
            for col in ['spend', 'impressions', 'clicks', 'conversions', 'revenue']:
                if col in df.columns:
                    agg_dict[col] = 'sum'
            if agg_dict:
                grouped = df.groupby(group_col).agg(agg_dict).reset_index()
                for _, row in grouped.iterrows():
                    metrics_text += f"\n{row[group_col]}:\n"
                    for col in agg_dict:
                        if col in ['spend', 'revenue']:
                            metrics_text += f"  {col.title()}: ${row[col]:,.2f}\n"
                        else:
                            metrics_text += f"  {col.title()}: {row[col]:,.0f}\n"

        # LLM analysis
        prompt = f"{ANALYST_TONE}\n\nHere is the ONLY data available — do not reference any metrics not listed here:\n\n{metrics_text}\n\nUsing ONLY the numbers above, provide detailed analysis of paid media performance. Focus on channel efficiency, optimization opportunities, and budget allocation recommendations.\n{WATERMARK}"
        analysis = query_llm(prompt)

        # Chart
        fig = go.Figure()
        if group_col and has_spend:
            grouped = df.groupby(group_col)['spend'].sum().reset_index()
            fig.add_trace(go.Bar(x=grouped[group_col], y=grouped['spend'], name='Spend'))
            fig.update_layout(title=f'Spend by {group_col.title()}', xaxis_title=group_col.title(), yaxis_title='Spend ($)')
        elif has_spend:
            fig.add_trace(go.Histogram(x=df['spend'], name='Spend Distribution'))
            fig.update_layout(title='Spend Distribution')

        return analysis, fig_to_dict(fig)
    except Exception as e:
        return f"Error in analysis: {str(e)}", None


# ============================================================================
# TAB 1: CREATIVE PERFORMANCE
# ============================================================================

def analyze_creative_performance(file):
    """Analyze creative/ad performance metrics."""
    df, error = load_data(file)
    if error:
        return f"Error: {error}", None

    df = normalize_columns(df)

    try:
        # Detect available columns
        has_creative_id = 'creative_id' in df.columns
        has_creative_name = 'creative_name' in df.columns
        has_impressions = 'impressions' in df.columns
        has_clicks = 'clicks' in df.columns
        has_conversions = 'conversions' in df.columns
        has_engagement = 'engagement' in df.columns
        has_video_views = 'video_views' in df.columns

        # Ensure numeric columns
        for col in ['impressions', 'clicks', 'conversions', 'engagement', 'video_views']:
            if col in df.columns:
                df[col] = pd.to_numeric(df[col], errors='coerce').fillna(0)

        # Find groupby column
        group_col = None
        if has_creative_id:
            group_col = 'creative_id'
        elif has_creative_name:
            group_col = 'creative_name'
        else:
            # Look for alternative grouping columns
            for candidate in ['campaign_name', 'campaign_id', 'segment']:
                if candidate in df.columns and df[candidate].nunique() > 1:
                    group_col = candidate
                    break

        metrics_text = "CREATIVE PERFORMANCE ANALYSIS\n"
        metrics_text += f"Dataset: {len(df)} rows, {len(df.columns)} columns\n"
        metrics_text += f"Columns detected: {', '.join(df.columns.tolist())}\n\n"

        if group_col:
            agg_dict = {}
            if has_impressions:
                agg_dict['impressions'] = 'sum'
            if has_clicks:
                agg_dict['clicks'] = 'sum'
            if has_conversions:
                agg_dict['conversions'] = 'sum'
            if has_engagement:
                agg_dict['engagement'] = 'sum'
            if has_video_views:
                agg_dict['video_views'] = 'sum'

            if agg_dict:
                creative_metrics = df.groupby(group_col).agg(agg_dict).reset_index()

                # Calculate rates
                if has_clicks and has_impressions and (creative_metrics['impressions'] > 0).any():
                    creative_metrics['ctr'] = (creative_metrics['clicks'] / creative_metrics['impressions'] * 100).round(2)
                if has_engagement and has_impressions and (creative_metrics['impressions'] > 0).any():
                    creative_metrics['engagement_rate'] = (creative_metrics['engagement'] / creative_metrics['impressions'] * 100).round(2)
                if has_conversions and has_clicks and (creative_metrics['clicks'] > 0).any():
                    creative_metrics['conv_rate'] = (creative_metrics['conversions'] / creative_metrics['clicks'] * 100).round(2)

                # Sort by performance
                sort_col = 'ctr' if 'ctr' in creative_metrics.columns else 'engagement'
                if sort_col in creative_metrics.columns:
                    creative_metrics_sorted = creative_metrics.sort_values(sort_col, ascending=False)
                else:
                    creative_metrics_sorted = creative_metrics

                metrics_text += f"Total {group_col.title()}s: {len(creative_metrics)}\n"
                if has_impressions:
                    metrics_text += f"Total Impressions: {creative_metrics['impressions'].sum():,.0f}\n"
                if has_clicks:
                    metrics_text += f"Total Clicks: {creative_metrics['clicks'].sum():,.0f}\n"
                if has_conversions:
                    metrics_text += f"Total Conversions: {creative_metrics['conversions'].sum():,.0f}\n\n"

                metrics_text += "TOP PERFORMERS:\n"
                for idx, row in creative_metrics_sorted.head(5).iterrows():
                    metrics_text += f"\n{group_col.title()} {row[group_col]}:\n"
                    for col in agg_dict:
                        if col in ['impressions', 'clicks', 'conversions']:
                            metrics_text += f"  {col.title()}: {row[col]:,.0f}\n"
                        if col in creative_metrics.columns and col not in agg_dict:
                            metrics_text += f"  {col}: {row[col]:.2f}%\n"

                metrics_text += "\nBOTTOM PERFORMERS:\n"
                for idx, row in creative_metrics_sorted.tail(5).iterrows():
                    metrics_text += f"\n{group_col.title()} {row[group_col]}:\n"
                    for col in agg_dict:
                        if col in ['impressions', 'clicks', 'conversions']:
                            metrics_text += f"  {col.title()}: {row[col]:,.0f}\n"
            else:
                metrics_text += "No aggregatable numeric columns found.\n"
                creative_metrics_sorted = pd.DataFrame()
        else:
            metrics_text += f"No grouping column available. Dataset has {len(df)} rows.\n"
            creative_metrics_sorted = pd.DataFrame()

        # LLM analysis
        prompt = f"{ANALYST_TONE}\n\nHere is the ONLY data available — do not reference any metrics not listed here:\n\n{metrics_text}\n\nUsing ONLY the numbers above, deliver a complete creative performance analysis. Rank all creatives by performance using the exact metrics provided. Identify the performance gap between top and bottom creatives with exact figures. Determine what the winning creatives have in common. Provide 5 specific optimization actions with the data evidence that justifies each one.\n{WATERMARK}"
        analysis = query_llm(prompt)

        # Create visualization
        fig = go.Figure()
        if not creative_metrics_sorted.empty and group_col:
            if 'ctr' in creative_metrics_sorted.columns:
                fig.add_trace(go.Bar(
                    x=creative_metrics_sorted[group_col].astype(str),
                    y=creative_metrics_sorted['ctr'],
                    name='CTR'
                ))
                fig.update_layout(title=f'Top Creatives by CTR', xaxis_title=group_col.title(), yaxis_title='CTR (%)')
            elif 'engagement_rate' in creative_metrics_sorted.columns:
                fig.add_trace(go.Bar(
                    x=creative_metrics_sorted[group_col].astype(str),
                    y=creative_metrics_sorted['engagement_rate'],
                    name='Engagement Rate'
                ))
                fig.update_layout(title='Top Creatives by Engagement', xaxis_title=group_col.title(), yaxis_title='Engagement Rate (%)')

        return analysis, fig_to_dict(fig)

    except Exception as e:
        return f"Error in analysis: {str(e)}", None


# ============================================================================
# TAB 2: BRAND LIFT & ATTRIBUTION
# ============================================================================

def analyze_brand_lift(file):
    """Analyze brand lift, MTA, and attribution."""
    df, error = load_data(file)
    if error:
        return f"Error: {error}", None

    df = normalize_columns(df)

    try:
        # Detect available columns
        has_brand_pre = 'brand_awareness_pre' in df.columns
        has_brand_post = 'brand_awareness_post' in df.columns
        has_channel = 'channel' in df.columns
        has_conversion_flag = 'conversion_flag' in df.columns
        has_ad_recall = 'ad_recall' in df.columns
        has_purchase_intent = 'purchase_intent' in df.columns
        has_favorability = 'favorability' in df.columns

        # Ensure numeric columns
        for col in ['brand_awareness_pre', 'brand_awareness_post', 'ad_recall', 'purchase_intent', 'favorability']:
            if col in df.columns:
                df[col] = pd.to_numeric(df[col], errors='coerce').fillna(0)

        # Brand lift analysis
        if has_brand_pre and has_brand_post:
            pre_avg = df['brand_awareness_pre'].mean()
            post_avg = df['brand_awareness_post'].mean()
            lift = ((post_avg - pre_avg) / pre_avg * 100) if pre_avg > 0 else 0
        else:
            pre_avg = post_avg = lift = 0

        # Channel attribution
        if has_channel and has_conversion_flag:
            df['conversion_flag'] = pd.to_numeric(df['conversion_flag'], errors='coerce').fillna(0)
            channel_attr = df[df['conversion_flag'] == 1].groupby('channel').size()
            total_conversions = (df['conversion_flag'] == 1).sum()
        else:
            channel_attr = pd.Series()
            total_conversions = 0

        # Build metrics text
        metrics_text = "BRAND LIFT & ATTRIBUTION ANALYSIS\n"
        metrics_text += f"Dataset: {len(df)} rows, {len(df.columns)} columns\n"
        metrics_text += f"Columns detected: {', '.join(df.columns.tolist())}\n\n"

        if has_brand_pre or has_brand_post:
            metrics_text += "BRAND AWARENESS:\n"
            metrics_text += f"Brand Awareness Lift: {lift:.2f}%\n"
            metrics_text += f"Pre-Campaign Awareness: {pre_avg:.2f}%\n"
            metrics_text += f"Post-Campaign Awareness: {post_avg:.2f}%\n\n"

        if has_ad_recall:
            metrics_text += f"Ad Recall: {df['ad_recall'].mean():.2f}%\n"
        if has_purchase_intent:
            metrics_text += f"Purchase Intent: {df['purchase_intent'].mean():.2f}%\n"
        if has_favorability:
            metrics_text += f"Brand Favorability: {df['favorability'].mean():.2f}%\n"

        metrics_text += f"\nTotal Conversions: {total_conversions:,.0f}\n"
        if has_channel and len(channel_attr) > 0:
            metrics_text += "CONVERSIONS BY CHANNEL:\n"
            for channel, count in channel_attr.items():
                pct = (count / total_conversions * 100) if total_conversions > 0 else 0
                metrics_text += f"  {channel}: {count:,.0f} ({pct:.1f}%)\n"

        # LLM analysis
        prompt = f"{ANALYST_TONE}\n\nHere is the ONLY data available — do not reference any metrics not listed here:\n\n{metrics_text}\n\nUsing ONLY the numbers above, deliver a complete brand lift and attribution analysis. Quantify lift by channel using exact figures. Identify which channels drive the strongest attribution signals. Calculate contribution percentages where the data allows. Provide 5 specific channel mix recommendations tied directly to the lift data, with clear before/after projections based on the actual numbers.\n{WATERMARK}"
        analysis = query_llm(prompt)

        # Create visualization
        fig = go.Figure()
        if isinstance(channel_attr, pd.Series) and len(channel_attr) > 0:
            fig = go.Figure(data=[go.Pie(labels=channel_attr.index.astype(str), values=channel_attr.values, title='Conversions by Channel')])

        return analysis, fig_to_dict(fig)

    except Exception as e:
        return f"Error in analysis: {str(e)}", None


# ============================================================================
# TAB 3: AUDIENCE & SEGMENTATION
# ============================================================================

def analyze_audience(file):
    """Analyze audience insights and targeting."""
    df, error = load_data(file)
    if error:
        return f"Error: {error}", None

    df = normalize_columns(df)

    try:
        # Detect available columns
        has_segment = 'segment' in df.columns
        has_impressions = 'impressions' in df.columns
        has_clicks = 'clicks' in df.columns
        has_conversions = 'conversions' in df.columns
        has_ltv = 'ltv' in df.columns

        # Ensure numeric columns
        for col in ['impressions', 'clicks', 'conversions', 'ltv']:
            if col in df.columns:
                df[col] = pd.to_numeric(df[col], errors='coerce').fillna(0)

        # Segment analysis
        if has_segment:
            agg_dict = {}
            if has_impressions:
                agg_dict['impressions'] = 'sum'
            if has_clicks:
                agg_dict['clicks'] = 'sum'
            if has_conversions:
                agg_dict['conversions'] = 'sum'
            if has_ltv:
                agg_dict['ltv'] = 'mean'

            if agg_dict:
                segment_metrics = df.groupby('segment').agg(agg_dict).reset_index()

                # Calculate rates
                if has_clicks and has_impressions and (segment_metrics['impressions'] > 0).any():
                    segment_metrics['ctr'] = (segment_metrics['clicks'] / segment_metrics['impressions'] * 100).round(2)
                if has_conversions and has_clicks and (segment_metrics['clicks'] > 0).any():
                    segment_metrics['conv_rate'] = (segment_metrics['conversions'] / segment_metrics['clicks'] * 100).round(2)
            else:
                segment_metrics = pd.DataFrame()
        else:
            segment_metrics = pd.DataFrame()

        # Demographics analysis
        demo_metrics = {}
        for demo_col in ['age_group', 'gender', 'geo']:
            if demo_col in df.columns:
                agg_dict = {}
                if has_impressions:
                    agg_dict['impressions'] = 'sum'
                if has_conversions:
                    agg_dict['conversions'] = 'sum'
                if agg_dict:
                    demo_metrics[demo_col] = df.groupby(demo_col).agg(agg_dict).reset_index()
                    if has_conversions and has_impressions and (demo_metrics[demo_col]['impressions'] > 0).any():
                        demo_metrics[demo_col]['conv_rate'] = (
                            demo_metrics[demo_col]['conversions'] / demo_metrics[demo_col]['impressions'] * 100
                        ).round(2)

        # Build metrics text
        metrics_text = "AUDIENCE & SEGMENTATION ANALYSIS\n"
        metrics_text += f"Dataset: {len(df)} rows, {len(df.columns)} columns\n"
        metrics_text += f"Columns detected: {', '.join(df.columns.tolist())}\n\n"

        if not segment_metrics.empty:
            metrics_text += f"Total Segments: {len(segment_metrics)}\n"
            for idx, row in segment_metrics.iterrows():
                metrics_text += f"\n{row['segment']}:\n"
                if has_impressions:
                    metrics_text += f"  Impressions: {row['impressions']:,.0f}\n"
                if 'ctr' in segment_metrics.columns:
                    metrics_text += f"  CTR: {row['ctr']:.2f}%\n"
                if 'conv_rate' in segment_metrics.columns:
                    metrics_text += f"  Conversion Rate: {row['conv_rate']:.2f}%\n"
                if has_ltv and row['ltv'] > 0:
                    metrics_text += f"  Avg LTV: ${row['ltv']:.2f}\n"

        for demo_col, demo_df in demo_metrics.items():
            if not demo_df.empty:
                metrics_text += f"\n{demo_col.upper()} PERFORMANCE:\n"
                for idx, row in demo_df.iterrows():
                    metrics_text += f"  {row.iloc[0]}: "
                    if 'conv_rate' in demo_df.columns:
                        metrics_text += f"{row['conv_rate']:.2f}% conversion rate"
                    if has_conversions:
                        metrics_text += f" ({row['conversions']:,.0f} conversions)"
                    metrics_text += "\n"

        # LLM analysis
        prompt = f"{ANALYST_TONE}\n\nHere is the ONLY data available — do not reference any metrics not listed here:\n\n{metrics_text}\n\nUsing ONLY the numbers above, deliver a complete audience analysis. Rank all segments by engagement and value metrics using exact figures. Calculate the performance gap between highest and lowest segments. Identify the highest-value audience segments with specific justification from the data. Provide 5 targeting recommendations with the exact metrics that support each recommendation.\n{WATERMARK}"
        analysis = query_llm(prompt)

        # Create visualization
        fig = go.Figure()
        if not segment_metrics.empty and has_segment:
            if 'conv_rate' in segment_metrics.columns:
                fig = px.bar(segment_metrics, x='segment', y='conv_rate', title='Conversion Rate by Segment',
                            labels={'conv_rate': 'Conversion Rate (%)'})
            elif 'ctr' in segment_metrics.columns:
                fig = px.bar(segment_metrics, x='segment', y='ctr', title='CTR by Segment',
                            labels={'ctr': 'CTR (%)'})

        return analysis, fig_to_dict(fig)

    except Exception as e:
        return f"Error in analysis: {str(e)}", None


# ============================================================================
# TAB 4: CAMPAIGN REPORTING
# ============================================================================

def analyze_campaign_reporting(file):
    """Campaign reporting with KPI tracking."""
    df, error = load_data(file)
    if error:
        return f"Error: {error}", None

    df = normalize_columns(df)

    try:
        # Detect available columns
        has_date = 'date' in df.columns
        has_campaign = 'campaign_name' in df.columns or 'campaign_id' in df.columns
        has_spend = 'spend' in df.columns
        has_impressions = 'impressions' in df.columns
        has_clicks = 'clicks' in df.columns
        has_conversions = 'conversions' in df.columns
        has_revenue = 'revenue' in df.columns

        # Ensure numeric columns
        for col in ['spend', 'impressions', 'clicks', 'conversions', 'revenue']:
            if col in df.columns:
                df[col] = pd.to_numeric(df[col], errors='coerce').fillna(0)

        # Time-based aggregation
        if has_date:
            try:
                df['date'] = pd.to_datetime(df['date'], errors='coerce')
                agg_dict = {}
                if has_spend:
                    agg_dict['spend'] = 'sum'
                if has_impressions:
                    agg_dict['impressions'] = 'sum'
                if has_clicks:
                    agg_dict['clicks'] = 'sum'
                if has_conversions:
                    agg_dict['conversions'] = 'sum'
                if has_revenue:
                    agg_dict['revenue'] = 'sum'

                if agg_dict:
                    daily_metrics = df.groupby('date').agg(agg_dict).reset_index().sort_values('date')
                else:
                    daily_metrics = pd.DataFrame()
            except:
                daily_metrics = pd.DataFrame()
        else:
            daily_metrics = pd.DataFrame()

        # Campaign summary
        campaign_col = None
        if 'campaign_name' in df.columns:
            campaign_col = 'campaign_name'
        elif 'campaign_id' in df.columns:
            campaign_col = 'campaign_id'

        if campaign_col and has_campaign:
            agg_dict = {}
            if has_spend:
                agg_dict['spend'] = 'sum'
            if has_impressions:
                agg_dict['impressions'] = 'sum'
            if has_clicks:
                agg_dict['clicks'] = 'sum'
            if has_conversions:
                agg_dict['conversions'] = 'sum'
            if has_revenue:
                agg_dict['revenue'] = 'sum'

            if agg_dict:
                campaign_summary = df.groupby(campaign_col).agg(agg_dict).reset_index()

                # Calculate derived metrics
                if has_clicks and has_impressions and (campaign_summary['impressions'] > 0).any():
                    campaign_summary['ctr'] = (campaign_summary['clicks'] / campaign_summary['impressions'] * 100).round(2)
                if has_revenue and has_spend and (campaign_summary['spend'] > 0).any():
                    campaign_summary['roas'] = (campaign_summary['revenue'] / campaign_summary['spend']).round(2)
            else:
                campaign_summary = pd.DataFrame()
        else:
            campaign_summary = pd.DataFrame()

        # Build metrics text
        metrics_text = "CAMPAIGN REPORTING\n"
        metrics_text += f"Dataset: {len(df)} rows, {len(df.columns)} columns\n"
        metrics_text += f"Columns detected: {', '.join(df.columns.tolist())}\n\n"

        if not campaign_summary.empty:
            metrics_text += f"Total Campaigns: {len(campaign_summary)}\n"
            if has_spend:
                metrics_text += f"Period Spend: ${campaign_summary['spend'].sum():,.2f}\n"
            if has_revenue:
                metrics_text += f"Period Revenue: ${campaign_summary['revenue'].sum():,.2f}\n"
            if has_revenue and has_spend and campaign_summary['spend'].sum() > 0:
                metrics_text += f"Overall ROAS: {campaign_summary['revenue'].sum() / campaign_summary['spend'].sum():.2f}x\n\n"

            for idx, row in campaign_summary.iterrows():
                metrics_text += f"\n{row[campaign_col]}:\n"
                if has_spend:
                    metrics_text += f"  Spend: ${row['spend']:,.2f}\n"
                if has_revenue:
                    metrics_text += f"  Revenue: ${row['revenue']:,.2f}\n"
                if 'roas' in campaign_summary.columns:
                    metrics_text += f"  ROAS: {row['roas']:.2f}x\n"
                if 'ctr' in campaign_summary.columns:
                    metrics_text += f"  CTR: {row['ctr']:.2f}%\n"
        else:
            metrics_text += "No campaign data available for summary.\n"

        # LLM analysis
        prompt = f"{ANALYST_TONE}\n\nHere is the ONLY data available — do not reference any metrics not listed here:\n\n{metrics_text}\n\nUsing ONLY the numbers above, deliver a complete campaign performance report. Rank all campaigns by key performance metrics with exact figures. Quantify the performance gap between top and bottom campaigns. Identify patterns in what makes campaigns succeed or fail based on the data. Provide 5 specific next steps with the data evidence that justifies each action, including budget reallocation recommendations.\n{WATERMARK}"
        analysis = query_llm(prompt)

        # Create visualization
        fig = go.Figure()
        if not daily_metrics.empty and has_date:
            if has_spend and has_revenue:
                fig = px.line(daily_metrics, x='date', y=['spend', 'revenue'], title='Daily Spend vs Revenue Trend',
                             labels={'value': 'Amount ($)', 'variable': 'Metric'})
            elif has_spend:
                fig = px.line(daily_metrics, x='date', y='spend', title='Daily Spend Trend',
                             labels={'spend': 'Spend ($)'})
        elif not campaign_summary.empty and campaign_col:
            if 'roas' in campaign_summary.columns:
                fig = px.bar(campaign_summary, x=campaign_col, y='roas', title='ROAS by Campaign')
            elif has_spend:
                fig = px.bar(campaign_summary, x=campaign_col, y='spend', title='Spend by Campaign')

        return analysis, fig_to_dict(fig)

    except Exception as e:
        return f"Error in analysis: {str(e)}", None


# ============================================================================
# TAB 5: ORGANIC SOCIAL
# ============================================================================

def analyze_organic_social(file):
    """Organic social media analytics."""
    df, error = load_data(file)
    if error:
        return f"Error: {error}", None

    df = normalize_columns(df)

    try:
        # Detect available columns
        has_post_type = 'post_type' in df.columns
        has_platform = 'platform' in df.columns or 'channel' in df.columns
        has_impressions = 'impressions' in df.columns
        has_reach = 'reach' in df.columns
        has_engagement = 'engagement' in df.columns
        has_likes = 'likes' in df.columns
        has_comments = 'comments' in df.columns
        has_shares = 'shares' in df.columns
        has_followers = 'followers' in df.columns

        # Ensure numeric columns (use .iloc[:,0] if duplicates somehow survive)
        for col in ['impressions', 'reach', 'engagement', 'likes', 'comments', 'shares', 'followers']:
            if col in df.columns:
                series = df[col]
                if isinstance(series, pd.DataFrame):
                    series = series.iloc[:, 0]
                df[col] = pd.to_numeric(series, errors='coerce').fillna(0)

        # Aggregate metrics
        post_type_metrics = pd.DataFrame()
        if has_post_type:
            agg_dict = {}
            if has_impressions:
                agg_dict['impressions'] = 'sum'
            if has_reach:
                agg_dict['reach'] = 'sum'
            if has_engagement:
                agg_dict['engagement'] = 'sum'
            if has_likes:
                agg_dict['likes'] = 'sum'
            if has_comments:
                agg_dict['comments'] = 'sum'
            if has_shares:
                agg_dict['shares'] = 'sum'

            if agg_dict:
                post_type_metrics = df.groupby('post_type').agg(agg_dict).reset_index()

        # Platform metrics
        platform_metrics = pd.DataFrame()
        platform_col = None
        if has_platform:
            if 'platform' in df.columns:
                platform_col = 'platform'
            elif 'channel' in df.columns:
                platform_col = 'channel'

            if platform_col:
                agg_dict = {}
                if has_impressions:
                    agg_dict['impressions'] = 'sum'
                if has_engagement:
                    agg_dict['engagement'] = 'sum'
                if has_followers:
                    agg_dict['followers'] = 'sum'

                if agg_dict:
                    platform_metrics = df.groupby(platform_col).agg(agg_dict).reset_index()

        # Calculate totals
        total_impressions = df['impressions'].sum() if has_impressions else 0
        total_engagement = df['engagement'].sum() if has_engagement else df.shape[0]
        total_followers_gained = df['followers'].sum() if has_followers else 0

        # Build metrics text
        metrics_text = "ORGANIC SOCIAL ANALYTICS\n"
        metrics_text += f"Dataset: {len(df)} rows, {len(df.columns)} columns\n"
        metrics_text += f"Columns detected: {', '.join(df.columns.tolist())}\n\n"

        metrics_text += f"Total Posts: {len(df)}\n"
        if has_impressions:
            metrics_text += f"Total Impressions: {total_impressions:,.0f}\n"
        metrics_text += f"Total Engagement: {total_engagement:,.0f}\n"
        if has_impressions and total_impressions > 0:
            metrics_text += f"Overall Engagement Rate: {(total_engagement / total_impressions * 100):.2f}%\n"
        if has_followers:
            metrics_text += f"Followers Gained: {total_followers_gained:,.0f}\n\n"

        if not post_type_metrics.empty:
            metrics_text += "BY POST TYPE:\n"
            for idx, row in post_type_metrics.iterrows():
                eng_rate = (row['engagement'] / row['impressions'] * 100) if has_impressions and row['impressions'] > 0 else 0
                metrics_text += f"\n{row['post_type']}:\n"
                if has_impressions:
                    metrics_text += f"  Impressions: {row['impressions']:,.0f}\n"
                metrics_text += f"  Engagement Rate: {eng_rate:.2f}%\n"
                if has_likes:
                    metrics_text += f"  Likes: {row['likes']:,.0f}\n"
                if has_comments:
                    metrics_text += f"  Comments: {row['comments']:,.0f}\n"
                if has_shares:
                    metrics_text += f"  Shares: {row['shares']:,.0f}\n"

        if not platform_metrics.empty and platform_col:
            metrics_text += f"\nBY {platform_col.upper()}:\n"
            for idx, row in platform_metrics.iterrows():
                eng_rate = (row['engagement'] / row['impressions'] * 100) if has_impressions and row['impressions'] > 0 else 0
                metrics_text += f"\n{row[platform_col]}:\n"
                if has_impressions:
                    metrics_text += f"  Impressions: {row['impressions']:,.0f}\n"
                metrics_text += f"  Engagement Rate: {eng_rate:.2f}%\n"
                if has_followers:
                    metrics_text += f"  Followers: {row['followers']:,.0f}\n"

        # LLM analysis
        prompt = f"{ANALYST_TONE}\n\nHere is the ONLY data available — do not reference any metrics not listed here:\n\n{metrics_text}\n\nUsing ONLY the numbers above, deliver a complete organic social analysis. Break down performance by all available dimensions with exact figures. Calculate engagement rates, reach efficiency, and growth metrics where the data allows. Identify top and bottom performing content with exact metrics. Provide 5 specific content strategy optimizations with the data evidence that justifies each one, including optimal posting frequency and content mix recommendations.\n{WATERMARK}"
        analysis = query_llm(prompt)

        # Create visualization
        fig = go.Figure()
        try:
            chart_created = False

            # Option 1: Bar chart by post type
            if not post_type_metrics.empty and has_post_type and len(post_type_metrics) > 0:
                if has_engagement and has_impressions:
                    safe_df = post_type_metrics[post_type_metrics['impressions'] > 0].copy()
                    if not safe_df.empty:
                        safe_df['engagement_rate'] = (safe_df['engagement'] / safe_df['impressions'] * 100).round(2)
                        fig = px.bar(safe_df, x='post_type', y='engagement_rate', title='Engagement Rate by Post Type')
                        chart_created = True
                elif has_engagement:
                    fig = px.bar(post_type_metrics, x='post_type', y='engagement', title='Engagement by Post Type')
                    chart_created = True

            # Option 2: Time-series line chart (most common for X/social analytics exports)
            if not chart_created and 'date' in df.columns:
                df_sorted = df.sort_values('date')
                if has_impressions:
                    fig = px.line(df_sorted, x='date', y='impressions', title='Impressions Over Time')
                    chart_created = True
                elif has_engagement:
                    fig = px.line(df_sorted, x='date', y='engagement', title='Engagement Over Time')
                    chart_created = True
                elif has_likes:
                    fig = px.line(df_sorted, x='date', y='likes', title='Likes Over Time')
                    chart_created = True

            # Option 3: Bar by any available groupby column
            if not chart_created and has_impressions:
                for gc in ['channel', 'platform', 'campaign_name', 'segment']:
                    if gc in df.columns and df[gc].nunique() > 1:
                        grouped = df.groupby(gc)['impressions'].sum().reset_index().sort_values('impressions', ascending=False).head(15)
                        fig = px.bar(grouped, x=gc, y='impressions', title=f'Impressions by {gc.title()}')
                        chart_created = True
                        break
        except Exception as chart_err:
            fig = go.Figure()

        return analysis, fig_to_dict(fig)

    except Exception as e:
        return f"Error in analysis: {str(e)}", None


# ============================================================================
# TAB 6: DATA STORYTELLING
# ============================================================================

def analyze_data_storytelling(file, context_text):
    """Generate executive narrative from marketing dataset."""
    df, error = load_data(file)
    if error:
        return f"Error: {error}"

    df = normalize_columns(df)

    try:
        # Extract key summary statistics
        numeric_cols = df.select_dtypes(include=['number']).columns.tolist()

        # Build basic narrative metrics
        metrics_text = f"Dataset Overview:\n"
        metrics_text += f"Total Records: {len(df):,}\n"
        metrics_text += f"Columns: {', '.join(df.columns)}\n"
        metrics_text += f"Numeric Columns: {', '.join(numeric_cols[:10])}\n\n"

        for col in numeric_cols[:10]:
            try:
                metrics_text += f"\n{col}:\n"
                metrics_text += f"  Total: {float(df[col].sum()):,.2f}\n"
                metrics_text += f"  Average: {float(df[col].mean()):,.2f}\n"
                metrics_text += f"  Min: {float(df[col].min()):,.2f}\n"
                metrics_text += f"  Max: {float(df[col].max()):,.2f}\n"
            except (ValueError, TypeError):
                metrics_text += f"\n{col}: non-numeric or mixed types\n"

        context_prompt = f"Additional Context: {context_text}\n\n" if context_text else ""

        # Query LLM for narrative
        prompt = f"{ANALYST_TONE}\n\nHere is the ONLY data available to you — do not reference any metrics not listed here:\n\n{metrics_text}\n\n{context_prompt}Create a compelling executive narrative about this marketing dataset using ONLY the statistics provided above. Every number you mention must come directly from the data above. Structure as: (1) Opening headline stat that captures the most important finding, (2) Performance story arc showing what happened and why the numbers moved the way they did, (3) Key turning points or anomalies in the data, (4) Business impact section quantifying what the results mean, (5) Forward-looking recommendations grounded in the data patterns. Do NOT mention any metrics not present in the data.\n{WATERMARK}"
        analysis_text = query_llm(prompt, max_tokens=8000)

        return analysis_text

    except Exception as e:
        return f"Error in analysis: {str(e)}"


# ============================================================================
# TAB 7: VISUALIZATIONS
# ============================================================================

def create_visualization(file, chart_type, x_col, y_col, group_col=None):
    """Create custom chart visualization."""
    df, error = load_data(file)
    if error:
        return None, f"Error: {error}"

    # Lowercase column names for consistency
    df.columns = df.columns.str.lower().str.strip()

    try:
        # Handle empty strings as None
        if not x_col or x_col.strip() == '':
            return None, f"Please select an X column. Available: {', '.join(df.columns)}"
        if y_col and y_col.strip() == '':
            y_col = None

        # Case-insensitive column matching
        x_col = x_col.strip().lower()
        if y_col:
            y_col = y_col.strip().lower()

        # Check if columns exist
        if x_col not in df.columns:
            return None, f"Column '{x_col}' not found. Available: {', '.join(df.columns)}"
        if y_col and y_col not in df.columns:
            return None, f"Column '{y_col}' not found. Available: {', '.join(df.columns)}"

        fig = go.Figure()

        try:
            if chart_type == "scatter" and y_col:
                fig = px.scatter(df, x=x_col, y=y_col, title=f"{y_col} vs {x_col}")
            elif chart_type == "bar" and y_col:
                fig = px.bar(df, x=x_col, y=y_col, title=f"{y_col} by {x_col}")
            elif chart_type == "line" and y_col:
                fig = px.line(df, x=x_col, y=y_col, title=f"{y_col} over {x_col}")
            elif chart_type == "histogram":
                fig = px.histogram(df, x=x_col, nbins=30, title=f"Distribution of {x_col}")
            elif chart_type == "pie" and y_col:
                fig = px.pie(df, names=x_col, values=y_col, title=f"{y_col} by {x_col}")
            else:
                return None, "Unsupported chart type or missing required columns"
        except Exception as chart_error:
            return None, f"Error creating chart: {str(chart_error)}"

        return fig_to_dict(fig), "Chart created successfully"

    except Exception as e:
        return None, f"Error creating chart: {str(e)}"


# ============================================================================
# TAB 8: REPORTS
# ============================================================================

def _build_report_metrics(df):
    """Build metrics text from DataFrame for report generation."""
    metrics_text = f"Dataset: {len(df)} rows, {len(df.columns)} columns\n"
    metrics_text += f"Columns: {', '.join(df.columns.tolist())}\n\n"
    numeric_cols = df.select_dtypes(include=['number']).columns.tolist()
    metrics_text += "KEY METRICS:\n"
    for col in numeric_cols[:10]:
        try:
            total = float(df[col].sum())
            avg = float(df[col].mean())
            mn = float(df[col].min())
            mx = float(df[col].max())
            metrics_text += f"\n{col}:\n"
            metrics_text += f"  Total: {total:,.2f}\n"
            metrics_text += f"  Average: {avg:,.2f}\n"
            metrics_text += f"  Min: {mn:,.2f} | Max: {mx:,.2f}\n"
        except (ValueError, TypeError):
            metrics_text += f"\n{col}: non-numeric or mixed types\n"
    return metrics_text


def _generate_pdf_report(df, report_text):
    """Generate a PDF report file and return the file path."""
    from reportlab.lib.pagesizes import letter
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle
    from reportlab.lib import colors
    from reportlab.lib.units import inch
    import tempfile

    tmp = tempfile.NamedTemporaryFile(suffix='.pdf', delete=False, dir='/tmp')
    doc = SimpleDocTemplate(tmp.name, pagesize=letter, topMargin=0.75*inch, bottomMargin=0.75*inch)
    styles = getSampleStyleSheet()

    # Custom styles
    title_style = ParagraphStyle('ReportTitle', parent=styles['Title'], fontSize=22, textColor=colors.HexColor('#0a1628'), spaceAfter=20)
    body_style = ParagraphStyle('ReportBody', parent=styles['Normal'], fontSize=10, leading=14, spaceAfter=8)
    section_style = ParagraphStyle('ReportSection', parent=styles['Heading2'], fontSize=14, textColor=colors.HexColor('#0a1628'), spaceBefore=16, spaceAfter=8)

    elements = []
    elements.append(Paragraph("Marketing Analytics Report", title_style))
    elements.append(Paragraph(f"Generated: {datetime.now().strftime('%B %d, %Y %I:%M %p')}", body_style))
    elements.append(Spacer(1, 20))

    # Dataset overview table
    numeric_cols = df.select_dtypes(include=['number']).columns.tolist()
    if numeric_cols:
        elements.append(Paragraph("Dataset Overview", section_style))
        table_data = [['Metric', 'Total', 'Average', 'Min', 'Max']]
        for col in numeric_cols[:10]:
            try:
                table_data.append([
                    col.replace('_', ' ').title(),
                    f"{float(df[col].sum()):,.2f}",
                    f"{float(df[col].mean()):,.2f}",
                    f"{float(df[col].min()):,.2f}",
                    f"{float(df[col].max()):,.2f}"
                ])
            except:
                pass
        if len(table_data) > 1:
            t = Table(table_data, repeatRows=1)
            t.setStyle(TableStyle([
                ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor('#0a1628')),
                ('TEXTCOLOR', (0, 0), (-1, 0), colors.HexColor('#c8a951')),
                ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
                ('FONTSIZE', (0, 0), (-1, -1), 9),
                ('GRID', (0, 0), (-1, -1), 0.5, colors.grey),
                ('ROWBACKGROUNDS', (0, 1), (-1, -1), [colors.white, colors.HexColor('#f5f5f5')]),
                ('ALIGN', (1, 0), (-1, -1), 'RIGHT'),
            ]))
            elements.append(t)
            elements.append(Spacer(1, 20))

    # AI Analysis
    elements.append(Paragraph("AI Analysis", section_style))
    for line in report_text.split('\n'):
        line = line.strip()
        if not line:
            elements.append(Spacer(1, 6))
        elif line.startswith('#'):
            elements.append(Paragraph(line.lstrip('#').strip(), section_style))
        else:
            # Escape HTML characters
            safe_line = line.replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;')
            elements.append(Paragraph(safe_line, body_style))

    # Footer
    elements.append(Spacer(1, 30))
    footer_style = ParagraphStyle('Footer', parent=styles['Normal'], fontSize=7, textColor=colors.grey, alignment=1)
    elements.append(Paragraph("Analysis powered by QuantusData.ai Marketing Analytics Suite", footer_style))

    doc.build(elements)
    return tmp.name


def _generate_pptx_report(df, report_text):
    """Generate a PowerPoint report file and return the file path."""
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    import tempfile

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Title slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = RGBColor(0x0a, 0x16, 0x28)

    txBox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11), Inches(2))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Marketing Analytics Report"
    p.font.size = Pt(44)
    p.font.color.rgb = RGBColor(0xc8, 0xa9, 0x51)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    p2 = tf.add_paragraph()
    p2.text = f"Generated {datetime.now().strftime('%B %d, %Y')}"
    p2.font.size = Pt(18)
    p2.font.color.rgb = RGBColor(0x88, 0x92, 0xa4)
    p2.alignment = PP_ALIGN.CENTER

    p3 = tf.add_paragraph()
    p3.text = "Powered by QuantusData.ai"
    p3.font.size = Pt(14)
    p3.font.color.rgb = RGBColor(0x88, 0x92, 0xa4)
    p3.alignment = PP_ALIGN.CENTER

    # Data overview slide
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    bg2 = slide2.background.fill
    bg2.solid()
    bg2.fore_color.rgb = RGBColor(0x0a, 0x16, 0x28)

    txBox2 = slide2.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    tf2 = txBox2.text_frame
    p = tf2.paragraphs[0]
    p.text = "Dataset Overview"
    p.font.size = Pt(32)
    p.font.color.rgb = RGBColor(0xc8, 0xa9, 0x51)
    p.font.bold = True

    numeric_cols = df.select_dtypes(include=['number']).columns.tolist()
    overview_text = f"Records: {len(df):,}  |  Columns: {len(df.columns)}\n\n"
    for col in numeric_cols[:8]:
        try:
            overview_text += f"{col.replace('_', ' ').title()}: Total {float(df[col].sum()):,.2f}  |  Avg {float(df[col].mean()):,.2f}\n"
        except:
            pass

    txBox3 = slide2.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(12), Inches(5.5))
    tf3 = txBox3.text_frame
    tf3.word_wrap = True
    p = tf3.paragraphs[0]
    p.text = overview_text
    p.font.size = Pt(16)
    p.font.color.rgb = RGBColor(0xf8, 0xf9, 0xfc)

    # Analysis slides — split report into chunks
    lines = report_text.split('\n')
    chunk_size = 12
    chunks = [lines[i:i+chunk_size] for i in range(0, len(lines), chunk_size)]

    for idx, chunk in enumerate(chunks[:8]):  # Max 8 analysis slides
        slide_n = prs.slides.add_slide(prs.slide_layouts[6])
        bg_n = slide_n.background.fill
        bg_n.solid()
        bg_n.fore_color.rgb = RGBColor(0x0a, 0x16, 0x28)

        txBox_title = slide_n.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
        tf_title = txBox_title.text_frame
        p = tf_title.paragraphs[0]
        p.text = f"Analysis {f'(continued)' if idx > 0 else ''}"
        p.font.size = Pt(28)
        p.font.color.rgb = RGBColor(0xc8, 0xa9, 0x51)
        p.font.bold = True

        txBox_body = slide_n.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(12), Inches(5.8))
        tf_body = txBox_body.text_frame
        tf_body.word_wrap = True
        p = tf_body.paragraphs[0]
        p.text = '\n'.join(chunk)
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(0xf8, 0xf9, 0xfc)
        p.line_spacing = Pt(20)

    tmp = tempfile.NamedTemporaryFile(suffix='.pptx', delete=False, dir='/tmp')
    prs.save(tmp.name)
    return tmp.name


def generate_report(file, report_type):
    """Generate executive summary report, PDF, or PowerPoint."""
    df, error = load_data(file)
    if error:
        return f"Error: {error}", None

    df = normalize_columns(df)

    try:
        metrics_text = f"EXECUTIVE SUMMARY REPORT\nReport Type: {report_type}\nGenerated: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}\n\n"
        metrics_text += _build_report_metrics(df)

        # Query LLM for analysis text
        prompt = f"{ANALYST_TONE}\n\nHere is the ONLY data available — do not reference any metrics not listed here:\n\n{metrics_text}\n\nUsing ONLY the numbers above, generate a professional executive report. Structure it with: (1) Executive Summary with the 3 most important findings and exact metrics, (2) Detailed Performance Analysis covering every major metric dimension in the data, (3) Trend Analysis identifying patterns and anomalies, (4) Conclusions — 5 definitive statements about what the data proves, (5) Strategic Recommendations — 6 specific, prioritized actions each tied to a specific data finding, (6) Risk Watch List — metrics that need monitoring with specific threshold values from the data. Do not invent any statistics not present in the data.\n{WATERMARK}"
        report_text = query_llm(prompt, max_tokens=8000)

        # Handle PDF export
        if report_type == 'pdf_export':
            try:
                pdf_path = _generate_pdf_report(df, report_text)
                return report_text + "\n\n📥 PDF report generated. Use the download button below.", pdf_path
            except Exception as pdf_err:
                return report_text + f"\n\n⚠️ PDF generation failed: {str(pdf_err)}", None

        # Handle PowerPoint export
        if report_type == 'presentation':
            try:
                pptx_path = _generate_pptx_report(df, report_text)
                return report_text + "\n\n📥 PowerPoint presentation generated. Use the download button below.", pptx_path
            except Exception as pptx_err:
                return report_text + f"\n\n⚠️ PowerPoint generation failed: {str(pptx_err)}", None

        return report_text, None

    except Exception as e:
        return f"Error generating report: {str(e)}", None


# ============================================================================
# TAB 9: DATA CLEANING
# ============================================================================

def clean_data(file, remove_duplicates, fill_method):
    """Clean and prepare data."""
    df, error = load_data(file)
    if error:
        return f"Error: {error}"

    try:
        initial_rows = len(df)
        initial_nulls = df.isnull().sum().sum()

        # Remove duplicates
        if remove_duplicates:
            df = df.drop_duplicates()
            removed_rows = initial_rows - len(df)
        else:
            removed_rows = 0

        # Fill missing values
        if fill_method == "drop":
            df = df.dropna()
        elif fill_method == "mean":
            numeric_cols = df.select_dtypes(include=['number']).columns
            df[numeric_cols] = df[numeric_cols].fillna(df[numeric_cols].mean())
        elif fill_method == "forward_fill":
            df = df.fillna(method='ffill')
        elif fill_method == "zero":
            df = df.fillna(0)

        final_nulls = df.isnull().sum().sum()

        report = f"""DATA CLEANING REPORT

Initial Rows: {initial_rows:,}
Final Rows: {len(df):,}
Rows Removed (Duplicates): {removed_rows}
Initial Missing Values: {initial_nulls}
Final Missing Values: {final_nulls}
Fill Method: {fill_method}
Columns: {len(df.columns)}

Data is ready for analysis."""

        return report

    except Exception as e:
        return f"Error cleaning data: {str(e)}"


# ============================================================================
# TAB 10: MULTI-DATASET ANALYSIS
# ============================================================================

def analyze_multi_dataset(file1, file2, file3, analysis_type, custom_instructions):
    """Compare 2-3 datasets with analysis."""
    try:
        dfs = []
        file_names = []

        for file_obj, name in [(file1, "Dataset 1"), (file2, "Dataset 2"), (file3, "Dataset 3")]:
            if file_obj is None:
                continue
            if isinstance(file_obj, str) and not str(file_obj).strip():
                continue
            if hasattr(file_obj, 'name') and file_obj.name in (None, ''):
                continue
            df, error = load_data(file_obj)
            if error:
                continue
            df = normalize_columns(df)
            dfs.append(df)
            file_names.append(name)

        if len(dfs) < 2:
            return "Error: Please upload at least 2 datasets"

        # Compare datasets
        num_datasets = len(dfs)
        comparison_text = f"MULTI-DATASET COMPARISON ({analysis_type})\n"
        comparison_text += f"Number of datasets: {num_datasets}\n\n"

        for i, (df, name) in enumerate(zip(dfs, file_names)):
            comparison_text += f"\n{name}:\n"
            comparison_text += f"  Rows: {len(df)}\n"
            comparison_text += f"  Columns: {', '.join(df.columns)}\n"

            numeric_cols = df.select_dtypes(include=['number']).columns
            if len(numeric_cols) > 0:
                comparison_text += f"  Numeric Metrics ({len(numeric_cols)} columns):\n"
                for col in numeric_cols[:15]:
                    try:
                        total = float(df[col].sum())
                        avg = float(df[col].mean())
                        mn = float(df[col].min())
                        mx = float(df[col].max())
                        comparison_text += f"    {col}: Total={total:,.2f}, Avg={avg:,.2f}, Min={mn:,.2f}, Max={mx:,.2f}\n"
                    except (ValueError, TypeError):
                        comparison_text += f"    {col}: non-numeric or mixed types\n"

        # Add custom instructions
        if custom_instructions:
            comparison_text += f"\n\nCustom Analysis Instructions:\n{custom_instructions}\n"

        # Query LLM
        prompt = f"{ANALYST_TONE}\n\nHere is the ONLY data available — do not reference any metrics not listed here:\n\n{comparison_text}\n\nUsing ONLY the numbers above, deliver a complete multi-dataset comparison analysis. For each metric present in multiple datasets, calculate the exact difference and percentage variance. Identify which dataset outperforms on each dimension. Highlight the 3 most significant performance gaps with exact figures. Draw 4-5 definitive conclusions from the comparison. Provide 5 specific recommendations for how to apply learnings from the stronger dataset to improve the weaker one. Do not invent any statistics not present in the data.\n{WATERMARK}"
        analysis_text = query_llm(prompt, max_tokens=8000)

        return analysis_text

    except Exception as e:
        return f"Error in multi-dataset analysis: {str(e)}"


# ============================================================================
# TAB 11: AI CHAT
# ============================================================================

def chat_with_data(file, message):
    """Chat with data - interactive analysis."""
    try:
        df, error = load_data(file)
        if error:
            return f"Error: {error}"

        df = normalize_columns(df)

        # Build context from data
        context = f"Dataset Info:\n"
        context += f"Rows: {len(df)}\n"
        context += f"Columns: {', '.join(df.columns)}\n\n"

        numeric_cols = df.select_dtypes(include=['number']).columns
        context += "Numeric Columns Summary:\n"
        for col in numeric_cols[:10]:
            try:
                col_min = float(df[col].min())
                col_max = float(df[col].max())
                col_avg = float(df[col].mean())
                context += f"{col}: min={col_min:.2f}, max={col_max:.2f}, avg={col_avg:.2f}\n"
            except (ValueError, TypeError):
                context += f"{col}: non-numeric or mixed types\n"

        # Chat with LLM
        prompt = f"{ANALYST_TONE}\n\nHere is the ONLY data available — do not reference any metrics not listed here:\n\n{context}\n\nUser Question: {message}\n\nUsing ONLY the data above, provide a helpful analysis response. Structure your response as follows:\n\n1. DIRECT ANSWER: Answer the question directly using only the available data.\n2. ANALYSIS: Provide deeper analytical insight about what the numbers reveal.\n3. PYTHON CODE: Provide a working pandas code snippet the analyst can use to explore this further. Use the actual column names from the dataset above. Format the code in a ```python code block.\n4. SQL EQUIVALENT: Provide the equivalent SQL query for analysts working in a database environment. Format in a ```sql code block.\n5. FOLLOW-UP: Suggest 2-3 related questions the analyst should investigate next.\n\nDo not invent any statistics not present in the data.\n{WATERMARK}"
        response = query_llm(prompt)

        return response

    except Exception as e:
        return f"Error in chat: {str(e)}"


# ============================================================================
# GRADIO APP CREATION
# ============================================================================

def create_app():
    """Create Gradio interface with all tabs."""
    with gr.Blocks(title="Marketing Analytics Suite") as demo:
        gr.Markdown("# Marketing Analytics Suite")
        gr.Markdown("Comprehensive marketing data analysis powered by QuantusData.ai")

        # TAB 0: Paid Media
        with gr.Tab("Paid Media Analytics"):
            with gr.Column():
                file_input_0 = gr.File(label="Upload CSV or Excel", file_types=[".csv", ".xlsx", ".xls"])
                analyze_btn_0 = gr.Button("Analyze Paid Media")
                output_text_0 = gr.Textbox(label="Analysis", lines=10)
                output_chart_0 = gr.JSON(label="Chart Data")
                analyze_btn_0.click(analyze_paid_media, inputs=[file_input_0], outputs=[output_text_0, output_chart_0])

        # TAB 1: Creative Performance
        with gr.Tab("Creative Performance"):
            with gr.Column():
                file_input_1 = gr.File(label="Upload CSV or Excel", file_types=[".csv", ".xlsx", ".xls"])
                analyze_btn_1 = gr.Button("Analyze Creatives")
                output_text_1 = gr.Textbox(label="Analysis", lines=10)
                output_chart_1 = gr.JSON(label="Chart Data")
                analyze_btn_1.click(analyze_creative_performance, inputs=[file_input_1], outputs=[output_text_1, output_chart_1])

        # TAB 2: Brand Lift & Attribution
        with gr.Tab("Brand Lift & Attribution"):
            with gr.Column():
                file_input_2 = gr.File(label="Upload CSV or Excel", file_types=[".csv", ".xlsx", ".xls"])
                analyze_btn_2 = gr.Button("Analyze Brand Lift")
                output_text_2 = gr.Textbox(label="Analysis", lines=10)
                output_chart_2 = gr.JSON(label="Chart Data")
                analyze_btn_2.click(analyze_brand_lift, inputs=[file_input_2], outputs=[output_text_2, output_chart_2])

        # TAB 3: Audience & Segmentation
        with gr.Tab("Audience & Segmentation"):
            with gr.Column():
                file_input_3 = gr.File(label="Upload CSV or Excel", file_types=[".csv", ".xlsx", ".xls"])
                analyze_btn_3 = gr.Button("Analyze Audience")
                output_text_3 = gr.Textbox(label="Analysis", lines=10)
                output_chart_3 = gr.JSON(label="Chart Data")
                analyze_btn_3.click(analyze_audience, inputs=[file_input_3], outputs=[output_text_3, output_chart_3])

        # TAB 4: Campaign Reporting
        with gr.Tab("Campaign Reporting"):
            with gr.Column():
                file_input_4 = gr.File(label="Upload CSV or Excel", file_types=[".csv", ".xlsx", ".xls"])
                analyze_btn_4 = gr.Button("Generate Campaign Report")
                output_text_4 = gr.Textbox(label="Analysis", lines=10)
                output_chart_4 = gr.JSON(label="Chart Data")
                analyze_btn_4.click(analyze_campaign_reporting, inputs=[file_input_4], outputs=[output_text_4, output_chart_4])

        # TAB 5: Organic Social
        with gr.Tab("Organic Social"):
            with gr.Column():
                file_input_5 = gr.File(label="Upload CSV or Excel", file_types=[".csv", ".xlsx", ".xls"])
                analyze_btn_5 = gr.Button("Analyze Organic Social")
                output_text_5 = gr.Textbox(label="Analysis", lines=10)
                output_chart_5 = gr.JSON(label="Chart Data")
                analyze_btn_5.click(analyze_organic_social, inputs=[file_input_5], outputs=[output_text_5, output_chart_5])

        # TAB 6: Data Storytelling
        with gr.Tab("Data Storytelling"):
            with gr.Column():
                file_input_6 = gr.File(label="Upload CSV or Excel", file_types=[".csv", ".xlsx", ".xls"])
                context_input = gr.Textbox(label="Additional Context (optional)", lines=3)
                analyze_btn_6 = gr.Button("Generate Narrative")
                output_text_6 = gr.Textbox(label="Executive Narrative", lines=15)
                analyze_btn_6.click(analyze_data_storytelling, inputs=[file_input_6, context_input], outputs=[output_text_6])

        # TAB 7: Visualizations
        with gr.Tab("Custom Visualizations"):
            with gr.Column():
                file_input_7 = gr.File(label="Upload CSV or Excel", file_types=[".csv", ".xlsx", ".xls"])
                x_col_input = gr.Textbox(label="X Column", placeholder="e.g., channel, date, segment")
                y_col_input = gr.Textbox(label="Y Column (optional)", placeholder="e.g., spend, conversions")
                chart_type_input = gr.Dropdown(
                    choices=["scatter", "bar", "line", "histogram", "pie"],
                    label="Chart Type"
                )
                group_col_input = gr.Textbox(label="Group By (optional)", placeholder="e.g., channel, segment")
                create_viz_btn = gr.Button("Create Chart")
                output_chart_7 = gr.JSON(label="Chart")
                output_text_7 = gr.Textbox(label="Message")
                create_viz_btn.click(create_visualization, inputs=[file_input_7, chart_type_input, x_col_input, y_col_input, group_col_input],
                                    outputs=[output_chart_7, output_text_7])

        # TAB 8: Reports
        with gr.Tab("Executive Reports"):
            with gr.Column():
                file_input_8 = gr.File(label="Upload CSV or Excel", file_types=[".csv", ".xlsx", ".xls"])
                report_type_input = gr.Dropdown(
                    choices=["executive_summary", "channel_deep_dive", "creative_audit", "budget_optimization", "pdf_export", "presentation"],
                    label="Report Type"
                )
                report_btn = gr.Textbox(label="Company Name (Business Plan — White Label)", placeholder="Leave blank for QuantusData branding", value=""),
                    gr.Button("Generate Report")
                output_text_8 = gr.Textbox(label="Report", lines=15)
                output_file_8 = gr.File(label="Download Report", visible=True)
                report_btn.click(generate_report, inputs=[file_input_8, report_type_input], outputs=[output_text_8, output_file_8])

        # TAB 9: Data Cleaning
        with gr.Tab("Data Cleaning"):
            with gr.Column():
                file_input_9 = gr.File(label="Upload CSV or Excel", file_types=[".csv", ".xlsx", ".xls"])
                remove_dup_check = gr.Checkbox(label="Remove Duplicates", value=True)
                fill_method_input = gr.Dropdown(
                    choices=["drop", "mean", "forward_fill", "zero"],
                    label="Fill Missing Values Method",
                    value="zero"
                )
                clean_btn = gr.Button("Clean Data")
                output_text_9 = gr.Textbox(label="Cleaning Report")
                clean_btn.click(clean_data, inputs=[file_input_9, remove_dup_check, fill_method_input], outputs=[output_text_9])

        # TAB 10: Multi-Dataset Analysis
        with gr.Tab("Multi-Dataset Analysis"):
            with gr.Column():
                file_input_10a = gr.File(label="Dataset 1", file_types=[".csv", ".xlsx", ".xls"])
                file_input_10b = gr.File(label="Dataset 2", file_types=[".csv", ".xlsx", ".xls"])
                file_input_10c = gr.File(label="Dataset 3 (optional)", file_types=[".csv", ".xlsx", ".xls"])
                analysis_type_input = gr.Dropdown(
                    choices=["Comparison", "Trend Analysis", "Cohort Analysis"],
                    label="Analysis Type"
                )
                custom_instr_input = gr.Textbox(label="Custom Instructions (optional)", lines=3)
                multi_btn = gr.Button("Compare Datasets")
                output_text_10 = gr.Textbox(label="Analysis", lines=15)
                multi_btn.click(analyze_multi_dataset, inputs=[file_input_10a, file_input_10b, file_input_10c, analysis_type_input, custom_instr_input],
                               outputs=[output_text_10])

        # TAB 11: AI Chat
        with gr.Tab("AI Data Chat"):
            with gr.Column():
                file_input_11 = gr.File(label="Upload CSV or Excel", file_types=[".csv", ".xlsx", ".xls"])
                message_input = gr.Textbox(label="Ask a question about your data", lines=3)
                chat_btn = gr.Button("Ask")
                output_text_11 = gr.Textbox(label="Response", lines=10)
                chat_btn.click(chat_with_data, inputs=[file_input_11, message_input], outputs=[output_text_11])

    return demo


if __name__ == "__main__":
    app = create_app()
    app.launch(share=False, server_name="0.0.0.0", server_port=7876)
