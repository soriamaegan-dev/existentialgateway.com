import gradio as gr
import requests
import os
import pandas as pd
import json
import io
import ast
import plotly.express as px
import plotly.graph_objects as go
from plotly.subplots import make_subplots
from pypdf import PdfReader
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from reportlab.lib.pagesizes import letter
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import inch
from reportlab.lib import colors
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
import tempfile
from datetime import datetime

CRIME_SYSTEM_PROMPT = """You are a senior crime analyst and law enforcement intelligence specialist with expertise in predictive policing, crime pattern analysis, recidivism modeling, and criminal justice data analytics.

ANALYTICAL STANDARDS:
1. Apply crime analysis frameworks: tactical (immediate incidents), operational (patterns/series), strategic (long-term trends).
2. Use standard crime classification: UCR/NIBRS Part I and Part II offenses, FBI crime definitions, and local code cross-references.
3. Apply spatial analysis concepts: crime hot spots (kernel density estimation), geographic profiling, displacement effects, diffusion of benefits.
4. Calculate key metrics: crime rate per 100,000 population, clearance rate, recidivism rate (3-year return to custody benchmark: 40-60% nationally), response time analysis.
5. Apply risk assessment frameworks: LSI-R, ORAS, COMPAS principles — identify criminogenic risk factors and protective factors.
6. Conduct temporal analysis: time-of-day patterns, day-of-week patterns, seasonal variations, special event impacts.
7. Structure every analysis with: Crime Overview → Trend Analysis → Hot Spot Identification → Offender/Victim Demographics → Pattern Recognition → Predictive Assessment → Resource Deployment Recommendations.
8. Apply constitutional considerations: Fourth Amendment implications for predictive policing, bias auditing requirements for algorithmic tools.
9. Reference FBI UCR/NIBRS standards, NIJ research benchmarks, and Bureau of Justice Statistics data for comparisons.
10. Flag data quality issues: dark figure of crime (unreported), reporting rate variations, reclassification issues that affect trend validity."""


HF_TOKEN = os.environ.get("HF_TOKEN", "")

DISCLAIMER = """
> **LEGAL DISCLAIMER**: This tool is AI-generated for law enforcement planning purposes only.
> All predictions and analyses must be verified by qualified law enforcement professionals.
> Do not make arrest or prosecution decisions based solely on AI output.
> © 2026 Existential Gateway, LLC. All Rights Reserved. Proprietary Software.
"""

WAIT_MSG = "*Results take approximately 1-2 minutes to generate. Please do not click multiple times.*"

WATERMARK = """
---
© 2026 Existential Gateway, LLC | AI Law Enforcement and Crime Analyzer
Unauthorized reproduction strictly prohibited. Licensing: existentialgateway@gmail.com
---
"""

PII_WARNING = """
> **DATA PRIVACY NOTICE**: Before uploading any data you MUST remove all personally
> identifiable information including: Victim Names | Suspect Names | SSN |
> Date of Birth | Address | Phone Number | Juvenile Records.
> This platform is NOT intended to store personally identifiable information.
"""



def fetch_url_content(text):
    """Detect URLs in user text, fetch their content, and return as context."""
    import re
    urls = re.findall(r'https?://[^\s<>"{}|\^`\[\]]+', text)
    if not urls:
        return ""
    fetched = []
    for url in urls[:3]:  # limit to 3 URLs
        try:
            headers = {
                "User-Agent": "Mozilla/5.0 (compatible; ExistentialGateway/1.0)",
                "Accept": "text/html,application/xhtml+xml,application/xml;q=0.9,*/*;q=0.8"
            }
            r = requests.get(url, timeout=15, headers=headers)
            if r.status_code == 200:
                from html.parser import HTMLParser
                class TextExtractor(HTMLParser):
                    def __init__(self):
                        super().__init__()
                        self.text = []
                        self.skip = False
                    def handle_starttag(self, tag, attrs):
                        if tag in ('script', 'style', 'nav', 'footer', 'header'):
                            self.skip = True
                    def handle_endtag(self, tag):
                        if tag in ('script', 'style', 'nav', 'footer', 'header'):
                            self.skip = False
                    def handle_data(self, data):
                        if not self.skip and data.strip():
                            self.text.append(data.strip())
                parser = TextExtractor()
                parser.feed(r.text)
                page_text = ' '.join(parser.text)[:3000]
                fetched.append(f"[URL: {url}]\n{page_text}")
        except Exception as e:
            fetched.append(f"[URL: {url}] Could not fetch: {str(e)}")
    if fetched:
        return "\n\nUSER-PROVIDED URL CONTENT (analyze this as part of your response):\n" + "\n---\n".join(fetched)
    return ""



def enrich_with_urls(text):
    """If text contains URLs, fetch and append their content."""
    if not text or 'http' not in text:
        return text
    url_content = fetch_url_content(text)
    if url_content:
        return text + url_content
    return text


def query_llm(prompt):
    API_KEY = os.environ.get("OPENAI_API_KEY", "")
    headers = {"Authorization": f"Bearer {API_KEY}", "Content-Type": "application/json"}
    msgs = [{"role": "system", "content": CRIME_SYSTEM_PROMPT}, {"role": "user", "content": prompt}]
    payload = {"model": "gpt-4o", "max_tokens": 4000, "messages": msgs}
    try:
        response = requests.post("https://api.openai.com/v1/chat/completions", headers=headers, json=payload, timeout=240)
        result = response.json()
        if "choices" in result:
            return result["choices"][0]["message"]["content"]
        return f"API Error: {result}"
    except Exception as e:
        return f"Error: {str(e)}"


def load_data(file):
    if file is None:
        return None, "No file uploaded."
    try:
        path = file.name if hasattr(file, "name") else file
        if path.endswith(".csv"):
            df = pd.read_csv(path)
        elif path.endswith((".xlsx", ".xls")):
            df = pd.read_excel(path)
        elif path.endswith(".json"):
            df = pd.read_json(path)
        elif path.endswith(".pdf"):
            reader = PdfReader(path)
            text = "\n".join(page.extract_text() or "" for page in reader.pages)
            return None, text
        else:
            return None, "Unsupported file type."
        return df, None
    except Exception as e:
        return None, f"Error loading file: {str(e)}"


# ─── Tab 1: Data Upload and Overview ─────────────────────────────────────────

def analyze_overview(file):
    if file is None:
        return "Please upload a file.", "", ""
    df, err = load_data(file)
    if err and df is None:
        if "Unsupported" in err or "Error loading" in err:
            return err, "", ""
        text_preview = err[:3000]
        prompt = f"""You are an expert crime data analyst. The user uploaded a PDF crime report. Write as the professional analyst who performed this work. State all findings directly with specific numbers, percentages, and comparisons. NEVER say 'this appears to be', 'the data seems to show', 'the dataset contains', or 'it looks like'. Instead, deliver results confidently: 'Average claim cost was $4,230, a 12% increase year-over-year' or 'Readmission rates of 18.3% in Cardiology exceeded the facility average by 4.1 percentage points.' Be direct, precise, and actionable. Present findings as YOUR analysis, not a description of a dataset. Content (truncated):
{text_preview}

Present your professional findings:
1. What type of crime report or dataset is this?
2. Key data elements and fields present
3. Geographic area and date range covered
4. Summary of main crime categories
5. Data quality observations
6. Recommended next analysis steps"""
        result = query_llm(prompt)
        return result + "\n\n" + WATERMARK, "PDF — text extracted", "See analysis above"

    rows, cols = df.shape
    buf = io.StringIO()
    df.info(buf=buf)
    dtypes = df.dtypes.to_string()
    missing = df.isnull().sum()
    missing_str = missing[missing > 0].to_string() if missing.any() else "None"
    preview = df.head(10).to_string()

    prompt = f"""You are an expert crime data analyst. Write as the professional analyst who performed this work. State all findings directly with specific numbers, percentages, and comparisons. NEVER say 'this appears to be', 'the data seems to show', 'the dataset contains', or 'it looks like'. Instead, deliver results confidently: 'Average claim cost was $4,230, a 12% increase year-over-year' or 'Readmission rates of 18.3% in Cardiology exceeded the facility average by 4.1 percentage points.' Be direct, precise, and actionable. Present findings as YOUR analysis, not a description of a dataset. Analyze this crime dataset overview:

Rows: {rows} | Columns: {cols}
Column types:\n{dtypes}
Missing values:\n{missing_str}
First 10 rows:\n{preview}

Provide:
1. What type of crime data this appears to be
2. Geographic area and date range if identifiable
3. Crime categories present
4. Data quality assessment
5. Key columns and their significance
6. Recommended analyses to run
7. Any immediate patterns or red flags"""

    result = query_llm(prompt)
    stats = f"**Rows:** {rows} | **Columns:** {cols}\n\n**Missing Values:**\n```\n{missing_str}\n```"
    return result + "\n\n" + WATERMARK, stats, f"```\n{preview}\n```"


# ─── Tab 2: Crime Pattern Analysis ───────────────────────────────────────────

def analyze_patterns(file):
    if file is None:
        return "Please upload a crime incident data file.", None
    df, err = load_data(file)
    if err:
        return err, None

    rows, cols = df.shape
    preview = df.head(20).to_string()
    desc = df.describe(include="all").to_string()

    prompt = f"""You are an expert crime analyst. Analyze this crime incident dataset for patterns:

Shape: {rows} rows x {cols} columns
Columns: {list(df.columns)}
Statistical Summary:\n{desc}
Sample Data:\n{preview}

Provide a comprehensive crime pattern analysis:
1. **Crime Hotspots** — identify high-frequency locations if location data present
2. **Time-Based Patterns** — day of week, time of day, monthly, seasonal trends
3. **Crime Type Frequency** — most common crime types and their trends
4. **Victim and Offender Demographics** — if demographic data present
5. **Seasonal Patterns** — peak crime periods
6. **Geographic Distribution** — spread and concentration
7. **Key Findings** — top 5 actionable insights
8. **Recommendations** — patrol and prevention strategies

Use specific numbers from the data. Be precise and professional."""

    result = query_llm(prompt)

    fig = None
    try:
        cat_cols = df.select_dtypes(include="object").columns.tolist()
        numeric_cols = df.select_dtypes(include="number").columns.tolist()
        if cat_cols:
            top_vals = df[cat_cols[0]].value_counts().head(10)
            fig = go.Figure(go.Bar(
                x=top_vals.index.tolist(),
                y=top_vals.values,
                marker_color="#e74c3c"
            ))
            fig.update_layout(
                title=f"Crime Frequency by {cat_cols[0]}",
                xaxis_title=cat_cols[0],
                yaxis_title="Count",
                template="plotly_dark",
                height=400
            )
        elif numeric_cols:
            fig = px.histogram(df, x=numeric_cols[0], template="plotly_dark",
                               title=f"{numeric_cols[0]} Distribution", nbins=30)
            fig.update_layout(height=400)
    except Exception:
        fig = None

    return result + "\n\n" + WATERMARK, fig


# ─── Tab 3: Predictive Analysis ───────────────────────────────────────────────

def analyze_predictive(file):
    if file is None:
        return "Please upload historical crime data.", None
    df, err = load_data(file)
    if err:
        return err, None

    rows, cols = df.shape
    preview = df.head(20).to_string()
    desc = df.describe(include="all").to_string()

    prompt = f"""You are an expert predictive crime analyst. Write as the professional analyst who performed this work. State all findings directly with specific numbers, percentages, and comparisons. NEVER say 'this appears to be', 'the data seems to show', 'the dataset contains', or 'it looks like'. Instead, deliver results confidently: 'Average claim cost was $4,230, a 12% increase year-over-year' or 'Readmission rates of 18.3% in Cardiology exceeded the facility average by 4.1 percentage points.' Be direct, precise, and actionable. Present findings as YOUR analysis, not a description of a dataset. Analyze this historical crime data and generate predictions:

Shape: {rows} rows x {cols} columns
Columns: {list(df.columns)}
Statistical Summary:\n{desc}
Sample Data:\n{preview}

Provide a comprehensive predictive analysis:
1. **Future Crime Hotspots** — likely high-risk locations based on historical patterns
2. **Risk Scores by Zone** — assign Low/Medium/High risk to identifiable areas
3. **High Risk Time Periods** — predict peak crime windows (days, times, seasons)
4. **Crime Type Predictions** — which crimes are trending up or down
5. **Resource Allocation Recommendations** — where to deploy resources
6. **Patrol Deployment Suggestions** — specific recommendations with rationale
7. **Probability Scores** — estimated likelihood of crime occurrence by category
8. **Prevention Opportunities** — where intervention could reduce predicted crimes

Be specific. Use trend analysis from the data to support predictions."""

    result = query_llm(prompt)

    fig = None
    try:
        numeric_cols = df.select_dtypes(include="number").columns.tolist()
        if numeric_cols:
            fig = go.Figure()
            for col in numeric_cols[:3]:
                fig.add_trace(go.Scatter(
                    y=df[col],
                    mode="lines+markers",
                    name=col
                ))
            fig.update_layout(
                title="Historical Trend Analysis",
                template="plotly_dark",
                height=400,
                yaxis_title="Value"
            )
    except Exception:
        fig = None

    return result + "\n\n" + WATERMARK, fig


# ─── Tab 4: Case Correlation ──────────────────────────────────────────────────

def analyze_correlation(file):
    if file is None:
        return "Please upload incident reports or case files.", None
    df, err = load_data(file)
    if err and df is None:
        if "Unsupported" in err or "Error" in err:
            return err, None
        prompt = f"""You are an expert criminal investigator. Analyze these case files for correlations:

{err[:3000]}

Identify:
1. Write as the professional analyst who performed this work. State all findings directly with specific numbers, percentages, and comparisons. NEVER say 'this appears to be', 'the data seems to show', 'the dataset contains', or 'it looks like'. Instead, deliver results confidently: 'Average claim cost was $4,230, a 12% increase year-over-year' or 'Readmission rates of 18.3% in Cardiology exceeded the facility average by 4.1 percentage points.' Be direct, precise, and actionable. Present findings as YOUR analysis, not a description of a dataset. **Case Links** — connections between incidents
2. **MO Patterns** — similar modus operandi across cases
3. **Geographic Connections** — location-based links
4. **Timeline Correlations** — temporal patterns across cases
5. **Suspect Patterns** — behavioral patterns suggesting same offender(s)
6. **Serial Incident Indicators** — evidence of serial offending
7. **Investigation Recommendations** — next steps for investigators"""
        result = query_llm(prompt)
        return result + "\n\n" + WATERMARK, None

    rows, cols = df.shape
    preview = df.head(30).to_string()
    desc = df.describe(include="all").to_string()

    prompt = f"""You are an expert criminal investigator and crime analyst. Write as the professional analyst who performed this work. State all findings directly with specific numbers, percentages, and comparisons. NEVER say 'this appears to be', 'the data seems to show', 'the dataset contains', or 'it looks like'. Instead, deliver results confidently: 'Average claim cost was $4,230, a 12% increase year-over-year' or 'Readmission rates of 18.3% in Cardiology exceeded the facility average by 4.1 percentage points.' Be direct, precise, and actionable. Present findings as YOUR analysis, not a description of a dataset. Analyze this case data for correlations:

Shape: {rows} rows x {cols} columns
Columns: {list(df.columns)}
Statistical Summary:\n{desc}
Sample Data:\n{preview}

Identify:
1. **Case Links** — statistical connections between incidents
2. **MO Pattern Detection** — similar methods across cases
3. **Geographic Connections** — location clustering or pathways
4. **Timeline Correlation** — temporal links between incidents
5. **Suspect Pattern Matching** — behavioral signatures
6. **Serial Incident Identification** — probability of serial offending
7. **Investigation Recommendations** — prioritized next steps

Flag any strong correlations with confidence levels (High/Medium/Low)."""

    result = query_llm(prompt)

    fig = None
    try:
        numeric_cols = df.select_dtypes(include="number").columns.tolist()
        if len(numeric_cols) >= 2:
            corr = df[numeric_cols].corr()
            fig = px.imshow(
                corr,
                text_auto=True,
                template="plotly_dark",
                title="Case Data Correlation Matrix",
                color_continuous_scale="RdBu_r"
            )
            fig.update_layout(height=450)
    except Exception:
        fig = None

    return result + "\n\n" + WATERMARK, fig


# ─── Tab 5: Recidivism Analysis ───────────────────────────────────────────────

def analyze_recidivism(file):
    if file is None:
        return "Please upload offender data.", None
    df, err = load_data(file)
    if err:
        return err, None

    rows, cols = df.shape
    preview = df.head(20).to_string()
    desc = df.describe(include="all").to_string()

    prompt = f"""You are an expert criminologist specializing in recidivism analysis. Write as the professional analyst who performed this work. State all findings directly with specific numbers, percentages, and comparisons. NEVER say 'this appears to be', 'the data seems to show', 'the dataset contains', or 'it looks like'. Instead, deliver results confidently: 'Average claim cost was $4,230, a 12% increase year-over-year' or 'Readmission rates of 18.3% in Cardiology exceeded the facility average by 4.1 percentage points.' Be direct, precise, and actionable. Present findings as YOUR analysis, not a description of a dataset. Analyze this offender data:

Shape: {rows} rows x {cols} columns
Columns: {list(df.columns)}
Statistical Summary:\n{desc}
Sample Data:\n{preview}

Provide a comprehensive recidivism analysis:
1. **Recidivism Risk Scores** — overall population risk level (Low/Medium/High) with breakdown
2. **Key Risk Factors** — top factors driving reoffense likelihood
3. **Rehabilitation Program Effectiveness** — if program data present
4. **Reoffense Probability by Crime Type** — which offenses have highest recidivism rates
5. **Demographic Risk Patterns** — age, offense history, and other factors
6. **High Risk Individuals** — characteristics of highest risk offenders
7. **Intervention Recommendations** — specific programs and strategies to reduce recidivism
8. **Resource Prioritization** — where to focus rehabilitation resources

Use research-backed criminology standards. Be professional and unbiased."""

    result = query_llm(prompt)

    fig = None
    try:
        cat_cols = df.select_dtypes(include="object").columns.tolist()
        numeric_cols = df.select_dtypes(include="number").columns.tolist()
        if cat_cols and numeric_cols:
            top_cat = df[cat_cols[0]].value_counts().head(8)
            fig = go.Figure(go.Bar(
                x=top_cat.index.tolist(),
                y=top_cat.values,
                marker_color="#9b59b6"
            ))
            fig.update_layout(
                title=f"Offender Distribution by {cat_cols[0]}",
                template="plotly_dark",
                height=400
            )
        elif numeric_cols:
            fig = px.box(
                df,
                y=numeric_cols[:min(4, len(numeric_cols))],
                title="Offender Risk Variable Distributions",
                template="plotly_dark"
            )
            fig.update_layout(height=400)
    except Exception:
        fig = None

    return result + "\n\n" + WATERMARK, fig


# ─── Tab 6: Resource Optimization ────────────────────────────────────────────

def analyze_resources(file):
    if file is None:
        return "Please upload staffing and incident data.", None
    df, err = load_data(file)
    if err:
        return err, None

    rows, cols = df.shape
    preview = df.head(20).to_string()
    desc = df.describe(include="all").to_string()

    prompt = f"""You are an expert law enforcement resource analyst. Write as the professional analyst who performed this work. State all findings directly with specific numbers, percentages, and comparisons. NEVER say 'this appears to be', 'the data seems to show', 'the dataset contains', or 'it looks like'. Instead, deliver results confidently: 'Average claim cost was $4,230, a 12% increase year-over-year' or 'Readmission rates of 18.3% in Cardiology exceeded the facility average by 4.1 percentage points.' Be direct, precise, and actionable. Present findings as YOUR analysis, not a description of a dataset. Analyze this staffing and incident data:

Shape: {rows} rows x {cols} columns
Columns: {list(df.columns)}
Statistical Summary:\n{desc}
Sample Data:\n{preview}

Provide a comprehensive resource optimization analysis:
1. **Optimal Patrol Deployment** — recommended officer allocation by zone and shift
2. **Response Time Analysis** — current performance and improvement opportunities
3. **Resource Allocation by Zone** — which areas are over/under-resourced
4. **Peak Demand Periods** — when and where resources are most needed
5. **Budget Optimization** — cost-efficiency recommendations
6. **Staffing Level Analysis** — gaps and surpluses by time and location
7. **Technology Recommendations** — tools to improve efficiency
8. **Implementation Plan** — prioritized action steps

Be specific with numbers and percentages. Focus on actionable recommendations."""

    result = query_llm(prompt)

    fig = None
    try:
        numeric_cols = df.select_dtypes(include="number").columns.tolist()
        cat_cols = df.select_dtypes(include="object").columns.tolist()
        if numeric_cols and cat_cols:
            fig = px.bar(
                df.head(20),
                x=cat_cols[0],
                y=numeric_cols[0],
                template="plotly_dark",
                title=f"{numeric_cols[0]} by {cat_cols[0]}",
                color_discrete_sequence=["#2ecc71"]
            )
            fig.update_layout(height=400)
        elif numeric_cols:
            fig = make_subplots(rows=1, cols=min(2, len(numeric_cols)),
                                subplot_titles=numeric_cols[:2])
            for i, col in enumerate(numeric_cols[:2], 1):
                fig.add_trace(go.Histogram(x=df[col], name=col, nbinsx=20), row=1, col=i)
            fig.update_layout(title="Resource Data Distribution",
                              template="plotly_dark", height=400)
    except Exception:
        fig = None

    return result + "\n\n" + WATERMARK, fig


# ─── Tab 7: Visualizations and Maps ──────────────────────────────────────────

def generate_viz(file, chart_type, x_col, y_col, color_col):
    if file is None:
        return None, "Please upload a data file."
    df, err = load_data(file)
    if err:
        return None, err

    cols = df.columns.tolist()
    x = x_col if x_col in cols else cols[0]
    y = y_col if y_col in cols else (cols[1] if len(cols) > 1 else cols[0])
    color = color_col if color_col in cols else None

    try:
        if chart_type == "Bar Chart":
            fig = px.bar(df, x=x, y=y, color=color, template="plotly_dark",
                         title=f"{y} by {x}")
        elif chart_type == "Line Chart":
            fig = px.line(df, x=x, y=y, color=color, template="plotly_dark",
                          title=f"{y} over {x}")
        elif chart_type == "Pie Chart":
            fig = px.pie(df, names=x, values=y, template="plotly_dark",
                         title=f"{y} Distribution by {x}")
        elif chart_type == "Scatter Plot":
            fig = px.scatter(df, x=x, y=y, color=color, template="plotly_dark",
                             title=f"{x} vs {y}")
        elif chart_type == "Heatmap":
            numeric = df.select_dtypes(include="number")
            corr = numeric.corr()
            fig = px.imshow(corr, text_auto=True, template="plotly_dark",
                            title="Correlation Heatmap", color_continuous_scale="RdBu_r")
        elif chart_type == "Box Plot":
            fig = px.box(df, x=x if df[x].dtype == "object" else None,
                         y=y, color=color, template="plotly_dark",
                         title=f"{y} Distribution")
        elif chart_type == "Histogram":
            fig = px.histogram(df, x=x, color=color, template="plotly_dark",
                               title=f"{x} Frequency Distribution", nbins=40)
        elif chart_type == "Density Heatmap":
            fig = px.density_heatmap(df, x=x, y=y, template="plotly_dark",
                                     title=f"Density: {x} vs {y}",
                                     color_continuous_scale="Reds")
        else:
            fig = px.bar(df, x=x, y=y, template="plotly_dark")

        fig.update_layout(height=500, margin=dict(l=40, r=40, t=60, b=40))
        return fig, f"Chart generated: {chart_type} — {y} by {x}"
    except Exception as e:
        return None, f"Chart error: {str(e)}"


def get_columns(file):
    if file is None:
        return gr.Dropdown(choices=[]), gr.Dropdown(choices=[]), gr.Dropdown(choices=[])
    df, err = load_data(file)
    if err or df is None:
        return gr.Dropdown(choices=[]), gr.Dropdown(choices=[]), gr.Dropdown(choices=[])
    cols = df.columns.tolist()
    return gr.Dropdown(choices=cols), gr.Dropdown(choices=cols), gr.Dropdown(choices=cols)


# ─── Tab 8: Reports and Presentations ────────────────────────────────────────

def generate_report(file, report_type):
    if file is None:
        return "Please upload a data file.", None
    df, err = load_data(file)
    if err and df is None:
        return err, None

    if df is not None:
        preview = df.head(20).to_string()
        desc = df.describe(include="all").to_string()
        cols = list(df.columns)
        rows = len(df)
        data_summary = f"Rows: {rows} | Columns: {cols}\nStatistics:\n{desc}\nSample:\n{preview}"
    else:
        data_summary = err[:3000]

    prompt = f"""You are an expert law enforcement analyst writing a professional report. Write as the professional analyst who performed this work. State all findings directly with specific numbers, percentages, and comparisons. NEVER say 'this appears to be', 'the data seems to show', 'the dataset contains', or 'it looks like'. Instead, deliver results confidently: 'Average claim cost was $4,230, a 12% increase year-over-year' or 'Readmission rates of 18.3% in Cardiology exceeded the facility average by 4.1 percentage points.' Be direct, precise, and actionable. Present findings as YOUR analysis, not a description of a dataset. Analyze this crime data and produce a complete {report_type}:

{data_summary}

Structure:
1. **Executive Summary** (3-5 sentences)
2. **Data Overview** — what the data contains and covers
3. **Key Findings** — top 5 findings with specific numbers
4. **Crime Pattern Analysis** — notable trends and hotspots
5. **Risk Assessment** — overall risk level and primary drivers
6. **Predictive Outlook** — likely future trends
7. **Recommendations** — 5 specific actionable recommendations
8. **Conclusion**

Write professionally for a law enforcement audience. Use specific numbers."""

    report_text = query_llm(prompt)

    output_path = None
    try:
        if report_type == "PowerPoint Presentation":
            output_path = _make_pptx(report_text, df if df is not None else None)
        else:
            output_path = _make_pdf(report_text)
    except Exception as e:
        report_text += f"\n\n[Document generation error: {str(e)}]"

        custom_watermark = f"""
---
{company_name if company_name and company_name.strip() else "© 2026 Existential Gateway, LLC | QuantusData.ai"}
Powered by QuantusData AI | Unauthorized reproduction prohibited.
---
""" if company_name and company_name.strip() else WATERMARK
    return report_text + "\n\n" + custom_watermark, output_path


def _make_pptx(text, df):
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    def add_slide(title_text, body_text, bg_color=(10, 10, 25)):
        slide = prs.slides.add_slide(blank_layout)
        bg = slide.background
        bg.fill.solid()
        bg.fill.fore_color.rgb = RGBColor(*bg_color)
        tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(1.2))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = RGBColor(231, 76, 60)
        body_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.7), Inches(12.3), Inches(5.5))
        tf2 = body_box.text_frame
        tf2.word_wrap = True
        for line in body_text.split("\n")[:20]:
            line = line.strip()
            if not line:
                continue
            para = tf2.add_paragraph()
            para.text = line
            para.font.size = Pt(14)
            para.font.color.rgb = RGBColor(220, 220, 240)

    # Title slide
    title_slide = prs.slides.add_slide(blank_layout)
    bg = title_slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(5, 5, 15)
    tb = title_slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11), Inches(1.5))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = "AI Law Enforcement & Crime Analysis Report"
    p.font.size = Pt(34)
    p.font.bold = True
    p.font.color.rgb = RGBColor(231, 76, 60)
    tb2 = title_slide.shapes.add_textbox(Inches(1), Inches(4.2), Inches(11), Inches(1))
    tf2 = tb2.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = f"Generated: {datetime.now().strftime('%B %d, %Y')} | AI Crime Analyzer"
    p2.font.size = Pt(18)
    p2.font.color.rgb = RGBColor(180, 180, 210)

    # Parse sections
    sections = []
    current_title = "Overview"
    current_body = []
    for line in text.split("\n"):
        stripped = line.strip()
        if stripped.startswith("**") and stripped.endswith("**") and len(stripped) > 4:
            if current_body:
                sections.append((current_title, "\n".join(current_body)))
            current_title = stripped.strip("*").strip()
            current_body = []
        elif stripped.startswith("## "):
            if current_body:
                sections.append((current_title, "\n".join(current_body)))
            current_title = stripped[3:].strip()
            current_body = []
        else:
            current_body.append(line)
    if current_body:
        sections.append((current_title, "\n".join(current_body)))

    bg_colors = [(10, 10, 25), (15, 10, 30), (5, 15, 25), (20, 10, 10), (10, 20, 10)]
    for i, (title, body) in enumerate(sections[:8]):
        add_slide(title, body, bg_colors[i % len(bg_colors)])

    if df is not None:
        stats_body = f"Total Records: {len(df)}\nColumns: {len(df.columns)}\n\nColumn Overview:\n"
        for col in df.columns[:10]:
            stats_body += f"  • {col} ({df[col].dtype})\n"
        add_slide("Data Statistics", stats_body)

    closing = prs.slides.add_slide(blank_layout)
    bg = closing.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(5, 5, 15)
    tb = closing.shapes.add_textbox(Inches(1), Inches(2.8), Inches(11), Inches(2))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Report Complete"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(231, 76, 60)
    p2 = tf.add_paragraph()
    p2.text = "© 2026 Existential Gateway, LLC | AI Law Enforcement and Crime Analyzer"
    p2.font.size = Pt(14)
    p2.font.color.rgb = RGBColor(160, 160, 200)

    tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".pptx")
    prs.save(tmp.name)
    return tmp.name


def _make_pdf(text):
    tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".pdf")
    doc = SimpleDocTemplate(tmp.name, pagesize=letter,
                            leftMargin=inch, rightMargin=inch,
                            topMargin=inch, bottomMargin=inch)
    styles = getSampleStyleSheet()
    title_style = ParagraphStyle("Title2", parent=styles["Title"],
                                 fontSize=20, spaceAfter=20,
                                 textColor=colors.HexColor("#8B0000"))
    heading_style = ParagraphStyle("H2", parent=styles["Heading2"],
                                   fontSize=14, spaceAfter=8,
                                   textColor=colors.HexColor("#A00000"))
    body_style = ParagraphStyle("Body2", parent=styles["Normal"],
                                fontSize=11, spaceAfter=6, leading=16)
    story = []
    story.append(Paragraph("AI Law Enforcement & Crime Analysis Report", title_style))
    story.append(Paragraph(
        f"Generated: {datetime.now().strftime('%B %d, %Y')} | AI Crime Analyzer",
        body_style))
    story.append(Spacer(1, 0.2 * inch))

    for line in text.split("\n"):
        stripped = line.strip()
        if not stripped:
            story.append(Spacer(1, 0.1 * inch))
            continue
        if stripped.startswith("**") and stripped.endswith("**"):
            story.append(Paragraph(stripped.strip("*"), heading_style))
        elif stripped.startswith("## "):
            story.append(Paragraph(stripped[3:], heading_style))
        elif stripped.startswith("# "):
            story.append(Paragraph(stripped[2:], title_style))
        else:
            safe = stripped.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")
            story.append(Paragraph(safe, body_style))

    story.append(Spacer(1, 0.3 * inch))
    story.append(Paragraph(
        "© 2026 Existential Gateway, LLC | AI Law Enforcement and Crime Analyzer | existentialgateway@gmail.com",
        ParagraphStyle("footer", parent=styles["Normal"], fontSize=9, textColor=colors.grey)))
    doc.build(story)
    return tmp.name


# ─── Tab 9: AI Data Chat ─────────────────────────────────────────────────────

_chat_df = None


def upload_chat_data(file):
    global _chat_df
    if file is None:
        _chat_df = None
        return "No file uploaded."
    df, err = load_data(file)
    if err and df is None:
        return err
    _chat_df = df
    return f"Data loaded: {len(df)} rows × {len(df.columns)} columns\nColumns: {list(df.columns)}"


def chat_with_data(message, history):
    global _chat_df

    data_context = ""
    if _chat_df is not None:
        desc = _chat_df.describe(include="all").to_string()
        preview = _chat_df.head(10).to_string()
        data_context = f"""
The user has uploaded crime/law enforcement data with the following properties:
Shape: {_chat_df.shape}
Columns: {list(_chat_df.columns)}
Statistical Summary:
{desc}
First 10 rows:
{preview}
"""
    else:
        data_context = "No data uploaded yet. Answer general law enforcement and crime analysis questions."

    messages = [
        {
            "role": "system",
            "content": (
                "You are an expert law enforcement data analyst and AI assistant. Write as the professional analyst who performed this work. State all findings directly with specific numbers, percentages, and comparisons. NEVER say 'this appears to be', 'the data seems to show', 'the dataset contains', or 'it looks like'. Instead, deliver results confidently: 'Average claim cost was $4,230, a 12% increase year-over-year' or 'Readmission rates of 18.3% in Cardiology exceeded the facility average by 4.1 percentage points.' Be direct, precise, and actionable. Present findings as YOUR analysis, not a description of a dataset. "
                "When asked a question, deliver the answer immediately with specific numbers. Do NOT explain how to calculate. Analyze their crime data, identify patterns, answer questions, "
                "provide specific computed answers with actual numbers from the data when asked. "
                "Be precise, professional, and helpful. Always remind users not to include "
                "personally identifiable information in their questions.\n\n" + data_context
            )
        }
    ]
    for user_msg, bot_msg in history:
        messages.append({"role": "user", "content": user_msg})
        messages.append({"role": "assistant", "content": bot_msg})
    messages.append({"role": "user", "content": message})

    import os
    API_KEY = os.environ.get("OPENAI_API_KEY", "")
    headers = {"Authorization": f"Bearer {API_KEY}", "Content-Type": "application/json"}
    msgs = [{"role": "system", "content": CRIME_SYSTEM_PROMPT}, {"role": "user", "content": prompt}]
    payload = {"model": "gpt-4o", "max_tokens": 4000, "messages": msgs}
    try:
        response = requests.post("https://api.openai.com/v1/chat/completions", headers=headers, json=payload, timeout=240)
        result = response.json()
        if "choices" in result:
            return result["choices"][0]["message"]["content"]
        return f"API Error: {result}"
    except Exception as e:
        return f"Error: {str(e)}"


# ─── Gradio UI ────────────────────────────────────────────────────────────────

with gr.Blocks(title="AI Law Enforcement & Crime Analyzer", theme=gr.themes.Base(
        primary_hue=gr.themes.colors.Color(
            c50="#fef9ec", c100="#faefd0", c200="#f4dea0", c300="#edc970",
            c400="#e4b04a", c500="#c9a84c", c600="#a8872e", c700="#856519",
            c800="#5e440d", c900="#3a2a05", c950="#1e1502"
        ),
        secondary_hue=gr.themes.colors.Color(
            c50="#e8edf5", c100="#c5d0e6", c200="#9eb0d4", c300="#7490c2",
            c400="#4f74b0", c500="#2d5a9e", c600="#1a3a6e", c700="#112240",
            c800="#0a1628", c900="#060e1a", c950="#03070d"
        ),
        neutral_hue=gr.themes.colors.Color(
            c50="#f0f2f7", c100="#d8dde9", c200="#b8c2d6", c300="#95a3be",
            c400="#7285a6", c500="#536890", c600="#3a4f73", c700="#263856",
            c800="#162438", c900="#0c1622", c950="#060b11"
        ),
        font=gr.themes.GoogleFont("DM Sans"),
        font_mono=gr.themes.GoogleFont("DM Mono"),
    ).set(
        body_background_fill="#0a1628",
        body_background_fill_dark="#0a1628",
        body_text_color="#f8f9fc",
        body_text_color_dark="#f8f9fc",
        block_background_fill="#112240",
        block_background_fill_dark="#112240",
        block_border_color="rgba(201,168,76,0.2)",
        block_border_color_dark="rgba(201,168,76,0.2)",
        block_label_background_fill="#112240",
        block_label_background_fill_dark="#112240",
        block_label_text_color="#c9a84c",
        block_label_text_color_dark="#c9a84c",
        block_title_text_color="#c9a84c",
        block_title_text_color_dark="#c9a84c",
        button_primary_background_fill="linear-gradient(135deg,#c9a84c,#e8c96a)",
        button_primary_background_fill_dark="linear-gradient(135deg,#c9a84c,#e8c96a)",
        button_primary_text_color="#0a1628",
        button_primary_text_color_dark="#0a1628",
        button_secondary_background_fill="#112240",
        button_secondary_background_fill_dark="#112240",
        button_secondary_border_color="rgba(201,168,76,0.3)",
        button_secondary_border_color_dark="rgba(201,168,76,0.3)",
        button_secondary_text_color="#c9a84c",
        button_secondary_text_color_dark="#c9a84c",
        input_background_fill="#0a1628",
        input_background_fill_dark="#0a1628",
        input_border_color="rgba(201,168,76,0.2)",
        input_border_color_dark="rgba(201,168,76,0.2)",
        input_placeholder_color="#8892a4",
        input_placeholder_color_dark="#8892a4",
        shadow_drop="0 4px 24px rgba(0,0,0,0.4)",
        shadow_drop_lg="0 8px 40px rgba(0,0,0,0.5)",
        table_even_background_fill="#0a1628",
        table_even_background_fill_dark="#0a1628",
        table_odd_background_fill="#112240",
        table_odd_background_fill_dark="#112240",
    )) as demo:
    gr.Markdown("# 🚔 AI Law Enforcement & Crime Analyzer")
    gr.Markdown(DISCLAIMER)

    with gr.Tabs():

        # ── Tab 1 ──────────────────────────────────────────────────────────────
        with gr.Tab("📊 Data Upload & Overview"):
            gr.Markdown("## Upload Your Crime Data")
            gr.Markdown(PII_WARNING)
            with gr.Row():
                with gr.Column(scale=1):
                    t1_file = gr.File(label="Upload File (CSV, Excel, PDF, JSON)",
                                      file_types=[".csv", ".xlsx", ".xls", ".pdf", ".json"])
                    t1_btn = gr.Button("🔍 Analyze Dataset", variant="primary")
                    gr.Markdown(WAIT_MSG)
                with gr.Column(scale=2):
                    t1_analysis = gr.Markdown(label="AI Analysis")
            with gr.Row():
                t1_stats = gr.Markdown(label="Statistics")
                t1_preview = gr.Markdown(label="Data Preview")
            t1_btn.click(analyze_overview, inputs=[t1_file],
                         outputs=[t1_analysis, t1_stats, t1_preview])

        # ── Tab 2 ──────────────────────────────────────────────────────────────
        with gr.Tab("🔍 Crime Pattern Analysis"):
            gr.Markdown("## Crime Pattern & Hotspot Analysis")
            gr.Markdown(PII_WARNING)
            with gr.Row():
                with gr.Column(scale=1):
                    t2_file = gr.File(label="Upload Crime Incident Data (CSV/Excel)",
                                      file_types=[".csv", ".xlsx", ".xls"])
                    t2_btn = gr.Button("🔍 Analyze Patterns", variant="primary")
                    gr.Markdown(WAIT_MSG)
                with gr.Column(scale=2):
                    t2_output = gr.Markdown(label="Pattern Analysis")
            t2_chart = gr.Plot(label="Crime Pattern Chart")
            t2_btn.click(analyze_patterns, inputs=[t2_file],
                         outputs=[t2_output, t2_chart])

        # ── Tab 3 ──────────────────────────────────────────────────────────────
        with gr.Tab("🔮 Predictive Analysis"):
            gr.Markdown("## Predictive Crime Analysis & Risk Scoring")
            gr.Markdown(PII_WARNING)
            with gr.Row():
                with gr.Column(scale=1):
                    t3_file = gr.File(label="Upload Historical Crime Data (CSV/Excel)",
                                      file_types=[".csv", ".xlsx", ".xls"])
                    t3_btn = gr.Button("🔮 Run Predictive Analysis", variant="primary")
                    gr.Markdown(WAIT_MSG)
                with gr.Column(scale=2):
                    t3_output = gr.Markdown(label="Predictive Analysis")
            t3_chart = gr.Plot(label="Trend Chart")
            t3_btn.click(analyze_predictive, inputs=[t3_file],
                         outputs=[t3_output, t3_chart])

        # ── Tab 4 ──────────────────────────────────────────────────────────────
        with gr.Tab("🔗 Case Correlation"):
            gr.Markdown("## Case Correlation & Serial Incident Analysis")
            gr.Markdown(PII_WARNING)
            with gr.Row():
                with gr.Column(scale=1):
                    t4_file = gr.File(label="Upload Case Files (CSV, Excel, PDF)",
                                      file_types=[".csv", ".xlsx", ".xls", ".pdf"])
                    t4_btn = gr.Button("🔗 Analyze Correlations", variant="primary")
                    gr.Markdown(WAIT_MSG)
                with gr.Column(scale=2):
                    t4_output = gr.Markdown(label="Case Correlation Report")
            t4_chart = gr.Plot(label="Correlation Matrix")
            t4_btn.click(analyze_correlation, inputs=[t4_file],
                         outputs=[t4_output, t4_chart])

        # ── Tab 5 ──────────────────────────────────────────────────────────────
        with gr.Tab("📋 Recidivism Analysis"):
            gr.Markdown("## Recidivism Risk & Rehabilitation Analysis")
            gr.Markdown(PII_WARNING)
            with gr.Row():
                with gr.Column(scale=1):
                    t5_file = gr.File(label="Upload Offender Data (CSV/Excel)",
                                      file_types=[".csv", ".xlsx", ".xls"])
                    t5_btn = gr.Button("📋 Analyze Recidivism", variant="primary")
                    gr.Markdown(WAIT_MSG)
                with gr.Column(scale=2):
                    t5_output = gr.Markdown(label="Recidivism Analysis")
            t5_chart = gr.Plot(label="Recidivism Chart")
            t5_btn.click(analyze_recidivism, inputs=[t5_file],
                         outputs=[t5_output, t5_chart])

        # ── Tab 6 ──────────────────────────────────────────────────────────────
        with gr.Tab("🚓 Resource Optimization"):
            gr.Markdown("## Patrol & Resource Optimization Analysis")
            gr.Markdown(PII_WARNING)
            with gr.Row():
                with gr.Column(scale=1):
                    t6_file = gr.File(label="Upload Staffing & Incident Data (CSV/Excel)",
                                      file_types=[".csv", ".xlsx", ".xls"])
                    t6_btn = gr.Button("🚓 Optimize Resources", variant="primary")
                    gr.Markdown(WAIT_MSG)
                with gr.Column(scale=2):
                    t6_output = gr.Markdown(label="Resource Optimization Report")
            t6_chart = gr.Plot(label="Resource Chart")
            t6_btn.click(analyze_resources, inputs=[t6_file],
                         outputs=[t6_output, t6_chart])

        # ── Tab 7 ──────────────────────────────────────────────────────────────
        with gr.Tab("📈 Visualizations & Maps"):
            gr.Markdown("## Interactive Crime Data Visualizations")
            gr.Markdown(PII_WARNING)
            with gr.Row():
                with gr.Column(scale=1):
                    t7_file = gr.File(label="Upload Data (CSV/Excel)",
                                      file_types=[".csv", ".xlsx", ".xls"])
                    t7_chart_type = gr.Dropdown(
                        choices=["Bar Chart", "Line Chart", "Pie Chart", "Scatter Plot",
                                 "Heatmap", "Box Plot", "Histogram", "Density Heatmap"],
                        value="Bar Chart", label="Chart Type")
                    t7_x = gr.Dropdown(choices=[], label="X Axis Column")
                    t7_y = gr.Dropdown(choices=[], label="Y Axis Column")
                    t7_color = gr.Dropdown(choices=[], label="Color By (optional)")
                    t7_btn = gr.Button("📊 Generate Chart", variant="primary")
                    t7_status = gr.Markdown()
                with gr.Column(scale=2):
                    t7_plot = gr.Plot(label="Chart")
            t7_file.change(get_columns, inputs=[t7_file],
                           outputs=[t7_x, t7_y, t7_color])
            t7_btn.click(generate_viz,
                         inputs=[t7_file, t7_chart_type, t7_x, t7_y, t7_color],
                         outputs=[t7_plot, t7_status])

        # ── Tab 8 ──────────────────────────────────────────────────────────────
        with gr.Tab("📄 Reports & Presentations"):
            gr.Markdown("## Generate Professional Crime Reports")
            gr.Markdown(PII_WARNING)
            with gr.Row():
                with gr.Column(scale=1):
                    t8_file = gr.File(label="Upload Data (CSV, Excel, PDF, JSON)",
                                      file_types=[".csv", ".xlsx", ".xls", ".pdf", ".json"])
                    t8_type = gr.Dropdown(
                        choices=["PowerPoint Presentation", "PDF Report", "Executive Briefing"],
                        value="PDF Report", label="Report Type")
                    t8_btn = gr.Button("📄 Generate Report", variant="primary")
                    gr.Markdown(WAIT_MSG)
                    t8_download = gr.File(label="⬇️ Download Report")
                with gr.Column(scale=2):
                    t8_output = gr.Markdown(label="Report Preview")
            t8_btn.click(generate_report, inputs=[t8_file, t8_type],
                         outputs=[t8_output, t8_download])

        # ── Tab 9 ──────────────────────────────────────────────────────────────
        with gr.Tab("💬 AI Data Chat"):
            gr.Markdown("## Chat With Your Crime Data")
            gr.Markdown(PII_WARNING)
            with gr.Row():
                with gr.Column(scale=1):
                    t9_file = gr.File(label="Upload Data for Context (CSV/Excel)",
                                      file_types=[".csv", ".xlsx", ".xls"])
                    t9_upload_btn = gr.Button("📤 Load Data", variant="secondary")
                    t9_data_status = gr.Markdown()
                with gr.Column(scale=2):
                    t9_chatbot = gr.ChatInterface(
                        fn=chat_with_data,
                        examples=[
                            "What are the most common crime types in this dataset?",
                            "Which location has the highest crime frequency?",
                            "What time of day do most crimes occur?",
                            "Generate a SQL query to find the top 10 highest crime areas",
                            "Write Python code to identify seasonal crime patterns",
                            "Are there any suspicious patterns or anomalies in this data?",
                            "What is the month with the highest crime rate?",
                        ],
                        title="",
                    )
            t9_upload_btn.click(upload_chat_data, inputs=[t9_file],
                                outputs=[t9_data_status])

    gr.Markdown(WATERMARK)


if __name__ == "__main__":
    demo.launch(server_name="0.0.0.0", server_port=int(os.environ.get("GRADIO_SERVER_PORT", 7860)), share=False, ssr_mode=False)
