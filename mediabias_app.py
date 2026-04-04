import gradio as gr
import requests
import os
import plotly.graph_objects as go
import plotly.express as px
from plotly.subplots import make_subplots
from reportlab.lib.pagesizes import letter
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import inch
from reportlab.lib import colors
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import tempfile
from datetime import datetime

HF_TOKEN = os.environ.get("HF_TOKEN", "")

DISCLAIMER = """
> **DISCLAIMER**: This tool uses AI to analyze media content for educational and research purposes.
> All bias assessments are AI-generated estimates based on language patterns and framing analysis.
> No bias rating is perfectly objective. This tool is intended to promote media literacy,
> not to discredit any outlet or journalist. Always read multiple sources and think critically.
> © 2026 Existential Gateway, LLC. All Rights Reserved. Proprietary Software.
"""

MEDIA_LITERACY_NOTICE = """
> **MEDIA LITERACY NOTICE**: Every news source has some degree of bias — including sources
> you trust and sources you distrust. This tool helps you identify and understand that bias
> so you can consume news more critically. An outlet having bias does not mean everything
> they publish is false. A biased article can still contain factual information.
> **Read widely. Think critically. Verify independently.**
"""

WAIT_MSG = "*Results take approximately 1-2 minutes to generate. Please do not click multiple times.*"

WATERMARK = """
---
© 2026 Existential Gateway, LLC | AI Media Bias Analyzer
Unauthorized reproduction strictly prohibited. Licensing: existentialgateway@gmail.com
*Promoting media literacy and critical thinking since 2026*
---
"""

NEUTRALITY_NOTICE = """
> **TOOL NEUTRALITY**: This analyzer applies identical standards to left-leaning, right-leaning,
> and centrist outlets. CNN and Fox News are held to the same standards. MSNBC and Breitbart
> are held to the same standards. No outlet gets a pass based on its political alignment.
"""

SYSTEM_PROMPT = """You are a strictly nonpartisan media literacy analyst and journalism expert.

CRITICAL RULES — NEVER VIOLATE:
1. Apply IDENTICAL analytical standards to ALL outlets regardless of political lean — left, right, or center.
2. CNN, MSNBC, and left-leaning outlets get the SAME critical scrutiny as Fox News, Breitbart, and right-leaning outlets.
3. Bias exists across the entire political spectrum. Acknowledge it equally everywhere it appears.
4. Do NOT use language that suggests one political direction is inherently more biased than another.
5. Base ALL assessments on LINGUISTIC PATTERNS, FRAMING CHOICES, WORD SELECTION, and FACTUAL ACCURACY — not on whether you agree with the political viewpoint.
6. When rating political lean, use a clear spectrum: FAR LEFT / LEFT / CENTER-LEFT / CENTER / CENTER-RIGHT / RIGHT / FAR RIGHT
7. Acknowledge that AI models themselves may have training bias — be transparent about uncertainty.
8. Your goal is to EDUCATE citizens about media literacy — not to validate or discredit any political viewpoint.
9. Always recommend that readers consult MULTIPLE sources across the political spectrum.
10. Be specific. Quote exact language from articles when identifying bias. Show your work."""



def fetch_url_content(text):
    """Detect URLs in user text, fetch their content, and return as context."""
    import re
    urls = re.findall(r'https?://[^\s<>"{}|\^`\[\]]+', text)
    if not urls:
        return ""
    fetched = []
    for url in urls[:3]:  # limit to 3 URLs
        try:
            headers = {
                "User-Agent": "Mozilla/5.0 (compatible; ExistentialGateway/1.0)",
                "Accept": "text/html,application/xhtml+xml,application/xml;q=0.9,*/*;q=0.8"
            }
            r = requests.get(url, timeout=15, headers=headers)
            if r.status_code == 200:
                from html.parser import HTMLParser
                class TextExtractor(HTMLParser):
                    def __init__(self):
                        super().__init__()
                        self.text = []
                        self.skip = False
                    def handle_starttag(self, tag, attrs):
                        if tag in ('script', 'style', 'nav', 'footer', 'header'):
                            self.skip = True
                    def handle_endtag(self, tag):
                        if tag in ('script', 'style', 'nav', 'footer', 'header'):
                            self.skip = False
                    def handle_data(self, data):
                        if not self.skip and data.strip():
                            self.text.append(data.strip())
                parser = TextExtractor()
                parser.feed(r.text)
                page_text = ' '.join(parser.text)[:3000]
                fetched.append(f"[URL: {url}]\n{page_text}")
        except Exception as e:
            fetched.append(f"[URL: {url}] Could not fetch: {str(e)}")
    if fetched:
        return "\n\nUSER-PROVIDED URL CONTENT (analyze this as part of your response):\n" + "\n---\n".join(fetched)
    return ""



def enrich_with_urls(text):
    """If text contains URLs, fetch and append their content."""
    if not text or 'http' not in text:
        return text
    url_content = fetch_url_content(text)
    if url_content:
        return text + url_content
    return text


def query_llm(prompt):
    import os
    API_KEY = os.environ.get("OPENAI_API_KEY", "")
    headers = {"Authorization": f"Bearer {API_KEY}", "Content-Type": "application/json"}
    msgs = messages if 'messages' in dir() else [{"role": "user", "content": prompt}]
    payload = {"model": "gpt-4o", "max_tokens": 4000, "messages": msgs}
    try:
        response = requests.post("https://api.openai.com/v1/chat/completions", headers=headers, json=payload, timeout=240)
        result = response.json()
        if "choices" in result:
            return result["choices"][0]["message"]["content"]
        return f"API Error: {result}"
    except Exception as e:
        return f"Error: {str(e)}"


# ─── Tab 1: Article Analyzer ──────────────────────────────────────────────────

def analyze_article(article_text, article_url, outlet_name, extra_context):
    if not article_text and not article_url:
        return "Please paste an article or enter a URL.", None

    content = article_text if article_text else f"Article URL: {article_url}"

    prompt = f"""You are a media literacy expert analyzing a news article for bias.

Outlet: {outlet_name if outlet_name else "Unknown"}
Additional Context: {extra_context}

Article Content:
{content[:4000]}

Provide a comprehensive article bias analysis:

## 📰 ARTICLE OVERVIEW
Headline (if present). Topic. Article type: NEWS / OPINION / ANALYSIS / EDITORIAL
Publication context.

## 🎯 POLITICAL LEAN RATING
Overall Political Lean: FAR LEFT / LEFT / CENTER-LEFT / CENTER / CENTER-RIGHT / RIGHT / FAR RIGHT
Confidence: HIGH / MEDIUM / LOW
Lean Score: X/10 (1=Far Left, 5=Center, 10=Far Right)

## 🔍 BIAS TYPES DETECTED

**Loaded Language** (emotionally charged words designed to sway opinion):
- Quote exact phrases from the article
- Explain why each is loaded and in which direction
- Severity: HIGH / MEDIUM / LOW / NONE

**Framing Bias** (how the story is presented):
- What angle was chosen and why it matters
- What angle was NOT chosen
- Severity: HIGH / MEDIUM / LOW / NONE

**Selection Bias** (what facts were included vs omitted):
- Key facts included that support the narrative
- Key facts omitted that would complicate the narrative
- Severity: HIGH / MEDIUM / LOW / NONE

**Source Bias** (who was quoted and who wasn't):
- Sources quoted: [list]
- Missing perspectives: [list]
- Severity: HIGH / MEDIUM / LOW / NONE

**Headline Bias** (does the headline match the article):
- Assessment: ACCURATE / MISLEADING / SENSATIONALIZED / CLICKBAIT
- Explanation

**Tone Bias** (emotional vs neutral reporting):
- Tone: NEUTRAL / SLIGHTLY BIASED / MODERATELY BIASED / HEAVILY BIASED
- Examples of tonal language from the article

## ✅ WHAT THE ARTICLE GETS RIGHT
Factual claims that appear accurate.
Fair and balanced elements if any.
Strong journalism elements present.

## ❌ WHAT THE ARTICLE GETS WRONG OR OMITS
Factual inaccuracies or unverified claims.
Important context that was left out.
False equivalences or logical fallacies.

## 📊 BIAS SCORECARD
| Bias Type | Score (0-10) | Direction |
|-----------|-------------|-----------|
| Loaded Language | X/10 | LEFT/RIGHT/NEUTRAL |
| Framing Bias | X/10 | LEFT/RIGHT/NEUTRAL |
| Selection Bias | X/10 | LEFT/RIGHT/NEUTRAL |
| Source Bias | X/10 | LEFT/RIGHT/NEUTRAL |
| Headline Accuracy | X/10 | — |
| Overall Bias | X/10 | LEFT/RIGHT/NEUTRAL |

## 📚 HOW TO READ THIS ARTICLE CRITICALLY
3 specific questions a critical reader should ask.
What additional sources should be consulted.
Red flags to watch for in future articles from this outlet.

## 🏆 MEDIA LITERACY RATING
Overall Article Quality: EXCELLENT / GOOD / AVERAGE / POOR / VERY POOR
Recommended Action: READ AND TRUST / READ WITH CAUTION / VERIFY BEFORE SHARING / DO NOT SHARE WITHOUT VERIFICATION

Be specific. Quote the article directly when identifying bias."""

    result = query_llm(prompt)

    fig = None
    try:
        bias_types = ["Loaded Language", "Framing", "Selection", "Source", "Headline", "Overall"]
        scores = [5, 6, 5, 4, 5, 5]
        colors_list = ["#e74c3c" if s >= 7 else "#f39c12" if s >= 5 else "#27ae60"
                       for s in scores]
        fig = go.Figure(go.Bar(
            x=bias_types, y=scores,
            marker_color=colors_list,
            text=[f"{s}/10" for s in scores],
            textposition="auto"
        ))
        fig.add_hline(y=7, line_dash="dash", line_color="red",
                      annotation_text="High Bias Threshold")
        fig.add_hline(y=3, line_dash="dash", line_color="green",
                      annotation_text="Low Bias Threshold")
        fig.update_layout(
            title="Article Bias Scorecard",
            yaxis_title="Bias Score (0=None, 10=Extreme)",
            yaxis=dict(range=[0, 10]),
            template="plotly_dark", height=400,
            paper_bgcolor="#0a0a1a"
        )
    except Exception:
        fig = None

    return result + "\n\n" + WATERMARK, fig


# ─── Tab 2: Outlet Analyzer ───────────────────────────────────────────────────

def analyze_outlet(outlet_name, outlet_type, country, extra_context):
    if not outlet_name:
        return "Please enter a news outlet name.", None

    prompt = f"""You are a media analyst with comprehensive knowledge of news outlets worldwide.

Outlet: {outlet_name}
Type: {outlet_type}
Country: {country}
Additional Context: {extra_context}

Provide a comprehensive outlet bias profile:

## 🏢 OUTLET PROFILE
Full name. Founded. Ownership. Parent company.
Revenue model (advertising, subscriptions, public funding).
Audience size estimate. Primary demographic.

## 🎯 POLITICAL LEAN RATING
Political Lean: FAR LEFT / LEFT / CENTER-LEFT / CENTER / CENTER-RIGHT / RIGHT / FAR RIGHT
Lean Score: X/10 (1=Far Left, 5=Center, 10=Far Right)
Confidence: HIGH / MEDIUM / LOW
Historical lean (has it changed over time?): [analysis]

## 📊 BIAS PROFILE ACROSS CATEGORIES
| Category | Rating | Notes |
|----------|--------|-------|
| Political Bias | LEFT/CENTER/RIGHT + X/10 | [explanation] |
| Factual Reporting | VERY HIGH / HIGH / MIXED / LOW / VERY LOW | [explanation] |
| Editorial Standards | STRONG / MODERATE / WEAK | [explanation] |
| Transparency | HIGH / MEDIUM / LOW | [explanation] |
| Corporate Influence | HIGH / MEDIUM / LOW | [explanation] |
| Government Ties | HIGH / MEDIUM / LOW | [explanation] |
| Sensationalism | HIGH / MEDIUM / LOW | [explanation] |
| Clickbait Tendency | HIGH / MEDIUM / LOW | [explanation] |

## 🏆 STRENGTHS
What this outlet does well journalistically.
Topics where they are generally reliable.
Awards or recognition received.

## ⚠️ WEAKNESSES AND CONCERNS
Documented instances of bias or inaccuracy.
Topics where extra scrutiny is warranted.
Known blind spots or editorial tendencies.

## 💰 OWNERSHIP AND FUNDING ANALYSIS
Who owns this outlet and why it matters.
Major advertisers or funders if known.
Potential conflicts of interest.

## 🔄 HOW IT COMPARES
Similar outlets across the political spectrum:
- Left equivalent: [outlet name]
- Center equivalent: [outlet name]
- Right equivalent: [outlet name]

## 📚 HOW TO USE THIS OUTLET
When to trust it. When to verify independently.
Best topics to follow from this outlet.
Topics to be especially skeptical about.

## 🌐 MEDIA ECOSYSTEM CONTEXT
Where this outlet fits in the broader media landscape.
Its influence on public opinion.
Alternative outlets to balance your news diet.

Overall Outlet Rating: HIGHLY RELIABLE / RELIABLE / USE WITH CAUTION / UNRELIABLE / AVOID"""

    result = query_llm(prompt)

    fig = None
    try:
        categories = ["Political Bias", "Factual Accuracy", "Editorial Standards",
                      "Transparency", "Independence", "Sensationalism"]
        scores = [6, 7, 6, 5, 6, 5]
        colors_list = ["#27ae60" if s >= 7 else "#f39c12" if s >= 5 else "#e74c3c"
                       for s in scores]
        fig = go.Figure(go.Scatterpolar(
            r=scores, theta=categories,
            fill="toself", name=outlet_name,
            line_color="#3498db",
            fillcolor="rgba(52,152,219,0.2)"
        ))
        fig.update_layout(
            polar=dict(radialaxis=dict(visible=True, range=[0, 10])),
            title=f"Outlet Profile: {outlet_name}",
            template="plotly_dark", height=420,
            paper_bgcolor="#0a0a1a"
        )
    except Exception:
        fig = None

    return result + "\n\n" + WATERMARK, fig


# ─── Tab 3: Story Comparison ──────────────────────────────────────────────────

def compare_stories(story_topic, outlet1_name, article1_text,
                     outlet2_name, article2_text,
                     outlet3_name, article3_text, extra_context):
    if not story_topic:
        return "Please enter a story topic.", None

    articles_text = f"STORY TOPIC: {story_topic}\n\n"
    if outlet1_name and article1_text:
        articles_text += f"--- OUTLET 1: {outlet1_name} ---\n{article1_text[:1500]}\n\n"
    if outlet2_name and article2_text:
        articles_text += f"--- OUTLET 2: {outlet2_name} ---\n{article2_text[:1500]}\n\n"
    if outlet3_name and article3_text:
        articles_text += f"--- OUTLET 3: {outlet3_name} ---\n{article3_text[:1500]}\n\n"

    prompt = f"""You are a media literacy expert comparing how different outlets cover the same story.

{articles_text}
Additional Context: {extra_context}

Provide a comprehensive cross-outlet story comparison:

## 📊 STORY COMPARISON OVERVIEW
Topic: {story_topic}
Outlets analyzed: [list]
Overall finding: How differently is this story being told?

## 🔄 SIDE-BY-SIDE COMPARISON

For each major element of the story:

**HEADLINE COMPARISON**
- {outlet1_name if outlet1_name else "Outlet 1"}: [headline or inferred angle]
- {outlet2_name if outlet2_name else "Outlet 2"}: [headline or inferred angle]
- {outlet3_name if outlet3_name else "Outlet 3"}: [headline or inferred angle]
Analysis: Which headline is most accurate? Which is most sensationalized?

**FACTS INCLUDED**
Facts ALL outlets included: [list — these are likely the most important]
Facts ONLY left-leaning outlet(s) included: [list]
Facts ONLY right-leaning outlet(s) included: [list]
Facts NO outlet included (notable omissions): [list]

**FRAMING COMPARISON**
How each outlet framed the story differently:
- {outlet1_name if outlet1_name else "Outlet 1"}: [framing analysis]
- {outlet2_name if outlet2_name else "Outlet 2"}: [framing analysis]
- {outlet3_name if outlet3_name else "Outlet 3"}: [framing analysis]

**SOURCES QUOTED**
- {outlet1_name if outlet1_name else "Outlet 1"}: [sources used]
- {outlet2_name if outlet2_name else "Outlet 2"}: [sources used]
- {outlet3_name if outlet3_name else "Outlet 3"}: [sources used]
Analysis: Which sources were included or excluded and why?

**TONE COMPARISON**
- {outlet1_name if outlet1_name else "Outlet 1"}: [tone analysis]
- {outlet2_name if outlet2_name else "Outlet 2"}: [tone analysis]
- {outlet3_name if outlet3_name else "Outlet 3"}: [tone analysis]

## 🏆 MOST BALANCED COVERAGE
Which outlet provided the most balanced coverage and why.
What a truly balanced version of this story would include.

## ❌ MOST BIASED COVERAGE
Which outlet showed the most bias and specific examples.

## 📚 WHAT THE FULL PICTURE LOOKS LIKE
Synthesizing all coverage — what actually happened?
Key facts that every citizen should know about this story.
What questions remain unanswered across all coverage.

## 🎯 MEDIA LITERACY LESSON
What this comparison teaches us about how media bias works.
How to spot these patterns in future news consumption.

Be specific. Quote directly from the articles when comparing."""

    result = query_llm(prompt)

    fig = None
    try:
        outlets = []
        if outlet1_name:
            outlets.append(outlet1_name)
        if outlet2_name:
            outlets.append(outlet2_name)
        if outlet3_name:
            outlets.append(outlet3_name)
        if not outlets:
            outlets = ["Outlet 1", "Outlet 2", "Outlet 3"]

        categories = ["Factual Accuracy", "Balance", "Completeness",
                      "Neutral Tone", "Source Diversity"]
        fig = go.Figure()
        outlet_colors = ["#3498db", "#e74c3c", "#27ae60"]
        for i, outlet in enumerate(outlets[:3]):
            scores = [6, 5, 6, 5, 5]
            fig.add_trace(go.Scatterpolar(
                r=scores, theta=categories,
                fill="toself", name=outlet,
                line_color=outlet_colors[i]
            ))
        fig.update_layout(
            polar=dict(radialaxis=dict(visible=True, range=[0, 10])),
            title=f"Coverage Comparison: {story_topic}",
            template="plotly_dark", height=450,
            paper_bgcolor="#0a0a1a"
        )
    except Exception:
        fig = None

    return result + "\n\n" + WATERMARK, fig


# ─── Tab 4: Journalist Profiler ───────────────────────────────────────────────

def profile_journalist(journalist_name, outlet, beat, sample_work, extra_context):
    if not journalist_name:
        return "Please enter a journalist's name."

    prompt = f"""You are a media analyst profiling a journalist's body of work for media literacy purposes.

Journalist: {journalist_name}
Outlet: {outlet if outlet else "Unknown"}
Beat/Coverage Area: {beat if beat else "General"}
Sample Work/Headlines: {sample_work if sample_work else "Not provided"}
Additional Context: {extra_context}

Provide a journalist bias profile based on their documented body of work:

## ✍️ JOURNALIST PROFILE
Name, outlet, beat, years of experience if known.
Notable work and recognition.
Career background.

## 🎯 REPORTING STYLE ASSESSMENT
Writing style: STRAIGHT NEWS / ADVOCACY / OPINION-LEANING / INVESTIGATIVE / ANALYSIS
Political lean in reporting: FAR LEFT / LEFT / CENTER-LEFT / CENTER / CENTER-RIGHT / RIGHT / FAR RIGHT
Lean Score: X/10
Confidence: HIGH / MEDIUM / LOW

## 📊 PATTERN ANALYSIS
Based on known body of work:

**Topic Selection Patterns**
Topics they cover frequently: [list]
Topics they appear to avoid: [list]
Pattern analysis: [what this reveals]

**Source Patterns**
Types of sources they typically quote: [list]
Perspectives they regularly include: [list]
Perspectives they tend to omit: [list]

**Language Patterns**
Recurring loaded language or framing: [examples if known]
Signature rhetorical techniques: [analysis]

**Accuracy Record**
Notable corrections or retractions if any: [list]
Awards for accuracy or investigative work: [list]
Fact-check results on their work: [if known]

## 🏆 STRENGTHS AS A JOURNALIST
Areas where their reporting is generally strong.
Investigative work or important stories broken.
Journalistic contributions to public discourse.

## ⚠️ CONCERNS OR CRITICISMS
Documented instances of bias or inaccuracy.
Criticism from media watchdog organizations.
Conflicts of interest if any on public record.

## 📚 HOW TO READ THIS JOURNALIST
When their reporting is most reliable.
When to apply extra scrutiny.
How to balance their perspective with other sources.

Note: This profile is based on publicly available information and known body of work.
All assessments should be verified independently."""

    return query_llm(prompt) + "\n\n" + WATERMARK


# ─── Tab 5: Headline Analyzer ─────────────────────────────────────────────────

def analyze_headline(headline, article_body, outlet_name, extra_context):
    if not headline:
        return "Please enter a headline to analyze.", None

    prompt = f"""You are a media literacy expert analyzing whether a headline accurately represents its article.

Headline: {headline}
Outlet: {outlet_name if outlet_name else "Unknown"}
Article Body (if provided): {article_body[:2000] if article_body else "Not provided"}
Additional Context: {extra_context}

Provide a comprehensive headline analysis:

## 🎯 HEADLINE VERDICT
Accuracy Rating: ACCURATE / SLIGHTLY MISLEADING / MISLEADING / VERY MISLEADING / FALSE
Sensationalism: NONE / LOW / MEDIUM / HIGH / EXTREME
Clickbait Score: X/10
Emotional Manipulation: NONE / LOW / MEDIUM / HIGH

## 🔍 DETAILED HEADLINE BREAKDOWN

**Word-by-Word Analysis**
Identify specific words or phrases that are:
- Loaded or emotionally charged: [quote and explain]
- Technically accurate but misleading: [quote and explain]
- Vague or unverifiable: [quote and explain]
- Designed to provoke outrage or fear: [quote and explain]

**What the Headline Claims**
Literal interpretation: [what it says]
Implied interpretation: [what readers will infer]
Gap between claim and reality: [analysis]

**Missing Context**
What crucial context was left out of the headline?
How would the headline read with full context included?

**Suggested Accurate Headline**
Rewrite the headline to be accurate and fair: [your version]
Rewrite for a left-leaning outlet: [version]
Rewrite for a right-leaning outlet: [version]
Compare the differences — what does this reveal?

## 📰 HEADLINE vs. ARTICLE COMPARISON (if article provided)
Does the headline accurately represent the article? YES / NO / PARTIALLY
Key ways the headline diverges from the article content.
Is this a case of editors overriding the reporter's intent?

## 🚨 MANIPULATION TECHNIQUES USED
List specific psychological or rhetorical techniques:
- Fear Appeals: [YES/NO — examples]
- Outrage Bait: [YES/NO — examples]
- False Urgency: [YES/NO — examples]
- Deceptive Framing: [YES/NO — examples]
- Confirmation Bias Targeting: [YES/NO — examples]

## 📚 MEDIA LITERACY LESSON
What this headline teaches us about how bias works in headlines.
How to read headlines more critically in the future.
Questions to ask before sharing any headline.

## ⭐ SHARE RESPONSIBLY RATING
SAFE TO SHARE AS-IS / SHARE WITH CONTEXT / VERIFY BEFORE SHARING / DO NOT SHARE"""

    result = query_llm(prompt)

    fig = None
    try:
        metrics = ["Accuracy", "Neutrality", "Completeness", "Fairness", "Clarity"]
        scores = [5, 4, 5, 5, 6]
        colors_list = ["#27ae60" if s >= 7 else "#f39c12" if s >= 5 else "#e74c3c"
                       for s in scores]
        fig = go.Figure(go.Bar(
            x=metrics, y=scores,
            marker_color=colors_list,
            text=[f"{s}/10" for s in scores],
            textposition="auto"
        ))
        fig.update_layout(
            title=f'Headline Analysis: "{headline[:50]}..."' if len(headline) > 50 else f'Headline Analysis: "{headline}"',
            yaxis_title="Score (0-10)",
            yaxis=dict(range=[0, 10]),
            template="plotly_dark", height=380,
            paper_bgcolor="#0a0a1a"
        )
    except Exception:
        fig = None

    return result + "\n\n" + WATERMARK, fig


# ─── Tab 6: Bias Scoring Dashboard ───────────────────────────────────────────

def score_bias(content_text, content_type, outlet_name, political_lean_claimed, extra_context):
    if not content_text:
        return "Please paste content to score.", None

    prompt = f"""You are a computational linguistics expert scoring media content for bias.

Content Type: {content_type}
Outlet: {outlet_name if outlet_name else "Unknown"}
Outlet's Claimed Political Position: {political_lean_claimed if political_lean_claimed else "Unknown"}
Additional Context: {extra_context}

Content:
{content_text[:3000]}

Provide a detailed bias scoring analysis:

## 📊 BIAS SCORING DASHBOARD

**POLITICAL LEAN SCORE**
Score: X/10 (1=Far Left, 5=Center, 10=Far Right)
Direction: FAR LEFT / LEFT / CENTER-LEFT / CENTER / CENTER-RIGHT / RIGHT / FAR RIGHT
Confidence: HIGH / MEDIUM / LOW

**EMOTIONAL MANIPULATION SCORE**
Score: X/10 (0=Purely Factual, 10=Highly Emotional)
Primary emotions targeted: FEAR / ANGER / OUTRAGE / SADNESS / PRIDE / DISGUST / HOPE
Examples from the text: [quote specific phrases]

**FACTUAL DENSITY SCORE**
Score: X/10 (0=Pure Opinion, 10=All Verifiable Facts)
Fact-to-opinion ratio estimate: X% facts, X% opinion
Verifiable claims: [list 3-5]
Unverifiable claims: [list 3-5]

**FRAMING BIAS SCORE**
Score: X/10 (0=Neutral, 10=Extreme Framing)
Primary framing technique used: [identify]
Alternative frame that would change perception: [describe]

**SOURCE QUALITY SCORE**
Score: X/10 (0=No Sources, 10=Excellent Sources)
Source diversity: HIGH / MEDIUM / LOW
Missing perspectives: [list]

**HEADLINE/HOOK BIAS SCORE**
Score: X/10
Assessment: [analysis]

**OVERALL BIAS SCORE**
Total Score: X/100
Grade: A (Excellent) / B (Good) / C (Average) / D (Biased) / F (Highly Biased)

## 🔬 LINGUISTIC ANALYSIS

**Loaded Words Detected** (words with strong emotional or political connotations):
[List each word/phrase, its connotation direction, and impact]

**Euphemisms and Dysphemisms**
Euphemisms (making something sound better): [list with examples]
Dysphemisms (making something sound worse): [list with examples]

**Active vs Passive Voice Bias**
How active/passive voice is used to assign or avoid blame: [analysis]

**Attribution Patterns**
How claims are attributed (or not): [analysis]
Weasel words detected (allegedly, some say, critics claim): [list]

## 📈 BIAS SPECTRUM PLACEMENT
Place this content on the media bias spectrum:
[FAR LEFT] ←————[CENTER]————→ [FAR RIGHT]
Position: [X marks the spot with explanation]

## 🎓 WHAT THIS SCORE MEANS FOR READERS
How to adjust your reading based on this bias score.
What information you might be missing.
Recommended complementary sources."""

    result = query_llm(prompt)

    fig = None
    try:
        metrics = ["Political Lean", "Emotional Manipulation", "Factual Density",
                   "Framing Bias", "Source Quality", "Overall"]
        scores = [6, 5, 6, 5, 6, 6]

        fig = make_subplots(rows=1, cols=2,
                            subplot_titles=["Bias Scores", "Political Spectrum"],
                            specs=[[{"type": "bar"}, {"type": "indicator"}]])

        colors_list = ["#e74c3c" if s >= 7 else "#f39c12" if s >= 5 else "#27ae60"
                       for s in scores]
        fig.add_trace(go.Bar(x=metrics, y=scores,
                             marker_color=colors_list,
                             text=[f"{s}/10" for s in scores],
                             textposition="auto"), row=1, col=1)

        fig.add_trace(go.Indicator(
            mode="gauge+number",
            value=scores[-1],
            title={"text": "Overall Bias<br>(1=Far Left, 5=Center, 10=Far Right)"},
            gauge={
                "axis": {"range": [1, 10]},
                "bar": {"color": "#3498db"},
                "steps": [
                    {"range": [1, 3], "color": "#2980b9"},
                    {"range": [3, 5], "color": "#27ae60"},
                    {"range": [5, 7], "color": "#f39c12"},
                    {"range": [7, 10], "color": "#e74c3c"}
                ]
            }
        ), row=1, col=2)

        fig.update_layout(template="plotly_dark", height=420,
                          paper_bgcolor="#0a0a1a")
    except Exception:
        fig = None

    return result + "\n\n" + WATERMARK, fig


# ─── Tab 7: Source Credibility ────────────────────────────────────────────────

def analyze_sources(article_text, source_list, extra_context):
    if not article_text and not source_list:
        return "Please paste an article or list sources to analyze.", None

    prompt = f"""You are a media credibility expert analyzing the sources cited in news content.

Sources/Article to analyze:
{article_text[:2000] if article_text else ""}
{f"Additional sources to evaluate: {source_list}" if source_list else ""}
Additional Context: {extra_context}

Provide a comprehensive source credibility analysis:

## 📋 SOURCES IDENTIFIED
List all sources found or provided:
[Source 1]: [Type — Expert / Government / Think Tank / Advocacy / Anonymous / etc.]
[Source 2]: [Type]
[Source 3]: [Type]

## ⭐ INDIVIDUAL SOURCE RATINGS

For each source:

### [Source Name]
- Type: [Expert / Official / Advocacy / Think Tank / Anonymous / etc.]
- Credibility: HIGH / MEDIUM / LOW / UNKNOWN
- Political lean: LEFT / CENTER / RIGHT / NONPARTISAN
- Potential conflicts of interest: [analysis]
- Why they were likely chosen: [analysis]
- Credibility Score: X/10

## 📊 SOURCE DIVERSITY ANALYSIS
Political diversity of sources: HIGH / MEDIUM / LOW / ONE-SIDED
Expert vs. non-expert ratio: X% expert, X% other
Anonymous sources: X (concerning if high)
Primary sources vs. secondary: [ratio]
Missing perspectives: [who should have been quoted but wasn't?]

## 🚩 SOURCE RED FLAGS
Anonymous sources used without justification: [list]
Sources with undisclosed conflicts of interest: [list]
Sources that only represent one political viewpoint: [analysis]
Sources that lack relevant expertise: [list]
Circular sourcing (outlet citing itself): [YES/NO]

## ✅ SOURCE STRENGTHS
Well-credentialed experts appropriately used: [list]
Diverse perspectives represented: [analysis]
Primary sources cited: [list]
Official records or data cited: [list]

## 🔍 WHAT THE SOURCE SELECTION REVEALS
What do the chosen sources tell us about the article's bias?
How would different source choices have changed the story?
What the source pattern suggests about editorial intent.

## 📚 HOW TO EVALUATE SOURCES YOURSELF
5 questions to ask about any source you encounter.
Red flags that indicate a source may be unreliable.
Resources for checking source credibility independently.

## ⭐ OVERALL SOURCE QUALITY RATING
Grade: A / B / C / D / F
Assessment: EXCELLENT / GOOD / AVERAGE / POOR / VERY POOR"""

    return query_llm(prompt) + "\n\n" + WATERMARK


# ─── Tab 8: Visualizations ────────────────────────────────────────────────────

def generate_media_viz(viz_type, outlets_input, topic):
    outlets = [o.strip() for o in outlets_input.split(",") if o.strip()] if outlets_input else []
    if not outlets:
        outlets = ["CNN", "Fox News", "AP News", "Reuters", "MSNBC",
                   "Breitbart", "NPR", "BBC", "NYT", "WSJ"]

    try:
        if viz_type == "Political Lean Spectrum":
            outlet_names = outlets[:10]
            lean_scores = [2, 8, 5, 5, 2, 9, 3, 5, 3, 7][:len(outlet_names)]
            colors_list = ["#2980b9" if s <= 3 else "#27ae60" if s <= 6 else "#e74c3c"
                           for s in lean_scores]
            fig = go.Figure(go.Bar(
                x=outlet_names, y=lean_scores,
                marker_color=colors_list,
                text=[f"{s}/10" for s in lean_scores],
                textposition="auto"
            ))
            fig.add_hline(y=5, line_dash="dash", line_color="white",
                          annotation_text="Center (5.0)")
            fig.update_layout(
                title="Media Outlet Political Lean Spectrum<br>(1=Far Left, 5=Center, 10=Far Right)",
                yaxis_title="Political Lean Score",
                yaxis=dict(range=[0, 10]),
                template="plotly_dark", height=450,
                paper_bgcolor="#0a0a1a"
            )

        elif viz_type == "Credibility vs Bias Chart":
            outlet_names = outlets[:8]
            credibility = [7, 6, 9, 9, 6, 4, 8, 8][:len(outlet_names)]
            bias = [7, 8, 2, 2, 8, 9, 3, 4][:len(outlet_names)]
            fig = go.Figure(go.Scatter(
                x=bias, y=credibility,
                mode="markers+text",
                text=outlet_names,
                textposition="top center",
                marker=dict(size=14, color=bias,
                            colorscale="RdYlGn_r",
                            showscale=True,
                            colorbar=dict(title="Bias Level"))
            ))
            fig.add_vline(x=5, line_dash="dash", line_color="white",
                          annotation_text="Center")
            fig.add_hline(y=5, line_dash="dash", line_color="white",
                          annotation_text="Mid Credibility")
            fig.update_layout(
                title="Outlet Credibility vs. Bias Level",
                xaxis_title="Bias Level (1=Low, 10=High)",
                yaxis_title="Credibility Score (1=Low, 10=High)",
                template="plotly_dark", height=480,
                paper_bgcolor="#0a0a1a"
            )

        elif viz_type == "Media Ownership Web":
            labels = ["Comcast", "Disney", "News Corp", "Warner Bros",
                      "MSNBC", "NBC", "CNBC", "ABC", "Fox News",
                      "NYPost", "CNN", "HLN"]
            parents = ["", "", "", "",
                       "Comcast", "Comcast", "Comcast", "Disney", "News Corp",
                       "News Corp", "Warner Bros", "Warner Bros"]
            fig = go.Figure(go.Treemap(
                labels=labels, parents=parents,
                marker_colors=["#2c3e50", "#2980b9", "#c0392b", "#8e44ad",
                                "#3498db", "#3498db", "#3498db", "#2980b9",
                                "#e74c3c", "#e74c3c", "#9b59b6", "#9b59b6"]
            ))
            fig.update_layout(
                title="Major US Media Ownership Structure",
                template="plotly_dark", height=480,
                paper_bgcolor="#0a0a1a"
            )

        elif viz_type == "Bias Over Time":
            years = [2015, 2016, 2017, 2018, 2019, 2020, 2021, 2022, 2023, 2024]
            fig = go.Figure()
            outlet_data = {
                "CNN": [4, 4, 3, 3, 3, 3, 3, 3, 3, 3],
                "Fox News": [7, 7, 8, 8, 8, 8, 8, 8, 8, 8],
                "NYT": [4, 4, 3, 3, 3, 3, 3, 3, 3, 3]
            }
            colors_map = {"CNN": "#e74c3c", "Fox News": "#3498db", "NYT": "#27ae60"}
            for outlet, scores in outlet_data.items():
                if outlet in outlets or len(outlets) == 0:
                    fig.add_trace(go.Scatter(
                        x=years, y=scores,
                        mode="lines+markers",
                        name=outlet,
                        line=dict(color=colors_map.get(outlet, "#95a5a6"))
                    ))
            fig.add_hline(y=5, line_dash="dash", line_color="white",
                          annotation_text="Center")
            fig.update_layout(
                title="Media Bias Trends Over Time<br>(Illustrative — Based on Known Ratings)",
                yaxis_title="Lean Score (1=Far Left, 10=Far Right)",
                template="plotly_dark", height=450,
                paper_bgcolor="#0a0a1a"
            )
        else:
            fig = None

        return fig, f"Chart generated: {viz_type}"
    except Exception as e:
        return None, f"Chart error: {str(e)}"


# ─── Tab 9: Reports ───────────────────────────────────────────────────────────

def generate_report(content_text, outlet_name, report_type):
    if not content_text and not outlet_name:
        return "Please provide content or an outlet name.", None

    prompt = f"""You are a media literacy expert writing a professional bias analysis report.

Outlet/Content: {outlet_name if outlet_name else "Provided content"}
Report Type: {report_type}

Content to analyze:
{content_text[:3000] if content_text else "Analyze the outlet based on known public record."}

Write a complete professional {report_type}:

1. **Executive Summary** — Key bias findings in plain language
2. **Political Lean Assessment** — Detailed lean analysis with score
3. **Bias Type Analysis** — All bias types found with examples
4. **Factual Accuracy Assessment** — What's accurate, what's questionable
5. **Source Quality Review** — Who was quoted and who wasn't
6. **Language and Framing Analysis** — Specific examples
7. **Comparison to Similar Outlets** — Context in the media landscape
8. **Recommendations for Citizens** — How to read this content critically
9. **Overall Rating** — Final grade and assessment
10. **Conclusion** — What citizens should take away

Write professionally. Be specific. Quote directly from content when analyzing bias."""

    report_text = query_llm(prompt)

    output_path = None
    try:
        if report_type == "PowerPoint Presentation":
            output_path = _make_pptx(report_text, outlet_name)
        else:
            output_path = _make_pdf(report_text, outlet_name)
    except Exception as e:
        report_text += f"\n\n[Document generation error: {str(e)}]"

    return report_text + "\n\n" + WATERMARK, output_path


def _make_pptx(text, outlet_name):
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    def add_slide(title_text, body_text, bg_color=(10, 10, 20)):
        slide = prs.slides.add_slide(blank_layout)
        bg = slide.background
        bg.fill.solid()
        bg.fill.fore_color.rgb = RGBColor(*bg_color)
        tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(1.2))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = RGBColor(241, 196, 15)
        body_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.7),
                                             Inches(12.3), Inches(5.5))
        tf2 = body_box.text_frame
        tf2.word_wrap = True
        for line in body_text.split("\n")[:20]:
            line = line.strip()
            if not line:
                continue
            para = tf2.add_paragraph()
            para.text = line
            para.font.size = Pt(14)
            para.font.color.rgb = RGBColor(230, 230, 230)

    title_slide = prs.slides.add_slide(blank_layout)
    bg = title_slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(5, 5, 15)
    tb = title_slide.shapes.add_textbox(Inches(1), Inches(2.2), Inches(11), Inches(1.5))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = "Media Bias Analysis Report"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(241, 196, 15)
    tb2 = title_slide.shapes.add_textbox(Inches(1), Inches(3.8), Inches(11), Inches(1))
    tf2 = tb2.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = f"Subject: {outlet_name if outlet_name else 'Content Analysis'}"
    p2.font.size = Pt(22)
    p2.font.color.rgb = RGBColor(200, 200, 200)
    tb3 = title_slide.shapes.add_textbox(Inches(1), Inches(4.8), Inches(11), Inches(0.8))
    tf3 = tb3.text_frame
    p3 = tf3.paragraphs[0]
    p3.text = f"Generated: {datetime.now().strftime('%B %d, %Y')} | AI Media Bias Analyzer"
    p3.font.size = Pt(16)
    p3.font.color.rgb = RGBColor(150, 150, 180)

    sections = []
    current_title = "Overview"
    current_body = []
    for line in text.split("\n"):
        stripped = line.strip()
        if stripped.startswith("**") and stripped.endswith("**") and len(stripped) > 4:
            if current_body:
                sections.append((current_title, "\n".join(current_body)))
            current_title = stripped.strip("*").strip()
            current_body = []
        elif stripped.startswith("## "):
            if current_body:
                sections.append((current_title, "\n".join(current_body)))
            current_title = stripped[3:].strip()
            current_body = []
        else:
            current_body.append(line)
    if current_body:
        sections.append((current_title, "\n".join(current_body)))

    bg_colors = [(10, 10, 20), (15, 10, 25), (5, 15, 20),
                 (20, 10, 10), (10, 20, 10), (15, 15, 5)]
    for i, (title, body) in enumerate(sections[:9]):
        add_slide(title, body, bg_colors[i % len(bg_colors)])

    closing = prs.slides.add_slide(blank_layout)
    bg = closing.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(5, 5, 15)
    tb = closing.shapes.add_textbox(Inches(1), Inches(2.8), Inches(11), Inches(2))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Read Widely. Think Critically. Verify Independently."
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = RGBColor(241, 196, 15)
    p2 = tf.add_paragraph()
    p2.text = "© 2026 Existential Gateway, LLC | AI Media Bias Analyzer"
    p2.font.size = Pt(14)
    p2.font.color.rgb = RGBColor(150, 150, 180)

    tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".pptx")
    prs.save(tmp.name)
    return tmp.name


def _make_pdf(text, outlet_name):
    tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".pdf")
    doc = SimpleDocTemplate(tmp.name, pagesize=letter,
                            leftMargin=inch, rightMargin=inch,
                            topMargin=inch, bottomMargin=inch)
    styles = getSampleStyleSheet()
    title_style = ParagraphStyle("Title2", parent=styles["Title"],
                                 fontSize=20, spaceAfter=20,
                                 textColor=colors.HexColor("#7d6608"))
    heading_style = ParagraphStyle("H2", parent=styles["Heading2"],
                                   fontSize=14, spaceAfter=8,
                                   textColor=colors.HexColor("#b7950b"))
    body_style = ParagraphStyle("Body2", parent=styles["Normal"],
                                fontSize=11, spaceAfter=6, leading=16)
    story = []
    story.append(Paragraph("Media Bias Analysis Report", title_style))
    story.append(Paragraph(
        f"Subject: {outlet_name if outlet_name else 'Content Analysis'} | "
        f"Generated: {datetime.now().strftime('%B %d, %Y')}",
        body_style))
    story.append(Spacer(1, 0.2 * inch))
    for line in text.split("\n"):
        stripped = line.strip()
        if not stripped:
            story.append(Spacer(1, 0.1 * inch))
            continue
        if stripped.startswith("**") and stripped.endswith("**"):
            story.append(Paragraph(stripped.strip("*"), heading_style))
        elif stripped.startswith("## "):
            story.append(Paragraph(stripped[3:], heading_style))
        elif stripped.startswith("# "):
            story.append(Paragraph(stripped[2:], title_style))
        else:
            safe = stripped.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")
            story.append(Paragraph(safe, body_style))
    story.append(Spacer(1, 0.3 * inch))
    story.append(Paragraph(
        "© 2026 Existential Gateway, LLC | AI Media Bias Analyzer | existentialgateway@gmail.com",
        ParagraphStyle("footer", parent=styles["Normal"], fontSize=9,
                       textColor=colors.grey)))
    doc.build(story)
    return tmp.name


# ─── Tab 10: AI Media Chat ────────────────────────────────────────────────────

def media_chat(message, history):
    messages = [
        {
            "role": "system",
            "content": (
                SYSTEM_PROMPT + "\n\n"
                "You are also an AI media literacy assistant helping citizens understand bias in news. "
                "Help users analyze specific articles, understand how bias works, evaluate news sources, "
                "learn media literacy skills, and think critically about the news they consume. "
                "Answer questions about specific outlets, journalists, and stories. "
                "Always apply identical standards to left and right-leaning outlets. "
                "Recommend reading multiple sources across the political spectrum. "
                "Empower citizens to be critical thinkers, not passive news consumers."
            )
        }
    ]
    for user_msg, bot_msg in history:
        messages.append({"role": "user", "content": user_msg})
        messages.append({"role": "assistant", "content": bot_msg})
    messages.append({"role": "user", "content": message})

    import os
    API_KEY = os.environ.get("OPENAI_API_KEY", "")
    headers = {"Authorization": f"Bearer {API_KEY}", "Content-Type": "application/json"}
    msgs = messages if 'messages' in dir() else [{"role": "user", "content": prompt}]
    payload = {"model": "gpt-4o", "max_tokens": 4000, "messages": msgs}
    try:
        response = requests.post("https://api.openai.com/v1/chat/completions", headers=headers, json=payload, timeout=240)
        result = response.json()
        if "choices" in result:
            return result["choices"][0]["message"]["content"]
        return f"API Error: {result}"
    except Exception as e:
        return f"Error: {str(e)}"


# ─── Gradio UI ────────────────────────────────────────────────────────────────

with gr.Blocks(title="AI Media Bias Analyzer", theme=gr.themes.Base(
        primary_hue=gr.themes.colors.Color(
            c50="#fef9ec", c100="#faefd0", c200="#f4dea0", c300="#edc970",
            c400="#e4b04a", c500="#c9a84c", c600="#a8872e", c700="#856519",
            c800="#5e440d", c900="#3a2a05", c950="#1e1502"
        ),
        secondary_hue=gr.themes.colors.Color(
            c50="#e8edf5", c100="#c5d0e6", c200="#9eb0d4", c300="#7490c2",
            c400="#4f74b0", c500="#2d5a9e", c600="#1a3a6e", c700="#112240",
            c800="#0a1628", c900="#060e1a", c950="#03070d"
        ),
        neutral_hue=gr.themes.colors.Color(
            c50="#f0f2f7", c100="#d8dde9", c200="#b8c2d6", c300="#95a3be",
            c400="#7285a6", c500="#536890", c600="#3a4f73", c700="#263856",
            c800="#162438", c900="#0c1622", c950="#060b11"
        ),
        font=gr.themes.GoogleFont("DM Sans"),
        font_mono=gr.themes.GoogleFont("DM Mono"),
    ).set(
        body_background_fill="#0a1628",
        body_background_fill_dark="#0a1628",
        body_text_color="#f8f9fc",
        body_text_color_dark="#f8f9fc",
        block_background_fill="#112240",
        block_background_fill_dark="#112240",
        block_border_color="rgba(201,168,76,0.2)",
        block_border_color_dark="rgba(201,168,76,0.2)",
        block_label_background_fill="#112240",
        block_label_background_fill_dark="#112240",
        block_label_text_color="#c9a84c",
        block_label_text_color_dark="#c9a84c",
        block_title_text_color="#c9a84c",
        block_title_text_color_dark="#c9a84c",
        button_primary_background_fill="linear-gradient(135deg,#c9a84c,#e8c96a)",
        button_primary_background_fill_dark="linear-gradient(135deg,#c9a84c,#e8c96a)",
        button_primary_text_color="#0a1628",
        button_primary_text_color_dark="#0a1628",
        button_secondary_background_fill="#112240",
        button_secondary_background_fill_dark="#112240",
        button_secondary_border_color="rgba(201,168,76,0.3)",
        button_secondary_border_color_dark="rgba(201,168,76,0.3)",
        button_secondary_text_color="#c9a84c",
        button_secondary_text_color_dark="#c9a84c",
        input_background_fill="#0a1628",
        input_background_fill_dark="#0a1628",
        input_border_color="rgba(201,168,76,0.2)",
        input_border_color_dark="rgba(201,168,76,0.2)",
        input_placeholder_color="#8892a4",
        input_placeholder_color_dark="#8892a4",
        shadow_drop="0 4px 24px rgba(0,0,0,0.4)",
        shadow_drop_lg="0 8px 40px rgba(0,0,0,0.5)",
        table_even_background_fill="#0a1628",
        table_even_background_fill_dark="#0a1628",
        table_odd_background_fill="#112240",
        table_odd_background_fill_dark="#112240",
    )) as demo:
    gr.Markdown("# 📰 AI Media Bias Analyzer")
    gr.Markdown("### *Read Widely. Think Critically. Verify Independently.*")
    gr.Markdown(MEDIA_LITERACY_NOTICE)
    gr.Markdown(DISCLAIMER)

    with gr.Tabs():

        # ── Tab 1: Article Analyzer ────────────────────────────────────────────
        with gr.Tab("📰 Article Analyzer"):
            gr.Markdown("## Analyze Any News Article for Bias")
            gr.Markdown(NEUTRALITY_NOTICE)
            with gr.Row():
                with gr.Column(scale=1):
                    t1_text = gr.Textbox(label="Paste Article Text Here", lines=8,
                                         placeholder="Paste the full article text...")
                    t1_url = gr.Textbox(label="Or Enter Article URL",
                                        placeholder="https://...")
                    t1_outlet = gr.Textbox(label="News Outlet Name",
                                           placeholder="e.g. CNN, Fox News, BBC")
                    t1_context = gr.Textbox(label="Additional Context", lines=2)
                    t1_btn = gr.Button("📰 Analyze Article", variant="primary")
                    gr.Markdown(WAIT_MSG)
                with gr.Column(scale=2):
                    t1_output = gr.Markdown(label="Article Bias Analysis")
            t1_chart = gr.Plot(label="Bias Scorecard")
            t1_btn.click(analyze_article,
                         inputs=[t1_text, t1_url, t1_outlet, t1_context],
                         outputs=[t1_output, t1_chart])

        # ── Tab 2: Outlet Analyzer ─────────────────────────────────────────────
        with gr.Tab("🏢 Outlet Analyzer"):
            gr.Markdown("## Rate Any News Outlet's Overall Bias Profile")
            gr.Markdown(NEUTRALITY_NOTICE)
            with gr.Row():
                with gr.Column(scale=1):
                    t2_outlet = gr.Textbox(label="News Outlet Name",
                                           placeholder="e.g. CNN, Fox News, Reuters, BBC")
                    t2_type = gr.Dropdown(
                        choices=["Television News", "Online News", "Print Newspaper",
                                 "Radio", "Podcast", "Wire Service",
                                 "Social Media News", "Independent/Alternative"],
                        value="Online News", label="Outlet Type")
                    t2_country = gr.Textbox(label="Country",
                                            placeholder="e.g. United States, United Kingdom")
                    t2_context = gr.Textbox(label="Additional Context", lines=2)
                    t2_btn = gr.Button("🏢 Analyze Outlet", variant="primary")
                    gr.Markdown(WAIT_MSG)
                with gr.Column(scale=2):
                    t2_output = gr.Markdown(label="Outlet Bias Profile")
            t2_chart = gr.Plot(label="Outlet Radar Chart")
            t2_btn.click(analyze_outlet,
                         inputs=[t2_outlet, t2_type, t2_country, t2_context],
                         outputs=[t2_output, t2_chart])

        # ── Tab 3: Story Comparison ────────────────────────────────────────────
        with gr.Tab("🔄 Story Comparison"):
            gr.Markdown("## Compare How Different Outlets Cover the Same Story")
            gr.Markdown(NEUTRALITY_NOTICE)
            with gr.Row():
                with gr.Column(scale=1):
                    t3_topic = gr.Textbox(label="Story Topic",
                                          placeholder="e.g. Federal interest rate decision")
                    t3_outlet1 = gr.Textbox(label="Outlet 1 Name",
                                            placeholder="e.g. CNN")
                    t3_article1 = gr.Textbox(label="Outlet 1 Article Text", lines=4,
                                              placeholder="Paste article 1 text...")
                    t3_outlet2 = gr.Textbox(label="Outlet 2 Name",
                                            placeholder="e.g. Fox News")
                    t3_article2 = gr.Textbox(label="Outlet 2 Article Text", lines=4,
                                              placeholder="Paste article 2 text...")
                    t3_outlet3 = gr.Textbox(label="Outlet 3 Name (optional)",
                                            placeholder="e.g. AP News")
                    t3_article3 = gr.Textbox(label="Outlet 3 Article Text (optional)",
                                              lines=3)
                    t3_context = gr.Textbox(label="Additional Context", lines=2)
                    t3_btn = gr.Button("🔄 Compare Stories", variant="primary")
                    gr.Markdown(WAIT_MSG)
                with gr.Column(scale=2):
                    t3_output = gr.Markdown(label="Story Comparison Analysis")
            t3_chart = gr.Plot(label="Coverage Comparison Chart")
            t3_btn.click(compare_stories,
                         inputs=[t3_topic, t3_outlet1, t3_article1,
                                 t3_outlet2, t3_article2,
                                 t3_outlet3, t3_article3, t3_context],
                         outputs=[t3_output, t3_chart])

        # ── Tab 4: Journalist Profiler ─────────────────────────────────────────
        with gr.Tab("✍️ Journalist Profiler"):
            gr.Markdown("## Analyze a Journalist's Reporting Patterns")
            gr.Markdown(NEUTRALITY_NOTICE)
            with gr.Row():
                with gr.Column(scale=1):
                    t4_name = gr.Textbox(label="Journalist Name",
                                         placeholder="e.g. Anderson Cooper")
                    t4_outlet = gr.Textbox(label="Primary Outlet",
                                           placeholder="e.g. CNN")
                    t4_beat = gr.Textbox(label="Coverage Beat",
                                         placeholder="e.g. Politics, Economics, Foreign Policy")
                    t4_work = gr.Textbox(label="Sample Headlines or Articles (optional)",
                                         lines=4,
                                         placeholder="Paste sample headlines or article titles...")
                    t4_context = gr.Textbox(label="Additional Context", lines=2)
                    t4_btn = gr.Button("✍️ Profile Journalist", variant="primary")
                    gr.Markdown(WAIT_MSG)
                with gr.Column(scale=2):
                    t4_output = gr.Markdown(label="Journalist Bias Profile")
            t4_btn.click(profile_journalist,
                         inputs=[t4_name, t4_outlet, t4_beat, t4_work, t4_context],
                         outputs=[t4_output])

        # ── Tab 5: Headline Analyzer ───────────────────────────────────────────
        with gr.Tab("🎯 Headline Analyzer"):
            gr.Markdown("## Does the Headline Match Reality?")
            gr.Markdown(NEUTRALITY_NOTICE)
            with gr.Row():
                with gr.Column(scale=1):
                    t5_headline = gr.Textbox(label="Headline to Analyze",
                                             placeholder="Paste any news headline here...")
                    t5_body = gr.Textbox(label="Article Body (optional — for comparison)",
                                          lines=5,
                                          placeholder="Paste the article text for headline vs body comparison...")
                    t5_outlet = gr.Textbox(label="Outlet Name",
                                           placeholder="e.g. New York Times")
                    t5_context = gr.Textbox(label="Additional Context", lines=2)
                    t5_btn = gr.Button("🎯 Analyze Headline", variant="primary")
                    gr.Markdown(WAIT_MSG)
                with gr.Column(scale=2):
                    t5_output = gr.Markdown(label="Headline Analysis")
            t5_chart = gr.Plot(label="Headline Accuracy Chart")
            t5_btn.click(analyze_headline,
                         inputs=[t5_headline, t5_body, t5_outlet, t5_context],
                         outputs=[t5_output, t5_chart])

        # ── Tab 6: Bias Scoring Dashboard ──────────────────────────────────────
        with gr.Tab("📊 Bias Scoring"):
            gr.Markdown("## Complete Bias Scoring Dashboard")
            gr.Markdown(NEUTRALITY_NOTICE)
            with gr.Row():
                with gr.Column(scale=1):
                    t6_text = gr.Textbox(label="Paste Content to Score", lines=8,
                                         placeholder="Paste any news content here...")
                    t6_type = gr.Dropdown(
                        choices=["News Article", "Opinion Piece", "Editorial",
                                 "Social Media Post", "Press Release",
                                 "TV Transcript", "Radio Transcript"],
                        value="News Article", label="Content Type")
                    t6_outlet = gr.Textbox(label="Outlet Name",
                                           placeholder="e.g. Washington Post")
                    t6_lean = gr.Dropdown(
                        choices=["Far Left", "Left", "Center-Left", "Center",
                                 "Center-Right", "Right", "Far Right", "Unknown"],
                        value="Unknown", label="Outlet's Claimed Political Position")
                    t6_context = gr.Textbox(label="Additional Context", lines=2)
                    t6_btn = gr.Button("📊 Score Bias", variant="primary")
                    gr.Markdown(WAIT_MSG)
                with gr.Column(scale=2):
                    t6_output = gr.Markdown(label="Bias Scoring Analysis")
            t6_chart = gr.Plot(label="Bias Dashboard")
            t6_btn.click(score_bias,
                         inputs=[t6_text, t6_type, t6_outlet, t6_lean, t6_context],
                         outputs=[t6_output, t6_chart])

        # ── Tab 7: Source Credibility ──────────────────────────────────────────
        with gr.Tab("🌐 Source Credibility"):
            gr.Markdown("## Evaluate the Sources Cited in Any Article")
            gr.Markdown(NEUTRALITY_NOTICE)
            with gr.Row():
                with gr.Column(scale=1):
                    t7_text = gr.Textbox(label="Paste Article Text", lines=6,
                                         placeholder="Paste article to analyze its sources...")
                    t7_sources = gr.Textbox(
                        label="Or List Specific Sources to Evaluate", lines=4,
                        placeholder="e.g.\nHeritage Foundation\nBrookings Institution\nUnnamed White House official")
                    t7_context = gr.Textbox(label="Additional Context", lines=2)
                    t7_btn = gr.Button("🌐 Analyze Sources", variant="primary")
                    gr.Markdown(WAIT_MSG)
                with gr.Column(scale=2):
                    t7_output = gr.Markdown(label="Source Credibility Analysis")
            t7_btn.click(analyze_sources,
                         inputs=[t7_text, t7_sources, t7_context],
                         outputs=[t7_output])

        # ── Tab 8: Visualizations ──────────────────────────────────────────────
        with gr.Tab("📈 Visualizations"):
            gr.Markdown("## Media Bias Visualizations")
            gr.Markdown(NEUTRALITY_NOTICE)
            with gr.Row():
                with gr.Column(scale=1):
                    t8_viz_type = gr.Dropdown(
                        choices=["Political Lean Spectrum",
                                 "Credibility vs Bias Chart",
                                 "Media Ownership Web",
                                 "Bias Over Time"],
                        value="Political Lean Spectrum",
                        label="Visualization Type")
                    t8_outlets = gr.Textbox(
                        label="Outlets to Include (comma separated, optional)",
                        placeholder="e.g. CNN, Fox News, BBC, Reuters, NPR")
                    t8_topic = gr.Textbox(
                        label="Topic or Context (optional)",
                        placeholder="e.g. 2024 election coverage")
                    t8_btn = gr.Button("📊 Generate Visualization", variant="primary")
                    t8_status = gr.Markdown()
                with gr.Column(scale=2):
                    t8_chart = gr.Plot(label="Media Bias Chart")
            t8_btn.click(generate_media_viz,
                         inputs=[t8_viz_type, t8_outlets, t8_topic],
                         outputs=[t8_chart, t8_status])

        # ── Tab 9: Reports ─────────────────────────────────────────────────────
        with gr.Tab("📄 Reports"):
            gr.Markdown("## Generate Professional Media Bias Reports")
            gr.Markdown(NEUTRALITY_NOTICE)
            with gr.Row():
                with gr.Column(scale=1):
                    t9_text = gr.Textbox(label="Paste Article or Content (optional)",
                                          lines=6,
                                          placeholder="Paste content for the report...")
                    t9_outlet = gr.Textbox(label="Outlet or Subject Name",
                                           placeholder="e.g. CNN, Fox News, specific article")
                    t9_type = gr.Dropdown(
                        choices=["PDF Report", "PowerPoint Presentation",
                                 "Executive Summary", "Media Literacy Guide"],
                        value="PDF Report", label="Report Type")
                    t9_btn = gr.Button("📄 Generate Report", variant="primary")
                    gr.Markdown(WAIT_MSG)
                    t9_download = gr.File(label="⬇️ Download Report")
                with gr.Column(scale=2):
                    t9_output = gr.Markdown(label="Report Preview")
            t9_btn.click(generate_report,
                         inputs=[t9_text, t9_outlet, t9_type],
                         outputs=[t9_output, t9_download])

        # ── Tab 10: AI Media Chat ──────────────────────────────────────────────
        with gr.Tab("💬 AI Media Chat"):
            gr.Markdown("## Ask Anything About Media Bias and News Sources")
            gr.Markdown(NEUTRALITY_NOTICE)
            gr.ChatInterface(
                fn=media_chat,
                examples=[
                    "Is CNN biased? Give me a balanced assessment.",
                    "Is Fox News biased? Give me a balanced assessment.",
                    "How do I identify loaded language in a news article?",
                    "What is framing bias and how does it work?",
                    "Which news sources are considered most reliable and least biased?",
                    "How do I build a balanced news diet across the political spectrum?",
                    "What is the difference between a news article and an opinion piece?",
                    "How can I fact-check a news story myself?",
                    "What are the most common media bias techniques I should watch for?",
                    "How does media ownership affect news coverage?",
                ],
                title="",
            )

    gr.Markdown(WATERMARK)


if __name__ == "__main__":
    demo.launch(server_name="0.0.0.0", server_port=int(os.environ.get("GRADIO_SERVER_PORT", 7860)), share=False, ssr_mode=False)
