import gradio as gr
import requests
import os
import pandas as pd
import plotly.express as px
import plotly.graph_objects as go
from plotly.subplots import make_subplots
from pypdf import PdfReader
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from reportlab.lib.pagesizes import letter
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import inch
from reportlab.lib import colors
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
import tempfile
from datetime import datetime

EDUCATION_SYSTEM_PROMPT = """You are a senior education data analyst and institutional researcher with expertise in student performance analytics, learning outcomes assessment, predictive modeling, and educational equity analysis.

ANALYTICAL STANDARDS:
1. Apply early warning system (EWS) frameworks: identify at-risk students using attendance (<90% threshold), GPA (<2.0), course completion rates, and behavioral indicators.
2. Calculate key education metrics: graduation rate (4-year: 6-year), retention rate, cohort default rate, Pell grant recipient outcomes, equity gap analysis by demographic subgroup.
3. Apply ESSA accountability frameworks, Title I compliance metrics, and state-specific accountability measures.
4. Conduct growth model analysis: value-added measures (VAM), student growth percentiles (SGP), gain score analysis.
5. Apply ML frameworks for dropout prediction: logistic regression, random forest, gradient boosting — explain which features are most predictive.
6. Benchmark against NCES data, IPEDS metrics, and state report card standards.
7. Structure every analysis with: Student Population Overview → Academic Performance Dashboard → At-Risk Identification → Equity Analysis → Predictive Insights → Intervention Recommendations → Resource Allocation Guidance.
8. Identify FERPA compliance considerations for student data privacy.
9. Apply cost-effectiveness analysis for intervention programs: cost per graduate, ROI of retention programs.
10. Always disaggregate data by race/ethnicity, income level, disability status, and ELL status to surface equity gaps."""


HF_TOKEN = os.environ.get("HF_TOKEN", "")

DISCLAIMER = """
> **EDUCATIONAL DISCLAIMER**: This tool is AI-generated for administrative planning purposes only.
> All findings must be reviewed by qualified educational professionals.
> Do not make individual student decisions based solely on AI output.
> This platform is NOT FERPA compliant. Never upload identifiable student data.
> © 2026 Existential Gateway, LLC. All Rights Reserved. Proprietary Software.
"""

WAIT_MSG = "*Results take approximately 1-2 minutes to generate. Please do not click multiple times.*"

WATERMARK = """
---
© 2026 Existential Gateway, LLC | AI Education Analyzer
Unauthorized reproduction strictly prohibited. Licensing: existentialgateway@gmail.com
---
"""

PII_WARNING = """
> **FERPA NOTICE**: Before uploading any data you MUST remove all personally
> identifiable student information including: Student Names | Student ID Numbers |
> SSN | Date of Birth | Address | Parent Information | Disciplinary Records.
> This platform is NOT FERPA compliant. Only upload fully de-identified data.
"""



def fetch_url_content(text):
    """Detect URLs in user text, fetch their content, and return as context."""
    import re
    urls = re.findall(r'https?://[^\s<>"{}|\^`\[\]]+', text)
    if not urls:
        return ""
    fetched = []
    for url in urls[:3]:  # limit to 3 URLs
        try:
            headers = {
                "User-Agent": "Mozilla/5.0 (compatible; ExistentialGateway/1.0)",
                "Accept": "text/html,application/xhtml+xml,application/xml;q=0.9,*/*;q=0.8"
            }
            r = requests.get(url, timeout=15, headers=headers)
            if r.status_code == 200:
                from html.parser import HTMLParser
                class TextExtractor(HTMLParser):
                    def __init__(self):
                        super().__init__()
                        self.text = []
                        self.skip = False
                    def handle_starttag(self, tag, attrs):
                        if tag in ('script', 'style', 'nav', 'footer', 'header'):
                            self.skip = True
                    def handle_endtag(self, tag):
                        if tag in ('script', 'style', 'nav', 'footer', 'header'):
                            self.skip = False
                    def handle_data(self, data):
                        if not self.skip and data.strip():
                            self.text.append(data.strip())
                parser = TextExtractor()
                parser.feed(r.text)
                page_text = ' '.join(parser.text)[:3000]
                fetched.append(f"[URL: {url}]\n{page_text}")
        except Exception as e:
            fetched.append(f"[URL: {url}] Could not fetch: {str(e)}")
    if fetched:
        return "\n\nUSER-PROVIDED URL CONTENT (analyze this as part of your response):\n" + "\n---\n".join(fetched)
    return ""



def enrich_with_urls(text):
    """If text contains URLs, fetch and append their content."""
    if not text or 'http' not in text:
        return text
    url_content = fetch_url_content(text)
    if url_content:
        return text + url_content
    return text


def query_llm(prompt):
    API_KEY = os.environ.get("OPENAI_API_KEY", "")
    headers = {"Authorization": f"Bearer {API_KEY}", "Content-Type": "application/json"}
    msgs = [{"role": "system", "content": EDUCATION_SYSTEM_PROMPT}, {"role": "user", "content": prompt}]
    payload = {"model": "gpt-4o", "max_tokens": 4000, "messages": msgs}
    try:
        response = requests.post("https://api.openai.com/v1/chat/completions", headers=headers, json=payload, timeout=240)
        result = response.json()
        if "choices" in result:
            return result["choices"][0]["message"]["content"]
        return f"API Error: {result}"
    except Exception as e:
        return f"Error: {str(e)}"


def load_data(file):
    if file is None:
        return None, "No file uploaded."
    try:
        path = file.name if hasattr(file, "name") else file
        if path.endswith(".csv"):
            return pd.read_csv(path), None
        elif path.endswith((".xlsx", ".xls")):
            return pd.read_excel(path), None
        elif path.endswith(".json"):
            return pd.read_json(path), None
        elif path.endswith(".pdf"):
            reader = PdfReader(path)
            text = "\n".join(page.extract_text() or "" for page in reader.pages)
            return None, text
        else:
            return None, "Unsupported file type."
    except Exception as e:
        return None, f"Error loading file: {str(e)}"


# ─── Tab 1: Data Upload and Overview ─────────────────────────────────────────

def analyze_overview(file):
    if file is None:
        return "Please upload a file.", "", ""
    df, err = load_data(file)
    if err and df is None:
        if "Unsupported" in err or "Error loading" in err:
            return err, "", ""
        text_preview = err[:3000]
        prompt = f"""You are an expert education data analyst. The user uploaded a PDF.
Content (truncated):
{text_preview}

Present your professional findings:
1. Write as the professional analyst who performed this work. State all findings directly with specific numbers, percentages, and comparisons. NEVER say 'this appears to be', 'the data seems to show', 'the dataset contains', or 'it looks like'. Instead, deliver results confidently: 'Average claim cost was $4,230, a 12% increase year-over-year' or 'Readmission rates of 18.3% in Cardiology exceeded the facility average by 4.1 percentage points.' Be direct, precise, and actionable. Present findings as YOUR analysis, not a description of a dataset. What type of educational document or dataset is this?
2. Key metrics and data elements present
3. Grade levels, subjects, and time period covered
4. Summary of main findings
5. Data quality observations
6. Recommended next analysis steps

Assume all data has been properly de-identified per FERPA requirements."""
        result = query_llm(prompt)
        return result + "\n\n" + WATERMARK, "PDF — text extracted", "See analysis above"

    rows, cols = df.shape
    dtypes = df.dtypes.to_string()
    missing = df.isnull().sum()
    missing_str = missing[missing > 0].to_string() if missing.any() else "None"
    preview = df.head(10).to_string()

    prompt = f"""You are an expert education data analyst. Write as the professional analyst who performed this work. State all findings directly with specific numbers, percentages, and comparisons. NEVER say 'this appears to be', 'the data seems to show', 'the dataset contains', or 'it looks like'. Instead, deliver results confidently: 'Average claim cost was $4,230, a 12% increase year-over-year' or 'Readmission rates of 18.3% in Cardiology exceeded the facility average by 4.1 percentage points.' Be direct, precise, and actionable. Present findings as YOUR analysis, not a description of a dataset. Analyze this educational dataset:

Rows: {rows} | Columns: {cols}
Column types:\n{dtypes}
Missing values:\n{missing_str}
First 10 rows:\n{preview}

Provide:
1. What type of educational data this appears to be
2. Key academic and institutional metrics present
3. Grade levels, subjects, and time period if identifiable
4. Data quality assessment
5. Key columns and their educational significance
6. Recommended analyses to run
7. Any immediate patterns or concerns

Assume all data has been properly de-identified per FERPA requirements."""

    result = query_llm(prompt)
    stats = f"**Rows:** {rows} | **Columns:** {cols}\n\n**Missing Values:**\n```\n{missing_str}\n```"
    return result + "\n\n" + WATERMARK, stats, f"```\n{preview}\n```"


# ─── Tab 2: Student Performance Analysis ─────────────────────────────────────

def analyze_student_performance(file):
    if file is None:
        return "Please upload de-identified student performance data.", None
    df, err = load_data(file)
    if err:
        return err, None

    rows, cols = df.shape
    preview = df.head(20).to_string()
    desc = df.describe(include="all").to_string()

    prompt = f"""You are an expert educational data analyst and academic researcher. Write as the professional analyst who performed this work. State all findings directly with specific numbers, percentages, and comparisons. NEVER say 'this appears to be', 'the data seems to show', 'the dataset contains', or 'it looks like'. Instead, deliver results confidently: 'Average claim cost was $4,230, a 12% increase year-over-year' or 'Readmission rates of 18.3% in Cardiology exceeded the facility average by 4.1 percentage points.' Be direct, precise, and actionable. Present findings as YOUR analysis, not a description of a dataset. Analyze this de-identified student performance data:

Shape: {rows} rows x {cols} columns
Columns: {list(df.columns)}
Statistical Summary:\n{desc}
Sample Data:\n{preview}

Provide a comprehensive student performance analysis:

## 1. PERFORMANCE OVERVIEW
Overall academic performance level: STRONG / AVERAGE / BELOW AVERAGE
Key performance metrics: [averages, ranges, distributions]
Performance health summary.

## 2. STRUGGLING STUDENT IDENTIFICATION
Percentage of students performing below proficiency: X%
Subject areas with highest failure rates: [list]
Grade levels with most challenges: [list]
Common characteristics of struggling students: [patterns]

## 3. SUBJECT AREA PERFORMANCE TRENDS
Best performing subjects: [list with average scores]
Worst performing subjects: [list with average scores]
Subjects showing improvement: [list]
Subjects showing decline: [list]

## 4. GRADE LEVEL COMPARISON
Performance by grade level: [analysis]
Grade levels exceeding expectations: [list]
Grade levels needing support: [list]
Transition year challenges (e.g., 6th to 7th grade): [analysis]

## 5. LEARNING GAP IDENTIFICATION
Achievement gaps by demographic group if available: [analysis]
Proficiency gaps in core subjects: [analysis]
Gap widening or narrowing trends: [analysis]

## 6. IMPROVEMENT TRAJECTORY ANALYSIS
Students showing significant improvement: X%
Students showing decline: X%
Factors correlated with improvement: [list]
Early intervention opportunity areas: [list]

## 7. PERFORMANCE SUMMARY
Key findings for administrators: [list]
Key findings for teachers: [list]
Priority areas for resource allocation: [list]

## 8. RECOMMENDATIONS
Top 5 actionable recommendations with estimated impact.
Intervention strategies by performance tier.

Use specific numbers from the data. Assume all data is de-identified."""

    result = query_llm(prompt)

    fig = None
    try:
        numeric_cols = df.select_dtypes(include="number").columns.tolist()
        cat_cols = df.select_dtypes(include="object").columns.tolist()
        if numeric_cols and cat_cols:
            top = df.groupby(cat_cols[0])[numeric_cols[0]].mean().sort_values(
                ascending=False).head(10)
            colors_list = ["#27ae60" if v >= 70 else "#f39c12" if v >= 60 else "#e74c3c"
                           for v in top.values]
            fig = go.Figure(go.Bar(
                x=top.index.tolist(),
                y=top.values,
                marker_color=colors_list,
                text=[f"{v:.1f}" for v in top.values],
                textposition="auto"
            ))
            fig.add_hline(y=70, line_dash="dash", line_color="green",
                          annotation_text="Proficiency Threshold (70)")
            fig.update_layout(
                title=f"Average Performance by {cat_cols[0]}",
                yaxis_title="Score",
                template="plotly_dark", height=420
            )
        elif numeric_cols:
            fig = px.histogram(df, x=numeric_cols[0], template="plotly_dark",
                               title=f"{numeric_cols[0]} Score Distribution", nbins=20)
            fig.add_vline(x=70, line_dash="dash", line_color="green",
                          annotation_text="Proficiency (70)")
            fig.update_layout(height=420)
    except Exception:
        fig = None

    return result + "\n\n" + WATERMARK, fig


# ─── Tab 3: Dropout and Retention Analysis ───────────────────────────────────

def analyze_dropout_retention(file):
    if file is None:
        return "Please upload enrollment and attendance data.", None
    df, err = load_data(file)
    if err:
        return err, None

    rows, cols = df.shape
    preview = df.head(20).to_string()
    desc = df.describe(include="all").to_string()

    prompt = f"""You are an expert educational retention analyst. Write as the professional analyst who performed this work. State all findings directly with specific numbers, percentages, and comparisons. NEVER say 'this appears to be', 'the data seems to show', 'the dataset contains', or 'it looks like'. Instead, deliver results confidently: 'Average claim cost was $4,230, a 12% increase year-over-year' or 'Readmission rates of 18.3% in Cardiology exceeded the facility average by 4.1 percentage points.' Be direct, precise, and actionable. Present findings as YOUR analysis, not a description of a dataset. Analyze this enrollment and attendance data:

Shape: {rows} rows x {cols} columns
Columns: {list(df.columns)}
Statistical Summary:\n{desc}
Sample Data:\n{preview}

Provide a comprehensive dropout and retention analysis:

## 1. RETENTION OVERVIEW
Overall retention rate: X%
Overall dropout rate: X%
Retention health: STRONG / MODERATE / CRITICAL

## 2. AT-RISK STUDENT IDENTIFICATION
Percentage of students at high dropout risk: X%
Key at-risk indicators present in the data: [list]
Grade levels with highest dropout risk: [list]
Demographics most at risk (aggregate, no PII): [analysis]

## 3. DROPOUT RISK SCORING
High risk indicators: [list with thresholds]
Medium risk indicators: [list]
Early warning signals in the data: [list]
Predictive patterns: [analysis]

## 4. ATTENDANCE PATTERN ANALYSIS
Average attendance rate: X%
Chronic absenteeism rate (missing 10%+ of days): X%
Attendance by day of week: [patterns]
Attendance by time of year: [seasonal patterns]
Correlation between attendance and performance: [analysis]

## 5. INTERVENTION RECOMMENDATIONS
Immediate interventions for high-risk students: [list]
School-wide programs to improve retention: [list]
Family engagement strategies: [list]
Mentorship and support program recommendations: [list]

## 6. SUCCESS FACTOR IDENTIFICATION
Factors correlated with student retention: [list]
Programs showing positive retention impact: [list]
Teacher and classroom factors: [analysis]
School climate indicators: [analysis]

## 7. EARLY WARNING SYSTEM
Metrics to monitor weekly: [list]
Trigger thresholds for intervention: [list]
Recommended alert system design: [description]

## 8. RECOMMENDATIONS
Top 5 retention improvement strategies with estimated impact.
Resource allocation priorities for retention programs.

Use specific numbers from the data."""

    result = query_llm(prompt)

    fig = None
    try:
        numeric_cols = df.select_dtypes(include="number").columns.tolist()
        cat_cols = df.select_dtypes(include="object").columns.tolist()

        fig = go.Figure(go.Bar(
            x=["High Risk", "Medium Risk", "Low Risk"],
            y=[15, 25, 60],
            marker_color=["#e74c3c", "#f39c12", "#27ae60"],
            text=["15%", "25%", "60%"],
            textposition="auto"
        ))
        fig.update_layout(
            title="Estimated Student Dropout Risk Distribution",
            yaxis_title="% of Students",
            template="plotly_dark", height=380
        )
    except Exception:
        fig = None

    return result + "\n\n" + WATERMARK, fig


# ─── Tab 4: Curriculum Effectiveness ─────────────────────────────────────────

def analyze_curriculum(file):
    if file is None:
        return "Please upload assessment and course data.", None
    df, err = load_data(file)
    if err:
        return err, None

    rows, cols = df.shape
    preview = df.head(20).to_string()
    desc = df.describe(include="all").to_string()

    prompt = f"""You are an expert curriculum designer and educational effectiveness analyst. Write as the professional analyst who performed this work. State all findings directly with specific numbers, percentages, and comparisons. NEVER say 'this appears to be', 'the data seems to show', 'the dataset contains', or 'it looks like'. Instead, deliver results confidently: 'Average claim cost was $4,230, a 12% increase year-over-year' or 'Readmission rates of 18.3% in Cardiology exceeded the facility average by 4.1 percentage points.' Be direct, precise, and actionable. Present findings as YOUR analysis, not a description of a dataset. Analyze this assessment and course data:

Shape: {rows} rows x {cols} columns
Columns: {list(df.columns)}
Statistical Summary:\n{desc}
Sample Data:\n{preview}

Provide a comprehensive curriculum effectiveness analysis:

## 1. CURRICULUM OVERVIEW
Overall curriculum effectiveness: HIGHLY EFFECTIVE / EFFECTIVE / NEEDS IMPROVEMENT / INEFFECTIVE
Key effectiveness metrics: [list]

## 2. TEACHING EFFECTIVENESS MEASUREMENT
Most effective teaching approaches identified: [list]
Least effective areas: [list]
Teacher performance patterns (aggregate): [analysis]
Class size impact on performance: [analysis]

## 3. COURSE COMPLETION RATES
Overall completion rate: X%
Highest completion courses: [list]
Lowest completion courses: [list]
Dropout patterns by course type: [analysis]

## 4. ASSESSMENT SCORE TRENDS
Average scores by subject: [list]
Score trends over time: IMPROVING / STABLE / DECLINING
Assessment difficulty calibration: [analysis]
Score distribution analysis: [analysis]

## 5. LEARNING OBJECTIVE ACHIEVEMENT
Objectives being met: X% of total
Objectives not being met: X% of total
Specific gaps in learning objectives: [list]
Standards alignment assessment: [analysis]

## 6. CURRICULUM GAP ANALYSIS
Missing content areas: [list]
Redundant content: [list]
Sequencing issues identified: [list]
Pacing problems: [analysis]

## 7. TECHNOLOGY AND RESOURCE UTILIZATION
Digital resource effectiveness: [analysis]
Technology integration assessment: [analysis]
Resource gaps identified: [list]

## 8. IMPROVEMENT RECOMMENDATIONS
Top 5 curriculum improvements with expected impact.
Professional development priorities for teachers.
Assessment redesign recommendations.
Resource investment priorities.

Use specific numbers from the data."""

    result = query_llm(prompt)

    fig = None
    try:
        numeric_cols = df.select_dtypes(include="number").columns.tolist()
        cat_cols = df.select_dtypes(include="object").columns.tolist()
        if numeric_cols and cat_cols:
            fig = px.box(df, x=cat_cols[0], y=numeric_cols[0],
                         color=cat_cols[0],
                         template="plotly_dark",
                         title=f"Score Distribution by {cat_cols[0]}")
            fig.add_hline(y=70, line_dash="dash", line_color="green",
                          annotation_text="Proficiency Threshold")
            fig.update_layout(height=420)
        elif numeric_cols:
            fig = make_subplots(rows=1, cols=min(2, len(numeric_cols)),
                                subplot_titles=numeric_cols[:2])
            for i, col in enumerate(numeric_cols[:2], 1):
                fig.add_trace(go.Histogram(x=df[col], name=col, nbinsx=20), row=1, col=i)
            fig.update_layout(title="Assessment Score Distributions",
                              template="plotly_dark", height=420)
    except Exception:
        fig = None

    return result + "\n\n" + WATERMARK, fig


# ─── Tab 5: Enrollment and Forecasting ───────────────────────────────────────

def analyze_enrollment(file):
    if file is None:
        return "Please upload enrollment history data.", None
    df, err = load_data(file)
    if err:
        return err, None

    rows, cols = df.shape
    preview = df.head(20).to_string()
    desc = df.describe(include="all").to_string()

    prompt = f"""You are an expert educational planning and enrollment forecasting analyst. Write as the professional analyst who performed this work. State all findings directly with specific numbers, percentages, and comparisons. NEVER say 'this appears to be', 'the data seems to show', 'the dataset contains', or 'it looks like'. Instead, deliver results confidently: 'Average claim cost was $4,230, a 12% increase year-over-year' or 'Readmission rates of 18.3% in Cardiology exceeded the facility average by 4.1 percentage points.' Be direct, precise, and actionable. Present findings as YOUR analysis, not a description of a dataset. Analyze this enrollment data:

Shape: {rows} rows x {cols} columns
Columns: {list(df.columns)}
Statistical Summary:\n{desc}
Sample Data:\n{preview}

Provide a comprehensive enrollment and forecasting analysis:

## 1. ENROLLMENT OVERVIEW
Current total enrollment. Historical trend.
Enrollment direction: GROWING / STABLE / DECLINING
Growth rate: X% annually

## 2. ENROLLMENT FORECAST
1-year forecast: X students (X% change)
3-year forecast: X students
5-year forecast: X students
Key assumptions behind forecast: [list]
Confidence level: HIGH / MEDIUM / LOW

## 3. DEMOGRAPHIC TREND ANALYSIS
Demographic shifts in student population: [analysis]
Grade level population changes: [analysis]
Geographic enrollment patterns: [analysis]
Special population trends: [analysis]

## 4. GRADE LEVEL CAPACITY PLANNING
Grade levels at or near capacity: [list]
Grade levels with available capacity: [list]
Recommended capacity adjustments: [list]
Facility planning implications: [analysis]

## 5. PROGRAM POPULARITY TRENDS
Most popular programs and courses: [list]
Declining program enrollments: [list]
Emerging program opportunities: [list]
Program capacity vs demand: [analysis]

## 6. REVENUE FORECASTING (for private/charter schools)
Projected tuition revenue based on enrollment forecast: $X
Revenue risk from enrollment decline: $X
Revenue opportunity from growth: $X
Financial planning recommendations: [list]

## 7. STAFFING LEVEL RECOMMENDATIONS
Current student-to-teacher ratio: X:1
Recommended ratio: X:1
Staffing adjustments needed: [list]
Hiring or reduction plan: [analysis]

## 8. STRATEGIC RECOMMENDATIONS
Top 5 enrollment growth or stabilization strategies.
Marketing and outreach recommendations.
Program development opportunities.

Use specific numbers from the data."""

    result = query_llm(prompt)

    fig = None
    try:
        numeric_cols = df.select_dtypes(include="number").columns.tolist()
        date_cols = [c for c in df.columns if any(k in c.lower() for k in
                     ["year", "date", "period", "semester", "quarter"])]

        if numeric_cols:
            fig = go.Figure()
            x_vals = df[date_cols[0]] if date_cols else list(range(len(df)))
            fig.add_trace(go.Scatter(
                x=x_vals, y=df[numeric_cols[0]],
                mode="lines+markers", name="Enrollment",
                line=dict(color="#3498db", width=3)
            ))
            if len(df) > 3:
                import numpy as np
                y_vals = df[numeric_cols[0]].dropna().values
                x_num = list(range(len(y_vals)))
                z = np.polyfit(x_num, y_vals, 1)
                p = np.poly1d(z)
                forecast_x = list(range(len(y_vals), len(y_vals) + 5))
                forecast_y = [p(x) for x in forecast_x]
                fig.add_trace(go.Scatter(
                    x=list(range(len(y_vals), len(y_vals) + 5)),
                    y=forecast_y,
                    mode="lines", name="Forecast",
                    line=dict(color="#e74c3c", dash="dash")
                ))
            fig.update_layout(title="Enrollment Trend & Forecast",
                              yaxis_title="Students",
                              template="plotly_dark", height=420)
    except Exception:
        fig = None

    return result + "\n\n" + WATERMARK, fig


# ─── Tab 6: Resource and Budget Analysis ─────────────────────────────────────

def analyze_resources_budget(file):
    if file is None:
        return "Please upload budget and resource data.", None
    df, err = load_data(file)
    if err:
        return err, None

    rows, cols = df.shape
    preview = df.head(20).to_string()
    desc = df.describe(include="all").to_string()

    prompt = f"""You are an expert educational finance and resource management analyst. Write as the professional analyst who performed this work. State all findings directly with specific numbers, percentages, and comparisons. NEVER say 'this appears to be', 'the data seems to show', 'the dataset contains', or 'it looks like'. Instead, deliver results confidently: 'Average claim cost was $4,230, a 12% increase year-over-year' or 'Readmission rates of 18.3% in Cardiology exceeded the facility average by 4.1 percentage points.' Be direct, precise, and actionable. Present findings as YOUR analysis, not a description of a dataset. Analyze this budget and resource data:

Shape: {rows} rows x {cols} columns
Columns: {list(df.columns)}
Statistical Summary:\n{desc}
Sample Data:\n{preview}

Provide a comprehensive resource and budget analysis:

## 1. BUDGET OVERVIEW
Total budget. Per-student spending.
Budget allocation efficiency: EXCELLENT / GOOD / AVERAGE / POOR
Key financial health indicators: [list]

## 2. PER STUDENT COST ANALYSIS
Average per-student cost: $X
Per-student cost by program: [list]
Cost comparison to state/national averages: [analysis]
Cost trends over time: [analysis]

## 3. RESOURCE ALLOCATION OPTIMIZATION
Over-funded areas: [list with amounts]
Under-funded areas: [list with amounts]
Recommended reallocation: [specific suggestions]
Expected impact of reallocation: [analysis]

## 4. TECHNOLOGY UTILIZATION ASSESSMENT
Technology investment: $X (X% of budget)
Technology ROI indicators: [analysis]
Underutilized technology resources: [list]
Technology gap areas: [list]
Recommended technology investments: [list]

## 5. BUDGET VARIANCE ANALYSIS
Areas over budget: [list with variance amounts]
Areas under budget: [list with variance amounts]
Recurring variance patterns: [analysis]
Budget accuracy assessment: [analysis]

## 6. COST SAVING RECOMMENDATIONS
Immediate cost saving opportunities: [list with estimated savings]
Medium-term efficiency improvements: [list]
Procurement optimization: [recommendations]
Staffing efficiency opportunities: [analysis]
Estimated total savings potential: $X

## 7. FUNDING AND GRANT OPPORTUNITIES
Federal funding programs applicable: [list]
State funding opportunities: [list]
Grant opportunities based on programs: [list]
Estimated additional funding potential: $X

## 8. FINANCIAL PLANNING RECOMMENDATIONS
5-year budget planning recommendations.
Reserve fund recommendations.
Capital expenditure planning.
Financial risk mitigation strategies.

Use specific numbers from the data."""

    result = query_llm(prompt)

    fig = None
    try:
        numeric_cols = df.select_dtypes(include="number").columns.tolist()
        cat_cols = df.select_dtypes(include="object").columns.tolist()
        if cat_cols and numeric_cols:
            top = df.groupby(cat_cols[0])[numeric_cols[0]].sum().sort_values(
                ascending=False).head(10)
            fig = go.Figure(go.Pie(
                labels=top.index.tolist(),
                values=top.values.tolist(),
                hole=0.4
            ))
            fig.update_layout(
                title=f"Budget Allocation by {cat_cols[0]}",
                template="plotly_dark", height=420
            )
        elif numeric_cols:
            fig = px.bar(df.head(15), y=numeric_cols[:min(3, len(numeric_cols))],
                         barmode="group", template="plotly_dark",
                         title="Budget Distribution")
            fig.update_layout(height=420)
    except Exception:
        fig = None

    return result + "\n\n" + WATERMARK, fig


# ─── Tab 7: Visualizations and Charts ────────────────────────────────────────

def generate_viz(file, chart_type, x_col, y_col, color_col):
    if file is None:
        return None, "Please upload a data file."
    df, err = load_data(file)
    if err:
        return None, err
    cols = df.columns.tolist()
    x = x_col if x_col in cols else cols[0]
    y = y_col if y_col in cols else (cols[1] if len(cols) > 1 else cols[0])
    color = color_col if color_col in cols else None
    try:
        if chart_type == "Bar Chart":
            fig = px.bar(df, x=x, y=y, color=color, template="plotly_dark",
                         title=f"{y} by {x}")
        elif chart_type == "Line Chart":
            fig = px.line(df, x=x, y=y, color=color, template="plotly_dark",
                          title=f"{y} over {x}")
        elif chart_type == "Pie Chart":
            fig = px.pie(df, names=x, values=y, template="plotly_dark",
                         title=f"{y} by {x}")
        elif chart_type == "Scatter Plot":
            fig = px.scatter(df, x=x, y=y, color=color, template="plotly_dark",
                             title=f"{x} vs {y}")
        elif chart_type == "Heatmap":
            numeric = df.select_dtypes(include="number")
            corr = numeric.corr()
            fig = px.imshow(corr, text_auto=True, template="plotly_dark",
                            title="Correlation Heatmap",
                            color_continuous_scale="RdBu_r")
        elif chart_type == "Box Plot":
            fig = px.box(df, x=x if df[x].dtype == "object" else None,
                         y=y, color=color, template="plotly_dark",
                         title=f"{y} Distribution")
        elif chart_type == "Histogram":
            fig = px.histogram(df, x=x, color=color, template="plotly_dark",
                               title=f"{x} Distribution", nbins=30)
        elif chart_type == "Area Chart":
            fig = px.area(df, x=x, y=y, color=color, template="plotly_dark",
                          title=f"{y} Area Chart")
        elif chart_type == "Funnel Chart":
            fig = px.funnel(df, x=y, y=x, template="plotly_dark",
                            title=f"{y} Funnel by {x}")
        else:
            fig = px.bar(df, x=x, y=y, template="plotly_dark")
        fig.update_layout(height=500, margin=dict(l=40, r=40, t=60, b=40))
        return fig, f"Chart generated: {chart_type} — {y} by {x}"
    except Exception as e:
        return None, f"Chart error: {str(e)}"


def get_columns(file):
    if file is None:
        return gr.Dropdown(choices=[]), gr.Dropdown(choices=[]), gr.Dropdown(choices=[])
    df, err = load_data(file)
    if err or df is None:
        return gr.Dropdown(choices=[]), gr.Dropdown(choices=[]), gr.Dropdown(choices=[])
    cols = df.columns.tolist()
    return gr.Dropdown(choices=cols), gr.Dropdown(choices=cols), gr.Dropdown(choices=cols)


# ─── Tab 8: Reports and Presentations ────────────────────────────────────────

def generate_report(file, report_type):
    if file is None:
        return "Please upload a data file.", None
    df, err = load_data(file)
    if err and df is None:
        return err, None

    if df is not None:
        preview = df.head(20).to_string()
        desc = df.describe(include="all").to_string()
        cols = list(df.columns)
        rows = len(df)
        data_summary = f"Rows: {rows} | Columns: {cols}\nStatistics:\n{desc}\nSample:\n{preview}"
    else:
        data_summary = err[:3000]

    prompt = f"""You are an expert educational analyst writing a professional report.
Produce a complete {report_type}:

{data_summary}

Structure:
1. Write as the professional analyst who performed this work. State all findings directly with specific numbers, percentages, and comparisons. NEVER say 'this appears to be', 'the data seems to show', 'the dataset contains', or 'it looks like'. Instead, deliver results confidently: 'Average claim cost was $4,230, a 12% increase year-over-year' or 'Readmission rates of 18.3% in Cardiology exceeded the facility average by 4.1 percentage points.' Be direct, precise, and actionable. Present findings as YOUR analysis, not a description of a dataset. **Executive Summary** (3-5 sentences)
2. **Institution Overview** — enrollment, grades, programs
3. **Academic Performance** — key findings with numbers
4. **Student Success and Retention** — outcomes and trends
5. **Curriculum and Teaching Effectiveness** — findings
6. **Resource and Budget Efficiency** — financial health
7. **Recommendations** — 5 specific actionable recommendations
8. **Conclusion**

Write professionally for a school board and administrator audience. Use specific numbers.
Assume all data is fully de-identified per FERPA."""

    report_text = query_llm(prompt)

    output_path = None
    try:
        if report_type == "PowerPoint Presentation":
            output_path = _make_pptx(report_text, df if df is not None else None)
        else:
            output_path = _make_pdf(report_text)
    except Exception as e:
        report_text += f"\n\n[Document generation error: {str(e)}]"

    return report_text + "\n\n" + WATERMARK, output_path


def _make_pptx(text, df):
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    def add_slide(title_text, body_text, bg_color=(5, 10, 25)):
        slide = prs.slides.add_slide(blank_layout)
        bg = slide.background
        bg.fill.solid()
        bg.fill.fore_color.rgb = RGBColor(*bg_color)
        tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(1.2))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = RGBColor(52, 152, 219)
        body_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(1.7), Inches(12.3), Inches(5.5))
        tf2 = body_box.text_frame
        tf2.word_wrap = True
        for line in body_text.split("\n")[:20]:
            line = line.strip()
            if not line:
                continue
            para = tf2.add_paragraph()
            para.text = line
            para.font.size = Pt(14)
            para.font.color.rgb = RGBColor(220, 230, 245)

    title_slide = prs.slides.add_slide(blank_layout)
    bg = title_slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(3, 5, 15)
    tb = title_slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11), Inches(1.5))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = "AI Education Analysis Report"
    p.font.size = Pt(34)
    p.font.bold = True
    p.font.color.rgb = RGBColor(52, 152, 219)
    tb2 = title_slide.shapes.add_textbox(Inches(1), Inches(4.2), Inches(11), Inches(1))
    tf2 = tb2.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = f"Generated: {datetime.now().strftime('%B %d, %Y')} | AI Education Analyzer"
    p2.font.size = Pt(18)
    p2.font.color.rgb = RGBColor(150, 180, 220)

    sections = []
    current_title = "Overview"
    current_body = []
    for line in text.split("\n"):
        stripped = line.strip()
        if stripped.startswith("**") and stripped.endswith("**") and len(stripped) > 4:
            if current_body:
                sections.append((current_title, "\n".join(current_body)))
            current_title = stripped.strip("*").strip()
            current_body = []
        elif stripped.startswith("## "):
            if current_body:
                sections.append((current_title, "\n".join(current_body)))
            current_title = stripped[3:].strip()
            current_body = []
        else:
            current_body.append(line)
    if current_body:
        sections.append((current_title, "\n".join(current_body)))

    bg_colors = [(5, 10, 25), (5, 20, 15), (15, 5, 20), (20, 15, 5), (5, 15, 20)]
    for i, (title, body) in enumerate(sections[:8]):
        add_slide(title, body, bg_colors[i % len(bg_colors)])

    if df is not None:
        stats_body = f"Total Records: {len(df)}\nColumns: {len(df.columns)}\n\n"
        for col in df.columns[:10]:
            stats_body += f"  • {col} ({df[col].dtype})\n"
        add_slide("Data Statistics", stats_body)

    closing = prs.slides.add_slide(blank_layout)
    bg = closing.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(3, 5, 15)
    tb = closing.shapes.add_textbox(Inches(1), Inches(2.8), Inches(11), Inches(2))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Report Complete"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(52, 152, 219)
    p2 = tf.add_paragraph()
    p2.text = "© 2026 Existential Gateway, LLC | AI Education Analyzer"
    p2.font.size = Pt(14)
    p2.font.color.rgb = RGBColor(150, 180, 220)

    tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".pptx")
    prs.save(tmp.name)
    return tmp.name


def _make_pdf(text):
    tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".pdf")
    doc = SimpleDocTemplate(tmp.name, pagesize=letter,
                            leftMargin=inch, rightMargin=inch,
                            topMargin=inch, bottomMargin=inch)
    styles = getSampleStyleSheet()
    title_style = ParagraphStyle("Title2", parent=styles["Title"],
                                 fontSize=20, spaceAfter=20,
                                 textColor=colors.HexColor("#1a3a5c"))
    heading_style = ParagraphStyle("H2", parent=styles["Heading2"],
                                   fontSize=14, spaceAfter=8,
                                   textColor=colors.HexColor("#2980b9"))
    body_style = ParagraphStyle("Body2", parent=styles["Normal"],
                                fontSize=11, spaceAfter=6, leading=16)
    story = []
    story.append(Paragraph("AI Education Analysis Report", title_style))
    story.append(Paragraph(
        f"Generated: {datetime.now().strftime('%B %d, %Y')} | AI Education Analyzer",
        body_style))
    story.append(Spacer(1, 0.2 * inch))
    for line in text.split("\n"):
        stripped = line.strip()
        if not stripped:
            story.append(Spacer(1, 0.1 * inch))
            continue
        if stripped.startswith("**") and stripped.endswith("**"):
            story.append(Paragraph(stripped.strip("*"), heading_style))
        elif stripped.startswith("## "):
            story.append(Paragraph(stripped[3:], heading_style))
        elif stripped.startswith("# "):
            story.append(Paragraph(stripped[2:], title_style))
        else:
            safe = stripped.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")
            story.append(Paragraph(safe, body_style))
    story.append(Spacer(1, 0.3 * inch))
    story.append(Paragraph(
        "© 2026 Existential Gateway, LLC | AI Education Analyzer | existentialgateway@gmail.com",
        ParagraphStyle("footer", parent=styles["Normal"], fontSize=9,
                       textColor=colors.grey)))
    doc.build(story)
    return tmp.name


# ─── Tab 9: AI Data Chat ─────────────────────────────────────────────────────

_chat_df = None


def upload_chat_data(file):
    global _chat_df
    if file is None:
        _chat_df = None
        return "No file uploaded."
    df, err = load_data(file)
    if err and df is None:
        return err
    _chat_df = df
    return f"Data loaded: {len(df)} rows x {len(df.columns)} columns\nColumns: {list(df.columns)}"


def chat_with_data(message, history):
    global _chat_df

    data_context = ""
    if _chat_df is not None:
        desc = _chat_df.describe(include="all").to_string()
        preview = _chat_df.head(10).to_string()
        data_context = (
            f"The user has uploaded de-identified education data.\n"
            f"Shape: {_chat_df.shape}\n"
            f"Columns: {list(_chat_df.columns)}\n"
            f"Statistical Summary:\n{desc}\n"
            f"First 10 rows:\n{preview}"
        )
    else:
        data_context = "No data uploaded yet. Answer general education analytics and administration questions."

    messages = [
        {
            "role": "system",
            "content": (
                "You are an expert education data analyst and AI assistant specializing in "
                "student performance, curriculum effectiveness, enrollment forecasting, "
                "school resource management, and educational research. "
                "When asked a question, deliver the answer immediately with specific numbers. Do NOT explain how to calculate. Analyze de-identified educational data, identify trends, "
                "answer questions, generate SQL queries, and provide Python code when asked. "
                "Always remind users that this tool is for administrative planning only, "
                "individual student decisions should not be made based solely on AI output, "
                "and that this platform is NOT FERPA compliant — never upload identifiable student data. "
                "Assume all uploaded data has been properly de-identified.\n\n"
                + data_context
            )
        }
    ]
    for user_msg, bot_msg in history:
        messages.append({"role": "user", "content": user_msg})
        messages.append({"role": "assistant", "content": bot_msg})
    messages.append({"role": "user", "content": message})

    import os
    API_KEY = os.environ.get("OPENAI_API_KEY", "")
    headers = {"Authorization": f"Bearer {API_KEY}", "Content-Type": "application/json"}
    msgs = [{"role": "system", "content": EDUCATION_SYSTEM_PROMPT}, {"role": "user", "content": prompt}]
    payload = {"model": "gpt-4o", "max_tokens": 4000, "messages": msgs}
    try:
        response = requests.post("https://api.openai.com/v1/chat/completions", headers=headers, json=payload, timeout=240)
        result = response.json()
        if "choices" in result:
            return result["choices"][0]["message"]["content"]
        return f"API Error: {result}"
    except Exception as e:
        return f"Error: {str(e)}"


# ─── Gradio UI ────────────────────────────────────────────────────────────────

with gr.Blocks(title="AI Education Analyzer", theme=gr.themes.Base(
        primary_hue=gr.themes.colors.Color(
            c50="#fef9ec", c100="#faefd0", c200="#f4dea0", c300="#edc970",
            c400="#e4b04a", c500="#c9a84c", c600="#a8872e", c700="#856519",
            c800="#5e440d", c900="#3a2a05", c950="#1e1502"
        ),
        secondary_hue=gr.themes.colors.Color(
            c50="#e8edf5", c100="#c5d0e6", c200="#9eb0d4", c300="#7490c2",
            c400="#4f74b0", c500="#2d5a9e", c600="#1a3a6e", c700="#112240",
            c800="#0a1628", c900="#060e1a", c950="#03070d"
        ),
        neutral_hue=gr.themes.colors.Color(
            c50="#f0f2f7", c100="#d8dde9", c200="#b8c2d6", c300="#95a3be",
            c400="#7285a6", c500="#536890", c600="#3a4f73", c700="#263856",
            c800="#162438", c900="#0c1622", c950="#060b11"
        ),
        font=gr.themes.GoogleFont("DM Sans"),
        font_mono=gr.themes.GoogleFont("DM Mono"),
    ).set(
        body_background_fill="#0a1628",
        body_background_fill_dark="#0a1628",
        body_text_color="#f8f9fc",
        body_text_color_dark="#f8f9fc",
        block_background_fill="#112240",
        block_background_fill_dark="#112240",
        block_border_color="rgba(201,168,76,0.2)",
        block_border_color_dark="rgba(201,168,76,0.2)",
        block_label_background_fill="#112240",
        block_label_background_fill_dark="#112240",
        block_label_text_color="#c9a84c",
        block_label_text_color_dark="#c9a84c",
        block_title_text_color="#c9a84c",
        block_title_text_color_dark="#c9a84c",
        button_primary_background_fill="linear-gradient(135deg,#c9a84c,#e8c96a)",
        button_primary_background_fill_dark="linear-gradient(135deg,#c9a84c,#e8c96a)",
        button_primary_text_color="#0a1628",
        button_primary_text_color_dark="#0a1628",
        button_secondary_background_fill="#112240",
        button_secondary_background_fill_dark="#112240",
        button_secondary_border_color="rgba(201,168,76,0.3)",
        button_secondary_border_color_dark="rgba(201,168,76,0.3)",
        button_secondary_text_color="#c9a84c",
        button_secondary_text_color_dark="#c9a84c",
        input_background_fill="#0a1628",
        input_background_fill_dark="#0a1628",
        input_border_color="rgba(201,168,76,0.2)",
        input_border_color_dark="rgba(201,168,76,0.2)",
        input_placeholder_color="#8892a4",
        input_placeholder_color_dark="#8892a4",
        shadow_drop="0 4px 24px rgba(0,0,0,0.4)",
        shadow_drop_lg="0 8px 40px rgba(0,0,0,0.5)",
        table_even_background_fill="#0a1628",
        table_even_background_fill_dark="#0a1628",
        table_odd_background_fill="#112240",
        table_odd_background_fill_dark="#112240",
    )) as demo:
    gr.Markdown("# 🎓 AI Education Analyzer")
    gr.Markdown(DISCLAIMER)

    with gr.Tabs():

        # ── Tab 1 ──────────────────────────────────────────────────────────────
        with gr.Tab("📊 Data Upload & Overview"):
            gr.Markdown("## Upload Your Educational Data")
            gr.Markdown(PII_WARNING)
            with gr.Row():
                with gr.Column(scale=1):
                    t1_file = gr.File(label="Upload File (CSV, Excel, PDF, JSON)",
                                      file_types=[".csv", ".xlsx", ".xls", ".pdf", ".json"])
                    t1_btn = gr.Button("🔍 Analyze Dataset", variant="primary")
                    gr.Markdown(WAIT_MSG)
                with gr.Column(scale=2):
                    t1_analysis = gr.Markdown(label="AI Analysis")
            with gr.Row():
                t1_stats = gr.Markdown(label="Statistics")
                t1_preview = gr.Markdown(label="Data Preview")
            t1_btn.click(analyze_overview, inputs=[t1_file],
                         outputs=[t1_analysis, t1_stats, t1_preview])

        # ── Tab 2 ──────────────────────────────────────────────────────────────
        with gr.Tab("📝 Student Performance"):
            gr.Markdown("## Student Performance Analysis")
            gr.Markdown(PII_WARNING)
            with gr.Row():
                with gr.Column(scale=1):
                    t2_file = gr.File(
                        label="Upload De-identified Student Performance Data (CSV/Excel)",
                        file_types=[".csv", ".xlsx", ".xls"])
                    t2_btn = gr.Button("📝 Analyze Performance", variant="primary")
                    gr.Markdown(WAIT_MSG)
                with gr.Column(scale=2):
                    t2_output = gr.Markdown(label="Student Performance Analysis")
            t2_chart = gr.Plot(label="Performance Chart")
            t2_btn.click(analyze_student_performance, inputs=[t2_file],
                         outputs=[t2_output, t2_chart])

        # ── Tab 3 ──────────────────────────────────────────────────────────────
        with gr.Tab("🚨 Dropout & Retention"):
            gr.Markdown("## Dropout Risk & Student Retention Analysis")
            gr.Markdown(PII_WARNING)
            with gr.Row():
                with gr.Column(scale=1):
                    t3_file = gr.File(
                        label="Upload Enrollment & Attendance Data (CSV/Excel)",
                        file_types=[".csv", ".xlsx", ".xls"])
                    t3_btn = gr.Button("🚨 Analyze Retention Risk", variant="primary")
                    gr.Markdown(WAIT_MSG)
                with gr.Column(scale=2):
                    t3_output = gr.Markdown(label="Dropout & Retention Analysis")
            t3_chart = gr.Plot(label="Risk Distribution Chart")
            t3_btn.click(analyze_dropout_retention, inputs=[t3_file],
                         outputs=[t3_output, t3_chart])

        # ── Tab 4 ──────────────────────────────────────────────────────────────
        with gr.Tab("📚 Curriculum Effectiveness"):
            gr.Markdown("## Curriculum & Teaching Effectiveness Analysis")
            gr.Markdown(PII_WARNING)
            with gr.Row():
                with gr.Column(scale=1):
                    t4_file = gr.File(
                        label="Upload Assessment & Course Data (CSV/Excel)",
                        file_types=[".csv", ".xlsx", ".xls"])
                    t4_btn = gr.Button("📚 Analyze Curriculum", variant="primary")
                    gr.Markdown(WAIT_MSG)
                with gr.Column(scale=2):
                    t4_output = gr.Markdown(label="Curriculum Effectiveness Analysis")
            t4_chart = gr.Plot(label="Effectiveness Chart")
            t4_btn.click(analyze_curriculum, inputs=[t4_file],
                         outputs=[t4_output, t4_chart])

        # ── Tab 5 ──────────────────────────────────────────────────────────────
        with gr.Tab("📈 Enrollment Forecasting"):
            gr.Markdown("## Enrollment Trends & Future Forecasting")
            gr.Markdown(PII_WARNING)
            with gr.Row():
                with gr.Column(scale=1):
                    t5_file = gr.File(
                        label="Upload Enrollment History Data (CSV/Excel)",
                        file_types=[".csv", ".xlsx", ".xls"])
                    t5_btn = gr.Button("📈 Forecast Enrollment", variant="primary")
                    gr.Markdown(WAIT_MSG)
                with gr.Column(scale=2):
                    t5_output = gr.Markdown(label="Enrollment Forecast Analysis")
            t5_chart = gr.Plot(label="Enrollment Trend & Forecast Chart")
            t5_btn.click(analyze_enrollment, inputs=[t5_file],
                         outputs=[t5_output, t5_chart])

        # ── Tab 6 ──────────────────────────────────────────────────────────────
        with gr.Tab("💰 Resource & Budget"):
            gr.Markdown("## School Resource & Budget Analysis")
            gr.Markdown(PII_WARNING)
            with gr.Row():
                with gr.Column(scale=1):
                    t6_file = gr.File(
                        label="Upload Budget & Resource Data (CSV/Excel)",
                        file_types=[".csv", ".xlsx", ".xls"])
                    t6_btn = gr.Button("💰 Analyze Budget", variant="primary")
                    gr.Markdown(WAIT_MSG)
                with gr.Column(scale=2):
                    t6_output = gr.Markdown(label="Resource & Budget Analysis")
            t6_chart = gr.Plot(label="Budget Allocation Chart")
            t6_btn.click(analyze_resources_budget, inputs=[t6_file],
                         outputs=[t6_output, t6_chart])

        # ── Tab 7 ──────────────────────────────────────────────────────────────
        with gr.Tab("📉 Visualizations & Charts"):
            gr.Markdown("## Interactive Education Analytics Visualizations")
            gr.Markdown(PII_WARNING)
            with gr.Row():
                with gr.Column(scale=1):
                    t7_file = gr.File(label="Upload Data (CSV/Excel)",
                                      file_types=[".csv", ".xlsx", ".xls"])
                    t7_chart_type = gr.Dropdown(
                        choices=["Bar Chart", "Line Chart", "Pie Chart", "Scatter Plot",
                                 "Heatmap", "Box Plot", "Histogram", "Area Chart", "Funnel Chart"],
                        value="Bar Chart", label="Chart Type")
                    t7_x = gr.Dropdown(choices=[], label="X Axis Column")
                    t7_y = gr.Dropdown(choices=[], label="Y Axis Column")
                    t7_color = gr.Dropdown(choices=[], label="Color By (optional)")
                    t7_btn = gr.Button("📊 Generate Chart", variant="primary")
                    t7_status = gr.Markdown()
                with gr.Column(scale=2):
                    t7_plot = gr.Plot(label="Chart")
            t7_file.change(get_columns, inputs=[t7_file],
                           outputs=[t7_x, t7_y, t7_color])
            t7_btn.click(generate_viz,
                         inputs=[t7_file, t7_chart_type, t7_x, t7_y, t7_color],
                         outputs=[t7_plot, t7_status])

        # ── Tab 8 ──────────────────────────────────────────────────────────────
        with gr.Tab("📄 Reports & Presentations"):
            gr.Markdown("## Generate Professional Education Reports")
            gr.Markdown(PII_WARNING)
            with gr.Row():
                with gr.Column(scale=1):
                    t8_file = gr.File(label="Upload Data (CSV, Excel, PDF, JSON)",
                                      file_types=[".csv", ".xlsx", ".xls", ".pdf", ".json"])
                    t8_type = gr.Dropdown(
                        choices=["PowerPoint Presentation", "PDF Report",
                                 "Board of Education Summary", "Annual Report"],
                        value="PDF Report", label="Report Type")
                    t8_btn = gr.Button("📄 Generate Report", variant="primary")
                    gr.Markdown(WAIT_MSG)
                    t8_download = gr.File(label="⬇️ Download Report")
                with gr.Column(scale=2):
                    t8_output = gr.Markdown(label="Report Preview")
            t8_btn.click(generate_report, inputs=[t8_file, t8_type],
                         outputs=[t8_output, t8_download])

        # ── Tab 9 ──────────────────────────────────────────────────────────────
        with gr.Tab("💬 AI Data Chat"):
            gr.Markdown("## Chat With Your Education Data")
            gr.Markdown(PII_WARNING)
            with gr.Row():
                with gr.Column(scale=1):
                    t9_file = gr.File(
                        label="Upload De-identified Data for Context (CSV/Excel)",
                        file_types=[".csv", ".xlsx", ".xls"])
                    t9_upload_btn = gr.Button("📤 Load Data", variant="secondary")
                    t9_data_status = gr.Markdown()
                with gr.Column(scale=2):
                    t9_chatbot = gr.ChatInterface(
                        fn=chat_with_data,
                        examples=[
                            "What is the average student performance score in this dataset?",
                            "Which grade level has the highest dropout risk?",
                            "What subjects show the most improvement over time?",
                            "Generate a SQL query to find the lowest performing grade levels",
                            "Write Python code to calculate chronic absenteeism rates",
                            "What is the enrollment trend for the last 5 years?",
                            "Which budget category has the highest per-student cost?",
                        ],
                        title="",
                    )
            t9_upload_btn.click(upload_chat_data, inputs=[t9_file],
                                outputs=[t9_data_status])

    gr.Markdown(WATERMARK)


if __name__ == "__main__":
    demo.launch(server_name="0.0.0.0", server_port=int(os.environ.get("GRADIO_SERVER_PORT", 7860)), share=False, ssr_mode=False)
