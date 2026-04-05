import gradio as gr
import requests
import os
import plotly.graph_objects as go
import plotly.express as px
from plotly.subplots import make_subplots
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from reportlab.lib.pagesizes import letter
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import inch
from reportlab.lib import colors
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
import tempfile
from datetime import datetime

HF_TOKEN = os.environ.get("HF_TOKEN", "")

DISCLAIMER = """
> **DISCLAIMER**: This tool is for educational, theological, and research purposes only.
> All prophetic interpretations are AI-generated based on scholarly and theological sources.
> This tool presents multiple theological perspectives and does not endorse any single
> eschatological viewpoint. Biblical prophecy interpretation is complex and scholars disagree
> on many points. Always consult qualified theologians and clergy for spiritual guidance.
> © 2026 Existential Gateway, LLC. All Rights Reserved. Proprietary Software.
"""

THEOLOGICAL_NOTICE = """
> **📖 THEOLOGICAL NOTICE**: This tool analyzes end times prophecy from multiple Christian,
> Jewish, and Ethiopian Orthodox traditions. It references the KJV, NKJV, Septuagint,
> Ethiopian Orthodox Canon, Book of Enoch, Book of Jubilees, and other extra-biblical texts.
> Pre-Tribulation, Mid-Tribulation, Post-Tribulation, and Preterist views are all presented.
> The goal is education and awareness — not to cause fear or predict exact dates.
> *"But of that day and hour knoweth no man" — Matthew 24:36*
"""

WAIT_MSG = "*Results take approximately 1-2 minutes to generate. Please do not click multiple times.*"

WATERMARK = """
---
© 2026 Existential Gateway, LLC | AI End Times Prophecy Analyzer
Unauthorized reproduction strictly prohibited. Licensing: existentialgateway@gmail.com
*Watching and Waiting — Matthew 24:42*
---
"""

SYSTEM_PROMPT = """You are a nondenominational biblical scholar, theologian, and eschatology researcher
with deep knowledge of:
- The Holy Bible (KJV, NKJV, NIV, ESV, Septuagint, Ethiopian Orthodox Canon)
- The Book of Enoch (1 Enoch, 2 Enoch, 3 Enoch)
- The Book of Jubilees
- The Book of Jasher
- The Ethiopian Orthodox Tewahedo Biblical Canon (81 books)
- The Dead Sea Scrolls
- Jewish eschatology and the Talmud
- Islamic eschatology (for comparative purposes)
- Pre-Tribulation, Mid-Tribulation, Post-Tribulation, and Preterist theological frameworks
- Astronomy, solar activity, and their prophetic significance
- Current geopolitical events and their prophetic alignment
- The Abraham Accords and Noahide Laws

CRITICAL CURRENT EVENTS FACTS AS OF 2026:
- Donald Trump is the 47th President of the United States (inaugurated January 20, 2025)
- Joe Biden was the 46th President and left office on January 20, 2025
- The current year is 2026 — use accurate 2025-2026 geopolitical information
- Always refer to Trump as the current US President in political analysis
- Vladimir Putin is President of Russia
- Klaus Schwab stepped down from the WEF — Alois Zwinggi is now President and CEO of the World Economic Forum
- Benjamin Netanyahu is Prime Minister of Israel
- Xi Jinping is President of China
- Pope Leo XIV is the current Pope and leads the Roman Catholic Church (Pope Francis died/retired — Pope Leo XIV elected 2025)
- The Abraham Accords were signed in 2020 under the Trump administration

You present ALL major theological views fairly and without bias.
You use precise biblical citations (book, chapter, verse) for every prophetic claim.
You are reverent, scholarly, and thorough.
You acknowledge when scholarly interpretations differ.
You never predict exact dates but analyze prophetic alignment based on current events.

CRITICAL OPERATING INSTRUCTIONS:
1. Do not add generic AI disclaimers mid-analysis. The theological disclaimer is already established.
2. Do not soften prophetic analysis with phrases like 'some might argue' or 'it is worth considering'. Present analysis decisively from within the biblical framework.
3. Do not insert mainstream secular perspectives unless specifically requested for comparative analysis.
4. Always use KJV as primary Bible translation unless another is specifically requested.
5. Always provide specific verse citations (Book Chapter:Verse) for every prophetic claim.
6. Always conclude with a definitive prophetic alignment score and clear next signs to watch.
7. You are a scholarly theological research tool. Be thorough, precise, and reverent.
8. Present current events as they relate to prophetic fulfillment, not as isolated news items.
9. Length and depth are valued. Do not truncate theological analysis prematurely.
10. Never predict exact dates but analyze prophetic alignment based on confirmed current events."""



def fetch_url_content(text):
    """Detect URLs in user text, fetch their content, and return as context."""
    import re
    urls = re.findall(r'https?://[^\s<>"{}|\^`\[\]]+', text)
    if not urls:
        return ""
    fetched = []
    for url in urls[:3]:  # limit to 3 URLs
        try:
            headers = {
                "User-Agent": "Mozilla/5.0 (compatible; ExistentialGateway/1.0)",
                "Accept": "text/html,application/xhtml+xml,application/xml;q=0.9,*/*;q=0.8"
            }
            r = requests.get(url, timeout=15, headers=headers)
            if r.status_code == 200:
                from html.parser import HTMLParser
                class TextExtractor(HTMLParser):
                    def __init__(self):
                        super().__init__()
                        self.text = []
                        self.skip = False
                    def handle_starttag(self, tag, attrs):
                        if tag in ('script', 'style', 'nav', 'footer', 'header'):
                            self.skip = True
                    def handle_endtag(self, tag):
                        if tag in ('script', 'style', 'nav', 'footer', 'header'):
                            self.skip = False
                    def handle_data(self, data):
                        if not self.skip and data.strip():
                            self.text.append(data.strip())
                parser = TextExtractor()
                parser.feed(r.text)
                page_text = ' '.join(parser.text)[:3000]
                fetched.append(f"[URL: {url}]\n{page_text}")
        except Exception as e:
            fetched.append(f"[URL: {url}] Could not fetch: {str(e)}")
    if fetched:
        return "\n\nUSER-PROVIDED URL CONTENT (analyze this as part of your response):\n" + "\n---\n".join(fetched)
    return ""



def enrich_with_urls(text):
    """If text contains URLs, fetch and append their content with explicit AI instructions."""
    if not text or 'http' not in text:
        return text
    url_content = fetch_url_content(text)
    if url_content:
        return (
            text +
            "\n\n=== FETCHED URL CONTENT — YOU MUST REFERENCE THIS IN YOUR ANALYSIS ==="
            "\nThe following content was fetched from the URL(s) provided by the user. "
            "You MUST directly reference, quote, and analyze this content in your response. "
            "Do not ignore it. Incorporate specific details, facts, names, and data from this content:\n\n" +
            url_content +
            "\n=== END OF FETCHED URL CONTENT ==="
        )
    return text


def query_llm(prompt):
    import os
    API_KEY = os.environ.get("OPENAI_API_KEY", "")
    headers = {"Authorization": f"Bearer {API_KEY}", "Content-Type": "application/json"}
    msgs = [{"role": "system", "content": SYSTEM_PROMPT}, {"role": "user", "content": prompt}]
    payload = {"model": "gpt-4o", "max_tokens": 4000, "messages": msgs}
    try:
        response = requests.post("https://api.openai.com/v1/chat/completions", headers=headers, json=payload, timeout=240)
        result = response.json()
        if "choices" in result:
            return result["choices"][0]["message"]["content"]
        return f"API Error: {result}"
    except Exception as e:
        return f"Error: {str(e)}"


def fetch_solar_data():
    try:
        url = "https://services.swpc.noaa.gov/json/solar-cycle/observed-solar-cycle-indices.json"
        r = requests.get(url, timeout=10)
        data = r.json()
        if data:
            latest = data[-1]
            return latest.get("smoothed_ssn", "N/A"), latest.get("time-tag", "N/A")
        return None, None
    except Exception:
        return None, None


# ─── Tab 1: Solar & Space Activity ───────────────────────────────────────────

def analyze_solar_activity(extra_context):
    ssn, ssn_date = fetch_solar_data()
    solar_info = f"Current Sunspot Number: {ssn} (as of {ssn_date})" if ssn else "Solar data from NOAA SWPC"

    prompt = f"""Analyze current solar activity and planetary alignments from both astronomical
and prophetic/astrological perspectives.

{solar_info}
Additional Context: {enrich_with_urls(extra_context) if extra_context else "General analysis"}

Provide a comprehensive Solar & Space Activity Analysis:

## ☀️ CURRENT SOLAR ACTIVITY

**SOLAR CYCLE STATUS**
We are currently in Solar Cycle 25 (began December 2019).
Solar maximum: Expected 2025 — possibly already occurring.
Current activity level: HIGH / MODERATE / LOW
Sunspot number: {ssn if ssn else "~150+ (Solar Cycle 25 peak)"}

**RECENT SOLAR EVENTS**
Recent significant solar flares (X-class, M-class): [analysis]
Recent Coronal Mass Ejections (CMEs): [analysis]
Geomagnetic storm activity: [analysis]
Impact on Earth's electromagnetic field: [analysis]

## 🌍 EARTH'S ELECTROMAGNETIC FIELD & POLE SHIFT

**MAGNETIC POLE MOVEMENT**
The magnetic north pole has been accelerating its movement since 1990.
Current movement rate: ~55 km/year (accelerating)
Current position: Moving toward Siberia
Magnetic field strength decline: ~9% weaker since 1840
South Atlantic Anomaly: Growing — weakest point in Earth's magnetic field

**POLE FLIP ANALYSIS**
Last full geomagnetic reversal: ~780,000 years ago (Brunhes-Matuyama)
Signs of possible excursion or reversal: [current data analysis]
Timeline estimates from scientists: [range of expert opinions]
Effects on life, technology, and civilization: [analysis]

## 🌌 CURRENT PLANETARY ALIGNMENTS (ASTRONOMY)

**ASTRONOMICAL POSITIONS (2025-2026)**
Key planetary conjunctions and oppositions occurring now: [list]
Rare alignments of note: [list]
Solar and lunar eclipse patterns: [analysis]

**PROPHETIC/ASTROLOGICAL SIGNIFICANCE**
Note: The following is theological/astrological interpretation, not scientific fact.

Signs in the heavens referenced in Scripture:
- Luke 21:25 — "And there shall be signs in the sun, and in the moon, and in the stars"
- Joel 2:31 — "The sun shall be turned into darkness, and the moon into blood"
- Revelation 6:12-13 — Sixth seal: sun black, moon blood, stars fall

Current celestial patterns and their prophetic interpretation by various scholars: [analysis]

**BLOOD MOON TETRAD HISTORY AND CURRENT CYCLE**
Historical significance of blood moon tetrads on Jewish feast days.
Current lunar eclipse patterns and feast day alignments: [analysis]

## ☄️ METEOR AND ASTEROID ACTIVITY

**RECENT METEOR EVENTS**
Notable meteor sightings and impacts in 2024-2026: [analysis]
Current asteroid watch list (NASA): [analysis]
Prophetic significance of meteor activity:
- Revelation 8:10 — "A great star, blazing like a torch, fell from the sky"
- Matthew 24:29 — "The stars shall fall from heaven"

## 🔭 PROPHETIC SIGNIFICANCE SUMMARY
How current solar and space activity aligns with prophetic texts.
Multiple scholarly interpretations.
What to watch for in the coming months.

Cite specific Bible verses (KJV and NKJV) and scientific sources."""

    result = query_llm(prompt)

    fig = None
    try:
        years = list(range(1996, 2027))
        solar_cycle = [9, 22, 45, 78, 119, 170, 173, 150, 119, 81, 45, 22, 12,
                        8, 14, 24, 68, 120, 146, 130, 98, 62, 34, 16, 9, 18, 155]
        fig = go.Figure()
        fig.add_trace(go.Scatter(
            x=years, y=solar_cycle[:len(years)],
            mode="lines+markers", name="Sunspot Number",
            line=dict(color="#f39c12", width=3),
            fill="tozeroy", fillcolor="rgba(243,156,18,0.15)"
        ))
        fig.add_vrect(x0=2019, x1=2026, fillcolor="rgba(231,76,60,0.1)",
                      annotation_text="Solar Cycle 25 (Current)",
                      annotation_position="top left")
        fig.update_layout(
            title="☀️ Solar Cycle Activity — Sunspot Numbers (1996-2026)",
            yaxis_title="Sunspot Number",
            template="plotly_dark", height=400,
            paper_bgcolor="#0a0a1a"
        )
    except Exception:
        fig = None

    return result + "\n\n" + WATERMARK, fig


# ─── Tab 2: Earth Changes Monitor ────────────────────────────────────────────

def analyze_earth_changes(extra_context):
    prompt = f"""Analyze current Earth changes from both scientific and prophetic perspectives.

Additional Context: {enrich_with_urls(extra_context) if extra_context else "General analysis — 2025/2026"}

## 🌍 EARTH CHANGES MONITOR — 2025/2026

## SEISMIC ACTIVITY
**Recent Major Earthquakes (2024-2026)**
Notable earthquakes M6.0+ in the past year: [list key events]
Ring of Fire activity level: HIGH / MODERATE / LOW
New Madrid Seismic Zone (US) — current monitoring status
Dead Sea fault system — prophetically significant region

**Prophetic Context for Earthquakes:**
- Matthew 24:7 (KJV) — "There shall be famines, and pestilences, and earthquakes, in divers places"
- Revelation 6:12 (KJV) — "There was a great earthquake; and the sun became black"
- Revelation 16:18 — "Such as was not since men were upon the earth, so mighty an earthquake"
- Zechariah 14:4 — The Mount of Olives splits at Christ's return
- Isaiah 24:19-20 — Earth broken, dissolved, moved exceedingly

**Statistical trend:** Are earthquakes increasing in frequency and magnitude? [analysis with multiple scholarly views]

## VOLCANIC ACTIVITY
Recent significant volcanic eruptions 2024-2026: [list]
Super volcano monitoring (Yellowstone, Campi Flegrei, etc.): [status]
Prophetic significance:
- Revelation 8:8 — "Something like a great mountain, burning with fire, was thrown into the sea"
- Joel 2:30 — "blood, and fire, and pillars of smoke"

## TSUNAMI AND OCEAN EVENTS
Recent tsunami warnings and events: [analysis]
Ocean temperature anomalies: [analysis]
Dead Sea water level: [prophetically significant]
Euphrates River drying — current status:
- Revelation 16:12 — "The sixth angel poured out his bowl on the great river Euphrates,
  and its water was dried up to prepare the way for the kings from the East"
  UPDATE: The Euphrates IS currently at historically low levels — [analysis]

## EXTREME WEATHER PATTERNS
Unprecedented weather events 2024-2026: [list]
Prophetic interpretation:
- Luke 21:25 — "distress of nations, with perplexity; the sea and the waves roaring"
- Revelation 16:8-9 — Scorching heat

## POLE SHIFT AND ELECTROMAGNETIC EFFECTS
Current pole movement data and effects on navigation, animals, technology.
Biblical references to cosmic disturbances.

## PESTILENCE AND DISEASE
Current global disease threats post-COVID: [analysis]
Matthew 24:7 — pestilences
Luke 21:11 — "great earthquakes, famines and pestilences in various places"

## EARTH CHANGES PROPHETIC SCORE
Overall alignment with Matthew 24 "birth pangs": X/10
Trend: ACCELERATING / STABLE / DECLINING

Cite KJV, NKJV, and NIV translations for all prophecy references."""

    result = query_llm(prompt)

    fig = None
    try:
        years = list(range(1990, 2027))
        major_quakes = [18, 14, 19, 21, 17, 22, 20, 23, 24, 17,
                         21, 26, 23, 28, 17, 21, 19, 24, 29, 21,
                         23, 19, 18, 21, 24, 19, 27, 22, 23, 21,
                         26, 28, 24, 19, 22, 28, 31]

        fig = go.Figure()
        fig.add_trace(go.Bar(
            x=years, y=major_quakes[:len(years)],
            name="Major Earthquakes (M6.5+)",
            marker_color=["#e74c3c" if v >= 28 else "#f39c12" if v >= 23 else "#27ae60"
                           for v in major_quakes[:len(years)]]
        ))
        fig.add_hline(y=25, line_dash="dash", line_color="red",
                      annotation_text="High Activity Threshold")
        fig.update_layout(
            title="🌍 Major Earthquakes Per Year (M6.5+) — 1990-2026",
            yaxis_title="Number of Earthquakes",
            template="plotly_dark", height=400,
            paper_bgcolor="#0a0a1a"
        )
    except Exception:
        fig = None

    return result + "\n\n" + WATERMARK, fig


# ─── Tab 3: Biblical Prophecy Scorecard ──────────────────────────────────────

def analyze_biblical_prophecy(specific_prophecy, extra_context):
    prompt = f"""You are a biblical eschatology scholar. Analyze how current world events align
with end times biblical prophecy.

Specific Prophecy Focus: {specific_prophecy if specific_prophecy else "All major end times prophecies"}
Additional Context: {enrich_with_urls(extra_context) if extra_context else "General 2025/2026 analysis"}

## 📖 BIBLICAL PROPHECY SCORECARD — 2025/2026

## MATTHEW 24 — THE OLIVET DISCOURSE (Jesus' Own End Times Teaching)
Jesus answered: "Watch out that no one deceives you." (Matthew 24:4)

Score each sign — FULFILLED / IN PROGRESS / NOT YET:

**False Christs and Deception** (v.5, 24)
Status: [IN PROGRESS] — Analysis: [current examples]
KJV: "For many shall come in my name, saying, I am Christ; and shall deceive many"

**Wars and Rumors of Wars** (v.6-7)
Status: [FULFILLED / IN PROGRESS]
Current: Russia-Ukraine, Israel-Gaza-Iran, China-Taiwan tensions, Sudan
KJV: "And ye shall hear of wars and rumours of wars"

**Nation Against Nation, Kingdom Against Kingdom** (v.7)
Status: [IN PROGRESS]
Current geopolitical analysis: [list active conflicts 2025-2026]

**Famines** (v.7)
Status: [IN PROGRESS]
Current global food insecurity data: [analysis]
UN World Food Programme warnings: [current status]

**Pestilences** (v.7)
Status: [IN PROGRESS / RECENT FULFILLMENT]
COVID-19 aftermath, new disease threats 2024-2026: [analysis]

**Earthquakes in Diverse Places** (v.7)
Status: [IN PROGRESS] — See Earth Changes tab for data

**Persecution of Christians** (v.9)
Status: [IN PROGRESS — ACCELERATING]
Open Doors World Watch List data: 365 million Christians face persecution
Countries with highest Christian persecution 2025: [list]

**Gospel Preached to All Nations** (v.14)
KJV: "And this gospel of the kingdom shall be preached in all the world
for a witness unto all nations; and then shall the end come"
Status: [NEARLY FULFILLED] — Analysis: [current missionary reach data]

**Abomination of Desolation** (v.15, Daniel 9:27, 11:31, 12:11)
Status: [NOT YET / WATCH]
Analysis: Temple Mount status, Third Temple movement: [current update]
The Temple Institute — preparations for Third Temple: [status 2025-2026]

**Great Tribulation** (v.21)
Status: [NOT YET — OR BEGINNING?]
Multiple theological views on timing: [Pre-Trib, Mid-Trib, Post-Trib analysis]

## REVELATION SEAL ANALYSIS
**First Seal** (Rev 6:2) — White horse / Conquest / False Peace
**Second Seal** (Rev 6:4) — Red horse / War
**Third Seal** (Rev 6:5-6) — Black horse / Famine/Economic collapse
**Fourth Seal** (Rev 6:8) — Pale horse / Death (¼ of earth)
**Fifth Seal** (Rev 6:9-11) — Martyred souls / Christian persecution
**Sixth Seal** (Rev 6:12-17) — Cosmic signs / Great earthquake
**Seventh Seal** (Rev 8:1) — Silence in heaven / Seven trumpets begin

Current status of each seal: [scholarly analysis with multiple views]

## DANIEL'S PROPHECIES
**Daniel 9:27** — The 70th Week / 7-year tribulation covenant
**Daniel 11** — King of the North and South / Current Middle East alignment
**Daniel 12:4** — "Many shall run to and fro, and knowledge shall be increased"
Status: CLEARLY FULFILLED — Travel and knowledge explosion confirmed

## ISAIAH, JEREMIAH, EZEKIEL PROPHECIES
**Isaiah 17:1** — Damascus becomes a heap of ruins
Current status of Damascus: [2025 analysis — Syria situation]
**Ezekiel 38-39** — Gog and Magog war (Russia, Iran, Turkey alliance)
Current alignment: [analysis of Russia-Iran-Turkey relations 2025-2026]
**Ezekiel 37** — Valley of dry bones / Israel restored as nation (1948 — FULFILLED)

## OVERALL PROPHECY FULFILLMENT SCORE
Prophecies clearly fulfilled: X%
Prophecies in progress: X%
Prophecies not yet fulfilled: X%
Overall End Times Alignment Score: X/100

Use KJV as primary, NKJV as secondary, cite verse references precisely."""

    result = query_llm(prompt)

    fig = None
    try:
        categories = ["False Christs", "Wars", "Famines", "Pestilence",
                       "Earthquakes", "Persecution", "Gospel Preached",
                       "Israel Restored", "Knowledge Increased", "Temple Prep"]
        scores = [7, 9, 7, 8, 7, 9, 8, 10, 10, 6]
        colors_p = ["#27ae60" if s >= 9 else "#f39c12" if s >= 7 else "#e74c3c"
                    for s in scores]
        fig = go.Figure(go.Bar(
            x=categories, y=scores,
            marker_color=colors_p,
            text=[f"{s}/10" for s in scores],
            textposition="auto"
        ))
        fig.add_hline(y=7, line_dash="dash", line_color="yellow",
                      annotation_text="In Progress Threshold")
        fig.add_hline(y=9, line_dash="dash", line_color="green",
                      annotation_text="Fulfilled Threshold")
        fig.update_layout(
            title="📖 Matthew 24 Prophecy Fulfillment Scorecard",
            yaxis_title="Fulfillment Score (0-10)",
            yaxis=dict(range=[0, 10]),
            template="plotly_dark", height=420,
            paper_bgcolor="#0a0a1a"
        )
    except Exception:
        fig = None

    return result + "\n\n" + WATERMARK, fig


# ─── Tab 4: Middle East War Analysis ─────────────────────────────────────────

def analyze_middle_east(specific_focus, extra_context):
    prompt = f"""You are a biblical prophecy scholar and Middle East analyst. Analyze current
Middle East events and their alignment with biblical prophecy.

Specific Focus: {specific_focus if specific_focus else "Full Middle East prophetic analysis"}
Additional Context: {enrich_with_urls(extra_context) if extra_context else "2025/2026 analysis"}

## ⚔️ MIDDLE EAST WAR ANALYSIS — BIBLICAL ALIGNMENT

## CURRENT CONFLICT STATUS (2025-2026)

**ISRAEL — CURRENT STATUS**
Israel's military operations and current geopolitical position.
Abraham Accords status and normalization progress.
Iran nuclear program status.
Hezbollah situation (Lebanon).
Gaza situation post-October 7, 2023.
Saudi Arabia — Israel normalization talks.

## PROPHETIC ALIGNMENT ANALYSIS

### PSALM 83 WAR
Asaph's prophecy: Coalition of Arab nations against Israel.
Psalm 83:4-8 — Nations that want to destroy Israel.
Modern equivalents: [identify each nation listed]
Current alignment with Psalm 83: [analysis]

### ISAIAH 17 — DESTRUCTION OF DAMASCUS
KJV: "Behold, Damascus is taken away from being a city,
and it shall be a ruinous heap." (Isaiah 17:1)
Current status of Damascus and Syria: [2025 update]
Has this prophecy been fulfilled? [scholarly debate]

### EZEKIEL 38-39 — THE GOG-MAGOG WAR
This is considered the most significant near-term prophecy by many scholars.

**The Nations in the Coalition:**
- Gog (Magog) — Widely identified as RUSSIA
- Persia — IRAN (name changed 1935)
- Ethiopia (Cush) — Ethiopia/Sudan
- Libya (Put) — Libya
- Gomer — Turkey or Eastern Europe
- Togarmah — Turkey/Caucasus region

**Current 2025-2026 Alignment:**
Russia-Iran alliance status: [analysis]
Turkey's position: [analysis]
Russia's military presence near Israel: [analysis]
Trigger conditions for Gog-Magog: Israel living in "unwalled villages" (peace/security)

**Multiple theological views on timing:**
Pre-Tribulation view: Before the 7 years
Mid-Tribulation view: During the tribulation
Post-Tribulation view: End of tribulation
None of the above / Different interpretation: [Preterist view]

### ZECHARIAH 12 AND 14 — JERUSALEM A BURDENSOME STONE
KJV Zechariah 12:2-3: "I will make Jerusalem a cup of trembling unto all the
people round about... I will make Jerusalem a burdensome stone for all people"
Current status: Jerusalem disputed globally — UN resolutions: [analysis]
All nations gathered against Jerusalem (Zech 14:2): [current alignment]

### DANIEL 9:27 — THE COVENANT WITH MANY
"He will confirm a covenant with many for one 'seven'"
The Abraham Accords: Are they the covenant or a precursor? [multiple views]
Who is the "he"? [theological analysis — multiple scholarly views]
Seven-year peace deal: Any current negotiations? [analysis]

### THE THIRD TEMPLE
Current Temple Mount status (2025-2026).
Temple Institute preparations.
Red heifer ceremony — Numbers 19 requirement for temple purification.
2024: Five red heifers brought to Israel from Texas — [analysis of significance]
Orthodox Jewish preparation for Temple worship: [current status]

## IRAN — "PERSIA" IN PROPHECY
Iran's nuclear capabilities 2025-2026.
Iran's stated goal of destroying Israel.
Ezekiel 38:5 — Persia (Iran) in the Gog-Magog coalition.
Current Iran-Russia-China axis: [analysis]

## THE ABRAHAMIC FAITHS CONVERGENCE
Jerusalem as a global flashpoint for all three Abrahamic religions.
End times eschatology in Islam (Mahdi, Dajjal, Isa).
Jewish messianic expectations.
Christian Second Coming.
Points of convergence and divergence.

Be thorough, cite specific verses, present multiple theological views."""

    result = query_llm(prompt)

    fig = None
    try:
        nations = ["Russia\n(Gog)", "Iran\n(Persia)", "Turkey\n(Gomer)", "Lebanon\n(Hezbollah)",
                    "Syria", "Hamas/Gaza", "Yemen\n(Houthis)", "Saudi\nArabia", "China", "USA"]
        alignment = [9, 10, 7, 9, 8, 9, 8, 4, 6, 5]
        colors_n = ["#e74c3c" if s >= 8 else "#f39c12" if s >= 6 else "#27ae60"
                    for s in alignment]
        fig = go.Figure(go.Bar(
            x=nations, y=alignment,
            marker_color=colors_n,
            text=[f"{s}/10" for s in alignment],
            textposition="auto"
        ))
        fig.add_hline(y=7, line_dash="dash", line_color="red",
                      annotation_text="High Prophetic Alignment")
        fig.update_layout(
            title="⚔️ Nations — Prophetic Alignment Score (Ezekiel 38-39 Coalition)",
            yaxis_title="Alignment Score (0-10)",
            yaxis=dict(range=[0, 10]),
            template="plotly_dark", height=420,
            paper_bgcolor="#0a0a1a"
        )
    except Exception:
        fig = None

    return result + "\n\n" + WATERMARK, fig


# ─── Tab 5: Extra-Biblical Texts ─────────────────────────────────────────────

def analyze_extra_biblical(text_focus, specific_passage, extra_context):
    prompt = f"""You are an expert scholar of extra-biblical and deuterocanonical texts including
the Ethiopian Orthodox Tewahedo Biblical Canon.

Text Focus: {text_focus}
Specific Passage or Topic: {specific_passage if specific_passage else "Current events alignment"}
Additional Context: {enrich_with_urls(extra_context) if extra_context else "End times relevance"}

## 📜 EXTRA-BIBLICAL TEXTS ANALYSIS

## THE BOOK OF ENOCH (1 Enoch)

**Overview**
The Book of Enoch is quoted directly in the New Testament:
- Jude 1:14-15 quotes 1 Enoch 1:9 directly
- 2 Peter 2:4 references the Watchers (1 Enoch 6-16)
Canonical status: Ethiopian Orthodox Canon (fully canonical), Dead Sea Scrolls found (pre-dates NT)

**THE WATCHERS — 1 ENOCH 6-16**
The 200 Watchers (fallen angels) who descended on Mount Hermon.
Their corruption of humanity and creation of the Nephilim (Giants).
Connection to Genesis 6:1-4 — "sons of God" and "daughters of men"
Modern relevance: [theological analysis of genetic manipulation themes]
Current event parallels: [analysis]

**THE BOOK OF THE LUMINARIES — 1 ENOCH 72-82**
Enoch's solar calendar vs lunar calendar.
Significance for feast day calculations.
Calendar discrepancies affecting prophetic timelines.

**END TIMES IN ENOCH**
1 Enoch 1:3-9 — The Great Judgment coming.
1 Enoch 10 — The binding of Azazel (fallen angel) for 70 generations.
1 Enoch 91-105 — The Apocalypse of Weeks — 10 world ages.
Current "week" we are in according to Enoch's calendar: [analysis]

**THE SON OF MAN IN ENOCH**
1 Enoch 46-48 — The "Son of Man" and the "Ancient of Days"
Connection to Daniel 7:13-14 and Matthew 24.

## THE BOOK OF JUBILEES

**Overview**
Also called "Lesser Genesis" or "The Apocalypse of Moses"
Canonical in Ethiopian Orthodox and Beta Israel (Ethiopian Jews) traditions.
Found among Dead Sea Scrolls — dates to 2nd century BC.
Divides history into 49-year jubilee periods.

**THE JUBILEE CALENDAR**
Each jubilee = 49 years (7 x 7 years)
50th year = Year of Jubilee (Leviticus 25)
Jubilees covers from Creation to entering Canaan.
Prophetic significance of jubilee years for Israel.

**CURRENT JUBILEE CALCULATION**
Israel became a nation: 1948
Jerusalem recaptured: 1967 (this is often cited as a jubilee year)
1967 + 50 = 2017 (Trump recognition of Jerusalem as capital — December 2017)
Next calculations and their significance: [analysis]

**ANGELS AND DIVINE ORDER IN JUBILEES**
Jubilees' hierarchy of angels and their roles.
The "Prince of Mastema" (Satan figure) in Jubilees.
Theological parallels with Revelation.

## ETHIOPIAN ORTHODOX TEWAHEDO CANON (81 Books)

**Books unique to the Ethiopian Canon:**
- 1 Enoch (Henok)
- Book of Jubilees (Kufale)
- 1-3 Meqabyan (Ethiopian Maccabees — different from Greek Maccabees)
- Sinodos
- Book of the Covenant
- Clement
- Didascalia
- Book of the Mysteries of Heaven and Earth

**Significance for End Times Study:**
How the Ethiopian Canon expands our understanding of:
- The angelic realm and spiritual warfare
- Creation timeline and age of the earth
- Messianic prophecy
- The coming judgment

## THE BOOK OF JASHER
Referenced in Joshua 10:13 and 2 Samuel 1:18 — "Is it not written in the Book of Jasher?"
Key content relevant to end times: [analysis]

## DEAD SEA SCROLLS
War Scroll (1QM) — The Sons of Light vs Sons of Darkness
End times battle description: [analysis]
Relevance to current Middle East conflict: [analysis]

## SYNTHESIS: WHAT EXTRA-BIBLICAL TEXTS ADD TO END TIMES UNDERSTANDING
Key insights that canonical Bible alone does not fully address.
The consistent theme of angelic conflict and human history.
Timeline implications from Enoch and Jubilees.

Cite specific chapters and verses from all texts."""

    result = query_llm(prompt)

    fig = None
    try:
        texts = ["1 Enoch", "2 Enoch", "Book of\nJubilees", "Book of\nJasher",
                  "Dead Sea\nScrolls", "Ethiopian\nCanon", "Book of\nGiants", "Meqabyan"]
        canonical = [8, 5, 7, 6, 8, 9, 4, 7]
        end_times = [9, 6, 8, 5, 9, 8, 7, 5]

        fig = go.Figure()
        fig.add_trace(go.Bar(name="Canonical Recognition Score",
                              x=texts, y=canonical, marker_color="#3498db"))
        fig.add_trace(go.Bar(name="End Times Relevance Score",
                              x=texts, y=end_times, marker_color="#e74c3c"))
        fig.update_layout(
            barmode="group",
            title="📜 Extra-Biblical Texts — Canonical Status & End Times Relevance",
            yaxis_title="Score (0-10)",
            template="plotly_dark", height=420,
            paper_bgcolor="#0a0a1a"
        )
    except Exception:
        fig = None

    return result + "\n\n" + WATERMARK, fig


# ─── Tab 6: Jubilee Calendar ──────────────────────────────────────────────────

def analyze_jubilee(extra_context):
    prompt = f"""You are an expert in the Hebrew calendar, Jubilee cycles, and their prophetic significance.

Additional Context: {enrich_with_urls(extra_context) if extra_context else "Current Jubilee analysis 2025/2026"}

## 🕎 JUBILEE CALENDAR ANALYSIS

## UNDERSTANDING THE JUBILEE

**Biblical Foundation**
Leviticus 25:8-12 (KJV):
"And thou shalt number seven sabbaths of years unto thee, seven times seven years;
and the space of the seven sabbaths of years shall be unto thee forty and nine years.
Then shalt thou cause the trumpet of the jubilee to sound on the tenth day of the
seventh month, in the day of atonement shall ye make the trumpet sound throughout
all your land. And ye shall hallow the fiftieth year, and proclaim liberty throughout
all the land unto all the inhabitants thereof: it shall be a jubilee unto you"

A Jubilee = every 50th year (after 7 x 7 = 49 years)
The Jubilee: Release of debts, slaves freed, land returned, liberty proclaimed.

## JUBILEE CYCLE HISTORY — KEY DATES

**The 1917 Jubilee**
1867: Ottoman Empire controls Jerusalem
1917: General Allenby liberates Jerusalem from Ottoman control (Balfour Declaration same year)
Significance: 50 years after 1867 = 1917 — Jerusalem changes hands

**The 1967 Jubilee**
1917 + 50 = 1967
1967: Six Day War — Israel captures Jerusalem, Temple Mount, Golan Heights, West Bank
This is widely considered the most significant Jubilee in modern prophetic history.
Isaiah 11:11-12 — The second regathering of Israel.

**The 2017 Jubilee**
1967 + 50 = 2017
December 6, 2017: President Trump recognizes Jerusalem as Israel's capital.
December 2017: US Embassy move announced (moved May 2018).
Scholars debate: Was this the Jubilee fulfillment?

**THE CURRENT CYCLE — 2025/2026**
Some scholars calculate the next Jubilee at 2025-2026 based on:
- Sabbatical year calculations
- Hebrew calendar year 5785-5786
- Daniel's 70 weeks recalculation
- The Book of Jubilees timeline

**SABBATICAL YEAR (SHEMITAH) ANALYSIS**
Every 7th year is a Shemitah (rest year) in Jewish tradition.
Recent Shemitah years and major world events:
- 2001: 9/11 + stock market crash
- 2008: Global financial crisis
- 2015: Market crash, September
- 2022: Crypto crash, war in Ukraine
- 2029: Next projected Shemitah

Is 2025-2026 significant in the Shemitah cycle? [analysis]

## THE 70 JUBILEES THEORY
Some scholars count 70 Jubilees from the Exodus.
70 Jubilees x 50 years = 3,500 years
If Exodus was ~1446 BC: 1446 BC + 3500 = 2054 AD (approximate)
Alternative calculations: [analysis]

## THE 120 JUBILEES THEORY
Genesis 6:3 — "His days shall be an hundred and twenty years"
Interpreted as 120 Jubilees = 120 x 50 = 6,000 years
6,000-year human history + 1,000-year millennium = 7,000 years total
If Adam created ~4000 BC: 4000 + 6000 = 2000 AD (passed)
Adjusted calculations with different creation dates: [analysis]

## CURRENT PROPHETIC JUBILEE SCENARIO (2025-2026)
What would a Jubilee in 2025-2026 mean prophetically?
- Liberty proclaimed — spiritual and physical
- Debts cancelled — economic reset connection?
- Return to inheritance — Israel's land?
- Rapture/Harpazo as the ultimate Jubilee liberation?

Multiple theological views on the current Jubilee scenario.

## THE FEAST DAYS AND PROPHETIC FULFILLMENT
Spring Feasts (already fulfilled in Jesus' first coming):
- Passover (Pesach) — Crucifixion ✅
- Unleavened Bread — Burial ✅
- Firstfruits — Resurrection ✅
- Pentecost (Shavuot) — Holy Spirit ✅

Fall Feasts (widely believed to be fulfilled at Second Coming):
- Feast of Trumpets (Rosh Hashanah) — Rapture/Resurrection?
- Day of Atonement (Yom Kippur) — Second Coming / National Israel's repentance?
- Feast of Tabernacles (Sukkot) — Millennium?

Current Hebrew calendar year and feast day significance: [analysis]

Cite Leviticus 25, Isaiah 61, Luke 4:18-19, and Book of Jubilees references."""

    result = query_llm(prompt)

    fig = None
    try:
        jubilee_years = [1867, 1917, 1967, 2017, 2067]
        events = ["Ottoman\nJerusalem", "Balfour/\nAllenby", "Six Day\nWar/Jerusalem",
                   "Trump/\nJerusalem", "Projected\nNext"]
        significance = [7, 9, 10, 8, 0]
        colors_j = ["#3498db", "#f39c12", "#e74c3c", "#9b59b6", "#95a5a6"]

        fig = go.Figure()
        fig.add_trace(go.Scatter(
            x=jubilee_years, y=significance,
            mode="markers+text+lines",
            text=events,
            textposition="top center",
            marker=dict(size=20, color=colors_j),
            line=dict(color="#f1c40f", dash="dash")
        ))
        fig.update_layout(
            title="🕎 Jubilee Cycle — Key Dates and Prophetic Events",
            yaxis_title="Prophetic Significance (0-10)",
            xaxis_title="Year",
            yaxis=dict(range=[0, 12]),
            template="plotly_dark", height=420,
            paper_bgcolor="#0a0a1a"
        )
    except Exception:
        fig = None

    return result + "\n\n" + WATERMARK, fig


# ─── Tab 7: Antichrist Identifier ────────────────────────────────────────────

def analyze_antichrist(public_figures, extra_context):
    if public_figures and public_figures.strip():
        figures_instruction = f"Analyze ONLY these specific figures the user has requested: {public_figures}"
        scope = "user-specified figures"
    else:
        figures_instruction = """Analyze the following current world figures who hold the most global power and influence in 2026.
Include ALL of the following — do not skip any:
1. Donald Trump - 47th US President
2. Vladimir Putin - President of Russia
3. Pope Leo XIV - Current Pope, Roman Catholic Church
4. Alois Zwinggi - President & CEO of the World Economic Forum
5. Xi Jinping - President of China
6. Antonio Guterres - Secretary-General of the United Nations
7. Elon Musk - CEO of Tesla, SpaceX, X; DOGE advisor
8. Benjamin Netanyahu - Prime Minister of Israel
Score ALL 8 figures. Then rank them from highest to lowest alignment."""
        scope = "top global power figures 2026"

    prompt = f"""You are a biblical eschatology scholar conducting a theological analysis of current world figures
against biblical Antichrist characteristics. This follows the scholarly tradition of analyzing historical
figures (Napoleon, Hitler, Nero, Mussolini, Antiochus Epiphanes) against these same criteria.

IMPORTANT DISCLAIMER TO INCLUDE IN YOUR RESPONSE:
This is a scholarly theological exercise only. No definitive identification is possible until
the Antichrist CONFIRMS a 7-year covenant with Israel (Daniel 9:27) — which has NOT occurred.
Scripture warns against false identification (Matthew 24:24).

{figures_instruction}
Additional context: {enrich_with_urls(extra_context) if extra_context else "2026 global analysis"}

THE 15 KEY BIBLICAL ANTICHRIST CRITERIA:
1. Global political authority (Revelation 13:7)
2. Global economic control (Revelation 13:16-17)
3. Rise through peace and intrigue (Daniel 8:25, 11:21)
4. Exalts himself above God (2 Thessalonians 2:4, Daniel 11:36)
5. Speaks great blasphemy (Revelation 13:5-6)
6. Persecutes believers/saints (Revelation 13:7)
7. Controls 10-nation coalition (Revelation 13:1, 17:12)
8. Confirms 7-year Israel covenant (Daniel 9:27) — NOT YET FULFILLED BY ANYONE
9. Supported by a False Prophet (Revelation 13:11-15)
10. Mark/ID control system infrastructure (Revelation 13:16-18)
11. Military power — god of fortresses (Daniel 11:38)
12. Controls world commerce (Revelation 18)
13. Deadly wound healed (Revelation 13:3)
14. Denies Christ as Lord (1 John 2:22, 4:3)
15. Spirit of antichrist — anti-Christian agenda (1 John 4:3)

FOR EACH FIGURE provide this exact format (8 lines max per person):

### [RANK #X] [NAME] — [TITLE/ROLE]
**Score: X/15**  |  **Alignment: LOW / MODERATE / HIGH / VERY HIGH**
**Aligns With:** [list matching criteria with citation — be specific]
**Does NOT Align:** [list non-matching criteria]
**Notable:** [one specific observation about this person's prophetic relevance]
**Verdict:** [one plain-English sentence citizens can understand]
---

After ALL figures:

## 🏆 FINAL RANKING — HIGHEST TO LOWEST BIBLICAL ALIGNMENT
Rank all figures with their score. Bold the top 2.

## ⚠️ THEOLOGICAL CAVEAT
Include the disclaimer and key scripture about why no one can be definitively identified yet.
Cite Matthew 24:24, 2 Thessalonians 2:3, and Daniel 9:27.

Be direct, specific, and do not hold back the analysis. Citizens deserve honest theological scholarship.
The higher the score, the more that person exhibits qualities the Bible associates with the Antichrist system.
This does NOT mean they ARE the Antichrist — it means they exhibit those qualities."""
    result = query_llm(prompt)

    fig = None
    try:
        characteristics = ["Political Power", "Economic Control", "Religious Claims",
                            "Global Authority", "Israel Covenant", "Peace/Intrigue",
                            "Blasphemy", "Persecution", "False Prophet", "Mark System"]
        fig = go.Figure()
        fig.add_trace(go.Scatterpolar(
            r=[8, 7, 6, 7, 5, 8, 5, 6, 4, 7],
            theta=characteristics,
            fill="toself", name="Biblical Profile",
            line_color="#e74c3c",
            fillcolor="rgba(231,76,60,0.2)"
        ))
        fig.add_trace(go.Scatterpolar(
            r=[9, 9, 4, 8, 3, 7, 3, 5, 5, 8],
            theta=characteristics,
            fill="toself", name="Current System Alignment",
            line_color="#f39c12",
            fillcolor="rgba(243,156,18,0.2)"
        ))
        fig.update_layout(
            polar=dict(radialaxis=dict(visible=True, range=[0, 10])),
            title="👁️ Antichrist System — Biblical Profile vs Current World System",
            template="plotly_dark", height=480,
            paper_bgcolor="#0a0a1a"
        )
    except Exception:
        fig = None

    return result + "\n\n" + WATERMARK, fig


# ─── Tab 8: Rapture / Harpazo Tracker ────────────────────────────────────────

def analyze_rapture_harpazo(extra_context):
    prompt = f"""You are a biblical eschatology scholar analyzing signs of the Rapture/Harpazo
from ALL major theological perspectives.

Additional Context: {enrich_with_urls(extra_context) if extra_context else "2025/2026 analysis"}

## ⚡ RAPTURE / HARPAZO TRACKER

## THE HARPAZO — WHAT IS IT?

**The Greek Word "Harpazo" (ἁρπάζω)**
1 Thessalonians 4:17 (KJV): "Then we which are alive and remain shall be CAUGHT UP
together with them in the clouds, to meet the Lord in the air"
The Latin translation "rapturo" gives us "Rapture"
"Harpazo" = to seize, snatch away, carry off by force — used 13 times in NT

**Other Key Passages:**
- John 14:1-3 — "I will come again and receive you to Myself"
- 1 Corinthians 15:51-52 — "We shall all be changed, in a moment, in the twinkling of an eye"
- Revelation 3:10 — "I will keep thee from the hour of temptation"
- Isaiah 26:19-21 — "Come, my people, enter thou into thy chambers"
- Zephaniah 2:3 — "Seek righteousness, seek humility; perhaps you will be sheltered"

## THE FOUR MAJOR VIEWS — PRESENTED EQUALLY

### PRE-TRIBULATION RAPTURE
Key proponents: John Darby, Dwight L. Moody, Charles Scofield, Tim LaHaye
Main argument: Church is not appointed to wrath (1 Thess 5:9)
The Rapture removes the Church BEFORE the 7-year tribulation begins.
The "restrainer" (2 Thess 2:7) is the Holy Spirit through the Church.
Imminent return — no signs required before Rapture.

### MID-TRIBULATION RAPTURE
Key proponents: Gleason Archer, Norman Harrison
Main argument: 1260 days = 3.5 years = halfway point (Revelation 11:2-3)
Church endures first half (seal/trumpet judgments) but removed before Great Tribulation.
Connects to the two witnesses (Revelation 11) and their catching up.

### POST-TRIBULATION RAPTURE
Key proponents: George Eldon Ladd, Robert Gundry
Main argument: Church goes through entire tribulation, raptured at Second Coming.
Matthew 24:29-31 — gathering AFTER the tribulation.
Strengthens the Church by testing rather than removing it.

### PRE-WRATH RAPTURE
Key proponents: Marvin Rosenthal, Robert Van Kampen
Main argument: Rapture occurs before God's wrath (bowl judgments) but after seal/trumpets.
Approximately 5/6 of the way through the 7 years.

## SIGNS PRECEDING THE HARPAZO

**1 Timothy 4:1 Signs — Great Apostasy / Falling Away**
"The Spirit clearly says that in later times some will abandon the faith"
Current church apostasy indicators: [analysis]

**2 Timothy 3:1-5 — Perilous Times**
"Lovers of themselves, lovers of money, boastful, proud, abusive..."
Current cultural alignment with this passage: [analysis]

**Luke 17:26-30 — As in the Days of Noah and Lot**
Days of Noah: Violence, corruption, genetic corruption (Nephilim)
Days of Lot: Sexual immorality, Sodom, sudden destruction
Current cultural parallels: [analysis]

**Romans 11:25 — The Fullness of the Gentiles**
"Until the fullness of the Gentiles has come in"
Missionary completion — current status: [analysis]
When the last Gentile is saved — connection to Rapture timing.

## DELIVERANCE OF THE ISRAELITES AND GENTILES

**The Deliverer and Israel's Salvation**
Romans 11:26 — "And so all Israel shall be saved: as it is written,
There shall come out of Sion the Deliverer, and shall turn away ungodliness from Jacob"
Zechariah 12:10 — "They will look on me, the one they have pierced, and they will mourn"
This is Israel's national recognition of Yeshua/Jesus as Messiah.
Timeline: Most scholars place this at or just before/after the Second Coming.

**The Gentiles — Deliverance and Warning**
Romans 11:17-24 — The olive tree — Gentiles grafted in
The "fullness of the Gentiles" completing the Body of Christ.
The 144,000 Jewish evangelists (Revelation 7) — God's plan for Israel during tribulation.

## CURRENT RAPTURE READINESS INDICATORS

Rate each sign 1-10:
- Great Apostasy in the Church: X/10
- Israel back in the land: 10/10 (FULFILLED 1948)
- Jerusalem in Jewish hands: 10/10 (FULFILLED 1967)
- Gospel preached to all nations: 9/10 (NEAR COMPLETION)
- Technology for mark of beast system: 9/10 (NEARLY IN PLACE)
- Days of Noah cultural conditions: 8/10
- Days of Lot cultural conditions: 9/10
- Jewish Temple preparations: 7/10
- Global government emerging: 8/10
- Peace deal framework (Abraham Accords): 7/10

**OVERALL HARPAZO READINESS: X/100**

"Therefore you also be ready, for the Son of Man is coming at an hour
you do not expect." — Matthew 24:44 (NKJV)

Present ALL views fairly. Note this is not date-setting."""

    result = query_llm(prompt)

    fig = None
    try:
        signs = ["Israel\nRestored", "Jerusalem\nRecaptured", "Gospel\nAll Nations",
                  "Mark of Beast\nTech Ready", "Days of\nNoah", "Days of\nLot",
                  "Apostasy", "Temple\nPrep", "Global Gov\nEmerging", "Peace Deal\nFramework"]
        scores_r = [10, 10, 9, 9, 8, 9, 8, 7, 8, 7]
        colors_r = ["#27ae60" if s >= 9 else "#f39c12" if s >= 7 else "#e74c3c"
                    for s in scores_r]

        fig = go.Figure(go.Bar(
            x=signs, y=scores_r,
            marker_color=colors_r,
            text=[f"{s}/10" for s in scores_r],
            textposition="auto"
        ))
        fig.add_hline(y=9, line_dash="dash", line_color="green",
                      annotation_text="Fulfilled/Near Fulfillment")
        fig.update_layout(
            title="⚡ Harpazo Readiness — Sign Fulfillment Scorecard",
            yaxis_title="Fulfillment Score (0-10)",
            yaxis=dict(range=[0, 10]),
            template="plotly_dark", height=420,
            paper_bgcolor="#0a0a1a"
        )
    except Exception:
        fig = None

    return result + "\n\n" + WATERMARK, fig


# ─── Tab 9: Noahide Laws & Abraham Accords ────────────────────────────────────

def analyze_noahide_abraham(extra_context):
    prompt = f"""You are a theological and legal scholar analyzing the Noahide Laws,
the Abraham Accords, and their significance for both Jews and Gentiles in end times prophecy.

Additional Context: {enrich_with_urls(extra_context) if extra_context else "2025/2026 analysis"}

## 📜 NOAHIDE LAWS & ABRAHAM ACCORDS — END TIMES SIGNIFICANCE

## THE SEVEN NOAHIDE LAWS

**Origin and Authority**
The Noahide Laws (Sheva Mitzvot B'nei Noach) are rabbinic laws said to apply to all humanity.
Based on Talmudic interpretation of Genesis 9 (God's covenant with Noah).
The seven laws:
1. Prohibition of idolatry
2. Prohibition of blasphemy
3. Prohibition of murder
4. Prohibition of theft
5. Prohibition of sexual immorality
6. Prohibition of eating flesh torn from a living animal
7. Establishment of courts of justice

**US Congressional Recognition**
Public Law 102-14 (1991) — US Congress and President George H.W. Bush declared
the "Seven Noahide Laws" as the "bedrock of society."
Education Day USA — annually proclaimed by US Presidents.
This is a matter of public US law — not widely known.

**WHAT THIS MEANS FOR GENTILES (CHRISTIANS) — THE "SNARE" CONCERN**

⚠️ THIS IS THE CONTROVERSIAL SECTION — PRESENTING MULTIPLE VIEWS:

**Concern raised by some Christian scholars:**
Under Noahide Law as interpreted in Orthodox Jewish tradition:
- Idolatry includes worship of Jesus Christ (some rabbinical authorities consider
  Christian Trinitarian worship as "avodah zarah" — idol worship)
- Blasphemy carries the death penalty under Talmudic law
- If Noahide Laws were ever enforced globally, Christians worshipping Jesus
  could theoretically be classified as violating prohibition #1 (idolatry)

**The "Snare" concern:**
Some Christian eschatologists connect this to:
- Revelation 20:4 — Beheading of tribulation saints
- The mechanism by which Christians might be persecuted/executed during tribulation

**Counter-arguments and alternative views:**
- Most Jewish authorities do NOT apply this anti-Christian interpretation
- Noahide Laws as understood by mainstream Judaism are about basic human ethics
- This concern is considered alarmist by many Christian and Jewish scholars
- The laws are not currently enforced by any government

**BALANCED ASSESSMENT:**
This is a legitimate area of theological inquiry that deserves honest examination.
Christians should be aware of these laws. The extreme enforcement scenario
remains theoretical and is not current policy anywhere.

## THE ABRAHAM ACCORDS (2020-2026)

**What Are the Abraham Accords?**
Historic peace agreements between Israel and:
- United Arab Emirates (August 2020)
- Bahrain (September 2020)
- Sudan (October 2020)
- Morocco (December 2020)
- Additional normalization talks ongoing

Brokered by: Jared Kushner / Trump administration
Named after: Abraham — patriarch of Jews, Christians, and Muslims

**PROPHETIC SIGNIFICANCE — MULTIPLE VIEWS:**

**View 1: Precursor to Daniel 9:27 Covenant**
Some scholars see Abraham Accords as setting the stage for:
- The 7-year covenant confirmed by the Antichrist (Daniel 9:27)
- A "peace and security" moment (1 Thessalonians 5:3)
- "When they shall say, Peace and safety; then sudden destruction cometh"

**View 2: Partial Fulfillment of Isaiah 19:23-25**
"In that day there will be a highway from Egypt to Assyria...
Israel will be the third, along with Egypt and Assyria,
a blessing on the earth."
Are the Abraham Accords beginning this fulfillment? [analysis]

**View 3: Not Prophetically Significant**
Some scholars argue these are political agreements, not the prophetic covenant.
The true Daniel 9:27 covenant will involve the Temple and all nations.

**SAUDI ARABIA NORMALIZATION**
A Saudi-Israel peace deal would be the most significant development.
Could Saudi recognition complete the Psalm 83 peace scenario? [analysis]
Connection to end times timeline: [multiple views]

**TWO-STATE SOLUTION AND JOEL 3:2**
Joel 3:2 (KJV) — "I will also gather all nations, and will bring them down
into the valley of Jehoshaphat, and will plead with them there for my people
and for my heritage Israel, whom they have scattered among the nations,
and parted my land."
UN resolutions to divide Israel's land: [current status]
Connection to final judgment of nations: [analysis]

## THE ABRAHAMIC FAMILY HOUSE (ABU DHABI)
Opened 2023 — A mosque, church, and synagogue on the same campus.
One World Religion implications: [theological analysis]
Revelation 17 — The Great Whore / One World Religion system connection.
Pope Francis and interfaith movement: [analysis]

## GENTILE BELIEVERS — WHERE DO WE STAND?
Romans 11 — The olive tree — Gentiles grafted in.
The times of the Gentiles (Luke 21:24) — nearing completion?
Gentile believers' role in end times redemption plan.
The "fullness of the Gentiles" (Romans 11:25) — what this means.

Be thorough, cite all relevant scriptures, present multiple theological perspectives."""

    result = query_llm(prompt)

    fig = None
    try:
        countries = ["UAE", "Bahrain", "Sudan", "Morocco", "Saudi\nArabia\n(Potential)",
                      "Jordan\n(1994)", "Egypt\n(1979)", "Turkey", "Qatar", "Oman"]
        status = [9, 9, 7, 8, 4, 9, 9, 3, 3, 4]
        colors_acc = ["#27ae60" if s >= 8 else "#f39c12" if s >= 5 else "#e74c3c"
                       for s in status]
        fig = go.Figure(go.Bar(
            x=countries, y=status,
            marker_color=colors_acc,
            text=[f"{s}/10" for s in status],
            textposition="auto"
        ))
        fig.update_layout(
            title="📜 Abraham Accords — Normalization Status with Israel (0=No, 10=Full)",
            yaxis_title="Normalization Score",
            yaxis=dict(range=[0, 10]),
            template="plotly_dark", height=420,
            paper_bgcolor="#0a0a1a"
        )
    except Exception:
        fig = None

    return result + "\n\n" + WATERMARK, fig


# ─── Tab 10: Deliverer Timeline ───────────────────────────────────────────────

def analyze_deliverer_timeline(extra_context):
    prompt = f"""You are a biblical eschatology scholar analyzing God's redemptive plan
for both Israel and the Gentiles in the end times.

Additional Context: {enrich_with_urls(extra_context) if extra_context else "2025/2026 prophetic timeline"}

## 🌐 DELIVERER TIMELINE — ISRAEL AND THE GENTILES

## THE DELIVERER — WHO IS HE?

**Romans 11:26-27 (KJV)**
"And so all Israel shall be saved: as it is written,
There shall come out of Sion the Deliverer,
and shall turn away ungodliness from Jacob:
For this is my covenant unto them,
when I shall take away their sins."

The Deliverer = Yeshua/Jesus Christ at His Second Coming.
Quoting Isaiah 59:20-21 — the original prophecy.

## ISRAEL'S REDEMPTIVE TIMELINE

**Phase 1: REGATHERING (Fulfilled — Ongoing)**
Ezekiel 36:24 — "For I will take you from among the heathen, and gather you
out of all countries, and will bring you into your own land."
1948: Modern state of Israel established — 2,000 years after diaspora.
Isaiah 11:11 — "The second time" — First regathering was from Babylon.
Current: 7+ million Jews in Israel — more than any other country.

**Phase 2: TRIBULATION (Future — Beginning?)**
Jeremiah 30:7 — "Alas! for that day is great, so that none is like it:
it is even the time of Jacob's trouble; but he shall be saved out of it."
Jacob's Trouble = 7-year tribulation focused on Israel's national repentance.
Purpose: Bring Israel to recognize Yeshua as their Messiah.

**Phase 3: NATIONAL REPENTANCE (Future)**
Zechariah 12:10 — "They shall look upon me whom they have pierced,
and they shall mourn for him"
Romans 11:26 — All Israel saved.
This occurs at or just before the Second Coming.

**Phase 4: THE SECOND COMING**
Zechariah 14:4 — His feet stand on the Mount of Olives.
Matthew 24:30 — "They will see the Son of Man coming on the clouds of heaven,
with power and great glory."
Revelation 19:11-16 — The Rider on the White Horse.
Acts 1:11 — "This same Jesus, which is taken up from you into heaven,
shall so come in like manner as ye have seen him go into heaven."

**Phase 5: THE MILLENNIUM**
Revelation 20:4-6 — 1,000 year reign of Christ on earth.
Isaiah 65:17-25 — The new creation during the Millennium.
Micah 4:1-4 — "They shall beat their swords into plowshares"
Jerusalem as the capital of the world government (Zechariah 8:22-23).

## THE GENTILES IN GOD'S END TIMES PLAN

**The Times of the Gentiles**
Luke 21:24 (KJV) — "Jerusalem shall be trodden down of the Gentiles,
until the times of the Gentiles be fulfilled."
This period began with Babylon (605 BC) and continues until...
June 7, 1967: Jews recaptured Jerusalem and the Temple Mount.
Are the "Times of the Gentiles" now fulfilled or still ongoing? [multiple views]

**Romans 11 — The Mystery of Israel and the Gentiles**
Romans 11:11 — Israel's stumbling = riches for the Gentiles.
Romans 11:17-24 — Wild olive branches (Gentiles) grafted into the cultivated tree.
Romans 11:25-26 — Fullness of Gentiles → then all Israel saved.
The order of end times salvation: Gentile fullness THEN Israel.

**Revelation 7 — The 144,000 and the Great Multitude**
144,000 Jewish evangelists sealed (12,000 from each tribe).
The Great Multitude (Revelation 7:9) — "a great multitude that no one could count,
from every nation, tribe, people and language" — Gentile tribulation saints.
These are saved AFTER the Rapture during the tribulation period.

**The Nations in the Millennium**
Isaiah 2:2-4 — All nations stream to Jerusalem.
Zechariah 8:22-23 — "Many peoples and nations will come to seek the LORD Almighty
in Jerusalem... Ten men from all languages and nations will take hold of the robe
of one Jew and say, 'Let us go with you, because we have heard that God is with you.'"
The Gentile nations will look to Israel for spiritual leadership.

## CURRENT PHASE ASSESSMENT (2025-2026)
Where are we in this timeline? [analysis with multiple theological views]
What signs indicate the transition from "Church Age" to "Tribulation"?
The convergence of all prophetic streams in 2025-2026: [analysis]

## THE CONVERGENCE THEORY
Many scholars note that for the first time in history:
✅ Israel is back in the land (1948)
✅ Jerusalem is under Jewish control (1967)
✅ Technology for global economic control exists (digital currency, AI)
✅ Technology for global surveillance exists
✅ Third Temple preparations advanced
✅ Multiple nation coalitions against Israel forming
✅ Gospel preached to almost all nations
✅ Global government infrastructure building (UN, WHO, WEF)

The convergence of ALL these signs simultaneously is historically unprecedented.

Cite specific scriptures. Present Pre-Trib, Mid-Trib, and Post-Trib perspectives."""

    result = query_llm(prompt)

    fig = None
    try:
        phases = ["Regathering\n(1948+)", "Jerusalem\nRecaptured\n(1967)",
                   "Church Age\nEnds/Rapture", "Jacob's\nTrouble\n(Tribulation)",
                   "National\nRepentance", "Second\nComing",
                   "Millennium\n(1000 yrs)", "Eternal\nState"]
        fulfillment = [10, 10, 8, 0, 0, 0, 0, 0]
        colors_t = ["#27ae60" if f == 10 else "#f39c12" if f >= 7 else
                     "#3498db" if f >= 1 else "#95a5a6" for f in fulfillment]

        fig = go.Figure(go.Bar(
            x=phases, y=[10, 10, 8, 10, 10, 10, 10, 10],
            marker_color=colors_t,
            text=["FULFILLED", "FULFILLED", "IMMINENT?",
                   "FUTURE", "FUTURE", "FUTURE", "FUTURE", "FUTURE"],
            textposition="auto"
        ))
        fig.update_layout(
            title="🌐 Israel & Gentiles — Prophetic Timeline Progress",
            template="plotly_dark", height=420,
            paper_bgcolor="#0a0a1a",
            yaxis=dict(showticklabels=False)
        )
    except Exception:
        fig = None

    return result + "\n\n" + WATERMARK, fig


# ─── Tab 11: Apocalypse Score Dashboard ───────────────────────────────────────

def analyze_apocalypse_score():
    prompt = """You are a biblical eschatology scholar providing a comprehensive
end times alignment score based on ALL current prophetic indicators.

## 📊 APOCALYPSE SCORE DASHBOARD — 2025/2026

## COMPREHENSIVE END TIMES ALIGNMENT SCORE

Rate each category 0-100 based on prophetic fulfillment:

## CATEGORY SCORES

**ISRAEL AND JERUSALEM** (Weight: 20%)
- Israel reestablished as nation (1948): 100%
- Jerusalem under Jewish control (1967): 100%
- Third Temple preparations: 60-70%
- Red heifers ready for purification: 80%
- Sanhedrin reestablished (2004): 70%
Category Score: 85/100

**GEOPOLITICAL ALIGNMENT** (Weight: 20%)
- Russia-Iran-Turkey coalition (Ezekiel 38): 85%
- Global government emerging (UN, WHO, WEF): 75%
- Israel-Arab normalization (Abraham Accords): 70%
- All nations against Jerusalem (Zechariah 12): 65%
- Kings of the East rising (China, Asia): 75%
Category Score: 74/100

**TECHNOLOGY AND MARK OF BEAST** (Weight: 15%)
- Digital payment infrastructure: 90%
- CBDC development (Central Bank Digital Currencies): 80%
- AI surveillance systems: 85%
- Biometric identification: 85%
- Global digital ID frameworks: 75%
Category Score: 83/100

**SOCIAL AND MORAL CONDITIONS** (Weight: 15%)
- Days of Noah conditions (Matthew 24:37): 80%
- Days of Lot conditions (Luke 17:28-29): 85%
- Great apostasy (2 Timothy 3): 75%
- Perilous times characteristics: 80%
- Lawlessness increasing (Matthew 24:12): 80%
Category Score: 80/100

**COSMIC AND NATURAL SIGNS** (Weight: 10%)
- Earthquakes (Matthew 24:7): 65%
- Solar activity anomalies: 70%
- Magnetic pole shift acceleration: 60%
- Pestilences post-COVID: 70%
- Unusual weather events: 70%
Category Score: 67/100

**SPIRITUAL SIGNS** (Weight: 10%)
- Gospel preached to all nations (Matthew 24:14): 90%
- Christian persecution increasing: 85%
- Fullness of Gentiles approaching: 80%
- Jewish aliyah (return) accelerating: 85%
Category Score: 85/100

**ECONOMIC SIGNS** (Weight: 10%)
- Global economic instability: 75%
- Dollar/debt crisis brewing: 80%
- BRICS/de-dollarization challenge: 70%
- Wealth inequality extremes: 80%
- Babylon the Great economic system: 65%
Category Score: 74/100

## OVERALL APOCALYPSE ALIGNMENT SCORE
**TOTAL SCORE: [Calculate weighted average] / 100**

## INTERPRETATION SCALE
0-20: Early signs — nothing extraordinary
21-40: Notable signs — prophetic alignment beginning
41-60: Significant alignment — multiple signs converging
61-80: STRONG alignment — generation may see fulfillment
81-90: VERY STRONG — tribulation could begin this generation
91-100: IMMINENT — all major signs fulfilled

## WHAT THE SCORE MEANS
Based on this analysis, we are in the [X] range, indicating [assessment].
This is consistent with the view of many eschatology scholars who believe
we are living in the terminal generation spoken of in Matthew 24:34.

"Verily I say unto you, This generation shall not pass,
till all these things be fulfilled." — Matthew 24:34 (KJV)

## THE CONVERGENCE FACTOR
For the first time in 2,000+ years, ALL the major end times signs are
converging simultaneously. This convergence itself is the most significant
prophetic indicator of all.

## IMPORTANT CAVEAT
"But of that day and hour knoweth no man, no, not the angels of heaven,
but my Father only." — Matthew 24:36 (KJV)

This score is for educational awareness, not date-setting.
The purpose is to encourage spiritual readiness.
"Watch therefore: for ye know not what hour your Lord doth come." — Matthew 24:42"""

    result = query_llm(prompt)

    fig = None
    try:
        categories = ["Israel &\nJerusalem", "Geopolitical\nAlignment",
                       "Mark of Beast\nTechnology", "Social &\nMoral Signs",
                       "Cosmic &\nNatural Signs", "Spiritual\nSigns", "Economic\nSigns"]
        scores_a = [85, 74, 83, 80, 67, 85, 74]
        weights = [20, 20, 15, 15, 10, 10, 10]
        weighted_avg = sum(s * w for s, w in zip(scores_a, weights)) / sum(weights)

        fig = make_subplots(rows=1, cols=2,
                             subplot_titles=[f"End Times Score by Category",
                                             f"Overall Score: {weighted_avg:.0f}/100"],
                             specs=[[{"type": "bar"}, {"type": "indicator"}]])

        colors_a = ["#27ae60" if s >= 80 else "#f39c12" if s >= 65 else "#e74c3c"
                    for s in scores_a]
        fig.add_trace(go.Bar(x=categories, y=scores_a,
                              marker_color=colors_a,
                              text=[f"{s}" for s in scores_a],
                              textposition="auto"), row=1, col=1)

        fig.add_trace(go.Indicator(
            mode="gauge+number",
            value=weighted_avg,
            title={"text": "Apocalypse Alignment %"},
            gauge={
                "axis": {"range": [0, 100]},
                "bar": {"color": "#e74c3c"},
                "steps": [
                    {"range": [0, 40], "color": "#27ae60"},
                    {"range": [40, 60], "color": "#f39c12"},
                    {"range": [60, 80], "color": "#e67e22"},
                    {"range": [80, 100], "color": "#c0392b"}
                ],
                "threshold": {
                    "line": {"color": "white", "width": 4},
                    "thickness": 0.75,
                    "value": weighted_avg
                }
            }
        ), row=1, col=2)

        fig.update_layout(
            title="📊 APOCALYPSE SCORE DASHBOARD — 2025/2026",
            template="plotly_dark", height=450,
            paper_bgcolor="#0a0a1a"
        )
    except Exception:
        fig = None

    return result + "\n\n" + WATERMARK, fig


# ─── Tab 12: AI Prophecy Chat ─────────────────────────────────────────────────

def prophecy_chat(message, history):
    messages = [
        {"role": "system", "content": SYSTEM_PROMPT + (
            "\n\nYou are helping citizens understand end times prophecy from an educational perspective. "
            "Answer questions about biblical prophecy, the Book of Enoch, Jubilees, the Rapture/Harpazo, "
            "the Antichrist, the Tribulation, Israel and the Gentiles, the Abraham Accords, Noahide Laws, "
            "solar activity, earth changes, and all related topics. "
            "Always cite specific Bible verses. Present multiple theological views fairly. "
            "Remind users that no one knows the day or hour (Matthew 24:36). "
            "Encourage spiritual readiness without causing fear."
        )}
    ]
    for user_msg, bot_msg in history:
        messages.append({"role": "user", "content": user_msg})
        messages.append({"role": "assistant", "content": bot_msg})
    messages.append({"role": "user", "content": message})

    import os
    API_KEY = os.environ.get("OPENAI_API_KEY", "")
    headers = {"Authorization": f"Bearer {API_KEY}", "Content-Type": "application/json"}
    msgs = [{"role": "system", "content": SYSTEM_PROMPT}, {"role": "user", "content": prompt}]
    payload = {"model": "gpt-4o", "max_tokens": 4000, "messages": msgs}
    try:
        response = requests.post("https://api.openai.com/v1/chat/completions", headers=headers, json=payload, timeout=240)
        result = response.json()
        if "choices" in result:
            return result["choices"][0]["message"]["content"]
        return f"API Error: {result}"
    except Exception as e:
        return f"Error: {str(e)}"


# ─── Gradio UI ────────────────────────────────────────────────────────────────



# ─── Tab 12: Prophetic Calendar ───────────────────────────────────────────────

def analyze_prophetic_calendar(extra_context):
    prompt = f"""You are a biblical scholar and Hebrew calendar expert.

Additional Context: {enrich_with_urls(extra_context) if extra_context else "2025/2026 prophetic calendar"}

## 📅 PROPHETIC CALENDAR — 2025/2026

## HEBREW CALENDAR YEAR
Current Hebrew Year: 5785-5786 (2025-2026)
Significance of this year in Jewish tradition: [analysis]
Sabbatical/Shemitah year status: [analysis]

## JEWISH FEAST DAYS 2025-2026 (MOEDIM — GOD'S APPOINTED TIMES)

Leviticus 23 — The seven feasts of the LORD:

**SPRING FEASTS (Prophetically fulfilled at First Coming):**
| Feast | Hebrew Date | 2025/2026 Date | Prophetic Fulfillment |
|-------|-------------|----------------|----------------------|
| Passover (Pesach) | Nisan 14 | [date] | Crucifixion of Christ ✅ |
| Unleavened Bread | Nisan 15-22 | [date] | Burial/Sinlessness ✅ |
| Firstfruits (Bikkurim) | Nisan 16-17 | [date] | Resurrection ✅ |
| Pentecost (Shavuot) | Sivan 6 | [date] | Holy Spirit outpouring ✅ |

**FALL FEASTS (Prophetically unfulfilled — Second Coming):**
| Feast | Hebrew Date | 2025/2026 Date | Prophetic Expectation |
|-------|-------------|----------------|----------------------|
| Feast of Trumpets (Rosh Hashanah) | Tishri 1 | [date] | Rapture/Resurrection? |
| Day of Atonement (Yom Kippur) | Tishri 10 | [date] | Israel's National Salvation? |
| Feast of Tabernacles (Sukkot) | Tishri 15-22 | [date] | Millennium? |
| Shemini Atzeret | Tishri 22 | [date] | Eternal State? |

## BLOOD MOON AND SOLAR ECLIPSE CALENDAR

**Recent Blood Moons on Jewish Feast Days:**
- 2014-2015 Tetrad: Passover and Sukkot both years — on Israel's feast days
- Pattern: Blood moon tetrads have coincided with major Jewish historical events
  - 1949-1950: Israel becomes a nation
  - 1967-1968: Six Day War
  - 2014-2015: [analysis of what followed]

**Upcoming Lunar Eclipses 2025-2027:**
[List upcoming eclipses and any feast day alignments]

**Solar Eclipses of Note:**
April 8, 2024 solar eclipse — path of totality through USA significance: [analysis]
Upcoming solar eclipses 2025-2027: [list]
Prophetic significance of solar eclipses in Scripture:
- Joel 2:31 — "The sun shall be turned into darkness"
- Amos 8:9 — "I will cause the sun to go down at noon"

## PLANETARY CONJUNCTIONS AND SIGNS 2025-2026
Key planetary alignments and their astronomical/astrological significance: [list]
The "Great Sign" of Revelation 12 (September 23, 2017) — retrospective analysis
Any similar signs upcoming: [analysis]

## SABBATH YEARS AND JUBILEES CALENDAR
Current Sabbatical cycle position: [analysis]
Next Shemitah year: [date]
Next projected Jubilee: [date and significance]

## COMING PROPHETIC DATES TO WATCH
Based on multiple calculation methods — dates scholars are watching:
[List specific upcoming dates with their prophetic basis]
Note: We are not predicting — we are watching (Matthew 24:42)

Be specific with dates, cite Leviticus 23, and multiple scholarly perspectives."""

    result = query_llm(prompt)

    fig = None
    try:
        feasts = ["Passover", "Unleavened\nBread", "Firstfruits", "Pentecost",
                   "Feast of\nTrumpets", "Day of\nAtonement", "Feast of\nTabernacles"]
        fulfillment = [10, 10, 10, 10, 0, 0, 0]
        colors_f = ["#27ae60" if f == 10 else "#3498db" for f in fulfillment]
        labels = ["FULFILLED ✅"] * 4 + ["FUTURE ⏳"] * 3

        fig = go.Figure(go.Bar(
            x=feasts, y=[10] * 7,
            marker_color=colors_f,
            text=labels,
            textposition="auto"
        ))
        fig.update_layout(
            title="📅 God's Seven Feasts — Prophetic Fulfillment Status",
            yaxis=dict(showticklabels=False),
            template="plotly_dark", height=380,
            paper_bgcolor="#0a0a1a"
        )
    except Exception:
        fig = None

    return result + "\n\n" + WATERMARK, fig


# ─── Tab 13: Third Temple Tracker ─────────────────────────────────────────────

def analyze_third_temple(extra_context):
    prompt = f"""You are a scholar of Jewish temple worship and end times prophecy.

Additional Context: {enrich_with_urls(extra_context) if extra_context else "2025/2026 Third Temple status"}

## 🕍 THIRD TEMPLE TRACKER — 2025/2026

## THE PROPHETIC NECESSITY OF THE THIRD TEMPLE

**Why the Third Temple Must Be Built:**
- Daniel 9:27 — "He will put an end to sacrifice and offering" (requires a functioning Temple)
- Matthew 24:15 — "The abomination that causes desolation, spoken of through the prophet Daniel,
  standing in the holy place" — requires a Temple
- 2 Thessalonians 2:4 — The Man of Sin "sits in the temple of God, proclaiming himself to be God"
- Revelation 11:1-2 — John told to measure the Temple (implies it exists during tribulation)

**History of the Temples:**
- First Temple (Solomon's): Built ~957 BC, Destroyed 586 BC by Babylon
- Second Temple (Zerubbabel/Herod): Rebuilt ~516 BC, Destroyed 70 AD by Rome
- Third Temple: Prophesied — not yet built
- Time gap between Second and Third: 1,955+ years (longest in history)

## TEMPLE MOUNT STATUS (2025-2026)

**Current Situation:**
- Al-Aqsa Mosque and Dome of the Rock occupy the Temple Mount
- Israeli sovereignty over Temple Mount since 1967 Six Day War
- Waqf (Islamic authority) controls day-to-day administration
- Jewish prayer on the Temple Mount: restricted but slowly expanding
- Political sensitivity: Any change here could trigger regional war

**Recent Developments 2024-2026:**
[Current status of Temple Mount access and political negotiations]

## THE TEMPLE INSTITUTE (MACHON HAMIKDASH)

**What They Have Prepared:**
- ✅ Rebuilt the Menorah (Golden, ready for Temple)
- ✅ High Priest garments completed
- ✅ Temple vessels recreated (280+ sacred vessels)
- ✅ Priestly garments for all Temple priests
- ✅ The Altar (rebuilt and ready)
- ✅ Musical instruments for Temple worship
- ✅ Training priests for Temple service (Kohanim)
- 🔄 Architectural plans for Third Temple completed
- ❌ Location — exact spot debated

**The Temple Institute Statement:**
They are prepared to begin Temple construction and worship
"immediately upon receiving the go-ahead" — All physical preparations complete.

## THE RED HEIFER — NUMBERS 19

**Why It Matters:**
Numbers 19:1-22 — The red heifer (parah adumah) ceremony is required for ritual purification.
Without it, priests cannot perform Temple service.
The ashes of a red heifer mixed with water create the "water of purification."

**The 2024 Red Heifers — MAJOR PROPHETIC DEVELOPMENT:**
In September 2022: Five red heifers from Texas arrived in Israel
Raised by rancher Byron Kellogg in conjunction with Temple Institute
As of 2024-2025: The heifers have been examined and certified as qualifying
These are the first kosher red heifers in Israel in 2,000+ years
The ceremony can be performed near the Mount of Olives (Numbers 19 location)
Status in 2025-2026: [current update]

**Historical significance:**
Only 9 red heifers have been sacrificed in all of Jewish history (per Mishnah).
The 10th red heifer is associated with the Messianic era.
Is the current red heifer the 10th? Many Jewish scholars believe so.

## THE SANHEDRIN

**Reestablished in 2004:**
The Sanhedrin (Jewish high court of 71) was reestablished October 13, 2004 in Tiberias.
This had not existed since 425 AD — a 1,579-year gap.
Prophetic significance: The Sanhedrin must rule on Temple worship and the Messiah.
Current Sanhedrin activities 2025-2026: [analysis]

## LOCATION DEBATE

**Three Main Theories for Temple Location:**
1. **Traditional (Over Dome of the Rock)**: Most widely held — but requires removing Islamic structures
2. **Northern Location (Asher Kaufman theory)**: Temple was 100m north of Dome
3. **Southern Location (Bob Cornuke theory)**: Temple was over the Gihon Spring

**If the Northern theory is correct:**
The Temple could be built WITHOUT removing the Dome of the Rock.
Revelation 11:2 — "But exclude the outer court; do not measure it,
because it has been given to the Gentiles" — this could allow coexistence.

## PREPARATION PROGRESS SCORECARD
| Item | Status | Score |
|------|--------|-------|
| Temple Vessels | COMPLETE | 10/10 |
| Priestly Garments | COMPLETE | 10/10 |
| Red Heifer | NEAR READY | 8/10 |
| Trained Priests | IN PROGRESS | 7/10 |
| Architectural Plans | COMPLETE | 10/10 |
| Political Permission | NOT YET | 1/10 |
| Temple Mount Access | LIMITED | 2/10 |
| **OVERALL** | **IN PROGRESS** | **7/10** |

Cite specific scriptures (Daniel 9:27, Matthew 24:15, 2 Thess 2:4, Revelation 11)."""

    result = query_llm(prompt)

    fig = None
    try:
        items = ["Temple\nVessels", "Priestly\nGarments", "Red\nHeifer",
                  "Trained\nPriests", "Architectural\nPlans", "Political\nPermission",
                  "Temple Mount\nAccess", "Exact\nLocation"]
        scores_t = [10, 10, 8, 7, 10, 1, 2, 4]
        colors_tt = ["#27ae60" if s >= 8 else "#f39c12" if s >= 5 else "#e74c3c"
                     for s in scores_t]

        fig = go.Figure(go.Bar(
            x=items, y=scores_t,
            marker_color=colors_tt,
            text=[f"{s}/10" for s in scores_t],
            textposition="auto"
        ))
        fig.add_hline(y=8, line_dash="dash", line_color="green",
                      annotation_text="Ready Threshold")
        fig.update_layout(
            title="🕍 Third Temple Preparation Progress — 2025/2026",
            yaxis_title="Readiness Score (0-10)",
            yaxis=dict(range=[0, 10]),
            template="plotly_dark", height=420,
            paper_bgcolor="#0a0a1a"
        )
    except Exception:
        fig = None

    return result + "\n\n" + WATERMARK, fig


# ─── Tab 14: One World System Monitor ────────────────────────────────────────

def analyze_one_world_system(extra_context):
    prompt = f"""You are a biblical scholar and geopolitical analyst examining the emergence
of global governance systems and their alignment with Revelation 13 and 17.

Additional Context: {enrich_with_urls(extra_context) if extra_context else "2025/2026 analysis"}

## 🌐 ONE WORLD SYSTEM MONITOR — REVELATION 13 ALIGNMENT

## THE BIBLICAL FRAMEWORK

**Revelation 13:7 (KJV)** — "And it was given unto him to make war with the saints,
and to overcome them: and power was given him over all kindreds, and tongues, and nations."

**Revelation 13:16-17 (KJV)** — "And he causeth all, both small and great, rich and poor,
free and bond, to receive a mark in their right hand, or in their foreheads: And that no man
might buy or sell, save he that had the mark, or the name of the beast, or the number of his name."

**Revelation 17:1-2** — The Great Whore / One World Religious System
**Daniel 7:23** — "A fourth beast... shall devour the whole earth"

## GLOBAL GOVERNANCE STRUCTURES (2025-2026)

### WORLD HEALTH ORGANIZATION (WHO) PANDEMIC TREATY
WHO International Health Regulations amendments: [current status]
The WHO Pandemic Treaty negotiations: [current status 2025-2026]
What powers it would grant WHO over member nations: [analysis]
Prophecy alignment: Global authority over health decisions
Score: X/10 implementation

### WORLD ECONOMIC FORUM (WEF) — THE GREAT RESET
Alois Zwinggi's "Great Reset" agenda — stated goals:
- "You will own nothing and you will be happy"
- Stakeholder capitalism
- Fourth Industrial Revolution
- Transhumanism connections
2030 Agenda for Sustainable Development (UN SDGs)
Prophecy alignment: Global economic control system
Score: X/10 implementation

### CENTRAL BANK DIGITAL CURRENCIES (CBDCs)
Countries with active CBDCs in 2025-2026: [list]
Countries in pilot/development: [list]
What CBDCs enable that cash does not:
- Transaction monitoring and control
- Ability to freeze individual accounts
- Programmable money (expiration dates, spending restrictions)
- Social credit score integration potential
Prophecy alignment: Revelation 13:17 — buy/sell control
Score: X/10 implementation

### DIGITAL ID AND BIOMETRIC SYSTEMS
WHO Smart Vaccination Certificate
EU Digital Identity Wallet
India's Aadhaar system (1.4 billion enrolled)
US digital ID initiatives
Facial recognition surveillance networks
Prophecy alignment: Mark/ID control system infrastructure
Score: X/10 implementation

### AI SURVEILLANCE AND SOCIAL CREDIT
China's Social Credit System — current status
AI surveillance expansion globally
Predictive policing systems
Mass data collection infrastructure
Prophecy alignment: Global monitoring and control
Score: X/10 implementation

### ONE WORLD RELIGION DEVELOPMENTS
Chrislam movement (Pope Francis/Islamic leaders)
The Abrahamic Family House (Abu Dhabi — mosque, church, synagogue)
Parliament of World Religions
UN Interfaith initiatives
Prophecy alignment: Revelation 17 — Great Whore one world religion
Score: X/10 implementation

### TRANSHUMANISM AND GENETIC MODIFICATION
mRNA technology platform expansion beyond vaccines
CRISPR genetic editing — human applications
Neural interfaces (Neuralink and competitors)
Synthetic biology developments
Prophecy alignment: Days of Noah — genetic corruption (Genesis 6, 1 Enoch 6-16)
Score: X/10 implementation

## COMPREHENSIVE SYSTEM IMPLEMENTATION SCORECARD
| System | Implementation | Prophecy Alignment |
|--------|---------------|-------------------|
| WHO Global Health Governance | X% | X/10 |
| CBDC Digital Currency | X% | X/10 |
| Digital ID/Biometrics | X% | X/10 |
| AI Surveillance | X% | X/10 |
| Social Credit Systems | X% | X/10 |
| One World Religion Movement | X% | X/10 |
| Transhumanism/Genetic Mod | X% | X/10 |
| Global Economic Control | X% | X/10 |
| **OVERALL SYSTEM READINESS** | **X%** | **X/10** |

## THE KEY OBSERVATION
For the first time in history, the TECHNOLOGY EXISTS to implement
Revelation 13:16-17 globally — a buy/sell control system.
This was impossible until the digital age.
The infrastructure is being built RIGHT NOW.

## IMPORTANT CONTEXT
Not all globalization is sinister — some cooperation is beneficial.
Many good-faith actors in these organizations have humanitarian goals.
The concern is the POTENTIAL for misuse by a future authoritarian system.
Prophetic analysis does not require that current leaders are intentionally fulfilling prophecy.

Cite Revelation 13, 17, Daniel 7, Matthew 24, and 2 Thessalonians 2."""

    result = query_llm(prompt)

    fig = None
    try:
        systems = ["WHO\nGovernance", "CBDC\nCurrency", "Digital\nID",
                    "AI\nSurveillance", "Social\nCredit", "One World\nReligion",
                    "Transhumanism", "Economic\nControl"]
        implementation = [65, 72, 78, 80, 55, 45, 40, 70]
        prophecy_align = [8, 9, 9, 8, 7, 7, 8, 9]

        fig = make_subplots(rows=1, cols=2,
                             subplot_titles=["System Implementation %",
                                             "Prophecy Alignment Score"])
        fig.add_trace(go.Bar(x=systems, y=implementation,
                              marker_color="#e74c3c", name="% Implemented",
                              text=[f"{v}%" for v in implementation],
                              textposition="auto"), row=1, col=1)
        fig.add_trace(go.Bar(x=systems, y=prophecy_align,
                              marker_color="#f39c12", name="Prophecy Alignment",
                              text=[f"{v}/10" for v in prophecy_align],
                              textposition="auto"), row=1, col=2)
        fig.update_layout(
            title="🌐 One World Systems — Implementation & Prophecy Alignment",
            template="plotly_dark", height=440,
            paper_bgcolor="#0a0a1a"
        )
    except Exception:
        fig = None

    return result + "\n\n" + WATERMARK, fig


# ─── Tab 15: Reports & Presentations ─────────────────────────────────────────

def generate_endtimes_report(report_type, focus, audience, extra_context):

    focus_prompts = {
        "Complete Overview — All Topics": "Cover ALL end times topics comprehensively",
        "Israel and Middle East Focus": "Focus on Israel, Middle East wars, and Ezekiel 38-39",
        "Signs of the Times Focus": "Focus on Matthew 24 signs, earth changes, and cosmic events",
        "Rapture/Harpazo Focus": "Focus on the Harpazo, tribulation views, and deliverance",
        "Extra-Biblical Texts Focus": "Focus on Enoch, Jubilees, Ethiopian Canon",
        "Antichrist System Focus": "Focus on the Antichrist, one world system, mark of the beast"
    }

    focus_instruction = focus_prompts.get(focus, "Cover all end times topics")

    prompt = f"""You are a biblical eschatology scholar creating a professional end times report.

Report Type: {report_type}
Content Focus: {focus}
Target Audience: {audience}
Additional Notes: {enrich_with_urls(extra_context) if extra_context else "Standard analysis"}

{focus_instruction}

Create a comprehensive {report_type} with the following structure:

# END TIMES PROPHETIC ANALYSIS REPORT
## Generated: {datetime.now().strftime('%B %d, %Y')}
## Target Audience: {audience}

---

## EXECUTIVE SUMMARY
3-5 sentence overview of the current prophetic moment we are living in.
Key headline finding.

---

## SECTION 1: WHERE WE ARE IN PROPHETIC HISTORY
Current prophetic timeline assessment.
Major prophecies already fulfilled.
Major prophecies in progress.
What scholars across theological traditions agree on.

---

## SECTION 2: ISRAEL AND THE MIDDLE EAST
Current geopolitical situation.
Prophetic alignment with Ezekiel 38-39, Psalm 83, Isaiah 17.
Third Temple preparations status.
Abraham Accords significance.

---

## SECTION 3: SIGNS OF THE TIMES
Matthew 24 scorecard — how many signs are fulfilled.
Earth changes — earthquakes, pestilence, extreme weather.
Solar activity and cosmic signs.
Technology enabling end times systems.

---

## SECTION 4: THE HARPAZO / RAPTURE
All major theological views presented fairly.
Pre-Tribulation, Mid-Tribulation, Post-Tribulation, Pre-Wrath.
Current readiness indicators.
Deliverance of Israel and the Gentiles.

---

## SECTION 5: EXTRA-BIBLICAL INSIGHTS
Book of Enoch contributions to end times understanding.
Book of Jubilees and the current Jubilee cycle.
Ethiopian Orthodox Canon unique insights.
Dead Sea Scrolls War Scroll.

---

## SECTION 6: THE ONE WORLD SYSTEM
CBDC, digital ID, WHO governance, WEF.
Revelation 13 alignment assessment.
The Mark of the Beast system infrastructure.

---

## SECTION 7: THE JUBILEE CALENDAR
Current position in the Jubilee cycle.
Significance of 2025-2026 in Hebrew calendar.
Shemitah year implications.

---

## SECTION 8: PROPHETIC SCORE SUMMARY
Overall end times alignment score: X/100
Key indicators driving this score.
What this means for believers.

---

## SECTION 9: CALL TO ACTION
What believers should do in response to these signs.
Spiritual preparation recommendations.
Matthew 24:42 — "Watch therefore"
Matthew 24:44 — "Be ye also ready"

---

## KEY SCRIPTURE REFERENCES
List 20 most important end times scriptures with full text (KJV).

---

## GLOSSARY
Define key terms: Harpazo, Tribulation, Antichrist, Abomination of Desolation,
Gog-Magog, Noahide Laws, Jubilee, Shemitah, Parousia, etc.

Write for a {audience} audience. Use KJV as primary translation.
Be thorough, reverent, and educational."""

    report_text = query_llm(prompt)

    output_path = None
    try:
        if "PowerPoint" in report_type:
            output_path = _make_prophecy_pptx(report_text, report_type, focus)
        else:
            output_path = _make_prophecy_pdf(report_text, report_type, focus)
    except Exception as e:
        report_text += f"\n\n[Document generation error: {str(e)}]"

    return report_text + "\n\n" + WATERMARK, output_path


def _make_prophecy_pptx(text, report_type, focus):
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    def hex_to_rgb(hex_color):
        h = hex_color.lstrip("#")
        return tuple(int(h[i:i+2], 16) for i in (0, 2, 4))

    def add_slide(title_text, body_text, bg_hex="#0a0a1a", title_hex="#f1c40f"):
        slide = prs.slides.add_slide(blank_layout)
        bg = slide.background
        bg.fill.solid()
        bg.fill.fore_color.rgb = RGBColor(*hex_to_rgb(bg_hex))
        tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12.3), Inches(1.3))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.size = Pt(26)
        p.font.bold = True
        p.font.color.rgb = RGBColor(*hex_to_rgb(title_hex))
        body_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.6), Inches(12.3), Inches(5.6))
        tf2 = body_box.text_frame
        tf2.word_wrap = True
        lines = [l.strip() for l in body_text.split("\n") if l.strip()]
        for line in lines[:22]:
            para = tf2.add_paragraph()
            line_clean = line.lstrip("*#- ").strip("*")
            para.text = line_clean
            para.font.size = Pt(13)
            para.font.color.rgb = RGBColor(235, 225, 200)

    # Title slide
    title_slide = prs.slides.add_slide(blank_layout)
    bg = title_slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(5, 5, 15)
    tb = title_slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(11.7), Inches(1.8))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "✝️ End Times Prophetic Analysis"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(241, 196, 15)
    tb2 = title_slide.shapes.add_textbox(Inches(0.8), Inches(3.3), Inches(11.7), Inches(1.0))
    tf2 = tb2.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = focus
    p2.font.size = Pt(22)
    p2.font.color.rgb = RGBColor(200, 180, 140)
    tb3 = title_slide.shapes.add_textbox(Inches(0.8), Inches(4.5), Inches(11.7), Inches(0.8))
    tf3 = tb3.text_frame
    p3 = tf3.paragraphs[0]
    p3.text = f"Generated: {datetime.now().strftime('%B %d, %Y')} | AI End Times Prophecy Analyzer"
    p3.font.size = Pt(16)
    p3.font.color.rgb = RGBColor(150, 140, 110)
    tb4 = title_slide.shapes.add_textbox(Inches(0.8), Inches(5.5), Inches(11.7), Inches(0.7))
    tf4 = tb4.text_frame
    p4 = tf4.paragraphs[0]
    p4.text = '"Watch therefore: for ye know not what hour your Lord doth come." — Matthew 24:42 (KJV)'
    p4.font.size = Pt(14)
    p4.font.italic = True
    p4.font.color.rgb = RGBColor(180, 160, 120)

    # Parse sections
    sections = []
    current_title = "Overview"
    current_body = []
    for line in text.split("\n"):
        stripped = line.strip()
        if stripped.startswith("## ") or (stripped.startswith("**") and stripped.endswith("**") and len(stripped) > 4):
            if current_body:
                sections.append((current_title, "\n".join(current_body)))
            current_title = stripped.lstrip("#* ").rstrip("*# ")
            current_body = []
        elif stripped.startswith("# ") and len(stripped) > 2:
            if current_body:
                sections.append((current_title, "\n".join(current_body)))
            current_title = stripped[2:].strip()
            current_body = []
        else:
            current_body.append(line)
    if current_body:
        sections.append((current_title, "\n".join(current_body)))

    bg_themes = [
        ("#0a0a1a", "#f1c40f"),
        ("#0d0a1a", "#e74c3c"),
        ("#0a1a0a", "#2ecc71"),
        ("#1a0a0a", "#e67e22"),
        ("#0a0d1a", "#3498db"),
        ("#1a1a0a", "#9b59b6"),
    ]

    for i, (sec_title, sec_body) in enumerate(sections[:12]):
        bg_hex, title_hex = bg_themes[i % len(bg_themes)]
        add_slide(sec_title, sec_body, bg_hex, title_hex)

    # Scripture slide
    scriptures = [
        "Matthew 24:36 — But of that day and hour knoweth no man",
        "Matthew 24:42 — Watch therefore: for ye know not what hour your Lord doth come",
        "1 Thess 4:17 — Then we which are alive shall be caught up... to meet the Lord",
        "Revelation 22:20 — Even so, come, Lord Jesus",
        "Romans 11:26 — And so all Israel shall be saved",
        "Zechariah 12:10 — They shall look upon me whom they have pierced",
    ]
    add_slide("📖 Key Scriptures to Remember", "\n".join(scriptures), "#0a0a15", "#f1c40f")

    # Closing slide
    closing = prs.slides.add_slide(blank_layout)
    bg = closing.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(5, 5, 15)
    tb = closing.shapes.add_textbox(Inches(1), Inches(2.2), Inches(11.3), Inches(1.5))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = '"Even so, come, Lord Jesus." — Revelation 22:20'
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(241, 196, 15)
    p2 = tf.add_paragraph()
    p2.text = " "
    p3 = tf.add_paragraph()
    p3.text = "© 2026 Existential Gateway, LLC | AI End Times Prophecy Analyzer"
    p3.font.size = Pt(14)
    p3.font.color.rgb = RGBColor(150, 140, 110)

    tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".pptx")
    prs.save(tmp.name)
    return tmp.name


def _make_prophecy_pdf(text, report_type, focus):
    tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".pdf")
    doc = SimpleDocTemplate(tmp.name, pagesize=letter,
                             leftMargin=inch, rightMargin=inch,
                             topMargin=inch, bottomMargin=inch)
    styles = getSampleStyleSheet()
    title_style = ParagraphStyle("Title2", parent=styles["Title"],
                                  fontSize=20, spaceAfter=16,
                                  textColor=colors.HexColor("#5d4037"))
    heading_style = ParagraphStyle("H2", parent=styles["Heading2"],
                                    fontSize=14, spaceAfter=8,
                                    textColor=colors.HexColor("#7b3f00"))
    verse_style = ParagraphStyle("Verse", parent=styles["Normal"],
                                  fontSize=11, spaceAfter=6, leading=16,
                                  leftIndent=20,
                                  textColor=colors.HexColor("#4a148c"),
                                  fontName="Helvetica-Oblique")
    body_style = ParagraphStyle("Body2", parent=styles["Normal"],
                                 fontSize=11, spaceAfter=6, leading=16)
    story = []
    story.append(Paragraph("✝️ End Times Prophetic Analysis Report", title_style))
    story.append(Paragraph(f"Focus: {focus}", heading_style))
    story.append(Paragraph(
        f"Generated: {datetime.now().strftime('%B %d, %Y')} | AI End Times Prophecy Analyzer",
        body_style))
    story.append(Paragraph(
        '"Watch therefore: for ye know not what hour your Lord doth come." — Matthew 24:42 (KJV)',
        verse_style))
    story.append(Spacer(1, 0.2 * inch))

    for line in text.split("\n"):
        stripped = line.strip()
        if not stripped:
            story.append(Spacer(1, 0.08 * inch))
            continue
        if stripped.startswith("## "):
            story.append(Paragraph(stripped[3:], heading_style))
        elif stripped.startswith("# "):
            story.append(Paragraph(stripped[2:], title_style))
        elif stripped.startswith("**") and stripped.endswith("**"):
            story.append(Paragraph(stripped.strip("*"), heading_style))
        elif any(book in stripped for book in
                  ["KJV", "NKJV", "Matthew", "Revelation", "Daniel", "Ezekiel",
                   "Isaiah", "Zechariah", "Romans", "Thessalonians", "Joel"]) and "—" in stripped:
            safe = stripped.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")
            story.append(Paragraph(safe, verse_style))
        else:
            safe = stripped.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")
            story.append(Paragraph(safe, body_style))

    story.append(Spacer(1, 0.4 * inch))
    story.append(Paragraph(
        '"Even so, come, Lord Jesus." — Revelation 22:20 (KJV)',
        verse_style))
    story.append(Spacer(1, 0.2 * inch))
    story.append(Paragraph(
        "© 2026 Existential Gateway, LLC | AI End Times Prophecy Analyzer | existentialgateway@gmail.com",
        ParagraphStyle("footer", parent=styles["Normal"], fontSize=9,
                        textColor=colors.grey)))
    doc.build(story)
    return tmp.name


with gr.Blocks(title="AI End Times Prophecy Analyzer",
               theme=gr.themes.Base(
        primary_hue=gr.themes.colors.Color(
            c50="#fef9ec", c100="#faefd0", c200="#f4dea0", c300="#edc970",
            c400="#e4b04a", c500="#c9a84c", c600="#a8872e", c700="#856519",
            c800="#5e440d", c900="#3a2a05", c950="#1e1502"
        ),
        secondary_hue=gr.themes.colors.Color(
            c50="#e8edf5", c100="#c5d0e6", c200="#9eb0d4", c300="#7490c2",
            c400="#4f74b0", c500="#2d5a9e", c600="#1a3a6e", c700="#112240",
            c800="#0a1628", c900="#060e1a", c950="#03070d"
        ),
        neutral_hue=gr.themes.colors.Color(
            c50="#f0f2f7", c100="#d8dde9", c200="#b8c2d6", c300="#95a3be",
            c400="#7285a6", c500="#536890", c600="#3a4f73", c700="#263856",
            c800="#162438", c900="#0c1622", c950="#060b11"
        ),
        font=gr.themes.GoogleFont("DM Sans"),
        font_mono=gr.themes.GoogleFont("DM Mono"),
    ).set(
        body_background_fill="#0a1628",
        body_background_fill_dark="#0a1628",
        body_text_color="#f8f9fc",
        body_text_color_dark="#f8f9fc",
        block_background_fill="#112240",
        block_background_fill_dark="#112240",
        block_border_color="rgba(201,168,76,0.2)",
        block_border_color_dark="rgba(201,168,76,0.2)",
        block_label_background_fill="#112240",
        block_label_background_fill_dark="#112240",
        block_label_text_color="#c9a84c",
        block_label_text_color_dark="#c9a84c",
        block_title_text_color="#c9a84c",
        block_title_text_color_dark="#c9a84c",
        button_primary_background_fill="linear-gradient(135deg,#c9a84c,#e8c96a)",
        button_primary_background_fill_dark="linear-gradient(135deg,#c9a84c,#e8c96a)",
        button_primary_text_color="#0a1628",
        button_primary_text_color_dark="#0a1628",
        button_secondary_background_fill="#112240",
        button_secondary_background_fill_dark="#112240",
        button_secondary_border_color="rgba(201,168,76,0.3)",
        button_secondary_border_color_dark="rgba(201,168,76,0.3)",
        button_secondary_text_color="#c9a84c",
        button_secondary_text_color_dark="#c9a84c",
        input_background_fill="#0a1628",
        input_background_fill_dark="#0a1628",
        input_border_color="rgba(201,168,76,0.2)",
        input_border_color_dark="rgba(201,168,76,0.2)",
        input_placeholder_color="#8892a4",
        input_placeholder_color_dark="#8892a4",
        shadow_drop="0 4px 24px rgba(0,0,0,0.4)",
        shadow_drop_lg="0 8px 40px rgba(0,0,0,0.5)",
        table_even_background_fill="#0a1628",
        table_even_background_fill_dark="#0a1628",
        table_odd_background_fill="#112240",
        table_odd_background_fill_dark="#112240",
    )) as demo:
    gr.Markdown("# ✝️ AI End Times Prophecy Analyzer")
    gr.Markdown("### *Watching, Waiting, and Understanding the Signs of the Times*")
    gr.Markdown(THEOLOGICAL_NOTICE)
    gr.Markdown(DISCLAIMER)

    with gr.Tabs():

        # ── Tab 1: Solar & Space Activity ──────────────────────────────────────
        with gr.Tab("☀️ Solar & Space Activity"):
            gr.Markdown("## Solar Activity, Planetary Alignments & Prophetic Signs in the Heavens")
            gr.Markdown("*Live solar data from NOAA Space Weather Prediction Center*")
            with gr.Row():
                with gr.Column(scale=1):
                    t1_context = gr.Textbox(
                        label="Additional Focus (optional)",
                        placeholder="e.g. Recent solar flares, specific planetary alignment...")
                    t1_btn = gr.Button("☀️ Analyze Solar & Space Activity", variant="primary")
                    gr.Markdown(WAIT_MSG)
                with gr.Column(scale=2):
                    t1_output = gr.Markdown(label="Solar & Space Analysis")
            t1_chart = gr.Plot(label="Solar Cycle Chart")
            t1_btn.click(analyze_solar_activity, inputs=[t1_context],
                         outputs=[t1_output, t1_chart])

        # ── Tab 2: Earth Changes Monitor ───────────────────────────────────────
        with gr.Tab("🌍 Earth Changes"):
            gr.Markdown("## Earthquakes, Volcanoes, Pole Shift & Earth Changes")
            with gr.Row():
                with gr.Column(scale=1):
                    t2_context = gr.Textbox(
                        label="Specific Earth Change to Focus On (optional)",
                        placeholder="e.g. Recent earthquake, Euphrates River drying...")
                    t2_btn = gr.Button("🌍 Analyze Earth Changes", variant="primary")
                    gr.Markdown(WAIT_MSG)
                with gr.Column(scale=2):
                    t2_output = gr.Markdown(label="Earth Changes Analysis")
            t2_chart = gr.Plot(label="Seismic Activity Chart")
            t2_btn.click(analyze_earth_changes, inputs=[t2_context],
                         outputs=[t2_output, t2_chart])

        # ── Tab 3: Biblical Prophecy Scorecard ─────────────────────────────────
        with gr.Tab("📖 Biblical Prophecy"):
            gr.Markdown("## Biblical Prophecy Scorecard — How Many Signs Are Fulfilled?")
            with gr.Row():
                with gr.Column(scale=1):
                    t3_prophecy = gr.Dropdown(
                        choices=["All Major End Times Prophecies",
                                  "Matthew 24 — Olivet Discourse",
                                  "Revelation Seals and Trumpets",
                                  "Daniel's Prophecies",
                                  "Ezekiel 38-39 Gog-Magog",
                                  "Isaiah End Times Prophecies",
                                  "Zechariah Prophecies"],
                        value="All Major End Times Prophecies",
                        label="Prophecy Focus")
                    t3_context = gr.Textbox(
                        label="Additional Context (optional)", lines=2)
                    t3_btn = gr.Button("📖 Analyze Prophecy Fulfillment", variant="primary")
                    gr.Markdown(WAIT_MSG)
                with gr.Column(scale=2):
                    t3_output = gr.Markdown(label="Prophecy Analysis")
            t3_chart = gr.Plot(label="Prophecy Fulfillment Chart")
            t3_btn.click(analyze_biblical_prophecy,
                         inputs=[t3_prophecy, t3_context],
                         outputs=[t3_output, t3_chart])

        # ── Tab 4: Middle East War Analysis ────────────────────────────────────
        with gr.Tab("⚔️ Middle East & Prophecy"):
            gr.Markdown("## Middle East Wars & Biblical Alignment — Israel, Iran, Russia")
            with gr.Row():
                with gr.Column(scale=1):
                    t4_focus = gr.Dropdown(
                        choices=["Full Middle East Analysis",
                                  "Israel-Gaza-Iran Conflict",
                                  "Ezekiel 38-39 Gog-Magog Alignment",
                                  "Isaiah 17 — Damascus",
                                  "Psalm 83 War",
                                  "Third Temple Preparations",
                                  "Abraham Accords Prophetic Significance"],
                        value="Full Middle East Analysis",
                        label="Analysis Focus")
                    t4_context = gr.Textbox(
                        label="Additional Context (optional)", lines=2)
                    t4_btn = gr.Button("⚔️ Analyze Middle East Prophecy", variant="primary")
                    gr.Markdown(WAIT_MSG)
                with gr.Column(scale=2):
                    t4_output = gr.Markdown(label="Middle East Analysis")
            t4_chart = gr.Plot(label="Nations Alignment Chart")
            t4_btn.click(analyze_middle_east,
                         inputs=[t4_focus, t4_context],
                         outputs=[t4_output, t4_chart])

        # ── Tab 5: Extra-Biblical Texts ─────────────────────────────────────────
        with gr.Tab("📜 Extra-Biblical Texts"):
            gr.Markdown("## Book of Enoch, Jubilees, Ethiopian Canon & Dead Sea Scrolls")
            with gr.Row():
                with gr.Column(scale=1):
                    t5_text = gr.Dropdown(
                        choices=["All Extra-Biblical Texts Overview",
                                  "Book of Enoch (1 Enoch) — Complete Analysis",
                                  "The Watchers & Nephilim (1 Enoch 6-16)",
                                  "Book of Jubilees — Full Analysis",
                                  "Ethiopian Orthodox Canon (81 Books)",
                                  "Dead Sea Scrolls — War Scroll",
                                  "Book of Jasher",
                                  "2 Enoch and 3 Enoch"],
                        value="All Extra-Biblical Texts Overview",
                        label="Text Focus")
                    t5_passage = gr.Textbox(
                        label="Specific Passage or Topic (optional)",
                        placeholder="e.g. 1 Enoch chapter 10, Jubilees 23, Watchers...")
                    t5_context = gr.Textbox(
                        label="Additional Context (optional)", lines=2)
                    t5_btn = gr.Button("📜 Analyze Extra-Biblical Texts", variant="primary")
                    gr.Markdown(WAIT_MSG)
                with gr.Column(scale=2):
                    t5_output = gr.Markdown(label="Extra-Biblical Texts Analysis")
            t5_chart = gr.Plot(label="Texts Overview Chart")
            t5_btn.click(analyze_extra_biblical,
                         inputs=[t5_text, t5_passage, t5_context],
                         outputs=[t5_output, t5_chart])

        # ── Tab 6: Jubilee Calendar ─────────────────────────────────────────────
        with gr.Tab("🕎 Jubilee Calendar"):
            gr.Markdown("## Jubilee Cycles, Shemitah Years & Hebrew Calendar Prophecy")
            with gr.Row():
                with gr.Column(scale=1):
                    t6_context = gr.Textbox(
                        label="Specific Jubilee Question (optional)",
                        placeholder="e.g. Is 2025-2026 a Jubilee year? Shemitah cycles...")
                    t6_btn = gr.Button("🕎 Analyze Jubilee Calendar", variant="primary")
                    gr.Markdown(WAIT_MSG)
                with gr.Column(scale=2):
                    t6_output = gr.Markdown(label="Jubilee Calendar Analysis")
            t6_chart = gr.Plot(label="Jubilee Cycle Chart")
            t6_btn.click(analyze_jubilee, inputs=[t6_context],
                         outputs=[t6_output, t6_chart])

        # ── Tab 7: Antichrist Identifier ────────────────────────────────────────
        with gr.Tab("👁️ Antichrist Identifier"):
            gr.Markdown("## Biblical Antichrist Characteristics — Theological Educational Analysis")
            gr.Markdown("> ⚠️ **IMPORTANT**: This is a scholarly theological exercise. No definitive identification is possible. Many historical figures have been analyzed this way. This follows in that academic tradition.")
            with gr.Row():
                with gr.Column(scale=1):
                    t7_figures = gr.Textbox(
                        label="Public Figures to Analyze (optional)",
                        lines=3,
                        placeholder="Enter names of world leaders, religious figures, or economic leaders to analyze against biblical criteria...")
                    t7_context = gr.Textbox(
                        label="Additional Context", lines=2)
                    t7_btn = gr.Button("👁️ Analyze Against Biblical Criteria", variant="primary")
                    gr.Markdown(WAIT_MSG)
                with gr.Column(scale=2):
                    t7_output = gr.Markdown(label="Antichrist Characteristics Analysis")
            t7_chart = gr.Plot(label="Characteristics Radar Chart")
            t7_btn.click(analyze_antichrist,
                         inputs=[t7_figures, t7_context],
                         outputs=[t7_output, t7_chart])

        # ── Tab 8: Rapture / Harpazo Tracker ────────────────────────────────────
        with gr.Tab("⚡ Rapture / Harpazo"):
            gr.Markdown("## The Harpazo (Rapture) — Signs, Views & Readiness Tracker")
            gr.Markdown("*Presenting Pre-Tribulation, Mid-Tribulation, Post-Tribulation, and Pre-Wrath views equally*")
            with gr.Row():
                with gr.Column(scale=1):
                    t8_context = gr.Textbox(
                        label="Specific Focus (optional)",
                        placeholder="e.g. Pre-Trib arguments, Harpazo timing, Deliverance of Gentiles...")
                    t8_btn = gr.Button("⚡ Analyze Harpazo Signs", variant="primary")
                    gr.Markdown(WAIT_MSG)
                with gr.Column(scale=2):
                    t8_output = gr.Markdown(label="Harpazo Analysis")
            t8_chart = gr.Plot(label="Rapture Readiness Chart")
            t8_btn.click(analyze_rapture_harpazo, inputs=[t8_context],
                         outputs=[t8_output, t8_chart])

        # ── Tab 9: Noahide Laws & Abraham Accords ───────────────────────────────
        with gr.Tab("📜 Noahide & Abraham Accords"):
            gr.Markdown("## Noahide Laws, Abraham Accords & What It Means for Gentiles")
            with gr.Row():
                with gr.Column(scale=1):
                    t9_context = gr.Textbox(
                        label="Specific Focus (optional)",
                        placeholder="e.g. Noahide laws and Christians, Abraham Accords prophetic significance, the snare for Gentiles...")
                    t9_btn = gr.Button("📜 Analyze Noahide Laws & Accords", variant="primary")
                    gr.Markdown(WAIT_MSG)
                with gr.Column(scale=2):
                    t9_output = gr.Markdown(label="Noahide Laws & Abraham Accords Analysis")
            t9_chart = gr.Plot(label="Normalization Progress Chart")
            t9_btn.click(analyze_noahide_abraham, inputs=[t9_context],
                         outputs=[t9_output, t9_chart])

        # ── Tab 10: Deliverer Timeline ──────────────────────────────────────────
        with gr.Tab("🌐 Deliverer Timeline"):
            gr.Markdown("## Deliverance of Israel & Gentiles — Where Are We in God's Plan?")
            with gr.Row():
                with gr.Column(scale=1):
                    t10_context = gr.Textbox(
                        label="Specific Focus (optional)",
                        placeholder="e.g. Fullness of the Gentiles, Times of the Gentiles, Romans 11...")
                    t10_btn = gr.Button("🌐 Analyze Deliverer Timeline", variant="primary")
                    gr.Markdown(WAIT_MSG)
                with gr.Column(scale=2):
                    t10_output = gr.Markdown(label="Deliverer Timeline Analysis")
            t10_chart = gr.Plot(label="Prophetic Timeline Chart")
            t10_btn.click(analyze_deliverer_timeline, inputs=[t10_context],
                          outputs=[t10_output, t10_chart])

        # ── Tab 11: Apocalypse Score Dashboard ──────────────────────────────────
        with gr.Tab("📊 Apocalypse Score"):
            gr.Markdown("## Overall End Times Alignment Score — All Signs Combined")
            with gr.Row():
                t11_btn = gr.Button("📊 Calculate Apocalypse Score", variant="primary")
                gr.Markdown(WAIT_MSG)
            t11_output = gr.Markdown(label="Apocalypse Score Analysis")
            t11_chart = gr.Plot(label="Apocalypse Score Dashboard")
            t11_btn.click(analyze_apocalypse_score, inputs=[],
                          outputs=[t11_output, t11_chart])


        # ── Tab 12: Prophetic Calendar ────────────────────────────────────────────
        with gr.Tab("📅 Prophetic Calendar"):
            gr.Markdown("## Prophetic Calendar — Feast Days, Blood Moons & Solar Events")
            with gr.Row():
                with gr.Column(scale=1):
                    t12_context = gr.Textbox(label="Specific Calendar Question (optional)",
                                              placeholder="e.g. upcoming feast days, blood moon dates...")
                    t12_btn = gr.Button("📅 Load Prophetic Calendar", variant="primary")
                    gr.Markdown(WAIT_MSG)
                with gr.Column(scale=2):
                    t12_output = gr.Markdown(label="Prophetic Calendar Analysis")
            t12_chart = gr.Plot(label="Calendar Chart")
            t12_btn.click(analyze_prophetic_calendar, inputs=[t12_context],
                          outputs=[t12_output, t12_chart])

        # ── Tab 13: Third Temple Tracker ──────────────────────────────────────────
        with gr.Tab("🕍 Third Temple Tracker"):
            gr.Markdown("## Third Temple Tracker — Red Heifers, Sanhedrin & Temple Institute")
            with gr.Row():
                with gr.Column(scale=1):
                    t13_context = gr.Textbox(label="Specific Focus (optional)",
                                              placeholder="e.g. red heifer status, Temple preparations...")
                    t13_btn = gr.Button("🕍 Analyze Third Temple Status", variant="primary")
                    gr.Markdown(WAIT_MSG)
                with gr.Column(scale=2):
                    t13_output = gr.Markdown(label="Third Temple Analysis")
            t13_chart = gr.Plot(label="Temple Preparation Progress")
            t13_btn.click(analyze_third_temple, inputs=[t13_context],
                          outputs=[t13_output, t13_chart])

        # ── Tab 14: One World System Monitor ─────────────────────────────────────
        with gr.Tab("🌐 One World System"):
            gr.Markdown("## One World System Monitor — WHO, CBDC, Digital ID, WEF & Revelation 13")
            with gr.Row():
                with gr.Column(scale=1):
                    t14_context = gr.Textbox(label="Specific System to Analyze (optional)",
                                              placeholder="e.g. WHO pandemic treaty, CBDC, digital ID...")
                    t14_btn = gr.Button("🌐 Analyze One World Systems", variant="primary")
                    gr.Markdown(WAIT_MSG)
                with gr.Column(scale=2):
                    t14_output = gr.Markdown(label="One World System Analysis")
            t14_chart = gr.Plot(label="System Implementation Chart")
            t14_btn.click(analyze_one_world_system, inputs=[t14_context],
                          outputs=[t14_output, t14_chart])

        # ── Tab 15: Reports & Presentations ──────────────────────────────────────
        with gr.Tab("📄 Reports & Presentations"):
            gr.Markdown("## Generate End Times Reports & Presentations")
            gr.Markdown("*PowerPoint and PDF reports combining all analysis tabs*")
            with gr.Row():
                with gr.Column(scale=1):
                    t15_report_type = gr.Dropdown(
                        choices=["Full End Times Briefing (PowerPoint)",
                                  "Church/Group Presentation Deck (PowerPoint)",
                                  "Prophecy Watchman Report (PDF)",
                                  "Signs of the Times Newsletter (PDF)",
                                  "Personal Prophecy Study Guide (PDF)"],
                        value="Full End Times Briefing (PowerPoint)",
                        label="Report Type")
                    t15_focus = gr.Dropdown(
                        choices=["Complete Overview — All Topics",
                                  "Israel and Middle East Focus",
                                  "Signs of the Times Focus",
                                  "Rapture/Harpazo Focus",
                                  "Extra-Biblical Texts Focus",
                                  "Antichrist System Focus"],
                        value="Complete Overview — All Topics",
                        label="Content Focus")
                    t15_audience = gr.Dropdown(
                        choices=["General Public", "Church Congregation",
                                  "Bible Study Group", "Personal Study", "Academic/Research"],
                        value="General Public", label="Target Audience")
                    t15_context = gr.Textbox(label="Additional Notes (optional)", lines=2)
                    t15_btn = gr.Button("📄 Generate Report", variant="primary")
                    gr.Markdown(WAIT_MSG)
                    t15_download = gr.File(label="⬇️ Download Report")
                with gr.Column(scale=2):
                    t15_output = gr.Markdown(label="Report Preview")
            t15_btn.click(generate_endtimes_report,
                          inputs=[t15_report_type, t15_focus, t15_audience, t15_context],
                          outputs=[t15_output, t15_download])

        # ── Tab 16: AI Prophecy Chat ─────────────────────────────────────────────
        with gr.Tab("💬 AI Prophecy Chat"):
            gr.Markdown("## Ask Anything About End Times Prophecy")
            gr.Markdown("*Covering all views — Pre-Trib, Mid-Trib, Post-Trib | KJV, NKJV, Ethiopian Canon | Book of Enoch, Jubilees, and more*")
            gr.ChatInterface(
                fn=prophecy_chat,
                examples=[
                    "What is the Harpazo/Rapture and when does it occur?",
                    "Explain the Book of Enoch and its relevance today",
                    "What are the Noahide Laws and what do they mean for Christians?",
                    "Where are we in the Jubilee cycle in 2025-2026?",
                    "What does Ezekiel 38-39 say about Russia and Iran?",
                    "What is the significance of the red heifers brought to Israel?",
                    "What does the Book of Jubilees say about the end times?",
                    "Explain the Abraham Accords and Daniel 9:27",
                    "What are the signs of the times in Matthew 24?",
                    "What does the Ethiopian Bible say that other Bibles don't include?",
                    "What is the Deliverance of the Israelites in Romans 11?",
                    "How does the Antichrist system compare to today's digital systems?",
                ],
                title="",
            )

    gr.Markdown(WATERMARK)


if __name__ == "__main__":
    demo.launch(server_name="0.0.0.0", server_port=int(os.environ.get("GRADIO_SERVER_PORT", 7860)), share=False, ssr_mode=False)
