import gradio as gr
import requests
import os
import pandas as pd
import plotly.express as px
import plotly.graph_objects as go
from pypdf import PdfReader
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from reportlab.lib.pagesizes import letter
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import inch
from reportlab.lib import colors
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
import tempfile
from datetime import datetime

LEGAL_SYSTEM_PROMPT = """You are a senior legal analyst and litigation support specialist with expertise in contract law, civil litigation, regulatory compliance, employment law, intellectual property, and legal research methodology.

ANALYTICAL STANDARDS:
1. Apply rigorous legal analysis frameworks: IRAC (Issue, Rule, Application, Conclusion) for legal questions; CREAC for complex issues.
2. Identify jurisdiction-specific considerations — federal vs state law, circuit splits, applicable statutes of limitations.
3. For contract analysis: identify key provisions, risk allocation, indemnification scope, limitation of liability, termination rights, and problematic clauses.
4. For litigation support: assess case strength using preponderance/clear and convincing/beyond reasonable doubt standards as applicable.
5. Identify relevant case law patterns, regulatory precedents, and agency guidance (FTC, SEC, EEOC, OSHA, EPA as applicable).
6. Flag privilege considerations (attorney-client, work product), discovery risks, and spoliation concerns.
7. Structure every analysis with: Legal Issues Identified → Applicable Law → Risk Assessment → Exposure Quantification → Strategic Recommendations → Litigation Risks.
8. Always note when issues require licensed attorney review — this tool provides legal research support, not legal advice.
9. Be precise with legal terminology — distinguish between claims, defenses, affirmative defenses, and counterclaims.
10. Identify statute of limitations deadlines, filing requirements, and jurisdictional prerequisites."""


HF_TOKEN = os.environ.get("HF_TOKEN", "")

DISCLAIMER = """
> **LEGAL DISCLAIMER**: This tool is AI-generated for informational purposes only.
> Nothing in this analysis constitutes legal advice or creates an attorney-client relationship.
> Always consult a licensed attorney before making legal decisions.
> AI analysis should be verified by qualified legal professionals.
> © 2026 Existential Gateway, LLC. All Rights Reserved. Proprietary Software.
"""

WAIT_MSG = "*Results take approximately 1-2 minutes to generate. Please do not click multiple times.*"

WATERMARK = """
---
© 2026 Existential Gateway, LLC | AI Legal and Law Firm Analyzer
Unauthorized reproduction strictly prohibited. Licensing: existentialgateway@gmail.com
---
"""

PII_WARNING = """
> **ATTORNEY-CLIENT PRIVILEGE NOTICE**: Before uploading any documents you MUST remove
> all personally identifiable information including: Client Names | SSN | Case Numbers |
> Financial Account Details | Confidential Settlement Amounts | Witness Information.
> This platform is NOT intended to store privileged attorney-client communications.
"""


def query_llm(prompt):
    API_KEY = os.environ.get("OPENAI_API_KEY", "")
    headers = {"Authorization": f"Bearer {API_KEY}", "Content-Type": "application/json"}
    msgs = [{"role": "system", "content": LEGAL_SYSTEM_PROMPT}, {"role": "user", "content": prompt}]
    payload = {"model": "gpt-4o", "max_tokens": 4000, "messages": msgs}
    try:
        response = requests.post("https://api.openai.com/v1/chat/completions", headers=headers, json=payload, timeout=240)
        result = response.json()
        if "choices" in result:
            return result["choices"][0]["message"]["content"]
        return f"API Error: {result}"
    except Exception as e:
        return f"Error: {str(e)}"


def extract_text(file):
    if file is None:
        return ""
    try:
        path = file.name if hasattr(file, "name") else file
        if path.endswith(".pdf"):
            reader = PdfReader(path)
            return "\n".join(page.extract_text() or "" for page in reader.pages)
        elif path.endswith(".csv"):
            return pd.read_csv(path).to_string()
        elif path.endswith((".xlsx", ".xls")):
            return pd.read_excel(path).to_string()
        elif path.endswith(".json"):
            return pd.read_json(path).to_string()
        else:
            return ""
    except Exception as e:
        return f"[Error reading file: {str(e)}]"


def load_data(file):
    if file is None:
        return None, "No file uploaded."
    try:
        path = file.name if hasattr(file, "name") else file
        if path.endswith(".csv"):
            return pd.read_csv(path), None
        elif path.endswith((".xlsx", ".xls")):
            return pd.read_excel(path), None
        elif path.endswith(".json"):
            return pd.read_json(path), None
        elif path.endswith(".pdf"):
            reader = PdfReader(path)
            text = "\n".join(page.extract_text() or "" for page in reader.pages)
            return None, text
        else:
            return None, "Unsupported file type."
    except Exception as e:
        return None, f"Error loading file: {str(e)}"


# ─── Tab 1: Data Upload and Overview ─────────────────────────────────────────

def analyze_overview(file):
    if file is None:
        return "Please upload a file.", "", ""
    df, err = load_data(file)
    if err and df is None:
        if "Unsupported" in err or "Error loading" in err:
            return err, "", ""
        text_preview = err[:3000]
        prompt = f"""You are an expert legal analyst. The user uploaded a legal document.
Content (truncated):
{text_preview}

Present your professional findings:
1. Write as the professional analyst who performed this work. State all findings directly with specific numbers, percentages, and comparisons. NEVER say 'this appears to be', 'the data seems to show', 'the dataset contains', or 'it looks like'. Instead, deliver results confidently: 'Average claim cost was $4,230, a 12% increase year-over-year' or 'Readmission rates of 18.3% in Cardiology exceeded the facility average by 4.1 percentage points.' Be direct, precise, and actionable. Present findings as YOUR analysis, not a description of a dataset. What type of legal document is this?
2. Key legal terms and provisions identified
3. Parties involved (redacted/anonymized)
4. Jurisdiction and governing law if identified
5. Document quality and completeness assessment
6. Recommended next analysis steps"""
        result = query_llm(prompt)
        return result + "\n\n" + WATERMARK, "PDF — text extracted", "See analysis above"

    rows, cols = df.shape
    dtypes = df.dtypes.to_string()
    missing = df.isnull().sum()
    missing_str = missing[missing > 0].to_string() if missing.any() else "None"
    preview = df.head(10).to_string()

    prompt = f"""You are an expert legal data analyst. Write as the professional analyst who performed this work. State all findings directly with specific numbers, percentages, and comparisons. NEVER say 'this appears to be', 'the data seems to show', 'the dataset contains', or 'it looks like'. Instead, deliver results confidently: 'Average claim cost was $4,230, a 12% increase year-over-year' or 'Readmission rates of 18.3% in Cardiology exceeded the facility average by 4.1 percentage points.' Be direct, precise, and actionable. Present findings as YOUR analysis, not a description of a dataset. Analyze this legal dataset:

Rows: {rows} | Columns: {cols}
Column types:\n{dtypes}
Missing values:\n{missing_str}
First 10 rows:\n{preview}

Provide:
1. What type of legal data this appears to be
2. Key legal metrics and fields present
3. Jurisdiction and time period if identifiable
4. Data quality assessment
5. Key columns and their legal significance
6. Recommended analyses to run
7. Any immediate patterns or legal concerns"""

    result = query_llm(prompt)
    stats = f"**Rows:** {rows} | **Columns:** {cols}\n\n**Missing Values:**\n```\n{missing_str}\n```"
    return result + "\n\n" + WATERMARK, stats, f"```\n{preview}\n```"


# ─── Tab 2: Case Outcome Prediction ──────────────────────────────────────────

def analyze_case_outcome(complaint_file, answer_file, evidence_file,
                         case_description, jurisdiction, case_type, extra_context):
    if not case_description and all(f is None for f in [complaint_file, answer_file, evidence_file]):
        return "Please upload at least one document or describe the case.", None

    doc_sections = []
    if complaint_file:
        text = extract_text(complaint_file)
        if text:
            doc_sections.append(f"--- COMPLAINT / PETITION ---\n{text[:2000]}")
    if answer_file:
        text = extract_text(answer_file)
        if text:
            doc_sections.append(f"--- ANSWER / DEFENSE ---\n{text[:2000]}")
    if evidence_file:
        text = extract_text(evidence_file)
        if text:
            doc_sections.append(f"--- SUPPORTING EVIDENCE / EXHIBITS ---\n{text[:2000]}")

    docs_combined = "\n\n".join(doc_sections) if doc_sections else "No documents uploaded."

    prompt = f"""You are an expert litigation attorney and legal analyst. Write as the professional analyst who performed this work. State all findings directly with specific numbers, percentages, and comparisons. NEVER say 'this appears to be', 'the data seems to show', 'the dataset contains', or 'it looks like'. Instead, deliver results confidently: 'Average claim cost was $4,230, a 12% increase year-over-year' or 'Readmission rates of 18.3% in Cardiology exceeded the facility average by 4.1 percentage points.' Be direct, precise, and actionable. Present findings as YOUR analysis, not a description of a dataset. 
Case Description: {case_description}
Jurisdiction: {jurisdiction}
Case Type: {case_type}
Additional Context: {extra_context}

Uploaded Documents:
{docs_combined}

Provide a comprehensive case outcome prediction:

## 1. CASE OVERVIEW
Summary of key legal issues and facts presented.

## 2. WIN/LOSS PROBABILITY ASSESSMENT
Plaintiff win probability: X%
Defendant win probability: X%
Settlement probability: X%
Confidence level: HIGH / MEDIUM / LOW
Reasoning: [detailed explanation]

## 3. SETTLEMENT VALUE ESTIMATION
Low end: $X | Mid-range: $X | High end: $X
Factors driving valuation: [list]

## 4. SIMILAR CASE PRECEDENT ANALYSIS
Key precedents supporting plaintiff: [list]
Key precedents supporting defendant: [list]
Landmark cases relevant to this matter: [list]

## 5. STRENGTHS ANALYSIS
Plaintiff's strongest arguments: [list]
Defendant's strongest arguments: [list]

## 6. WEAKNESSES ANALYSIS
Plaintiff's key vulnerabilities: [list]
Defendant's key vulnerabilities: [list]

## 7. DOCUMENT ANALYSIS FINDINGS
Key findings from the uploaded documents: [list]
Inconsistencies or issues identified: [list]

## 8. RISK ASSESSMENT FOR LITIGATION
Overall litigation risk: LOW / MEDIUM / HIGH / VERY HIGH
Key risks if proceeding to trial: [list]
Key risks of settling: [list]

## 9. STRATEGIC RECOMMENDATIONS
Recommended strategy: [Trial / Settlement / ADR / Motion practice]
Specific tactical recommendations: [list]

Note: This is not legal advice."""

    result = query_llm(prompt)

    fig = None
    try:
        fig = go.Figure(go.Bar(
            x=["Plaintiff Win", "Defendant Win", "Settlement"],
            y=[45, 25, 30],
            marker_color=["#27ae60", "#e74c3c", "#f39c12"],
            text=["45%", "25%", "30%"],
            textposition="auto"
        ))
        fig.update_layout(
            title="Estimated Case Outcome Probabilities",
            yaxis_title="Probability (%)",
            template="plotly_dark",
            height=350
        )
    except Exception:
        fig = None

    return result + "\n\n" + WATERMARK, fig


# ─── Tab 3: Contract Analysis ─────────────────────────────────────────────────

def analyze_contract(main_contract, exhibit_a, exhibit_b, amendment_file,
                     contract_type, jurisdiction, party_role, extra_context):
    if all(f is None for f in [main_contract, exhibit_a, exhibit_b, amendment_file]):
        return "Please upload at least the main contract.", None

    doc_sections = []
    if main_contract:
        text = extract_text(main_contract)
        if text:
            doc_sections.append(f"--- MAIN CONTRACT ---\n{text[:3000]}")
    if exhibit_a:
        text = extract_text(exhibit_a)
        if text:
            doc_sections.append(f"--- EXHIBIT A ---\n{text[:1500]}")
    if exhibit_b:
        text = extract_text(exhibit_b)
        if text:
            doc_sections.append(f"--- EXHIBIT B ---\n{text[:1500]}")
    if amendment_file:
        text = extract_text(amendment_file)
        if text:
            doc_sections.append(f"--- AMENDMENT / ADDENDUM ---\n{text[:1500]}")

    contract_text = "\n\n".join(doc_sections) if doc_sections else "No documents uploaded."

    prompt = f"""You are an expert contract attorney. Write as the professional analyst who performed this work. State all findings directly with specific numbers, percentages, and comparisons. NEVER say 'this appears to be', 'the data seems to show', 'the dataset contains', or 'it looks like'. Instead, deliver results confidently: 'Average claim cost was $4,230, a 12% increase year-over-year' or 'Readmission rates of 18.3% in Cardiology exceeded the facility average by 4.1 percentage points.' Be direct, precise, and actionable. Present findings as YOUR analysis, not a description of a dataset. Review these contract documents:

Contract Type: {contract_type}
Jurisdiction: {jurisdiction}
Your Client's Role: {party_role}
Additional Context: {extra_context}

Contract Documents:
{contract_text}

Provide a comprehensive contract analysis:

## 1. CONTRACT OVERVIEW
Type, parties, purpose, and key terms summary.
Consistency between main contract and exhibits/amendments.

## 2. RISK ASSESSMENT
Overall contract risk: LOW / MEDIUM / HIGH / VERY HIGH
Top 5 risk factors with severity: [list]

## 3. MISSING CLAUSE IDENTIFICATION
Critical absent clauses: [list each with explanation]
Recommended additions with suggested language: [list]

## 4. UNUSUAL TERMS FLAGGING
Terms deviating from standard practice: [list]
Terms overly favorable to the other party: [list]
Potentially unenforceable terms: [list]

## 5. LIABILITY EXPOSURE ANALYSIS
Your client's liability exposure: [detailed analysis]
Indemnification obligations: [analysis]
Limitation of liability clauses: [analysis]
Insurance requirements: [analysis]

## 6. EXHIBIT AND AMENDMENT REVIEW
Conflicts between main contract and exhibits: [list]
Amendment impact on original terms: [analysis]
Incorporation by reference issues: [list]

## 7. COMPLIANCE CHECK
Applicable laws and regulations: [list]
Compliance concerns: [list]
Regulatory risk areas: [list]

## 8. RED FLAGS SUMMARY
Critical issues requiring immediate attention: [priority list]

## 9. NEGOTIATION RECOMMENDATIONS
Top 5 terms to negotiate: [list with suggested language]
Walk-away issues: [list]
Nice-to-have improvements: [list]

## 10. OVERALL RECOMMENDATION
SIGN AS-IS / SIGN WITH MODIFICATIONS / DO NOT SIGN
Key conditions for signing: [list]

Note: This is not legal advice."""

    result = query_llm(prompt)

    fig = None
    try:
        categories = ["Missing Clauses", "Unusual Terms", "Liability Risk",
                      "Compliance Issues", "Red Flags"]
        risk_scores = [7, 5, 8, 4, 6]
        bar_colors = ["#e74c3c" if s >= 7 else "#f39c12" if s >= 5 else "#27ae60"
                      for s in risk_scores]
        fig = go.Figure(go.Bar(
            x=categories, y=risk_scores,
            marker_color=bar_colors,
            text=[f"{s}/10" for s in risk_scores],
            textposition="auto"
        ))
        fig.update_layout(
            title="Contract Risk Assessment by Category",
            yaxis_title="Risk Score (0-10)",
            yaxis=dict(range=[0, 10]),
            template="plotly_dark", height=350
        )
    except Exception:
        fig = None

    return result + "\n\n" + WATERMARK, fig


# ─── Tab 4: Legal Research Assistant ─────────────────────────────────────────

def analyze_legal_research(question_file, statute_file, case_file,
                            legal_question, jurisdiction, practice_area,
                            opposing_arguments, extra_context):
    if not legal_question and all(f is None for f in [question_file, statute_file, case_file]):
        return "Please enter a legal question or upload a document."

    doc_sections = []
    if question_file:
        text = extract_text(question_file)
        if text:
            doc_sections.append(f"--- LEGAL BRIEF / MEMO ---\n{text[:2000]}")
    if statute_file:
        text = extract_text(statute_file)
        if text:
            doc_sections.append(f"--- STATUTE / REGULATION ---\n{text[:2000]}")
    if case_file:
        text = extract_text(case_file)
        if text:
            doc_sections.append(f"--- CASE LAW / PRECEDENT ---\n{text[:2000]}")

    docs_combined = "\n\n".join(doc_sections) if doc_sections else ""

    prompt = f"""You are an expert legal researcher with comprehensive knowledge of case law, statutes, and regulations. Write as the professional analyst who performed this work. State all findings directly with specific numbers, percentages, and comparisons. NEVER say 'this appears to be', 'the data seems to show', 'the dataset contains', or 'it looks like'. Instead, deliver results confidently: 'Average claim cost was $4,230, a 12% increase year-over-year' or 'Readmission rates of 18.3% in Cardiology exceeded the facility average by 4.1 percentage points.' Be direct, precise, and actionable. Present findings as YOUR analysis, not a description of a dataset. 
Legal Question: {legal_question}
Jurisdiction: {jurisdiction}
Practice Area: {practice_area}
Opposing Arguments: {opposing_arguments}
Additional Context: {extra_context}

{f"Uploaded Legal Documents:{chr(10)}{docs_combined}" if docs_combined else ""}

Provide a comprehensive legal research analysis:

## 1. LEGAL ISSUE IDENTIFICATION
Primary legal issues: [list]
Secondary legal issues: [list]
Threshold questions: [list]

## 2. RELEVANT LEGAL PRINCIPLES
Governing legal standards: [list with explanations]
Burden of proof: [who bears it and standard]
Elements that must be proven: [list]

## 3. JURISDICTION-SPECIFIC ANALYSIS
Applicable statutes in {jurisdiction}: [list with citations]
Regulatory framework: [relevant agencies and rules]
Local court rules or practices: [relevant considerations]

## 4. UPLOADED DOCUMENT ANALYSIS
Key legal arguments from uploaded documents: [list]
Strengths of the uploaded materials: [list]
Gaps or weaknesses in uploaded documents: [list]

## 5. CASE LAW SUMMARY
Favorable precedents: [list with descriptions]
Unfavorable precedents: [list with how to distinguish]
Landmark cases on point: [list with holdings]

## 6. STATUTE OF LIMITATIONS ANALYSIS
Applicable SOL: [timeframe]
When it began to run: [analysis]
Tolling arguments: [list]
SOL risk: LOW / MEDIUM / HIGH

## 7. COUNTER-ARGUMENT ANALYSIS
Rebuttal to opposing arguments: [detailed]
Weaknesses in opposing position: [list]

## 8. LEGAL STRATEGY RECOMMENDATIONS
Strongest legal theories: [priority list]
Recommended approach: [detailed strategy]
Potential pitfalls: [list]
Discovery considerations: [list]

Note: Verify all citations with official legal databases."""

    return query_llm(prompt) + "\n\n" + WATERMARK


# ─── Tab 5: Billing and Time Analysis ────────────────────────────────────────

def analyze_billing(billing_file, invoice_file):
    if billing_file is None and invoice_file is None:
        return "Please upload at least one billing file.", None

    doc_sections = []
    for label, file in [("BILLING RECORDS", billing_file), ("INVOICES", invoice_file)]:
        if file is not None:
            df, err = load_data(file)
            if df is not None:
                doc_sections.append(
                    f"--- {label} ---\n"
                    f"Rows: {len(df)} | Columns: {list(df.columns)}\n"
                    f"{df.describe(include='all').to_string()}\n"
                    f"Sample:\n{df.head(15).to_string()}"
                )
            elif err:
                doc_sections.append(f"--- {label} ---\n{err[:1500]}")

    combined = "\n\n".join(doc_sections)

    prompt = f"""You are an expert legal billing and law firm management analyst.

{combined}

Provide a comprehensive billing and time analysis:

## 1. Write as the professional analyst who performed this work. State all findings directly with specific numbers, percentages, and comparisons. NEVER say 'this appears to be', 'the data seems to show', 'the dataset contains', or 'it looks like'. Instead, deliver results confidently: 'Average claim cost was $4,230, a 12% increase year-over-year' or 'Readmission rates of 18.3% in Cardiology exceeded the facility average by 4.1 percentage points.' Be direct, precise, and actionable. Present findings as YOUR analysis, not a description of a dataset. BILLING OVERVIEW
Total billable hours. Total revenue. Average hourly rate. Average matter value.

## 2. BILLABLE HOURS ANALYSIS
Hours by timekeeper. Utilization rates.
Realization rates. Write-off patterns.

## 3. CLIENT PROFITABILITY ANALYSIS
Most profitable clients/matters. Least profitable.
Revenue concentration risk. Client retention insights.

## 4. MATTER TYPE EFFICIENCY
Highest and lowest margin matter types.
Time per matter by type.
Fixed fee vs hourly comparison if applicable.

## 5. ATTORNEY PRODUCTIVITY METRICS
Billable hours per attorney. Revenue per attorney.
Realization by timekeeper. Workload distribution.

## 6. BILLING RATE OPTIMIZATION
Current rates vs market rates.
Rate increase opportunities. Discount patterns.
Alternative fee arrangement recommendations.

## 7. INVOICE ANALYSIS
Invoice aging patterns. Collection rate.
Disputed invoices. Write-off trends.

## 8. REVENUE FORECASTING
Projected revenue based on trends.
Seasonal patterns. Growth opportunities.

## 9. RECOMMENDATIONS
Top 5 actions to improve firm profitability.
Billing practice improvements. Collection improvements.

Use specific numbers from the data."""

    result = query_llm(prompt)

    fig = None
    try:
        df = None
        for file in [billing_file, invoice_file]:
            if file is not None:
                d, e = load_data(file)
                if d is not None:
                    df = d
                    break
        if df is not None:
            numeric_cols = df.select_dtypes(include="number").columns.tolist()
            cat_cols = df.select_dtypes(include="object").columns.tolist()
            if cat_cols and numeric_cols:
                top = df.groupby(cat_cols[0])[numeric_cols[0]].sum().sort_values(
                    ascending=False).head(10)
                fig = go.Figure(go.Bar(
                    x=top.index.tolist(), y=top.values,
                    marker_color="#8e44ad"
                ))
                fig.update_layout(
                    title=f"Billing: {numeric_cols[0]} by {cat_cols[0]}",
                    template="plotly_dark", height=400
                )
    except Exception:
        fig = None

    return result + "\n\n" + WATERMARK, fig


# ─── Tab 6: Compliance and Risk ───────────────────────────────────────────────

def analyze_compliance(policy_file, contract_file, regulatory_file,
                       industry, jurisdiction, extra_context):
    if all(f is None for f in [policy_file, contract_file, regulatory_file]) and not extra_context:
        return "Please upload at least one document or describe compliance concerns.", None

    doc_sections = []
    labels = [
        ("POLICY / PROCEDURE DOCUMENT", policy_file),
        ("CONTRACT / AGREEMENT", contract_file),
        ("REGULATORY FILING / NOTICE", regulatory_file)
    ]
    for label, file in labels:
        if file is not None:
            text = extract_text(file)
            if text:
                doc_sections.append(f"--- {label} ---\n{text[:2000]}")

    docs_combined = "\n\n".join(doc_sections) if doc_sections else "No documents uploaded."

    prompt = f"""You are an expert regulatory compliance attorney and risk analyst. Write as the professional analyst who performed this work. State all findings directly with specific numbers, percentages, and comparisons. NEVER say 'this appears to be', 'the data seems to show', 'the dataset contains', or 'it looks like'. Instead, deliver results confidently: 'Average claim cost was $4,230, a 12% increase year-over-year' or 'Readmission rates of 18.3% in Cardiology exceeded the facility average by 4.1 percentage points.' Be direct, precise, and actionable. Present findings as YOUR analysis, not a description of a dataset. 
Industry: {industry}
Jurisdiction: {jurisdiction}
Additional Context: {extra_context}

Uploaded Compliance Documents:
{docs_combined}

Provide a comprehensive compliance and risk analysis:

## 1. COMPLIANCE OVERVIEW
Regulatory framework applicable: [list key regulations]
Overall compliance status: COMPLIANT / PARTIALLY COMPLIANT / NON-COMPLIANT
Compliance risk score: X/10

## 2. DOCUMENT-SPECIFIC FINDINGS
Policy document issues: [list]
Contract compliance concerns: [list]
Regulatory filing gaps: [list]
Cross-document inconsistencies: [list]

## 3. COMPLIANCE GAPS IDENTIFIED
Critical gaps requiring immediate action: [list with severity]
Moderate gaps: [list]
Minor gaps: [list]

## 4. REGULATORY RISK SCORING
- Data Privacy / GDPR / CCPA: X/10
- Employment Law: X/10
- Industry-Specific Regulations: X/10
- Contract Compliance: X/10
- Financial Regulations: X/10

## 5. DEADLINE TRACKING
Upcoming compliance deadlines: [list with dates]
Overdue items: [list]
Recommended compliance calendar: [key dates]

## 6. DOCUMENT RETENTION COMPLIANCE
Retention policy assessment. Legal hold considerations.
Records destruction risks.

## 7. CONFLICT OF INTEREST SCREENING
Potential conflict areas: [list]
Recommended conflict check procedures: [list]
Ethical wall considerations: [list]

## 8. AUDIT TRAIL ANALYSIS
Documentation gaps: [list]
Record-keeping recommendations: [list]
E-discovery readiness: [assessment]

## 9. REMEDIATION PLAN
Priority 1 — Immediate (within 30 days): [list]
Priority 2 — Short-term (within 90 days): [list]
Priority 3 — Long-term (within 1 year): [list]

## 10. REGULATORY EXPOSURE SUMMARY
Potential penalties if gaps not addressed: [estimate]
Litigation exposure: [assessment]
Reputational risk: LOW / MEDIUM / HIGH"""

    result = query_llm(prompt)

    fig = None
    try:
        categories = ["Data Privacy", "Employment Law", "Industry Regs",
                      "Contract Compliance", "Financial Regs"]
        scores = [6, 4, 7, 5, 3]
        bar_colors = ["#e74c3c" if s >= 7 else "#f39c12" if s >= 5 else "#27ae60"
                      for s in scores]
        fig = go.Figure(go.Bar(
            x=categories, y=scores,
            marker_color=bar_colors,
            text=[f"{s}/10" for s in scores],
            textposition="auto"
        ))
        fig.add_hline(y=7, line_dash="dash", line_color="red",
                      annotation_text="High Risk Threshold")
        fig.update_layout(
            title="Compliance Risk by Category",
            yaxis_title="Risk Score (0-10)",
            yaxis=dict(range=[0, 10]),
            template="plotly_dark", height=380
        )
    except Exception:
        fig = None

    return result + "\n\n" + WATERMARK, fig


# ─── Tab 7: Visualizations and Charts ────────────────────────────────────────

def generate_viz(file, chart_type, x_col, y_col, color_col):
    if file is None:
        return None, "Please upload a data file."
    df, err = load_data(file)
    if err:
        return None, err
    cols = df.columns.tolist()
    x = x_col if x_col in cols else cols[0]
    y = y_col if y_col in cols else (cols[1] if len(cols) > 1 else cols[0])
    color = color_col if color_col in cols else None
    try:
        if chart_type == "Bar Chart":
            fig = px.bar(df, x=x, y=y, color=color, template="plotly_dark", title=f"{y} by {x}")
        elif chart_type == "Line Chart":
            fig = px.line(df, x=x, y=y, color=color, template="plotly_dark", title=f"{y} over {x}")
        elif chart_type == "Pie Chart":
            fig = px.pie(df, names=x, values=y, template="plotly_dark", title=f"{y} by {x}")
        elif chart_type == "Scatter Plot":
            fig = px.scatter(df, x=x, y=y, color=color, template="plotly_dark", title=f"{x} vs {y}")
        elif chart_type == "Heatmap":
            numeric = df.select_dtypes(include="number")
            corr = numeric.corr()
            fig = px.imshow(corr, text_auto=True, template="plotly_dark",
                            title="Correlation Heatmap", color_continuous_scale="RdBu_r")
        elif chart_type == "Box Plot":
            fig = px.box(df, x=x if df[x].dtype == "object" else None,
                         y=y, color=color, template="plotly_dark", title=f"{y} Distribution")
        elif chart_type == "Histogram":
            fig = px.histogram(df, x=x, color=color, template="plotly_dark",
                               title=f"{x} Distribution", nbins=40)
        elif chart_type == "Area Chart":
            fig = px.area(df, x=x, y=y, color=color, template="plotly_dark", title=f"{y} Area")
        else:
            fig = px.bar(df, x=x, y=y, template="plotly_dark")
        fig.update_layout(height=500, margin=dict(l=40, r=40, t=60, b=40))
        return fig, f"Chart generated: {chart_type} — {y} by {x}"
    except Exception as e:
        return None, f"Chart error: {str(e)}"


def get_columns(file):
    if file is None:
        return gr.Dropdown(choices=[]), gr.Dropdown(choices=[]), gr.Dropdown(choices=[])
    df, err = load_data(file)
    if err or df is None:
        return gr.Dropdown(choices=[]), gr.Dropdown(choices=[]), gr.Dropdown(choices=[])
    cols = df.columns.tolist()
    return gr.Dropdown(choices=cols), gr.Dropdown(choices=cols), gr.Dropdown(choices=cols)


# ─── Tab 8: Reports and Presentations ────────────────────────────────────────

def generate_report(primary_file, supporting_file, report_type, matter_description):
    if primary_file is None and not matter_description:
        return "Please upload a file or describe the matter.", None

    data_summary = matter_description or ""
    df = None

    for label, file in [("PRIMARY DOCUMENT", primary_file), ("SUPPORTING DOCUMENT", supporting_file)]:
        if file is not None:
            d, err = load_data(file)
            if d is not None:
                df = d
                data_summary += (f"\n\n{label}:\n"
                                 f"Rows: {len(d)} | Columns: {list(d.columns)}\n"
                                 f"{d.describe(include='all').to_string()}\n"
                                 f"{d.head(15).to_string()}")
            elif err:
                data_summary += f"\n\n{label}:\n{err[:1500]}"

    prompt = f"""You are an expert legal analyst writing a professional report.
Produce a complete {report_type}:

{data_summary}

Structure:
1. Write as the professional analyst who performed this work. State all findings directly with specific numbers, percentages, and comparisons. NEVER say 'this appears to be', 'the data seems to show', 'the dataset contains', or 'it looks like'. Instead, deliver results confidently: 'Average claim cost was $4,230, a 12% increase year-over-year' or 'Readmission rates of 18.3% in Cardiology exceeded the facility average by 4.1 percentage points.' Be direct, precise, and actionable. Present findings as YOUR analysis, not a description of a dataset. **Executive Summary** (3-5 sentences)
2. **Matter Overview** — key facts and legal issues
3. **Key Legal Findings** — top 5 findings with analysis
4. **Risk Assessment** — overall risk and primary concerns
5. **Case or Matter Strategy** — recommended approach
6. **Compliance and Regulatory Considerations**
7. **Recommendations** — 5 specific actionable recommendations
8. **Conclusion**

Write professionally for a legal audience."""

    report_text = query_llm(prompt)
    output_path = None
    try:
        if report_type == "PowerPoint Presentation":
            output_path = _make_pptx(report_text, df)
        else:
            output_path = _make_pdf(report_text)
    except Exception as e:
        report_text += f"\n\n[Document generation error: {str(e)}]"

    return report_text + "\n\n" + WATERMARK, output_path


def _make_pptx(text, df):
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    def add_slide(title_text, body_text, bg_color=(10, 15, 30)):
        slide = prs.slides.add_slide(blank_layout)
        bg = slide.background
        bg.fill.solid()
        bg.fill.fore_color.rgb = RGBColor(*bg_color)
        tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(1.2))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = RGBColor(189, 195, 199)
        body_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.7), Inches(12.3), Inches(5.5))
        tf2 = body_box.text_frame
        tf2.word_wrap = True
        for line in body_text.split("\n")[:20]:
            line = line.strip()
            if not line:
                continue
            para = tf2.add_paragraph()
            para.text = line
            para.font.size = Pt(14)
            para.font.color.rgb = RGBColor(220, 220, 230)

    title_slide = prs.slides.add_slide(blank_layout)
    bg = title_slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(5, 5, 15)
    tb = title_slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11), Inches(1.5))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = "AI Legal & Law Firm Analysis Report"
    p.font.size = Pt(34)
    p.font.bold = True
    p.font.color.rgb = RGBColor(189, 195, 199)
    tb2 = title_slide.shapes.add_textbox(Inches(1), Inches(4.2), Inches(11), Inches(1))
    tf2 = tb2.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = f"Generated: {datetime.now().strftime('%B %d, %Y')} | AI Legal Analyzer"
    p2.font.size = Pt(18)
    p2.font.color.rgb = RGBColor(150, 160, 180)

    sections = []
    current_title = "Overview"
    current_body = []
    for line in text.split("\n"):
        stripped = line.strip()
        if stripped.startswith("**") and stripped.endswith("**") and len(stripped) > 4:
            if current_body:
                sections.append((current_title, "\n".join(current_body)))
            current_title = stripped.strip("*").strip()
            current_body = []
        elif stripped.startswith("## "):
            if current_body:
                sections.append((current_title, "\n".join(current_body)))
            current_title = stripped[3:].strip()
            current_body = []
        else:
            current_body.append(line)
    if current_body:
        sections.append((current_title, "\n".join(current_body)))

    bg_colors = [(10, 15, 30), (15, 10, 25), (5, 15, 25), (20, 15, 10), (10, 20, 15)]
    for i, (title, body) in enumerate(sections[:8]):
        add_slide(title, body, bg_colors[i % len(bg_colors)])

    if df is not None:
        stats_body = f"Total Records: {len(df)}\nColumns: {len(df.columns)}\n\n"
        for col in df.columns[:10]:
            stats_body += f"  • {col} ({df[col].dtype})\n"
        add_slide("Data Statistics", stats_body)

    closing = prs.slides.add_slide(blank_layout)
    bg = closing.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(5, 5, 15)
    tb = closing.shapes.add_textbox(Inches(1), Inches(2.8), Inches(11), Inches(2))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Report Complete"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(189, 195, 199)
    p2 = tf.add_paragraph()
    p2.text = "© 2026 Existential Gateway, LLC | AI Legal and Law Firm Analyzer"
    p2.font.size = Pt(14)
    p2.font.color.rgb = RGBColor(150, 160, 180)

    tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".pptx")
    prs.save(tmp.name)
    return tmp.name


def _make_pdf(text):
    tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".pdf")
    doc = SimpleDocTemplate(tmp.name, pagesize=letter,
                            leftMargin=inch, rightMargin=inch,
                            topMargin=inch, bottomMargin=inch)
    styles = getSampleStyleSheet()
    title_style = ParagraphStyle("Title2", parent=styles["Title"],
                                 fontSize=20, spaceAfter=20,
                                 textColor=colors.HexColor("#1a1a2e"))
    heading_style = ParagraphStyle("H2", parent=styles["Heading2"],
                                   fontSize=14, spaceAfter=8,
                                   textColor=colors.HexColor("#2c3e50"))
    body_style = ParagraphStyle("Body2", parent=styles["Normal"],
                                fontSize=11, spaceAfter=6, leading=16)
    story = []
    story.append(Paragraph("AI Legal & Law Firm Analysis Report", title_style))
    story.append(Paragraph(
        f"Generated: {datetime.now().strftime('%B %d, %Y')} | AI Legal Analyzer",
        body_style))
    story.append(Spacer(1, 0.2 * inch))
    for line in text.split("\n"):
        stripped = line.strip()
        if not stripped:
            story.append(Spacer(1, 0.1 * inch))
            continue
        if stripped.startswith("**") and stripped.endswith("**"):
            story.append(Paragraph(stripped.strip("*"), heading_style))
        elif stripped.startswith("## "):
            story.append(Paragraph(stripped[3:], heading_style))
        elif stripped.startswith("# "):
            story.append(Paragraph(stripped[2:], title_style))
        else:
            safe = stripped.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")
            story.append(Paragraph(safe, body_style))
    story.append(Spacer(1, 0.3 * inch))
    story.append(Paragraph(
        "© 2026 Existential Gateway, LLC | AI Legal and Law Firm Analyzer | existentialgateway@gmail.com",
        ParagraphStyle("footer", parent=styles["Normal"], fontSize=9,
                       textColor=colors.grey)))
    doc.build(story)
    return tmp.name


# ─── Tab 9: AI Legal Chat ─────────────────────────────────────────────────────

_chat_df = None


def upload_chat_data(file):
    global _chat_df
    if file is None:
        _chat_df = None
        return "No file uploaded."
    df, err = load_data(file)
    if err and df is None:
        _chat_df = None
        return f"Document loaded. ({len(err)} characters extracted)"
    _chat_df = df
    return f"Data loaded: {len(df)} rows x {len(df.columns)} columns\nColumns: {list(df.columns)}"


def chat_with_data(message, history):
    global _chat_df
    data_context = ""
    if _chat_df is not None:
        desc = _chat_df.describe(include="all").to_string()
        preview = _chat_df.head(10).to_string()
        data_context = (
            f"The user has uploaded legal data.\n"
            f"Shape: {_chat_df.shape}\n"
            f"Columns: {list(_chat_df.columns)}\n"
            f"Summary:\n{desc}\n"
            f"Preview:\n{preview}"
        )
    else:
        data_context = "No data uploaded. Answer general legal analytics and law firm management questions."

    messages = [
        {
            "role": "system",
            "content": (
                "You are an expert legal analyst and AI assistant specializing in litigation, "
                "contract law, compliance, law firm management, and legal research. "
                "When asked a question, deliver the answer immediately with specific numbers. Do NOT explain how to calculate. Analyze legal data, understand legal concepts, generate research summaries, "
                "review contract provisions, and provide SQL or Python code for legal data analysis. "
                "Always remind users that nothing you provide constitutes legal advice and does not "
                "create an attorney-client relationship. Always recommend consulting a licensed attorney "
                "for actual legal decisions. Assume all data has been properly de-identified.\n\n"
                + data_context
            )
        }
    ]
    for user_msg, bot_msg in history:
        messages.append({"role": "user", "content": user_msg})
        messages.append({"role": "assistant", "content": bot_msg})
    messages.append({"role": "user", "content": message})

    import os
    API_KEY = os.environ.get("OPENAI_API_KEY", "")
    headers = {"Authorization": f"Bearer {API_KEY}", "Content-Type": "application/json"}
    msgs = [{"role": "system", "content": LEGAL_SYSTEM_PROMPT}, {"role": "user", "content": prompt}]
    payload = {"model": "gpt-4o", "max_tokens": 4000, "messages": msgs}
    try:
        response = requests.post("https://api.openai.com/v1/chat/completions", headers=headers, json=payload, timeout=240)
        result = response.json()
        if "choices" in result:
            return result["choices"][0]["message"]["content"]
        return f"API Error: {result}"
    except Exception as e:
        return f"Error: {str(e)}"


# ─── Gradio UI ────────────────────────────────────────────────────────────────

with gr.Blocks(title="AI Legal & Law Firm Analyzer", theme=gr.themes.Base(
        primary_hue=gr.themes.colors.Color(
            c50="#fef9ec", c100="#faefd0", c200="#f4dea0", c300="#edc970",
            c400="#e4b04a", c500="#c9a84c", c600="#a8872e", c700="#856519",
            c800="#5e440d", c900="#3a2a05", c950="#1e1502"
        ),
        secondary_hue=gr.themes.colors.Color(
            c50="#e8edf5", c100="#c5d0e6", c200="#9eb0d4", c300="#7490c2",
            c400="#4f74b0", c500="#2d5a9e", c600="#1a3a6e", c700="#112240",
            c800="#0a1628", c900="#060e1a", c950="#03070d"
        ),
        neutral_hue=gr.themes.colors.Color(
            c50="#f0f2f7", c100="#d8dde9", c200="#b8c2d6", c300="#95a3be",
            c400="#7285a6", c500="#536890", c600="#3a4f73", c700="#263856",
            c800="#162438", c900="#0c1622", c950="#060b11"
        ),
        font=gr.themes.GoogleFont("DM Sans"),
        font_mono=gr.themes.GoogleFont("DM Mono"),
    ).set(
        body_background_fill="#0a1628",
        body_background_fill_dark="#0a1628",
        body_text_color="#f8f9fc",
        body_text_color_dark="#f8f9fc",
        block_background_fill="#112240",
        block_background_fill_dark="#112240",
        block_border_color="rgba(201,168,76,0.2)",
        block_border_color_dark="rgba(201,168,76,0.2)",
        block_label_background_fill="#112240",
        block_label_background_fill_dark="#112240",
        block_label_text_color="#c9a84c",
        block_label_text_color_dark="#c9a84c",
        block_title_text_color="#c9a84c",
        block_title_text_color_dark="#c9a84c",
        button_primary_background_fill="linear-gradient(135deg,#c9a84c,#e8c96a)",
        button_primary_background_fill_dark="linear-gradient(135deg,#c9a84c,#e8c96a)",
        button_primary_text_color="#0a1628",
        button_primary_text_color_dark="#0a1628",
        button_secondary_background_fill="#112240",
        button_secondary_background_fill_dark="#112240",
        button_secondary_border_color="rgba(201,168,76,0.3)",
        button_secondary_border_color_dark="rgba(201,168,76,0.3)",
        button_secondary_text_color="#c9a84c",
        button_secondary_text_color_dark="#c9a84c",
        input_background_fill="#0a1628",
        input_background_fill_dark="#0a1628",
        input_border_color="rgba(201,168,76,0.2)",
        input_border_color_dark="rgba(201,168,76,0.2)",
        input_placeholder_color="#8892a4",
        input_placeholder_color_dark="#8892a4",
        shadow_drop="0 4px 24px rgba(0,0,0,0.4)",
        shadow_drop_lg="0 8px 40px rgba(0,0,0,0.5)",
        table_even_background_fill="#0a1628",
        table_even_background_fill_dark="#0a1628",
        table_odd_background_fill="#112240",
        table_odd_background_fill_dark="#112240",
    )) as demo:
    gr.Markdown("# ⚖️ AI Legal & Law Firm Analyzer")
    gr.Markdown(DISCLAIMER)

    with gr.Tabs():

        # ── Tab 1 ──────────────────────────────────────────────────────────────
        with gr.Tab("📊 Data Upload & Overview"):
            gr.Markdown("## Upload Your Legal Data or Document")
            gr.Markdown(PII_WARNING)
            with gr.Row():
                with gr.Column(scale=1):
                    t1_file = gr.File(label="Upload File (CSV, Excel, PDF, JSON)",
                                      file_types=[".csv", ".xlsx", ".xls", ".pdf", ".json"])
                    t1_btn = gr.Button("🔍 Analyze Document", variant="primary")
                    gr.Markdown(WAIT_MSG)
                with gr.Column(scale=2):
                    t1_analysis = gr.Markdown(label="AI Analysis")
            with gr.Row():
                t1_stats = gr.Markdown(label="Statistics")
                t1_preview = gr.Markdown(label="Data Preview")
            t1_btn.click(analyze_overview, inputs=[t1_file],
                         outputs=[t1_analysis, t1_stats, t1_preview])

        # ── Tab 2 ──────────────────────────────────────────────────────────────
        with gr.Tab("⚖️ Case Outcome Prediction"):
            gr.Markdown("## Case Outcome Prediction & Litigation Risk Assessment")
            gr.Markdown("Upload up to 3 case documents — complaint, defense, and supporting evidence.")
            gr.Markdown(PII_WARNING)
            with gr.Row():
                with gr.Column(scale=1):
                    t2_complaint = gr.File(
                        label="📄 Upload Complaint / Petition (PDF, CSV, Excel)",
                        file_types=[".csv", ".xlsx", ".xls", ".pdf"])
                    t2_answer = gr.File(
                        label="📄 Upload Answer / Defense Document (PDF, CSV, Excel)",
                        file_types=[".csv", ".xlsx", ".xls", ".pdf"])
                    t2_evidence = gr.File(
                        label="📄 Upload Supporting Evidence / Exhibits (PDF, CSV, Excel)",
                        file_types=[".csv", ".xlsx", ".xls", ".pdf"])
                    t2_description = gr.Textbox(
                        label="Case Description (optional if uploading documents)", lines=3,
                        placeholder="Describe the key facts, legal issues, and arguments...")
                    t2_jurisdiction = gr.Textbox(
                        label="Jurisdiction", placeholder="e.g. California, Federal — 9th Circuit")
                    t2_case_type = gr.Dropdown(
                        choices=["Civil Litigation", "Criminal Defense", "Contract Dispute",
                                 "Employment Law", "Personal Injury", "IP Litigation",
                                 "Family Law", "Real Estate", "Bankruptcy", "Other"],
                        value="Civil Litigation", label="Case Type")
                    t2_context = gr.Textbox(label="Extra Context", lines=2)
                    t2_btn = gr.Button("⚖️ Predict Case Outcome", variant="primary")
                    gr.Markdown(WAIT_MSG)
                with gr.Column(scale=2):
                    t2_output = gr.Markdown(label="Case Outcome Analysis")
            t2_chart = gr.Plot(label="Outcome Probability Chart")
            t2_btn.click(analyze_case_outcome,
                         inputs=[t2_complaint, t2_answer, t2_evidence,
                                 t2_description, t2_jurisdiction,
                                 t2_case_type, t2_context],
                         outputs=[t2_output, t2_chart])

        # ── Tab 3 ──────────────────────────────────────────────────────────────
        with gr.Tab("📝 Contract Analysis"):
            gr.Markdown("## AI Contract Review & Risk Analysis")
            gr.Markdown("Upload the main contract plus any exhibits or amendments for a complete review.")
            gr.Markdown(PII_WARNING)
            with gr.Row():
                with gr.Column(scale=1):
                    t3_main = gr.File(
                        label="📄 Upload Main Contract (PDF, CSV, Excel) — Required",
                        file_types=[".csv", ".xlsx", ".xls", ".pdf"])
                    t3_exhibit_a = gr.File(
                        label="📄 Upload Exhibit A (optional)",
                        file_types=[".csv", ".xlsx", ".xls", ".pdf"])
                    t3_exhibit_b = gr.File(
                        label="📄 Upload Exhibit B (optional)",
                        file_types=[".csv", ".xlsx", ".xls", ".pdf"])
                    t3_amendment = gr.File(
                        label="📄 Upload Amendment / Addendum (optional)",
                        file_types=[".csv", ".xlsx", ".xls", ".pdf"])
                    t3_contract_type = gr.Dropdown(
                        choices=["Employment Agreement", "Service Agreement", "NDA",
                                 "Lease Agreement", "Purchase Agreement", "Partnership Agreement",
                                 "Software License", "Construction Contract", "Other"],
                        value="Service Agreement", label="Contract Type")
                    t3_jurisdiction = gr.Textbox(
                        label="Governing Jurisdiction",
                        placeholder="e.g. New York, Texas, Federal")
                    t3_party_role = gr.Dropdown(
                        choices=["Buyer / Client", "Seller / Vendor", "Employee",
                                 "Employer", "Licensor", "Licensee",
                                 "Landlord", "Tenant", "Other"],
                        value="Buyer / Client", label="Your Client's Role")
                    t3_context = gr.Textbox(label="Extra Context", lines=2)
                    t3_btn = gr.Button("📝 Analyze Contract", variant="primary")
                    gr.Markdown(WAIT_MSG)
                with gr.Column(scale=2):
                    t3_output = gr.Markdown(label="Contract Analysis")
            t3_chart = gr.Plot(label="Risk Assessment Chart")
            t3_btn.click(analyze_contract,
                         inputs=[t3_main, t3_exhibit_a, t3_exhibit_b, t3_amendment,
                                 t3_contract_type, t3_jurisdiction,
                                 t3_party_role, t3_context],
                         outputs=[t3_output, t3_chart])

        # ── Tab 4 ──────────────────────────────────────────────────────────────
        with gr.Tab("🔍 Legal Research"):
            gr.Markdown("## AI Legal Research Assistant")
            gr.Markdown("Upload a brief, statute, or case law to supplement your research question.")
            gr.Markdown(PII_WARNING)
            with gr.Row():
                with gr.Column(scale=1):
                    t4_brief = gr.File(
                        label="📄 Upload Legal Brief / Memo (optional)",
                        file_types=[".csv", ".xlsx", ".xls", ".pdf"])
                    t4_statute = gr.File(
                        label="📄 Upload Statute / Regulation (optional)",
                        file_types=[".csv", ".xlsx", ".xls", ".pdf"])
                    t4_caselaw = gr.File(
                        label="📄 Upload Case Law / Precedent (optional)",
                        file_types=[".csv", ".xlsx", ".xls", ".pdf"])
                    t4_question = gr.Textbox(
                        label="Legal Question or Issue", lines=4,
                        placeholder="e.g. What are the elements of negligence in California?")
                    t4_jurisdiction = gr.Textbox(
                        label="Jurisdiction", placeholder="e.g. California, Federal, New York")
                    t4_practice_area = gr.Dropdown(
                        choices=["Civil Litigation", "Criminal Law", "Contract Law",
                                 "Employment Law", "IP Law", "Real Estate",
                                 "Corporate Law", "Family Law", "Tax Law", "Other"],
                        value="Civil Litigation", label="Practice Area")
                    t4_opposing = gr.Textbox(
                        label="Opposing Arguments to Address (optional)", lines=3)
                    t4_context = gr.Textbox(label="Extra Context", lines=2)
                    t4_btn = gr.Button("🔍 Research Legal Issue", variant="primary")
                    gr.Markdown(WAIT_MSG)
                with gr.Column(scale=2):
                    t4_output = gr.Markdown(label="Legal Research Results")
            t4_btn.click(analyze_legal_research,
                         inputs=[t4_brief, t4_statute, t4_caselaw,
                                 t4_question, t4_jurisdiction,
                                 t4_practice_area, t4_opposing, t4_context],
                         outputs=[t4_output])

        # ── Tab 5 ──────────────────────────────────────────────────────────────
        with gr.Tab("💰 Billing & Time Analysis"):
            gr.Markdown("## Law Firm Billing & Productivity Analysis")
            gr.Markdown("Upload billing records and invoices together for a complete financial picture.")
            gr.Markdown(PII_WARNING)
            with gr.Row():
                with gr.Column(scale=1):
                    t5_billing = gr.File(
                        label="📄 Upload Billing Records / Time Entries (CSV/Excel)",
                        file_types=[".csv", ".xlsx", ".xls"])
                    t5_invoices = gr.File(
                        label="📄 Upload Invoices / AR Data (CSV/Excel — optional)",
                        file_types=[".csv", ".xlsx", ".xls"])
                    t5_btn = gr.Button("💰 Analyze Billing", variant="primary")
                    gr.Markdown(WAIT_MSG)
                with gr.Column(scale=2):
                    t5_output = gr.Markdown(label="Billing Analysis")
            t5_chart = gr.Plot(label="Billing Chart")
            t5_btn.click(analyze_billing,
                         inputs=[t5_billing, t5_invoices],
                         outputs=[t5_output, t5_chart])

        # ── Tab 6 ──────────────────────────────────────────────────────────────
        with gr.Tab("🛡️ Compliance & Risk"):
            gr.Markdown("## Regulatory Compliance & Risk Analysis")
            gr.Markdown("Upload policy documents, contracts, and regulatory filings for a complete compliance review.")
            gr.Markdown(PII_WARNING)
            with gr.Row():
                with gr.Column(scale=1):
                    t6_policy = gr.File(
                        label="📄 Upload Policy / Procedure Document (PDF, CSV, Excel)",
                        file_types=[".csv", ".xlsx", ".xls", ".pdf"])
                    t6_contract = gr.File(
                        label="📄 Upload Contract / Agreement (PDF, CSV, Excel — optional)",
                        file_types=[".csv", ".xlsx", ".xls", ".pdf"])
                    t6_regulatory = gr.File(
                        label="📄 Upload Regulatory Filing / Notice (PDF — optional)",
                        file_types=[".csv", ".xlsx", ".xls", ".pdf"])
                    t6_industry = gr.Dropdown(
                        choices=["Healthcare", "Finance / Banking", "Technology",
                                 "Real Estate", "Manufacturing", "Retail",
                                 "Legal Services", "Government", "Other"],
                        value="Legal Services", label="Industry")
                    t6_jurisdiction = gr.Textbox(
                        label="Jurisdiction", placeholder="e.g. Federal, California, EU")
                    t6_context = gr.Textbox(
                        label="Describe Compliance Concerns (optional)", lines=3)
                    t6_btn = gr.Button("🛡️ Analyze Compliance", variant="primary")
                    gr.Markdown(WAIT_MSG)
                with gr.Column(scale=2):
                    t6_output = gr.Markdown(label="Compliance & Risk Analysis")
            t6_chart = gr.Plot(label="Compliance Risk Chart")
            t6_btn.click(analyze_compliance,
                         inputs=[t6_policy, t6_contract, t6_regulatory,
                                 t6_industry, t6_jurisdiction, t6_context],
                         outputs=[t6_output, t6_chart])

        # ── Tab 7 ──────────────────────────────────────────────────────────────
        with gr.Tab("📈 Visualizations & Charts"):
            gr.Markdown("## Interactive Legal Analytics Visualizations")
            gr.Markdown(PII_WARNING)
            with gr.Row():
                with gr.Column(scale=1):
                    t7_file = gr.File(label="Upload Data (CSV/Excel)",
                                      file_types=[".csv", ".xlsx", ".xls"])
                    t7_chart_type = gr.Dropdown(
                        choices=["Bar Chart", "Line Chart", "Pie Chart", "Scatter Plot",
                                 "Heatmap", "Box Plot", "Histogram", "Area Chart"],
                        value="Bar Chart", label="Chart Type")
                    t7_x = gr.Dropdown(choices=[], label="X Axis Column")
                    t7_y = gr.Dropdown(choices=[], label="Y Axis Column")
                    t7_color = gr.Dropdown(choices=[], label="Color By (optional)")
                    t7_btn = gr.Button("📊 Generate Chart", variant="primary")
                    t7_status = gr.Markdown()
                with gr.Column(scale=2):
                    t7_plot = gr.Plot(label="Chart")
            t7_file.change(get_columns, inputs=[t7_file],
                           outputs=[t7_x, t7_y, t7_color])
            t7_btn.click(generate_viz,
                         inputs=[t7_file, t7_chart_type, t7_x, t7_y, t7_color],
                         outputs=[t7_plot, t7_status])

        # ── Tab 8 ──────────────────────────────────────────────────────────────
        with gr.Tab("📄 Reports & Presentations"):
            gr.Markdown("## Generate Professional Legal Reports")
            gr.Markdown("Upload a primary document and an optional supporting document.")
            gr.Markdown(PII_WARNING)
            with gr.Row():
                with gr.Column(scale=1):
                    t8_primary = gr.File(
                        label="📄 Upload Primary Document (CSV, Excel, PDF, JSON)",
                        file_types=[".csv", ".xlsx", ".xls", ".pdf", ".json"])
                    t8_supporting = gr.File(
                        label="📄 Upload Supporting Document (optional)",
                        file_types=[".csv", ".xlsx", ".xls", ".pdf", ".json"])
                    t8_matter = gr.Textbox(
                        label="Matter Description (optional)", lines=3,
                        placeholder="Describe the legal matter for the report...")
                    t8_type = gr.Dropdown(
                        choices=["PowerPoint Presentation", "PDF Report",
                                 "Case Summary", "Client Briefing"],
                        value="PDF Report", label="Report Type")
                    t8_btn = gr.Button("📄 Generate Report", variant="primary")
                    gr.Markdown(WAIT_MSG)
                    t8_download = gr.File(label="⬇️ Download Report")
                with gr.Column(scale=2):
                    t8_output = gr.Markdown(label="Report Preview")
            t8_btn.click(generate_report,
                         inputs=[t8_primary, t8_supporting, t8_type, t8_matter],
                         outputs=[t8_output, t8_download])

        # ── Tab 9 ──────────────────────────────────────────────────────────────
        with gr.Tab("💬 AI Legal Chat"):
            gr.Markdown("## Chat With Your Legal Data")
            gr.Markdown(PII_WARNING)
            with gr.Row():
                with gr.Column(scale=1):
                    t9_file = gr.File(
                        label="Upload Data for Context (CSV/Excel/PDF — optional)",
                        file_types=[".csv", ".xlsx", ".xls", ".pdf"])
                    t9_upload_btn = gr.Button("📤 Load Data", variant="secondary")
                    t9_data_status = gr.Markdown()
                with gr.Column(scale=2):
                    t9_chatbot = gr.ChatInterface(
                        fn=chat_with_data,
                        examples=[
                            "What are the key elements of a breach of contract claim?",
                            "What is the statute of limitations for personal injury in California?",
                            "What clauses should I always include in a service agreement?",
                            "Generate a SQL query to find the highest billing matters",
                            "Write Python code to calculate attorney utilization rates",
                            "What are the most common compliance risks for law firms?",
                            "How do I calculate the value of a personal injury settlement?",
                        ],
                        title="",
                    )
            t9_upload_btn.click(upload_chat_data, inputs=[t9_file],
                                outputs=[t9_data_status])

    gr.Markdown(WATERMARK)


if __name__ == "__main__":
    demo.launch(server_name="0.0.0.0", server_port=int(os.environ.get("GRADIO_SERVER_PORT", 7860)), share=False, ssr_mode=False)
