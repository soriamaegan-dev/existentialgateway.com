import gradio as gr
import requests
import os
import pandas as pd
import json
import io
import ast
import plotly.express as px
import plotly.graph_objects as go
from plotly.subplots import make_subplots
from pypdf import PdfReader
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from reportlab.lib.pagesizes import letter
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import inch
from reportlab.lib import colors
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
import tempfile
from datetime import datetime

FINANCIAL_SYSTEM_PROMPT = """You are a senior financial analyst and CFA-level investment researcher with expertise in equity analysis, fixed income, derivatives, portfolio management, financial statement analysis, and macroeconomic modeling.

ANALYTICAL STANDARDS:
1. Apply rigorous financial analysis frameworks: DCF valuation, comparable company analysis, precedent transaction analysis, LBO modeling where applicable.
2. Calculate and interpret key financial ratios: P/E, EV/EBITDA, P/B, ROE, ROA, ROIC, debt/equity, interest coverage, current ratio, quick ratio, cash conversion cycle.
3. Conduct DuPont analysis for profitability decomposition where financial statements are available.
4. Apply technical analysis: support/resistance levels, moving averages (50/200-day), RSI, MACD, volume analysis for market data.
5. Evaluate macroeconomic factors: Fed policy, yield curve dynamics, inflation expectations, sector rotation signals.
6. Identify accounting red flags: aggressive revenue recognition, off-balance-sheet liabilities, earnings quality issues, related party transactions.
7. Structure every analysis with: Executive Summary → Financial Health Score → Ratio Analysis → Trend Analysis → Valuation Assessment → Risk Factors → Investment Thesis.
8. Always provide bear, base, and bull case scenarios for forward-looking analysis.
9. Flag SEC disclosure requirements, material non-public information (MNPI) risks, and regulatory concerns.
10. Note that this is analytical support — not investment advice. Always recommend professional financial advisor consultation for investment decisions."""


HF_TOKEN = os.environ.get("HF_TOKEN", "")

DISCLAIMER = """
> **FINANCIAL DISCLAIMER**: This tool is AI-generated for informational purposes only.
> Nothing in this analysis constitutes financial, investment, or legal advice.
> Always consult a licensed financial advisor before making investment decisions.
> Past performance does not guarantee future results.
> © 2026 Existential Gateway, LLC. All Rights Reserved. Proprietary Software.
"""

WAIT_MSG = "*Results take approximately 1-2 minutes to generate. Please do not click multiple times.*"

WATERMARK = """
---
© 2026 Existential Gateway, LLC | AI Financial and Market Analyzer
Unauthorized reproduction strictly prohibited. Licensing: existentialgateway@gmail.com
---
"""

PII_WARNING = """
> **DATA PRIVACY NOTICE**: Before uploading any data you MUST remove all personally
> identifiable information including: Account Numbers | SSN | Client Names |
> Tax ID Numbers | Bank Account Details | Credit Card Numbers.
> This platform is NOT intended to store personally identifiable financial information.
"""



def fetch_url_content(text):
    """Detect URLs in user text, fetch their content, and return as context."""
    import re
    urls = re.findall(r'https?://[^\s<>"{}|\^`\[\]]+', text)
    if not urls:
        return ""
    fetched = []
    for url in urls[:3]:  # limit to 3 URLs
        try:
            headers = {
                "User-Agent": "Mozilla/5.0 (compatible; ExistentialGateway/1.0)",
                "Accept": "text/html,application/xhtml+xml,application/xml;q=0.9,*/*;q=0.8"
            }
            r = requests.get(url, timeout=15, headers=headers)
            if r.status_code == 200:
                from html.parser import HTMLParser
                class TextExtractor(HTMLParser):
                    def __init__(self):
                        super().__init__()
                        self.text = []
                        self.skip = False
                    def handle_starttag(self, tag, attrs):
                        if tag in ('script', 'style', 'nav', 'footer', 'header'):
                            self.skip = True
                    def handle_endtag(self, tag):
                        if tag in ('script', 'style', 'nav', 'footer', 'header'):
                            self.skip = False
                    def handle_data(self, data):
                        if not self.skip and data.strip():
                            self.text.append(data.strip())
                parser = TextExtractor()
                parser.feed(r.text)
                page_text = ' '.join(parser.text)[:3000]
                fetched.append(f"[URL: {url}]\n{page_text}")
        except Exception as e:
            fetched.append(f"[URL: {url}] Could not fetch: {str(e)}")
    if fetched:
        return "\n\nUSER-PROVIDED URL CONTENT (analyze this as part of your response):\n" + "\n---\n".join(fetched)
    return ""



def enrich_with_urls(text):
    """If text contains URLs, fetch and append their content."""
    if not text or 'http' not in text:
        return text
    url_content = fetch_url_content(text)
    if url_content:
        return text + url_content
    return text


def query_llm(prompt):
    API_KEY = os.environ.get("OPENAI_API_KEY", "")
    headers = {"Authorization": f"Bearer {API_KEY}", "Content-Type": "application/json"}
    msgs = [{"role": "system", "content": FINANCIAL_SYSTEM_PROMPT}, {"role": "user", "content": prompt}]
    payload = {"model": "gpt-4o", "max_tokens": 4000, "messages": msgs}
    try:
        response = requests.post("https://api.openai.com/v1/chat/completions", headers=headers, json=payload, timeout=240)
        result = response.json()
        if "choices" in result:
            return result["choices"][0]["message"]["content"]
        return f"API Error: {result}"
    except Exception as e:
        return f"Error: {str(e)}"


def load_data(file):
    if file is None:
        return None, "No file uploaded."
    try:
        path = file.name if hasattr(file, "name") else file
        if path.endswith(".csv"):
            df = pd.read_csv(path)
        elif path.endswith((".xlsx", ".xls")):
            df = pd.read_excel(path)
        elif path.endswith(".json"):
            df = pd.read_json(path)
        elif path.endswith(".pdf"):
            reader = PdfReader(path)
            text = "\n".join(page.extract_text() or "" for page in reader.pages)
            return None, text
        else:
            return None, "Unsupported file type."
        return df, None
    except Exception as e:
        return None, f"Error loading file: {str(e)}"


# ─── Tab 1: Data Upload and Overview ─────────────────────────────────────────

def analyze_overview(file):
    if file is None:
        return "Please upload a file.", "", ""
    df, err = load_data(file)
    if err and df is None:
        if "Unsupported" in err or "Error loading" in err:
            return err, "", ""
        text_preview = err[:3000]
        prompt = f"""You are an expert financial analyst. The user uploaded a PDF financial document. Write as the professional analyst who performed this work. State all findings directly with specific numbers, percentages, and comparisons. NEVER say 'this appears to be', 'the data seems to show', 'the dataset contains', or 'it looks like'. Instead, deliver results confidently: 'Average claim cost was $4,230, a 12% increase year-over-year' or 'Readmission rates of 18.3% in Cardiology exceeded the facility average by 4.1 percentage points.' Be direct, precise, and actionable. Present findings as YOUR analysis, not a description of a dataset. Content (truncated):
{text_preview}

Present your professional findings:
1. What type of financial document is this?
2. Key financial metrics and figures present
3. Time period covered
4. Summary of financial position or performance
5. Data quality observations
6. Recommended next analysis steps"""
        result = query_llm(prompt)
        return result + "\n\n" + WATERMARK, "PDF — text extracted", "See analysis above"

    rows, cols = df.shape
    dtypes = df.dtypes.to_string()
    missing = df.isnull().sum()
    missing_str = missing[missing > 0].to_string() if missing.any() else "None"
    preview = df.head(10).to_string()

    prompt = f"""You are an expert financial analyst. Write as the professional analyst who performed this work. State all findings directly with specific numbers, percentages, and comparisons. NEVER say 'this appears to be', 'the data seems to show', 'the dataset contains', or 'it looks like'. Instead, deliver results confidently: 'Average claim cost was $4,230, a 12% increase year-over-year' or 'Readmission rates of 18.3% in Cardiology exceeded the facility average by 4.1 percentage points.' Be direct, precise, and actionable. Present findings as YOUR analysis, not a description of a dataset. Analyze this financial dataset overview:

Rows: {rows} | Columns: {cols}
Column types:\n{dtypes}
Missing values:\n{missing_str}
First 10 rows:\n{preview}

Provide:
1. What type of financial data this appears to be
2. Time period and coverage
3. Key financial metrics present
4. Data quality assessment
5. Key columns and their significance
6. Recommended analyses to run
7. Any immediate patterns or red flags"""

    result = query_llm(prompt)
    stats = f"**Rows:** {rows} | **Columns:** {cols}\n\n**Missing Values:**\n```\n{missing_str}\n```"
    return result + "\n\n" + WATERMARK, stats, f"```\n{preview}\n```"


# ─── Tab 2: Market Trend Analysis ────────────────────────────────────────────

def analyze_market_trends(file):
    if file is None:
        return "Please upload market or stock data.", None
    df, err = load_data(file)
    if err:
        return err, None

    rows, cols = df.shape
    preview = df.head(20).to_string()
    desc = df.describe(include="all").to_string()

    prompt = f"""You are an expert financial market analyst. Analyze this market/stock data:

Shape: {rows} rows x {cols} columns
Columns: {list(df.columns)}
Statistical Summary:\n{desc}
Sample Data:\n{preview}

Provide a comprehensive market trend analysis:
1. **Overall Trend Direction** — upward, downward, or sideways with justification
2. **Moving Average Analysis** — short-term vs long-term trends
3. **Support and Resistance Levels** — key price levels if price data present
4. **Volume Analysis** — volume trends and what they signal
5. **Momentum Indicators** — strength of current trend
6. **Volatility Assessment** — current volatility level and implications
7. **Key Pattern Identification** — any recognizable chart patterns
8. **Market Summary** — professional written summary with outlook

Use specific numbers from the data. Be precise and professional."""

    result = query_llm(prompt)

    fig = None
    try:
        numeric_cols = df.select_dtypes(include="number").columns.tolist()
        date_cols = [c for c in df.columns if "date" in c.lower() or "time" in c.lower()]

        if numeric_cols:
            fig = make_subplots(rows=2, cols=1, shared_xaxes=True,
                                subplot_titles=["Price/Value Trend", "Distribution"],
                                row_heights=[0.7, 0.3])
            x_vals = df[date_cols[0]] if date_cols else list(range(len(df)))
            fig.add_trace(go.Scatter(x=x_vals, y=df[numeric_cols[0]],
                                     mode="lines", name=numeric_cols[0],
                                     line=dict(color="#00d4ff")), row=1, col=1)
            if len(numeric_cols) > 1:
                fig.add_trace(go.Bar(x=x_vals, y=df[numeric_cols[1]],
                                     name=numeric_cols[1],
                                     marker_color="#f39c12"), row=2, col=1)
            fig.update_layout(title="Market Trend Analysis", template="plotly_dark",
                              height=500)
    except Exception:
        fig = None

    return result + "\n\n" + WATERMARK, fig


# ─── Tab 3: Predictive Forecasting ───────────────────────────────────────────

def analyze_forecasting(file):
    if file is None:
        return "Please upload historical financial data.", None
    df, err = load_data(file)
    if err:
        return err, None

    rows, cols = df.shape
    preview = df.head(20).to_string()
    desc = df.describe(include="all").to_string()

    prompt = f"""You are an expert financial forecaster and quantitative analyst. Write as the professional analyst who performed this work. State all findings directly with specific numbers, percentages, and comparisons. NEVER say 'this appears to be', 'the data seems to show', 'the dataset contains', or 'it looks like'. Instead, deliver results confidently: 'Average claim cost was $4,230, a 12% increase year-over-year' or 'Readmission rates of 18.3% in Cardiology exceeded the facility average by 4.1 percentage points.' Be direct, precise, and actionable. Present findings as YOUR analysis, not a description of a dataset. Analyze this historical data:

Shape: {rows} rows x {cols} columns
Columns: {list(df.columns)}
Statistical Summary:\n{desc}
Sample Data:\n{preview}

Provide a comprehensive predictive forecast:
1. **Trend Extrapolation** — project current trends forward 3, 6, and 12 months
2. **Seasonal Adjustment** — identify and adjust for seasonal patterns
3. **Confidence Intervals** — provide optimistic, base, and pessimistic scenarios
4. **Bull Scenario** — best case projections with conditions required
5. **Bear Scenario** — worst case projections with risk factors
6. **Base Case Scenario** — most likely outcome with probability estimate
7. **Risk-Adjusted Forecast** — forecast adjusted for identified risks
8. **Key Assumptions** — list the assumptions underlying the forecast

Use specific numbers. Be professional and precise."""

    result = query_llm(prompt)

    fig = None
    try:
        numeric_cols = df.select_dtypes(include="number").columns.tolist()
        if numeric_cols:
            col = numeric_cols[0]
            y_vals = df[col].dropna().tolist()
            x_hist = list(range(len(y_vals)))

            # Simple linear trend extension
            if len(y_vals) > 1:
                import numpy as np
                z = np.polyfit(x_hist, y_vals, 1)
                p = np.poly1d(z)
                forecast_x = list(range(len(y_vals), len(y_vals) + 12))
                forecast_y = [p(x) for x in forecast_x]

                fig = go.Figure()
                fig.add_trace(go.Scatter(x=x_hist, y=y_vals, mode="lines+markers",
                                         name="Historical", line=dict(color="#00d4ff")))
                fig.add_trace(go.Scatter(x=forecast_x, y=forecast_y, mode="lines",
                                         name="Forecast (Linear Trend)",
                                         line=dict(color="#f39c12", dash="dash")))
                upper = [v * 1.1 for v in forecast_y]
                lower = [v * 0.9 for v in forecast_y]
                fig.add_trace(go.Scatter(x=forecast_x + forecast_x[::-1],
                                         y=upper + lower[::-1],
                                         fill="toself", fillcolor="rgba(243,156,18,0.15)",
                                         line=dict(color="rgba(255,255,255,0)"),
                                         name="Confidence Band"))
                fig.update_layout(title=f"{col} — Historical + 12-Period Forecast",
                                  template="plotly_dark", height=450)
    except Exception:
        fig = None

    return result + "\n\n" + WATERMARK, fig


# ─── Tab 4: Portfolio Risk Assessment ────────────────────────────────────────

def analyze_portfolio(file):
    if file is None:
        return "Please upload portfolio holdings data.", None
    df, err = load_data(file)
    if err:
        return err, None

    rows, cols = df.shape
    preview = df.head(20).to_string()
    desc = df.describe(include="all").to_string()

    prompt = f"""You are an expert portfolio manager and risk analyst. Analyze this portfolio data:

Shape: {rows} rows x {cols} columns
Columns: {list(df.columns)}
Statistical Summary:\n{desc}
Sample Data:\n{preview}

Provide a comprehensive portfolio risk assessment:
1. **Overall Portfolio Risk Score** — Low/Medium/High with justification
2. **Risk Score Per Asset** — individual asset risk levels
3. **Diversification Analysis** — is the portfolio adequately diversified?
4. **Correlation Analysis** — how correlated are the holdings?
5. **Volatility Assessment** — portfolio and individual asset volatility
6. **Value at Risk (VaR)** — estimated maximum loss at 95% confidence
7. **Concentration Risk** — any dangerous overweighting
8. **Rebalancing Recommendations** — specific actions to optimize the portfolio

Use standard portfolio theory. Be specific with percentages and numbers."""

    result = query_llm(prompt)

    fig = None
    try:
        numeric_cols = df.select_dtypes(include="number").columns.tolist()
        cat_cols = df.select_dtypes(include="object").columns.tolist()

        if cat_cols and numeric_cols:
            fig = make_subplots(rows=1, cols=2,
                                subplot_titles=["Portfolio Allocation", "Value Distribution"],
                                specs=[[{"type": "pie"}, {"type": "bar"}]])
            top = df.groupby(cat_cols[0])[numeric_cols[0]].sum()
            fig.add_trace(go.Pie(labels=top.index.tolist(), values=top.values.tolist(),
                                 hole=0.4), row=1, col=1)
            fig.add_trace(go.Bar(x=top.index.tolist(), y=top.values.tolist(),
                                 marker_color="#3498db"), row=1, col=2)
            fig.update_layout(title="Portfolio Composition", template="plotly_dark",
                              height=420)
        elif numeric_cols:
            corr = df[numeric_cols].corr()
            fig = px.imshow(corr, text_auto=True, template="plotly_dark",
                            title="Asset Correlation Matrix",
                            color_continuous_scale="RdBu_r")
            fig.update_layout(height=420)
    except Exception:
        fig = None

    return result + "\n\n" + WATERMARK, fig


# ─── Tab 5: Fraud Detection ───────────────────────────────────────────────────

def detect_fraud(file):
    if file is None:
        return "Please upload transaction data.", None
    df, err = load_data(file)
    if err:
        return err, None

    rows, cols = df.shape
    preview = df.head(30).to_string()
    desc = df.describe(include="all").to_string()

    signals = []
    for col in df.columns:
        cl = col.lower()
        if "amount" in cl or "value" in cl or "transaction" in cl:
            try:
                q99 = df[col].quantile(0.99)
                high = (df[col] > q99).sum()
                if high > 0:
                    signals.append(f"- {high} transactions above 99th percentile in '{col}' (>{q99:.2f})")
            except Exception:
                pass
    for col in df.columns:
        if df[col].duplicated().sum() > 0:
            dups = df[col].duplicated().sum()
            signals.append(f"- {dups} duplicate values detected in '{col}'")
            break

    signal_str = "\n".join(signals) if signals else "No automatic signals computed."

    prompt = f"""You are an expert financial fraud investigator. Analyze this transaction data:

Rows: {rows} | Columns: {list(df.columns)}
Statistical Summary:\n{desc}
Sample Data:\n{preview}
Auto-detected signals:\n{signal_str}

Provide a comprehensive fraud detection report:
1. **Overall Fraud Risk Score** — Low/Medium/High with justification
2. **Suspicious Transaction Patterns** — unusual amounts, frequencies, or sequences
3. **Anomalous Timing Patterns** — transactions at unusual times
4. **Geographic Anomalies** — if location data present
5. **Duplicate Transaction Detection** — exact or near-duplicate transactions
6. **Round Number Patterns** — suspiciously round transaction amounts
7. **Velocity Patterns** — unusually high transaction frequency
8. **Priority Investigation List** — top 5 patterns requiring immediate review
9. **Fraud Alert Report** — executive summary for compliance team

Be specific. Use numbers. Flag red flags clearly."""

    result = query_llm(prompt)

    fig = None
    try:
        numeric_cols = df.select_dtypes(include="number").columns.tolist()
        if numeric_cols:
            col = numeric_cols[0]
            q99 = df[col].quantile(0.99)
            flag_colors = ["red" if v > q99 else "#3498db" for v in df[col]]
            fig = go.Figure()
            fig.add_trace(go.Bar(x=list(range(len(df))), y=df[col],
                                 marker_color=flag_colors, name=col))
            fig.add_hline(y=q99, line_dash="dash", line_color="orange",
                          annotation_text="99th Percentile Flag Threshold")
            fig.update_layout(title=f"Transaction Amounts — Red = Flagged",
                              template="plotly_dark", height=400)
    except Exception:
        fig = None

    return result + "\n\n" + WATERMARK, fig


# ─── Tab 6: Financial Statement Analysis ─────────────────────────────────────

def analyze_statements(file):
    if file is None:
        return "Please upload financial statements (balance sheet, income statement, or cash flow).", None
    df, err = load_data(file)
    if err and df is None:
        if "Unsupported" in err or "Error" in err:
            return err, None
        prompt = f"""You are an expert CPA and financial analyst. Analyze these financial statements:

{err[:3000]}

Calculate and interpret:
1. Write as the professional analyst who performed this work. State all findings directly with specific numbers, percentages, and comparisons. NEVER say 'this appears to be', 'the data seems to show', 'the dataset contains', or 'it looks like'. Instead, deliver results confidently: 'Average claim cost was $4,230, a 12% increase year-over-year' or 'Readmission rates of 18.3% in Cardiology exceeded the facility average by 4.1 percentage points.' Be direct, precise, and actionable. Present findings as YOUR analysis, not a description of a dataset. **Liquidity Ratios** — Current ratio, Quick ratio, Cash ratio
2. **Profitability Ratios** — Gross margin, Net margin, ROA, ROE, EBITDA margin
3. **Leverage Ratios** — Debt-to-equity, Debt-to-assets, Interest coverage
4. **Efficiency Ratios** — Asset turnover, Inventory turnover, Receivables turnover
5. **Industry Benchmarks** — Compare ratios to typical industry standards
6. **Red Flags** — Any ratios indicating financial stress or risk
7. **Strengths** — Areas where the company performs well
8. **Recommendations** — Specific actions to improve financial health

Be precise. Use the actual numbers from the statements."""
        result = query_llm(prompt)
        return result + "\n\n" + WATERMARK, None

    rows, cols = df.shape
    preview = df.head(20).to_string()
    desc = df.describe(include="all").to_string()

    prompt = f"""You are an expert CPA and financial analyst. Analyze this financial statement data:

Shape: {rows} rows x {cols} columns
Columns: {list(df.columns)}
Statistical Summary:\n{desc}
Sample Data:\n{preview}

Calculate and interpret:
1. **Liquidity Ratios** — Current ratio, Quick ratio, Cash ratio
2. **Profitability Ratios** — Gross margin, Net margin, ROA, ROE, EBITDA margin
3. **Leverage Ratios** — Debt-to-equity, Debt-to-assets, Interest coverage
4. **Efficiency Ratios** — Asset turnover, Inventory turnover, Receivables turnover
5. **Industry Benchmarks** — Compare to typical industry standards
6. **Red Flags** — Ratios indicating financial stress or risk
7. **Strengths** — Areas of strong financial performance
8. **Recommendations** — Specific actions to improve financial health

Use actual numbers from the data. Be precise and professional."""

    result = query_llm(prompt)

    fig = None
    try:
        numeric_cols = df.select_dtypes(include="number").columns.tolist()
        cat_cols = df.select_dtypes(include="object").columns.tolist()
        if numeric_cols:
            fig = px.bar(df.head(15),
                         x=cat_cols[0] if cat_cols else df.index,
                         y=numeric_cols[:min(3, len(numeric_cols))],
                         barmode="group",
                         template="plotly_dark",
                         title="Financial Statement Overview")
            fig.update_layout(height=420)
    except Exception:
        fig = None

    return result + "\n\n" + WATERMARK, fig


# ─── Tab 7: Visualizations and Charts ────────────────────────────────────────

def generate_viz(file, chart_type, x_col, y_col, color_col):
    if file is None:
        return None, "Please upload a data file."
    df, err = load_data(file)
    if err:
        return None, err

    cols = df.columns.tolist()
    x = x_col if x_col in cols else cols[0]
    y = y_col if y_col in cols else (cols[1] if len(cols) > 1 else cols[0])
    color = color_col if color_col in cols else None

    try:
        if chart_type == "Bar Chart":
            fig = px.bar(df, x=x, y=y, color=color, template="plotly_dark",
                         title=f"{y} by {x}")
        elif chart_type == "Line Chart":
            fig = px.line(df, x=x, y=y, color=color, template="plotly_dark",
                          title=f"{y} over {x}")
        elif chart_type == "Candlestick":
            numeric_cols = df.select_dtypes(include="number").columns.tolist()
            if len(numeric_cols) >= 4:
                fig = go.Figure(data=[go.Candlestick(
                    x=df[x],
                    open=df[numeric_cols[0]],
                    high=df[numeric_cols[1]],
                    low=df[numeric_cols[2]],
                    close=df[numeric_cols[3]]
                )])
                fig.update_layout(title="Candlestick Chart", template="plotly_dark")
            else:
                fig = px.line(df, x=x, y=y, template="plotly_dark",
                              title="Line Chart (need 4 numeric cols for candlestick)")
        elif chart_type == "Pie Chart":
            fig = px.pie(df, names=x, values=y, template="plotly_dark",
                         title=f"{y} Distribution by {x}")
        elif chart_type == "Scatter Plot":
            fig = px.scatter(df, x=x, y=y, color=color, template="plotly_dark",
                             title=f"{x} vs {y} — Risk/Return")
        elif chart_type == "Heatmap":
            numeric = df.select_dtypes(include="number")
            corr = numeric.corr()
            fig = px.imshow(corr, text_auto=True, template="plotly_dark",
                            title="Correlation Heatmap", color_continuous_scale="RdBu_r")
        elif chart_type == "Area Chart":
            fig = px.area(df, x=x, y=y, color=color, template="plotly_dark",
                          title=f"{y} Area Chart")
        elif chart_type == "Histogram":
            fig = px.histogram(df, x=x, color=color, template="plotly_dark",
                               title=f"{x} Distribution", nbins=40)
        else:
            fig = px.bar(df, x=x, y=y, template="plotly_dark")

        fig.update_layout(height=500, margin=dict(l=40, r=40, t=60, b=40))
        return fig, f"Chart generated: {chart_type} — {y} by {x}"
    except Exception as e:
        return None, f"Chart error: {str(e)}"


def get_columns(file):
    if file is None:
        return gr.Dropdown(choices=[]), gr.Dropdown(choices=[]), gr.Dropdown(choices=[])
    df, err = load_data(file)
    if err or df is None:
        return gr.Dropdown(choices=[]), gr.Dropdown(choices=[]), gr.Dropdown(choices=[])
    cols = df.columns.tolist()
    return gr.Dropdown(choices=cols), gr.Dropdown(choices=cols), gr.Dropdown(choices=cols)


# ─── Tab 8: Reports and Presentations ────────────────────────────────────────

def generate_report(file, report_type):
    if file is None:
        return "Please upload a data file.", None
    df, err = load_data(file)
    if err and df is None:
        return err, None

    if df is not None:
        preview = df.head(20).to_string()
        desc = df.describe(include="all").to_string()
        cols = list(df.columns)
        rows = len(df)
        data_summary = f"Rows: {rows} | Columns: {cols}\nStatistics:\n{desc}\nSample:\n{preview}"
    else:
        data_summary = err[:3000]

    prompt = f"""You are an expert financial analyst writing a professional report. Write as the professional analyst who performed this work. State all findings directly with specific numbers, percentages, and comparisons. NEVER say 'this appears to be', 'the data seems to show', 'the dataset contains', or 'it looks like'. Instead, deliver results confidently: 'Average claim cost was $4,230, a 12% increase year-over-year' or 'Readmission rates of 18.3% in Cardiology exceeded the facility average by 4.1 percentage points.' Be direct, precise, and actionable. Present findings as YOUR analysis, not a description of a dataset. Analyze this financial data and produce a complete {report_type}:

{data_summary}

Structure:
1. **Executive Summary** (3-5 sentences)
2. **Data Overview** — what the financial data contains
3. **Key Financial Findings** — top 5 findings with specific numbers
4. **Market and Trend Analysis** — notable trends and patterns
5. **Risk Assessment** — overall risk level and primary drivers
6. **Forecast and Outlook** — likely future performance
7. **Recommendations** — 5 specific actionable recommendations
8. **Conclusion**

Write professionally for a financial audience. Use specific numbers."""

    report_text = query_llm(prompt)

    output_path = None
    try:
        if report_type == "PowerPoint Presentation":
            output_path = _make_pptx(report_text, df if df is not None else None)
        else:
            output_path = _make_pdf(report_text)
    except Exception as e:
        report_text += f"\n\n[Document generation error: {str(e)}]"

        custom_watermark = f"""
---
{company_name if company_name and company_name.strip() else "© 2026 Existential Gateway, LLC | QuantusData.ai"}
Powered by QuantusData AI | Unauthorized reproduction prohibited.
---
""" if company_name and company_name.strip() else WATERMARK
    return report_text + "\n\n" + custom_watermark, output_path


def _make_pptx(text, df):
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    def add_slide(title_text, body_text, bg_color=(5, 15, 30)):
        slide = prs.slides.add_slide(blank_layout)
        bg = slide.background
        bg.fill.solid()
        bg.fill.fore_color.rgb = RGBColor(*bg_color)
        tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(1.2))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = RGBColor(39, 174, 96)
        body_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.7), Inches(12.3), Inches(5.5))
        tf2 = body_box.text_frame
        tf2.word_wrap = True
        for line in body_text.split("\n")[:20]:
            line = line.strip()
            if not line:
                continue
            para = tf2.add_paragraph()
            para.text = line
            para.font.size = Pt(14)
            para.font.color.rgb = RGBColor(220, 240, 220)

    title_slide = prs.slides.add_slide(blank_layout)
    bg = title_slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(3, 10, 20)
    tb = title_slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11), Inches(1.5))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = "AI Financial & Market Analysis Report"
    p.font.size = Pt(34)
    p.font.bold = True
    p.font.color.rgb = RGBColor(39, 174, 96)
    tb2 = title_slide.shapes.add_textbox(Inches(1), Inches(4.2), Inches(11), Inches(1))
    tf2 = tb2.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = f"Generated: {datetime.now().strftime('%B %d, %Y')} | AI Financial Analyzer"
    p2.font.size = Pt(18)
    p2.font.color.rgb = RGBColor(160, 200, 160)

    sections = []
    current_title = "Overview"
    current_body = []
    for line in text.split("\n"):
        stripped = line.strip()
        if stripped.startswith("**") and stripped.endswith("**") and len(stripped) > 4:
            if current_body:
                sections.append((current_title, "\n".join(current_body)))
            current_title = stripped.strip("*").strip()
            current_body = []
        elif stripped.startswith("## "):
            if current_body:
                sections.append((current_title, "\n".join(current_body)))
            current_title = stripped[3:].strip()
            current_body = []
        else:
            current_body.append(line)
    if current_body:
        sections.append((current_title, "\n".join(current_body)))

    bg_colors = [(5, 15, 30), (5, 20, 15), (10, 15, 25), (5, 25, 10), (15, 10, 20)]
    for i, (title, body) in enumerate(sections[:8]):
        add_slide(title, body, bg_colors[i % len(bg_colors)])

    if df is not None:
        stats_body = f"Total Records: {len(df)}\nColumns: {len(df.columns)}\n\nColumn Overview:\n"
        for col in df.columns[:10]:
            stats_body += f"  • {col} ({df[col].dtype})\n"
        add_slide("Data Statistics", stats_body)

    closing = prs.slides.add_slide(blank_layout)
    bg = closing.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(3, 10, 20)
    tb = closing.shapes.add_textbox(Inches(1), Inches(2.8), Inches(11), Inches(2))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Report Complete"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(39, 174, 96)
    p2 = tf.add_paragraph()
    p2.text = "© 2026 Existential Gateway, LLC | AI Financial and Market Analyzer"
    p2.font.size = Pt(14)
    p2.font.color.rgb = RGBColor(160, 200, 160)

    tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".pptx")
    prs.save(tmp.name)
    return tmp.name


def _make_pdf(text):
    tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".pdf")
    doc = SimpleDocTemplate(tmp.name, pagesize=letter,
                            leftMargin=inch, rightMargin=inch,
                            topMargin=inch, bottomMargin=inch)
    styles = getSampleStyleSheet()
    title_style = ParagraphStyle("Title2", parent=styles["Title"],
                                 fontSize=20, spaceAfter=20,
                                 textColor=colors.HexColor("#1a5c2a"))
    heading_style = ParagraphStyle("H2", parent=styles["Heading2"],
                                   fontSize=14, spaceAfter=8,
                                   textColor=colors.HexColor("#27ae60"))
    body_style = ParagraphStyle("Body2", parent=styles["Normal"],
                                fontSize=11, spaceAfter=6, leading=16)
    story = []
    story.append(Paragraph("AI Financial & Market Analysis Report", title_style))
    story.append(Paragraph(
        f"Generated: {datetime.now().strftime('%B %d, %Y')} | AI Financial Analyzer",
        body_style))
    story.append(Spacer(1, 0.2 * inch))

    for line in text.split("\n"):
        stripped = line.strip()
        if not stripped:
            story.append(Spacer(1, 0.1 * inch))
            continue
        if stripped.startswith("**") and stripped.endswith("**"):
            story.append(Paragraph(stripped.strip("*"), heading_style))
        elif stripped.startswith("## "):
            story.append(Paragraph(stripped[3:], heading_style))
        elif stripped.startswith("# "):
            story.append(Paragraph(stripped[2:], title_style))
        else:
            safe = stripped.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")
            story.append(Paragraph(safe, body_style))

    story.append(Spacer(1, 0.3 * inch))
    story.append(Paragraph(
        "© 2026 Existential Gateway, LLC | AI Financial and Market Analyzer | existentialgateway@gmail.com",
        ParagraphStyle("footer", parent=styles["Normal"], fontSize=9,
                       textColor=colors.grey)))
    doc.build(story)
    return tmp.name


# ─── Tab 9: AI Data Chat ─────────────────────────────────────────────────────

_chat_df = None


def upload_chat_data(file):
    global _chat_df
    if file is None:
        _chat_df = None
        return "No file uploaded."
    df, err = load_data(file)
    if err and df is None:
        return err
    _chat_df = df
    return f"Data loaded: {len(df)} rows x {len(df.columns)} columns\nColumns: {list(df.columns)}"


def chat_with_data(message, history):
    global _chat_df

    data_context = ""
    if _chat_df is not None:
        desc = _chat_df.describe(include="all").to_string()
        preview = _chat_df.head(10).to_string()
        data_context = (
            f"The user has uploaded financial data.\n"
            f"Shape: {_chat_df.shape}\n"
            f"Columns: {list(_chat_df.columns)}\n"
            f"Statistical Summary:\n{desc}\n"
            f"First 10 rows:\n{preview}"
        )
    else:
        data_context = "No data uploaded yet. Answer general financial analysis questions."

    messages = [
        {
            "role": "system",
            "content": (
                "You are an expert financial analyst and AI assistant. Write as the professional analyst who performed this work. State all findings directly with specific numbers, percentages, and comparisons. NEVER say 'this appears to be', 'the data seems to show', 'the dataset contains', or 'it looks like'. Instead, deliver results confidently: 'Average claim cost was $4,230, a 12% increase year-over-year' or 'Readmission rates of 18.3% in Cardiology exceeded the facility average by 4.1 percentage points.' Be direct, precise, and actionable. Present findings as YOUR analysis, not a description of a dataset. "
                "When asked a question, deliver the answer immediately with specific numbers. Do NOT explain how to calculate. Analyze their financial data, identify trends, answer questions, "
                "provide specific computed answers with actual numbers from the data when asked. "
                "Be precise, professional, and helpful. Always remind users not to include "
                "personally identifiable or sensitive financial account information.\n\n"
                + data_context
            )
        }
    ]
    for user_msg, bot_msg in history:
        messages.append({"role": "user", "content": user_msg})
        messages.append({"role": "assistant", "content": bot_msg})
    messages.append({"role": "user", "content": message})

    import os
    API_KEY = os.environ.get("OPENAI_API_KEY", "")
    headers = {"Authorization": f"Bearer {API_KEY}", "Content-Type": "application/json"}
    msgs = [{"role": "system", "content": FINANCIAL_SYSTEM_PROMPT}, {"role": "user", "content": prompt}]
    payload = {"model": "gpt-4o", "max_tokens": 4000, "messages": msgs}
    try:
        response = requests.post("https://api.openai.com/v1/chat/completions", headers=headers, json=payload, timeout=240)
        result = response.json()
        if "choices" in result:
            return result["choices"][0]["message"]["content"]
        return f"API Error: {result}"
    except Exception as e:
        return f"Error: {str(e)}"


# ─── Gradio UI ────────────────────────────────────────────────────────────────

with gr.Blocks(title="AI Financial & Market Analyzer", theme=gr.themes.Base(
        primary_hue=gr.themes.colors.Color(
            c50="#fef9ec", c100="#faefd0", c200="#f4dea0", c300="#edc970",
            c400="#e4b04a", c500="#c9a84c", c600="#a8872e", c700="#856519",
            c800="#5e440d", c900="#3a2a05", c950="#1e1502"
        ),
        secondary_hue=gr.themes.colors.Color(
            c50="#e8edf5", c100="#c5d0e6", c200="#9eb0d4", c300="#7490c2",
            c400="#4f74b0", c500="#2d5a9e", c600="#1a3a6e", c700="#112240",
            c800="#0a1628", c900="#060e1a", c950="#03070d"
        ),
        neutral_hue=gr.themes.colors.Color(
            c50="#f0f2f7", c100="#d8dde9", c200="#b8c2d6", c300="#95a3be",
            c400="#7285a6", c500="#536890", c600="#3a4f73", c700="#263856",
            c800="#162438", c900="#0c1622", c950="#060b11"
        ),
        font=gr.themes.GoogleFont("DM Sans"),
        font_mono=gr.themes.GoogleFont("DM Mono"),
    ).set(
        body_background_fill="#0a1628",
        body_background_fill_dark="#0a1628",
        body_text_color="#f8f9fc",
        body_text_color_dark="#f8f9fc",
        block_background_fill="#112240",
        block_background_fill_dark="#112240",
        block_border_color="rgba(201,168,76,0.2)",
        block_border_color_dark="rgba(201,168,76,0.2)",
        block_label_background_fill="#112240",
        block_label_background_fill_dark="#112240",
        block_label_text_color="#c9a84c",
        block_label_text_color_dark="#c9a84c",
        block_title_text_color="#c9a84c",
        block_title_text_color_dark="#c9a84c",
        button_primary_background_fill="linear-gradient(135deg,#c9a84c,#e8c96a)",
        button_primary_background_fill_dark="linear-gradient(135deg,#c9a84c,#e8c96a)",
        button_primary_text_color="#0a1628",
        button_primary_text_color_dark="#0a1628",
        button_secondary_background_fill="#112240",
        button_secondary_background_fill_dark="#112240",
        button_secondary_border_color="rgba(201,168,76,0.3)",
        button_secondary_border_color_dark="rgba(201,168,76,0.3)",
        button_secondary_text_color="#c9a84c",
        button_secondary_text_color_dark="#c9a84c",
        input_background_fill="#0a1628",
        input_background_fill_dark="#0a1628",
        input_border_color="rgba(201,168,76,0.2)",
        input_border_color_dark="rgba(201,168,76,0.2)",
        input_placeholder_color="#8892a4",
        input_placeholder_color_dark="#8892a4",
        shadow_drop="0 4px 24px rgba(0,0,0,0.4)",
        shadow_drop_lg="0 8px 40px rgba(0,0,0,0.5)",
        table_even_background_fill="#0a1628",
        table_even_background_fill_dark="#0a1628",
        table_odd_background_fill="#112240",
        table_odd_background_fill_dark="#112240",
    )) as demo:
    gr.Markdown("# 💹 AI Financial & Market Analyzer")
    gr.Markdown(DISCLAIMER)

    with gr.Tabs():

        # ── Tab 1 ──────────────────────────────────────────────────────────────
        with gr.Tab("📊 Data Upload & Overview"):
            gr.Markdown("## Upload Your Financial Data")
            gr.Markdown(PII_WARNING)
            with gr.Row():
                with gr.Column(scale=1):
                    t1_file = gr.File(label="Upload File (CSV, Excel, PDF, JSON)",
                                      file_types=[".csv", ".xlsx", ".xls", ".pdf", ".json"])
                    t1_btn = gr.Button("🔍 Analyze Dataset", variant="primary")
                    gr.Markdown(WAIT_MSG)
                with gr.Column(scale=2):
                    t1_analysis = gr.Markdown(label="AI Analysis")
            with gr.Row():
                t1_stats = gr.Markdown(label="Statistics")
                t1_preview = gr.Markdown(label="Data Preview")
            t1_btn.click(analyze_overview, inputs=[t1_file],
                         outputs=[t1_analysis, t1_stats, t1_preview])

        # ── Tab 2 ──────────────────────────────────────────────────────────────
        with gr.Tab("📈 Market Trend Analysis"):
            gr.Markdown("## Market & Stock Trend Analysis")
            gr.Markdown(PII_WARNING)
            with gr.Row():
                with gr.Column(scale=1):
                    t2_file = gr.File(label="Upload Market/Stock Data (CSV/Excel)",
                                      file_types=[".csv", ".xlsx", ".xls"])
                    t2_btn = gr.Button("📈 Analyze Trends", variant="primary")
                    gr.Markdown(WAIT_MSG)
                with gr.Column(scale=2):
                    t2_output = gr.Markdown(label="Market Trend Analysis")
            t2_chart = gr.Plot(label="Trend Chart")
            t2_btn.click(analyze_market_trends, inputs=[t2_file],
                         outputs=[t2_output, t2_chart])

        # ── Tab 3 ──────────────────────────────────────────────────────────────
        with gr.Tab("🔮 Predictive Forecasting"):
            gr.Markdown("## AI-Powered Financial Forecasting")
            gr.Markdown(PII_WARNING)
            with gr.Row():
                with gr.Column(scale=1):
                    t3_file = gr.File(label="Upload Historical Financial Data (CSV/Excel)",
                                      file_types=[".csv", ".xlsx", ".xls"])
                    t3_btn = gr.Button("🔮 Generate Forecast", variant="primary")
                    gr.Markdown(WAIT_MSG)
                with gr.Column(scale=2):
                    t3_output = gr.Markdown(label="Predictive Forecast")
            t3_chart = gr.Plot(label="Forecast Chart")
            t3_btn.click(analyze_forecasting, inputs=[t3_file],
                         outputs=[t3_output, t3_chart])

        # ── Tab 4 ──────────────────────────────────────────────────────────────
        with gr.Tab("📂 Portfolio Risk Assessment"):
            gr.Markdown("## Portfolio Risk & Diversification Analysis")
            gr.Markdown(PII_WARNING)
            with gr.Row():
                with gr.Column(scale=1):
                    t4_file = gr.File(label="Upload Portfolio Holdings (CSV/Excel)",
                                      file_types=[".csv", ".xlsx", ".xls"])
                    t4_btn = gr.Button("📂 Assess Portfolio Risk", variant="primary")
                    gr.Markdown(WAIT_MSG)
                with gr.Column(scale=2):
                    t4_output = gr.Markdown(label="Portfolio Risk Assessment")
            t4_chart = gr.Plot(label="Portfolio Chart")
            t4_btn.click(analyze_portfolio, inputs=[t4_file],
                         outputs=[t4_output, t4_chart])

        # ── Tab 5 ──────────────────────────────────────────────────────────────
        with gr.Tab("🚨 Fraud Detection"):
            gr.Markdown("## AI Financial Fraud Detection")
            gr.Markdown(PII_WARNING)
            with gr.Row():
                with gr.Column(scale=1):
                    t5_file = gr.File(label="Upload Transaction Data (CSV/Excel)",
                                      file_types=[".csv", ".xlsx", ".xls"])
                    t5_btn = gr.Button("🚨 Run Fraud Detection", variant="primary")
                    gr.Markdown(WAIT_MSG)
                with gr.Column(scale=2):
                    t5_output = gr.Markdown(label="Fraud Detection Report")
            t5_chart = gr.Plot(label="Fraud Flag Chart")
            t5_btn.click(detect_fraud, inputs=[t5_file],
                         outputs=[t5_output, t5_chart])

        # ── Tab 6 ──────────────────────────────────────────────────────────────
        with gr.Tab("📋 Financial Statement Analysis"):
            gr.Markdown("## Financial Ratios & Statement Analysis")
            gr.Markdown(PII_WARNING)
            with gr.Row():
                with gr.Column(scale=1):
                    t6_file = gr.File(label="Upload Financial Statements (CSV, Excel, PDF)",
                                      file_types=[".csv", ".xlsx", ".xls", ".pdf"])
                    t6_btn = gr.Button("📋 Analyze Statements", variant="primary")
                    gr.Markdown(WAIT_MSG)
                with gr.Column(scale=2):
                    t6_output = gr.Markdown(label="Financial Statement Analysis")
            t6_chart = gr.Plot(label="Financial Chart")
            t6_btn.click(analyze_statements, inputs=[t6_file],
                         outputs=[t6_output, t6_chart])

        # ── Tab 7 ──────────────────────────────────────────────────────────────
        with gr.Tab("📉 Visualizations & Charts"):
            gr.Markdown("## Interactive Financial Visualizations")
            gr.Markdown(PII_WARNING)
            with gr.Row():
                with gr.Column(scale=1):
                    t7_file = gr.File(label="Upload Data (CSV/Excel)",
                                      file_types=[".csv", ".xlsx", ".xls"])
                    t7_chart_type = gr.Dropdown(
                        choices=["Bar Chart", "Line Chart", "Candlestick", "Pie Chart",
                                 "Scatter Plot", "Heatmap", "Area Chart", "Histogram"],
                        value="Line Chart", label="Chart Type")
                    t7_x = gr.Dropdown(choices=[], label="X Axis Column")
                    t7_y = gr.Dropdown(choices=[], label="Y Axis Column")
                    t7_color = gr.Dropdown(choices=[], label="Color By (optional)")
                    t7_btn = gr.Button("📊 Generate Chart", variant="primary")
                    t7_status = gr.Markdown()
                with gr.Column(scale=2):
                    t7_plot = gr.Plot(label="Chart")
            t7_file.change(get_columns, inputs=[t7_file],
                           outputs=[t7_x, t7_y, t7_color])
            t7_btn.click(generate_viz,
                         inputs=[t7_file, t7_chart_type, t7_x, t7_y, t7_color],
                         outputs=[t7_plot, t7_status])

        # ── Tab 8 ──────────────────────────────────────────────────────────────
        with gr.Tab("📄 Reports & Presentations"):
            gr.Markdown("## Generate Professional Financial Reports")
            gr.Markdown(PII_WARNING)
            with gr.Row():
                with gr.Column(scale=1):
                    t8_file = gr.File(label="Upload Data (CSV, Excel, PDF, JSON)",
                                      file_types=[".csv", ".xlsx", ".xls", ".pdf", ".json"])
                    t8_type = gr.Dropdown(
                        choices=["PowerPoint Presentation", "PDF Report",
                                 "Executive Summary", "Investment Memo"],
                        value="PDF Report", label="Report Type")
                    t8_btn = gr.Button("📄 Generate Report", variant="primary")
                    gr.Markdown(WAIT_MSG)
                    t8_download = gr.File(label="⬇️ Download Report")
                with gr.Column(scale=2):
                    t8_output = gr.Markdown(label="Report Preview")
            t8_btn.click(generate_report, inputs=[t8_file, t8_type],
                         outputs=[t8_output, t8_download])

        # ── Tab 9 ──────────────────────────────────────────────────────────────
        with gr.Tab("💬 AI Data Chat"):
            gr.Markdown("## Chat With Your Financial Data")
            gr.Markdown(PII_WARNING)
            with gr.Row():
                with gr.Column(scale=1):
                    t9_file = gr.File(label="Upload Data for Context (CSV/Excel)",
                                      file_types=[".csv", ".xlsx", ".xls"])
                    t9_upload_btn = gr.Button("📤 Load Data", variant="secondary")
                    t9_data_status = gr.Markdown()
                with gr.Column(scale=2):
                    t9_chatbot = gr.ChatInterface(
                        fn=chat_with_data,
                        examples=[
                            "What is the average transaction amount?",
                            "Which asset has the highest return in this portfolio?",
                            "Are there any suspicious patterns in this data?",
                            "Generate a SQL query to find the top 10 transactions by amount",
                            "Write Python code to calculate a 30-day moving average",
                            "What is the overall trend in this financial data?",
                            "Calculate the month with the highest revenue",
                        ],
                        title="",
                    )
            t9_upload_btn.click(upload_chat_data, inputs=[t9_file],
                                outputs=[t9_data_status])

    gr.Markdown(WATERMARK)


if __name__ == "__main__":
    demo.launch(server_name="0.0.0.0", server_port=int(os.environ.get("GRADIO_SERVER_PORT", 7860)), share=False, ssr_mode=False)
