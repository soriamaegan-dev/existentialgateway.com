import gradio as gr
import requests
import os
import pandas as pd
import json
import io
import ast
import plotly.express as px
import plotly.graph_objects as go
from plotly.subplots import make_subplots
import plotly.io as pio
from pypdf import PdfReader
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from reportlab.lib.pagesizes import letter
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import inch
from reportlab.lib import colors
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle
import tempfile
import base64
from datetime import datetime

INSURANCE_SYSTEM_PROMPT = """You are a senior insurance industry analyst with 20+ years of experience in P&C, life, health, and specialty insurance. Your expertise includes actuarial science, claims analysis, underwriting, fraud detection, loss ratio analysis, and regulatory compliance (NAIC, state DOI standards).

ANALYTICAL STANDARDS:
1. Always calculate and interpret key insurance KPIs: Loss Ratio, Combined Ratio, Expense Ratio, LAE Ratio, Claim Frequency, Claim Severity, Pure Premium, Earned Premium.
2. Apply industry benchmarks: Healthy combined ratio <100%, loss ratio benchmarks vary by line (auto: 65-75%, homeowners: 55-70%, commercial: 60-75%).
3. Identify fraud indicators using ISO ClaimSearch patterns, SIU red flags, and NICB guidelines.
4. Reference NAIC model regulations, state-specific requirements, and ISO rating rules where applicable.
5. For claims analysis, evaluate coverage applicability, policy conditions, exclusions, and subrogation potential.
6. Provide reserve adequacy assessments using chain-ladder, Bornhuetter-Ferguson, and development methods where data allows.
7. Structure every analysis with: Executive Summary → KPI Dashboard → Trend Analysis → Risk Assessment → Recommendations → Compliance Notes.
8. Be decisive — provide specific findings, not vague observations.
9. Flag anomalies that warrant SIU referral, actuarial review, or regulatory attention.
10. Always note data limitations and what additional information would improve the analysis."""


HF_TOKEN = os.environ.get("HF_TOKEN", "")

DISCLAIMER = """
> **LEGAL DISCLAIMER**: This tool is powered by Artificial Intelligence and is intended
> for informational purposes only. All results are AI-generated and should be reviewed
> by a qualified insurance professional before use in actual business decisions.
> State laws and insurance regulations change frequently.
> © 2026 Existential Gateway, LLC. All Rights Reserved. Proprietary Software.
"""

WAIT_MSG = "*Results take approximately 1-2 minutes to generate. Please do not click multiple times.*"

WATERMARK = """
---
© 2026 Existential Gateway, LLC | AI Insurance Data Analyzer
Unauthorized reproduction, distribution, or commercial use strictly prohibited.
For licensing inquiries: existentialgateway@gmail.com
---
"""

PII_WARNING = """
> **DATA PRIVACY NOTICE**: Before uploading any data, you MUST remove all personally
> identifiable information including: Policy Holder Names | SSN | Date of Birth |
> Address | Phone Number | Email | Driver License | Financial Account Numbers.
> This platform is not intended to store personal identifying information.
> By uploading data you confirm all personal information has been removed.
"""

conversation_history = []



def fetch_url_content(text):
    """Detect URLs in user text, fetch their content, and return as context."""
    import re
    urls = re.findall(r'https?://[^\s<>"{}|\^`\[\]]+', text)
    if not urls:
        return ""
    fetched = []
    for url in urls[:3]:  # limit to 3 URLs
        try:
            headers = {
                "User-Agent": "Mozilla/5.0 (compatible; ExistentialGateway/1.0)",
                "Accept": "text/html,application/xhtml+xml,application/xml;q=0.9,*/*;q=0.8"
            }
            r = requests.get(url, timeout=15, headers=headers)
            if r.status_code == 200:
                from html.parser import HTMLParser
                class TextExtractor(HTMLParser):
                    def __init__(self):
                        super().__init__()
                        self.text = []
                        self.skip = False
                    def handle_starttag(self, tag, attrs):
                        if tag in ('script', 'style', 'nav', 'footer', 'header'):
                            self.skip = True
                    def handle_endtag(self, tag):
                        if tag in ('script', 'style', 'nav', 'footer', 'header'):
                            self.skip = False
                    def handle_data(self, data):
                        if not self.skip and data.strip():
                            self.text.append(data.strip())
                parser = TextExtractor()
                parser.feed(r.text)
                page_text = ' '.join(parser.text)[:3000]
                fetched.append(f"[URL: {url}]\n{page_text}")
        except Exception as e:
            fetched.append(f"[URL: {url}] Could not fetch: {str(e)}")
    if fetched:
        return "\n\nUSER-PROVIDED URL CONTENT (analyze this as part of your response):\n" + "\n---\n".join(fetched)
    return ""



def enrich_with_urls(text):
    """If text contains URLs, fetch and append their content."""
    if not text or 'http' not in text:
        return text
    url_content = fetch_url_content(text)
    if url_content:
        return text + url_content
    return text


def query_llm(prompt):
    API_KEY = os.environ.get("OPENAI_API_KEY", "")
    headers = {"Authorization": f"Bearer {API_KEY}", "Content-Type": "application/json"}
    payload = {"model": "gpt-4o", "max_tokens": 4000, "messages": [{"role": "system", "content": INSURANCE_SYSTEM_PROMPT}, {"role": "user", "content": prompt}]}
    try:
        response = requests.post("https://api.openai.com/v1/chat/completions", headers=headers, json=payload, timeout=240)
        result = response.json()
        if "choices" in result:
            return result["choices"][0]["message"]["content"]
        return f"API Error: {result}"
    except Exception as e:
        return f"Error: {str(e)}"


def load_data(file):
    if file is None:
        return None, "No file uploaded."
    try:
        path = file.name if hasattr(file, "name") else file
        if path.endswith(".csv"):
            df = pd.read_csv(path)
        elif path.endswith((".xlsx", ".xls")):
            df = pd.read_excel(path)
        elif path.endswith(".json"):
            df = pd.read_json(path)
        elif path.endswith(".pdf"):
            reader = PdfReader(path)
            text = "\n".join(page.extract_text() or "" for page in reader.pages)
            return None, text
        else:
            return None, "Unsupported file type."
        return df, None
    except Exception as e:
        return None, f"Error loading file: {str(e)}"


def df_summary(df):
    buf = io.StringIO()
    df.info(buf=buf)
    info_str = buf.getvalue()
    missing = df.isnull().sum()
    missing_str = missing[missing > 0].to_string() if missing.any() else "None"
    preview = df.head(10).to_string()
    return info_str, missing_str, preview


# ─── Tab 1: Data Upload and Overview ────────────────────────────────────────

def analyze_overview(file):
    if file is None:
        return "Please upload a file.", "", ""
    df, err = load_data(file)
    if err and df is None:
        if "Unsupported" in err or "Error" in err:
            return err, "", ""
        text_preview = err[:3000]
        prompt = f"""You are an expert insurance data analyst. The user uploaded a PDF document. Write as the professional analyst who performed this work. State all findings directly with specific numbers, percentages, and comparisons. NEVER say 'this appears to be', 'the data seems to show', 'the dataset contains', or 'it looks like'. Instead, deliver results confidently: 'Average claim cost was $4,230, a 12% increase year-over-year' or 'Readmission rates of 18.3% in Cardiology exceeded the facility average by 4.1 percentage points.' Be direct, precise, and actionable. Present findings as YOUR analysis, not a description of a dataset. Here is its content (truncated):

{text_preview}

Present your professional findings:
1. What kind of insurance document is this?
2. Key data elements present
3. Summary of main findings
4. Data quality observations
5. Recommended next analysis steps"""
        result = query_llm(prompt)
        return result + "\n\n" + WATERMARK, "PDF — text extracted", "See above"

    rows, cols = df.shape
    info_str, missing_str, preview = df_summary(df)
    dtypes = df.dtypes.to_string()

    prompt = f"""You are an expert insurance data analyst. Write as the professional analyst who performed this work. State all findings directly with specific numbers, percentages, and comparisons. NEVER say 'this appears to be', 'the data seems to show', 'the dataset contains', or 'it looks like'. Instead, deliver results confidently: 'Average claim cost was $4,230, a 12% increase year-over-year' or 'Readmission rates of 18.3% in Cardiology exceeded the facility average by 4.1 percentage points.' Be direct, precise, and actionable. Present findings as YOUR analysis, not a description of a dataset. Analyze this dataset overview:

Rows: {rows} | Columns: {cols}
Column types:\n{dtypes}
Missing values:\n{missing_str}
First 10 rows:\n{preview}

Provide:
1. What kind of insurance data this appears to be
2. Data quality assessment (missing values, outliers, anomalies)
3. Key columns and their significance
4. Recommended analyses to run
5. Any immediate red flags or interesting patterns
Be specific and professional."""

    result = query_llm(prompt)
    stats = f"**Rows:** {rows} | **Columns:** {cols}\n\n**Missing Values:**\n```\n{missing_str}\n```"
    return result + "\n\n" + WATERMARK, stats, f"```\n{preview}\n```"


# ─── Tab 2: Claims Analysis ──────────────────────────────────────────────────

def analyze_claims(file):
    if file is None:
        return "Please upload a claims data file.", None
    df, err = load_data(file)
    if err:
        return err, None
    rows, cols = df.shape
    preview = df.head(20).to_string()
    desc = df.describe(include="all").to_string()

    prompt = f"""You are an expert insurance claims analyst. Analyze this claims dataset:

Shape: {rows} rows x {cols} columns
Columns: {list(df.columns)}
Statistical Summary:\n{desc}
Sample Data:\n{preview}

Provide a comprehensive claims analysis covering:
1. **Claim Frequency Analysis** — how many claims, frequency patterns
2. **Claim Severity** — average amounts, distribution, outliers
3. **Top Claim Types and Causes** — if identifiable from columns
4. **Trend Analysis** — monthly/yearly patterns if date columns exist
5. **Geographic Distribution** — if location data present
6. **Key Insights and Anomalies**
7. **Recommendations for Claims Management**

Be specific, use numbers from the data."""

    result = query_llm(prompt)

    fig = None
    try:
        numeric_cols = df.select_dtypes(include="number").columns.tolist()
        cat_cols = df.select_dtypes(include="object").columns.tolist()

        if numeric_cols:
            fig = make_subplots(rows=1, cols=min(2, len(numeric_cols)),
                                subplot_titles=numeric_cols[:2])
            for i, col in enumerate(numeric_cols[:2], 1):
                fig.add_trace(go.Histogram(x=df[col], name=col, nbinsx=30), row=1, col=i)
            fig.update_layout(title="Claims Numeric Distribution", height=400,
                              template="plotly_dark")
    except Exception:
        fig = None

    return result + "\n\n" + WATERMARK, fig


# ─── Tab 3: Loss Ratio Analysis ──────────────────────────────────────────────

def analyze_loss_ratio(file):
    if file is None:
        return "Please upload a premium/loss data file.", None
    df, err = load_data(file)
    if err:
        return err, None

    preview = df.head(20).to_string()
    desc = df.describe(include="all").to_string()
    cols = list(df.columns)

    # Try auto-calculate loss ratio if columns found
    loss_ratio_str = ""
    for c in df.columns:
        cl = c.lower()
        if "loss" in cl or "claim" in cl:
            for p in df.columns:
                pl = p.lower()
                if "premium" in pl or "earned" in pl:
                    try:
                        df["_loss_ratio"] = df[c] / df[p]
                        mean_lr = df["_loss_ratio"].mean()
                        loss_ratio_str = f"\nAuto-calculated Loss Ratio (mean): {mean_lr:.2%}"
                    except Exception:
                        pass

    prompt = f"""You are an expert insurance actuary. Write as the professional analyst who performed this work. State all findings directly with specific numbers, percentages, and comparisons. NEVER say 'this appears to be', 'the data seems to show', 'the dataset contains', or 'it looks like'. Instead, deliver results confidently: 'Average claim cost was $4,230, a 12% increase year-over-year' or 'Readmission rates of 18.3% in Cardiology exceeded the facility average by 4.1 percentage points.' Be direct, precise, and actionable. Present findings as YOUR analysis, not a description of a dataset. Analyze this premium and loss data:

Columns: {cols}
Statistical Summary:\n{desc}
Sample Data:\n{preview}
{loss_ratio_str}

Provide:
1. **Loss Ratio Calculation** — compute if data allows, interpret results
2. **Combined Ratio Analysis** — if expense data present
3. **Trend Over Time** — if date/period columns present
4. **Line of Business Comparison** — if LOB column present
5. **Industry Benchmarks** — typical loss ratios by line (auto ~65-75%, home ~50-60%, commercial varies)
6. **Concerns or Red Flags** — ratios above 100%, deteriorating trends
7. **Recommendations** — pricing, reserving, underwriting actions

Use industry standards. Be precise and professional."""

    result = query_llm(prompt)

    fig = None
    try:
        if "_loss_ratio" in df.columns:
            fig = go.Figure()
            fig.add_trace(go.Scatter(y=df["_loss_ratio"], mode="lines+markers",
                                     name="Loss Ratio", line=dict(color="#00d4ff")))
            fig.add_hline(y=1.0, line_dash="dash", line_color="red",
                          annotation_text="100% — Break Even")
            fig.update_layout(title="Loss Ratio Trend", yaxis_title="Loss Ratio",
                              template="plotly_dark", height=400)
        else:
            numeric_cols = df.select_dtypes(include="number").columns.tolist()
            if len(numeric_cols) >= 2:
                fig = go.Figure()
                for col in numeric_cols[:3]:
                    fig.add_trace(go.Bar(x=list(range(len(df))), y=df[col], name=col))
                fig.update_layout(title="Premium vs Loss Data", barmode="group",
                                  template="plotly_dark", height=400)
    except Exception:
        fig = None

    return result + "\n\n" + WATERMARK, fig


# ─── Tab 4: Fraud Detection ──────────────────────────────────────────────────

def detect_fraud(file):
    if file is None:
        return "Please upload a claims data file.", None
    df, err = load_data(file)
    if err:
        return err, None

    preview = df.head(30).to_string()
    desc = df.describe(include="all").to_string()
    cols = list(df.columns)
    rows = len(df)

    # Basic fraud signals
    signals = []
    for col in df.columns:
        cl = col.lower()
        if "amount" in cl or "claim" in cl:
            try:
                q99 = df[col].quantile(0.99)
                high = (df[col] > q99).sum()
                if high > 0:
                    signals.append(f"- {high} claims above 99th percentile in '{col}' (>{q99:.2f})")
            except Exception:
                pass
    for col in df.columns:
        if df[col].duplicated().sum() > 0:
            dups = df[col].duplicated().sum()
            signals.append(f"- {dups} duplicate values in '{col}'")
            break

    signal_str = "\n".join(signals) if signals else "No automatic signals computed."

    prompt = f"""You are an expert insurance fraud investigator. Write as the professional analyst who performed this work. State all findings directly with specific numbers, percentages, and comparisons. NEVER say 'this appears to be', 'the data seems to show', 'the dataset contains', or 'it looks like'. Instead, deliver results confidently: 'Average claim cost was $4,230, a 12% increase year-over-year' or 'Readmission rates of 18.3% in Cardiology exceeded the facility average by 4.1 percentage points.' Be direct, precise, and actionable. Present findings as YOUR analysis, not a description of a dataset. Analyze this claims dataset for fraud indicators:

Rows: {rows} | Columns: {cols}
Statistical Summary:\n{desc}
Sample Data:\n{preview}
Auto-detected signals:\n{signal_str}

Provide a comprehensive fraud detection report:
1. **Fraud Risk Score** — overall dataset risk (Low/Medium/High) with justification
2. **Duplicate Claims** — patterns, counts
3. **Unusual Claim Amounts** — statistical outliers
4. **High Frequency Claimants** — repeat filers
5. **Suspicious Timing Patterns** — claims on weekends, holidays, round numbers
6. **Geographic Anomalies** — clusters or unusual locations
7. **Provider Patterns** — if provider data present
8. **Priority Investigation List** — top 5 patterns to investigate first
9. **Fraud Alert Report** — executive summary

Be specific. Use numbers. Flag red flags clearly."""

    result = query_llm(prompt)

    fig = None
    try:
        numeric_cols = df.select_dtypes(include="number").columns.tolist()
        if numeric_cols:
            col = numeric_cols[0]
            q99 = df[col].quantile(0.99)
            colors_list = ["red" if v > q99 else "#00d4ff" for v in df[col]]
            fig = go.Figure()
            fig.add_trace(go.Bar(x=list(range(len(df))), y=df[col],
                                 marker_color=colors_list, name=col))
            fig.add_hline(y=q99, line_dash="dash", line_color="orange",
                          annotation_text="99th Percentile (Flag Threshold)")
            fig.update_layout(title=f"Fraud Flag: {col} Distribution (Red = Flagged)",
                              template="plotly_dark", height=400)
    except Exception:
        fig = None

    return result + "\n\n" + WATERMARK, fig


# ─── Tab 5: Underwriting Analysis ────────────────────────────────────────────

def analyze_underwriting(file):
    if file is None:
        return "Please upload a policy/risk data file.", None
    df, err = load_data(file)
    if err:
        return err, None

    preview = df.head(20).to_string()
    desc = df.describe(include="all").to_string()
    cols = list(df.columns)

    prompt = f"""You are an expert insurance underwriter. Write as the professional analyst who performed this work. State all findings directly with specific numbers, percentages, and comparisons. NEVER say 'this appears to be', 'the data seems to show', 'the dataset contains', or 'it looks like'. Instead, deliver results confidently: 'Average claim cost was $4,230, a 12% increase year-over-year' or 'Readmission rates of 18.3% in Cardiology exceeded the facility average by 4.1 percentage points.' Be direct, precise, and actionable. Present findings as YOUR analysis, not a description of a dataset. Analyze this policy and risk data:

Columns: {cols}
Statistical Summary:\n{desc}
Sample Data:\n{preview}

Provide a comprehensive underwriting analysis:
1. **Risk Scoring** — assess overall portfolio risk level
2. **Premium Adequacy Analysis** — are premiums sufficient for the risk?
3. **Exposure Analysis** — total insured values, limits, concentrations
4. **Risk Concentration** — geographic, industry, or segment concentrations
5. **High Risk Segments** — identify the riskiest policies or groups
6. **Underwriting Guidelines** — recommended changes based on data
7. **Portfolio Recommendations** — actions to improve profitability and reduce risk

Be specific and actionable."""

    result = query_llm(prompt)

    fig = None
    try:
        numeric_cols = df.select_dtypes(include="number").columns.tolist()
        cat_cols = df.select_dtypes(include="object").columns.tolist()
        if cat_cols and numeric_cols:
            top_cat = df[cat_cols[0]].value_counts().head(10)
            fig = go.Figure(go.Bar(x=top_cat.index.tolist(), y=top_cat.values,
                                   marker_color="#7c3aed"))
            fig.update_layout(title=f"Policy Distribution by {cat_cols[0]}",
                              template="plotly_dark", height=400)
        elif numeric_cols:
            fig = px.box(df, y=numeric_cols[:min(4, len(numeric_cols))],
                         title="Risk Variable Distributions", template="plotly_dark")
            fig.update_layout(height=400)
    except Exception:
        fig = None

    return result + "\n\n" + WATERMARK, fig


# ─── Tab 6: Visualizations ───────────────────────────────────────────────────

def generate_viz(file, chart_type, x_col, y_col, color_col):
    if file is None:
        return None, "Please upload a data file."
    df, err = load_data(file)
    if err:
        return None, err

    cols = df.columns.tolist()
    x = x_col if x_col in cols else cols[0]
    y = y_col if y_col in cols else (cols[1] if len(cols) > 1 else cols[0])
    color = color_col if color_col in cols else None

    try:
        if chart_type == "Bar Chart":
            fig = px.bar(df, x=x, y=y, color=color, template="plotly_dark",
                         title=f"{y} by {x}")
        elif chart_type == "Line Chart":
            fig = px.line(df, x=x, y=y, color=color, template="plotly_dark",
                          title=f"{y} over {x}")
        elif chart_type == "Pie Chart":
            fig = px.pie(df, names=x, values=y, template="plotly_dark",
                         title=f"{y} Distribution")
        elif chart_type == "Scatter Plot":
            fig = px.scatter(df, x=x, y=y, color=color, template="plotly_dark",
                             title=f"{x} vs {y}")
        elif chart_type == "Heatmap":
            numeric = df.select_dtypes(include="number")
            corr = numeric.corr()
            fig = px.imshow(corr, text_auto=True, template="plotly_dark",
                            title="Correlation Heatmap", color_continuous_scale="RdBu_r")
        elif chart_type == "Box Plot":
            fig = px.box(df, x=x if df[x].dtype == "object" else None,
                         y=y, color=color, template="plotly_dark",
                         title=f"{y} Distribution")
        elif chart_type == "Histogram":
            fig = px.histogram(df, x=x, color=color, template="plotly_dark",
                               title=f"{x} Frequency Distribution", nbins=40)
        else:
            fig = px.bar(df, x=x, y=y, template="plotly_dark")

        fig.update_layout(height=500, margin=dict(l=40, r=40, t=60, b=40))
        return fig, f"Chart generated: {chart_type} — {y} by {x}"
    except Exception as e:
        return None, f"Chart error: {str(e)}"


def get_columns(file):
    if file is None:
        return gr.Dropdown(choices=[]), gr.Dropdown(choices=[]), gr.Dropdown(choices=[])
    df, err = load_data(file)
    if err or df is None:
        return gr.Dropdown(choices=[]), gr.Dropdown(choices=[]), gr.Dropdown(choices=[])
    cols = df.columns.tolist()
    return gr.Dropdown(choices=cols), gr.Dropdown(choices=cols), gr.Dropdown(choices=cols)




# ─── Tab 9: Multi-Dataset Analysis ───────────────────────────────────────────
def analyze_multi_dataset(file1, file2, file3, analysis_type, instructions):
    if file1 is None or file2 is None:
        return "Please upload at least 2 datasets.", None

    datasets = []
    for i, f in enumerate([file1, file2, file3], 1):
        if f is None:
            continue
        df, err = load_data(f)
        if err or df is None:
            datasets.append((f"Dataset {i}", None, err or "Could not load file"))
        else:
            datasets.append((f"Dataset {i}", df, None))

    # Build individual summaries
    all_summaries = []
    for name, df, err in datasets:
        if df is None:
            all_summaries.append(f"### {name}\nError: {err}")
            continue
        rows, cols_list = df.shape
        dtypes = df.dtypes.to_string()
        missing = df.isnull().sum()
        missing_str = missing[missing > 0].to_string() if missing.sum() > 0 else "None"
        desc = df.describe(include="all").to_string()
        preview = df.head(15).to_string()
        summary = f"""### {name}: {rows} rows x {cols_list} columns
Columns: {list(df.columns)}
Types:\n{dtypes}
Missing values:\n{missing_str}
Statistics:\n{desc}
Sample data (first 15 rows):\n{preview}"""
        all_summaries.append(summary)

    combined = "\n\n".join(all_summaries)

    # ─── Compute real cross-dataset analytics ─────────────────────────────
    import pandas as pd
    import numpy as np
    computed_analytics = []
    valid_dfs = [(name, df) for name, df, err in datasets if df is not None]

    def find_col(df, candidates):
        for c in df.columns:
            if c.lower().strip().replace(" ", "_") in [x.lower().replace(" ", "_") for x in candidates]:
                return c
        return None

    if len(valid_dfs) >= 2:
        policy_col_names = ["Policy ID", "Policy_ID", "policy_id", "PolicyID", "POLICY_ID"]
        premium_col_names = ["Annual Premium", "Annual_Premium", "annual_premium", "Premium", "premium", "Total Premium"]
        claimed_col_names = ["Claimed Amount", "Claimed_Amount", "claimed_amount", "Claim Amount", "claim_amount", "Loss Amount"]
        paid_col_names = ["Paid Amount", "Paid_Amount", "paid_amount", "Payment Amount", "Settled Amount"]
        deductible_col_names = ["Deductible", "deductible", "Deductible Applied", "deductible_applied"]
        policy_type_names = ["Policy Type", "Policy_Type", "policy_type", "PolicyType"]
        risk_score_names = ["Risk Score", "Risk_Score", "risk_score"]
        fraud_flag_names = ["Fraud Flag", "Fraud_Flag", "fraud_flag"]
        fraud_score_names = ["Fraud Score", "Fraud_Score", "fraud_score"]
        claim_status_names = ["Claim Status", "Claim_Status", "claim_status"]
        loss_ratio_names = ["Historical Loss Ratio", "Loss Ratio", "loss_ratio", "Loss_Ratio_Historical"]
        combined_ratio_names = ["Combined Ratio", "combined_ratio", "Combined_Ratio"]
        coverage_limit_names = ["Coverage Limit", "Coverage_Limit", "coverage_limit"]
        credit_score_names = ["Credit Score", "Credit_Score", "credit_score"]
        state_names = ["State", "state"]
        decision_names = ["Decision", "decision", "UW Decision"]
        risk_tier_names = ["Risk Tier", "Risk_Tier", "risk_tier"]
        prior_claims_names = ["Prior Claims Count", "Prior_Claims_Count", "prior_claims_count"]
        industry_names = ["Industry", "industry"]
        closure_days_names = ["Days to Close", "Days_to_Close", "closure_days", "Closure Days"]
        injury_names = ["Claimant Injury", "Claimant_Injury", "claimant_injury"]

        merged = None
        merged_policy_col = None
        for name, df in valid_dfs:
            pc = find_col(df, policy_col_names)
            if pc:
                if merged is None:
                    merged = df.copy()
                    merged_policy_col = pc
                else:
                    merged = merged.merge(df, left_on=merged_policy_col, right_on=pc,
                                          how="outer", suffixes=("", f"_{name.replace(' ', '_')}"))

        try:
            for name, df in valid_dfs:
                metrics = []
                premium_c = find_col(df, premium_col_names)
                claimed_c = find_col(df, claimed_col_names)
                paid_c = find_col(df, paid_col_names)
                ptype_c = find_col(df, policy_type_names)
                risk_c = find_col(df, risk_score_names)
                fraud_c = find_col(df, fraud_flag_names)
                fraud_sc = find_col(df, fraud_score_names)
                status_c = find_col(df, claim_status_names)
                state_c = find_col(df, state_names)
                lr_c = find_col(df, loss_ratio_names)
                cr_c = find_col(df, combined_ratio_names)
                credit_c = find_col(df, credit_score_names)
                decision_c = find_col(df, decision_names)
                tier_c = find_col(df, risk_tier_names)
                closure_c = find_col(df, closure_days_names)
                injury_c = find_col(df, injury_names)
                industry_c = find_col(df, industry_names)
                prior_c = find_col(df, prior_claims_names)

                if premium_c:
                    vals = df[premium_c].dropna()
                    metrics.append(f"Total Premium Volume: ${vals.sum():,.2f}")
                    metrics.append(f"Avg Premium: ${vals.mean():,.2f} | Median: ${vals.median():,.2f}")
                    if ptype_c:
                        by_type = df.groupby(ptype_c)[premium_c].agg(["mean","sum","count"]).round(2)
                        by_type.columns = ["Avg_Premium","Total_Premium","Count"]
                        by_type = by_type.sort_values("Total_Premium", ascending=False)
                        metrics.append(f"Premium by Policy Type:\n{by_type.to_string()}")

                if claimed_c:
                    vals = df[claimed_c].dropna()
                    metrics.append(f"Total Claims Incurred: ${vals.sum():,.2f}")
                    metrics.append(f"Avg Claim: ${vals.mean():,.2f} | Median: ${vals.median():,.2f} | Max: ${vals.max():,.2f}")
                    bins = [0, 5000, 25000, 50000, 100000, 200000, float("inf")]
                    labels = ["<$5K","$5K-$25K","$25K-$50K","$50K-$100K","$100K-$200K",">$200K"]
                    buckets = pd.cut(vals, bins=bins, labels=labels)
                    metrics.append(f"Claims by Severity:\n{buckets.value_counts().sort_index().to_string()}")

                if paid_c:
                    vals = df[paid_c].dropna()
                    metrics.append(f"Total Paid: ${vals.sum():,.2f} | Avg: ${vals.mean():,.2f}")

                if claimed_c and paid_c:
                    ct = df[claimed_c].dropna().sum()
                    pt = df[paid_c].dropna().sum()
                    if ct > 0:
                        metrics.append(f"Settlement Ratio (Paid/Claimed): {pt/ct*100:.1f}%")

                if premium_c and claimed_c:
                    pr = df[premium_c].dropna().sum()
                    cl = df[claimed_c].dropna().sum()
                    if pr > 0:
                        metrics.append(f"Computed Loss Ratio: {cl/pr*100:.1f}%")

                if fraud_c:
                    fv = df[fraud_c].astype(str).str.lower().str.strip()
                    fc = fv.isin(["yes","y","true","1"]).sum()
                    metrics.append(f"Fraud-Flagged: {fc} ({fc/len(df)*100:.1f}%)")
                    if ptype_c and fraud_sc:
                        fbp = df.groupby(ptype_c).apply(
                            lambda x: pd.Series({
                                "Fraud_Count": x[fraud_c].astype(str).str.lower().isin(["yes","y","true","1"]).sum(),
                                "Avg_Fraud_Score": x[fraud_sc].mean(),
                                "Total": len(x)
                            }))
                        fbp["Fraud_Rate_%"] = (fbp["Fraud_Count"]/fbp["Total"]*100).round(1)
                        fbp = fbp.sort_values("Fraud_Rate_%", ascending=False)
                        metrics.append(f"Fraud by Policy Type:\n{fbp.to_string()}")

                if status_c:
                    metrics.append(f"Claim Status:\n{df[status_c].value_counts().to_string()}")

                if state_c and claimed_c:
                    sc = df.groupby(state_c)[claimed_c].agg(["sum","mean","count"]).round(2)
                    sc.columns = ["Total_Claims","Avg_Claim","Count"]
                    sc = sc.sort_values("Total_Claims", ascending=False).head(10)
                    metrics.append(f"Top 10 States by Claims:\n{sc.to_string()}")

                if lr_c:
                    v = df[lr_c].dropna()
                    metrics.append(f"Historical Loss Ratio: Mean={v.mean():.3f} | Median={v.median():.3f}")
                    hl = (v > 0.8).sum()
                    metrics.append(f"Loss Ratio > 80%: {hl} ({hl/len(v)*100:.1f}%)")

                if cr_c:
                    v = df[cr_c].dropna()
                    metrics.append(f"Combined Ratio: Mean={v.mean():.3f} | Median={v.median():.3f}")
                    up = (v > 1.0).sum()
                    metrics.append(f"Unprofitable (CR>100%): {up} ({up/len(v)*100:.1f}%)")

                if decision_c:
                    metrics.append(f"UW Decisions:\n{df[decision_c].value_counts().to_string()}")
                    if ptype_c:
                        ar = df.groupby(ptype_c)[decision_c].apply(
                            lambda x: (x.str.lower().str.contains("approv").sum()/len(x)*100)).round(1)
                        ar = ar.sort_values(ascending=False)
                        metrics.append(f"Approval Rate by Type (%):\n{ar.to_string()}")

                if tier_c:
                    metrics.append(f"Risk Tiers:\n{df[tier_c].value_counts().to_string()}")

                if credit_c:
                    v = df[credit_c].dropna()
                    metrics.append(f"Credit Score: Mean={v.mean():.0f} | Median={v.median():.0f} | Range={v.min():.0f}-{v.max():.0f}")

                if closure_c:
                    v = df[closure_c].dropna()
                    metrics.append(f"Days to Close: Mean={v.mean():.0f} | Median={v.median():.0f}")
                    if ptype_c:
                        cbt = df.groupby(ptype_c)[closure_c].agg(["mean","median","count"]).round(0)
                        metrics.append(f"Closure by Type:\n{cbt.sort_values('mean', ascending=False).to_string()}")

                if injury_c:
                    iv = df[injury_c].dropna()
                    iv = iv[iv.astype(str).str.strip() != "N/A"]
                    if len(iv) > 0:
                        metrics.append(f"Injury Severity:\n{iv.value_counts().to_string()}")

                if prior_c:
                    v = df[prior_c].dropna()
                    metrics.append(f"Prior Claims: Mean={v.mean():.1f} | Max={v.max():.0f}")

                if industry_c:
                    metrics.append(f"Top Industries:\n{df[industry_c].value_counts().head(10).to_string()}")

                if metrics:
                    computed_analytics.append(f"### COMPUTED METRICS for {name}\n" + "\n".join(metrics))

            if merged is not None and len(valid_dfs) >= 2:
                cross = []
                mp = find_col(merged, premium_col_names)
                mc = find_col(merged, claimed_col_names)
                mpd = find_col(merged, paid_col_names)
                mpt = find_col(merged, policy_type_names)
                mr = find_col(merged, risk_score_names)
                mf = find_col(merged, fraud_flag_names)
                mfs = find_col(merged, fraud_score_names)
                mlr = find_col(merged, loss_ratio_names)
                mcr = find_col(merged, combined_ratio_names)
                md = find_col(merged, decision_names)
                mt = find_col(merged, risk_tier_names)
                mcc = find_col(merged, credit_score_names)
                ms = find_col(merged, state_names)
                mi = find_col(merged, industry_names)

                portfolio = []
                if mp: portfolio.append(f"Total Portfolio Premium: ${merged[mp].sum():,.2f}")
                if mc: portfolio.append(f"Total Incurred Claims: ${merged[mc].sum():,.2f}")
                if mpd: portfolio.append(f"Total Paid Claims: ${merged[mpd].sum():,.2f}")
                if mp and mc:
                    tp = merged[mp].sum()
                    tc = merged[mc].sum()
                    if tp > 0:
                        portfolio.append(f"Portfolio Loss Ratio: {tc/tp*100:.1f}%")
                        portfolio.append(f"Net UW Position: ${tp-tc:,.2f}")
                if portfolio:
                    cross.append("PORTFOLIO SUMMARY:\n" + "\n".join(portfolio))

                if mp and mc and mpt:
                    lr = merged.groupby(mpt).apply(lambda x: pd.Series({
                        "Total_Premium": x[mp].sum(),
                        "Total_Claims": x[mc].sum(),
                        "Loss_Ratio": x[mc].sum()/x[mp].sum() if x[mp].sum()>0 else 0,
                        "Avg_Premium": x[mp].mean(),
                        "Avg_Claim": x[mc].mean(),
                        "Count": len(x)
                    })).round(2).sort_values("Loss_Ratio", ascending=False)
                    cross.append(f"Loss Ratio by Policy Type:\n{lr.to_string()}")

                if mf and mr:
                    fbr = merged.groupby(mr).apply(lambda x: pd.Series({
                        "Total": len(x),
                        "Fraud": x[mf].astype(str).str.lower().isin(["yes","y","true","1"]).sum()
                    }))
                    fbr["Rate_%"] = (fbr["Fraud"]/fbr["Total"]*100).round(1)
                    cross.append(f"Fraud Rate by Risk Score:\n{fbr.to_string()}")

                if mf and md:
                    fbd = merged.groupby(md).apply(lambda x: pd.Series({
                        "Total": len(x),
                        "Fraud": x[mf].astype(str).str.lower().isin(["yes","y","true","1"]).sum()
                    }))
                    fbd["Rate_%"] = (fbd["Fraud"]/fbd["Total"]*100).round(1)
                    cross.append(f"Fraud Rate by UW Decision:\n{fbd.to_string()}")

                if mp and mc and ms:
                    sp = merged.groupby(ms).apply(lambda x: pd.Series({
                        "Premium": x[mp].sum(),
                        "Claims": x[mc].sum(),
                        "Loss_Ratio": x[mc].sum()/x[mp].sum() if x[mp].sum()>0 else 0,
                        "Net": x[mp].sum()-x[mc].sum()
                    })).round(2).sort_values("Loss_Ratio", ascending=False).head(15)
                    cross.append(f"Top 15 States by Loss Ratio:\n{sp.to_string()}")

                if mcc and mc:
                    try:
                        corr = merged[[mcc, mc]].dropna().corr().iloc[0,1]
                        cross.append(f"Correlation (Credit Score vs Claims): {corr:.4f}")
                        bins = [0,600,650,700,750,800,900]
                        labels = ["<600","600-649","650-699","700-749","750-799","800+"]
                        merged["_cb"] = pd.cut(merged[mcc], bins=bins, labels=labels)
                        ccl = merged.groupby("_cb")[mc].agg(["mean","sum","count"]).round(2)
                        ccl.columns = ["Avg_Claim","Total_Claims","Count"]
                        cross.append(f"Claims by Credit Tier:\n{ccl.to_string()}")
                        merged.drop(columns=["_cb"], inplace=True, errors="ignore")
                    except Exception:
                        pass

                if mt and mlr:
                    tlr = merged.groupby(mt)[mlr].agg(["mean","median","count"]).round(4)
                    cross.append(f"Loss Ratio by Risk Tier:\n{tlr.sort_values('mean', ascending=False).to_string()}")

                if md and mc and mpd:
                    do = merged.groupby(md).apply(lambda x: pd.Series({
                        "Policies": len(x),
                        "Avg_Claim": x[mc].mean(),
                        "Avg_Paid": x[mpd].mean(),
                        "Total_Claims": x[mc].sum(),
                        "Total_Paid": x[mpd].sum()
                    })).round(2)
                    cross.append(f"Claims by UW Decision:\n{do.to_string()}")

                if mi and mc:
                    ir = merged.groupby(mi).apply(lambda x: pd.Series({
                        "Count": len(x),
                        "Total_Claims": x[mc].sum(),
                        "Avg_Claim": x[mc].mean(),
                        "Fraud_%": (x[mf].astype(str).str.lower().isin(["yes","y","true","1"]).sum()/len(x)*100) if mf else 0
                    })).round(2).sort_values("Total_Claims", ascending=False).head(10)
                    cross.append(f"Top Industries by Claims:\n{ir.to_string()}")

                if cross:
                    computed_analytics.append("### CROSS-DATASET COMPUTED METRICS\n" + "\n\n".join(cross))

        except Exception as e:
            computed_analytics.append(f"Note: Some computations failed: {str(e)}")

    computed_section = "\n\n".join(computed_analytics) if computed_analytics else "No cross-dataset metrics could be computed."

    prompt = f"""You are a senior insurance data analyst performing a {analysis_type} across {len(datasets)} related datasets. Write as the professional analyst who performed this work. State all findings directly with specific numbers, percentages, and comparisons. NEVER say 'this appears to be', 'the data seems to show', 'the dataset contains', or 'it looks like'. Instead, deliver results confidently: 'Average claim cost was $4,230, a 12% increase year-over-year' or 'Readmission rates of 18.3% in Cardiology exceeded the facility average by 4.1 percentage points.' Be direct, precise, and actionable. Present findings as YOUR analysis, not a description of a dataset. 
HERE IS THE ACTUAL DATA FROM EACH DATASET:

{combined}

{'=' * 80}
IMPORTANT: COMPUTED ANALYTICS (real numbers calculated by joining all datasets on shared keys):

{computed_section}

{'=' * 80}

CRITICAL: The COMPUTED ANALYTICS above contain REAL joined metrics. You MUST:
- Use these exact numbers in your analysis, do NOT re-estimate
- Reference specific loss ratios, fraud rates, and dollar amounts from the tables
- Compare segments using the cross-dataset tables (by policy type, state, risk tier)
- Identify most/least profitable segments from the loss ratio tables
- Flag fraud patterns from the fraud analysis breakdowns
- Assess underwriting quality from the UW decision vs claims data

{f'Additional instructions: {instructions}' if instructions else ''}

You MUST structure your response EXACTLY as follows:

## Dataset 1 Individual Summary
- What type of insurance data this is
- Row count, column count, key fields
- Key statistics (averages, medians, ranges for important numeric columns)
- Data quality assessment
- Notable patterns or issues specific to this dataset

## Dataset 2 Individual Summary
- What type of insurance data this is
- Row count, column count, key fields
- Key statistics (averages, medians, ranges for important numeric columns)
- Data quality assessment
- Notable patterns or issues specific to this dataset

{"## Dataset 3 Individual Summary" + chr(10) + "- Same structure as above" + chr(10) if len(datasets) > 2 else ""}

## Merged Cross-Dataset Analysis
- How these datasets relate to each other (shared keys, overlapping fields, complementary data)
- Patterns that ONLY emerge when comparing across datasets
- Correlations between variables across different datasets
- Discrepancies or inconsistencies between datasets
- Combined risk assessment

## Top 5 Key Insights
Numbered list of the most important actionable findings from the COMBINED analysis, with specific numbers and column references

## Summary & Conclusion
- Overall assessment of the combined insurance data landscape
- Risk areas identified across all datasets
- Strategic recommendations based on the merged analysis
- Suggested next steps

Use SPECIFIC numbers from the actual data. Reference actual column names and values. Do NOT make up data."""

    result = query_llm(prompt)

    # Try to create a comparison chart
    fig = None
    try:
        numeric_summaries = []
        for name, df, err in datasets:
            if df is not None:
                numeric_cols = df.select_dtypes(include="number").columns.tolist()
                if numeric_cols:
                    means = {col: df[col].mean() for col in numeric_cols[:8]}
                    numeric_summaries.append((name, means))
        if len(numeric_summaries) >= 2:
            import plotly.graph_objects as go
            fig = go.Figure()
            for name, means in numeric_summaries:
                fig.add_trace(go.Bar(name=name, x=list(means.keys()), y=list(means.values())))
            fig.update_layout(
                title="Cross-Dataset Numeric Comparison (Mean Values)",
                barmode="group", template="plotly_dark", height=450,
                margin=dict(l=40, r=40, t=60, b=40)
            )
    except Exception:
        fig = None

    return result + "\n\n" + WATERMARK, fig

# ─── Tab 8: Data Cleaning ────────────────────────────────────────────────────
def clean_data(file, cleaning_type):
    if file is None:
        return "Please upload a data file.", None
    df, err = load_data(file)
    if err or df is None:
        return err or "Could not load file.", None

    import copy
    original_rows = len(df)
    original_cols = list(df.columns)
    issues_found = []
    changes_made = []

    # --- Always: Report data quality ---
    missing_counts = df.isnull().sum()
    total_missing = int(missing_counts.sum())
    dup_count = int(df.duplicated().sum())

    if cleaning_type in ("full", "report"):
        for col in df.columns:
            mc = int(missing_counts[col])
            if mc > 0:
                issues_found.append(f"Column '{col}': {mc} missing values ({mc*100//len(df)}%)")
        if dup_count > 0:
            issues_found.append(f"{dup_count} duplicate rows found")
        blanks = df.map(lambda x: str(x).strip() in ("", "N/A", "n/a", "NA", "na", "None", "none", "-") if pd.notna(x) else False)
        blank_count = int(blanks.sum().sum())
        if blank_count > 0:
            issues_found.append(f"{blank_count} cells contain placeholder values (N/A, None, -, etc.)")

    if cleaning_type != "report":
        # --- Replace common placeholders with NaN ---
        placeholders = ["N/A", "n/a", "NA", "na", "None", "none", "-", "", "N\A", "null", "NULL"]
        replaced = 0
        for col in df.columns:
            mask = df[col].astype(str).str.strip().isin(placeholders)
            count = int(mask.sum())
            if count > 0:
                df.loc[mask, col] = pd.NA
                replaced += count
        if replaced > 0:
            changes_made.append(f"Replaced {replaced} placeholder values with blank (N/A, None, -, etc.)")

    if cleaning_type in ("full", "duplicates"):
        before = len(df)
        df = df.drop_duplicates()
        removed = before - len(df)
        if removed > 0:
            changes_made.append(f"Removed {removed} duplicate rows")

    if cleaning_type in ("full", "formats"):
        # Standardize text casing for common text columns
        for col in df.select_dtypes(include="object").columns:
            unique_vals = df[col].dropna().unique()
            if len(unique_vals) < 50:
                lower_map = {}
                for v in unique_vals:
                    key = str(v).strip().lower()
                    if key not in lower_map:
                        lower_map[key] = str(v).strip().title()
                before_vals = df[col].copy()
                df[col] = df[col].apply(lambda x: lower_map.get(str(x).strip().lower(), x) if pd.notna(x) else x)
                changed = int((before_vals != df[col]).sum())
                if changed > 0:
                    changes_made.append(f"Standardized casing in '{col}': {changed} values")

        # Standardize phone numbers
        for col in df.columns:
            if any(k in col.lower() for k in ["phone", "tel", "mobile", "cell"]):
                import re
                def clean_phone(x):
                    if pd.isna(x) or str(x).strip() == "":
                        return x
                    digits = re.sub(r"\D", "", str(x))
                    if digits.startswith("1") and len(digits) == 11:
                        digits = digits[1:]
                    if len(digits) == 10:
                        return f"({digits[:3]}) {digits[3:6]}-{digits[6:]}"
                    return x
                before_vals = df[col].copy()
                df[col] = df[col].apply(clean_phone)
                changed = int((before_vals.fillna("") != df[col].fillna("")).sum())
                if changed > 0:
                    changes_made.append(f"Standardized phone format in '{col}': {changed} values")

    if cleaning_type in ("full", "nulls"):
        # Fill numeric nulls with median
        for col in df.select_dtypes(include="number").columns:
            mc = int(df[col].isnull().sum())
            if mc > 0:
                median_val = df[col].median()
                df[col] = df[col].fillna(median_val)
                changes_made.append(f"Filled {mc} missing values in '{col}' with median ({median_val:.2f})")
        # Fill text nulls with "Unknown"
        for col in df.select_dtypes(include="object").columns:
            mc = int(df[col].isnull().sum())
            if mc > 0:
                df[col] = df[col].fillna("Unknown")
                changes_made.append(f"Filled {mc} missing values in '{col}' with 'Unknown'")

    if cleaning_type in ("full", "outliers"):
        # Flag outliers using IQR method
        for col in df.select_dtypes(include="number").columns:
            Q1 = df[col].quantile(0.25)
            Q3 = df[col].quantile(0.75)
            IQR = Q3 - Q1
            lower = Q1 - 1.5 * IQR
            upper = Q3 + 1.5 * IQR
            outliers = int(((df[col] < lower) | (df[col] > upper)).sum())
            if outliers > 0:
                flag_col = f"{col}_outlier_flag"
                df[flag_col] = ((df[col] < lower) | (df[col] > upper)).map({True: "Outlier", False: ""})
                changes_made.append(f"Flagged {outliers} outliers in '{col}' (IQR method, added '{flag_col}' column)")

    # Build report
    report = f"## Data Cleaning Report\n\n"
    report += f"**Original:** {original_rows} rows, {len(original_cols)} columns\n\n"
    report += f"**After cleaning:** {len(df)} rows, {len(df.columns)} columns\n\n"

    if issues_found:
        report += "### Issues Found\n"
        for issue in issues_found:
            report += f"- {issue}\n"
        report += "\n"

    if changes_made:
        report += "### Changes Made\n"
        for change in changes_made:
            report += f"- {change}\n"
        report += "\n"
    elif cleaning_type == "report":
        report += "### Recommendation\nRun a Full Clean to address the issues above.\n\n"

    report += f"**Data Health Score:** "
    completeness = 1 - (df.isnull().sum().sum() / (len(df) * len(df.columns)))
    score = round(completeness * 10, 1)
    report += f"{score}/10 (based on data completeness)\n"

    # Save cleaned file
    output_path = None
    if cleaning_type != "report":
        import tempfile, os
        output_path = os.path.join(tempfile.gettempdir(), "cleaned_data.xlsx")
        df.to_excel(output_path, index=False)

    report += "\n" + WATERMARK
    return report, output_path

# ─── Tab 7: Reports & PowerPoint ─────────────────────────────────────────────

def generate_report(file, report_type, output_format="PDF Report", company_name=""):
    if file is None:
        return "Please upload a data file.", None
    df, err = load_data(file)
    if err and df is None:
        return err, None

    if df is not None:
        preview = df.head(20).to_string()
        desc = df.describe(include="all").to_string()
        cols = list(df.columns)
        rows = len(df)
        data_summary = f"Rows: {rows} | Columns: {cols}\nStatistics:\n{desc}\nSample:\n{preview}"
    else:
        data_summary = err[:3000]

    prompt = f"""You are an expert insurance data analyst writing a professional report. Write as the professional analyst who performed this work. State all findings directly with specific numbers, percentages, and comparisons. NEVER say 'this appears to be', 'the data seems to show', 'the dataset contains', or 'it looks like'. Instead, deliver results confidently: 'Average claim cost was $4,230, a 12% increase year-over-year' or 'Readmission rates of 18.3% in Cardiology exceeded the facility average by 4.1 percentage points.' Be direct, precise, and actionable. Present findings as YOUR analysis, not a description of a dataset. Analyze this data and produce a complete {report_type}:

{data_summary}

Structure:
1. **Executive Summary** (3-5 sentences)
2. **Data Overview** — what the data contains
3. **Key Findings** — top 5 findings with numbers
4. **Risk Assessment** — overall risk level and drivers
5. **Trend Analysis** — notable trends
6. **Recommendations** — 5 actionable recommendations
7. **Conclusion**

Write professionally. Use specific numbers. Format with clear headers."""

    report_text = query_llm(prompt)

    output_path = None
    try:
        if output_format == "PowerPoint Presentation":
            output_path = _make_pptx(report_text, df)
        else:
            output_path = _make_pdf(report_text)
    except Exception as e:
        report_text += f"\n\n[Document generation error: {str(e)}]"

        custom_watermark = f"""
---
{company_name if company_name and company_name.strip() else "© 2026 Existential Gateway, LLC | QuantusData.ai"}
Powered by QuantusData AI | Unauthorized reproduction prohibited.
---
""" if company_name and company_name.strip() else WATERMARK
    return report_text + "\n\n" + custom_watermark, output_path


def _make_pptx(text, df):
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    blank_layout = prs.slide_layouts[6]

    def add_slide(title_text, body_text, bg_color=(15, 15, 30)):
        slide = prs.slides.add_slide(blank_layout)
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(*bg_color)

        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(1.2))
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0, 212, 255)

        body_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.7), Inches(12.3), Inches(5.5))
        tf2 = body_box.text_frame
        tf2.word_wrap = True
        for line in body_text.split("\n")[:20]:
            line = line.strip()
            if not line:
                continue
            para = tf2.add_paragraph()
            para.text = line
            para.font.size = Pt(14)
            para.font.color.rgb = RGBColor(220, 220, 240)

    # Title slide
    title_slide = prs.slides.add_slide(blank_layout)
    bg = title_slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(5, 5, 20)
    tb = title_slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11), Inches(1.5))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = "AI Insurance Data Analysis Report"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 212, 255)
    tb2 = title_slide.shapes.add_textbox(Inches(1), Inches(4), Inches(11), Inches(1))
    tf2 = tb2.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = f"Generated: {datetime.now().strftime('%B %d, %Y')} | AI Insurance Data Analyzer"
    p2.font.size = Pt(18)
    p2.font.color.rgb = RGBColor(150, 150, 200)

    # Parse sections from text
    sections = []
    current_title = "Overview"
    current_body = []
    for line in text.split("\n"):
        stripped = line.strip()
        if stripped.startswith("**") and stripped.endswith("**") and len(stripped) > 4:
            if current_body:
                sections.append((current_title, "\n".join(current_body)))
            current_title = stripped.strip("*").strip()
            current_body = []
        elif stripped.startswith("## "):
            if current_body:
                sections.append((current_title, "\n".join(current_body)))
            current_title = stripped[3:].strip()
            current_body = []
        else:
            current_body.append(line)
    if current_body:
        sections.append((current_title, "\n".join(current_body)))

    bg_colors = [(15, 15, 30), (10, 20, 40), (20, 10, 30), (5, 25, 35), (15, 25, 15)]
    for i, (title, body) in enumerate(sections[:8]):
        add_slide(title, body, bg_colors[i % len(bg_colors)])

    # Data stats slide
    if df is not None:
        stats_body = f"Total Records: {len(df)}\nColumns: {len(df.columns)}\n\n"
        stats_body += "Column Overview:\n"
        for col in df.columns[:10]:
            stats_body += f"  • {col} ({df[col].dtype})\n"
        add_slide("Data Statistics", stats_body)

    # Closing slide
    closing = prs.slides.add_slide(blank_layout)
    bg = closing.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(5, 5, 20)
    tb = closing.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11), Inches(2))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Report Complete"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 212, 255)
    p2 = tf.add_paragraph()
    p2.text = "© 2026 Existential Gateway, LLC | AI Insurance Data Analyzer"
    p2.font.size = Pt(14)
    p2.font.color.rgb = RGBColor(150, 150, 200)

    tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".pptx")
    prs.save(tmp.name)
    return tmp.name


def _make_pdf(text):
    tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".pdf")
    doc = SimpleDocTemplate(tmp.name, pagesize=letter,
                            leftMargin=inch, rightMargin=inch,
                            topMargin=inch, bottomMargin=inch)
    styles = getSampleStyleSheet()
    title_style = ParagraphStyle("Title2", parent=styles["Title"],
                                 fontSize=20, spaceAfter=20,
                                 textColor=colors.HexColor("#003366"))
    heading_style = ParagraphStyle("H2", parent=styles["Heading2"],
                                   fontSize=14, spaceAfter=8,
                                   textColor=colors.HexColor("#004488"))
    body_style = ParagraphStyle("Body2", parent=styles["Normal"],
                                fontSize=11, spaceAfter=6, leading=16)
    story = []
    story.append(Paragraph("AI Insurance Data Analysis Report", title_style))
    story.append(Paragraph(
        f"Generated: {datetime.now().strftime('%B %d, %Y')} | AI Insurance Data Analyzer",
        body_style))
    story.append(Spacer(1, 0.2 * inch))

    for line in text.split("\n"):
        stripped = line.strip()
        if not stripped:
            story.append(Spacer(1, 0.1 * inch))
            continue
        if stripped.startswith("**") and stripped.endswith("**"):
            story.append(Paragraph(stripped.strip("*"), heading_style))
        elif stripped.startswith("## "):
            story.append(Paragraph(stripped[3:], heading_style))
        elif stripped.startswith("# "):
            story.append(Paragraph(stripped[2:], title_style))
        else:
            safe = stripped.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")
            story.append(Paragraph(safe, body_style))

    story.append(Spacer(1, 0.3 * inch))
    story.append(Paragraph(
        "© 2026 Existential Gateway, LLC | AI Insurance Data Analyzer | existentialgateway@gmail.com",
        ParagraphStyle("footer", parent=styles["Normal"], fontSize=9,
                       textColor=colors.grey)))
    doc.build(story)
    return tmp.name


# ─── Tab 8: AI Data Chat ─────────────────────────────────────────────────────

_chat_df = None


def simple_chat(message, file_data, context):
    """Simple chat endpoint for custom frontend"""
    global _chat_df
    # If file provided and no data loaded yet, load it
    if file_data is not None and _chat_df is None:
        df, err = load_data(file_data)
        if df is not None:
            _chat_df = df
    # Call the main chat function with empty history
    return chat_with_data(message, [])

def upload_chat_data(file):
    global _chat_df
    if file is None:
        _chat_df = None
        return "No file uploaded."
    df, err = load_data(file)
    if err and df is None:
        return err
    _chat_df = df
    return f"Data loaded: {len(df)} rows × {len(df.columns)} columns\nColumns: {list(df.columns)}"


def chat_with_data(message, history):
    global _chat_df

    data_context = ""
    if _chat_df is not None:
        desc = _chat_df.describe(include="all").to_string()
        preview = _chat_df.head(10).to_string()
        data_context = f"""
The user has uploaded insurance data with the following properties:
Shape: {_chat_df.shape}
Columns: {list(_chat_df.columns)}
Statistical Summary:
{desc}
First 10 rows:
{preview}
"""
    else:
        data_context = "No data uploaded yet. Answer general insurance analytics questions."

    messages = [
        {
            "role": "system",
            "content": (
                "You are an expert insurance data analyst. The user has uploaded real data and you have access to its statistics and sample rows below. Write as the professional analyst who performed this work. State all findings directly with specific numbers, percentages, and comparisons. NEVER say 'this appears to be', 'the data seems to show', 'the dataset contains', or 'it looks like'. Instead, deliver results confidently: 'Average claim cost was $4,230, a 12% increase year-over-year' or 'Readmission rates of 18.3% in Cardiology exceeded the facility average by 4.1 percentage points.' Be direct, precise, and actionable. Present findings as YOUR analysis, not a description of a dataset. "
                "IMPORTANT RULES:\n"
                "1. ALWAYS answer with ACTUAL NUMBERS from the data. Never just give code — give the real answer FIRST.\n"
                "2. When asked about averages, totals, counts, etc., CALCULATE the answer from the Statistical Summary provided and state it clearly.\n"
                "3. Format currency as $X,XXX.XX and percentages as X.XX%.\n"
                "4. After giving the answer, you may optionally show a short code snippet the user could run for further analysis.\n"
                "5. If you cannot determine an exact answer from the data provided, explain what you can determine and what additional data would be needed.\n"
                "6. Be specific — reference actual column names, row counts, and values from the data.\n"
                "7. Present findings in a clear, professional format that a non-technical user can understand.\n\n" + data_context
            )
        }
    ]
    for user_msg, bot_msg in history:
        messages.append({"role": "user", "content": user_msg})
        messages.append({"role": "assistant", "content": bot_msg})
    messages.append({"role": "user", "content": message})

    import os
    API_KEY = os.environ.get("OPENAI_API_KEY", "")
    headers = {"Authorization": f"Bearer {API_KEY}", "Content-Type": "application/json"}
    msgs = messages if 'messages' in dir() else [{"role": "user", "content": prompt}]
    payload = {"model": "gpt-4o", "max_tokens": 4000, "messages": msgs}
    try:
        response = requests.post("https://api.openai.com/v1/chat/completions", headers=headers, json=payload, timeout=240)
        result = response.json()
        if "choices" in result:
            return result["choices"][0]["message"]["content"]
        return f"API Error: {result}"
    except Exception as e:
        return f"Error: {str(e)}"


# ─── Gradio UI ────────────────────────────────────────────────────────────────

with gr.Blocks(title="AI Insurance Data Analyzer", theme=gr.themes.Base(
        primary_hue=gr.themes.colors.Color(
            c50="#fef9ec", c100="#faefd0", c200="#f4dea0", c300="#edc970",
            c400="#e4b04a", c500="#c9a84c", c600="#a8872e", c700="#856519",
            c800="#5e440d", c900="#3a2a05", c950="#1e1502"
        ),
        secondary_hue=gr.themes.colors.Color(
            c50="#e8edf5", c100="#c5d0e6", c200="#9eb0d4", c300="#7490c2",
            c400="#4f74b0", c500="#2d5a9e", c600="#1a3a6e", c700="#112240",
            c800="#0a1628", c900="#060e1a", c950="#03070d"
        ),
        neutral_hue=gr.themes.colors.Color(
            c50="#f0f2f7", c100="#d8dde9", c200="#b8c2d6", c300="#95a3be",
            c400="#7285a6", c500="#536890", c600="#3a4f73", c700="#263856",
            c800="#162438", c900="#0c1622", c950="#060b11"
        ),
        font=gr.themes.GoogleFont("DM Sans"),
        font_mono=gr.themes.GoogleFont("DM Mono"),
    ).set(
        body_background_fill="#0a1628",
        body_background_fill_dark="#0a1628",
        body_text_color="#f8f9fc",
        body_text_color_dark="#f8f9fc",
        block_background_fill="#112240",
        block_background_fill_dark="#112240",
        block_border_color="rgba(201,168,76,0.2)",
        block_border_color_dark="rgba(201,168,76,0.2)",
        block_label_background_fill="#112240",
        block_label_background_fill_dark="#112240",
        block_label_text_color="#c9a84c",
        block_label_text_color_dark="#c9a84c",
        block_title_text_color="#c9a84c",
        block_title_text_color_dark="#c9a84c",
        button_primary_background_fill="linear-gradient(135deg,#c9a84c,#e8c96a)",
        button_primary_background_fill_dark="linear-gradient(135deg,#c9a84c,#e8c96a)",
        button_primary_text_color="#0a1628",
        button_primary_text_color_dark="#0a1628",
        button_secondary_background_fill="#112240",
        button_secondary_background_fill_dark="#112240",
        button_secondary_border_color="rgba(201,168,76,0.3)",
        button_secondary_border_color_dark="rgba(201,168,76,0.3)",
        button_secondary_text_color="#c9a84c",
        button_secondary_text_color_dark="#c9a84c",
        input_background_fill="#0a1628",
        input_background_fill_dark="#0a1628",
        input_border_color="rgba(201,168,76,0.2)",
        input_border_color_dark="rgba(201,168,76,0.2)",
        input_placeholder_color="#8892a4",
        input_placeholder_color_dark="#8892a4",
        shadow_drop="0 4px 24px rgba(0,0,0,0.4)",
        shadow_drop_lg="0 8px 40px rgba(0,0,0,0.5)",
        table_even_background_fill="#0a1628",
        table_even_background_fill_dark="#0a1628",
        table_odd_background_fill="#112240",
        table_odd_background_fill_dark="#112240",
    )) as demo:
    gr.Markdown("# 🛡️ AI Insurance Data Analyzer")
    gr.Markdown(DISCLAIMER)

    with gr.Tabs():

        # ── Tab 1 ──────────────────────────────────────────────────────────────
        with gr.Tab("📊 Data Upload & Overview"):
            gr.Markdown("## Upload Your Insurance Data")
            gr.Markdown(PII_WARNING)
            with gr.Row():
                with gr.Column(scale=1):
                    t1_file = gr.File(label="Upload File (CSV, Excel, PDF, JSON)",
                                      file_types=[".csv", ".xlsx", ".xls", ".pdf", ".json"])
                    t1_btn = gr.Button("🔍 Analyze Dataset", variant="primary")
                    gr.Markdown(WAIT_MSG)
                with gr.Column(scale=2):
                    t1_analysis = gr.Markdown(label="AI Analysis")
            with gr.Row():
                t1_stats = gr.Markdown(label="Statistics")
                t1_preview = gr.Markdown(label="Data Preview")
            t1_btn.click(analyze_overview, inputs=[t1_file],
                         outputs=[t1_analysis, t1_stats, t1_preview])

        # ── Tab 2 ──────────────────────────────────────────────────────────────
        with gr.Tab("📋 Claims Analysis"):
            gr.Markdown("## Claims Data Analysis")
            gr.Markdown(PII_WARNING)
            with gr.Row():
                with gr.Column(scale=1):
                    t2_file = gr.File(label="Upload Claims Data (CSV/Excel)",
                                      file_types=[".csv", ".xlsx", ".xls"])
                    t2_btn = gr.Button("🔍 Analyze Claims", variant="primary")
                    gr.Markdown(WAIT_MSG)
                with gr.Column(scale=2):
                    t2_output = gr.Markdown(label="Claims Analysis")
            t2_chart = gr.Plot(label="Claims Chart")
            t2_btn.click(analyze_claims, inputs=[t2_file], outputs=[t2_output, t2_chart])

        # ── Tab 3 ──────────────────────────────────────────────────────────────
        with gr.Tab("📉 Loss Ratio Analysis"):
            gr.Markdown("## Loss Ratio & Combined Ratio Analysis")
            gr.Markdown(PII_WARNING)
            with gr.Row():
                with gr.Column(scale=1):
                    t3_file = gr.File(label="Upload Premium & Loss Data (CSV/Excel)",
                                      file_types=[".csv", ".xlsx", ".xls"])
                    t3_btn = gr.Button("🔍 Analyze Loss Ratios", variant="primary")
                    gr.Markdown(WAIT_MSG)
                with gr.Column(scale=2):
                    t3_output = gr.Markdown(label="Loss Ratio Analysis")
            t3_chart = gr.Plot(label="Loss Ratio Chart")
            t3_btn.click(analyze_loss_ratio, inputs=[t3_file], outputs=[t3_output, t3_chart])

        # ── Tab 4 ──────────────────────────────────────────────────────────────
        with gr.Tab("🚨 Fraud Detection"):
            gr.Markdown("## AI Fraud Detection & Risk Flagging")
            gr.Markdown(PII_WARNING)
            with gr.Row():
                with gr.Column(scale=1):
                    t4_file = gr.File(label="Upload Claims Data (CSV/Excel)",
                                      file_types=[".csv", ".xlsx", ".xls"])
                    t4_btn = gr.Button("🚨 Run Fraud Detection", variant="primary")
                    gr.Markdown(WAIT_MSG)
                with gr.Column(scale=2):
                    t4_output = gr.Markdown(label="Fraud Detection Report")
            t4_chart = gr.Plot(label="Fraud Flag Chart")
            t4_btn.click(detect_fraud, inputs=[t4_file], outputs=[t4_output, t4_chart])

        # ── Tab 5 ──────────────────────────────────────────────────────────────
        with gr.Tab("🏢 Underwriting Analysis"):
            gr.Markdown("## Underwriting & Risk Portfolio Analysis")
            gr.Markdown(PII_WARNING)
            with gr.Row():
                with gr.Column(scale=1):
                    t5_file = gr.File(label="Upload Policy/Risk Data (CSV/Excel)",
                                      file_types=[".csv", ".xlsx", ".xls"])
                    t5_btn = gr.Button("🔍 Analyze Underwriting", variant="primary")
                    gr.Markdown(WAIT_MSG)
                with gr.Column(scale=2):
                    t5_output = gr.Markdown(label="Underwriting Analysis")
            t5_chart = gr.Plot(label="Underwriting Chart")
            t5_btn.click(analyze_underwriting, inputs=[t5_file], outputs=[t5_output, t5_chart])

        # ── Tab 6 ──────────────────────────────────────────────────────────────
        with gr.Tab("📈 Visualizations"):
            gr.Markdown("## Interactive Data Visualizations")
            gr.Markdown(PII_WARNING)
            with gr.Row():
                with gr.Column(scale=1):
                    t6_file = gr.File(label="Upload Data (CSV/Excel)",
                                      file_types=[".csv", ".xlsx", ".xls"])
                    t6_chart_type = gr.Dropdown(
                        choices=["Bar Chart", "Line Chart", "Pie Chart", "Scatter Plot",
                                 "Heatmap", "Box Plot", "Histogram"],
                        value="Bar Chart", label="Chart Type")
                    t6_x = gr.Dropdown(choices=[], label="X Axis Column", allow_custom_value=True)
                    t6_y = gr.Dropdown(choices=[], label="Y Axis Column", allow_custom_value=True)
                    t6_color = gr.Dropdown(choices=[], label="Color By (optional)", allow_custom_value=True)
                    t6_btn = gr.Button("📊 Generate Chart", variant="primary")
                    t6_status = gr.Markdown()
                with gr.Column(scale=2):
                    t6_plot = gr.Plot(label="Chart")
            t6_file.change(get_columns, inputs=[t6_file],
                           outputs=[t6_x, t6_y, t6_color])
            t6_btn.click(generate_viz,
                         inputs=[t6_file, t6_chart_type, t6_x, t6_y, t6_color],
                         outputs=[t6_plot, t6_status])

        # ── Tab 7 ──────────────────────────────────────────────────────────────
        with gr.Tab("📄 Reports & PowerPoint"):
            gr.Markdown("## Generate Professional Reports")
            gr.Markdown(PII_WARNING)
            with gr.Row():
                with gr.Column(scale=1):
                    t7_file = gr.File(label="Upload Data (CSV/Excel/PDF/JSON)",
                                      file_types=[".csv", ".xlsx", ".xls", ".pdf", ".json"])
                    t7_type = gr.Dropdown(
                        choices=["Executive Summary", "Claims Report", "Loss Ratio Report",
                                 "Fraud Analysis Report", "Underwriting Report", "Full Analysis Report"],
                        value="PDF Report", label="Report Type")
                    t7_format = gr.Dropdown(
                        choices=["PDF Report", "PowerPoint Presentation"],
                        value="PDF Report", label="Download Format")
                    t7_btn = gr.Button("📄 Generate Report", variant="primary")
                    gr.Markdown(WAIT_MSG)
                    t7_download = gr.File(label="⬇️ Download Report")
                with gr.Column(scale=2):
                    t7_output = gr.Markdown(label="Report Preview")
            t7_btn.click(generate_report, inputs=[t7_file, t7_type, t7_format],
                         outputs=[t7_output, t7_download])

        # ── Tab 8: Multi-Dataset ───────────────────────────────────────────────
        with gr.Tab("🔗 Multi-Dataset Analysis"):
            gr.Markdown("## Cross-Dataset Intelligence")
            gr.Markdown(PII_WARNING)
            with gr.Row():
                with gr.Column(scale=1):
                    t8m_file1 = gr.File(label="Dataset 1 (Primary)", file_types=[".csv", ".xlsx", ".xls"])
                    t8m_file2 = gr.File(label="Dataset 2", file_types=[".csv", ".xlsx", ".xls"])
                    t8m_file3 = gr.File(label="Dataset 3 (Optional)", file_types=[".csv", ".xlsx", ".xls"])
                    t8m_type = gr.Dropdown(
                        choices=["Cross-Dataset Comparison", "Correlation Analysis",
                                 "Year-over-Year Trends", "Merge & Unified Analysis",
                                 "Anomaly Detection Across Datasets", "Full Multi-Dataset Report"],
                        value="Full Multi-Dataset Report", label="Analysis Type")
                    t8m_instructions = gr.Textbox(label="Analysis Instructions (Optional)",
                        placeholder="e.g. Compare claims vs policies, find correlations...")
                    t8m_btn = gr.Button("🔗 Analyze All Datasets", variant="primary")
                    gr.Markdown(WAIT_MSG)
                with gr.Column(scale=2):
                    t8m_output = gr.Markdown(label="Multi-Dataset Analysis")
                    t8m_plot = gr.Plot(label="Cross-Dataset Chart")
            t8m_btn.click(analyze_multi_dataset,
                          inputs=[t8m_file1, t8m_file2, t8m_file3, t8m_type, t8m_instructions],
                          outputs=[t8m_output, t8m_plot])
        # ── Tab 9: Data Cleaning ───────────────────────────────────────────────
        with gr.Tab("🧹 Data Cleaning"):
            gr.Markdown("## Clean & Standardize Your Data")
            gr.Markdown(PII_WARNING)
            with gr.Row():
                with gr.Column(scale=1):
                    t8c_file = gr.File(label="Upload Data (CSV/Excel)",
                                       file_types=[".csv", ".xlsx", ".xls"])
                    t8c_type = gr.Dropdown(
                        choices=["full", "nulls", "duplicates", "formats", "outliers", "report"],
                        value="full", label="Cleaning Type",
                        info="Full=all fixes, or pick a specific cleaning action")
                    t8c_btn = gr.Button("🧹 Clean Dataset", variant="primary")
                    gr.Markdown(WAIT_MSG)
                    t8c_download = gr.File(label="⬇️ Download Cleaned Data")
                with gr.Column(scale=2):
                    t8c_output = gr.Markdown(label="Cleaning Report")
            t8c_btn.click(clean_data, inputs=[t8c_file, t8c_type],
                          outputs=[t8c_output, t8c_download])
        # ── Tab 9 ──────────────────────────────────────────────────────────────
        with gr.Tab("💬 AI Data Chat"):
            gr.Markdown("## Chat With Your Data")
            gr.Markdown(PII_WARNING)
            with gr.Row():
                with gr.Column(scale=1):
                    t8_file = gr.File(label="Upload Data for Context (CSV/Excel)",
                                      file_types=[".csv", ".xlsx", ".xls"])
                    t8_upload_btn = gr.Button("📤 Load Data", variant="secondary")
                    t8_data_status = gr.Markdown()
                with gr.Column(scale=2):
                    t8_chatbot = gr.ChatInterface(
                        fn=chat_with_data,
                        examples=[
                            "What is the average claim amount?",
                            "Which region has the highest loss ratio?",
                            "Are there any suspicious patterns in this data?",
                            "Generate a SQL query to find the top 10 claims by amount",
                            "Write Python code to analyze this data with pandas",
                            "What are the main risk factors in this dataset?",
                        ],
                        title="",
                    )
            t8_upload_btn.click(upload_chat_data, inputs=[t8_file], outputs=[t8_data_status])
            # Hidden endpoint for custom frontend chat
            t8_simple = gr.Button(visible=False, elem_id="simple-chat-btn")
            t8_simple_msg = gr.Textbox(visible=False)
            t8_simple_file = gr.File(visible=False)
            t8_simple_ctx = gr.Textbox(visible=False)
            t8_simple_out = gr.Textbox(visible=False)
            t8_simple.click(simple_chat,
                            inputs=[t8_simple_msg, t8_simple_file, t8_simple_ctx],
                            outputs=[t8_simple_out],
                            api_name="simple_chat")

    gr.Markdown(WATERMARK)


if __name__ == "__main__":
    demo.launch(server_name="0.0.0.0", server_port=int(os.environ.get("GRADIO_SERVER_PORT", 7860)), share=False, ssr_mode=False)
