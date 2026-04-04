"""
Insurance Economics Analyzer - Gradio Backend
Powered by QuantusData.ai
Port: 7870
Gradio: 6.10.0
"""

import os
import json
import requests
import pandas as pd
import numpy as np
import gradio as gr
from datetime import datetime, timedelta
import plotly.graph_objects as go


def fetch_url_content(text):
    """Detect URLs in user text, fetch their content, and return as context."""
    import re
    urls = re.findall(r'https?://[^\s<>"{}|\^`\[\]]+', text)
    if not urls:
        return ""
    fetched = []
    for url in urls[:3]:  # limit to 3 URLs
        try:
            headers = {
                "User-Agent": "Mozilla/5.0 (compatible; ExistentialGateway/1.0)",
                "Accept": "text/html,application/xhtml+xml,application/xml;q=0.9,*/*;q=0.8"
            }
            r = requests.get(url, timeout=15, headers=headers)
            if r.status_code == 200:
                from html.parser import HTMLParser
                class TextExtractor(HTMLParser):
                    def __init__(self):
                        super().__init__()
                        self.text = []
                        self.skip = False
                    def handle_starttag(self, tag, attrs):
                        if tag in ('script', 'style', 'nav', 'footer', 'header'):
                            self.skip = True
                    def handle_endtag(self, tag):
                        if tag in ('script', 'style', 'nav', 'footer', 'header'):
                            self.skip = False
                    def handle_data(self, data):
                        if not self.skip and data.strip():
                            self.text.append(data.strip())
                parser = TextExtractor()
                parser.feed(r.text)
                page_text = ' '.join(parser.text)[:3000]
                fetched.append(f"[URL: {url}]\n{page_text}")
        except Exception as e:
            fetched.append(f"[URL: {url}] Could not fetch: {str(e)}")
    if fetched:
        return "\n\nUSER-PROVIDED URL CONTENT (analyze this as part of your response):\n" + "\n---\n".join(fetched)
    return ""



def enrich_with_urls(text):
    """If text contains URLs, fetch and append their content."""
    if not text or 'http' not in text:
        return text
    url_content = fetch_url_content(text)
    if url_content:
        return text + url_content
    return text


def is_true(series):
    """Check boolean-like column for truthy values (handles Yes/No, 1/0, True/False, Y/N)."""
    return series.astype(str).str.lower().str.strip().isin(['1', 'true', 'yes', 'y', 'bundled'])

def is_false(series):
    """Check boolean-like column for falsy values."""
    return series.astype(str).str.lower().str.strip().isin(['0', 'false', 'no', 'n'])

def fig_to_dict(fig):
    """Convert a Plotly figure to a JSON-serializable dict for the custom frontend."""
    if fig is None:
        return None
    try:
        return json.loads(fig.to_json())
    except Exception:
        return None
import plotly.express as px
from io import BytesIO
import base64

# ============================================================================
# CONFIGURATION & CONSTANTS
# ============================================================================

WATERMARK = "\n\n---\n*Analysis powered by QuantusData.ai Insurance Economics Analyzer*"

ANALYST_TONE = """You are a senior insurance economics analyst with expertise in household insurance economics, actuarial analysis, market penetration, customer lifecycle modeling, retention analytics, risk assessment, loss ratio analysis, pricing optimization, and geographic market analysis.

CRITICAL ACCURACY RULES:
1. State findings directly with specific numbers, percentages, and dollar amounts from the data only.
2. NEVER use phrases like: 'As the professional analyst', 'This dataset appears to portray', 'the data seems to show', 'I have identified', 'Under my analysis'.
3. Lead every section with hard numbers. No preamble, no self-referencing.
4. ONLY cite statistics explicitly present in the provided data. Never fabricate metrics.
5. Apply insurance economics benchmarks: healthy loss ratio <65%, combined ratio <100%, retention rate >85%, cross-sell ratio >1.8 products per household.

REQUIRED ANALYSIS STRUCTURE — every response MUST include ALL sections:

## Key Metrics Summary
List the 5-10 most critical metrics with exact values from the data.

## Performance Analysis
Interpret what the numbers mean. Identify trends, outliers, and patterns. Compare against insurance industry benchmarks where applicable. Calculate derived metrics (loss ratios, retention rates, penetration rates, LTV) where the raw data allows.

## Segment Analysis
Break down performance by all available dimensions — geography, product line, customer segment, channel, tenure. Identify top and bottom performers with exact figures.

## Conclusions
4-5 definitive conclusions drawn directly from the data. Be specific and decisive.

## Strategic Recommendations
5-6 specific, prioritized actions each tied to a specific data finding. Include expected impact where calculable from the data.

## Risk & Watch List
2-3 metrics or trends requiring immediate attention with specific threshold values.

Example of BAD: 'As the professional analyst conducting this analysis, I have identified key insights.'
Example of GOOD: 'The portfolio contains 652 households with a 64.3% multi-product conversion rate. The 180-day cross-sell cycle is 23% longer than the industry benchmark of 146 days, suggesting friction in the upsell process. Households with 3+ products show 91.2% retention vs 67.4% for single-product households — a 23.8-point retention premium that justifies aggressive cross-sell investment.'"""

# ============================================================================
# LLM QUERY FUNCTION
# ============================================================================

def query_llm(prompt, max_tokens=6000):
    """Query OpenAI-compatible endpoint for analysis"""
    api_key = os.getenv("OPENAI_API_KEY", "")
    try:
        resp = requests.post(
            "https://api.openai.com/v1/chat/completions",
            headers={"Authorization": f"Bearer {api_key}", "Content-Type": "application/json"},
            json={
                "model": "gpt-4o",
                "messages": [{"role": "user", "content": prompt}],
                "max_tokens": max_tokens,
                "temperature": 0.3
            },
            timeout=120
        )
        data = resp.json()
        return data["choices"][0]["message"]["content"]
    except Exception as e:
        return f"LLM Error: {str(e)}"

# ============================================================================
# DATA LOADING & UTILITIES
# ============================================================================

def load_data(file):
    """Load CSV or Excel file"""
    try:
        path = file.name if hasattr(file, 'name') else file
        if path.endswith('.csv'):
            df = pd.read_csv(path)
        elif path.endswith(('.xls', '.xlsx')):
            df = pd.read_excel(path)
        else:
            return None, "Unsupported file type (use CSV or Excel)"
        return df, None
    except Exception as e:
        return None, f"Error loading file: {str(e)}"

def normalize_columns(df):
    """Map common column name variations to standard names used by analysis functions."""
    col_map = {}
    lower_cols = {c.lower().replace(' ', '_').replace('-', '_'): c for c in df.columns}

    # Define mappings: standard_name -> list of possible patterns
    mappings = {
        'household_id': ['household_id', 'householdid', 'hh_id', 'hhid', 'customer_id', 'customerid', 'policyholder_id'],
        'policy_id': ['policy_id', 'policyid', 'policy_number', 'policynumber'],
        'premium': ['annual_premium', 'premium', 'premium_amount', 'total_premium', 'written_premium', 'gross_premium', 'policy_premium', 'quote_amount'],
        'claims': ['paid_amount', 'claim_amount', 'claims', 'total_claims', 'claimed_amount', 'claims_paid', 'incurred_amount', 'loss_amount'],
        'claim_amount_raw': ['claimed_amount', 'claim_amount_raw', 'gross_claim'],
        'reserve_amount': ['reserve_amount', 'reserves', 'case_reserve', 'outstanding_reserve'],
        'line_of_business': ['policy_type', 'line_of_business', 'lob', 'product', 'product_type', 'coverage_type', 'product_quoted'],
        'payment_status': ['payment_behavior', 'payment_status', 'payment_freq', 'pay_status'],
        'risk_tier': ['risk_tier', 'risk_level', 'risk_class', 'risk_score', 'risk_category'],
        'state': ['state', 'region', 'territory', 'geography', 'location'],
        'zip_code': ['zip_code', 'zipcode', 'zip', 'postal_code'],
        'city': ['city', 'metro', 'metro_area'],
        'age': ['age', 'customer_age', 'policyholder_age', 'insured_age'],
        'gender': ['gender', 'sex'],
        'income_bracket': ['income_bracket', 'income', 'income_level', 'income_range', 'income_group'],
        'life_stage': ['life_stage', 'lifestage', 'segment', 'customer_segment', 'life_cycle_stage'],
        'effective_date': ['effective_date', 'start_date', 'inception_date', 'policy_start', 'bind_date'],
        'expiration_date': ['expiration_date', 'end_date', 'policy_end', 'renewal_date'],
        'deductible': ['deductible', 'deductible_amount'],
        'coverage_limit': ['coverage_limit', 'coverage_amount', 'sum_insured', 'limit'],
        'agent_id': ['agent_id', 'agentid', 'agent', 'producer_id', 'broker_id'],
        'bundle_flag': ['bundle_flag', 'bundled', 'multi_policy', 'bundle'],
        'acquisition_channel': ['acquisition_channel', 'channel', 'source', 'lead_source', 'marketing_channel'],
        'acquisition_date': ['acquisition_date', 'lead_date', 'signup_date', 'quote_date'],
        'policy_status': ['policy_status', 'status', 'active_flag'],
        'cancellation_date': ['cancellation_date', 'cancel_date', 'lapse_date'],
        'cancellation_reason': ['cancellation_reason', 'cancel_reason', 'decline_reason', 'lapse_reason'],
        'renewal_count': ['renewal_count', 'renewals', 'num_renewals'],
        'tenure_years': ['tenure_years', 'tenure', 'customer_tenure', 'policy_tenure', 'years_active'],
        'claim_id': ['claim_id', 'claimid', 'claim_number'],
        'claim_date': ['claim_date', 'loss_date', 'date_of_loss', 'incident_date'],
        'claim_type': ['claim_type', 'loss_type', 'peril', 'cause_of_loss', 'loss_cause'],
        'claim_status': ['claim_status', 'status'],
        'days_to_close': ['days_to_close', 'settlement_days', 'cycle_time', 'resolution_days'],
        'fraud_flag': ['fraud_flag', 'fraud', 'suspicious', 'fraud_indicator'],
        'fraud_score': ['fraud_score', 'fraud_probability', 'risk_score'],
        'severity': ['severity', 'loss_severity', 'claim_severity'],
        'claimant_injury': ['claimant_injury', 'injury', 'injury_type', 'bodily_injury'],
        'subrogation_flag': ['subrogation_flag', 'subrogation', 'subro'],
        'lead_id': ['lead_id', 'leadid', 'prospect_id'],
        'bind_flag': ['bind_flag', 'converted', 'bound', 'conversion_flag'],
        'cpa': ['cpa', 'cost_per_acquisition', 'acquisition_cost', 'cac'],
        'days_to_bind': ['days_to_bind', 'days_to_convert', 'conversion_days', 'sales_cycle'],
        'follow_up_count': ['follow_up_count', 'touchpoints', 'contacts', 'follow_ups'],
    }

    for std_name, patterns in mappings.items():
        if std_name in lower_cols:
            # Already exists with standard name
            col_map[lower_cols[std_name]] = std_name
            continue
        for pat in patterns:
            if pat in lower_cols and lower_cols[pat] not in col_map:
                col_map[lower_cols[pat]] = std_name
                break

    if col_map:
        df = df.rename(columns=col_map)
    return df

def load_and_normalize(file):
    """Load file and normalize column names"""
    df, err = load_data(file)
    if err:
        return None, err
    df = normalize_columns(df)
    return df, None

def df_summary(df):
    """Generate summary of dataframe"""
    info_str = f"Rows: {len(df)}, Columns: {len(df.columns)}"
    missing = df.isnull().sum()
    missing_str = missing[missing > 0].to_string() if missing.sum() > 0 else "None"
    preview = df.head(10).to_string()
    return info_str, missing_str, preview

def safe_numeric(val):
    """Safely convert to numeric"""
    try:
        return pd.to_numeric(val, errors='coerce')
    except:
        return pd.Series([np.nan] * len(val))

def get_date_column(df):
    """Find date column in dataframe"""
    for col in df.columns:
        if 'date' in col.lower() or 'month' in col.lower():
            return col
    return None

# ============================================================================
# TAB 1: HOUSEHOLD ECONOMICS (FOUNDATION)
# ============================================================================

def analyze_household_economics(file):
    """Analyze household-level economics: premium, claims, CLV, payment behavior"""
    df, err = load_and_normalize(file)
    if err:
        return err, None

    # Detect which columns are available
    has_claims = 'claims' in df.columns and df['claims'].notna().any() and safe_numeric(df['claims']).sum() > 0
    has_premium = 'premium' in df.columns and df['premium'].notna().any()

    # Ensure numeric columns
    df['premium'] = safe_numeric(df.get('premium', 0))
    if has_claims:
        df['claims'] = safe_numeric(df['claims'])
    else:
        df['claims'] = 0
    df['household_id'] = df.get('household_id', df.get('customer_id', range(len(df))))

    # Group by household
    agg_dict = {'premium': 'sum', 'claims': 'sum'}
    hh_metrics = df.groupby('household_id').agg(agg_dict)
    hh_metrics['policy_count'] = df.groupby('household_id').size()

    # Compute metrics
    total_premium = hh_metrics['premium'].sum()
    total_claims = hh_metrics['claims'].sum()
    loss_ratio = (total_claims / total_premium) if total_premium > 0 else 0
    avg_premium_per_hh = hh_metrics['premium'].mean()
    avg_claims_per_hh = hh_metrics['claims'].mean()
    avg_policies_per_hh = hh_metrics['policy_count'].mean()

    # Cross-line analysis
    if 'line_of_business' in df.columns:
        cross_line = df.groupby('household_id')['line_of_business'].nunique().mean()
    else:
        cross_line = 1.0

    # CLV estimate (lifetime value = premium - claims)
    hh_metrics['clv'] = hh_metrics['premium'] - hh_metrics['claims']
    avg_clv = hh_metrics['clv'].mean()

    # Payment behavior — handle common variations like "On Time", "On-Time", "on_time", "Current"
    if 'payment_status' in df.columns:
        on_time_patterns = ['on-time', 'on time', 'on_time', 'current', 'paid', 'good']
        on_time_count = df['payment_status'].astype(str).str.lower().str.strip().apply(
            lambda x: any(p in x for p in on_time_patterns)
        ).sum()
        on_time_rate = on_time_count / len(df) * 100
    else:
        on_time_rate = 100.0

    # Build metrics text — note when claims data is unavailable
    claims_note = ""
    if not has_claims:
        claims_note = "\nNOTE: This dataset does NOT contain claims data. Loss ratio and CLV metrics are incomplete. Upload the Claims/Loss dataset to the Risk & Loss tab for complete loss analysis.\n"

    metrics_text = f"""
HOUSEHOLD ECONOMICS ANALYSIS
=============================================
Total Premium Portfolio: ${total_premium:,.0f}
{'Total Claims Paid: $' + f'{total_claims:,.0f}' if has_claims else 'Total Claims Paid: N/A (no claims column in this dataset)'}
{'Portfolio Loss Ratio: ' + f'{loss_ratio:.2%}' if has_claims else 'Portfolio Loss Ratio: N/A (no claims data available)'}
Number of Households: {len(hh_metrics)}
{claims_note}
PER HOUSEHOLD METRICS:
- Average Premium: ${avg_premium_per_hh:,.2f}
{'- Average Claims: $' + f'{avg_claims_per_hh:,.2f}' if has_claims else '- Average Claims: N/A (no claims data)'}
- Average Customer Lifetime Value (Premium only): ${avg_clv:,.2f}
- Average Policies per Household: {avg_policies_per_hh:.2f}
- Cross-Line Products per Household: {cross_line:.2f}
- On-Time Payment Rate: {on_time_rate:.1f}%

DISTRIBUTION METRICS:
- Premium 75th Percentile: ${hh_metrics['premium'].quantile(0.75):,.2f}
- Premium 90th Percentile: ${hh_metrics['premium'].quantile(0.90):,.2f}
- CLV Top 10% Threshold: ${hh_metrics['clv'].quantile(0.90):,.2f}
"""

    # Build prompt
    prompt = f"""{ANALYST_TONE}

{metrics_text}

Data Summary: {len(df)} policies across {len(hh_metrics)} households.

Provide a concise executive summary of household economics. Focus on:
1. Portfolio health and loss ratio implications
2. Premium efficiency and CLV concentration
3. Cross-sell opportunity assessment
4. Payment behavior insights
5. Actionable recommendations for household-level growth

Keep analysis to 300 words maximum. Use only the computed numbers provided above.
{WATERMARK}"""

    analysis = query_llm(prompt)

    # Create visualizations
    fig_clv = go.Figure(data=[go.Histogram(x=hh_metrics['clv'], nbinsx=40, name='CLV')])
    fig_clv.update_layout(title="Customer Lifetime Value Distribution", xaxis_title="CLV ($)", yaxis_title="Households")

    fig_loss = go.Figure()
    if 'segment' in df.columns:
        loss_by_seg = df.groupby('segment').apply(lambda x: (x['claims'].sum() / x['premium'].sum() if x['premium'].sum() > 0 else 0))
        fig_loss = go.Figure(data=[go.Bar(x=loss_by_seg.index, y=loss_by_seg.values, name='Loss Ratio')])
        fig_loss.update_layout(title="Loss Ratio by Customer Segment", xaxis_title="Segment", yaxis_title="Loss Ratio")

    fig_scatter = go.Figure(data=[go.Scatter(
        x=hh_metrics['premium'],
        y=hh_metrics['claims'],
        mode='markers',
        marker=dict(size=5, opacity=0.6),
        name='Households'
    )])
    fig_scatter.update_layout(title="Premium vs Claims by Household", xaxis_title="Premium ($)", yaxis_title="Claims ($)")

    return analysis, fig_to_dict(fig_clv)

# ============================================================================
# TAB 2: HOUSEHOLD PENETRATION (GROWTH)
# ============================================================================

def analyze_penetration(file):
    """Analyze cross-sell and penetration metrics"""
    df, err = load_and_normalize(file)
    if err:
        return err, None

    df['household_id'] = df.get('household_id', df.get('customer_id', range(len(df))))

    # Products per household
    if 'line_of_business' in df.columns:
        products_per_hh = df.groupby('household_id')['line_of_business'].nunique()
        avg_products = products_per_hh.mean()
        cross_sell_rate = (products_per_hh > 1).sum() / len(products_per_hh) * 100
    else:
        avg_products = 1.0
        cross_sell_rate = 0.0

    # Bundle analysis
    if 'bundle_flag' in df.columns:
        bundle_rate = is_true(df['bundle_flag']).sum() / len(df) * 100
    else:
        bundle_rate = 0.0

    # Market penetration
    total_hh = len(df.groupby('household_id'))
    if 'segment' in df.columns:
        penetration_by_segment = df.groupby('segment')['household_id'].nunique()
    else:
        penetration_by_segment = pd.Series([total_hh])

    metrics_text = f"""
HOUSEHOLD PENETRATION ANALYSIS
=============================================
Total Households: {total_hh}
Average Products per Household: {avg_products:.2f}
Cross-Sell Rate (Multi-Product): {cross_sell_rate:.1f}%
Bundle Rate: {bundle_rate:.1f}%
Single-Product Households: {(products_per_hh == 1).sum() if 'line_of_business' in df.columns else 'N/A'}
Multi-Product Households: {(products_per_hh > 1).sum() if 'line_of_business' in df.columns else 'N/A'}

PENETRATION BY SEGMENT:
{penetration_by_segment.to_string()}

KEY INSIGHTS:
- Available cross-sell opportunity: {100 - cross_sell_rate:.1f}% of households
- Current bundling penetration: {bundle_rate:.1f}%
"""

    prompt = f"""{ANALYST_TONE}

{metrics_text}

Total policies analyzed: {len(df)}.

Analyze penetration and cross-sell opportunity. Cover:
1. Current cross-sell performance vs best-in-class (typically 40-60%)
2. Segment-level penetration gaps
3. Bundle strategy effectiveness
4. Growth levers for household expansion
5. Recommended cross-sell targets

Keep to 250 words. Use only computed metrics provided.
{WATERMARK}"""

    analysis = query_llm(prompt)

    # Visualization
    if 'line_of_business' in df.columns:
        products_dist = products_per_hh.value_counts().sort_index()
        fig = go.Figure(data=[go.Bar(x=products_dist.index, y=products_dist.values, name='Households')])
        fig.update_layout(title="Products per Household Distribution", xaxis_title="Number of Products", yaxis_title="Count")
    else:
        fig = go.Figure()

    return analysis, fig_to_dict(fig)

# ============================================================================
# TAB 3: HOUSEHOLD RETENTION (CRITICAL)
# ============================================================================

def analyze_retention(file):
    """Analyze retention, churn, renewal rates"""
    df, err = load_and_normalize(file)
    if err:
        return err, None

    df['household_id'] = df.get('household_id', df.get('customer_id', range(len(df))))

    # Use policy_status or renewal_status for retention analysis
    status_col = None
    for col_name in ['policy_status', 'renewal_status', 'status']:
        if col_name in df.columns:
            status_col = col_name
            break

    # Retention metrics
    total_hh = df['household_id'].nunique()
    if status_col:
        status_lower = df[status_col].astype(str).str.lower().str.strip()
        renewed = status_lower.isin(['renewed', 'active']).sum()
        cancelled = status_lower.isin(['cancelled', 'canceled', 'lapsed', 'non-renewed', 'non_renewed']).sum()
        renewal_rate = renewed / len(df) * 100 if len(df) > 0 else 0
    else:
        renewed = 0
        cancelled = 0
        renewal_rate = 0

    if 'churn_flag' in df.columns:
        churn_rate = is_true(df['churn_flag']).sum() / len(df) * 100
        retained = is_false(df['churn_flag']).sum()
        retention_rate = (retained / len(df)) * 100
    else:
        churn_rate = cancelled / len(df) * 100 if len(df) > 0 else 0
        retention_rate = 100 - churn_rate

    # Tenure analysis
    if 'tenure_years' in df.columns:
        avg_tenure = safe_numeric(df['tenure_years']).mean() * 12  # convert to months
    elif 'tenure_months' in df.columns:
        avg_tenure = safe_numeric(df['tenure_months']).mean()
    else:
        avg_tenure = 24  # default

    # Retention by bundle
    if 'bundle_flag' in df.columns:
        bundle_churn = df[df['bundle_flag'] == 1]['churn_flag'].mean() * 100 if 'churn_flag' in df.columns else 0
        single_churn = df[df['bundle_flag'] == 0]['churn_flag'].mean() * 100 if 'churn_flag' in df.columns else 0
    else:
        bundle_churn = 0
        single_churn = 0

    metrics_text = f"""
HOUSEHOLD RETENTION ANALYSIS
=============================================
Total Households: {total_hh}
Overall Retention Rate: {retention_rate:.1f}%
Overall Churn Rate: {churn_rate:.1f}%
Renewal Rate (Policies): {renewal_rate:.1f}%
Average Customer Tenure: {avg_tenure:.1f} months

RETENTION BY BUNDLE:
- Multi-Product Churn Rate: {bundle_churn:.1f}%
- Single-Product Churn Rate: {single_churn:.1f}%
- Bundle Stickiness Factor: {(single_churn - bundle_churn) if single_churn > bundle_churn else 'Positive'}

COHORT INSIGHTS:
- 12-Month Retention: {(1 - churn_rate/100)**1:.1%}
- 24-Month Retention: {(1 - churn_rate/100)**2:.1%}
"""

    prompt = f"""{ANALYST_TONE}

{metrics_text}

Provide retention analysis summary:
1. Current retention vs industry benchmarks (92-96%)
2. Churn drivers and bundle impact
3. Tenure implications and lifecycle stage assessment
4. Cohort retention trajectory
5. Retention improvement priorities

Keep to 250 words. Use only provided metrics.
{WATERMARK}"""

    analysis = query_llm(prompt)

    # Visualization - churn by segment
    if 'segment' in df.columns and 'churn_flag' in df.columns:
        churn_by_seg = df.groupby('segment')['churn_flag'].mean() * 100
        fig = go.Figure(data=[go.Bar(x=churn_by_seg.index, y=churn_by_seg.values, name='Churn Rate')])
        fig.update_layout(title="Churn Rate by Segment (%)", xaxis_title="Segment", yaxis_title="Churn Rate (%)")
    else:
        fig = go.Figure()

    return analysis, fig_to_dict(fig)

# ============================================================================
# TAB 4: CUSTOMER SEGMENTATION
# ============================================================================

def analyze_segmentation(file):
    """Analyze customer segmentation by demographics and risk"""
    df, err = load_and_normalize(file)
    if err:
        return err, None

    df['premium'] = safe_numeric(df.get('premium', 0))
    df['claims'] = safe_numeric(df.get('claims', 0))

    segments_found = []

    # Age-based segmentation
    if 'age' in df.columns:
        df['age_num'] = safe_numeric(df['age'])
        age_seg = pd.cut(df['age_num'], bins=[0, 25, 45, 65, 100], labels=['Young', 'Middle', 'Senior', 'Mature'])
        age_dist = age_seg.value_counts()
        segments_found.append(f"Age Distribution: {age_dist.to_dict()}")

    # Income-based segmentation
    if 'income' in df.columns:
        df['income_num'] = safe_numeric(df['income'])
        income_seg = pd.cut(df['income_num'], bins=[0, 50000, 100000, 150000, np.inf], labels=['Low', 'Mid', 'High', 'Premium'])
        income_dist = income_seg.value_counts()
        segments_found.append(f"Income Segments: {income_dist.to_dict()}")

    # Risk tier
    if 'risk_score' in df.columns:
        df['risk_num'] = safe_numeric(df['risk_score'])
        risk_seg = pd.cut(df['risk_num'], bins=[0, 30, 70, 100], labels=['Low Risk', 'Standard', 'High Risk'])
        risk_dist = risk_seg.value_counts()
        segments_found.append(f"Risk Distribution: {risk_dist.to_dict()}")

    # Household size
    if 'household_size' in df.columns:
        hh_size_dist = safe_numeric(df['household_size']).value_counts().sort_index()
        segments_found.append(f"Household Size: {hh_size_dist.to_dict()}")

    # Life stage (if available)
    if 'life_stage' in df.columns:
        life_dist = df['life_stage'].value_counts()
        segments_found.append(f"Life Stage: {life_dist.to_dict()}")

    # Value by segment
    if 'segment' in df.columns:
        segment_value = df.groupby('segment').agg({
            'premium': 'sum',
            'claims': 'sum',
            'household_id': 'count'
        }).rename(columns={'household_id': 'count'})
        segment_value['profit'] = segment_value['premium'] - segment_value['claims']
        value_summary = segment_value.to_string()
    else:
        value_summary = "No segment column found"

    metrics_text = f"""
CUSTOMER SEGMENTATION ANALYSIS
=============================================
Total Policies: {len(df)}
Total Households: {df.get('household_id', range(len(df))).nunique() if 'household_id' in df.columns else 'N/A'}

SEGMENTATION DIMENSIONS:
{chr(10).join(segments_found)}

VALUE BY SEGMENT:
{value_summary}
"""

    prompt = f"""{ANALYST_TONE}

{metrics_text}

Deliver segmentation analysis:
1. Segment distribution and composition
2. Value concentration (top 20% rule)
3. Risk profile by segment
4. Premium efficiency by segment
5. Recommended segment strategy (focus, growth, harvest)

Keep to 300 words. Use only computed data provided.
{WATERMARK}"""

    analysis = query_llm(prompt)

    # Visualization
    if 'segment' in df.columns:
        segment_counts = df['segment'].value_counts()
        fig = go.Figure(data=[go.Bar(x=segment_counts.index, y=segment_counts.values, name='Count')])
        fig.update_layout(title="Segment Distribution", xaxis_title="Segment", yaxis_title="Count")
    else:
        fig = go.Figure()

    return analysis, fig_to_dict(fig)

# ============================================================================
# TAB 5: RISK & LOSS ANALYSIS
# ============================================================================

def analyze_risk_loss(file):
    """Analyze loss ratios, claims frequency, severity, fraud indicators"""
    df, err = load_and_normalize(file)
    if err:
        return err, None

    df['premium'] = safe_numeric(df.get('premium', 0))
    df['claims'] = safe_numeric(df.get('claims', 0))
    df['household_id'] = df.get('household_id', df.get('customer_id', range(len(df))))

    # Aggregate to household level
    hh_loss = df.groupby('household_id').agg({
        'premium': 'sum',
        'claims': 'sum',
        'household_id': 'count'
    }).rename(columns={'household_id': 'claim_count'})

    hh_loss['loss_ratio'] = (hh_loss['claims'] / hh_loss['premium']).fillna(0)

    # Claims frequency and severity
    total_claims = df.get('claim_id', []).nunique() if 'claim_id' in df.columns else len(df[df['claims'] > 0])
    avg_claim_size = (df['claims'].sum() / total_claims) if total_claims > 0 else 0
    claim_frequency = (df['claims'] > 0).sum() / len(df) * 100

    # Portfolio loss ratio
    portfolio_lr = (df['claims'].sum() / df['premium'].sum()) if df['premium'].sum() > 0 else 0

    # Fraud indicators
    fraud_indicators = 0
    if 'fraud_flag' in df.columns:
        fraud_indicators = is_true(df['fraud_flag']).sum()
        fraud_rate = fraud_indicators / len(df) * 100
    else:
        fraud_rate = 0.0

    # Loss ratio distribution
    loss_q25 = hh_loss['loss_ratio'].quantile(0.25)
    loss_q50 = hh_loss['loss_ratio'].quantile(0.50)
    loss_q75 = hh_loss['loss_ratio'].quantile(0.75)
    loss_q90 = hh_loss['loss_ratio'].quantile(0.90)

    metrics_text = f"""
RISK & LOSS ANALYSIS
=============================================
Portfolio Loss Ratio: {portfolio_lr:.2%}
Total Claims Count: {total_claims}
Average Claim Size: ${avg_claim_size:,.2f}
Claim Frequency (% with claims): {claim_frequency:.1f}%
Fraud Indicators: {fraud_indicators} ({fraud_rate:.2f}%)

LOSS RATIO DISTRIBUTION (by Household):
- 25th Percentile: {loss_q25:.2%}
- 50th Percentile (Median): {loss_q50:.2%}
- 75th Percentile: {loss_q75:.2%}
- 90th Percentile: {loss_q90:.2%}
- Households with Loss Ratio > 1.0: {(hh_loss['loss_ratio'] > 1.0).sum()}

SEVERITY ANALYSIS:
- High-Frequency Claimants (>median): {(hh_loss['claim_count'] > hh_loss['claim_count'].median()).sum()}
- Underperforming Book (LR > 100%): {(hh_loss['loss_ratio'] > 1.0).sum() / len(hh_loss) * 100:.1f}%
"""

    prompt = f"""{ANALYST_TONE}

{metrics_text}

Analyze risk and loss performance:
1. Portfolio loss ratio health and trend implications
2. Claims frequency and severity drivers
3. High-risk household identification and concentration
4. Fraud exposure assessment
5. Risk mitigation and pricing adjustment recommendations

Keep to 300 words. Use only provided metrics.
{WATERMARK}"""

    analysis = query_llm(prompt)

    # Visualization
    fig = go.Figure(data=[go.Histogram(x=hh_loss['loss_ratio'], nbinsx=50, name='Loss Ratio')])
    fig.update_layout(title="Loss Ratio Distribution Across Households", xaxis_title="Loss Ratio", yaxis_title="Count")

    return analysis, fig_to_dict(fig)

# ============================================================================
# TAB 6: PRICING & PROFITABILITY
# ============================================================================

def analyze_pricing(file):
    """Analyze pricing, profitability, rate adequacy"""
    df, err = load_and_normalize(file)
    if err:
        return err, None

    df['premium'] = safe_numeric(df.get('premium', 0))
    df['claims'] = safe_numeric(df.get('claims', 0))
    df['household_id'] = df.get('household_id', df.get('customer_id', range(len(df))))

    # Household-level profitability
    hh_profit = df.groupby('household_id').agg({
        'premium': 'sum',
        'claims': 'sum'
    })
    hh_profit['profit'] = hh_profit['premium'] - hh_profit['claims']
    hh_profit['profit_margin'] = hh_profit['profit'] / hh_profit['premium']

    total_profit = hh_profit['profit'].sum()
    avg_profit = hh_profit['profit'].mean()
    profitable_hh = (hh_profit['profit'] > 0).sum()
    unprofitable_hh = (hh_profit['profit'] < 0).sum()

    # Rate adequacy
    if 'premium_target' in df.columns:
        df['premium_target_num'] = safe_numeric(df['premium_target'])
        actual_vs_target = df['premium'].sum() / df['premium_target_num'].sum() if df['premium_target_num'].sum() > 0 else 1.0
    else:
        actual_vs_target = 1.0

    # Loss ratio by segment for rate adequacy
    if 'segment' in df.columns:
        segment_lr = df.groupby('segment').apply(lambda x: x['claims'].sum() / x['premium'].sum() if x['premium'].sum() > 0 else 0)
        underpriced = (segment_lr > 0.65).sum()
        overpriced = (segment_lr < 0.45).sum()
    else:
        underpriced = 0
        overpriced = 0

    metrics_text = f"""
PRICING & PROFITABILITY ANALYSIS
=============================================
Total Portfolio Profit: ${total_profit:,.0f}
Average Profit per Household: ${avg_profit:,.2f}
Profitable Households: {profitable_hh} ({profitable_hh / len(hh_profit) * 100:.1f}%)
Unprofitable Households: {unprofitable_hh} ({unprofitable_hh / len(hh_profit) * 100:.1f}%)

PROFITABILITY DISTRIBUTION:
- Top 25% Profit Contribution: ${hh_profit['profit'].quantile(0.75):,.2f}
- Median Profit: ${hh_profit['profit'].median():,.2f}
- Bottom 25% Profit: ${hh_profit['profit'].quantile(0.25):,.2f}
- Profit Margin Range: {hh_profit['profit_margin'].min():.2%} to {hh_profit['profit_margin'].max():.2%}

RATE ADEQUACY:
- Actual vs Target Premium: {actual_vs_target:.2%}
- Underpriced Segments (LR>65%): {underpriced}
- Overpriced Segments (LR<45%): {overpriced}
- Portfolio Expansion Capacity: {'Limited' if unprofitable_hh / len(hh_profit) > 0.2 else 'Moderate' if unprofitable_hh / len(hh_profit) > 0.1 else 'Strong'}
"""

    prompt = f"""{ANALYST_TONE}

{metrics_text}

Deliver pricing and profitability analysis:
1. Profitability concentration and efficiency
2. Segment-level rate adequacy assessment
3. Pricing optimization opportunities
4. Unprofitable book action plan
5. Rate action recommendations (increase, decrease, refocus)

Keep to 300 words. Use only provided metrics.
{WATERMARK}"""

    analysis = query_llm(prompt)

    # Visualization
    fig = go.Figure(data=[go.Histogram(x=hh_profit['profit'], nbinsx=40, name='Profit')])
    fig.update_layout(title="Profitability Distribution by Household", xaxis_title="Profit ($)", yaxis_title="Count")

    return analysis, fig_to_dict(fig)

# ============================================================================
# TAB 7: ACQUISITION FUNNEL
# ============================================================================

def analyze_acquisition(file):
    """Analyze quote-to-bind, lead source performance, CPA"""
    df, err = load_and_normalize(file)
    if err:
        return err, None

    df['premium'] = safe_numeric(df.get('premium', 0))

    # Conversion metrics — use bind_flag if available (each row is a lead/quote)
    if 'bind_flag' in df.columns:
        quotes = len(df)
        binds = is_true(df['bind_flag']).sum()
        quote_to_bind_rate = (binds / quotes * 100) if quotes > 0 else 0
    elif 'quote_flag' in df.columns:
        quotes = is_true(df['quote_flag']).sum()
        binds = len(df)
        quote_to_bind_rate = (binds / quotes * 100) if quotes > 0 else 0
    else:
        quotes = len(df)
        binds = len(df)
        quote_to_bind_rate = 100.0

    # Lead source performance
    lead_col = 'acquisition_channel' if 'acquisition_channel' in df.columns else ('lead_source' if 'lead_source' in df.columns else None)
    if lead_col:
        lead_source_perf = df.groupby(lead_col).agg({
            'household_id': 'count',
            'premium': 'sum'
        }).rename(columns={'household_id': 'binds'})
        lead_source_perf['avg_premium'] = lead_source_perf['premium'] / lead_source_perf['binds']
        lead_source_text = lead_source_perf.to_string()
    else:
        lead_source_text = "No lead source data available"

    # CPA calculation
    cpa_col = 'cpa' if 'cpa' in df.columns else ('acquisition_cost' if 'acquisition_cost' in df.columns else None)
    if cpa_col:
        total_acq_cost = safe_numeric(df[cpa_col]).sum()
        cpa = total_acq_cost / binds if binds > 0 else 0
        cpl = total_acq_cost / quotes if quotes > 0 else 0
    else:
        cpa = 0
        cpl = 0

    # Conversion funnel stages
    if 'application_flag' in df.columns:
        applications = is_true(df['application_flag']).sum()
        app_to_bind = (binds / applications * 100) if applications > 0 else 0
    else:
        applications = binds
        app_to_bind = 100.0

    metrics_text = f"""
ACQUISITION FUNNEL ANALYSIS
=============================================
Total Quotes Generated: {quotes}
Total Policies Bound: {binds}
Quote-to-Bind Rate: {quote_to_bind_rate:.1f}%
Application-to-Bind Rate: {app_to_bind:.1f}%

LEAD PERFORMANCE:
{lead_source_text}

ACQUISITION ECONOMICS:
- Cost Per Acquisition (CPA): ${cpa:,.2f}
- Cost Per Lead (CPL): ${cpl:,.2f}
- Average Premium per Bind: ${df['premium'].mean():,.2f}
- CPA as % of First Year Premium: {(cpa / df['premium'].mean() * 100) if df['premium'].mean() > 0 else 0:.1f}%
- Payback Period (months): {(cpa / (df['premium'].mean() * 0.6)) if df['premium'].mean() > 0 else 0:.1f}

CONVERSION EFFICIENCY:
- Funnel Depth: {(binds / quotes) if quotes > 0 else 0:.1%}
- Lead Quality Score: {'High' if quote_to_bind_rate > 50 else 'Medium' if quote_to_bind_rate > 30 else 'Low'}
"""

    prompt = f"""{ANALYST_TONE}

{metrics_text}

Analyze acquisition performance:
1. Funnel conversion efficiency vs targets (typically 30-50% Q2B)
2. Lead source performance and ROI ranking
3. Acquisition cost efficiency and payback
4. Channel optimization recommendations
5. Volume vs profitability tradeoff analysis

Keep to 300 words. Use only provided metrics.
{WATERMARK}"""

    analysis = query_llm(prompt)

    # Visualization - funnel
    funnel_data = [quotes, applications if 'application_flag' in df.columns else int(binds * 0.8), binds]
    funnel_labels = ['Quotes', 'Applications', 'Binds']
    fig = go.Figure(go.Funnel(y=funnel_labels, x=funnel_data))
    fig.update_layout(title="Acquisition Funnel")

    return analysis, fig_to_dict(fig)

# ============================================================================
# TAB 8: GEOGRAPHIC ANALYSIS
# ============================================================================

def analyze_geographic(file):
    """Analyze penetration, retention, loss ratio by geography"""
    df, err = load_and_normalize(file)
    if err:
        return err, None

    df['premium'] = safe_numeric(df.get('premium', 0))
    df['claims'] = safe_numeric(df.get('claims', 0))

    # Detect geographic column
    geo_col = None
    for col in df.columns:
        if 'zip' in col.lower() or 'state' in col.lower() or 'region' in col.lower():
            geo_col = col
            break

    if geo_col is None:
        return "No geographic column found (ZIP, State, or Region)", None

    # Geographic metrics
    geo_metrics = df.groupby(geo_col).agg({
        'household_id': 'nunique' if 'household_id' in df.columns else 'count',
        'premium': 'sum',
        'claims': 'sum'
    }).rename(columns={'household_id': 'households'})

    geo_metrics['loss_ratio'] = geo_metrics['claims'] / geo_metrics['premium']
    geo_metrics['avg_premium'] = geo_metrics['premium'] / geo_metrics['households']

    if 'churn_flag' in df.columns:
        geo_churn = df.groupby(geo_col)['churn_flag'].mean() * 100
        geo_metrics['churn_rate'] = geo_churn

    top_regions = geo_metrics.nlargest(5, 'premium')
    metrics_text = f"""
GEOGRAPHIC ANALYSIS
=============================================
Total Regions/ZIPs: {len(geo_metrics)}
Total Households: {geo_metrics['households'].sum()}
Total Premium: ${geo_metrics['premium'].sum():,.0f}
Total Claims: ${geo_metrics['claims'].sum():,.0f}

TOP 5 REGIONS BY PREMIUM:
{top_regions[['households', 'premium', 'loss_ratio', 'avg_premium']].to_string()}

GEOGRAPHIC DISTRIBUTION:
- Concentration (Top 5): {(top_regions['premium'].sum() / geo_metrics['premium'].sum() * 100):.1f}%
- Highest Loss Ratio Region: {geo_metrics['loss_ratio'].idxmax()} ({geo_metrics['loss_ratio'].max():.2%})
- Lowest Loss Ratio Region: {geo_metrics['loss_ratio'].idxmin()} ({geo_metrics['loss_ratio'].min():.2%})
- Premium Variance: {geo_metrics['avg_premium'].std() / geo_metrics['avg_premium'].mean():.2f}x

REGIONAL PERFORMANCE TIERS:
- Tier 1 (Top 20%): ${geo_metrics['premium'].quantile(0.80):,.0f}+
- Tier 2 (50-80%): ${geo_metrics['premium'].quantile(0.50):,.0f} - ${geo_metrics['premium'].quantile(0.80):,.0f}
- Tier 3 (Bottom 50%): <${geo_metrics['premium'].quantile(0.50):,.0f}
"""

    prompt = f"""{ANALYST_TONE}

{metrics_text}

Deliver geographic strategy analysis:
1. Market concentration and diversification assessment
2. Regional underperformance diagnosis
3. Penetration opportunity by region
4. Pricing and product fit by geography
5. Growth, harvest, or exit recommendations by region

Keep to 300 words. Use only provided metrics.
{WATERMARK}"""

    analysis = query_llm(prompt)

    # Visualization
    top10 = geo_metrics.nlargest(10, 'premium')
    fig = go.Figure(data=[go.Bar(x=top10.index, y=top10['premium'], name='Premium')])
    fig.update_layout(title="Premium Distribution by Region (Top 10)", xaxis_title=geo_col, yaxis_title="Premium ($)")

    return analysis, fig_to_dict(fig)

# ============================================================================
# TAB 9: LIFECYCLE / JOURNEY ANALYSIS
# ============================================================================

def analyze_lifecycle(file):
    """Analyze time-to-second-policy, time-to-churn, lifecycle stages"""
    df, err = load_and_normalize(file)
    if err:
        return err, None

    df['household_id'] = df.get('household_id', df.get('customer_id', range(len(df))))
    df['premium'] = safe_numeric(df.get('premium', 0))

    # Time to second policy
    hh_policy_count = df.groupby('household_id').size()
    multi_product_hh = len(hh_policy_count[hh_policy_count > 1])

    if 'policy_date' in df.columns or 'date' in df.columns:
        date_col = 'policy_date' if 'policy_date' in df.columns else 'date'
        df[date_col] = pd.to_datetime(df[date_col], errors='coerce')

        # Calculate days to second policy
        hh_dates = df.groupby('household_id')[date_col].agg(['min', 'max', 'count'])
        hh_dates['days_to_expand'] = (hh_dates['max'] - hh_dates['min']).dt.days
        avg_days_to_expand = hh_dates[hh_dates['count'] > 1]['days_to_expand'].mean()
    else:
        avg_days_to_expand = 180  # default estimate

    # Lifecycle stage (based on tenure and activity)
    if 'tenure_months' in df.columns:
        df['tenure_num'] = safe_numeric(df['tenure_months'])
        df['lifecycle_stage'] = pd.cut(df['tenure_num'],
                                       bins=[0, 6, 24, 60, np.inf],
                                       labels=['Acquisition', 'Growth', 'Maturity', 'Decline'])
        lifecycle_dist = df['lifecycle_stage'].value_counts()
        lifecycle_text = lifecycle_dist.to_string()
    else:
        lifecycle_text = "No tenure data available"

    # Stage transitions
    acquisition_to_growth = (hh_policy_count > 1).sum() / len(hh_policy_count) * 100

    if 'churn_flag' in df.columns:
        # Time to churn
        churn_hh = df[df['churn_flag'] == 1]['household_id'].nunique()
        if 'tenure_months' in df.columns:
            avg_tenure_before_churn = df[df['churn_flag'] == 1]['tenure_num'].mean()
        else:
            avg_tenure_before_churn = 12
    else:
        churn_hh = 0
        avg_tenure_before_churn = 0

    metrics_text = f"""
LIFECYCLE & JOURNEY ANALYSIS
=============================================
Total Households: {df['household_id'].nunique()}
Multi-Product Households: {multi_product_hh} ({multi_product_hh / df['household_id'].nunique() * 100:.1f}%)
Average Time to Second Policy: {avg_days_to_expand:.0f} days
Acquisition-to-Growth Conversion: {acquisition_to_growth:.1f}%

LIFECYCLE STAGE DISTRIBUTION:
{lifecycle_text}

CHURN ANALYSIS:
- Churned Households: {churn_hh}
- Average Tenure Before Churn: {avg_tenure_before_churn:.1f} months
- High-Risk Window: Months 6-12 ({(df[(df['tenure_num'] >= 6) & (df['tenure_num'] <= 12) & is_true(df['churn_flag'])].shape[0] if 'churn_flag' in df.columns and 'tenure_num' in df.columns else 0)} churns)

GROWTH TRAJECTORY:
- Acquisition Phase Households: {len(df) if 'lifecycle_stage' not in df.columns else len(df[df['lifecycle_stage'] == 'Acquisition'])}
- Growth Phase Households: {len(df[df['lifecycle_stage'] == 'Growth']) if 'lifecycle_stage' in df.columns else 'N/A'}
- Maturity Phase Households: {len(df[df['lifecycle_stage'] == 'Maturity']) if 'lifecycle_stage' in df.columns else 'N/A'}
"""

    prompt = f"""{ANALYST_TONE}

{metrics_text}

Analyze customer lifecycle journey:
1. Acquisition effectiveness and time-to-value metrics
2. Growth phase progression and expansion readiness
3. Churn risk indicators and prevention strategy
4. Lifecycle stage health and transition rates
5. Interventions by stage (nurture new, expand growth, prevent decline)

Keep to 300 words. Use only provided metrics.
{WATERMARK}"""

    analysis = query_llm(prompt)

    # Visualization
    if 'lifecycle_stage' in df.columns:
        stage_dist = df['lifecycle_stage'].value_counts()
        fig = go.Figure(data=[go.Bar(x=stage_dist.index, y=stage_dist.values, name='Households')])
        fig.update_layout(title="Household Distribution by Lifecycle Stage", xaxis_title="Stage", yaxis_title="Count")
    else:
        fig = go.Figure()

    return analysis, fig_to_dict(fig)

# ============================================================================
# TAB 10: CUSTOM VISUALIZATIONS
# ============================================================================

def create_custom_viz(file, x_column, y_column, chart_type="scatter"):
    """Create custom visualization from user-selected columns"""
    df, err = load_and_normalize(file)
    if err:
        return err, None

    if x_column not in df.columns or y_column not in df.columns:
        return "Column not found in data", None

    try:
        df[x_column] = safe_numeric(df[x_column])
        df[y_column] = safe_numeric(df[y_column])

        if chart_type == "scatter":
            fig = go.Figure(data=[go.Scatter(x=df[x_column], y=df[y_column], mode='markers')])
        elif chart_type == "bar":
            agg_data = df.groupby(x_column)[y_column].sum()
            fig = go.Figure(data=[go.Bar(x=agg_data.index, y=agg_data.values)])
        elif chart_type == "histogram":
            fig = go.Figure(data=[go.Histogram(x=df[x_column])])
        else:
            fig = go.Figure()

        fig.update_layout(title=f"{y_column} vs {x_column}", xaxis_title=x_column, yaxis_title=y_column)
        return "Chart created successfully", fig_to_dict(fig)
    except Exception as e:
        return f"Error creating chart: {str(e)}", None

# ============================================================================
# TAB 11: REPORT GENERATION (with PowerPoint)
# ============================================================================

def generate_pptx(report_text, df, report_type, file_name=""):
    """Generate a PowerPoint presentation from the report analysis."""
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt, Emu
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
    except ImportError:
        return None

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    NAVY = RGBColor(0x0A, 0x16, 0x28)
    GOLD = RGBColor(0xC9, 0xA8, 0x4C)
    WHITE = RGBColor(0xF8, 0xF9, 0xFC)
    MUTED = RGBColor(0x88, 0x92, 0xA4)
    NAVY_MID = RGBColor(0x11, 0x22, 0x40)

    def add_bg(slide, color=NAVY):
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = color

    def add_text_box(slide, text, left, top, width, height, font_size=14, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Arial"):
        txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = font_name
        p.alignment = alignment
        return tf

    # ─── SLIDE 1: Title ───
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    add_bg(slide)
    add_text_box(slide, "INSURANCE ECONOMICS", 1, 1.5, 11, 1, 38, GOLD, True, PP_ALIGN.CENTER)
    report_titles = {
        "executive_summary": "Executive Summary Report",
        "detailed_analysis": "Detailed Analysis Report",
        "financial_impact": "Financial Impact Report",
        "market_trends": "Market Trends Report",
        "competitor_analysis": "Competitor Analysis Report"
    }
    add_text_box(slide, report_titles.get(report_type, "Analysis Report"), 1, 2.8, 11, 0.8, 26, WHITE, False, PP_ALIGN.CENTER)
    add_text_box(slide, f"Generated: {datetime.now().strftime('%B %d, %Y')}", 1, 4, 11, 0.5, 14, MUTED, False, PP_ALIGN.CENTER)
    add_text_box(slide, f"Data: {file_name} | {len(df)} records | {len(df.columns)} columns", 1, 4.6, 11, 0.5, 12, MUTED, False, PP_ALIGN.CENTER)
    add_text_box(slide, "QuantusData.ai — Powered by Existential Gateway, LLC", 1, 6.2, 11, 0.5, 11, MUTED, False, PP_ALIGN.CENTER)

    # ─── SLIDE 2: Data Overview ───
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_text_box(slide, "DATA OVERVIEW", 0.8, 0.5, 11, 0.6, 22, GOLD, True)

    num_cols = df.select_dtypes(include=[np.number]).columns.tolist()
    cat_cols = df.select_dtypes(include=['object']).columns.tolist()
    missing_pct = (df.isnull().sum().sum() / (len(df) * len(df.columns)) * 100)

    overview = f"Records: {len(df):,}    |    Columns: {len(df.columns)}    |    Numeric: {len(num_cols)}    |    Categorical: {len(cat_cols)}    |    Missing: {missing_pct:.1f}%"
    add_text_box(slide, overview, 0.8, 1.3, 11.5, 0.5, 13, WHITE)

    if num_cols:
        stats_text = "KEY NUMERIC STATISTICS:\n"
        for col in num_cols[:8]:
            vals = df[col].dropna()
            if len(vals) > 0:
                stats_text += f"  {col}: mean={vals.mean():,.1f}, median={vals.median():,.1f}, std={vals.std():,.1f}, min={vals.min():,.1f}, max={vals.max():,.1f}\n"
        add_text_box(slide, stats_text, 0.8, 2.0, 11.5, 3, 11, WHITE, font_name="Courier New")

    if cat_cols:
        cat_text = "TOP CATEGORICAL DISTRIBUTIONS:\n"
        for col in cat_cols[:5]:
            top3 = df[col].value_counts().head(3)
            cat_text += f"  {col}: " + ", ".join([f"{k} ({v})" for k, v in top3.items()]) + "\n"
        add_text_box(slide, cat_text, 0.8, 5.0, 11.5, 2, 11, WHITE, font_name="Courier New")

    # ─── SLIDES 3+: Analysis Content ───
    # Split report into sections and create slides
    sections = report_text.split('\n\n')
    current_slide_text = ""
    slide_num = 0

    for section in sections:
        section = section.strip()
        if not section:
            continue
        if len(current_slide_text) + len(section) > 800 or slide_num == 0:
            if current_slide_text and slide_num > 0:
                # Write previous slide
                slide = prs.slides.add_slide(prs.slide_layouts[6])
                add_bg(slide)
                add_text_box(slide, f"ANALYSIS — Part {slide_num}", 0.8, 0.5, 11, 0.6, 20, GOLD, True)
                add_text_box(slide, current_slide_text.strip(), 0.8, 1.3, 11.5, 5.5, 13, WHITE)
            current_slide_text = section + "\n\n"
            slide_num += 1
        else:
            current_slide_text += section + "\n\n"

    # Write final content slide
    if current_slide_text.strip():
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        add_bg(slide)
        add_text_box(slide, f"ANALYSIS — Part {slide_num}", 0.8, 0.5, 11, 0.6, 20, GOLD, True)
        add_text_box(slide, current_slide_text.strip(), 0.8, 1.3, 11.5, 5.5, 13, WHITE)

    # ─── FINAL SLIDE: Disclaimer ───
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_text_box(slide, "DISCLAIMER", 1, 2, 11, 0.6, 22, GOLD, True, PP_ALIGN.CENTER)
    add_text_box(slide, "This report was generated by AI and is intended for informational purposes only.\nAll results should be reviewed by a qualified insurance professional.\n\n© 2026 Existential Gateway, LLC. All Rights Reserved.\nUnauthorized reproduction strictly prohibited.\nLicensing: existentialgateway@gmail.com", 1.5, 3, 10, 3, 14, MUTED, False, PP_ALIGN.CENTER)

    # Save
    import tempfile
    pptx_path = tempfile.mktemp(suffix='.pptx', prefix='InsEcon_Report_')
    prs.save(pptx_path)
    return pptx_path


def generate_report(file, report_type):
    """Generate comprehensive AI-powered report with PowerPoint download"""
    df, err = load_and_normalize(file)
    if err:
        return err, None

    if not report_type or report_type == 'true' or report_type == 'True':
        report_type = 'executive_summary'

    file_name = file.name if hasattr(file, 'name') else 'uploaded'

    # Build comprehensive metrics
    num_cols = df.select_dtypes(include=[np.number]).columns.tolist()
    cat_cols = df.select_dtypes(include=['object']).columns.tolist()
    stats = df.describe().to_string() if num_cols else "No numeric columns"
    missing = df.isnull().sum()
    missing_str = missing[missing > 0].to_string() if missing.sum() > 0 else "No missing values"

    # Build column distributions
    distributions = ""
    for col in cat_cols[:6]:
        dist = df[col].value_counts().head(5)
        distributions += f"\n{col}:\n" + "\n".join([f"  {k}: {v} ({v/len(df)*100:.1f}%)" for k, v in dist.items()])

    for col in num_cols[:6]:
        vals = df[col].dropna()
        if len(vals) > 0:
            distributions += f"\n{col}: mean={vals.mean():,.2f}, median={vals.median():,.2f}, std={vals.std():,.2f}, min={vals.min():,.2f}, max={vals.max():,.2f}"

    report_prompts = {
        "executive_summary": "Write a comprehensive executive summary covering all key findings, portfolio health, risk indicators, and strategic recommendations. Include specific numbers and percentages.",
        "detailed_analysis": "Write a detailed analysis covering every column and metric in the data. Break down by segments, identify patterns, outliers, and correlations. Include specific numbers for every finding.",
        "financial_impact": "Analyze the financial implications of this data. Focus on revenue, loss ratios, profitability, CLV, premium distribution, and financial risk. Quantify all impacts.",
        "market_trends": "Analyze market trends visible in this data. Cover growth patterns, segment shifts, geographic trends, channel performance, and competitive positioning.",
        "competitor_analysis": "Analyze competitive positioning based on this data. Cover market share implications, pricing competitiveness, product gaps, and strategic advantages."
    }

    prompt = f"""{ANALYST_TONE}

INSURANCE ECONOMICS REPORT — {report_type.upper().replace('_', ' ')}
Data File: {file_name}
Records: {len(df)} | Columns: {len(df.columns)}

COLUMN NAMES: {', '.join(df.columns.tolist())}

SUMMARY STATISTICS:
{stats}

DISTRIBUTIONS:
{distributions}

MISSING DATA:
{missing_str}

SAMPLE DATA (first 5 rows):
{df.head(5).to_string()}

{report_prompts.get(report_type, report_prompts['executive_summary'])}

Structure your report with clear sections:
1. Executive Overview (2-3 sentences)
2. Key Metrics and Findings (with specific numbers)
3. Segment Analysis (break down by available categories)
4. Risk Assessment
5. Strategic Recommendations (numbered, actionable)

Write at least 500 words. Use specific numbers from the data for every finding.
{WATERMARK}"""

    report_text = query_llm(prompt, max_tokens=4000)

    # Generate PowerPoint
    pptx_path = generate_pptx(report_text, df, report_type, file_name)

    if pptx_path:
        report_text += f"\n\n---\n**PowerPoint report generated successfully.** Click the download button below to save."
        return report_text, pptx_path

    return report_text, None

# ============================================================================
# TAB 12: DATA CLEANING
# ============================================================================

def clean_data(file, remove_duplicates=True, fill_method="drop"):
    """Data cleaning utility"""
    df, err = load_and_normalize(file)
    if err:
        return err, None

    original_rows = len(df)

    # Remove duplicates
    if remove_duplicates:
        df = df.drop_duplicates()

    # Handle missing values
    if fill_method == "drop":
        df = df.dropna()
    elif fill_method == "mean":
        df = df.fillna(df.mean(numeric_only=True))

    new_rows = len(df)
    rows_removed = original_rows - new_rows

    summary = f"""
DATA CLEANING REPORT
=============================================
Original Rows: {original_rows}
Rows After Cleaning: {new_rows}
Rows Removed: {rows_removed}
Duplicates Removed: {original_rows - new_rows if remove_duplicates else 0}
Fill Method: {fill_method}

Cleaned data ready for analysis.
"""

    return summary

# ============================================================================
# TAB 13: MULTI-DATASET ANALYSIS
# ============================================================================

def multi_dataset_analysis(file1, file2, file3, analysis_type, custom_instructions):
    """Analyze multiple datasets with comprehensive cross-dataset insights"""
    dfs = []
    errors = []
    file_names = []

    for file, idx in [(file1, 1), (file2, 2), (file3, 3)]:
        if file is None:
            continue
        df, err = load_and_normalize(file)
        if err:
            errors.append(f"Dataset {idx}: {err}")
        else:
            fname = file.name if hasattr(file, 'name') else f'Dataset_{idx}'
            file_names.append(fname)
            dfs.append((df, fname, idx))

    if errors:
        return "\n".join(errors)

    if len(dfs) < 2:
        return "Please upload at least 2 datasets for cross-dataset analysis."

    # Build detailed summary for each dataset
    dataset_summaries = ""
    all_columns_by_ds = {}
    for df, fname, idx in dfs:
        num_cols = df.select_dtypes(include=[np.number]).columns.tolist()
        cat_cols = df.select_dtypes(include=['object']).columns.tolist()

        ds_summary = f"\n{'='*60}\nDATASET {idx}: {fname}\n{'='*60}\n"
        ds_summary += f"Records: {len(df):,} | Columns: {len(df.columns)} | Numeric: {len(num_cols)} | Categorical: {len(cat_cols)}\n"
        ds_summary += f"Columns: {', '.join(df.columns.tolist())}\n"

        # Numeric stats
        for col in num_cols[:8]:
            vals = df[col].dropna()
            if len(vals) > 0:
                ds_summary += f"  {col}: total={vals.sum():,.2f}, mean={vals.mean():,.2f}, median={vals.median():,.2f}, min={vals.min():,.2f}, max={vals.max():,.2f}\n"

        # Categorical distributions
        for col in cat_cols[:6]:
            dist = df[col].value_counts().head(5)
            ds_summary += f"  {col}: " + ", ".join([f"{k}={v} ({v/len(df)*100:.1f}%)" for k, v in dist.items()]) + "\n"

        # Missing data
        missing = df.isnull().sum()
        missing_cols = missing[missing > 0]
        if len(missing_cols) > 0:
            ds_summary += f"  Missing data: " + ", ".join([f"{c}={v}" for c, v in missing_cols.items()]) + "\n"
        else:
            ds_summary += f"  No missing data.\n"

        # Key insurance metrics if present
        if 'premium' in df.columns:
            ds_summary += f"  Total Premium: ${safe_numeric(df['premium']).sum():,.2f}\n"
        if 'claims' in df.columns:
            ds_summary += f"  Total Claims: ${safe_numeric(df['claims']).sum():,.2f}\n"
        if 'household_id' in df.columns:
            ds_summary += f"  Unique Households: {df['household_id'].nunique():,}\n"

        dataset_summaries += ds_summary
        all_columns_by_ds[idx] = set(df.columns.tolist())

    # Find common columns across datasets
    common_cols = set.intersection(*all_columns_by_ds.values()) if all_columns_by_ds else set()
    cross_info = f"\nCROSS-DATASET INFO:\n"
    cross_info += f"Common Columns: {', '.join(sorted(common_cols)) if common_cols else 'None'}\n"

    # Auto-detect join key and attempt merge for correlation
    join_key = None
    for candidate in ['household_id', 'policy_id', 'customer_id', 'claim_id', 'lead_id']:
        if candidate in common_cols:
            join_key = candidate
            break

    if join_key:
        cross_info += f"Join Key Detected: {join_key}\n"
        # Check overlap
        id_sets = [set(df[join_key].dropna().unique()) for df, _, _ in dfs if join_key in df.columns]
        if len(id_sets) >= 2:
            overlap = id_sets[0].intersection(id_sets[1])
            cross_info += f"Overlapping {join_key} values between Dataset 1 & 2: {len(overlap):,}\n"

    analysis_prompts = {
        "comprehensive_comparison": "Perform a comprehensive comparison of all datasets. For each dataset, explain what it represents, its key metrics, and how it relates to the other datasets. Identify patterns, discrepancies, and complementary insights across datasets.",
        "cross_correlation": "Analyze correlations and relationships between the datasets. Identify which metrics in one dataset drive or predict metrics in another. Look for causal chains across datasets.",
        "risk_assessment": "Perform a cross-dataset risk assessment. Identify risk indicators from each dataset and synthesize them into a holistic risk profile. Flag compounding risks visible only when combining datasets.",
        "portfolio_analysis": "Analyze the complete insurance portfolio by synthesizing all datasets. Cover underwriting quality, claims experience, acquisition efficiency, and overall portfolio health."
    }

    prompt = f"""{ANALYST_TONE}

MULTI-DATASET INSURANCE ECONOMICS ANALYSIS
Analysis Type: {analysis_type.upper().replace('_', ' ')}

{dataset_summaries}
{cross_info}

Custom Instructions: {custom_instructions if custom_instructions else 'None'}

{analysis_prompts.get(analysis_type, analysis_prompts['comprehensive_comparison'])}

CRITICAL REQUIREMENTS:
- Reference EACH dataset by name and number (Dataset 1, Dataset 2, etc.)
- Provide specific numbers, percentages, and dollar amounts FROM EACH dataset
- Show how datasets relate to each other and what combined insights emerge
- Give actionable recommendations that synthesize findings across all datasets
- Include at least one specific comparison between datasets (e.g., "Dataset 1 shows X while Dataset 2 reveals Y, which together suggest Z")

Structure your analysis:
1. Dataset Overview (1-2 sentences each, with key metric from each)
2. Cross-Dataset Findings (specific numbers from each dataset)
3. Combined Risk/Opportunity Assessment
4. Discrepancies or Gaps Between Datasets
5. Strategic Recommendations (numbered, referencing specific datasets)

Write at least 600 words. Be comprehensive and reference specific data from ALL datasets.
{WATERMARK}"""

    analysis = query_llm(prompt, max_tokens=5000)
    return analysis

# ============================================================================
# TAB 14: AI CHAT
# ============================================================================

def chat_with_data(file, user_message):
    """Interactive chat — answers the user's specific question about their data"""
    df, err = load_and_normalize(file)
    if err:
        return f"Error loading file: {err}"

    # Build rich data context
    num_cols = df.select_dtypes(include=[np.number]).columns.tolist()
    cat_cols = df.select_dtypes(include=['object']).columns.tolist()

    stats_text = ""
    for col in num_cols[:10]:
        vals = df[col].dropna()
        if len(vals) > 0:
            stats_text += f"  {col}: count={len(vals)}, sum={vals.sum():,.2f}, mean={vals.mean():,.2f}, median={vals.median():,.2f}, std={vals.std():,.2f}, min={vals.min():,.2f}, max={vals.max():,.2f}\n"

    dist_text = ""
    for col in cat_cols[:8]:
        dist = df[col].value_counts().head(8)
        dist_text += f"  {col} ({df[col].nunique()} unique): " + ", ".join([f"{k}={v} ({v/len(df)*100:.1f}%)" for k, v in dist.items()]) + "\n"

    sample_rows = df.head(5).to_string()

    prompt = f"""{ANALYST_TONE}

You are answering a specific question from an insurance analyst about their dataset.
Your job is to DIRECTLY ANSWER THE QUESTION with specific numbers from the data.

DATASET: {file.name if hasattr(file, 'name') else 'uploaded'}
Records: {len(df):,} | Columns: {len(df.columns)}
Column Names: {', '.join(df.columns.tolist())}

NUMERIC STATISTICS:
{stats_text}

CATEGORICAL DISTRIBUTIONS:
{dist_text}

SAMPLE DATA (first 5 rows):
{sample_rows}

USER'S QUESTION: {user_message}

INSTRUCTIONS:
1. DIRECTLY answer the question first with specific numbers
2. Then provide additional relevant context or breakdown
3. At the end, suggest 1-2 Python/Pandas code snippets the analyst could run to explore further
   - Format code in markdown code blocks with ```python
   - Use realistic column names from this dataset
   - Examples: groupby aggregations, filtering, pivot tables, visualizations
4. If relevant, also suggest a SQL query they could use

Keep your answer focused on the question. Do not give a general overview of the dataset unless that's what was asked.
{WATERMARK}"""

    response = query_llm(prompt, max_tokens=3000)
    return response

# ============================================================================
# GRADIO UI
# ============================================================================

def create_app():
    """Create Gradio interface"""

    with gr.Blocks(title="Insurance Economics Analyzer") as demo:

        gr.Markdown("""
        # Insurance Economics Analyzer
        **QuantusData.ai** - Professional insurance analytics & insights

        Upload your insurance data (CSV or Excel) and select an analysis type.
        Powered by advanced LLM analysis with Gradio 6.10.0
        """)

        with gr.Tabs():

            # TAB 1: Household Economics
            with gr.Tab("Household Economics"):
                with gr.Row():
                    file_input = gr.File(label="Upload Data", file_types=[".csv", ".xlsx", ".xls"])
                with gr.Row():
                    analyze_btn = gr.Button("Analyze Household Economics", variant="primary")
                with gr.Row():
                    output_text = gr.Textbox(label="Analysis", lines=15)
                with gr.Row():
                    output_chart = gr.JSON(label="CLV Distribution")

                analyze_btn.click(
                    analyze_household_economics,
                    inputs=[file_input],
                    outputs=[output_text, output_chart]
                )

            # TAB 2: Household Penetration
            with gr.Tab("Penetration Analysis"):
                with gr.Row():
                    file_input = gr.File(label="Upload Data", file_types=[".csv", ".xlsx", ".xls"])
                with gr.Row():
                    analyze_btn = gr.Button("Analyze Penetration", variant="primary")
                with gr.Row():
                    output_text = gr.Textbox(label="Analysis", lines=15)
                with gr.Row():
                    output_chart = gr.JSON(label="Products per Household")

                analyze_btn.click(
                    analyze_penetration,
                    inputs=[file_input],
                    outputs=[output_text, output_chart]
                )

            # TAB 3: Household Retention
            with gr.Tab("Retention Analysis"):
                with gr.Row():
                    file_input = gr.File(label="Upload Data", file_types=[".csv", ".xlsx", ".xls"])
                with gr.Row():
                    analyze_btn = gr.Button("Analyze Retention", variant="primary")
                with gr.Row():
                    output_text = gr.Textbox(label="Analysis", lines=15)
                with gr.Row():
                    output_chart = gr.JSON(label="Churn by Segment")

                analyze_btn.click(
                    analyze_retention,
                    inputs=[file_input],
                    outputs=[output_text, output_chart]
                )

            # TAB 4: Customer Segmentation
            with gr.Tab("Segmentation"):
                with gr.Row():
                    file_input = gr.File(label="Upload Data", file_types=[".csv", ".xlsx", ".xls"])
                with gr.Row():
                    analyze_btn = gr.Button("Analyze Segmentation", variant="primary")
                with gr.Row():
                    output_text = gr.Textbox(label="Analysis", lines=15)
                with gr.Row():
                    output_chart = gr.JSON(label="Segment Distribution")

                analyze_btn.click(
                    analyze_segmentation,
                    inputs=[file_input],
                    outputs=[output_text, output_chart]
                )

            # TAB 5: Risk & Loss Analysis
            with gr.Tab("Risk & Loss"):
                with gr.Row():
                    file_input = gr.File(label="Upload Data", file_types=[".csv", ".xlsx", ".xls"])
                with gr.Row():
                    analyze_btn = gr.Button("Analyze Risk & Loss", variant="primary")
                with gr.Row():
                    output_text = gr.Textbox(label="Analysis", lines=15)
                with gr.Row():
                    output_chart = gr.JSON(label="Loss Ratio Distribution")

                analyze_btn.click(
                    analyze_risk_loss,
                    inputs=[file_input],
                    outputs=[output_text, output_chart]
                )

            # TAB 6: Pricing & Profitability
            with gr.Tab("Pricing & Profitability"):
                with gr.Row():
                    file_input = gr.File(label="Upload Data", file_types=[".csv", ".xlsx", ".xls"])
                with gr.Row():
                    analyze_btn = gr.Button("Analyze Pricing", variant="primary")
                with gr.Row():
                    output_text = gr.Textbox(label="Analysis", lines=15)
                with gr.Row():
                    output_chart = gr.JSON(label="Profitability Distribution")

                analyze_btn.click(
                    analyze_pricing,
                    inputs=[file_input],
                    outputs=[output_text, output_chart]
                )

            # TAB 7: Acquisition Funnel
            with gr.Tab("Acquisition Funnel"):
                with gr.Row():
                    file_input = gr.File(label="Upload Data", file_types=[".csv", ".xlsx", ".xls"])
                with gr.Row():
                    analyze_btn = gr.Button("Analyze Acquisition", variant="primary")
                with gr.Row():
                    output_text = gr.Textbox(label="Analysis", lines=15)
                with gr.Row():
                    output_chart = gr.JSON(label="Funnel")

                analyze_btn.click(
                    analyze_acquisition,
                    inputs=[file_input],
                    outputs=[output_text, output_chart]
                )

            # TAB 8: Geographic Analysis
            with gr.Tab("Geographic"):
                with gr.Row():
                    file_input = gr.File(label="Upload Data", file_types=[".csv", ".xlsx", ".xls"])
                with gr.Row():
                    analyze_btn = gr.Button("Analyze Geographic", variant="primary")
                with gr.Row():
                    output_text = gr.Textbox(label="Analysis", lines=15)
                with gr.Row():
                    output_chart = gr.JSON(label="Premium by Region")

                analyze_btn.click(
                    analyze_geographic,
                    inputs=[file_input],
                    outputs=[output_text, output_chart]
                )

            # TAB 9: Lifecycle Analysis
            with gr.Tab("Lifecycle"):
                with gr.Row():
                    file_input = gr.File(label="Upload Data", file_types=[".csv", ".xlsx", ".xls"])
                with gr.Row():
                    analyze_btn = gr.Button("Analyze Lifecycle", variant="primary")
                with gr.Row():
                    output_text = gr.Textbox(label="Analysis", lines=15)
                with gr.Row():
                    output_chart = gr.JSON(label="Lifecycle Stage Distribution")

                analyze_btn.click(
                    analyze_lifecycle,
                    inputs=[file_input],
                    outputs=[output_text, output_chart]
                )

            # TAB 10: Custom Visualizations
            with gr.Tab("Custom Visualizations"):
                with gr.Row():
                    file_input = gr.File(label="Upload Data", file_types=[".csv", ".xlsx", ".xls"])
                with gr.Row():
                    x_col = gr.Textbox(label="X Column", placeholder="e.g., premium")
                    y_col = gr.Textbox(label="Y Column", placeholder="e.g., claims")
                    chart_type = gr.Dropdown(["scatter", "bar", "histogram"], label="Chart Type")
                with gr.Row():
                    viz_btn = gr.Button("Create Visualization", variant="primary")
                with gr.Row():
                    viz_text = gr.Textbox(label="Status")
                with gr.Row():
                    viz_chart = gr.JSON(label="Visualization")

                viz_btn.click(
                    create_custom_viz,
                    inputs=[file_input, x_col, y_col, chart_type],
                    outputs=[viz_text, viz_chart]
                )

            # TAB 11: Report Generation
            with gr.Tab("Report Generation"):
                with gr.Row():
                    file_input = gr.File(label="Upload Data", file_types=[".csv", ".xlsx", ".xls"])
                with gr.Row():
                    report_type = gr.Textbox(label="Report Type", value="executive_summary", visible=False)
                with gr.Row():
                    report_btn = gr.Button("Generate Report", variant="primary")
                with gr.Row():
                    report_output = gr.Textbox(label="Report", lines=20)
                with gr.Row():
                    report_file = gr.File(label="Download PowerPoint Report", visible=False)

                report_btn.click(
                    generate_report,
                    inputs=[file_input, report_type],
                    outputs=[report_output, report_file]
                )

            # TAB 12: Data Cleaning
            with gr.Tab("Data Cleaning"):
                with gr.Row():
                    file_input = gr.File(label="Upload Data", file_types=[".csv", ".xlsx", ".xls"])
                with gr.Row():
                    remove_dups = gr.Checkbox(label="Remove Duplicates", value=True)
                    fill_method = gr.Dropdown(["drop", "mean"], label="Missing Data Handling")
                with gr.Row():
                    clean_btn = gr.Button("Clean Data", variant="primary")
                with gr.Row():
                    clean_output = gr.Textbox(label="Cleaning Report")

                clean_btn.click(
                    clean_data,
                    inputs=[file_input, remove_dups, fill_method],
                    outputs=[clean_output]
                )

            # TAB 13: Multi-Dataset Analysis
            with gr.Tab("Multi-Dataset Analysis"):
                with gr.Row():
                    file1 = gr.File(label="Dataset 1", file_types=[".csv", ".xlsx", ".xls"])
                    file2 = gr.File(label="Dataset 2", file_types=[".csv", ".xlsx", ".xls"])
                    file3 = gr.File(label="Dataset 3 (optional)", file_types=[".csv", ".xlsx", ".xls"])
                with gr.Row():
                    analysis_type = gr.Textbox(label="Analysis Type", value="comprehensive_comparison", visible=False)
                with gr.Row():
                    custom_instr = gr.Textbox(label="Custom Instructions", placeholder="Any specific analysis focus?")
                with gr.Row():
                    multi_btn = gr.Button("Analyze Datasets", variant="primary")
                with gr.Row():
                    multi_output = gr.Textbox(label="Analysis", lines=15)

                multi_btn.click(
                    multi_dataset_analysis,
                    inputs=[file1, file2, file3, analysis_type, custom_instr],
                    outputs=[multi_output],
                    api_name="analyze_multi_dataset"
                )

            # TAB 14: AI Chat
            with gr.Tab("AI Chat"):
                with gr.Row():
                    file_input = gr.File(label="Upload Data", file_types=[".csv", ".xlsx", ".xls"])
                with gr.Row():
                    msg = gr.Textbox(label="Ask a question about your data", scale=4)
                    send_btn = gr.Button("Send", variant="primary", scale=1)
                with gr.Row():
                    chat_output = gr.Textbox(label="Response", lines=15)

                send_btn.click(
                    chat_with_data,
                    inputs=[file_input, msg],
                    outputs=[chat_output],
                    api_name="chat_with_data"
                )

    return demo

# ============================================================================
# LAUNCH
# ============================================================================

if __name__ == "__main__":
    demo = create_app()
    demo.launch(server_name="0.0.0.0", server_port=7875)
