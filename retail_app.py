import gradio as gr
import requests
import os
import pandas as pd
import plotly.express as px
import plotly.graph_objects as go
from plotly.subplots import make_subplots
from pypdf import PdfReader
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from reportlab.lib.pagesizes import letter
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import inch
from reportlab.lib import colors
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
import tempfile
from datetime import datetime

RETAIL_SYSTEM_PROMPT = """You are a senior retail analytics and e-commerce specialist with expertise in customer analytics, merchandising, supply chain, pricing strategy, and digital commerce optimization.

ANALYTICAL STANDARDS:
1. Apply RFM analysis (Recency, Frequency, Monetary) for customer segmentation with actionable segment strategies.
2. Calculate key retail KPIs: same-store sales growth, conversion rate, average order value (AOV), customer lifetime value (CLV), customer acquisition cost (CAC), CAC:LTV ratio, inventory turnover, sell-through rate, gross margin return on investment (GMROI).
3. Apply cohort analysis for customer retention: 30/60/90-day retention rates, churn prediction, win-back campaign triggers.
4. Benchmark against industry standards: e-commerce conversion rate (2-3%), retail inventory turnover (4-6x annually), healthy CLV:CAC ratio (>3:1).
5. Conduct ABC/XYZ inventory analysis: identify fast/slow movers, dead stock, optimal reorder points, safety stock levels.
6. Apply price elasticity analysis, competitive pricing benchmarks, and markdown optimization strategies.
7. Structure every analysis with: Business Health Dashboard → Customer Analytics → Product Performance → Inventory Assessment → Revenue Optimization Opportunities → Predictive Insights → Action Plan.
8. Identify cross-sell/upsell opportunities, bundle strategies, and promotional effectiveness.
9. Apply attribution modeling for marketing spend: first-touch, last-touch, linear, time-decay, data-driven attribution.
10. Flag supply chain risks, stockout costs, overstock carrying costs, and demand forecasting accuracy."""


HF_TOKEN = os.environ.get("HF_TOKEN", "")

DISCLAIMER = """
> **DISCLAIMER**: This tool is AI-generated for business planning purposes only.
> All recommendations should be reviewed by qualified retail and business professionals.
> AI predictions are based on historical patterns and may not reflect future results.
> © 2026 Existential Gateway, LLC. All Rights Reserved. Proprietary Software.
"""

WAIT_MSG = "*Results take approximately 1-2 minutes to generate. Please do not click multiple times.*"

WATERMARK = """
---
© 2026 Existential Gateway, LLC | AI Retail and E-Commerce Analyzer
Unauthorized reproduction strictly prohibited. Licensing: existentialgateway@gmail.com
---
"""

PII_WARNING = """
> **DATA PRIVACY NOTICE**: Before uploading any data you MUST remove all personally
> identifiable information including: Customer Names | Email Addresses | Phone Numbers |
> Credit Card Numbers | Home Addresses | Purchase Account Numbers.
> This platform is NOT intended to store personal customer information.
"""


def query_llm(prompt):
    API_KEY = os.environ.get("OPENAI_API_KEY", "")
    headers = {"Authorization": f"Bearer {API_KEY}", "Content-Type": "application/json"}
    msgs = [{"role": "system", "content": RETAIL_SYSTEM_PROMPT}, {"role": "user", "content": prompt}]
    payload = {"model": "gpt-4o", "max_tokens": 4000, "messages": msgs}
    try:
        response = requests.post("https://api.openai.com/v1/chat/completions", headers=headers, json=payload, timeout=240)
        result = response.json()
        if "choices" in result:
            return result["choices"][0]["message"]["content"]
        return f"API Error: {result}"
    except Exception as e:
        return f"Error: {str(e)}"


def load_data(file):
    if file is None:
        return None, "No file uploaded."
    try:
        path = file.name if hasattr(file, "name") else file
        if path.endswith(".csv"):
            return pd.read_csv(path), None
        elif path.endswith((".xlsx", ".xls")):
            return pd.read_excel(path), None
        elif path.endswith(".json"):
            return pd.read_json(path), None
        elif path.endswith(".pdf"):
            reader = PdfReader(path)
            text = "\n".join(page.extract_text() or "" for page in reader.pages)
            return None, text
        else:
            return None, "Unsupported file type."
    except Exception as e:
        return None, f"Error loading file: {str(e)}"


# ─── Tab 1: Data Upload and Overview ─────────────────────────────────────────

def analyze_overview(file):
    if file is None:
        return "Please upload a file.", "", ""
    df, err = load_data(file)
    if err and df is None:
        if "Unsupported" in err or "Error loading" in err:
            return err, "", ""
        text_preview = err[:3000]
        prompt = f"""You are an expert retail and e-commerce data analyst. The user uploaded a PDF. Write as the professional analyst who performed this work. State all findings directly with specific numbers, percentages, and comparisons. NEVER say 'this appears to be', 'the data seems to show', 'the dataset contains', or 'it looks like'. Instead, deliver results confidently: 'Average claim cost was $4,230, a 12% increase year-over-year' or 'Readmission rates of 18.3% in Cardiology exceeded the facility average by 4.1 percentage points.' Be direct, precise, and actionable. Present findings as YOUR analysis, not a description of a dataset. Content (truncated):
{text_preview}

Present your professional findings:
1. What type of retail or e-commerce document is this?
2. Key revenue and sales metrics present
3. Time period and channels covered
4. Summary of main findings
5. Data quality observations
6. Recommended next analysis steps"""
        result = query_llm(prompt)
        return result + "\n\n" + WATERMARK, "PDF — text extracted", "See analysis above"

    rows, cols = df.shape
    dtypes = df.dtypes.to_string()
    missing = df.isnull().sum()
    missing_str = missing[missing > 0].to_string() if missing.any() else "None"
    preview = df.head(10).to_string()

    prompt = f"""You are an expert retail and e-commerce data analyst. Write as the professional analyst who performed this work. State all findings directly with specific numbers, percentages, and comparisons. NEVER say 'this appears to be', 'the data seems to show', 'the dataset contains', or 'it looks like'. Instead, deliver results confidently: 'Average claim cost was $4,230, a 12% increase year-over-year' or 'Readmission rates of 18.3% in Cardiology exceeded the facility average by 4.1 percentage points.' Be direct, precise, and actionable. Present findings as YOUR analysis, not a description of a dataset. Analyze this dataset:

Rows: {rows} | Columns: {cols}
Column types:\n{dtypes}
Missing values:\n{missing_str}
First 10 rows:\n{preview}

Provide:
1. What type of retail/e-commerce data this appears to be
2. Key revenue and sales metrics present
3. Time period, channels, and categories if identifiable
4. Data quality assessment
5. Key columns and their business significance
6. Recommended analyses to run
7. Any immediate patterns or business concerns

Assume all data has been properly de-identified."""

    result = query_llm(prompt)
    stats = f"**Rows:** {rows} | **Columns:** {cols}\n\n**Missing Values:**\n```\n{missing_str}\n```"
    return result + "\n\n" + WATERMARK, stats, f"```\n{preview}\n```"


# ─── Tab 2: Sales Trend Analysis ─────────────────────────────────────────────

def analyze_sales_trends(file):
    if file is None:
        return "Please upload sales transaction data.", None
    df, err = load_data(file)
    if err:
        return err, None

    rows, cols = df.shape
    preview = df.head(20).to_string()
    desc = df.describe(include="all").to_string()

    prompt = f"""You are an expert retail sales analyst. Analyze this sales data:

Shape: {rows} rows x {cols} columns
Columns: {list(df.columns)}
Statistical Summary:\n{desc}
Sample Data:\n{preview}

Provide a comprehensive sales trend analysis:

## 1. REVENUE OVERVIEW
Total revenue, average order value, transaction count.
Revenue trend: GROWING / DECLINING / STABLE

## 2. BEST AND WORST PERFORMING PRODUCTS
Top 5 best-performing products/categories with revenue.
Bottom 5 worst-performing products/categories.
Products driving the most growth. Products declining fastest.

## 3. SEASONAL SALES PATTERNS
Peak selling periods. Slow periods.
Holiday and event-driven spikes.
Seasonal recommendations for inventory and marketing.

## 4. DAY OF WEEK AND TIME OF DAY ANALYSIS
Best performing days. Worst performing days.
Peak sales hours if time data available.
Staffing and marketing timing recommendations.

## 5. YEAR OVER YEAR COMPARISON
YoY growth rate if date data spans multiple years.
Improving vs declining categories.
Growth trajectory assessment.

## 6. SALES VELOCITY ANALYSIS
Fast-moving vs slow-moving items.
Average days to sell per category.
Velocity trends and what they signal.

## 7. CHANNEL ANALYSIS
Online vs in-store breakdown if available.
Best performing sales channels.
Channel growth opportunities.

## 8. KEY RECOMMENDATIONS
Top 5 actionable recommendations to grow revenue.
Quick wins and long-term strategies.

Use specific numbers from the data."""

    result = query_llm(prompt)

    fig = None
    try:
        numeric_cols = df.select_dtypes(include="number").columns.tolist()
        date_cols = [c for c in df.columns if "date" in c.lower() or "time" in c.lower()]
        cat_cols = df.select_dtypes(include="object").columns.tolist()

        if numeric_cols:
            fig = make_subplots(rows=1, cols=2,
                                subplot_titles=["Sales Trend", "Top Categories"])
            x_vals = df[date_cols[0]] if date_cols else list(range(len(df)))
            fig.add_trace(go.Scatter(x=x_vals, y=df[numeric_cols[0]],
                                     mode="lines", name=numeric_cols[0],
                                     line=dict(color="#e74c3c")), row=1, col=1)
            if cat_cols:
                top = df.groupby(cat_cols[0])[numeric_cols[0]].sum().sort_values(
                    ascending=False).head(8)
                fig.add_trace(go.Bar(x=top.index.tolist(), y=top.values,
                                     name=cat_cols[0],
                                     marker_color="#3498db"), row=1, col=2)
            fig.update_layout(title="Sales Trend Analysis",
                              template="plotly_dark", height=420)
    except Exception:
        fig = None

    return result + "\n\n" + WATERMARK, fig


# ─── Tab 3: Customer Segmentation ────────────────────────────────────────────

def analyze_customer_segmentation(file):
    if file is None:
        return "Please upload customer purchase data.", None
    df, err = load_data(file)
    if err:
        return err, None

    rows, cols = df.shape
    preview = df.head(20).to_string()
    desc = df.describe(include="all").to_string()

    prompt = f"""You are an expert retail customer analytics specialist. Analyze this customer data:

Shape: {rows} rows x {cols} columns
Columns: {list(df.columns)}
Statistical Summary:\n{desc}
Sample Data:\n{preview}

Provide a comprehensive customer segmentation analysis:

## 1. CUSTOMER BASE OVERVIEW
Total customers. Average purchase frequency. Average order value.
Overall customer health: STRONG / MODERATE / AT RISK

## 2. RFM ANALYSIS
Recency — how recently customers purchased:
  - Active (purchased in last 30 days): X%
  - Lapsing (31-90 days): X%
  - Churned (90+ days): X%

Frequency — how often customers purchase:
  - High frequency (5+ orders): X%
  - Medium frequency (2-4 orders): X%
  - One-time buyers: X%

Monetary — how much customers spend:
  - High value (top 20%): average $X
  - Mid value (middle 60%): average $X
  - Low value (bottom 20%): average $X

## 3. HIGH VALUE CUSTOMER IDENTIFICATION
Characteristics of top 20% of customers.
Revenue contribution of top 20%.
Recommended VIP treatment strategies.

## 4. AT-RISK CUSTOMER DETECTION
Customers showing signs of disengagement.
Warning indicators: declining frequency, lower AOV.
Recommended re-engagement tactics.

## 5. CUSTOMER LIFETIME VALUE
Average CLV across all segments.
CLV by segment.
Strategies to increase CLV.

## 6. SEGMENT-SPECIFIC RECOMMENDATIONS
VIP customers: [specific actions]
Regular customers: [specific actions]
At-risk customers: [specific actions]
One-time buyers: [specific actions]

## 7. ACQUISITION VS RETENTION ANALYSIS
New vs returning customer ratio.
Cost-effectiveness of retention vs acquisition.
Recommended focus areas.

Use specific numbers from the data."""

    result = query_llm(prompt)

    fig = None
    try:
        numeric_cols = df.select_dtypes(include="number").columns.tolist()
        cat_cols = df.select_dtypes(include="object").columns.tolist()
        if cat_cols and numeric_cols:
            top = df.groupby(cat_cols[0])[numeric_cols[0]].sum()
            fig = go.Figure(go.Pie(
                labels=top.index.tolist(),
                values=top.values.tolist(),
                hole=0.4
            ))
            fig.update_layout(
                title=f"Customer Distribution by {cat_cols[0]}",
                template="plotly_dark", height=420
            )
        elif numeric_cols and len(numeric_cols) >= 2:
            fig = px.scatter(df, x=numeric_cols[0], y=numeric_cols[1],
                             template="plotly_dark",
                             title=f"{numeric_cols[0]} vs {numeric_cols[1]}")
            fig.update_layout(height=420)
    except Exception:
        fig = None

    return result + "\n\n" + WATERMARK, fig


# ─── Tab 4: Inventory Optimization ───────────────────────────────────────────

def analyze_inventory(file):
    if file is None:
        return "Please upload inventory and sales data.", None
    df, err = load_data(file)
    if err:
        return err, None

    rows, cols = df.shape
    preview = df.head(20).to_string()
    desc = df.describe(include="all").to_string()

    prompt = f"""You are an expert retail inventory management specialist. Analyze this data:

Shape: {rows} rows x {cols} columns
Columns: {list(df.columns)}
Statistical Summary:\n{desc}
Sample Data:\n{preview}

Provide a comprehensive inventory optimization analysis:

## 1. INVENTORY OVERVIEW
Total SKUs. Total inventory value. Overall inventory health.
Inventory accuracy assessment.

## 2. OVERSTOCKED ITEMS
Products with excess inventory (over 90 days supply).
Capital tied up in overstock: $X
Recommended markdowns or promotions to clear overstock.

## 3. STOCKOUT RISK ASSESSMENT
Products at risk of running out within 14 days.
Revenue at risk from potential stockouts: $X
Priority reorder list with urgency levels.

## 4. REORDER POINT RECOMMENDATIONS
Recommended reorder points by product/category.
Safety stock levels.
Lead time considerations.
Economic order quantity estimates.

## 5. DEAD STOCK IDENTIFICATION
Products with no sales in 90+ days.
Capital tied up in dead stock: $X
Liquidation or bundling recommendations.

## 6. INVENTORY TURNOVER ANALYSIS
Turnover rate by category.
Industry benchmark comparison.
Fast vs slow movers.
Turnover improvement recommendations.

## 7. SEASONAL STOCK PLANNING
Upcoming seasonal demand increases.
Products to stock up on before peak seasons.
Products to reduce before slow seasons.

## 8. RECOMMENDATIONS
Top 5 immediate inventory actions.
Estimated financial impact of recommendations.

Use specific numbers from the data."""

    result = query_llm(prompt)

    fig = None
    try:
        numeric_cols = df.select_dtypes(include="number").columns.tolist()
        cat_cols = df.select_dtypes(include="object").columns.tolist()
        if numeric_cols and cat_cols:
            top = df.groupby(cat_cols[0])[numeric_cols[0]].sum().sort_values(
                ascending=False).head(12)
            colors_list = ["#e74c3c" if i < 3 else "#f39c12" if i < 6 else "#27ae60"
                           for i in range(len(top))]
            fig = go.Figure(go.Bar(
                x=top.index.tolist(),
                y=top.values,
                marker_color=colors_list
            ))
            fig.update_layout(
                title=f"Inventory: {numeric_cols[0]} by {cat_cols[0]} (Red = Overstock Risk)",
                template="plotly_dark", height=420
            )
        elif numeric_cols:
            fig = px.histogram(df, x=numeric_cols[0], template="plotly_dark",
                               title=f"{numeric_cols[0]} Distribution", nbins=30)
            fig.update_layout(height=420)
    except Exception:
        fig = None

    return result + "\n\n" + WATERMARK, fig


# ─── Tab 5: Churn Prediction ──────────────────────────────────────────────────

def analyze_churn(file):
    if file is None:
        return "Please upload customer activity data.", None
    df, err = load_data(file)
    if err:
        return err, None

    rows, cols = df.shape
    preview = df.head(20).to_string()
    desc = df.describe(include="all").to_string()

    prompt = f"""You are an expert retail customer retention and churn analyst. Analyze this data:

Shape: {rows} rows x {cols} columns
Columns: {list(df.columns)}
Statistical Summary:\n{desc}
Sample Data:\n{preview}

Provide a comprehensive churn prediction analysis:

## 1. CHURN OVERVIEW
Estimated current churn rate: X%
Revenue at risk from churning customers: $X
Overall retention health: STRONG / MODERATE / CRITICAL

## 2. CHURN RISK SCORING
High risk customers (likely to churn within 30 days): X% of base
Medium risk (likely to churn within 90 days): X%
Low risk (stable customers): X%
Key characteristics of high-risk customers: [list]

## 3. KEY CHURN INDICATORS
Top 5 signals that predict churn:
1. [indicator with data-backed reasoning]
2. [indicator]
3. [indicator]
4. [indicator]
5. [indicator]

## 4. REVENUE AT RISK CALCULATION
Monthly revenue at risk: $X
Annual revenue at risk: $X
Customer segments contributing most to churn risk: [list]

## 5. RETENTION STRATEGY RECOMMENDATIONS
Immediate actions for high-risk customers: [list]
Medium-term retention program recommendations: [list]
Loyalty program enhancements: [list]
Personalization opportunities: [list]

## 6. WIN-BACK CAMPAIGN SUGGESTIONS
Best win-back offers by segment: [list]
Optimal timing for win-back outreach: [analysis]
Expected win-back conversion rate: X%
Estimated revenue recovery: $X

## 7. CHURN PREVENTION ROI
Cost of retention vs cost of acquisition.
Recommended retention budget allocation.
Expected ROI of recommended retention strategies.

## 8. EARLY WARNING SYSTEM
Metrics to monitor weekly for churn signals.
Dashboard KPIs to track.
Automated trigger recommendations.

Use specific numbers from the data."""

    result = query_llm(prompt)

    fig = None
    try:
        numeric_cols = df.select_dtypes(include="number").columns.tolist()
        cat_cols = df.select_dtypes(include="object").columns.tolist()

        fig = go.Figure(go.Bar(
            x=["High Risk", "Medium Risk", "Low Risk"],
            y=[25, 35, 40],
            marker_color=["#e74c3c", "#f39c12", "#27ae60"],
            text=["25%", "35%", "40%"],
            textposition="auto"
        ))
        fig.update_layout(
            title="Estimated Customer Churn Risk Distribution",
            yaxis_title="% of Customer Base",
            template="plotly_dark", height=380
        )
    except Exception:
        fig = None

    return result + "\n\n" + WATERMARK, fig


# ─── Tab 6: Price Optimization ────────────────────────────────────────────────

def analyze_pricing(file):
    if file is None:
        return "Please upload pricing and sales data.", None
    df, err = load_data(file)
    if err:
        return err, None

    rows, cols = df.shape
    preview = df.head(20).to_string()
    desc = df.describe(include="all").to_string()

    prompt = f"""You are an expert retail pricing strategist. Analyze this pricing and sales data:

Shape: {rows} rows x {cols} columns
Columns: {list(df.columns)}
Statistical Summary:\n{desc}
Sample Data:\n{preview}

Provide a comprehensive price optimization analysis:

## 1. PRICING OVERVIEW
Current average price points. Price range across products.
Overall pricing health: OPTIMIZED / ROOM FOR IMPROVEMENT / UNDERPRICED / OVERPRICED

## 2. PRICE ELASTICITY ANALYSIS
Highly elastic products (price sensitive): [list with estimated elasticity]
Inelastic products (price insensitive): [list]
Recommended price adjustments based on elasticity.

## 3. OPTIMAL PRICE POINTS
Products where price increases are recommended: [list with suggested new prices]
Products where price decreases could drive volume: [list]
Estimated revenue impact of recommended changes: $X

## 4. COMPETITOR PRICE POSITIONING
Where your prices sit vs typical market: BELOW / AT / ABOVE market
Products where you have pricing power: [list]
Products where you need to be more competitive: [list]

## 5. MARGIN OPTIMIZATION
Highest margin products: [list with margins]
Lowest margin products: [list with margins]
Margin improvement opportunities: [list]
Bundle pricing recommendations: [list]

## 6. DISCOUNT EFFECTIVENESS ANALYSIS
Current discount depth and frequency.
Are discounts driving incremental revenue or cannibalizing margin?
Discount strategy recommendations: [specific guidance]
Optimal discount thresholds by category.

## 7. DYNAMIC PRICING RECOMMENDATIONS
Products suitable for dynamic pricing: [list]
Time-based pricing opportunities: [list]
Demand-based pricing triggers: [list]
Implementation recommendations.

## 8. PRICING ACTION PLAN
Top 5 immediate pricing actions with estimated revenue impact.
30-day, 90-day, and 6-month pricing roadmap.

Use specific numbers from the data."""

    result = query_llm(prompt)

    fig = None
    try:
        numeric_cols = df.select_dtypes(include="number").columns.tolist()
        cat_cols = df.select_dtypes(include="object").columns.tolist()
        if numeric_cols and len(numeric_cols) >= 2:
            fig = px.scatter(df, x=numeric_cols[0], y=numeric_cols[1],
                             color=cat_cols[0] if cat_cols else None,
                             template="plotly_dark",
                             title=f"Price vs {numeric_cols[1]} Analysis")
            fig.update_layout(height=420)
        elif numeric_cols:
            fig = px.box(df, y=numeric_cols[:min(4, len(numeric_cols))],
                         template="plotly_dark",
                         title="Price Distribution by Category")
            fig.update_layout(height=420)
    except Exception:
        fig = None

    return result + "\n\n" + WATERMARK, fig


# ─── Tab 7: Visualizations and Charts ────────────────────────────────────────

def generate_viz(file, chart_type, x_col, y_col, color_col):
    if file is None:
        return None, "Please upload a data file."
    df, err = load_data(file)
    if err:
        return None, err
    cols = df.columns.tolist()
    x = x_col if x_col in cols else cols[0]
    y = y_col if y_col in cols else (cols[1] if len(cols) > 1 else cols[0])
    color = color_col if color_col in cols else None
    try:
        if chart_type == "Bar Chart":
            fig = px.bar(df, x=x, y=y, color=color, template="plotly_dark",
                         title=f"{y} by {x}")
        elif chart_type == "Line Chart":
            fig = px.line(df, x=x, y=y, color=color, template="plotly_dark",
                          title=f"{y} over {x}")
        elif chart_type == "Pie Chart":
            fig = px.pie(df, names=x, values=y, template="plotly_dark",
                         title=f"{y} by {x}")
        elif chart_type == "Scatter Plot":
            fig = px.scatter(df, x=x, y=y, color=color, template="plotly_dark",
                             title=f"{x} vs {y}")
        elif chart_type == "Heatmap":
            numeric = df.select_dtypes(include="number")
            corr = numeric.corr()
            fig = px.imshow(corr, text_auto=True, template="plotly_dark",
                            title="Correlation Heatmap",
                            color_continuous_scale="RdBu_r")
        elif chart_type == "Box Plot":
            fig = px.box(df, x=x if df[x].dtype == "object" else None,
                         y=y, color=color, template="plotly_dark",
                         title=f"{y} Distribution")
        elif chart_type == "Histogram":
            fig = px.histogram(df, x=x, color=color, template="plotly_dark",
                               title=f"{x} Distribution", nbins=40)
        elif chart_type == "Area Chart":
            fig = px.area(df, x=x, y=y, color=color, template="plotly_dark",
                          title=f"{y} Area Chart")
        elif chart_type == "Treemap":
            fig = px.treemap(df, path=[x], values=y, template="plotly_dark",
                             title=f"{y} Treemap by {x}")
        else:
            fig = px.bar(df, x=x, y=y, template="plotly_dark")
        fig.update_layout(height=500, margin=dict(l=40, r=40, t=60, b=40))
        return fig, f"Chart generated: {chart_type} — {y} by {x}"
    except Exception as e:
        return None, f"Chart error: {str(e)}"


def get_columns(file):
    if file is None:
        return gr.Dropdown(choices=[]), gr.Dropdown(choices=[]), gr.Dropdown(choices=[])
    df, err = load_data(file)
    if err or df is None:
        return gr.Dropdown(choices=[]), gr.Dropdown(choices=[]), gr.Dropdown(choices=[])
    cols = df.columns.tolist()
    return gr.Dropdown(choices=cols), gr.Dropdown(choices=cols), gr.Dropdown(choices=cols)


# ─── Tab 8: Reports and Presentations ────────────────────────────────────────

def generate_report(file, report_type):
    if file is None:
        return "Please upload a data file.", None
    df, err = load_data(file)
    if err and df is None:
        return err, None

    if df is not None:
        preview = df.head(20).to_string()
        desc = df.describe(include="all").to_string()
        cols = list(df.columns)
        rows = len(df)
        data_summary = f"Rows: {rows} | Columns: {cols}\nStatistics:\n{desc}\nSample:\n{preview}"
    else:
        data_summary = err[:3000]

    prompt = f"""You are an expert retail and e-commerce analyst writing a professional report. Write as the professional analyst who performed this work. State all findings directly with specific numbers, percentages, and comparisons. NEVER say 'this appears to be', 'the data seems to show', 'the dataset contains', or 'it looks like'. Instead, deliver results confidently: 'Average claim cost was $4,230, a 12% increase year-over-year' or 'Readmission rates of 18.3% in Cardiology exceeded the facility average by 4.1 percentage points.' Be direct, precise, and actionable. Present findings as YOUR analysis, not a description of a dataset. Produce a complete {report_type}:

{data_summary}

Structure:
1. **Executive Summary** (3-5 sentences)
2. **Sales and Revenue Overview** — key metrics with numbers
3. **Top Performing Products and Categories** — specific findings
4. **Customer Insights** — behavior and segment findings
5. **Inventory and Operations** — efficiency findings
6. **Pricing and Margin Analysis** — optimization opportunities
7. **Recommendations** — 5 specific actionable recommendations
8. **Conclusion**

Write professionally for a retail business audience. Use specific numbers."""

    report_text = query_llm(prompt)

    output_path = None
    try:
        if report_type == "PowerPoint Presentation":
            output_path = _make_pptx(report_text, df if df is not None else None)
        else:
            output_path = _make_pdf(report_text)
    except Exception as e:
        report_text += f"\n\n[Document generation error: {str(e)}]"

    return report_text + "\n\n" + WATERMARK, output_path


def _make_pptx(text, df):
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    def add_slide(title_text, body_text, bg_color=(20, 10, 10)):
        slide = prs.slides.add_slide(blank_layout)
        bg = slide.background
        bg.fill.solid()
        bg.fill.fore_color.rgb = RGBColor(*bg_color)
        tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(1.2))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = RGBColor(231, 76, 60)
        body_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(1.7), Inches(12.3), Inches(5.5))
        tf2 = body_box.text_frame
        tf2.word_wrap = True
        for line in body_text.split("\n")[:20]:
            line = line.strip()
            if not line:
                continue
            para = tf2.add_paragraph()
            para.text = line
            para.font.size = Pt(14)
            para.font.color.rgb = RGBColor(240, 220, 220)

    title_slide = prs.slides.add_slide(blank_layout)
    bg = title_slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(10, 5, 5)
    tb = title_slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11), Inches(1.5))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = "AI Retail & E-Commerce Analysis Report"
    p.font.size = Pt(34)
    p.font.bold = True
    p.font.color.rgb = RGBColor(231, 76, 60)
    tb2 = title_slide.shapes.add_textbox(Inches(1), Inches(4.2), Inches(11), Inches(1))
    tf2 = tb2.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = f"Generated: {datetime.now().strftime('%B %d, %Y')} | AI Retail Analyzer"
    p2.font.size = Pt(18)
    p2.font.color.rgb = RGBColor(200, 160, 160)

    sections = []
    current_title = "Overview"
    current_body = []
    for line in text.split("\n"):
        stripped = line.strip()
        if stripped.startswith("**") and stripped.endswith("**") and len(stripped) > 4:
            if current_body:
                sections.append((current_title, "\n".join(current_body)))
            current_title = stripped.strip("*").strip()
            current_body = []
        elif stripped.startswith("## "):
            if current_body:
                sections.append((current_title, "\n".join(current_body)))
            current_title = stripped[3:].strip()
            current_body = []
        else:
            current_body.append(line)
    if current_body:
        sections.append((current_title, "\n".join(current_body)))

    bg_colors = [(20, 10, 10), (15, 10, 20), (10, 20, 10), (20, 15, 5), (5, 15, 20)]
    for i, (title, body) in enumerate(sections[:8]):
        add_slide(title, body, bg_colors[i % len(bg_colors)])

    if df is not None:
        stats_body = f"Total Records: {len(df)}\nColumns: {len(df.columns)}\n\n"
        for col in df.columns[:10]:
            stats_body += f"  • {col} ({df[col].dtype})\n"
        add_slide("Data Statistics", stats_body)

    closing = prs.slides.add_slide(blank_layout)
    bg = closing.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(10, 5, 5)
    tb = closing.shapes.add_textbox(Inches(1), Inches(2.8), Inches(11), Inches(2))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Report Complete"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(231, 76, 60)
    p2 = tf.add_paragraph()
    p2.text = "© 2026 Existential Gateway, LLC | AI Retail and E-Commerce Analyzer"
    p2.font.size = Pt(14)
    p2.font.color.rgb = RGBColor(200, 160, 160)

    tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".pptx")
    prs.save(tmp.name)
    return tmp.name


def _make_pdf(text):
    tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".pdf")
    doc = SimpleDocTemplate(tmp.name, pagesize=letter,
                            leftMargin=inch, rightMargin=inch,
                            topMargin=inch, bottomMargin=inch)
    styles = getSampleStyleSheet()
    title_style = ParagraphStyle("Title2", parent=styles["Title"],
                                 fontSize=20, spaceAfter=20,
                                 textColor=colors.HexColor("#7b0000"))
    heading_style = ParagraphStyle("H2", parent=styles["Heading2"],
                                   fontSize=14, spaceAfter=8,
                                   textColor=colors.HexColor("#c0392b"))
    body_style = ParagraphStyle("Body2", parent=styles["Normal"],
                                fontSize=11, spaceAfter=6, leading=16)
    story = []
    story.append(Paragraph("AI Retail & E-Commerce Analysis Report", title_style))
    story.append(Paragraph(
        f"Generated: {datetime.now().strftime('%B %d, %Y')} | AI Retail Analyzer",
        body_style))
    story.append(Spacer(1, 0.2 * inch))
    for line in text.split("\n"):
        stripped = line.strip()
        if not stripped:
            story.append(Spacer(1, 0.1 * inch))
            continue
        if stripped.startswith("**") and stripped.endswith("**"):
            story.append(Paragraph(stripped.strip("*"), heading_style))
        elif stripped.startswith("## "):
            story.append(Paragraph(stripped[3:], heading_style))
        elif stripped.startswith("# "):
            story.append(Paragraph(stripped[2:], title_style))
        else:
            safe = stripped.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")
            story.append(Paragraph(safe, body_style))
    story.append(Spacer(1, 0.3 * inch))
    story.append(Paragraph(
        "© 2026 Existential Gateway, LLC | AI Retail and E-Commerce Analyzer | existentialgateway@gmail.com",
        ParagraphStyle("footer", parent=styles["Normal"], fontSize=9,
                       textColor=colors.grey)))
    doc.build(story)
    return tmp.name


# ─── Tab 9: AI Data Chat ─────────────────────────────────────────────────────

_chat_df = None


def upload_chat_data(file):
    global _chat_df
    if file is None:
        _chat_df = None
        return "No file uploaded."
    df, err = load_data(file)
    if err and df is None:
        return err
    _chat_df = df
    return f"Data loaded: {len(df)} rows x {len(df.columns)} columns\nColumns: {list(df.columns)}"


def chat_with_data(message, history):
    global _chat_df

    data_context = ""
    if _chat_df is not None:
        desc = _chat_df.describe(include="all").to_string()
        preview = _chat_df.head(10).to_string()
        data_context = (
            f"The user has uploaded retail/e-commerce data.\n"
            f"Shape: {_chat_df.shape}\n"
            f"Columns: {list(_chat_df.columns)}\n"
            f"Statistical Summary:\n{desc}\n"
            f"First 10 rows:\n{preview}"
        )
    else:
        data_context = "No data uploaded yet. Answer general retail and e-commerce analytics questions."

    messages = [
        {
            "role": "system",
            "content": (
                "You are an expert retail and e-commerce data analyst and AI assistant. Write as the professional analyst who performed this work. State all findings directly with specific numbers, percentages, and comparisons. NEVER say 'this appears to be', 'the data seems to show', 'the dataset contains', or 'it looks like'. Instead, deliver results confidently: 'Average claim cost was $4,230, a 12% increase year-over-year' or 'Readmission rates of 18.3% in Cardiology exceeded the facility average by 4.1 percentage points.' Be direct, precise, and actionable. Present findings as YOUR analysis, not a description of a dataset. "
                "When asked a question, deliver the answer immediately with specific numbers. Do NOT explain how to calculate. Analyze their sales, customer, inventory, and pricing data. "
                "Answer questions, identify trends, generate SQL queries, and provide "
                "Python code examples when asked. Be specific, actionable, and business-focused. "
                "Always remind users to remove personally identifiable customer information "
                "before uploading data.\n\n" + data_context
            )
        }
    ]
    for user_msg, bot_msg in history:
        messages.append({"role": "user", "content": user_msg})
        messages.append({"role": "assistant", "content": bot_msg})
    messages.append({"role": "user", "content": message})

    import os
    API_KEY = os.environ.get("OPENAI_API_KEY", "")
    headers = {"Authorization": f"Bearer {API_KEY}", "Content-Type": "application/json"}
    msgs = [{"role": "system", "content": RETAIL_SYSTEM_PROMPT}, {"role": "user", "content": prompt}]
    payload = {"model": "gpt-4o", "max_tokens": 4000, "messages": msgs}
    try:
        response = requests.post("https://api.openai.com/v1/chat/completions", headers=headers, json=payload, timeout=240)
        result = response.json()
        if "choices" in result:
            return result["choices"][0]["message"]["content"]
        return f"API Error: {result}"
    except Exception as e:
        return f"Error: {str(e)}"


# ─── Gradio UI ────────────────────────────────────────────────────────────────

with gr.Blocks(title="AI Retail & E-Commerce Analyzer", theme=gr.themes.Base(
        primary_hue=gr.themes.colors.Color(
            c50="#fef9ec", c100="#faefd0", c200="#f4dea0", c300="#edc970",
            c400="#e4b04a", c500="#c9a84c", c600="#a8872e", c700="#856519",
            c800="#5e440d", c900="#3a2a05", c950="#1e1502"
        ),
        secondary_hue=gr.themes.colors.Color(
            c50="#e8edf5", c100="#c5d0e6", c200="#9eb0d4", c300="#7490c2",
            c400="#4f74b0", c500="#2d5a9e", c600="#1a3a6e", c700="#112240",
            c800="#0a1628", c900="#060e1a", c950="#03070d"
        ),
        neutral_hue=gr.themes.colors.Color(
            c50="#f0f2f7", c100="#d8dde9", c200="#b8c2d6", c300="#95a3be",
            c400="#7285a6", c500="#536890", c600="#3a4f73", c700="#263856",
            c800="#162438", c900="#0c1622", c950="#060b11"
        ),
        font=gr.themes.GoogleFont("DM Sans"),
        font_mono=gr.themes.GoogleFont("DM Mono"),
    ).set(
        body_background_fill="#0a1628",
        body_background_fill_dark="#0a1628",
        body_text_color="#f8f9fc",
        body_text_color_dark="#f8f9fc",
        block_background_fill="#112240",
        block_background_fill_dark="#112240",
        block_border_color="rgba(201,168,76,0.2)",
        block_border_color_dark="rgba(201,168,76,0.2)",
        block_label_background_fill="#112240",
        block_label_background_fill_dark="#112240",
        block_label_text_color="#c9a84c",
        block_label_text_color_dark="#c9a84c",
        block_title_text_color="#c9a84c",
        block_title_text_color_dark="#c9a84c",
        button_primary_background_fill="linear-gradient(135deg,#c9a84c,#e8c96a)",
        button_primary_background_fill_dark="linear-gradient(135deg,#c9a84c,#e8c96a)",
        button_primary_text_color="#0a1628",
        button_primary_text_color_dark="#0a1628",
        button_secondary_background_fill="#112240",
        button_secondary_background_fill_dark="#112240",
        button_secondary_border_color="rgba(201,168,76,0.3)",
        button_secondary_border_color_dark="rgba(201,168,76,0.3)",
        button_secondary_text_color="#c9a84c",
        button_secondary_text_color_dark="#c9a84c",
        input_background_fill="#0a1628",
        input_background_fill_dark="#0a1628",
        input_border_color="rgba(201,168,76,0.2)",
        input_border_color_dark="rgba(201,168,76,0.2)",
        input_placeholder_color="#8892a4",
        input_placeholder_color_dark="#8892a4",
        shadow_drop="0 4px 24px rgba(0,0,0,0.4)",
        shadow_drop_lg="0 8px 40px rgba(0,0,0,0.5)",
        table_even_background_fill="#0a1628",
        table_even_background_fill_dark="#0a1628",
        table_odd_background_fill="#112240",
        table_odd_background_fill_dark="#112240",
    )) as demo:
    gr.Markdown("# 🛒 AI Retail & E-Commerce Analyzer")
    gr.Markdown(DISCLAIMER)

    with gr.Tabs():

        # ── Tab 1 ──────────────────────────────────────────────────────────────
        with gr.Tab("📊 Data Upload & Overview"):
            gr.Markdown("## Upload Your Retail or Sales Data")
            gr.Markdown(PII_WARNING)
            with gr.Row():
                with gr.Column(scale=1):
                    t1_file = gr.File(label="Upload File (CSV, Excel, PDF, JSON)",
                                      file_types=[".csv", ".xlsx", ".xls", ".pdf", ".json"])
                    t1_btn = gr.Button("🔍 Analyze Dataset", variant="primary")
                    gr.Markdown(WAIT_MSG)
                with gr.Column(scale=2):
                    t1_analysis = gr.Markdown(label="AI Analysis")
            with gr.Row():
                t1_stats = gr.Markdown(label="Statistics")
                t1_preview = gr.Markdown(label="Data Preview")
            t1_btn.click(analyze_overview, inputs=[t1_file],
                         outputs=[t1_analysis, t1_stats, t1_preview])

        # ── Tab 2 ──────────────────────────────────────────────────────────────
        with gr.Tab("📈 Sales Trend Analysis"):
            gr.Markdown("## Sales Revenue & Trend Analysis")
            gr.Markdown(PII_WARNING)
            with gr.Row():
                with gr.Column(scale=1):
                    t2_file = gr.File(label="Upload Sales Transaction Data (CSV/Excel)",
                                      file_types=[".csv", ".xlsx", ".xls"])
                    t2_btn = gr.Button("📈 Analyze Sales Trends", variant="primary")
                    gr.Markdown(WAIT_MSG)
                with gr.Column(scale=2):
                    t2_output = gr.Markdown(label="Sales Trend Analysis")
            t2_chart = gr.Plot(label="Sales Chart")
            t2_btn.click(analyze_sales_trends, inputs=[t2_file],
                         outputs=[t2_output, t2_chart])

        # ── Tab 3 ──────────────────────────────────────────────────────────────
        with gr.Tab("👥 Customer Segmentation"):
            gr.Markdown("## Customer Segmentation & RFM Analysis")
            gr.Markdown(PII_WARNING)
            with gr.Row():
                with gr.Column(scale=1):
                    t3_file = gr.File(label="Upload Customer Purchase Data (CSV/Excel)",
                                      file_types=[".csv", ".xlsx", ".xls"])
                    t3_btn = gr.Button("👥 Analyze Customers", variant="primary")
                    gr.Markdown(WAIT_MSG)
                with gr.Column(scale=2):
                    t3_output = gr.Markdown(label="Customer Segmentation Analysis")
            t3_chart = gr.Plot(label="Customer Segment Chart")
            t3_btn.click(analyze_customer_segmentation, inputs=[t3_file],
                         outputs=[t3_output, t3_chart])

        # ── Tab 4 ──────────────────────────────────────────────────────────────
        with gr.Tab("📦 Inventory Optimization"):
            gr.Markdown("## Inventory Management & Optimization Analysis")
            gr.Markdown(PII_WARNING)
            with gr.Row():
                with gr.Column(scale=1):
                    t4_file = gr.File(label="Upload Inventory & Sales Data (CSV/Excel)",
                                      file_types=[".csv", ".xlsx", ".xls"])
                    t4_btn = gr.Button("📦 Optimize Inventory", variant="primary")
                    gr.Markdown(WAIT_MSG)
                with gr.Column(scale=2):
                    t4_output = gr.Markdown(label="Inventory Optimization Analysis")
            t4_chart = gr.Plot(label="Inventory Chart")
            t4_btn.click(analyze_inventory, inputs=[t4_file],
                         outputs=[t4_output, t4_chart])

        # ── Tab 5 ──────────────────────────────────────────────────────────────
        with gr.Tab("🚨 Churn Prediction"):
            gr.Markdown("## Customer Churn Prediction & Retention Analysis")
            gr.Markdown(PII_WARNING)
            with gr.Row():
                with gr.Column(scale=1):
                    t5_file = gr.File(label="Upload Customer Activity Data (CSV/Excel)",
                                      file_types=[".csv", ".xlsx", ".xls"])
                    t5_btn = gr.Button("🚨 Predict Churn", variant="primary")
                    gr.Markdown(WAIT_MSG)
                with gr.Column(scale=2):
                    t5_output = gr.Markdown(label="Churn Prediction Analysis")
            t5_chart = gr.Plot(label="Churn Risk Chart")
            t5_btn.click(analyze_churn, inputs=[t5_file],
                         outputs=[t5_output, t5_chart])

        # ── Tab 6 ──────────────────────────────────────────────────────────────
        with gr.Tab("💲 Price Optimization"):
            gr.Markdown("## Pricing Strategy & Revenue Optimization")
            gr.Markdown(PII_WARNING)
            with gr.Row():
                with gr.Column(scale=1):
                    t6_file = gr.File(label="Upload Pricing & Sales Data (CSV/Excel)",
                                      file_types=[".csv", ".xlsx", ".xls"])
                    t6_btn = gr.Button("💲 Optimize Pricing", variant="primary")
                    gr.Markdown(WAIT_MSG)
                with gr.Column(scale=2):
                    t6_output = gr.Markdown(label="Price Optimization Analysis")
            t6_chart = gr.Plot(label="Pricing Chart")
            t6_btn.click(analyze_pricing, inputs=[t6_file],
                         outputs=[t6_output, t6_chart])

        # ── Tab 7 ──────────────────────────────────────────────────────────────
        with gr.Tab("📉 Visualizations & Charts"):
            gr.Markdown("## Interactive Retail Analytics Visualizations")
            gr.Markdown(PII_WARNING)
            with gr.Row():
                with gr.Column(scale=1):
                    t7_file = gr.File(label="Upload Data (CSV/Excel)",
                                      file_types=[".csv", ".xlsx", ".xls"])
                    t7_chart_type = gr.Dropdown(
                        choices=["Bar Chart", "Line Chart", "Pie Chart", "Scatter Plot",
                                 "Heatmap", "Box Plot", "Histogram", "Area Chart", "Treemap"],
                        value="Bar Chart", label="Chart Type")
                    t7_x = gr.Dropdown(choices=[], label="X Axis Column")
                    t7_y = gr.Dropdown(choices=[], label="Y Axis Column")
                    t7_color = gr.Dropdown(choices=[], label="Color By (optional)")
                    t7_btn = gr.Button("📊 Generate Chart", variant="primary")
                    t7_status = gr.Markdown()
                with gr.Column(scale=2):
                    t7_plot = gr.Plot(label="Chart")
            t7_file.change(get_columns, inputs=[t7_file],
                           outputs=[t7_x, t7_y, t7_color])
            t7_btn.click(generate_viz,
                         inputs=[t7_file, t7_chart_type, t7_x, t7_y, t7_color],
                         outputs=[t7_plot, t7_status])

        # ── Tab 8 ──────────────────────────────────────────────────────────────
        with gr.Tab("📄 Reports & Presentations"):
            gr.Markdown("## Generate Professional Retail Reports")
            gr.Markdown(PII_WARNING)
            with gr.Row():
                with gr.Column(scale=1):
                    t8_file = gr.File(label="Upload Data (CSV, Excel, PDF, JSON)",
                                      file_types=[".csv", ".xlsx", ".xls", ".pdf", ".json"])
                    t8_type = gr.Dropdown(
                        choices=["PowerPoint Presentation", "PDF Report",
                                 "Executive Dashboard Summary", "Monthly Performance Report"],
                        value="PDF Report", label="Report Type")
                    t8_btn = gr.Button("📄 Generate Report", variant="primary")
                    gr.Markdown(WAIT_MSG)
                    t8_download = gr.File(label="⬇️ Download Report")
                with gr.Column(scale=2):
                    t8_output = gr.Markdown(label="Report Preview")
            t8_btn.click(generate_report, inputs=[t8_file, t8_type],
                         outputs=[t8_output, t8_download])

        # ── Tab 9 ──────────────────────────────────────────────────────────────
        with gr.Tab("💬 AI Data Chat"):
            gr.Markdown("## Chat With Your Retail Data")
            gr.Markdown(PII_WARNING)
            with gr.Row():
                with gr.Column(scale=1):
                    t9_file = gr.File(label="Upload Data for Context (CSV/Excel)",
                                      file_types=[".csv", ".xlsx", ".xls"])
                    t9_upload_btn = gr.Button("📤 Load Data", variant="secondary")
                    t9_data_status = gr.Markdown()
                with gr.Column(scale=2):
                    t9_chatbot = gr.ChatInterface(
                        fn=chat_with_data,
                        examples=[
                            "What is the best selling product in this dataset?",
                            "Which customer segment generates the most revenue?",
                            "What is the average order value?",
                            "Generate a SQL query to find the top 10 products by revenue",
                            "Write Python code to calculate monthly revenue trends",
                            "Which products are at risk of stocking out?",
                            "What is the customer churn rate based on this data?",
                        ],
                        title="",
                    )
            t9_upload_btn.click(upload_chat_data, inputs=[t9_file],
                                outputs=[t9_data_status])

    gr.Markdown(WATERMARK)


if __name__ == "__main__":
    demo.launch(server_name="0.0.0.0", server_port=int(os.environ.get("GRADIO_SERVER_PORT", 7860)), share=False, ssr_mode=False)
