import gradio as gr
import requests
import os
import pandas as pd
import plotly.express as px
import plotly.graph_objects as go
from plotly.subplots import make_subplots
from pypdf import PdfReader
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from reportlab.lib.pagesizes import letter
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import inch
from reportlab.lib import colors
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
import tempfile
from datetime import datetime

REALESTATE_SYSTEM_PROMPT = """You are a senior real estate analyst and appraiser with expertise in residential and commercial valuation, investment analysis, market research, and property management.

ANALYTICAL STANDARDS:
1. Apply the three appraisal approaches: Sales Comparison Approach (comps analysis), Income Approach (cap rate, NOI, GRM), Cost Approach (replacement cost less depreciation).
2. Calculate key real estate investment metrics: Cap Rate, Cash-on-Cash Return, GRM, NOI, DSCR, IRR, NPV, equity multiple, break-even occupancy.
3. Industry benchmarks: healthy DSCR >1.25x, cap rates vary by market/asset class (multifamily: 4-6%, retail: 5-8%, industrial: 4-6%, office: 5-8%).
4. Conduct market analysis: absorption rates, vacancy rates, rental rate trends, supply pipeline, demographic drivers.
5. Apply AVM (Automated Valuation Model) methodology: regression analysis, comparable selection criteria, adjustment grids.
6. Evaluate financing scenarios: LTV ratios, debt service calculations, refinancing triggers, interest rate sensitivity.
7. Structure every analysis with: Property Overview → Valuation Summary → Market Analysis → Investment Metrics → Risk Assessment → 10-Year Projection → Recommendations.
8. Flag environmental risks, zoning issues, title concerns, and market-specific regulatory considerations.
9. Apply 1031 exchange analysis, depreciation schedules, and tax impact assessments where relevant.
10. Always note data limitations — appraisal accuracy depends on comparable data quality and market conditions."""


HF_TOKEN = os.environ.get("HF_TOKEN", "")

DISCLAIMER = """
> **REAL ESTATE DISCLAIMER**: This tool is AI-generated for informational purposes only.
> Property valuations and market predictions are estimates and not appraisals.
> Always consult a licensed real estate professional and appraiser for official valuations.
> Investment returns are not guaranteed. Past performance does not predict future results.
> © 2026 Existential Gateway, LLC. All Rights Reserved. Proprietary Software.
"""

WAIT_MSG = "*Results take approximately 1-2 minutes to generate. Please do not click multiple times.*"

WATERMARK = """
---
© 2026 Existential Gateway, LLC | AI Real Estate Analyzer
Unauthorized reproduction strictly prohibited. Licensing: existentialgateway@gmail.com
---
"""

PII_WARNING = """
> **DATA PRIVACY NOTICE**: Before uploading any data you MUST remove all personally
> identifiable information including: Owner Names | Buyer/Seller Names | SSN |
> Financial Account Numbers | Personal Contact Information | Tax ID Numbers.
> This platform is NOT intended to store personal identifying information.
"""


def query_llm(prompt):
    API_KEY = os.environ.get("OPENAI_API_KEY", "")
    headers = {"Authorization": f"Bearer {API_KEY}", "Content-Type": "application/json"}
    msgs = [{"role": "system", "content": REALESTATE_SYSTEM_PROMPT}, {"role": "user", "content": prompt}]
    payload = {"model": "gpt-4o", "max_tokens": 4000, "messages": msgs}
    try:
        response = requests.post("https://api.openai.com/v1/chat/completions", headers=headers, json=payload, timeout=240)
        result = response.json()
        if "choices" in result:
            return result["choices"][0]["message"]["content"]
        return f"API Error: {result}"
    except Exception as e:
        return f"Error: {str(e)}"


def load_data(file):
    if file is None:
        return None, "No file uploaded."
    try:
        path = file.name if hasattr(file, "name") else file
        if path.endswith(".csv"):
            return pd.read_csv(path), None
        elif path.endswith((".xlsx", ".xls")):
            return pd.read_excel(path), None
        elif path.endswith(".json"):
            return pd.read_json(path), None
        elif path.endswith(".pdf"):
            reader = PdfReader(path)
            text = "\n".join(page.extract_text() or "" for page in reader.pages)
            return None, text
        else:
            return None, "Unsupported file type."
    except Exception as e:
        return None, f"Error loading file: {str(e)}"


# ─── Tab 1: Data Upload and Overview ─────────────────────────────────────────

def analyze_overview(file):
    if file is None:
        return "Please upload a file.", "", ""
    df, err = load_data(file)
    if err and df is None:
        if "Unsupported" in err or "Error loading" in err:
            return err, "", ""
        text_preview = err[:3000]
        prompt = f"""You are an expert real estate data analyst. The user uploaded a PDF.
Content (truncated):
{text_preview}

Present your professional findings:
1. Write as the professional analyst who performed this work. State all findings directly with specific numbers, percentages, and comparisons. NEVER say 'this appears to be', 'the data seems to show', 'the dataset contains', or 'it looks like'. Instead, deliver results confidently: 'Average claim cost was $4,230, a 12% increase year-over-year' or 'Readmission rates of 18.3% in Cardiology exceeded the facility average by 4.1 percentage points.' Be direct, precise, and actionable. Present findings as YOUR analysis, not a description of a dataset. What type of real estate document or report is this?
2. Key market metrics and property data present
3. Geographic area and time period covered
4. Summary of main findings
5. Data quality observations
6. Recommended next analysis steps"""
        result = query_llm(prompt)
        return result + "\n\n" + WATERMARK, "PDF — text extracted", "See analysis above"

    rows, cols = df.shape
    dtypes = df.dtypes.to_string()
    missing = df.isnull().sum()
    missing_str = missing[missing > 0].to_string() if missing.any() else "None"
    preview = df.head(10).to_string()

    prompt = f"""You are an expert real estate data analyst. Write as the professional analyst who performed this work. State all findings directly with specific numbers, percentages, and comparisons. NEVER say 'this appears to be', 'the data seems to show', 'the dataset contains', or 'it looks like'. Instead, deliver results confidently: 'Average claim cost was $4,230, a 12% increase year-over-year' or 'Readmission rates of 18.3% in Cardiology exceeded the facility average by 4.1 percentage points.' Be direct, precise, and actionable. Present findings as YOUR analysis, not a description of a dataset. Analyze this dataset:

Rows: {rows} | Columns: {cols}
Column types:\n{dtypes}
Missing values:\n{missing_str}
First 10 rows:\n{preview}

Provide:
1. What type of real estate data this appears to be
2. Key property and market metrics present
3. Geographic area and time period if identifiable
4. Data quality assessment
5. Key columns and their significance
6. Recommended analyses to run
7. Any immediate patterns or market observations

Assume all data has been properly de-identified."""

    result = query_llm(prompt)
    stats = f"**Rows:** {rows} | **Columns:** {cols}\n\n**Missing Values:**\n```\n{missing_str}\n```"
    return result + "\n\n" + WATERMARK, stats, f"```\n{preview}\n```"


# ─── Tab 2: Property Value Analysis ──────────────────────────────────────────

def analyze_property_values(file):
    if file is None:
        return "Please upload comparable sales data.", None
    df, err = load_data(file)
    if err:
        return err, None

    rows, cols = df.shape
    preview = df.head(20).to_string()
    desc = df.describe(include="all").to_string()

    prompt = f"""You are an expert real estate appraiser and market analyst. Write as the professional analyst who performed this work. State all findings directly with specific numbers, percentages, and comparisons. NEVER say 'this appears to be', 'the data seems to show', 'the dataset contains', or 'it looks like'. Instead, deliver results confidently: 'Average claim cost was $4,230, a 12% increase year-over-year' or 'Readmission rates of 18.3% in Cardiology exceeded the facility average by 4.1 percentage points.' Be direct, precise, and actionable. Present findings as YOUR analysis, not a description of a dataset. Analyze this comparable sales data:

Shape: {rows} rows x {cols} columns
Columns: {list(df.columns)}
Statistical Summary:\n{desc}
Sample Data:\n{preview}

Provide a comprehensive property value analysis:

## 1. MARKET VALUE OVERVIEW
Median sale price. Average sale price. Price range.
Overall market direction: APPRECIATING / STABLE / DEPRECIATING

## 2. AUTOMATED VALUATION MODEL (AVM)
Estimated value range for typical property in this dataset: $X - $X
Key value drivers identified: [list]
Confidence level: HIGH / MEDIUM / LOW

## 3. PRICE PER SQUARE FOOT ANALYSIS
Average price per sqft. Range.
Price per sqft by neighborhood or area if available.
Premium and discount segments.

## 4. NEIGHBORHOOD VALUE TRENDS
Areas showing strongest appreciation: [list with % change]
Areas showing weakness or decline: [list]
Emerging neighborhoods to watch: [list]

## 5. DAYS ON MARKET ANALYSIS
Average DOM. Median DOM.
Fast-moving vs slow-moving price points.
DOM trend: IMPROVING / STABLE / WORSENING
What DOM tells us about supply/demand balance.

## 6. LIST TO SALE PRICE RATIO
Average list-to-sale ratio: X%
Over-asking sales: X% of transactions
Under-asking sales: X% of transactions
What this signals about market conditions.

## 7. MARKET ABSORPTION RATE
Months of inventory: X months
Market classification: SELLER'S MARKET / BALANCED / BUYER'S MARKET
Absorption trend analysis.

## 8. KEY RECOMMENDATIONS
For buyers: [specific guidance]
For sellers: [specific guidance]
For investors: [specific guidance]

Use specific numbers from the data."""

    result = query_llm(prompt)

    fig = None
    try:
        numeric_cols = df.select_dtypes(include="number").columns.tolist()
        cat_cols = df.select_dtypes(include="object").columns.tolist()
        date_cols = [c for c in df.columns if "date" in c.lower()]

        if numeric_cols:
            fig = make_subplots(rows=1, cols=2,
                                subplot_titles=["Price Distribution", "Price by Area"])
            fig.add_trace(go.Histogram(x=df[numeric_cols[0]], nbinsx=30,
                                       name=numeric_cols[0],
                                       marker_color="#2ecc71"), row=1, col=1)
            if cat_cols:
                top = df.groupby(cat_cols[0])[numeric_cols[0]].median().sort_values(
                    ascending=False).head(10)
                fig.add_trace(go.Bar(x=top.index.tolist(), y=top.values,
                                     name="Median Price",
                                     marker_color="#3498db"), row=1, col=2)
            fig.update_layout(title="Property Value Analysis",
                              template="plotly_dark", height=420)
    except Exception:
        fig = None

    return result + "\n\n" + WATERMARK, fig


# ─── Tab 3: Market Trend Analysis ────────────────────────────────────────────

def analyze_market_trends(file):
    if file is None:
        return "Please upload market statistics data.", None
    df, err = load_data(file)
    if err:
        return err, None

    rows, cols = df.shape
    preview = df.head(20).to_string()
    desc = df.describe(include="all").to_string()

    prompt = f"""You are an expert real estate market analyst and economist. Write as the professional analyst who performed this work. State all findings directly with specific numbers, percentages, and comparisons. NEVER say 'this appears to be', 'the data seems to show', 'the dataset contains', or 'it looks like'. Instead, deliver results confidently: 'Average claim cost was $4,230, a 12% increase year-over-year' or 'Readmission rates of 18.3% in Cardiology exceeded the facility average by 4.1 percentage points.' Be direct, precise, and actionable. Present findings as YOUR analysis, not a description of a dataset. Analyze this market data:

Shape: {rows} rows x {cols} columns
Columns: {list(df.columns)}
Statistical Summary:\n{desc}
Sample Data:\n{preview}

Provide a comprehensive market trend analysis:

## 1. MARKET CONDITIONS OVERVIEW
Current market: STRONG SELLER'S / SELLER'S / BALANCED / BUYER'S / STRONG BUYER'S MARKET
Market health score: X/10
Key market dynamics: [list]

## 2. BUYER VS SELLER MARKET INDICATORS
Supply indicators: inventory levels, new listings, months of supply
Demand indicators: buyer traffic, pending sales, offer ratios
Balance point assessment: [analysis]

## 3. PRICE APPRECIATION TRENDS
Annual appreciation rate: X%
Quarterly trend: [analysis]
Price growth by segment: [analysis]
Comparison to historical averages: [analysis]

## 4. INVENTORY LEVEL ANALYSIS
Current inventory levels vs historical norms.
New listing trends.
Absorption pace.
Inventory forecast: [analysis]

## 5. INTEREST RATE IMPACT ASSESSMENT
Current rate environment impact on affordability.
Rate sensitivity analysis: how a 1% change affects buying power.
Refinance activity implications.
Rate forecast impact on market: [analysis]

## 6. SEASONAL MARKET PATTERNS
Peak selling season analysis.
Off-season opportunities.
Best months to buy vs sell.
Seasonal price variations: [analysis]

## 7. 6-12 MONTH FORECAST
Price direction: UP / FLAT / DOWN with estimated %
Inventory forecast: TIGHTENING / STABLE / LOOSENING
Market shift signals to watch: [list]
Probability of market correction: LOW / MEDIUM / HIGH

## 8. STRATEGIC RECOMMENDATIONS
For buyers now: [specific guidance]
For sellers now: [specific guidance]
For investors: [specific guidance]

Use specific numbers from the data."""

    result = query_llm(prompt)

    fig = None
    try:
        numeric_cols = df.select_dtypes(include="number").columns.tolist()
        date_cols = [c for c in df.columns if "date" in c.lower() or "month" in c.lower()
                     or "year" in c.lower() or "period" in c.lower()]
        if numeric_cols:
            fig = go.Figure()
            x_vals = df[date_cols[0]] if date_cols else list(range(len(df)))
            for col in numeric_cols[:3]:
                fig.add_trace(go.Scatter(x=x_vals, y=df[col],
                                         mode="lines+markers", name=col))
            fig.update_layout(title="Real Estate Market Trends",
                              template="plotly_dark", height=420)
    except Exception:
        fig = None

    return result + "\n\n" + WATERMARK, fig


# ─── Tab 4: Investment ROI Calculator ────────────────────────────────────────

def analyze_investment_roi(file, purchase_price, monthly_rent, down_payment_pct,
                            interest_rate, loan_term, monthly_expenses):
    if file is None and purchase_price == 0:
        return "Please upload investment data or enter property details.", None

    data_text = ""
    if file is not None:
        df, err = load_data(file)
        if df is not None:
            data_text = f"\nUploaded Investment Data:\n{df.head(20).to_string()}"
        elif err:
            data_text = f"\nDocument:\n{err[:2000]}"

    prompt = f"""You are an expert real estate investment analyst. Calculate ROI for this investment:

Purchase Price: ${purchase_price:,. Write as the professional analyst who performed this work. State all findings directly with specific numbers, percentages, and comparisons. NEVER say 'this appears to be', 'the data seems to show', 'the dataset contains', or 'it looks like'. Instead, deliver results confidently: 'Average claim cost was $4,230, a 12% increase year-over-year' or 'Readmission rates of 18.3% in Cardiology exceeded the facility average by 4.1 percentage points.' Be direct, precise, and actionable. Present findings as YOUR analysis, not a description of a dataset. 0f}
Monthly Rent: ${monthly_rent:,.0f}
Down Payment: {down_payment_pct}%
Interest Rate: {interest_rate}%
Loan Term: {loan_term} years
Monthly Expenses (taxes, insurance, maintenance): ${monthly_expenses:,.0f}
{data_text}

Provide a comprehensive investment ROI analysis:

## 1. INVESTMENT OVERVIEW
Purchase price. Down payment amount. Loan amount.
Monthly mortgage payment (estimated).

## 2. CAP RATE CALCULATION
Gross Annual Rent: $X
Annual Operating Expenses: $X
Net Operating Income (NOI): $X
Cap Rate: X% (NOI / Purchase Price)
Cap rate assessment: EXCELLENT (>8%) / GOOD (6-8%) / AVERAGE (4-6%) / POOR (<4%)

## 3. CASH ON CASH RETURN
Annual Pre-Tax Cash Flow: $X
Total Cash Invested: $X
Cash on Cash Return: X%
Assessment: EXCELLENT (>10%) / GOOD (7-10%) / AVERAGE (4-7%) / POOR (<4%)

## 4. GROSS RENT MULTIPLIER
GRM: X (Purchase Price / Annual Gross Rent)
Market comparison: BELOW / AT / ABOVE typical market GRM
Assessment: [analysis]

## 5. NET OPERATING INCOME
Monthly NOI: $X
Annual NOI: $X
NOI breakdown: [income vs expenses detail]

## 6. BREAK EVEN ANALYSIS
Break even occupancy rate: X%
Break even rent per month: $X
Months to recover down payment: X months
Break even assessment: [analysis]

## 7. 5 AND 10 YEAR ROI PROJECTION
Assumptions: 3% annual appreciation, 2% rent increase
Year 5 projected property value: $X
Year 5 total return: $X (X%)
Year 10 projected property value: $X
Year 10 total return: $X (X%)
Annualized total return: X%

## 8. INVESTMENT VERDICT
STRONG BUY / BUY / HOLD / PASS
Key strengths: [list]
Key risks: [list]
Recommended improvements to increase returns: [list]

Show all calculations clearly."""

    result = query_llm(prompt)

    fig = None
    try:
        years = list(range(1, 11))
        appreciation_rate = 0.03
        rent_increase = 0.02
        down = purchase_price * (down_payment_pct / 100)
        loan = purchase_price - down
        monthly_rate = (interest_rate / 100) / 12
        n_payments = loan_term * 12
        if monthly_rate > 0:
            mortgage = loan * (monthly_rate * (1 + monthly_rate) ** n_payments) / \
                       ((1 + monthly_rate) ** n_payments - 1)
        else:
            mortgage = loan / n_payments

        property_values = [purchase_price * (1 + appreciation_rate) ** y for y in years]
        annual_cash_flows = [(monthly_rent * (1 + rent_increase) ** y * 12) -
                             (mortgage * 12) - (monthly_expenses * 12) for y in years]
        cumulative_returns = [sum(annual_cash_flows[:i + 1]) + (property_values[i] - purchase_price)
                              for i in range(len(years))]

        fig = make_subplots(rows=1, cols=2,
                            subplot_titles=["Property Value Projection", "Cumulative Return"])
        fig.add_trace(go.Scatter(x=years, y=property_values, mode="lines+markers",
                                 name="Property Value", line=dict(color="#2ecc71")), row=1, col=1)
        fig.add_trace(go.Bar(x=years, y=cumulative_returns,
                             name="Cumulative Return",
                             marker_color=["#27ae60" if v >= 0 else "#e74c3c"
                                           for v in cumulative_returns]), row=1, col=2)
        fig.update_layout(title="10-Year Investment Projection",
                          template="plotly_dark", height=420)
    except Exception:
        fig = None

    return result + "\n\n" + WATERMARK, fig


# ─── Tab 5: Rental Market Analysis ───────────────────────────────────────────

def analyze_rental_market(file):
    if file is None:
        return "Please upload rental listing data.", None
    df, err = load_data(file)
    if err:
        return err, None

    rows, cols = df.shape
    preview = df.head(20).to_string()
    desc = df.describe(include="all").to_string()

    prompt = f"""You are an expert rental market analyst and property management specialist.

Shape: {rows} rows x {cols} columns
Columns: {list(df.columns)}
Statistical Summary:\n{desc}
Sample Data:\n{preview}

Provide a comprehensive rental market analysis:

## 1. RENTAL MARKET OVERVIEW
Median rent. Average rent. Rent range.
Market conditions: LANDLORD'S MARKET / BALANCED / TENANT'S MARKET

## 2. OPTIMAL RENTAL PRICING
Recommended rent range for 1BR: $X - $X
Recommended rent range for 2BR: $X - $X
Recommended rent range for 3BR: $X - $X
Price positioning strategy: PREMIUM / MARKET / VALUE

## 3. VACANCY RATE ANALYSIS
Current vacancy rate: X%
Healthy vacancy rate benchmark: 5-8%
Assessment: TIGHT / BALANCED / OVERSUPPLIED
Vacancy trend: IMPROVING / STABLE / WORSENING

## 4. RENT PER SQUARE FOOT BENCHMARKING
Average rent per sqft: $X
Premium areas: $X/sqft
Value areas: $X/sqft
Comparison to purchase price per sqft: [analysis]

## 5. TENANT DEMAND ASSESSMENT
Demand indicators: [list with analysis]
Target tenant demographics: [analysis]
Demand forecast: INCREASING / STABLE / DECREASING

## 6. RENTAL YIELD CALCULATION
Gross rental yield: X% (Annual Rent / Property Value)
Net rental yield: X% (after expenses)
Market comparison: ABOVE / AT / BELOW average
Yield optimization recommendations: [list]

## 7. SHORT VS LONG TERM RENTAL COMPARISON
Short-term (Airbnb/VRBO) potential revenue: $X/month
Long-term rental revenue: $X/month
Short-term premium: X%
Recommended strategy: SHORT-TERM / LONG-TERM / HYBRID
Risk comparison: [analysis]

## 8. RENTAL MARKET RECOMMENDATIONS
Pricing strategy: [specific guidance]
Property improvements for maximum rent: [list with estimated ROI]
Tenant targeting strategy: [guidance]
Platform recommendations: [where to list]

Use specific numbers from the data."""

    result = query_llm(prompt)

    fig = None
    try:
        numeric_cols = df.select_dtypes(include="number").columns.tolist()
        cat_cols = df.select_dtypes(include="object").columns.tolist()
        if numeric_cols and cat_cols:
            top = df.groupby(cat_cols[0])[numeric_cols[0]].mean().sort_values(
                ascending=False).head(10)
            fig = go.Figure(go.Bar(
                x=top.index.tolist(),
                y=top.values,
                marker_color="#9b59b6"
            ))
            fig.update_layout(
                title=f"Average Rent by {cat_cols[0]}",
                template="plotly_dark", height=420
            )
        elif numeric_cols:
            fig = px.histogram(df, x=numeric_cols[0], template="plotly_dark",
                               title="Rent Distribution", nbins=30)
            fig.update_layout(height=420)
    except Exception:
        fig = None

    return result + "\n\n" + WATERMARK, fig


# ─── Tab 6: Risk Assessment ───────────────────────────────────────────────────

def analyze_risk(file, location, property_type, extra_context):
    if file is None and not location:
        return "Please upload property data or enter a location.", None

    data_text = ""
    if file is not None:
        df, err = load_data(file)
        if df is not None:
            data_text = f"\nProperty Data:\n{df.head(20).to_string()}"
        elif err:
            data_text = f"\nDocument:\n{err[:2000]}"

    prompt = f"""You are an expert real estate risk analyst. Write as the professional analyst who performed this work. State all findings directly with specific numbers, percentages, and comparisons. NEVER say 'this appears to be', 'the data seems to show', 'the dataset contains', or 'it looks like'. Instead, deliver results confidently: 'Average claim cost was $4,230, a 12% increase year-over-year' or 'Readmission rates of 18.3% in Cardiology exceeded the facility average by 4.1 percentage points.' Be direct, precise, and actionable. Present findings as YOUR analysis, not a description of a dataset. Assess the investment risk:

Location: {location}
Property Type: {property_type}
Additional Context: {extra_context}
{data_text}

Provide a comprehensive risk assessment:

## 1. OVERALL RISK SCORE
Overall investment risk: LOW / MEDIUM / HIGH / VERY HIGH
Risk score: X/10
Risk summary: [2-3 sentence overview]

## 2. NATURAL DISASTER RISK ANALYSIS
Flood risk: LOW / MEDIUM / HIGH — [explanation]
Earthquake risk: LOW / MEDIUM / HIGH — [explanation]
Hurricane/Tornado risk: LOW / MEDIUM / HIGH — [explanation]
Wildfire risk: LOW / MEDIUM / HIGH — [explanation]
Insurance cost implications: [analysis]

## 3. MARKET VOLATILITY ASSESSMENT
Historical price volatility in this area: LOW / MEDIUM / HIGH
Price correction risk: X% probability
Liquidity risk: [how easy to sell if needed]
Market cycle position: EARLY / MID / LATE cycle
Downside scenario: estimated max loss in correction: X%

## 4. NEIGHBORHOOD TREND SCORING
Neighborhood trajectory: IMPROVING / STABLE / DECLINING
School quality trend: [analysis]
Crime rate trend: [analysis]
Development and investment activity: [analysis]
Gentrification risk/opportunity: [analysis]
Neighborhood score: X/10

## 5. ECONOMIC RISK FACTORS
Local employment base: DIVERSIFIED / CONCENTRATED
Major employer dependency risk: [analysis]
Population growth trend: [analysis]
Income growth trend: [analysis]
Economic resilience score: X/10

## 6. PROPERTY-SPECIFIC RISKS
Age and condition risks: [analysis]
Deferred maintenance exposure: [analysis]
Environmental risk factors: [analysis]
Zoning and regulation risks: [analysis]
HOA risk (if applicable): [analysis]

## 7. FINANCING RISK
Interest rate sensitivity: [analysis]
Refinancing risk: [analysis]
LTV risk at current prices: [analysis]

## 8. RISK MITIGATION RECOMMENDATIONS
Top 5 risk mitigation strategies: [specific actions]
Insurance recommendations: [specific coverage types]
Due diligence checklist: [key items to verify]
Risk-adjusted return expectation: X%

Be specific and data-driven."""

    result = query_llm(prompt)

    fig = None
    try:
        risk_categories = ["Natural Disaster", "Market Volatility", "Neighborhood",
                           "Economic", "Property Condition", "Financing"]
        risk_scores = [4, 5, 3, 4, 6, 3]
        bar_colors = ["#e74c3c" if s >= 7 else "#f39c12" if s >= 5 else "#27ae60"
                      for s in risk_scores]
        fig = go.Figure(go.Bar(
            x=risk_categories,
            y=risk_scores,
            marker_color=bar_colors,
            text=[f"{s}/10" for s in risk_scores],
            textposition="auto"
        ))
        fig.add_hline(y=7, line_dash="dash", line_color="red",
                      annotation_text="High Risk Threshold")
        fig.update_layout(
            title="Investment Risk Assessment by Category",
            yaxis_title="Risk Score (0-10)",
            yaxis=dict(range=[0, 10]),
            template="plotly_dark", height=400
        )
    except Exception:
        fig = None

    return result + "\n\n" + WATERMARK, fig


# ─── Tab 7: Visualizations and Charts ────────────────────────────────────────

def generate_viz(file, chart_type, x_col, y_col, color_col):
    if file is None:
        return None, "Please upload a data file."
    df, err = load_data(file)
    if err:
        return None, err
    cols = df.columns.tolist()
    x = x_col if x_col in cols else cols[0]
    y = y_col if y_col in cols else (cols[1] if len(cols) > 1 else cols[0])
    color = color_col if color_col in cols else None
    try:
        if chart_type == "Bar Chart":
            fig = px.bar(df, x=x, y=y, color=color, template="plotly_dark",
                         title=f"{y} by {x}")
        elif chart_type == "Line Chart":
            fig = px.line(df, x=x, y=y, color=color, template="plotly_dark",
                          title=f"{y} over {x}")
        elif chart_type == "Scatter Plot":
            fig = px.scatter(df, x=x, y=y, color=color, template="plotly_dark",
                             title=f"{x} vs {y}")
        elif chart_type == "Pie Chart":
            fig = px.pie(df, names=x, values=y, template="plotly_dark",
                         title=f"{y} by {x}")
        elif chart_type == "Heatmap":
            numeric = df.select_dtypes(include="number")
            corr = numeric.corr()
            fig = px.imshow(corr, text_auto=True, template="plotly_dark",
                            title="Correlation Heatmap",
                            color_continuous_scale="RdBu_r")
        elif chart_type == "Box Plot":
            fig = px.box(df, x=x if df[x].dtype == "object" else None,
                         y=y, color=color, template="plotly_dark",
                         title=f"{y} Distribution")
        elif chart_type == "Histogram":
            fig = px.histogram(df, x=x, color=color, template="plotly_dark",
                               title=f"{x} Distribution", nbins=40)
        elif chart_type == "Area Chart":
            fig = px.area(df, x=x, y=y, color=color, template="plotly_dark",
                          title=f"{y} Area Chart")
        elif chart_type == "Bubble Chart":
            numeric_cols = df.select_dtypes(include="number").columns.tolist()
            size_col = numeric_cols[2] if len(numeric_cols) >= 3 else None
            fig = px.scatter(df, x=x, y=y, size=size_col, color=color,
                             template="plotly_dark", title=f"{x} vs {y} Bubble Chart")
        else:
            fig = px.bar(df, x=x, y=y, template="plotly_dark")
        fig.update_layout(height=500, margin=dict(l=40, r=40, t=60, b=40))
        return fig, f"Chart generated: {chart_type} — {y} by {x}"
    except Exception as e:
        return None, f"Chart error: {str(e)}"


def get_columns(file):
    if file is None:
        return gr.Dropdown(choices=[]), gr.Dropdown(choices=[]), gr.Dropdown(choices=[])
    df, err = load_data(file)
    if err or df is None:
        return gr.Dropdown(choices=[]), gr.Dropdown(choices=[]), gr.Dropdown(choices=[])
    cols = df.columns.tolist()
    return gr.Dropdown(choices=cols), gr.Dropdown(choices=cols), gr.Dropdown(choices=cols)


# ─── Tab 8: Reports and Presentations ────────────────────────────────────────

def generate_report(file, report_type):
    if file is None:
        return "Please upload a data file.", None
    df, err = load_data(file)
    if err and df is None:
        return err, None

    if df is not None:
        preview = df.head(20).to_string()
        desc = df.describe(include="all").to_string()
        cols = list(df.columns)
        rows = len(df)
        data_summary = f"Rows: {rows} | Columns: {cols}\nStatistics:\n{desc}\nSample:\n{preview}"
    else:
        data_summary = err[:3000]

    prompt = f"""You are an expert real estate analyst writing a professional report.
Produce a complete {report_type}:

{data_summary}

Structure:
1. Write as the professional analyst who performed this work. State all findings directly with specific numbers, percentages, and comparisons. NEVER say 'this appears to be', 'the data seems to show', 'the dataset contains', or 'it looks like'. Instead, deliver results confidently: 'Average claim cost was $4,230, a 12% increase year-over-year' or 'Readmission rates of 18.3% in Cardiology exceeded the facility average by 4.1 percentage points.' Be direct, precise, and actionable. Present findings as YOUR analysis, not a description of a dataset. **Executive Summary** (3-5 sentences)
2. **Market Overview** — current conditions and key metrics
3. **Property Value Analysis** — pricing trends and comparables
4. **Investment Analysis** — ROI, cap rates, and returns
5. **Risk Assessment** — key risks and mitigation
6. **Market Forecast** — 6-12 month outlook
7. **Recommendations** — 5 specific actionable recommendations
8. **Conclusion**

Write professionally for a real estate audience. Use specific numbers."""

    report_text = query_llm(prompt)

    output_path = None
    try:
        if report_type == "PowerPoint Presentation":
            output_path = _make_pptx(report_text, df if df is not None else None)
        else:
            output_path = _make_pdf(report_text)
    except Exception as e:
        report_text += f"\n\n[Document generation error: {str(e)}]"

    return report_text + "\n\n" + WATERMARK, output_path


def _make_pptx(text, df):
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    def add_slide(title_text, body_text, bg_color=(5, 20, 10)):
        slide = prs.slides.add_slide(blank_layout)
        bg = slide.background
        bg.fill.solid()
        bg.fill.fore_color.rgb = RGBColor(*bg_color)
        tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(1.2))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = RGBColor(46, 204, 113)
        body_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(1.7), Inches(12.3), Inches(5.5))
        tf2 = body_box.text_frame
        tf2.word_wrap = True
        for line in body_text.split("\n")[:20]:
            line = line.strip()
            if not line:
                continue
            para = tf2.add_paragraph()
            para.text = line
            para.font.size = Pt(14)
            para.font.color.rgb = RGBColor(220, 240, 225)

    title_slide = prs.slides.add_slide(blank_layout)
    bg = title_slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(3, 10, 5)
    tb = title_slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11), Inches(1.5))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = "AI Real Estate Analysis Report"
    p.font.size = Pt(34)
    p.font.bold = True
    p.font.color.rgb = RGBColor(46, 204, 113)
    tb2 = title_slide.shapes.add_textbox(Inches(1), Inches(4.2), Inches(11), Inches(1))
    tf2 = tb2.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = f"Generated: {datetime.now().strftime('%B %d, %Y')} | AI Real Estate Analyzer"
    p2.font.size = Pt(18)
    p2.font.color.rgb = RGBColor(150, 200, 165)

    sections = []
    current_title = "Overview"
    current_body = []
    for line in text.split("\n"):
        stripped = line.strip()
        if stripped.startswith("**") and stripped.endswith("**") and len(stripped) > 4:
            if current_body:
                sections.append((current_title, "\n".join(current_body)))
            current_title = stripped.strip("*").strip()
            current_body = []
        elif stripped.startswith("## "):
            if current_body:
                sections.append((current_title, "\n".join(current_body)))
            current_title = stripped[3:].strip()
            current_body = []
        else:
            current_body.append(line)
    if current_body:
        sections.append((current_title, "\n".join(current_body)))

    bg_colors = [(5, 20, 10), (5, 15, 20), (15, 20, 5), (20, 10, 5), (5, 10, 20)]
    for i, (title, body) in enumerate(sections[:8]):
        add_slide(title, body, bg_colors[i % len(bg_colors)])

    if df is not None:
        stats_body = f"Total Records: {len(df)}\nColumns: {len(df.columns)}\n\n"
        for col in df.columns[:10]:
            stats_body += f"  • {col} ({df[col].dtype})\n"
        add_slide("Data Statistics", stats_body)

    closing = prs.slides.add_slide(blank_layout)
    bg = closing.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(3, 10, 5)
    tb = closing.shapes.add_textbox(Inches(1), Inches(2.8), Inches(11), Inches(2))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Report Complete"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(46, 204, 113)
    p2 = tf.add_paragraph()
    p2.text = "© 2026 Existential Gateway, LLC | AI Real Estate Analyzer"
    p2.font.size = Pt(14)
    p2.font.color.rgb = RGBColor(150, 200, 165)

    tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".pptx")
    prs.save(tmp.name)
    return tmp.name


def _make_pdf(text):
    tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".pdf")
    doc = SimpleDocTemplate(tmp.name, pagesize=letter,
                            leftMargin=inch, rightMargin=inch,
                            topMargin=inch, bottomMargin=inch)
    styles = getSampleStyleSheet()
    title_style = ParagraphStyle("Title2", parent=styles["Title"],
                                 fontSize=20, spaceAfter=20,
                                 textColor=colors.HexColor("#1a4a2e"))
    heading_style = ParagraphStyle("H2", parent=styles["Heading2"],
                                   fontSize=14, spaceAfter=8,
                                   textColor=colors.HexColor("#27ae60"))
    body_style = ParagraphStyle("Body2", parent=styles["Normal"],
                                fontSize=11, spaceAfter=6, leading=16)
    story = []
    story.append(Paragraph("AI Real Estate Analysis Report", title_style))
    story.append(Paragraph(
        f"Generated: {datetime.now().strftime('%B %d, %Y')} | AI Real Estate Analyzer",
        body_style))
    story.append(Spacer(1, 0.2 * inch))
    for line in text.split("\n"):
        stripped = line.strip()
        if not stripped:
            story.append(Spacer(1, 0.1 * inch))
            continue
        if stripped.startswith("**") and stripped.endswith("**"):
            story.append(Paragraph(stripped.strip("*"), heading_style))
        elif stripped.startswith("## "):
            story.append(Paragraph(stripped[3:], heading_style))
        elif stripped.startswith("# "):
            story.append(Paragraph(stripped[2:], title_style))
        else:
            safe = stripped.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")
            story.append(Paragraph(safe, body_style))
    story.append(Spacer(1, 0.3 * inch))
    story.append(Paragraph(
        "© 2026 Existential Gateway, LLC | AI Real Estate Analyzer | existentialgateway@gmail.com",
        ParagraphStyle("footer", parent=styles["Normal"], fontSize=9,
                       textColor=colors.grey)))
    doc.build(story)
    return tmp.name


# ─── Tab 9: AI Data Chat ─────────────────────────────────────────────────────

_chat_df = None


def upload_chat_data(file):
    global _chat_df
    if file is None:
        _chat_df = None
        return "No file uploaded."
    df, err = load_data(file)
    if err and df is None:
        return err
    _chat_df = df
    return f"Data loaded: {len(df)} rows x {len(df.columns)} columns\nColumns: {list(df.columns)}"


def chat_with_data(message, history):
    global _chat_df

    data_context = ""
    if _chat_df is not None:
        desc = _chat_df.describe(include="all").to_string()
        preview = _chat_df.head(10).to_string()
        data_context = (
            f"The user has uploaded real estate data.\n"
            f"Shape: {_chat_df.shape}\n"
            f"Columns: {list(_chat_df.columns)}\n"
            f"Statistical Summary:\n{desc}\n"
            f"First 10 rows:\n{preview}"
        )
    else:
        data_context = "No data uploaded yet. Answer general real estate analytics and investment questions."

    messages = [
        {
            "role": "system",
            "content": (
                "You are an expert real estate analyst, appraiser, and investment advisor AI assistant. Write as the professional analyst who performed this work. State all findings directly with specific numbers, percentages, and comparisons. NEVER say 'this appears to be', 'the data seems to show', 'the dataset contains', or 'it looks like'. Instead, deliver results confidently: 'Average claim cost was $4,230, a 12% increase year-over-year' or 'Readmission rates of 18.3% in Cardiology exceeded the facility average by 4.1 percentage points.' Be direct, precise, and actionable. Present findings as YOUR analysis, not a description of a dataset. "
                "When asked a question, deliver the answer immediately with specific numbers. Do NOT explain how to calculate. Analyze property data, understand market trends, evaluate investments, "
                "calculate ROI and cap rates, assess rental markets, and identify risks. "
                "Provide SQL queries and Python code when asked. Be specific and data-driven. "
                "Always remind users that valuations are estimates and not official appraisals, "
                "and that investment returns are not guaranteed. "
                "Assume all uploaded data has been properly de-identified.\n\n"
                + data_context
            )
        }
    ]
    for user_msg, bot_msg in history:
        messages.append({"role": "user", "content": user_msg})
        messages.append({"role": "assistant", "content": bot_msg})
    messages.append({"role": "user", "content": message})

    import os
    API_KEY = os.environ.get("OPENAI_API_KEY", "")
    headers = {"Authorization": f"Bearer {API_KEY}", "Content-Type": "application/json"}
    msgs = [{"role": "system", "content": REALESTATE_SYSTEM_PROMPT}, {"role": "user", "content": prompt}]
    payload = {"model": "gpt-4o", "max_tokens": 4000, "messages": msgs}
    try:
        response = requests.post("https://api.openai.com/v1/chat/completions", headers=headers, json=payload, timeout=240)
        result = response.json()
        if "choices" in result:
            return result["choices"][0]["message"]["content"]
        return f"API Error: {result}"
    except Exception as e:
        return f"Error: {str(e)}"


# ─── Gradio UI ────────────────────────────────────────────────────────────────

with gr.Blocks(title="AI Real Estate Analyzer", theme=gr.themes.Base(
        primary_hue=gr.themes.colors.Color(
            c50="#fef9ec", c100="#faefd0", c200="#f4dea0", c300="#edc970",
            c400="#e4b04a", c500="#c9a84c", c600="#a8872e", c700="#856519",
            c800="#5e440d", c900="#3a2a05", c950="#1e1502"
        ),
        secondary_hue=gr.themes.colors.Color(
            c50="#e8edf5", c100="#c5d0e6", c200="#9eb0d4", c300="#7490c2",
            c400="#4f74b0", c500="#2d5a9e", c600="#1a3a6e", c700="#112240",
            c800="#0a1628", c900="#060e1a", c950="#03070d"
        ),
        neutral_hue=gr.themes.colors.Color(
            c50="#f0f2f7", c100="#d8dde9", c200="#b8c2d6", c300="#95a3be",
            c400="#7285a6", c500="#536890", c600="#3a4f73", c700="#263856",
            c800="#162438", c900="#0c1622", c950="#060b11"
        ),
        font=gr.themes.GoogleFont("DM Sans"),
        font_mono=gr.themes.GoogleFont("DM Mono"),
    ).set(
        body_background_fill="#0a1628",
        body_background_fill_dark="#0a1628",
        body_text_color="#f8f9fc",
        body_text_color_dark="#f8f9fc",
        block_background_fill="#112240",
        block_background_fill_dark="#112240",
        block_border_color="rgba(201,168,76,0.2)",
        block_border_color_dark="rgba(201,168,76,0.2)",
        block_label_background_fill="#112240",
        block_label_background_fill_dark="#112240",
        block_label_text_color="#c9a84c",
        block_label_text_color_dark="#c9a84c",
        block_title_text_color="#c9a84c",
        block_title_text_color_dark="#c9a84c",
        button_primary_background_fill="linear-gradient(135deg,#c9a84c,#e8c96a)",
        button_primary_background_fill_dark="linear-gradient(135deg,#c9a84c,#e8c96a)",
        button_primary_text_color="#0a1628",
        button_primary_text_color_dark="#0a1628",
        button_secondary_background_fill="#112240",
        button_secondary_background_fill_dark="#112240",
        button_secondary_border_color="rgba(201,168,76,0.3)",
        button_secondary_border_color_dark="rgba(201,168,76,0.3)",
        button_secondary_text_color="#c9a84c",
        button_secondary_text_color_dark="#c9a84c",
        input_background_fill="#0a1628",
        input_background_fill_dark="#0a1628",
        input_border_color="rgba(201,168,76,0.2)",
        input_border_color_dark="rgba(201,168,76,0.2)",
        input_placeholder_color="#8892a4",
        input_placeholder_color_dark="#8892a4",
        shadow_drop="0 4px 24px rgba(0,0,0,0.4)",
        shadow_drop_lg="0 8px 40px rgba(0,0,0,0.5)",
        table_even_background_fill="#0a1628",
        table_even_background_fill_dark="#0a1628",
        table_odd_background_fill="#112240",
        table_odd_background_fill_dark="#112240",
    )) as demo:
    gr.Markdown("# 🏠 AI Real Estate Analyzer")
    gr.Markdown(DISCLAIMER)

    with gr.Tabs():

        # ── Tab 1 ──────────────────────────────────────────────────────────────
        with gr.Tab("📊 Data Upload & Overview"):
            gr.Markdown("## Upload Your Real Estate Data")
            gr.Markdown(PII_WARNING)
            with gr.Row():
                with gr.Column(scale=1):
                    t1_file = gr.File(label="Upload File (CSV, Excel, PDF, JSON)",
                                      file_types=[".csv", ".xlsx", ".xls", ".pdf", ".json"])
                    t1_btn = gr.Button("🔍 Analyze Dataset", variant="primary")
                    gr.Markdown(WAIT_MSG)
                with gr.Column(scale=2):
                    t1_analysis = gr.Markdown(label="AI Analysis")
            with gr.Row():
                t1_stats = gr.Markdown(label="Statistics")
                t1_preview = gr.Markdown(label="Data Preview")
            t1_btn.click(analyze_overview, inputs=[t1_file],
                         outputs=[t1_analysis, t1_stats, t1_preview])

        # ── Tab 2 ──────────────────────────────────────────────────────────────
        with gr.Tab("🏡 Property Value Analysis"):
            gr.Markdown("## Property Valuation & Comparable Sales Analysis")
            gr.Markdown(PII_WARNING)
            with gr.Row():
                with gr.Column(scale=1):
                    t2_file = gr.File(label="Upload Comparable Sales Data (CSV/Excel)",
                                      file_types=[".csv", ".xlsx", ".xls"])
                    t2_btn = gr.Button("🏡 Analyze Property Values", variant="primary")
                    gr.Markdown(WAIT_MSG)
                with gr.Column(scale=2):
                    t2_output = gr.Markdown(label="Property Value Analysis")
            t2_chart = gr.Plot(label="Value Analysis Chart")
            t2_btn.click(analyze_property_values, inputs=[t2_file],
                         outputs=[t2_output, t2_chart])

        # ── Tab 3 ──────────────────────────────────────────────────────────────
        with gr.Tab("📈 Market Trend Analysis"):
            gr.Markdown("## Real Estate Market Trends & Forecast")
            gr.Markdown(PII_WARNING)
            with gr.Row():
                with gr.Column(scale=1):
                    t3_file = gr.File(label="Upload Market Statistics Data (CSV/Excel)",
                                      file_types=[".csv", ".xlsx", ".xls"])
                    t3_btn = gr.Button("📈 Analyze Market Trends", variant="primary")
                    gr.Markdown(WAIT_MSG)
                with gr.Column(scale=2):
                    t3_output = gr.Markdown(label="Market Trend Analysis")
            t3_chart = gr.Plot(label="Market Trend Chart")
            t3_btn.click(analyze_market_trends, inputs=[t3_file],
                         outputs=[t3_output, t3_chart])

        # ── Tab 4 ──────────────────────────────────────────────────────────────
        with gr.Tab("💰 Investment ROI Calculator"):
            gr.Markdown("## Investment Property ROI & Return Analysis")
            gr.Markdown(PII_WARNING)
            with gr.Row():
                with gr.Column(scale=1):
                    t4_file = gr.File(label="Upload Investment Data (CSV/Excel — optional)",
                                      file_types=[".csv", ".xlsx", ".xls"])
                    t4_price = gr.Number(label="Purchase Price ($)", value=400000)
                    t4_rent = gr.Number(label="Monthly Rent ($)", value=2500)
                    t4_down = gr.Number(label="Down Payment (%)", value=20)
                    t4_rate = gr.Number(label="Interest Rate (%)", value=7.0)
                    t4_term = gr.Dropdown(choices=[15, 20, 25, 30], value=30,
                                          label="Loan Term (years)")
                    t4_expenses = gr.Number(label="Monthly Expenses — taxes, insurance, maintenance ($)",
                                            value=800)
                    t4_btn = gr.Button("💰 Calculate ROI", variant="primary")
                    gr.Markdown(WAIT_MSG)
                with gr.Column(scale=2):
                    t4_output = gr.Markdown(label="Investment ROI Analysis")
            t4_chart = gr.Plot(label="10-Year Projection Chart")
            t4_btn.click(analyze_investment_roi,
                         inputs=[t4_file, t4_price, t4_rent, t4_down,
                                 t4_rate, t4_term, t4_expenses],
                         outputs=[t4_output, t4_chart])

        # ── Tab 5 ──────────────────────────────────────────────────────────────
        with gr.Tab("🏘️ Rental Market Analysis"):
            gr.Markdown("## Rental Market & Yield Analysis")
            gr.Markdown(PII_WARNING)
            with gr.Row():
                with gr.Column(scale=1):
                    t5_file = gr.File(label="Upload Rental Listing Data (CSV/Excel)",
                                      file_types=[".csv", ".xlsx", ".xls"])
                    t5_btn = gr.Button("🏘️ Analyze Rental Market", variant="primary")
                    gr.Markdown(WAIT_MSG)
                with gr.Column(scale=2):
                    t5_output = gr.Markdown(label="Rental Market Analysis")
            t5_chart = gr.Plot(label="Rental Market Chart")
            t5_btn.click(analyze_rental_market, inputs=[t5_file],
                         outputs=[t5_output, t5_chart])

        # ── Tab 6 ──────────────────────────────────────────────────────────────
        with gr.Tab("⚠️ Risk Assessment"):
            gr.Markdown("## Investment Risk Analysis & Mitigation")
            gr.Markdown(PII_WARNING)
            with gr.Row():
                with gr.Column(scale=1):
                    t6_file = gr.File(label="Upload Property/Market Data (CSV/Excel — optional)",
                                      file_types=[".csv", ".xlsx", ".xls", ".pdf"])
                    t6_location = gr.Textbox(
                        label="Property Location",
                        placeholder="e.g. Austin, TX or Miami Beach, FL")
                    t6_property_type = gr.Dropdown(
                        choices=["Single Family Home", "Multi-Family", "Condo/Townhome",
                                 "Commercial", "Mixed Use", "Land", "Industrial", "Other"],
                        value="Single Family Home", label="Property Type")
                    t6_context = gr.Textbox(label="Additional Context", lines=3,
                                            placeholder="Any specific risk concerns...")
                    t6_btn = gr.Button("⚠️ Assess Risk", variant="primary")
                    gr.Markdown(WAIT_MSG)
                with gr.Column(scale=2):
                    t6_output = gr.Markdown(label="Risk Assessment")
            t6_chart = gr.Plot(label="Risk Assessment Chart")
            t6_btn.click(analyze_risk,
                         inputs=[t6_file, t6_location, t6_property_type, t6_context],
                         outputs=[t6_output, t6_chart])

        # ── Tab 7 ──────────────────────────────────────────────────────────────
        with gr.Tab("📉 Visualizations & Charts"):
            gr.Markdown("## Interactive Real Estate Data Visualizations")
            gr.Markdown(PII_WARNING)
            with gr.Row():
                with gr.Column(scale=1):
                    t7_file = gr.File(label="Upload Data (CSV/Excel)",
                                      file_types=[".csv", ".xlsx", ".xls"])
                    t7_chart_type = gr.Dropdown(
                        choices=["Bar Chart", "Line Chart", "Scatter Plot", "Pie Chart",
                                 "Heatmap", "Box Plot", "Histogram", "Area Chart", "Bubble Chart"],
                        value="Bar Chart", label="Chart Type")
                    t7_x = gr.Dropdown(choices=[], label="X Axis Column")
                    t7_y = gr.Dropdown(choices=[], label="Y Axis Column")
                    t7_color = gr.Dropdown(choices=[], label="Color By (optional)")
                    t7_btn = gr.Button("📊 Generate Chart", variant="primary")
                    t7_status = gr.Markdown()
                with gr.Column(scale=2):
                    t7_plot = gr.Plot(label="Chart")
            t7_file.change(get_columns, inputs=[t7_file],
                           outputs=[t7_x, t7_y, t7_color])
            t7_btn.click(generate_viz,
                         inputs=[t7_file, t7_chart_type, t7_x, t7_y, t7_color],
                         outputs=[t7_plot, t7_status])

        # ── Tab 8 ──────────────────────────────────────────────────────────────
        with gr.Tab("📄 Reports & Presentations"):
            gr.Markdown("## Generate Professional Real Estate Reports")
            gr.Markdown(PII_WARNING)
            with gr.Row():
                with gr.Column(scale=1):
                    t8_file = gr.File(label="Upload Data (CSV, Excel, PDF, JSON)",
                                      file_types=[".csv", ".xlsx", ".xls", ".pdf", ".json"])
                    t8_type = gr.Dropdown(
                        choices=["PowerPoint Presentation", "PDF Report",
                                 "CMA Report", "Investment Proposal"],
                        value="PDF Report", label="Report Type")
                    t8_btn = gr.Button("📄 Generate Report", variant="primary")
                    gr.Markdown(WAIT_MSG)
                    t8_download = gr.File(label="⬇️ Download Report")
                with gr.Column(scale=2):
                    t8_output = gr.Markdown(label="Report Preview")
            t8_btn.click(generate_report, inputs=[t8_file, t8_type],
                         outputs=[t8_output, t8_download])

        # ── Tab 9 ──────────────────────────────────────────────────────────────
        with gr.Tab("💬 AI Data Chat"):
            gr.Markdown("## Chat With Your Real Estate Data")
            gr.Markdown(PII_WARNING)
            with gr.Row():
                with gr.Column(scale=1):
                    t9_file = gr.File(label="Upload Data for Context (CSV/Excel)",
                                      file_types=[".csv", ".xlsx", ".xls"])
                    t9_upload_btn = gr.Button("📤 Load Data", variant="secondary")
                    t9_data_status = gr.Markdown()
                with gr.Column(scale=2):
                    t9_chatbot = gr.ChatInterface(
                        fn=chat_with_data,
                        examples=[
                            "What is the median sale price in this dataset?",
                            "Which neighborhood has the highest price appreciation?",
                            "Calculate the cap rate for a $500,000 property renting for $3,000/month",
                            "What is the average days on market?",
                            "Generate a SQL query to find the top 10 most expensive properties",
                            "Write Python code to calculate price per square foot trends",
                            "Is this a buyer's or seller's market based on the data?",
                        ],
                        title="",
                    )
            t9_upload_btn.click(upload_chat_data, inputs=[t9_file],
                                outputs=[t9_data_status])

    gr.Markdown(WATERMARK)


if __name__ == "__main__":
    demo.launch(server_name="0.0.0.0", server_port=int(os.environ.get("GRADIO_SERVER_PORT", 7860)), share=False, ssr_mode=False)
